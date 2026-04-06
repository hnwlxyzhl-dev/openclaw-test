#!/usr/bin/env python3
"""
SD PPT Fix: Rewrite slides 4, 5, 6 to resolve:
  - Problem 1: CLIP detail page (slide 5) — positional encoding & Transformer too brief
  - Problem 2: SD architecture (slide 4) — VAE encoder/decoder must be separated
  - Problem 3: VAE detail page (slide 6) — encoder vs decoder content混乱

Output: /home/admin/.openclaw/workspace-weaver/output/sd_fix_456.pptx (3 slides)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MAX_X = 13.333
MAX_Y = 7.5

# Colors
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_LT  = RGBColor(0xEB,0xF5,0xFB)
C_W   = RGBColor(0xFF,0xFF,0xFF)
C_D   = RGBColor(0x33,0x33,0x33)
C_G   = RGBColor(0x88,0x88,0x88)
C_GRN = RGBColor(0x27,0xAE,0x60)
C_ORG = RGBColor(0xE6,0x7E,0x22)
C_RED = RGBColor(0xE7,0x4C,0x3C)
C_PUR = RGBColor(0x8E,0x44,0xAD)
C_TEL = RGBColor(0x00,0x96,0x88)

FC = "微软雅黑"
FE = "Arial"
OUT = "/home/admin/.openclaw/workspace-weaver/output/sd_fix_456.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ── Helpers (same style as v3) ──
def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(fs); p.font.color.rgb=fc
    p.font.bold=b; p.font.name=fn; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld
        p.font.name=fn; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def chev(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.CHEVRON,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Stable Diffusion 技术解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    """Verify all shapes are within bounds."""
    violations = []
    for sh in s.shapes:
        r = sh.left + sh.width
        b = sh.top + sh.height
        if r > SLIDE_W + Emu(5000):  # 5px tolerance
            violations.append(f"  OVERFLOW X: right={r/914400:.3f} shape={sh.shape_type} text={sh.text_frame.text[:30] if sh.has_text_frame else ''}")
        if b > SLIDE_H + Emu(5000):
            violations.append(f"  OVERFLOW Y: bottom={b/914400:.3f} shape={sh.shape_type} text={sh.text_frame.text[:30] if sh.has_text_frame else ''}")
    if violations:
        print(f"⚠️  {msg} — {len(violations)} violations:")
        for v in violations:
            print(v)
    else:
        print(f"✅ {msg} — all shapes within bounds")

# ── Create presentation ──
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# SLIDE 4: SD 整体架构图 — VAE编码器/解码器分离
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "03  整体架构概览", "CLIP + U-Net + VAE 三大组件协同工作 — 编码器与解码器在流程中位置不同")

# ── 推理流程（上半部分）──
# 说明文字
tb(sl, Inches(0.5), Inches(1.2), Inches(12), Inches(0.25),
   "推理流程（从文字到图像）", fs=13, fc=C_PRI, b=True)

# Step 1: 文字输入
rrect(sl, Inches(0.3), Inches(1.55), Inches(1.4), Inches(0.8), fc=C_W, lc=C_ACC)
tb(sl, Inches(0.35), Inches(1.57), Inches(1.3), Inches(0.22),
   "① 文字输入", fs=10, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(1.82), Inches(1.3), Inches(0.5),
   '"一只橘猫\n坐在沙发上"', fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# Arrow 1
arrow_r(sl, Inches(1.75), Inches(1.8), Inches(0.3), Inches(0.2), C_ACC)

# Step 2: CLIP
rrect(sl, Inches(2.1), Inches(1.45), Inches(2.2), Inches(1.0), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_ACC)
tb(sl, Inches(2.15), Inches(1.47), Inches(2.1), Inches(0.22),
   "② CLIP 文本编码器", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
# CLIP header bar
rect(sl, Inches(2.1), Inches(1.45), Inches(2.2), Inches(0.22), fc=C_ACC)
tb(sl, Inches(2.15), Inches(1.47), Inches(2.1), Inches(0.22),
   "② CLIP 文本编码器", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(2.15), Inches(1.7), Inches(2.1), Inches(0.7),
   "123M参数\n12层Transformer\n输出: 77×768语义向量", fs=9, fc=C_D, ls=1.1)

# Arrow 2
arrow_r(sl, Inches(4.35), Inches(1.8), Inches(0.3), Inches(0.2), C_ACC)

# Step 3: U-Net (核心, 稍大)
rrect(sl, Inches(4.7), Inches(1.35), Inches(3.0), Inches(1.2), fc=RGBColor(0xFD,0xED,0xEC), lc=C_RED)
rect(sl, Inches(4.7), Inches(1.35), Inches(3.0), Inches(0.22), fc=C_RED)
tb(sl, Inches(4.75), Inches(1.37), Inches(2.9), Inches(0.22),
   "③ U-Net 去噪网络（核心）", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(4.75), Inches(1.6), Inches(2.9), Inches(0.9),
   "860M参数（占80%）\n输入: 噪声潜空间 + 文本向量\n⟳ 循环20-50步逐步去噪", fs=9, fc=C_D, ls=1.1)

# Cross-Attention 标注
rect(sl, Inches(4.3), Inches(1.55), Inches(0.45), Inches(0.2), fc=C_ACC)
tb(sl, Inches(4.3), Inches(1.56), Inches(0.45), Inches(0.18),
   "CA", fs=8, fc=C_W, b=True, a=PP_ALIGN.CENTER, fn=FE)

# Arrow 3
arrow_r(sl, Inches(7.75), Inches(1.8), Inches(0.3), Inches(0.2), C_RED)

# Step 4: VAE 解码器
rrect(sl, Inches(8.1), Inches(1.45), Inches(2.2), Inches(1.0), fc=RGBColor(0xE8,0xF8,0xF5), lc=C_GRN)
rect(sl, Inches(8.1), Inches(1.45), Inches(2.2), Inches(0.22), fc=C_GRN)
tb(sl, Inches(8.15), Inches(1.47), Inches(2.1), Inches(0.22),
   "④ VAE 解码器", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(8.15), Inches(1.7), Inches(2.1), Inches(0.7),
   "~49M参数\n64×64×4 → 512×512×3\n4层上采样还原高清", fs=9, fc=C_D, ls=1.1)

# Arrow 4
arrow_r(sl, Inches(10.35), Inches(1.8), Inches(0.3), Inches(0.2), C_GRN)

# Step 5: 输出图像
rrect(sl, Inches(10.7), Inches(1.55), Inches(1.4), Inches(0.8), fc=C_W, lc=C_D)
tb(sl, Inches(10.75), Inches(1.57), Inches(1.3), Inches(0.22),
   "⑤ 图像输出", fs=10, fc=C_D, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(10.75), Inches(1.82), Inches(1.3), Inches(0.5),
   "512×512×3\n高清 RGB 图像", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# 随机噪声输入（从下方进入U-Net）
rrect(sl, Inches(5.3), Inches(2.65), Inches(1.8), Inches(0.5), fc=C_W, lc=C_G)
tb(sl, Inches(5.35), Inches(2.67), Inches(1.7), Inches(0.45),
   "随机噪声 N(0,I)\n64×64×4 纯噪声", fs=9, fc=C_G, a=PP_ALIGN.CENTER, ls=1.1)
arrow_d(sl, Inches(6.05), Inches(2.55), Inches(0.2), Inches(0.15), C_G)

# ── 训练流程（中间）──
rrect(sl, Inches(0.3), Inches(3.3), Inches(6.3), Inches(1.5), fc=RGBColor(0xFE,0xF9,0xE7), lc=C_ORG)
tb(sl, Inches(0.5), Inches(3.33), Inches(5.8), Inches(0.22),
   "训练流程（VAE编码器在前面！）", fs=12, fc=C_ORG, b=True)
mtb(sl, Inches(0.5), Inches(3.6), Inches(5.8), Inches(1.1), [
    ("① 原始图像(512×512×3) → [VAE 编码器] → 潜空间(64×64×4)   ← 编码器先压缩", False, C_D, 10),
    ("② 对潜空间添加噪声（前向扩散）", False, C_D, 10),
    ("③ CLIP编码文本 + 带噪声潜空间 → U-Net → 预测噪声", False, C_D, 10),
    ("④ 计算损失 ||ε - ε_pred||²，反向传播更新U-Net参数", False, C_D, 10),
], ls=1.15)

# 推理流程关键标注
rrect(sl, Inches(6.8), Inches(3.3), Inches(6.2), Inches(1.5), fc=C_LT, lc=C_ACC)
tb(sl, Inches(7.0), Inches(3.33), Inches(5.8), Inches(0.22),
   "推理流程（VAE解码器在最后！）", fs=12, fc=C_ACC, b=True)
mtb(sl, Inches(7.0), Inches(3.6), Inches(5.8), Inches(1.1), [
    ("① 文本 → CLIP → 77×768向量（只执行一次）", False, C_D, 10),
    ("② 随机噪声(64×64×4) → U-Net逐步去噪(20-50步)", False, C_D, 10),
    ("③ 去噪完成 → [VAE 解码器] → 高清图像(512×512×3)  ← 解码器最后还原", False, C_D, 10),
    ("注意：推理时不需要VAE编码器！只用解码器还原图像", True, C_RED, 10),
], ls=1.15)

# ── 三个组件卡片（下半部分）──
dims = [
    ("1 CLIP 文本编码器", "翻译官 — 123M参数\n训练: 4亿张图文对\n12层Transformer, 768维\n作用: 把人类语言翻译成AI数字语言", C_ACC),
    ("2 U-Net 去噪网络", "核心画师 — 860M参数(80%)\nU型编码器-解码器结构\nCross-Attention + Skip Connections\n作用: 一步步去掉噪声", C_RED),
    ("3 VAE 编解码器", "编码器(34M): 图像→潜空间(训练时)\n解码器(49M): 潜空间→图像(推理时)\n48倍压缩: 786K→16K数值\n注意: 编码器和解码器位置不同!", C_GRN),
]
for i, (t, d, c) in enumerate(dims):
    x = Inches(0.3 + i * 4.3)
    rrect(sl, x, Inches(5.0), Inches(4.1), Inches(1.9), fc=C_W, lc=c)
    rect(sl, x, Inches(5.0), Inches(4.1), Inches(0.28), fc=c)
    tb(sl, x + Inches(0.1), Inches(5.03), Inches(3.9), Inches(0.24),
       t, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(5.35), Inches(3.9), Inches(1.5),
       d, fs=9.5, fc=C_D, ls=1.2)

chk(sl, "Slide 4 (architecture)")

# ============================================================
# SLIDE 5: CLIP 文本编码器详解 — 大幅增强位置编码和Transformer
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "04  CLIP 文本编码器详解", "Contrastive Language-Image Pre-training — 从文字到语义向量的完整旅程")

# ── 顶部: 四维度速览 ──
cd = [
    ("输入", "文本描述字符串\n经分词器最多77个token\n约50-60个汉字\n超出部分截断", C_ACC),
    ("输出", "77×768维张量\n共59,136个浮点数\n文本的语义指纹\n包含所有含义信息", C_GRN),
    ("参数", "约1.23亿(123M)\n12层Transformer\n隐藏维度768\n12个注意力头", C_ORG),
    ("作用", "人类文字→数字向量\n通过Cross-Attention注入U-Net\n指导去噪方向\n没有它AI不知道画什么", C_PUR),
]
for i, (t, d, c) in enumerate(cd):
    x = Inches(0.3 + i * 3.25)
    rrect(sl, x, Inches(1.2), Inches(3.05), Inches(1.2), fc=C_W, lc=c)
    rect(sl, x, Inches(1.2), Inches(3.05), Inches(0.25), fc=c)
    tb(sl, x + Inches(0.08), Inches(1.22), Inches(2.85), Inches(0.22),
       t, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(1.5), Inches(2.8), Inches(0.85),
       d, fs=9, fc=C_D, ls=1.15)

# ── CLIP 完整工作流程（核心区域）──
rrect(sl, Inches(0.3), Inches(2.55), Inches(12.7), Inches(4.45), fc=C_LT)
tb(sl, Inches(0.5), Inches(2.58), Inches(8), Inches(0.25),
   "CLIP 内部工作流程（以\"一只橘猫坐在沙发上\"为例）", fs=13, fc=C_PRI, b=True)

# ── Step 1: 分词 ──
chev(sl, Inches(0.5), Inches(2.95), Inches(1.1), Inches(0.3), C_ACC)
tb(sl, Inches(0.62), Inches(2.97), Inches(0.9), Inches(0.26),
   "① 分词", fs=10, fc=C_W, b=True)
tb(sl, Inches(0.5), Inches(3.3), Inches(1.2), Inches(1.0),
   "文本拆成最小语义单元(token)\n\n一只橘猫坐在沙发上\n\n→ 一只 / 橘 / 猫 /\n  坐在 / 沙发 / 上\n\n每个token映射成768维向量",
   fs=8.5, fc=C_D, ls=1.05)

arrow_r(sl, Inches(1.65), Inches(3.0), Inches(0.15), Inches(0.15), C_ACC)

# ── Step 2: 位置编码（大幅增强！）──
chev(sl, Inches(1.85), Inches(2.95), Inches(1.1), Inches(0.3), C_GRN)
tb(sl, Inches(1.97), Inches(2.97), Inches(0.9), Inches(0.26),
   "② 位置编码", fs=10, fc=C_W, b=True)

# 位置编码详解框
rrect(sl, Inches(1.85), Inches(3.3), Inches(3.4), Inches(3.5), fc=C_W, lc=C_GRN)
tb(sl, Inches(1.95), Inches(3.33), Inches(3.2), Inches(0.2),
   "为什么需要位置编码？", fs=10, fc=C_GRN, b=True)
mtb(sl, Inches(1.95), Inches(3.55), Inches(3.2), Inches(1.0), [
    ("\"猫坐在沙发上\" ≠ \"沙发坐在猫上\"", True, C_RED, 10),
    ("词一样，顺序不同，含义天差地别", False, C_D, 9),
    ("Transformer本身不区分顺序，必须额外注入位置信息", False, C_D, 9),
], ls=1.1)

tb(sl, Inches(1.95), Inches(4.5), Inches(3.2), Inches(0.2),
   "正弦位置编码原理", fs=10, fc=C_GRN, b=True)
mtb(sl, Inches(1.95), Inches(4.72), Inches(3.2), Inches(2.0), [
    ("对每个位置pos(0,1,2...)和维度i:", False, C_D, 9),
    ("  PE(pos,2i) = sin(pos / 10000^(2i/d))", False, C_D, 9),
    ("  PE(pos,2i+1) = cos(pos / 10000^(2i/d))", False, C_D, 9),
    ("", False, C_D, 4),
    ("• 不同频率的正弦/余弦波编码位置", False, C_D, 9),
    ("• 低频波表示宏观位置(开头/结尾)", False, C_D, 9),
    ("• 高频波表示微观位置(相邻词距离)", False, C_D, 9),
    ("• 不需要训练，天然支持任意长度", False, C_D, 9),
    ("• 最终: 含义向量 + 位置向量 = 完整表示", True, C_PUR, 9),
], ls=1.05)

arrow_r(sl, Inches(5.3), Inches(3.0), Inches(0.15), Inches(0.15), C_GRN)

# ── Step 3: Transformer（大幅增强！）──
chev(sl, Inches(5.5), Inches(2.95), Inches(1.1), Inches(0.3), C_ORG)
tb(sl, Inches(5.62), Inches(2.97), Inches(0.9), Inches(0.26),
   "③ Transformer", fs=10, fc=C_W, b=True)

# Transformer详解框
rrect(sl, Inches(5.5), Inches(3.3), Inches(3.5), Inches(3.5), fc=C_W, lc=C_ORG)
tb(sl, Inches(5.6), Inches(3.33), Inches(3.3), Inches(0.2),
   "自注意力机制 (Self-Attention)", fs=10, fc=C_ORG, b=True)
mtb(sl, Inches(5.6), Inches(3.55), Inches(3.3), Inches(1.2), [
    ("每个词都能\"看到\"其他所有词", True, C_RED, 10),
    ("", False, C_D, 3),
    ("\"橘猫\" → 注意到\"坐在\"(动作) + \"沙发\"(位置)", False, C_D, 9),
    ("\"坐在\" → 连接\"橘猫\"(主语) + \"沙发\"(宾语)", False, C_D, 9),
    ("模型理解: 这是一只猫在沙发上的场景", False, C_D, 9),
    ("", False, C_D, 3),
    ("Q*K^T/sqrt(d) → Softmax → *V = 注意力权重", False, C_PUR, 9),
    ("每个词获得其他所有词的加权信息", False, C_D, 9),
], ls=1.05)

tb(sl, Inches(5.6), Inches(4.9), Inches(3.3), Inches(0.2),
   "为什么需要12层？", fs=10, fc=C_ORG, b=True)
mtb(sl, Inches(5.6), Inches(5.12), Inches(3.3), Inches(1.6), [
    ("每层提取不同级别的语义信息:", False, C_D, 9),
    ("  第1-3层: 词性、语法关系(猫=名词, 坐=动词)", False, C_D, 9),
    ("  第4-7层: 局部语义(橘猫+沙发→室内场景)", False, C_D, 9),
    ("  第8-12层: 全局理解(完整场景+风格暗示)", False, C_D, 9),
    ("", False, C_D, 3),
    ("类比: 第1层认字→第5层理解句子→第12层理解整段", True, C_ACC, 9),
    ("每层还有FFN(前馈网络)进一步变换特征", False, C_D, 9),
], ls=1.05)

arrow_r(sl, Inches(9.05), Inches(3.0), Inches(0.15), Inches(0.15), C_ORG)

# ── Step 4: 输出向量 ──
chev(sl, Inches(9.25), Inches(2.95), Inches(1.1), Inches(0.3), C_PUR)
tb(sl, Inches(9.37), Inches(2.97), Inches(0.9), Inches(0.26),
   "④ 输出", fs=10, fc=C_W, b=True)

rrect(sl, Inches(9.25), Inches(3.3), Inches(3.5), Inches(3.5), fc=C_W, lc=C_PUR)
tb(sl, Inches(9.35), Inches(3.33), Inches(3.3), Inches(0.2),
   "输出: 77×768 语义向量", fs=10, fc=C_PUR, b=True)
mtb(sl, Inches(9.35), Inches(3.55), Inches(3.3), Inches(3.0), [
    ("每个token变成768维向量", False, C_D, 9),
    ("共77×768 = 59,136个浮点数", False, C_D, 9),
    ("", False, C_D, 3),
    ("不仅含每个词的独立含义:", False, C_D, 9),
    ("  \"橘\" → 橙色、毛茸茸、猫科", False, C_D, 9),
    ("  \"沙发\" → 家具、室内、坐具", False, C_D, 9),
    ("", False, C_D, 3),
    ("还含词与词之间的关系:", False, C_D, 9),
    ("  \"橘猫\"和\"沙发\"的共现关系", False, C_D, 9),
    ("  \"坐在\"连接主语和宾语", False, C_D, 9),
    ("", False, C_D, 3),
    ("这是对整段文本的完整语义编码", True, C_PUR, 9),
    ("后续通过Cross-Attention注入U-Net", True, C_ACC, 9),
], ls=1.05)

chk(sl, "Slide 5 (CLIP)")

# ============================================================
# SLIDE 6: VAE 编解码器详解 — 编码器和解码器清晰分离
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "05  VAE 编解码器详解", "Variational Autoencoder — 压缩48倍但肉眼几乎看不出差别")

# ── 顶部: 核心参数速览 ──
vd = [
    ("编码器输入", "原始图像 512×512×3\n786,432个像素值\n值域 0-255", C_ACC),
    ("编码器输出", "潜空间 64×64×4\n16,384个浮点数\n值域约[-5, 5]", C_GRN),
    ("总参数量", "约8300万(83M)\n编码器~34M 解码器~49M\n解码器更大(还原>压缩)", C_ORG),
    ("为什么用VAE?", "正则化确保潜空间连续平滑\n普通AE潜空间不规则\n扩散操作需要平滑空间", C_PUR),
]
for i, (t, d, c) in enumerate(vd):
    x = Inches(0.3 + i * 3.25)
    rrect(sl, x, Inches(1.2), Inches(3.05), Inches(1.1), fc=C_W, lc=c)
    rect(sl, x, Inches(1.2), Inches(3.05), Inches(0.22), fc=c)
    tb(sl, x + Inches(0.08), Inches(1.22), Inches(2.85), Inches(0.2),
       t, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(1.45), Inches(2.8), Inches(0.8),
       d, fs=9, fc=C_D, ls=1.1)

# ── 左半: 编码器 ──
rrect(sl, Inches(0.3), Inches(2.45), Inches(6.2), Inches(3.0), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_ACC)
tb(sl, Inches(0.5), Inches(2.48), Inches(5.8), Inches(0.22),
   "▼ 编码器 (Encoder): 4层下采样 — 图像→潜空间（训练时使用）", fs=11, fc=C_ACC, b=True)

# 编码器流程
enc_steps = [
    ("512×512×3", "输入图像（RGB三通道）", C_ACC),
    ("  ↓ 卷积 + 下采样 (stride=2)"),
    ("256×256×128", "通道从3→128，提取低级特征（边缘、色块）", C_ACC),
    ("  ↓ 卷积 + 下采样 (stride=2)"),
    ("128×128×256", "提取中级特征（纹理、局部结构）", C_ACC),
    ("  ↓ 卷积 + 下采样 (stride=2)"),
    ("64×64×512", "提取高级特征（物体部件、语义信息）", C_ACC),
    ("  ↓ 卷积 + 下采样 (stride=2)"),
    ("32×32×512", "最深层特征表示", C_ACC),
    ("  ↓ 特殊卷积调整维度"),
    ("64×64×4", "输出潜空间（48倍压缩！）", C_GRN),
]
y = Inches(2.78)
for item in enc_steps:
    if isinstance(item, tuple) and len(item) == 3:
        dim, desc, color = item
        rrect(sl, Inches(0.5), y, Inches(1.4), Inches(0.22), fc=color)
        tb(sl, Inches(0.55), y + Inches(0.01), Inches(1.3), Inches(0.2),
           dim, fs=8, fc=C_W, b=True, a=PP_ALIGN.CENTER, fn=FE)
        tb(sl, Inches(2.0), y + Inches(0.01), Inches(4.3), Inches(0.2),
           desc, fs=8.5, fc=C_D)
        y += Inches(0.22)
    else:
        tb(sl, Inches(0.6), y, Inches(5.8), Inches(0.15), item, fs=8, fc=C_G, a=PP_ALIGN.CENTER)
        y += Inches(0.15)

# 编码器参数标注
tb(sl, Inches(0.5), y + Inches(0.05), Inches(5.8), Inches(0.15),
   "编码器参数量: ~34M | 宽高每次缩小一半，通道数逐步增加", fs=8, fc=C_ACC, b=True)

# ── 右半: 解码器 ──
rrect(sl, Inches(6.8), Inches(2.45), Inches(6.2), Inches(3.0), fc=RGBColor(0xE8,0xF8,0xF5), lc=C_GRN)
tb(sl, Inches(7.0), Inches(2.48), Inches(5.8), Inches(0.22),
   "▲ 解码器 (Decoder): 4层上采样 — 潜空间→图像（推理时使用）", fs=11, fc=C_GRN, b=True)

# 解码器流程
dec_steps = [
    ("64×64×4", "输入潜空间", C_GRN),
    ("  ↑ 卷积 + 上采样 (ConvTranspose2d)"),
    ("128×128×512", "通道从4→512，开始重建空间细节", C_GRN),
    ("  ↑ 卷积 + 上采样"),
    ("256×256×256", "通道减半，分辨率翻倍", C_GRN),
    ("  ↑ 卷积 + 上采样"),
    ("512×512×128", "继续增加分辨率，减少通道", C_GRN),
    ("  ↑ 卷积 + 上采样"),
    ("512×512×64", "接近最终分辨率", C_GRN),
    ("  ↑ 最终卷积层 (输出RGB)"),
    ("512×512×3", "输出高清图像（肉眼几乎无损！）", C_ACC),
]
y = Inches(2.78)
for item in dec_steps:
    if isinstance(item, tuple) and len(item) == 3:
        dim, desc, color = item
        rrect(sl, Inches(7.0), y, Inches(1.4), Inches(0.22), fc=color)
        tb(sl, Inches(7.05), y + Inches(0.01), Inches(1.3), Inches(0.2),
           dim, fs=8, fc=C_W, b=True, a=PP_ALIGN.CENTER, fn=FE)
        tb(sl, Inches(8.5), y + Inches(0.01), Inches(4.3), Inches(0.2),
           desc, fs=8.5, fc=C_D)
        y += Inches(0.22)
    else:
        tb(sl, Inches(7.1), y, Inches(5.8), Inches(0.15), item, fs=8, fc=C_G, a=PP_ALIGN.CENTER)
        y += Inches(0.15)

# 解码器参数标注
tb(sl, Inches(7.0), y + Inches(0.05), Inches(5.8), Inches(0.15),
   "解码器参数量: ~49M（比编码器大，因为\"还原\"比\"压缩\"更复杂）", fs=8, fc=C_GRN, b=True)

# ── 底部: 编码器 vs 解码器对比 + VAE vs AE ──
rrect(sl, Inches(0.3), Inches(5.55), Inches(6.2), Inches(1.45), fc=RGBColor(0xFE,0xF9,0xE7), lc=C_ORG)
tb(sl, Inches(0.5), Inches(5.58), Inches(5.8), Inches(0.2),
   "编码器 vs 解码器对比", fs=11, fc=C_ORG, b=True)
mtb(sl, Inches(0.5), Inches(5.82), Inches(5.8), Inches(1.1), [
    ("          编码器          |          解码器", True, C_D, 9),
    ("方向:    512→64 下采样    |    64→512 上采样", False, C_D, 9),
    ("参数:    ~34M            |    ~49M (更大)", False, C_D, 9),
    ("操作:    卷积+stride=2    |    转置卷积上采样", False, C_D, 9),
    ("通道:  3→128→256→512→4 |  4→512→256→128→3", False, C_D, 9),
    ("使用:    训练时(压缩图像)  |    推理时(还原图像)", False, C_D, 9),
], ls=1.05)

rrect(sl, Inches(6.8), Inches(5.55), Inches(6.2), Inches(1.45), fc=RGBColor(0xF5,0xEE,0xF8), lc=C_PUR)
tb(sl, Inches(7.0), Inches(5.58), Inches(5.8), Inches(0.2),
   "为什么用VAE而不是普通自编码器(AE)？", fs=11, fc=C_PUR, b=True)
mtb(sl, Inches(7.0), Inches(5.82), Inches(5.8), Inches(1.1), [
    ("普通AE的问题: 潜空间\"不规则\"", True, C_RED, 9),
    ("• 两个相近的潜空间点可能映射到完全不相关的图像", False, C_D, 9),
    ("• 潜空间中存在大量\"空白区域\"(没有对应图像)", False, C_D, 9),
    ("• 在不规则空间做扩散(加噪声/去噪声)数学上不成立", False, C_D, 9),
    ("", False, C_D, 3),
    ("VAE的解决方案: KL散度正则化", True, C_GRN, 9),
    ("• 约束潜空间服从标准正态分布，确保连续平滑", False, C_D, 9),
    ("• 任意潜空间点都能解码成有意义的图像", False, C_D, 9),
], ls=1.05)

chk(sl, "Slide 6 (VAE)")

# ── Save ──
prs.save(OUT)
print(f"\n✅ Saved: {OUT} | {len(prs.slides)} slides")

# ── Final boundary check ──
print("\n═══ Final Boundary Verification ═══")
all_ok = True
for idx, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        r_val = (sh.left + sh.width) / 914400
        b_val = (sh.top + sh.height) / 914400
        if r_val > MAX_X + 0.01:
            print(f"❌ Slide {idx+1}: X overflow {r_val:.3f} — {sh.text_frame.text[:40] if sh.has_text_frame else sh.shape_type}")
            all_ok = False
        if b_val > MAX_Y + 0.01:
            print(f"❌ Slide {idx+1}: Y overflow {b_val:.3f} — {sh.text_frame.text[:40] if sh.has_text_frame else sh.shape_type}")
            all_ok = False
if all_ok:
    print("✅ All shapes within bounds across all 3 slides!")
