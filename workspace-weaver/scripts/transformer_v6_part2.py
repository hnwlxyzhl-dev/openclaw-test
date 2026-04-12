#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT v6 Part2 - Pages 8-13 (全中文版)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# === Constants ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xED, 0xC6)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)

FONT = "Arial"

def add_rect(slide, left, top, width, height, fill_color, text="", font_size=Pt(9), font_color=C_TEXT, bold=False):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(0.5)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_line(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5)):
    """Add a connector line."""
    c = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
    c.line.color.rgb = color
    c.line.width = width
    return c

def add_tb(slide, left, top, width, height, text, font_size=Pt(10), font_color=C_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Add text box."""
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return txBox

def add_title(slide, text):
    return add_tb(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5),
                  text, Pt(16), C_TITLE, bold=True)

def add_points(slide, left, top, width, height, points, tsz=Pt(10), bsz=Pt(9)):
    """Add bullet points. points = [(title, body), ...]"""
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (title, body) in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = title
        r.font.size = tsz
        r.font.bold = True
        r.font.color.rgb = C_TITLE
        r.font.name = FONT
        p.space_before = Pt(3) if i > 0 else Pt(0)
        p.space_after = Pt(1)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = bsz
        r2.font.color.rgb = C_TEXT
        r2.font.name = FONT
        p2.space_before = Pt(0)
        p2.space_after = Pt(2)
    return txBox

def count_cn(t):
    return len(re.findall(r'[\u4e00-\u9fff]', t))

def count_en(t):
    return len(re.findall(r'[a-zA-Z]', t))

# === Create Presentation ===
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ============================================================
# 第8页：训练细节
# ============================================================
s8 = prs.slides.add_slide(blank_layout)
add_title(s8, "训练细节：安全网、掩码与损失函数")

# 因果掩码矩阵
ml = Inches(0.4)
mt = Inches(0.9)
cs = Inches(0.55)
labels = ["[开始]", "我", "爱", "你"]
for j, lb in enumerate(labels):
    add_tb(s8, ml + cs*(j+1), mt - Inches(0.22), cs, Inches(0.22), lb, Pt(9), C_TEXT, align=PP_ALIGN.CENTER)
for i, lb in enumerate(labels):
    add_tb(s8, ml, mt + cs*i, cs, cs, lb, Pt(9), C_TEXT, align=PP_ALIGN.RIGHT)
for i in range(4):
    for j in range(4):
        cx = ml + cs*(j+1)
        cy = mt + cs*i
        if i >= j:
            add_rect(s8, cx, cy, cs-Inches(0.04), cs-Inches(0.04), C_ENCODER, "可以看", Pt(8), C_WHITE, True)
        else:
            add_rect(s8, cx, cy, cs-Inches(0.04), cs-Inches(0.04), C_DECODER, "-∞", Pt(8), C_WHITE, True)
add_tb(s8, ml, mt + cs*4 + Inches(0.05), Inches(3.5), Inches(0.3),
       "因果掩码矩阵：蓝色=可见，红色=被遮蔽（Softmax后为0）", Pt(8), C_GRAY)

# 损失计算流程
ft = Inches(3.6)
fl = Inches(0.4)
bw = Inches(1.15)
bh = Inches(0.42)
fg = Inches(0.12)
flow = [
    ("512维输出", C_LIGHT_BLUE, C_TEXT),
    ("线性变换\n512×37000", C_LIGHT_ORANGE, C_TEXT),
    ("Softmax", C_LIGHT_GREEN, C_TEXT),
    ("预测概率分布", C_LIGHT_ORANGE, C_DECODER),
]
for idx, (txt, bg, fc) in enumerate(flow):
    x = fl + idx*(bw + fg)
    add_rect(s8, x, ft, bw, bh, bg, txt, Pt(8), fc, True)
    if idx < len(flow)-1:
        add_line(s8, x+bw, ft+bh/2, x+bw+fg, ft+bh/2, C_GRAY)

y2 = ft + bh + Inches(0.2)
add_rect(s8, fl, y2, bw, bh, C_LIGHT_RED, "真实目标词", Pt(8), C_DECODER, True)
add_tb(s8, fl+bw+Inches(0.02), y2, Inches(0.4), bh, "↔", Pt(11), C_DECODER, True, PP_ALIGN.CENTER)
add_rect(s8, fl+bw+Inches(0.5), y2, Inches(2.2), bh, C_LIGHT_RED, "交叉熵损失=−log(p正确词)", Pt(8), C_DECODER, True)

y3 = y2 + bh + Inches(0.15)
add_rect(s8, fl, y3, Inches(2.2), bh, C_LIGHT_GRAY, "反向传播（∂Loss/∂θ）", Pt(8), C_GRAY, True)
add_line(s8, fl+Inches(1.1), y2+bh, fl+Inches(1.1), y3, C_GRAY)
add_rect(s8, fl+Inches(2.4), y3, Inches(2.2), bh, C_LIGHT_GRAY, "更新所有6500万参数", Pt(8), C_GRAY, True)
add_line(s8, fl+Inches(2.2), y3+bh/2, fl+Inches(2.4), y3+bh/2, C_GRAY)

# 右侧4个要点
pts8 = [
    ("① 残差连接的数学原理：",
     "残差连接的公式为输出等于层归一化后x加上子层处理x的结果，其中x是层的输入。从梯度角度看，总梯度等于输出对损失的梯度乘以子层梯度加一。加号后面的常数一意味着无论子层的梯度多小，总梯度中至少有一个一的贡献，梯度永远不会消失。这就是残差网络能训练超过一百层深度网络的核心原因。编码器和解码器的每个子层都配备残差连接，确保深度堆叠时训练过程稳定。"),
    ("② 因果掩码杜绝作弊：",
     "因果掩码是一个下三角矩阵，对角线及以下位置为零（可以看），对角线以上位置为负无穷（不能看）。分数矩阵加上掩码后，被禁止位置的分数变成负无穷，经过Softmax运算后指数负无穷等于零，完全不可见。例如处理\"爱\"字时（第3行），只能看到开始、我、爱三个词，\"你\"的注意力权重为零。训练和推理都必须执行这一操作——保证行为一致，避免暴露偏差。如果训练时偷看答案，推理时没有答案可抄就会崩溃。"),
    ("③ 交叉熵损失衡量预测质量：",
     "交叉熵公式为损失等于负的对数乘以正确词的预测概率。如果模型对正确词的预测概率为0.99，损失约等于0.01（很好）；如果概率为0.01，损失约等于4.6（很差）。对于整个序列，总损失等于所有位置损失的平均值。训练初期损失约等于10.5（随机猜测水平），充分训练后降到1.5到3.0之间。损失越小，模型预测越准确。"),
    ("④ 训练配置与优化：",
     "原始Transformer的训练配置为——WMT 2014英德翻译数据集，450万句对，批量大小25000个词元，Adam优化器(β₁=0.9,β₂=0.98)，学习率峰值0.0003（前4000步线性增温后按步数的−0.5次方衰减），标签平滑ε=0.1（防止模型对正确答案过于自信），dropout率0.3（训练时随机丢弃30%的神经元防止过拟合），8块P100 GPU训练12小时（30万步），最终BLEU得分28.4。"),
]
add_points(s8, Inches(8.4), Inches(0.75), Inches(4.7), Inches(6.4), pts8, Pt(10), Pt(9))

# ============================================================
# 第9页：推理阶段+KV缓存
# ============================================================
s9 = prs.slides.add_slide(blank_layout)
add_title(s9, "推理阶段：闭卷考试与KV缓存")

add_rect(s9, Inches(0.4), Inches(0.85), Inches(2.8), Inches(0.38),
         C_LIGHT_BLUE, "编码器只运行一次→输出缓存", Pt(9), C_ENCODER, True)

st = Inches(1.5)
sl = Inches(0.4)
sh = Inches(0.45)
sw = Inches(7.2)
steps = [
    ("步骤1: [开始]→解码器→\"我\"(概率0.85)", C_LIGHT_BLUE),
    ("步骤2: [开始,我]→解码器→\"爱\"(概率0.78)", C_LIGHT_BLUE),
    ("步骤3: [开始,我,爱]→解码器→\"你\"(概率0.91)", C_LIGHT_BLUE),
    ("步骤4: [开始,我,爱,你]→解码器→\"结束\"(概率0.96)", C_LIGHT_BLUE),
]
for idx, (txt, bg) in enumerate(steps):
    y = st + idx*(sh+Inches(0.12))
    add_rect(s9, sl, y, sw, sh, bg, txt, Pt(9), C_TEXT, True)
    if idx < len(steps)-1:
        add_line(s9, sl+sw/2, y+sh, sl+sw/2, y+sh+Inches(0.12), C_GRAY)

# 无缓存 vs 有缓存
ct = Inches(4.0)
cw = Inches(3.3)
add_rect(s9, Inches(0.4), ct, cw, Inches(0.35), C_LIGHT_RED, "无缓存（慢）", Pt(10), C_DECODER, True)
nc_items = [
    ("步骤2: 重新计算[开始]的K,V", C_TEXT),
    ("→ 红色标注：\"重复计算！\"", C_DECODER),
    ("步骤3: 又重新计算全部K,V", C_TEXT),
    ("→ 红色标注：\"又重复！\"", C_DECODER),
]
for i, (item, c) in enumerate(nc_items):
    add_tb(s9, Inches(0.5), ct+Inches(0.4)+i*Inches(0.4), cw-Inches(0.2), Inches(0.35), item, Pt(8), c)

cl = Inches(0.4) + cw + Inches(0.3)
add_rect(s9, cl, ct, cw, Inches(0.35), C_LIGHT_GREEN, "有缓存（快）", Pt(10), C_FFN, True)
wc_items = [
    ("步骤2: 只计算[我]的新K,V", C_TEXT),
    ("→ 绿色标注：\"只算新词！\"", C_FFN),
    ("步骤3: 只算[爱]的新K,V", C_TEXT),
    ("→ 绿色标注：\"只算新词！\"", C_FFN),
]
for i, (item, c) in enumerate(wc_items):
    add_tb(s9, cl+Inches(0.1), ct+Inches(0.4)+i*Inches(0.4), cw-Inches(0.2), Inches(0.35), item, Pt(8), c)

add_rect(s9, cl+Inches(0.2), ct+Inches(2.05), Inches(2.8), Inches(0.55),
         C_LIGHT_PURPLE, "缓存:{K:[k_开始,k_我,k_爱],\nV:[v_开始,v_我,v_爱]}", Pt(8), C_EMBED, True)

pts9 = [
    ("① 推理=闭卷考试：",
     "与训练的\"开卷考试\"不同，推理阶段模型必须完全靠自己之前生成的词来预测下一个词。第一步只有[开始]，凭空生成第一个词\"我\"；然后将\"我\"加入输入序列，生成\"爱\"；再加入\"爱\"，生成\"你\"；最后生成\"结束\"标记停止。每一步都依赖前一步的正确性——如果第一步生成了\"他\"而不是\"我\"，后续的预测可能完全不同。这就像闭卷考试，第一题做错了可能影响后面所有题。"),
    ("② 编码器只运行一次：",
     "编码器处理输入后，输出6层编码的向量表示被完整缓存（K和V向量），在整个生成过程中反复使用。解码器的每一层在计算交叉注意力时，都会查阅这份缓存，无需重新计算编码器的输出。这使得编码器的计算成本只发生一次，与生成序列的长度无关。"),
    ("③ 串行瓶颈与KV缓存：",
     "与训练的完全并行不同，推理时解码器必须逐词串行生成。没有缓存时，每步重算全部词的QKV，总计算量O(T²)，生成长度翻倍计算量增加4倍。KV缓存的核心思想是：每步只计算新词的K和V追加到缓存，新词的Q与缓存中所有K做点积。每步O(d×T)而非O(T×d×T)，就像开会只需查之前纪要，只讨论新增内容。"),
    ("④ KV缓存的内存代价：",
     "每步缓存所有已生成词的K和V向量。对于12层、64头、128维的模型，生成100个词的缓存约：12层×100词×2(K+V)×64头×128维×2字节≈375MB。GPT-4等大模型生成500词时KV缓存可达数GB，成为推理的内存瓶颈。"),
]
add_points(s9, Inches(8.4), Inches(0.75), Inches(4.7), Inches(6.4), pts9, Pt(10), Pt(9))

# ============================================================
# 第10页：词选择策略
# ============================================================
s10 = prs.slides.add_slide(blank_layout)
add_title(s10, "推理·词选择策略：从贪心到创意")

cl2 = Inches(0.4)
ct2 = Inches(0.85)
cw2 = Inches(7.0)
ch2 = Inches(0.82)
cg = Inches(0.1)

strats = [
    ("① 贪心搜索（灰色）：每步选概率最高的词。确定性，无随机性。缺点：可能错过全局最优，容易生成重复内容如\"今天天气今天天气\"",
     C_LIGHT_GRAY, C_GRAY),
    ("② 束搜索（蓝色）：同时维护k条最佳路径。beam_size=4时每步保留4个候选。翻译效果好但结果偏平庸，缺乏多样性",
     C_LIGHT_BLUE, C_ENCODER),
    ("③ Top-K采样（绿色）：从概率最高的K个词中随机采样。K=50在质量和多样性间取得平衡",
     C_LIGHT_GREEN, C_FFN),
    ("④ Top-P核采样（橙色）：从累计概率达P的最小词集中采样。P=0.9时动态调整候选集大小，比Top-K更灵活",
     C_LIGHT_ORANGE, C_ATTENTION),
    ("⑤ 温度缩放（红色）：Softmax前对分数缩放。T<1分布尖锐更确定，T>1分布平坦更随机",
     C_LIGHT_RED, C_DECODER),
]
for idx, (txt, bg, tc) in enumerate(strats):
    y = ct2 + idx*(ch2+cg)
    add_rect(s10, cl2, y, cw2, ch2, bg)
    colon_idx = txt.find('：')
    if colon_idx > 0:
        txBox = s10.shapes.add_textbox(cl2+Inches(0.15), y+Inches(0.05), cw2-Inches(0.3), ch2-Inches(0.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = txt[:colon_idx+1]
        r1.font.size = Pt(9)
        r1.font.bold = True
        r1.font.color.rgb = tc
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = txt[colon_idx+1:]
        r2.font.size = Pt(9)
        r2.font.color.rgb = C_TEXT
        r2.font.name = FONT

# 温度分布对比
dl = Inches(7.8)
dt = Inches(0.9)
add_tb(s10, dl, dt, Inches(5), Inches(0.3), "概率分布对比：", Pt(10), C_TITLE, True)
dists = [
    ("低温 T=0.3（尖锐峰值）", C_DECODER, "[|||||||||||...........]"),
    ("中温 T=0.7（正常分布）", C_ATTENTION, "[|||||......|||||....]"),
    ("高温 T=1.5（平坦均匀）", C_ENCODER, "[..||..||..||..||..||.]"),
]
for i, (title, color, bar) in enumerate(dists):
    y = dt + Inches(0.35) + i*Inches(0.9)
    add_tb(s10, dl, y, Inches(3), Inches(0.22), title, Pt(9), color, True)
    add_tb(s10, dl, y+Inches(0.22), Inches(5), Inches(0.3), bar, Pt(9), C_TEXT)

pts10 = [
    ("① 实际组合策略：",
     "当前主流大语言模型通常使用温度系数0.7加上核采样概率阈值0.9的组合策略，同时叠加重复惩罚机制来避免生成内容陷入循环重复。不同应用场景推荐不同的温度设定：编写代码时建议温度设为0.2（需要精确和确定性），创意写作和故事生成时温度设为0.7到1.0（需要多样性和想象力），日常对话时温度设为0.7（在准确性和自然流畅之间取得平衡）。这种组合策略在实践中表现最为稳定可靠。"),
    ("② 为什么不能总选最高概率词：",
     "贪心搜索看似最优（每步选最好的），但可能错过全局最优序列。例如\"我喜欢在公园里___\"，贪心选\"散步\"(p=0.3)，但\"跑步\"(p=0.25)+\"之后喝杯咖啡\"(p=0.4)的整体序列概率更高。在开放域对话中，\"最优\"本身就不存在——同一个问题可以有无数种合理的回答。我们需要的是多样性和创造性，而非单一最优解。贪心搜索生成的文本往往机械、乏味、缺乏个性。采样策略引入了受控的随机性，让模型每次都能产出不同的高质量内容，这正是对话和创作任务所必需的特质。"),
]
add_points(s10, Inches(7.8), Inches(3.8), Inches(5.2), Inches(3.5), pts10, Pt(10), Pt(9))

# ============================================================
# 第11页：训练vs推理对比
# ============================================================
s11 = prs.slides.add_slide(blank_layout)
add_title(s11, "训练vs推理：开卷考试与闭卷考试")

tl = Inches(0.4)
tt = Inches(0.85)
cws = [Inches(2.2), Inches(5.0), Inches(5.0)]
rh = Inches(0.38)
hh = Inches(0.38)

headers = ["维度", "训练（开卷考试）", "推理（闭卷考试）"]
rows = [
    ["解码器输入", "完整目标序列（教师强制策略）", "模型自己逐词生成的序列"],
    ["计算方式", "所有位置完全并行（一次前向传播）", "逐词串行（T次前向传播）"],
    ["损失函数", "交叉熵损失（每个位置计算并取平均）", "无损失计算，只做预测"],
    ["梯度计算", "反向传播更新所有6500万参数", "无梯度计算（参数完全冻结）"],
    ["KV缓存", "不需要（并行计算无需缓存中间结果）", "必须使用（节省重复计算，降低延迟）"],
    ["GPU利用率", "高（矩阵运算充分并行，GPU跑满）", "低（串行生成，大量时间在等待上一步）"],
    ["速度", "快（处理一个批次只需一次前向+反向）", "慢（生成长度为T需要T次前向传播）"],
    ["核心目标", "学习参数（学会预测下一个词的模式）", "生成结果（产出可用的文本内容）"],
]

for j, h in enumerate(headers):
    x = tl + sum(cws[:j])
    w = cws[j]
    bg = C_ENCODER if j == 0 else (C_ENCODER if j == 1 else C_FFN)
    add_rect(s11, x, tt, w-Inches(0.04), hh, bg, h, Pt(10), C_WHITE, True)

for i, row in enumerate(rows):
    y = tt + hh + Inches(0.04) + i*(rh+Inches(0.03))
    for j, cell in enumerate(row):
        x = tl + sum(cws[:j])
        w = cws[j]
        if j == 0:
            bg, fc, bld = C_LIGHT_BLUE, C_ENCODER, True
        elif j == 1:
            bg, fc, bld = C_LIGHT_BLUE, C_TEXT, False
        else:
            bg, fc, bld = C_LIGHT_GREEN, C_TEXT, False
        add_rect(s11, x, y, w-Inches(0.04), rh, bg, cell, Pt(9), fc, bld)

st2 = tt + hh + Inches(0.04) + 8*(rh+Inches(0.03)) + Inches(0.1)
summary = ("比喻总结：训练就像新员工培训期——师傅手把手教，每一步都有标准答案参考（教师强制），"
           "所有题目同时做（并行），效率很高但需要大量资源（算力、数据、时间）。"
           "推理就像正式上岗——全靠自己，每一步都参考之前的结果（自回归），必须一步步来不能跳步（串行），"
           "虽然只需一台电脑但速度远慢于训练。两者的核心矛盾是\"暴露偏差\"：培训时看标准答案，"
           "上岗时看自己的输出，场景不匹配。KV缓存和Scheduled Sampling等技术都是为了缩小这个差距，"
           "让推理表现尽可能接近训练水平。")
add_tb(s11, tl, st2, Inches(12.5), Inches(1.6), summary, Pt(10), C_TEXT)

# ============================================================
# 第12页：应用与总结
# ============================================================
s12 = prs.slides.add_slide(blank_layout)
add_title(s12, "从翻译工具到AI通用引擎")

rl = Inches(0.4)
rt = Inches(0.85)

# 根节点
add_rect(s12, rl+Inches(1.5), rt, Inches(3), Inches(0.45), C_TITLE, "Transformer 2017", Pt(12), C_WHITE, True)

by = rt + Inches(0.65)
bh = Inches(0.35)

# 左分支：仅解码器
lx = rl
add_rect(s12, lx, by, Inches(2.5), bh, C_ENCODER, "仅解码器", Pt(10), C_WHITE, True)
add_line(s12, lx+Inches(1.25), rt+Inches(0.45), lx+Inches(1.25), by, C_ENCODER)
gpts = ["GPT-1", "GPT-2", "GPT-3", "GPT-4", "对话模型"]
for i, m in enumerate(gpts):
    my = by+bh+Inches(0.08)+i*Inches(0.3)
    add_rect(s12, lx, my, Inches(2.5), Inches(0.26), C_LIGHT_BLUE, m, Pt(8), C_ENCODER, True)
add_tb(s12, lx, by+bh+Inches(0.08)+5*Inches(0.3), Inches(2.5), Inches(0.22), "生成之王", Pt(9), C_ENCODER, True, PP_ALIGN.CENTER)

# 中分支：完整架构
mx = lx + Inches(2.8)
add_rect(s12, mx, by, Inches(2.5), bh, C_FFN, "完整架构", Pt(10), C_WHITE, True)
add_line(s12, mx+Inches(1.25), rt+Inches(0.45), mx+Inches(1.25), by, C_FFN)
fmods = ["T5", "BART", "Whisper"]
for i, m in enumerate(fmods):
    my = by+bh+Inches(0.08)+i*Inches(0.3)
    add_rect(s12, mx, my, Inches(2.5), Inches(0.26), C_LIGHT_GREEN, m, Pt(8), C_FFN, True)
add_tb(s12, mx, by+bh+Inches(0.08)+3*Inches(0.3), Inches(2.5), Inches(0.22), "转换之王", Pt(9), C_FFN, True, PP_ALIGN.CENTER)

# 右分支：仅编码器
rx = mx + Inches(2.8)
add_rect(s12, rx, by, Inches(2.5), bh, C_EMBED, "仅编码器", Pt(10), C_WHITE, True)
add_line(s12, rx+Inches(1.25), rt+Inches(0.45), rx+Inches(1.25), by, C_EMBED)
emods = ["基础版", "增强版", "进阶版"]
for i, m in enumerate(emods):
    my = by+bh+Inches(0.08)+i*Inches(0.3)
    add_rect(s12, rx, my, Inches(2.5), Inches(0.26), C_LIGHT_PURPLE, m, Pt(8), C_EMBED, True)
add_tb(s12, rx, by+bh+Inches(0.08)+3*Inches(0.3), Inches(2.5), Inches(0.22), "理解之王", Pt(9), C_EMBED, True, PP_ALIGN.CENTER)

# 参数量条
add_tb(s12, rl, Inches(3.9), Inches(8), Inches(0.25),
       "参数量递增：六千五百万→一亿一千万→十五亿→一千七百五十亿→一万八千亿", Pt(9), C_TEXT, True)
add_tb(s12, rl, Inches(4.15), Inches(8), Inches(0.25),
       "时间线：2017原始论文→2020视觉模型→2022语音模型→2023多模态→2024视频模型", Pt(9), C_GRAY, True)

pts12 = [
    ("① GPT系列（仅解码器）——生成之王：",
     "GPT系列只保留解码器部分，去掉了编码器和交叉注意力。核心能力是\"预测下一个词\"。主流对话模型的三步训练流程：预训练（在海量文本上预测下一个词，训练了三千亿词元）→指令微调（在人工标注的指令回复对上微调，学会理解人类指令）→人类反馈强化学习（用人类偏好训练奖励模型，再用策略优化算法改进生成策略）。代表模型包括千亿参数和万亿参数级别的大语言模型。"),
    ("② BERT系列（仅编码器）——理解之王：",
     "核心方法是\"掩码语言建模\"——随机遮盖输入中百分之十五的词元，让模型根据上下文预测被遮盖的内容。不能直接生成文本，但极其擅长文本理解任务：分类、问答、命名实体识别、语义相似度等。搜索引擎从2019年开始使用该技术来理解用户查询意图。代表模型从小型到大型多个版本。"),
    ("③ T5系列（完整架构）——转换之王：",
     "核心思想是将所有自然语言处理任务统一为\"文本到文本\"格式。翻译任务、摘要任务、分类任务都可以用统一格式处理。翻译服务基于该架构的变体实现。代表模型从六千万到一百一十亿参数不等。"),
    ("④ 超越语言——征服所有模态：",
     "视觉模型（2020年）将图像切成小块当\"词\"处理，达到与卷积神经网络相当的效果。语音识别模型（2022年）用编码器解码器做语音识别，支持99种语言。多模态模型（2023年）实现视觉语言多模态理解。蛋白质折叠模型用注意力机制预测蛋白质三维结构。只要能表示为序列，就能处理。"),
    ("⑤ 核心要点：",
     "本质是让信息自由流动——自注意力使任意元素直接交流，彻底解决循环神经网络的长距离依赖、串行计算和梯度消失三大问题。编码器负责理解，解码器负责生成。训练是开卷考试（并行高效），推理是闭卷考试（串行，缓存加速）。三大家族（生成、理解、转换）共享同一基因却走上不同道路。从六千五百万到万亿参数，架构从未改变——不是一个模型，而是人工智能时代的基石。"),
]
add_points(s12, Inches(8.6), Inches(0.75), Inches(4.5), Inches(6.4), pts12, Pt(10), Pt(8))

# ============================================================
# 第13页：感谢页
# ============================================================
s13 = prs.slides.add_slide(blank_layout)

txBox = s13.shapes.add_textbox(Inches(0), Inches(2.2), SLIDE_W, Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "感谢聆听"
r.font.size = Pt(36)
r.font.bold = True
r.font.color.rgb = C_TITLE
r.font.name = FONT
p.alignment = PP_ALIGN.CENTER

txBox2 = s13.shapes.add_textbox(Inches(0), Inches(3.5), SLIDE_W, Inches(0.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
r2 = p2.add_run()
r2.text = "Transformer: 不是一个模型，而是一个时代"
r2.font.size = Pt(20)
r2.font.color.rgb = C_TEXT
r2.font.name = FONT
p2.alignment = PP_ALIGN.CENTER

txBox3 = s13.shapes.add_textbox(Inches(0), Inches(5.5), SLIDE_W, Inches(0.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
r3 = p3.add_run()
r3.text = "Attention Is All You Need — Vaswani et al., 2017"
r3.font.size = Pt(12)
r3.font.color.rgb = C_GRAY
r3.font.name = FONT
p3.alignment = PP_ALIGN.CENTER

# ============================================================
# Save & Stats
# ============================================================
out = "/home/admin/.openclaw/workspace-weaver/output/transformer_v6_part2.pptx"
prs.save(out)
print("Saved: " + out)

print("\n=== 每页字数统计 ===")
pages = [
    ("第8页 训练细节", pts8),
    ("第9页 推理+KV缓存", pts9),
    ("第10页 词选择策略", pts10),
    ("第11页 训练vs推理", [(None, summary)]),
    ("第12页 应用与总结", pts12),
    ("第13页 感谢页", []),
]
for name, pts in pages:
    txt = ""
    for t, b in pts:
        if t: txt += t + " "
        if b: txt += b + " "
    cn = count_cn(txt)
    en = count_en(txt)
    total = cn + en
    ratio = (cn/total*100) if total > 0 else 0
    print("  %s: 中文=%d, 英文=%d, 总计=%d, 中文占比=%.1f%%" % (name, cn, en, total, ratio))

print("\n完成！")
