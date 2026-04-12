#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT V4 - 第8-13页 (训练细节/推理/词选择/对比/应用/总结)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 常量 ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_RIGHT = 13.2
SAFE_BOTTOM = 7.3

# 配色
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTN = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
C_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
C_LIGHT_RED = RGBColor(0xFD, 0xED, 0xED)
C_LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xE8)
C_LIGHT_ORANGE = RGBColor(0xFE, 0xF5, 0xE7)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_DARK = RGBColor(0x34, 0x49, 0x5E)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # 空白布局


# ── 辅助函数 ──
def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape


def add_box(slide, left, top, width, height, text, fill_color, font_size=10, font_color=C_TEXT, bold=False, align=PP_ALIGN.CENTER):
    shape = add_rect(slide, left, top, width, height, fill_color=fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_CN
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=10, font_color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_CN
    return txBox


def add_arrow_line(slide, x1, y1, x2, y2, color=C_DARK, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.fill.solid()
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_multi_text(slide, left, top, width, height, paragraphs_data, default_size=10, default_color=C_TEXT, line_spacing=1.15):
    """paragraphs_data: list of (text, font_size, bold, color, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, pdata in enumerate(paragraphs_data):
        text = pdata[0]
        fs = pdata[1] if len(pdata) > 1 else default_size
        b = pdata[2] if len(pdata) > 2 else False
        c = pdata[3] if len(pdata) > 3 else default_color
        a = pdata[4] if len(pdata) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = a
        p.space_after = Pt(2)
        p.space_before = Pt(1)
        p.line_spacing = Pt(fs * line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = FONT_CN
    return txBox


def add_title(slide, text):
    return add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5),
                       text, font_size=16, font_color=C_TITLE, bold=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# 第8页: 训练细节——残差/归一化/掩码/损失
# ════════════════════════════════════════════════════
def slide_08():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "训练细节：安全网、掩码与损失函数")

    # ── 左侧上方: 因果掩码矩阵 (4x4) ──
    mask_left = Inches(0.3)
    mask_top = Inches(0.85)
    cell_w = Inches(0.8)
    cell_h = Inches(0.55)
    add_textbox(slide, mask_left, mask_top - Inches(0.35), Inches(4), Inches(0.3),
                "因果掩码矩阵 (Causal Mask)", font_size=11, font_color=C_TITLE, bold=True)

    # Column headers
    headers = ["[START]", "我", "爱", "你"]
    for j, h in enumerate(headers):
        add_textbox(slide, mask_left + Inches(1.0) + j * cell_w, mask_top, cell_w, cell_h,
                    h, font_size=9, font_color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
    # Row labels + cells
    for i, h in enumerate(headers):
        add_textbox(slide, mask_left, mask_top + (i + 1) * cell_h, Inches(0.95), cell_h,
                    h, font_size=9, font_color=C_DARK, bold=True, align=PP_ALIGN.RIGHT)
        for j in range(4):
            x = mask_left + Inches(1.0) + j * cell_w
            y = mask_top + (i + 1) * cell_h
            if j <= i:
                # visible - blue
                add_box(slide, x, y, cell_w, cell_h, str(0.1 + (3 - i) * 0.2 if j == i else 0.05),
                        fill_color=C_LIGHT_BLUE, font_size=9, font_color=C_ENCODER)
            else:
                # masked - red
                add_box(slide, x, y, cell_w, cell_h, "-inf",
                        fill_color=C_LIGHT_RED, font_size=9, font_color=C_DECODER)

    add_textbox(slide, mask_left, mask_top + 5.2 * cell_h, Inches(4), Inches(0.25),
                "下三角(蓝色)=可见, 上三角(红色)=-inf被屏蔽", font_size=9, font_color=C_GRAY)

    # ── 左侧下方: 损失计算流程 ──
    flow_top = Inches(4.0)
    add_textbox(slide, mask_left, flow_top, Inches(4), Inches(0.3),
                "损失计算流程", font_size=11, font_color=C_TITLE, bold=True)

    flow_items = [
        ("解码器输出(512维)", C_ENCODER),
        ("Linear(512 x 37000)", C_EMBED),
        ("Softmax -> 概率分布", C_ATTN),
        ("与目标词对比 -> Loss", C_DECODER),
        ("反向传播 -> 更新6500万参数", C_FFN),
    ]
    for i, (txt, color) in enumerate(flow_items):
        y = flow_top + Inches(0.4) + i * Inches(0.55)
        add_box(slide, mask_left + Inches(0.2), y, Inches(3.2), Inches(0.42), txt, fill_color=None,
                font_size=9, font_color=color, align=PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            add_arrow_line(slide, mask_left + Inches(1.8), y + Inches(0.42),
                           mask_left + Inches(1.8), y + Inches(0.55), color=color, width=Pt(1.5))

    # ── 左侧中部: 残差连接示意 ──
    res_top = Inches(1.5)
    add_textbox(slide, Inches(4.5), res_top, Inches(4), Inches(0.3),
                "残差连接 + 层归一化", font_size=11, font_color=C_TITLE, bold=True)

    # Main path
    boxes_res = [
        ("输入 x(512维)", Inches(0.15)),
        ("Multi-Head\nSelf-Attention", Inches(0.9)),
        ("Add(x+Attn)", Inches(1.7)),
        ("LayerNorm", Inches(2.4)),
        ("FFN", Inches(3.1)),
        ("Add + Norm", Inches(3.8)),
    ]
    for txt, dx in boxes_res:
        add_box(slide, Inches(4.5) + dx, res_top + Inches(0.4), Inches(0.65), Inches(0.5),
                txt, fill_color=None, font_size=8, font_color=C_TEXT)
    # Bypass arrow (residual)
    add_arrow_line(slide, Inches(4.72), res_top + Inches(0.95),
                   Inches(4.72), res_top + Inches(0.95), color=C_FFN, width=Pt(2))

    # Residual bypass label
    add_textbox(slide, Inches(5.3), res_top + Inches(1.0), Inches(2.5), Inches(0.3),
                "残差旁路: x 直接跳接到 Add", font_size=8, font_color=C_FFN, bold=True)

    # LayerNorm formula
    add_textbox(slide, Inches(4.5), res_top + Inches(1.35), Inches(4), Inches(0.5),
                "LayerNorm: y = gamma * (x - mu) / sigma + beta\nmu=均值, sigma=标准差, gamma/beta=可学习参数",
                font_size=9, font_color=C_DARK)

    # ── 右侧文字 ──
    text_left = Inches(8.6)
    text_w = Inches(4.5)
    text_top = Inches(0.85)
    text_h = Inches(6.3)

    paras = [
        ("1. 残差连接的数学原理", 10.5, True, C_TITLE),
        ("Output = LayerNorm(x + SubLayer(x))。从梯度角度分析: "
         "dLoss/dx = dLoss/dOutput x (dSubLayer/dx + 1), 加号后面的\"1\"确保梯度永远不会消失。"
         "这是ResNet能训练100+层深度网络的核心秘诀。一层完整的数据流为: "
         "输入 -> 注意力(词间交流) -> +残差 -> 归一化(数值稳定) -> FFN(特征提取) -> +残差 -> 归一化 -> 输出。"
         "6层重复执行，每轮让词的表示向量更加丰富。", 9.5, False, C_TEXT),
        ("2. 因果掩码 = 杜绝作弊", 10.5, True, C_TITLE),
        ("掩码矩阵M[i,j]的规则: 当i>=j时M[i,j]=0(可以看), "
         "当i<j时M[i,j]=-inf(不可看)。经过Score_masked = Score + M处理后, "
         "Softmax会将-inf位置的概率压缩为0, 模型完全无法看到未来信息。"
         "训练和推理都必须严格执行这个掩码——保证训练和推理行为一致, 避免\"暴露偏差\"问题。", 9.5, False, C_TEXT),
        ("3. 交叉熵损失 = 衡量预测与答案的距离", 10.5, True, C_TITLE),
        ("Loss = -log(p(正确词))。当p=0.99时Loss约0.01(预测很好), "
         "当p=0.01时Loss约4.6(预测很差)。总Loss = -(1/T) x sum(log(p(正确词_t))), "
         "对所有位置取平均。训练初期Loss约10.5(接近log(37000)的随机猜测水平), "
         "充分训练后降到1.5-3.0(模型已经学会大部分语言规律)。", 9.5, False, C_TEXT),
        ("4. 训练配置与优化策略", 10.5, True, C_TITLE),
        ("WMT 2014英德翻译, 450万句对, 批量大小25000 token, "
         "Adam优化器(beta1=0.9, beta2=0.98), 学习率峰值0.0003(前4000步线性增温后余弦衰减), "
         "标签平滑0.1(防止过度自信), dropout 0.3(防止过拟合), "
         "8块P100 GPU训练12小时约30万步, 最终BLEU得分28.4。"
         "反向传播通过自动微分计算dLoss/dtheta(6500万个梯度), Adam按减小损失方向更新全部参数。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, text_left, text_top, text_w, text_h, paras, line_spacing=1.15)


# ════════════════════════════════════════════════════
# 第9页: 推理总览 + KV缓存
# ════════════════════════════════════════════════════
def slide_09():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "推理阶段：模型如何独立生成回答？")

    # ── 左侧上方: 自回归步骤图 ──
    ar_left = Inches(0.3)
    ar_top = Inches(0.85)
    add_textbox(slide, ar_left, ar_top - Inches(0.25), Inches(8), Inches(0.3),
                "自回归生成流程 (Autoregressive)", font_size=11, font_color=C_TITLE, bold=True)

    steps = [
        ("Step 1", "[START]", "->", "我", C_ENCODER),
        ("Step 2", "[START, 我]", "->", "爱", C_ENCODER),
        ("Step 3", "[START, 我, 爱]", "->", "你", C_ENCODER),
        ("Step 4", "[START, 我, 爱, 你]", "->", "<END>", C_DECODER),
    ]
    for i, (step, inp, arrow, out, color) in enumerate(steps):
        y = ar_top + Inches(0.15) + i * Inches(0.55)
        add_box(slide, ar_left, y, Inches(0.9), Inches(0.42), step, fill_color=C_LIGHT_BG,
                font_size=9, font_color=C_DARK, bold=True)
        add_box(slide, ar_left + Inches(1.0), y, Inches(2.8), Inches(0.42), inp, fill_color=None,
                font_size=9, font_color=C_TEXT, align=PP_ALIGN.LEFT)
        add_textbox(slide, ar_left + Inches(3.9), y, Inches(0.4), Inches(0.42),
                    arrow, font_size=12, font_color=color, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, ar_left + Inches(4.3), y, Inches(1.0), Inches(0.42), out, fill_color=color,
                font_size=10, font_color=C_WHITE, bold=True)

    # Encoder only once
    enc_y = ar_top + Inches(2.4)
    add_box(slide, ar_left + Inches(5.5), ar_top + Inches(0.15), Inches(2.8), Inches(0.42),
            "编码器: 只运行一次 -> 缓存输出", fill_color=C_ENCODER, font_size=9, font_color=C_WHITE, bold=True)
    add_box(slide, ar_left + Inches(5.5), ar_top + Inches(0.65), Inches(2.8), Inches(0.42),
            "K, V 向量供所有Step复用", fill_color=C_LIGHT_BLUE, font_size=9, font_color=C_ENCODER)

    # ── 左侧下方: 无缓存 vs 有缓存对比 ──
    cmp_top = Inches(3.4)
    add_textbox(slide, ar_left, cmp_top, Inches(8), Inches(0.3),
                "KV缓存: 从重复计算到增量计算", font_size=11, font_color=C_TITLE, bold=True)

    # 无缓存
    add_textbox(slide, ar_left, cmp_top + Inches(0.35), Inches(3.5), Inches(0.25),
                "无缓存 (每次全部重算)", font_size=10, font_color=C_DECODER, bold=True)
    no_cache_steps = [
        "Step1: 算[START]的Q,K,V",
        "Step2: 重算[START]+新算[我]的Q,K,V",
        "Step3: 重算[START,我]+新算[爱]的Q,K,V",
        "计算量: O(T^2 x d) - 严重浪费",
    ]
    for i, txt in enumerate(no_cache_steps):
        add_textbox(slide, ar_left + Inches(0.1), cmp_top + Inches(0.65) + i * Inches(0.35),
                    Inches(3.8), Inches(0.3), txt, font_size=9, font_color=C_DECODER)

    # 有缓存
    add_textbox(slide, Inches(4.3), cmp_top + Inches(0.35), Inches(3.5), Inches(0.25),
                "有缓存 (只算新词)", font_size=10, font_color=C_FFN, bold=True)
    cache_steps = [
        "Step1: 算[START]的Q,K,V -> 存入Cache",
        "Step2: 只算[我]的新K,V -> 追加到Cache",
        "Step3: 只算[爱]的新K,V -> 追加到Cache",
        "计算量: O(T x d) - 大幅节省",
    ]
    for i, txt in enumerate(cache_steps):
        add_textbox(slide, Inches(4.4), cmp_top + Inches(0.65) + i * Inches(0.35),
                    Inches(3.8), Inches(0.3), txt, font_size=9, font_color=C_FFN)

    # Cache visualization
    cache_y = cmp_top + Inches(2.1)
    add_box(slide, ar_left, cache_y, Inches(8), Inches(0.7),
            "KV Cache: { K: [k_START, k_我, k_爱, ...], V: [v_START, v_我, v_爱, ...] }\n"
            "12层 x T词 x 2(K+V) x 8头 x 64维 x 2字节 = 内存开销随生成长度线性增长",
            fill_color=C_LIGHT_GREEN, font_size=9, font_color=C_TEXT, align=PP_ALIGN.LEFT)

    # ── 右侧文字 ──
    paras = [
        ("1. 推理 = 闭卷考试", 10.5, True, C_TITLE),
        ("与训练的\"开卷考试\"完全不同, 推理时模型只能靠自己之前生成的词来预测下一个词。"
         "第一步只有[START], 凭空生成第一个词\"我\"; 加入\"我\"后生成\"爱\"; "
         "再加入\"爱\"生成\"你\"; 直到生成\"<END>\"停止。"
         "每一步都依赖前一步的正确性——第一步错了, 后续所有词可能全部偏离。"
         "这就是为什么推理质量对输入如此敏感。", 9.5, False, C_TEXT),
        ("2. 编码器只运行一次", 10.5, True, C_TITLE),
        ("编码器处理完源语言输入后, 输出的K和V向量被缓存, 整个生成过程中反复使用。"
         "编码器的全部计算成本只发生一次, 与目标语言的生成长度完全无关。"
         "这也是Encoder-Decoder架构的一个重要优势。", 9.5, False, C_TEXT),
        ("3. 串行瓶颈与KV缓存", 10.5, True, C_TITLE),
        ("与训练时所有位置完全并行不同, 推理必须逐词串行——生成100个词需要100次前向传播, "
         "而训练只需要1次。无缓存时每步重新计算全部已有词的K和V, 总计算量O(T^2 x d)。"
         "KV缓存的思路: 每步只算新词的K和V追加到缓存, 新词的Q与缓存中所有K做点积, "
         "每步仅需O(d x T)而非O(T x d x T)。", 9.5, False, C_TEXT),
        ("4. KV缓存的内存代价", 10.5, True, C_TITLE),
        ("典型配置下: 12层 x 100词 x 2(K+V) x 8头 x 64维 x 2字节(float16)约375MB。"
         "GPT-4等大模型的KV缓存可达数GB。这是推理速度和内存之间的核心权衡, "
         "也是PagedAttention(MQ推理框架vLLM的核心技术)等优化技术要解决的问题。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, Inches(8.6), Inches(0.85), Inches(4.5), Inches(6.3), paras, line_spacing=1.15)


# ════════════════════════════════════════════════════
# 第10页: 词选择策略
# ════════════════════════════════════════════════════
def slide_10():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "推理·词选择：从贪心到创意")

    # ── 左侧: 5种策略卡片 ──
    card_left = Inches(0.3)
    card_top = Inches(0.85)
    card_w = Inches(3.8)
    card_h = Inches(1.1)

    strategies = [
        ("1. 贪心搜索 (Greedy)", "每步选概率最高的词", "确定性强 | 容易循环 | 可能错过全局最优", C_ENCODER, C_LIGHT_BLUE),
        ("2. 束搜索 (Beam Search)", "保留k条最佳路径(beam=4-10)", "翻译效果好 | 偏平庸 | 不适合开放域对话", C_EMBED, RGBColor(0xF4, 0xEC, 0xF7)),
        ("3. Top-K采样", "从前K个候选词中随机选择(K=50)", "质量与多样性平衡 | K值固定不够灵活", C_ATTN, C_LIGHT_ORANGE),
        ("4. Top-P采样 (核采样)", "累计概率达P(0.9)的最小词集", "动态调整候选数 | ChatGPT默认策略 | 比Top-K更智能", C_FFN, C_LIGHT_GREEN),
        ("5. 温度 (Temperature)", "Softmax前对logits缩放: logits/T", "T<1(0.3)确定 | T=1原始 | T>1(1.5)随机多样", C_DECODER, C_LIGHT_RED),
    ]

    for i, (title, desc, note, title_color, bg_color) in enumerate(strategies):
        y = card_top + i * (card_h + Inches(0.08))
        # Card background
        add_rect(slide, card_left, y, card_w, card_h, fill_color=bg_color, line_color=title_color, line_width=Pt(1.5))
        # Title
        add_textbox(slide, card_left + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.3),
                    title, font_size=10, font_color=title_color, bold=True)
        # Description
        add_textbox(slide, card_left + Inches(0.1), y + Inches(0.35), card_w - Inches(0.2), Inches(0.3),
                    desc, font_size=9, font_color=C_TEXT)
        # Note
        add_textbox(slide, card_left + Inches(0.1), y + Inches(0.65), card_w - Inches(0.2), Inches(0.35),
                    note, font_size=8, font_color=C_GRAY)

    # ── 左侧右侧: 温度效果示意 ──
    temp_left = Inches(4.3)
    temp_top = Inches(0.85)
    add_textbox(slide, temp_left, temp_top, Inches(4), Inches(0.3),
                "温度对概率分布的影响", font_size=11, font_color=C_TITLE, bold=True)

    # Three distribution bars
    dists = [
        ("T=0.3 (低温/精确)", [(0.7, "散步"), (0.15, "跑步"), (0.1, "休息"), (0.05, "其他")], C_ENCODER),
        ("T=1.0 (中温/平衡)", [(0.35, "散步"), (0.3, "跑步"), (0.2, "休息"), (0.15, "其他")], C_ATTN),
        ("T=1.5 (高温/创意)", [(0.2, "散步"), (0.22, "跑步"), (0.23, "休息"), (0.35, "其他")], C_DECODER),
    ]
    bar_max_w = Inches(2.8)
    for i, (label, bars, color) in enumerate(dists):
        y = temp_top + Inches(0.4) + i * Inches(1.2)
        add_textbox(slide, temp_left, y, Inches(4), Inches(0.25),
                    label, font_size=9, font_color=color, bold=True)
        x_offset = temp_left
        for prob, name in bars:
            w = Inches(prob * 3.5)
            add_box(slide, x_offset, y + Inches(0.28), w, Inches(0.3),
                    f"{name} {prob:.0%}" if prob >= 0.1 else f"{prob:.0%}",
                    fill_color=color, font_size=8, font_color=C_WHITE, align=PP_ALIGN.CENTER)
            x_offset += w

    # 组合策略说明
    combo_top = Inches(4.1)
    add_rect(slide, temp_left, combo_top, Inches(4), Inches(1.5), fill_color=C_LIGHT_BG, line_color=C_ENCODER, line_width=Pt(1))
    add_textbox(slide, temp_left + Inches(0.1), combo_top + Inches(0.05), Inches(3.8), Inches(0.25),
                "ChatGPT组合策略", font_size=10, font_color=C_TITLE, bold=True)
    combo_text = ("Temperature=0.7 + Top-P=0.9 + 重复惩罚\n"
                  "写代码时 T=0.2 (精确), 创意写作 T=0.7-1.0 (多样性)\n"
                  "翻译任务: Beam Search + 长度惩罚\n"
                  "数学推理: T=0.0 (确定性输出)")
    add_textbox(slide, temp_left + Inches(0.1), combo_top + Inches(0.3), Inches(3.8), Inches(1.1),
                combo_text, font_size=9, font_color=C_TEXT)

    # 贪心陷阱
    trap_top = Inches(5.8)
    add_rect(slide, temp_left, trap_top, Inches(4), Inches(1.3), fill_color=C_LIGHT_RED, line_color=C_DECODER, line_width=Pt(1))
    add_textbox(slide, temp_left + Inches(0.1), trap_top + Inches(0.05), Inches(3.8), Inches(0.25),
                "贪心陷阱: 局部最优 != 全局最优", font_size=10, font_color=C_DECODER, bold=True)
    add_textbox(slide, temp_left + Inches(0.1), trap_top + Inches(0.3), Inches(3.8), Inches(0.9),
                "\"我喜欢在公园里___\"\n"
                "贪心选\"散步\"(p=0.3) -> \"然后回家\"(p=0.2) -> 总概率低\n"
                "最优: \"跑步\"(p=0.25) + \"然后喝杯咖啡\"(p=0.4) -> 总概率更高",
                font_size=9, font_color=C_TEXT)

    # ── 右侧文字 ──
    paras = [
        ("1. 为什么不能总选最高概率词？", 10.5, True, C_TITLE),
        ("贪心搜索看似每步最优, 但可能错过全局最优路径。"
         "就像下棋, 每步吃子最多不一定能赢——局部最优不等于全局最优。"
         "在开放域对话中, \"最优\"本身就不存在, 因为同样的问题有无数种合理的回答方式。"
         "贪心还容易陷入循环: 模型反复生成相同的几个词, 无法跳出局部最优。", 9.5, False, C_TEXT),
        ("2. 束搜索的适用场景", 10.5, True, C_TITLE),
        ("束搜索保留k条最优路径(beam size通常4-10), 兼顾确定性和搜索空间。"
         "在机器翻译等任务中效果很好, 因为翻译通常有一个相对明确的标准答案。"
         "但在开放域对话和创意写作中, 束搜索生成的文本往往\"安全但平庸\"——"
         "语法正确但缺乏创意, 像官方发言稿。ChatGPT等对话模型不使用束搜索。", 9.5, False, C_TEXT),
        ("3. Top-K与Top-P的本质区别", 10.5, True, C_TITLE),
        ("Top-K固定选前K个候选词(K=50), 问题在于: 有时只有2-3个词概率很高(不需要50个), "
         "有时有100个词概率相当(50个不够)。Top-P(核采样)更智能: "
         "从概率最高的词开始累加, 直到累计概率达到P(通常0.9), 候选数量动态调整。"
         "这是ChatGPT的默认采样策略。", 9.5, False, C_TEXT),
        ("4. 温度: 控制创造力的旋钮", 10.5, True, C_TITLE),
        ("温度T通过Softmax前缩放logits来调节分布的\"尖锐程度\"。"
         "T<1(如0.3): 高概率词被进一步放大, 输出更确定、更保守, 适合代码生成和数学推理。"
         "T=1: 保持原始概率分布, 不加任何偏好。"
         "T>1(如1.5): 概率分布被\"压平\", 低概率词也有机会被选中, 输出更随机、更有创意。"
         "T=0等价于贪心搜索, 总是选概率最高的词。", 9.5, False, C_TEXT),
        ("5. 实际应用中的组合策略", 10.5, True, C_TITLE),
        ("现代大语言模型通常组合使用多种策略: ChatGPT用Temperature=0.7+Top-P=0.9+重复惩罚。"
         "不同任务需要不同的\"创造力温度\": 写代码T=0.2(精确), 写故事T=0.9(创意), "
         "做翻译用束搜索(准确), 数学推理T=0(确定性)。"
         "这些看似简单的选择策略, 直接决定了AI输出的质量和风格。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, Inches(8.6), Inches(0.85), Inches(4.5), Inches(6.3), paras, line_spacing=1.15)


# ════════════════════════════════════════════════════
# 第11页: 训练vs推理全面对比
# ════════════════════════════════════════════════════
def slide_11():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "训练vs推理：开卷考试与闭卷考试")

    # ── 左侧: 对比表格 ──
    table_left = Inches(0.3)
    table_top = Inches(0.85)
    col_w = Inches(3.5)
    col1_w = Inches(1.5)
    row_h = Inches(0.55)
    header_h = Inches(0.45)

    # Headers
    add_box(slide, table_left, table_top, col1_w, header_h, "对比维度",
            fill_color=C_TITLE, font_size=10, font_color=C_WHITE, bold=True)
    add_box(slide, table_left + col1_w, table_top, col_w, header_h, "训练 (开卷考试)",
            fill_color=C_ENCODER, font_size=10, font_color=C_WHITE, bold=True)
    add_box(slide, table_left + col1_w + col_w, table_top, col_w, header_h, "推理 (闭卷考试)",
            fill_color=C_FFN, font_size=10, font_color=C_WHITE, bold=True)

    rows = [
        ("解码器输入", "完整目标序列 (Teacher Forcing)", "模型逐词生成 (自回归)"),
        ("计算方式", "所有位置并行 (1次前向传播)", "逐词串行 (T次前向传播)"),
        ("损失函数", "交叉熵 (每个位置都计算)", "无损失函数 (只做预测)"),
        ("梯度计算", "反向传播更新全部参数", "无梯度 (参数完全冻结)"),
        ("KV Cache", "不需要 (每次全量计算)", "必须使用 (增量计算)"),
        ("GPU利用率", "高 (矩阵运算充分并行)", "低 (串行等待, 难以饱和)"),
        ("处理速度", "快 (1次前向+1次反向)", "慢 (T次前向传播)"),
        ("核心目标", "学习参数 (学会预测下一个词)", "生成结果 (产出实际文本)"),
    ]

    for i, (dim, train, infer) in enumerate(rows):
        y = table_top + header_h + i * row_h
        bg = C_LIGHT_BG if i % 2 == 0 else C_WHITE
        add_box(slide, table_left, y, col1_w, row_h, dim,
                fill_color=bg, font_size=9, font_color=C_DARK, bold=True, align=PP_ALIGN.LEFT)
        add_box(slide, table_left + col1_w, y, col_w, row_h, train,
                fill_color=C_LIGHT_BLUE, font_size=9, font_color=C_ENCODER, align=PP_ALIGN.LEFT)
        add_box(slide, table_left + col1_w + col_w, y, col_w, row_h, infer,
                fill_color=C_LIGHT_GREEN, font_size=9, font_color=C_TEXT, align=PP_ALIGN.LEFT)

    # ── 左侧下方: 比喻 + 暴露偏差 ──
    meta_top = table_top + header_h + len(rows) * row_h + Inches(0.2)
    add_rect(slide, table_left, meta_top, Inches(8.5), Inches(2.0), fill_color=C_LIGHT_BG, line_color=C_ATTN, line_width=Pt(1.5))
    add_textbox(slide, table_left + Inches(0.1), meta_top + Inches(0.05), Inches(8), Inches(0.3),
                "核心矛盾: 暴露偏差 (Exposure Bias)", font_size=11, font_color=C_DECODER, bold=True)
    meta_text = (
        "训练像新员工培训——师傅手把手教, 有标准答案对照, 效率高但需要大量计算资源(GPU集群)。"
        "推理像正式上岗——全靠自己积累的经验, 每步参考自己之前的结果, 只需一台电脑但必须一步步来。\n"
        "暴露偏差: 训练时解码器看到的是\"标准答案\"(Teacher Forcing), 但推理时看到的是\"自己的输出\"。"
        "如果模型在第3步犯了小错, 后续步骤基于错误输入, 错误会累积放大(误差级联)。\n"
        "缓解方案: Scheduled Sampling(训练时随机用模型输出替换标准答案), "
        "以及RLHF(基于人类反馈的强化学习, 让模型学会从自身输出继续生成)。"
    )
    add_textbox(slide, table_left + Inches(0.1), meta_top + Inches(0.35), Inches(8.2), Inches(1.5),
                meta_text, font_size=9, font_color=C_TEXT)

    # ── 右侧文字 ──
    paras = [
        ("1. 并行 vs 串行: 根本性差异", 10.5, True, C_TITLE),
        ("训练时, 解码器同时看到完整的目标序列[START, 我, 爱, 你, END], "
         "所有位置并行计算注意力, 一次前向传播得到所有位置的预测。"
         "GPU的数千个核心同时工作, 矩阵运算效率极高。推理时, 每步只能看到之前生成的词, "
         "必须逐词串行——生成100个词需要100次前向传播, GPU大量时间在等待。"
         "这是训练和推理最本质的差异, 也是推理优化的核心方向。", 9.5, False, C_TEXT),
        ("2. Teacher Forcing的利与弊", 10.5, True, C_TITLE),
        ("Teacher Forcing是训练的标准做法: 用真实答案(而非模型预测)作为解码器的输入。"
         "好处是训练稳定、收敛快——每步输入都是\"正确的\", 梯度信号清晰。"
         "坏处是造成暴露偏差: 训练时从未见过\"错误输入\", 推理时一旦犯错就不知所措。"
         "就像考试时做错一道题, 后面的题全依赖这道题的答案, 错误会级联放大。", 9.5, False, C_TEXT),
        ("3. KV缓存: 推理加速的关键", 10.5, True, C_TITLE),
        ("训练时每步计算所有词的K和V(反正只有一步), 不需要缓存。"
         "推理时每步只算新词的K和V, 追加到缓存中, 大幅减少重复计算。"
         "但KV缓存消耗大量显存: 序列越长, 缓存越大。GPT-4生成10000个token时, "
         "KV缓存可达数十GB, 这是长文本推理的主要瓶颈。"
         "PagedAttention(vLLM)通过内存分页管理优化KV缓存的利用效率。", 9.5, False, C_TEXT),
        ("4. 从训练到部署的完整链路", 10.5, True, C_TITLE),
        ("预训练(大规模语料, 学会语言规律) -> 指令微调(学会遵循指令) -> RLHF(学会人类偏好) -> "
         "量化压缩(降低推理成本) -> 部署服务(KV缓存管理, 批处理调度)。"
         "每个环节都在解决训练-推理之间差距带来的问题。"
         "训练追求\"学会\", 推理追求\"用好\", 两者的优化目标和技术路线截然不同。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, Inches(8.6), Inches(0.85), Inches(4.5), Inches(6.3), paras, line_spacing=1.15)


# ════════════════════════════════════════════════════
# 第12页: 应用与总结 (上) - 技术演进树
# ════════════════════════════════════════════════════
def slide_12():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "从翻译工具到AI通用引擎：技术演进全景")

    # ── 左侧: 演进树形图 ──
    tree_left = Inches(0.3)
    tree_top = Inches(0.85)

    # Root node
    root_w = Inches(2.2)
    root_h = Inches(0.5)
    root_x = Inches(3.5)
    add_box(slide, root_x, tree_top, root_w, root_h, "Transformer 2017\nVaswani et al.",
            fill_color=C_TITLE, font_size=10, font_color=C_WHITE, bold=True)

    # Three branches
    branch_y = tree_top + Inches(0.8)
    branch_gap = Inches(3.0)
    branch_starts = [tree_left + Inches(0.3), tree_left + Inches(3.3), tree_left + Inches(6.3)]
    branch_colors = [C_ENCODER, C_FFN, C_EMBED]
    branch_labels = ["仅解码器 (Decoder-Only)", "完整架构 (Encoder-Decoder)", "仅编码器 (Encoder-Only)"]
    branch_tags = ["生成之王", "转换之王", "理解之王"]

    for i, (bx, color, label, tag) in enumerate(zip(branch_starts, branch_colors, branch_labels, branch_tags)):
        add_box(slide, bx, branch_y, Inches(2.6), Inches(0.55), f"{label}\n[{tag}]",
                fill_color=color, font_size=9, font_color=C_WHITE, bold=True)
        # Arrow from root to branch
        add_arrow_line(slide, root_x + root_w / 2, tree_top + root_h,
                       bx + Inches(1.3), branch_y, color=color, width=Pt(1.5))

    # Sub-branches
    sub_top = branch_y + Inches(0.75)
    sub_data = [
        # (branch_idx, models, color, bg)
        (0, ["GPT-1 (2018, 117M)", "GPT-2 (2019, 1.5B)", "GPT-3 (2020, 175B)", "ChatGPT/GPT-4 (2023)", "GLM/Qwen/LLaMA"], C_ENCODER, C_LIGHT_BLUE),
        (1, ["T5 (2019, 60M-11B)", "BART (2019)", "mT5 (多语言翻译)", "Whisper (2022, 语音)", "Universal Transformer"], C_FFN, C_LIGHT_GREEN),
        (2, ["BERT (2018, 110M-340M)", "RoBERTa (2019)", "DeBERTa (2020)", "ALBERT (轻量化)", "ELECTRA (高效预训练)"], C_EMBED, RGBColor(0xF4, 0xEC, 0xF7)),
    ]

    for branch_idx, models, color, bg in sub_data:
        bx = branch_starts[branch_idx]
        for j, model in enumerate(models):
            y = sub_top + j * Inches(0.4)
            add_box(slide, bx, y, Inches(2.6), Inches(0.35), model, fill_color=bg,
                    font_size=8, font_color=color, align=PP_ALIGN.CENTER)

    # 参数量时间线
    tl_top = sub_top + Inches(2.3)
    add_textbox(slide, tree_left, tl_top, Inches(8.5), Inches(0.25),
                "参数量演进: 65M -> 110M -> 1.5B -> 175B -> 1.8T (增长27000倍)", font_size=9, font_color=C_TITLE, bold=True)
    add_textbox(slide, tree_left, tl_top + Inches(0.3), Inches(8.5), Inches(0.25),
                "时间线: 2017 Transformer -> 2018 BERT/GPT -> 2020 GPT-3 -> 2022 ChatGPT -> 2023 GPT-4 -> 2024 Sora", font_size=9, font_color=C_DARK)

    # 超越语言
    beyond_top = tl_top + Inches(0.65)
    add_rect(slide, tree_left, beyond_top, Inches(8.5), Inches(0.8), fill_color=C_LIGHT_ORANGE, line_color=C_ATTN, line_width=Pt(1.5))
    add_textbox(slide, tree_left + Inches(0.1), beyond_top + Inches(0.05), Inches(8), Inches(0.25),
                "超越自然语言: Transformer的通用性", font_size=10, font_color=C_ATTN, bold=True)
    add_textbox(slide, tree_left + Inches(0.1), beyond_top + Inches(0.3), Inches(8.2), Inches(0.4),
                "ViT(2020): 图像切块当token -> 图像分类 | Whisper(2022): 音频频谱 -> 语音识别 | "
                "GPT-4V(2023): 图像+文本 -> 多模态理解 | Sora(2024): 视频生成 | AlphaFold2: 蛋白质结构预测",
                font_size=9, font_color=C_TEXT)

    # ── 右侧文字 ──
    paras = [
        ("1. GPT系列: 生成之王的崛起", 10.5, True, C_TITLE),
        ("仅保留解码器部分, 通过\"预测下一个词\"的任务, 在海量文本上预训练。"
         "ChatGPT的三阶段训练: 预训练(3000亿token学习语言规律) -> 指令微调(学会遵循人类指令) -> "
         "RLHF(通过人类反馈的强化学习对齐人类偏好)。代表模型: "
         "GPT-3(1750亿参数), GPT-4(约1.8万亿, 多模态), "
         "以及开源阵营的LLaMA、Qwen、GLM-4等。GPT系列证明了\"规模即能力\"的缩放定律。", 9.5, False, C_TEXT),
        ("2. BERT系列: 理解之王的双向编码", 10.5, True, C_TITLE),
        ("仅保留编码器, 通过掩码语言建模(随机遮盖15%的token让模型预测)学习双向上下文。"
         "擅长文本分类、问答系统、命名实体识别等\"理解型\"任务。"
         "Google搜索从2019年开始使用BERT, 显著提升了搜索结果的相关性。"
         "代表模型: BERT-base(1.1亿参数), BERT-large(3.4亿), "
         "RoBERTa(优化训练策略), DeBERTa(解耦注意力机制)。", 9.5, False, C_TEXT),
        ("3. T5系列: 万物皆可Text-to-Text", 10.5, True, C_TITLE),
        ("保留完整的编码器-解码器架构, 将所有NLP任务统一为\"输入文本 -> 输出文本\"。"
         "翻译是\"translate English to German: ...\", 摘要是\"summarize: ...\", "
         "分类是\"classify: ...\"。Google Translate基于mT5模型。"
         "代表: T5-Small(6000万)到T5-XXL(110亿)。Whisper(OpenAI的语音识别模型)也采用编码器-解码器架构。", 9.5, False, C_TEXT),
        ("4. 超越语言: 万物皆序列", 10.5, True, C_TITLE),
        ("ViT(2020)将图像切成16x16的小块当作token处理, 证明了Transformer在视觉任务上的有效性。"
         "Whisper将音频频谱图当作序列输入, 实现了多语言语音识别。"
         "GPT-4V实现了图像和文本的联合理解。Sora将视频生成建模为时空token的预测。"
         "AlphaFold2用Transformer预测蛋白质三维结构(2024年诺贝尔化学奖)。"
         "核心洞察: 只要能表示为序列, Transformer就能处理。", 9.5, False, C_TEXT),
        ("5. 缩放定律: 规模即能力", 10.5, True, C_TITLE),
        ("Kaplan等(2020)发现: 模型性能随参数量、数据量、计算量的幂律增长——"
         "参数增加10倍, loss稳定下降。从原始Transformer的6500万参数到GPT-4的约1.8万亿, "
         "增长了27000倍, 但核心架构(自注意力+FFN)几乎没有改变。"
         "Transformer不是某个具体的模型, 而是AI时代的计算范式和基础设施。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, Inches(8.6), Inches(0.85), Inches(4.5), Inches(6.3), paras, line_spacing=1.15)


# ════════════════════════════════════════════════════
# 第13页: 总结 - Transformer的核心本质
# ════════════════════════════════════════════════════
def slide_13():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "总结：Transformer的核心本质与深远影响")

    # ── 左侧: 核心架构回顾图 ──
    arch_left = Inches(0.3)
    arch_top = Inches(0.85)

    add_textbox(slide, arch_left, arch_top, Inches(8.5), Inches(0.3),
                "Transformer核心架构回顾", font_size=11, font_color=C_TITLE, bold=True)

    # 完整架构简图 (水平布局)
    # Input Embedding
    y_base = arch_top + Inches(0.4)
    comp_h = Inches(0.45)
    comp_gap = Inches(0.08)

    # Row 1: Input pipeline
    components_r1 = [
        ("输入文本", Inches(0.9), C_GRAY),
        ("BPE分词", Inches(1.0), C_EMBED),
        ("词嵌入\n37000x512", Inches(1.2), C_EMBED),
        ("位置编码\nsin/cos", Inches(1.1), C_ATTN),
    ]
    x = arch_left
    for label, w, color in components_r1:
        add_box(slide, x, y_base, w, comp_h, label, fill_color=color, font_size=8, font_color=C_WHITE, bold=True)
        if x > arch_left:
            add_arrow_line(slide, x - comp_gap, y_base + comp_h / 2, x, y_base + comp_h / 2, color=color, width=Pt(1.2))
        x += w + comp_gap

    # Arrow down to encoder
    enc_start_x = arch_left
    enc_top = y_base + comp_h + Inches(0.15)

    # Encoder block
    enc_w = Inches(4.0)
    add_box(slide, enc_start_x, enc_top, enc_w, Inches(0.4), "编码器 x 6层", fill_color=C_ENCODER, font_size=10, font_color=C_WHITE, bold=True)
    # Encoder internals
    enc_items = ["Multi-Head Self-Attention (8头x64维)", "Add & LayerNorm (残差+归一化)",
                 "Feed-Forward Network (512->2048->512)", "Add & LayerNorm"]
    for i, item in enumerate(enc_items):
        add_box(slide, enc_start_x + Inches(0.1), enc_top + Inches(0.45) + i * Inches(0.35),
                enc_w - Inches(0.2), Inches(0.3), item, fill_color=C_LIGHT_BLUE, font_size=8, font_color=C_ENCODER)

    # Decoder block
    dec_start_x = enc_start_x + enc_w + Inches(0.3)
    dec_w = Inches(4.2)
    add_box(slide, dec_start_x, enc_top, dec_w, Inches(0.4), "解码器 x 6层", fill_color=C_DECODER, font_size=10, font_color=C_WHITE, bold=True)
    dec_items = ["Masked Self-Attention (因果掩码)", "Add & LayerNorm",
                 "Cross-Attention (Q来自解码器, K/V来自编码器)", "Add & LayerNorm",
                 "Feed-Forward Network (512->2048->512)", "Add & LayerNorm"]
    for i, item in enumerate(dec_items):
        add_box(slide, dec_start_x + Inches(0.1), enc_top + Inches(0.45) + i * Inches(0.35),
                dec_w - Inches(0.2), Inches(0.3), item, fill_color=C_LIGHT_RED, font_size=8, font_color=C_DECODER)

    # K,V arrow from encoder to decoder cross-attention
    kv_y = enc_top + Inches(0.45) + 2 * Inches(0.35) + Inches(0.15)
    add_arrow_line(slide, enc_start_x + enc_w, kv_y, dec_start_x, kv_y,
                   color=C_EMBED, width=Pt(2))
    add_textbox(slide, enc_start_x + enc_w + Inches(0.05), kv_y - Inches(0.2), Inches(0.25), Inches(0.2),
                "K,V", font_size=8, font_color=C_EMBED, bold=True)

    # Output pipeline
    out_top = enc_top + Inches(0.45) + 6 * Inches(0.35) + Inches(0.15)
    out_items = [
        ("Linear\n512x37000", Inches(1.2), C_DARK),
        ("Softmax", Inches(0.8), C_ATTN),
        ("概率分布\n-> 选词", Inches(1.0), C_FFN),
        ("输出文本", Inches(0.9), C_GRAY),
    ]
    x = dec_start_x
    for label, w, color in out_items:
        add_box(slide, x, out_top, w, comp_h, label, fill_color=color, font_size=8, font_color=C_WHITE, bold=True)
        if x > dec_start_x:
            add_arrow_line(slide, x - comp_gap, out_top + comp_h / 2, x, out_top + comp_h / 2, color=color, width=Pt(1.2))
        x += w + comp_gap

    # ── 左侧下方: 五大核心要点 ──
    key_top = out_top + comp_h + Inches(0.25)
    add_textbox(slide, arch_left, key_top, Inches(8.5), Inches(0.25),
                "五大核心要点", font_size=11, font_color=C_TITLE, bold=True)

    key_points = [
        ("1. 自注意力 = 信息自由流动", "任意token直接交互, 不受距离限制, 解决RNN的顺序瓶颈和长距离遗忘", C_ENCODER),
        ("2. 并行化 = 效率革命", "所有位置同时计算, 充分利用GPU并行能力, 训练速度比RNN快数十倍", C_FFN),
        ("3. 残差+归一化 = 深度保障", "梯度永不消失, 数值稳定, 支撑6层(原始)到数百层(现代)的深度堆叠", C_ATTN),
        ("4. 缩放定律 = 规模即能力", "更多参数+更多数据+更多计算 = 更强性能, 从6500万到万亿参数", C_EMBED),
        ("5. 通用架构 = 万物皆序列", "语言、图像、语音、视频、蛋白质... 只要能表示为序列, Transformer就能处理", C_DECODER),
    ]
    for i, (title, desc, color) in enumerate(key_points):
        y = key_top + Inches(0.3) + i * Inches(0.4)
        add_textbox(slide, arch_left + Inches(0.1), y, Inches(2.8), Inches(0.35),
                    title, font_size=9, font_color=color, bold=True)
        add_textbox(slide, arch_left + Inches(3.0), y, Inches(5.5), Inches(0.35),
                    desc, font_size=9, font_color=C_TEXT)

    # ── 右侧文字 ──
    paras = [
        ("1. Transformer的本质: 让信息自由流动", 10.5, True, C_TITLE),
        ("RNN像排队传递消息——第100个人要等前99个人传完才能收到第一条消息。"
         "Transformer像圆桌会议——任何人可以随时直接与任何人交流。"
         "自注意力机制(Q x K^T)让序列中任意两个位置的token直接建立关联, "
         "距离不再是信息传递的障碍。这个看似简单的设计, 彻底改变了序列建模的范式。", 9.5, False, C_TEXT),
        ("2. 编码器理解, 解码器生成", 10.5, True, C_TITLE),
        ("编码器的职责是\"理解\": 通过6层自注意力+FFN, 让每个token的表示融合整个句子的上下文。"
         "解码器的职责是\"生成\": 在掩码自注意力(保证因果性)的基础上, "
         "通过交叉注意力\"查阅\"编码器的理解结果, 逐步生成目标序列。"
         "这种\"先理解后生成\"的分工, 是Encoder-Decoder架构的核心设计哲学。", 9.5, False, C_TEXT),
        ("3. 训练是开卷考试, 推理是闭卷考试", 10.5, True, C_TITLE),
        ("训练时, 解码器看到完整目标序列(Teacher Forcing), 所有位置并行计算, 效率极高。"
         "推理时, 模型只能依赖自己之前生成的词, 逐词串行生成, 速度远慢于训练。"
         "KV缓存通过存储已计算的Key和Value向量, 避免重复计算, 是推理加速的核心技术。"
         "暴露偏差(训练看答案、推理看自己输出)是连接两个阶段的核心挑战。", 9.5, False, C_TEXT),
        ("4. 从6500万到万亿: 架构未变, 世界已变", 10.5, True, C_TITLE),
        ("2017年原始Transformer: 6500万参数, 机器翻译, 8块P100 GPU, 12小时训练。"
         "2024年GPT-4: 约1.8万亿参数, 多模态理解与生成, 数万块GPU, 数月训练。"
         "参数增长27000倍, 但核心架构(自注意力+前馈网络+残差连接+层归一化)几乎没有改变。"
         "Transformer不是某个具体的模型, 而是AI时代的计算范式和基础设施, "
         "就像TCP/IP之于互联网, 关系型数据库之于软件工程。", 9.5, False, C_TEXT),
        ("5. 展望: 下一个突破在哪里？", 10.5, True, C_TITLE),
        ("当前挑战: 推理成本(O(T^2)注意力计算)、长上下文窗口(百万token)、"
         "多模态融合(文本/图像/音频/视频统一架构)、"
         "高效推理(量化、蒸馏、MoE混合专家)。"
         "潜在方向: 线性注意力(降低复杂度)、状态空间模型(Mamba/RWKV)、"
         "稀疏注意力(长距离选择性连接)。"
         "但无论下一个突破是什么, Transformer已经为AI奠定了不可动摇的基石。", 9.5, False, C_TEXT),
    ]
    add_multi_text(slide, Inches(8.6), Inches(0.85), Inches(4.5), Inches(6.3), paras, line_spacing=1.15)


# ── 生成所有页面 ──
slide_08()
slide_09()
slide_10()
slide_11()
slide_12()
slide_13()

# ── 保存 ──
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v4_part2.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved: {output_path}")

# ── 统计 ──
print("\n" + "=" * 60)
print("每页统计:")
print("=" * 60)
for i, slide in enumerate(prs.slides, start=8):
    shape_count = len(slide.shapes)
    total_text = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    total_text += len(r.text)
    print(f"  第{i}页: {shape_count} shapes, {total_text} chars")
print("=" * 60)
print(f"共 {len(prs.slides)} 页")
