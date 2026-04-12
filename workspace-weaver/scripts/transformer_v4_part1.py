#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer PPT V4 - Part 1 (Pages 1-7)
全中文内容, 翔实专业, 文字填满页面
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ============================================================
# Constants
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_TITLE   = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN     = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED   = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT    = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY    = RGBColor(0x95, 0xA5, 0xA6)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
C_LIGHT_RED  = RGBColor(0xFD, 0xED, 0xEC)
C_LIGHT_GREEN = RGBColor(0xEA, 0xFE, 0xF1)
C_LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xE7)
C_LIGHT_PURPLE = RGBColor(0xF4, 0xEC, 0xF7)
C_LIGHT_GRAY  = RGBColor(0xF2, 0xF3, 0xF4)
C_BG_ENCODER_FILL = RGBColor(0xD6, 0xEA, 0xF8)
C_BG_DECODER_FILL = RGBColor(0xFA, 0xDB, 0xD8)
C_BG_ATTENTION_FILL = RGBColor(0xFC, 0xF3, 0xCF)
C_BG_FFN_FILL = RGBColor(0xD5, 0xF5, 0xE3)
C_BG_EMBED_FILL = RGBColor(0xEB, 0xDE, 0xFB)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank


# ============================================================
# Helper Functions
# ============================================================
def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rect_straight(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=10, color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_CN
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox

def add_rich_text_box(slide, left, top, width, height, paragraphs_data):
    """paragraphs_data: list of (text, font_size, color, bold, alignment, space_after, font_name)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pdata in enumerate(paragraphs_data):
        text, fs, color, bold, align, sa = pdata[0], pdata[1], pdata[2], pdata[3], pdata[4], pdata[5] if len(pdata) > 5 else Pt(2)
        fn = pdata[6] if len(pdata) > 6 else None
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = fn or FONT_CN
        p.alignment = align
        p.space_after = sa
        p.space_before = Pt(1)
        p.line_spacing = Pt(fs * 1.15)
    return txBox

def add_line(slide, x1, y1, x2, y2, color=C_TEXT, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_arrow_line(slide, x1, y1, x2, y2, color=C_TEXT, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead via XML
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)
    return connector

def set_shape_text(shape, text, font_size=9, color=C_TEXT, bold=False, alignment=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.name = FONT_CN
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    return shape


# ============================================================
# PAGE 1: Cover
# ============================================================
def make_page1():
    slide = prs.slides.add_slide(blank_layout)
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE

    # Decorative lines
    add_rect_straight(slide, Inches(0), Inches(2.3), SLIDE_W, Pt(3), C_ENCODER)
    add_rect_straight(slide, Inches(0), Inches(2.38), SLIDE_W, Pt(2), C_EMBED)
    add_rect_straight(slide, Inches(0), Inches(2.44), SLIDE_W, Pt(3), C_ATTENTION)

    # Main title
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11.333), Inches(1.2),
                 "Transformer 架构深度解析", 32, C_TITLE, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.6),
                 "从直觉到原理，一篇看懂AI的核心引擎", 16, C_TEXT, False, PP_ALIGN.CENTER)

    # Paper info
    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.5),
                 "Attention Is All You Need - Vaswani et al., 2017  |  Google Brain  |  6500万参数",
                 11, C_GRAY, False, PP_ALIGN.CENTER)

    # Bottom text
    add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10.333), Inches(0.5),
                 "从训练到推理，完整拆解Transformer的工作原理", 11, C_TEXT, False, PP_ALIGN.CENTER)

    # Decorative bottom lines
    add_rect_straight(slide, Inches(0), Inches(6.5), SLIDE_W, Pt(3), C_ENCODER)
    add_rect_straight(slide, Inches(0), Inches(6.58), SLIDE_W, Pt(2), C_EMBED)
    add_rect_straight(slide, Inches(0), Inches(6.64), SLIDE_W, Pt(3), C_ATTENTION)

    # Small architectural sketch at bottom-right
    # Simple encoder-decoder mini diagram
    bx = Inches(9.5)
    by = Inches(0.3)
    add_rect(slide, bx, by, Inches(1.2), Inches(0.6), C_BG_ENCODER_FILL, C_ENCODER, Pt(1.5))
    set_shape_text(slide.shapes[-1], "Encoder", 9, C_ENCODER, True)
    add_rect(slide, bx + Inches(1.5), by, Inches(1.2), Inches(0.6), C_BG_DECODER_FILL, C_DECODER, Pt(1.5))
    set_shape_text(slide.shapes[-1], "Decoder", 9, C_DECODER, True)
    add_arrow_line(slide, bx + Inches(1.2), by + Inches(0.3), bx + Inches(1.5), by + Inches(0.3), C_GRAY, Pt(1))


# ============================================================
# PAGE 2: Why Transformer
# ============================================================
def make_page2():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "为什么需要Transformer：RNN的三大痛点与注意力突破", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === LEFT SIDE: RNN sequential diagram ===
    lx = Inches(0.3)
    ly = Inches(0.85)

    # RNN sequential flow label
    add_text_box(slide, lx, ly, Inches(3), Inches(0.35),
                 "RNN顺序处理流程:", 11, C_TITLE, True)

    # RNN boxes: The -> cat -> sat -> on -> the -> mat
    words_rnn = ["The", "cat", "sat", "on", "the", "mat"]
    box_w = Inches(1.05)
    box_h = Inches(0.45)
    gap = Inches(0.15)
    start_x = lx
    start_y = ly + Inches(0.35)

    for i, w in enumerate(words_rnn):
        bx = start_x + i * (box_w + gap)
        s = add_rect(slide, bx, start_y, box_w, box_h, C_LIGHT_BLUE, C_ENCODER, Pt(1))
        set_shape_text(s, w, 10, C_ENCODER, True)
        if i < len(words_rnn) - 1:
            ax = bx + box_w
            add_arrow_line(slide, ax, start_y + box_h / 2, ax + gap, start_y + box_h / 2, C_DECODER, Pt(2))

    # Pain point annotations
    annotations = [
        ("顺序瓶颈", "串行处理,无法并行"),
        ("长距离遗忘", "早期信息指数级稀释"),
        ("信息瓶颈", "固定维度压缩所有信息"),
    ]
    ann_y = start_y + box_h + Inches(0.2)
    for i, (title, desc) in enumerate(annotations):
        ax = start_x + i * Inches(2.4)
        s = add_rect(slide, ax, ann_y, Inches(2.2), Inches(0.55), C_BG_DECODER_FILL, C_DECODER, Pt(1))
        tf = s.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = C_DECODER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_CN
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(8)
        p2.font.color.rgb = C_TEXT
        p2.font.name = FONT_CN
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(1)

    # === Self-attention visualization ===
    att_y = ann_y + Inches(0.75)
    add_text_box(slide, lx, att_y, Inches(4), Inches(0.35),
                 "自注意力可视化(以\"it\"为例):", 11, C_TITLE, True)

    att_words = ["The", "cat", "sat", "on", "the", "it", "..."]
    att_weights = [0.02, 0.60, 0.05, 0.01, 0.03, 0.00, 0.29]
    aw_w = Inches(0.85)
    aw_h = Inches(0.4)
    aw_gap = Inches(0.08)
    aw_start_x = lx + Inches(0.3)
    aw_start_y = att_y + Inches(0.35)

    # Draw attention lines first (behind boxes)
    it_index = 5
    it_center_x = aw_start_x + it_index * (aw_w + aw_gap) + aw_w / 2
    it_center_y = aw_start_y + aw_h / 2

    for i, w in enumerate(att_words):
        if i == it_index:
            continue
        w_center_x = aw_start_x + i * (aw_w + aw_gap) + aw_w / 2
        weight = att_weights[i]
        lw = Pt(1) if weight < 0.05 else Pt(2.5) if weight < 0.2 else Pt(4)
        color = C_GRAY if weight < 0.05 else C_ATTENTION if weight < 0.2 else C_DECODER
        # Draw curved line from it to word
        add_line(slide, w_center_x, aw_start_y + aw_h, it_center_x, aw_start_y + aw_h + Inches(0.4), color, lw)

    for i, w in enumerate(att_words):
        bx = aw_start_x + i * (aw_w + aw_gap)
        fill = C_BG_ATTENTION_FILL if i == it_index else C_LIGHT_BLUE
        lc = C_ATTENTION if i == it_index else C_ENCODER
        s = add_rect(slide, bx, aw_start_y, aw_w, aw_h, fill, lc, Pt(1.5))
        set_shape_text(s, w, 9, lc, i == it_index)

    # Weight labels below
    for i, w in enumerate(att_words):
        bx = aw_start_x + i * (aw_w + aw_gap)
        if att_weights[i] >= 0.05:
            add_text_box(slide, bx, aw_start_y + aw_h + Inches(0.42), aw_w, Inches(0.25),
                         str(att_weights[i]), 8, C_ATTENTION if att_weights[i] >= 0.2 else C_GRAY, True, PP_ALIGN.CENTER)

    # Legend box
    leg_y = aw_start_y + aw_h + Inches(0.75)
    add_rect(slide, lx, leg_y, Inches(6.5), Inches(0.7), C_LIGHT_GRAY, C_GRAY, Pt(0.5))
    add_text_box(slide, lx + Inches(0.1), leg_y + Inches(0.05), Inches(6.3), Inches(0.6),
                 "连线粗细 = 注意力权重: 粗线(cat=0.6)表示\"it\"与\"cat\"关联最强, 模型学会\"it\"指代\"cat\"",
                 9, C_TEXT, False, PP_ALIGN.LEFT)

    # === RIGHT SIDE: Text ===
    rx = Inches(7.8)
    ry = Inches(0.85)
    rw = Inches(5.2)
    rh = Inches(6.3)

    paragraphs = [
        ("\u2460 顺序处理的效率陷阱:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("RNN必须一个词一个词按顺序处理,每个词的处理都依赖前一个词的隐藏状态h_t。这意味着即使GPU有数千个并行计算单元,RNN也只能串行使用其中一个。处理一个1000词的句子需要1000个顺序步骤。就像雇了1000个工人搬砖,却规定每次只能1个人搬,效率浪费了99.9%。LSTM虽然引入了门控机制缓解梯度消失,但顺序处理限制完全没有改变。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 长距离遗忘(梯度消失):", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("处理到句子末尾时,早期信息经过数百次矩阵乘法被指数级稀释。每步更新h_t=tanh(W*h_{t-1}+U*x_t),如果W的特征值小于1,信息衰减速度极快。100步后可能只剩10^(-30)的权重。LSTM的门控机制部分缓解了这个问题,但超过100词的长文本中遗忘仍然严重。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 信息瓶颈:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("整个输入序列被压缩成一个固定长度的向量(256或512维),从这个向量生成输出。就像用100字摘要概括500页书,再根据摘要回答细节问题,不可能完成。信息瓶颈直接限制了模型处理长文本和理解复杂语义的能力。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 核心突破 - 自注意力机制:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("Transformer让每个词都能与序列中所有其他词直接交流。这像从\"排队发言\"变成\"圆桌会议\" - 任何人可以直接听到所有人。以\"it\"为例,自注意力让\"it\"与\"cat\"建立最强关联(权重0.6),\"it\"的表示向量融合了\"cat\"的信息,模型\"知道\"it指代的是cat。这个设计同时解决了三个问题:并行化解决顺序瓶颈,直接连接解决长距离遗忘,不压缩信息解决信息瓶颈。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 数学本质:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("自注意力的核心是Q*K^T(查询*键的转置),本质是高维空间中的点积/余弦相似度。点积越大,两个词的语义关联越强。这个操作一次性计算所有词对的关联度,形成T*T注意力矩阵,完全并行,无需顺序处理。这是Transformer的灵魂。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# PAGE 3: Architecture Overview
# ============================================================
def make_page3():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "架构总览:左边理解,右边生成", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === LEFT: Full architecture diagram ===
    lx = Inches(0.2)
    ly = Inches(0.8)

    # Input embedding area
    add_rect(slide, lx, ly, Inches(7.5), Inches(0.5), C_BG_EMBED_FILL, C_EMBED, Pt(1.5))
    set_shape_text(slide.shapes[-1], "Input Embedding (37000 x 512) + Positional Encoding (sin/cos)", 9, C_EMBED, True)

    # Encoder x6 big box
    enc_y = ly + Inches(0.65)
    enc_h = Inches(1.8)
    enc_outer = add_rect(slide, lx, enc_y, Inches(7.5), enc_h, None, C_ENCODER, Pt(2.5))
    enc_outer.fill.background()
    add_text_box(slide, lx + Inches(0.1), enc_y + Inches(0.02), Inches(2), Inches(0.25),
                 "Encoder x 6", 10, C_ENCODER, True)

    # Inside encoder: one layer detail
    inner_y = enc_y + Inches(0.3)
    inner_h = Inches(1.3)
    comp_w = Inches(1.6)
    comp_h = Inches(1.1)
    comp_gap = Inches(0.12)

    # Multi-Head Self-Attention
    bx = lx + Inches(0.3)
    s = add_rect(slide, bx, inner_y, comp_w, comp_h, C_BG_ATTENTION_FILL, C_ATTENTION, Pt(1))
    set_shape_text(s, "Multi-Head\nSelf-Attention\n(8 heads x 64d)", 8, C_TEXT, True)

    # Add & Norm
    bx2 = bx + comp_w + comp_gap
    s2 = add_rect(slide, bx2, inner_y, Inches(0.7), comp_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s2, "Add &\nLayer\nNorm", 8, C_FFN, True)
    add_arrow_line(slide, bx + comp_w, inner_y + comp_h/2, bx2, inner_y + comp_h/2, C_GRAY, Pt(1))

    # FFN
    bx3 = bx2 + Inches(0.7) + comp_gap
    s3 = add_rect(slide, bx3, inner_y, comp_w, comp_h, C_BG_FFN_FILL, C_FFN, Pt(1))
    set_shape_text(s3, "Feed-Forward\nNetwork\n(512->2048->512)", 8, C_TEXT, True)
    add_arrow_line(slide, bx2 + Inches(0.7), inner_y + comp_h/2, bx3, inner_y + comp_h/2, C_GRAY, Pt(1))

    # Add & Norm 2
    bx4 = bx3 + comp_w + comp_gap
    s4 = add_rect(slide, bx4, inner_y, Inches(0.7), comp_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s4, "Add &\nLayer\nNorm", 8, C_FFN, True)
    add_arrow_line(slide, bx3 + comp_w, inner_y + comp_h/2, bx4, inner_y + comp_h/2, C_GRAY, Pt(1))

    # Residual connections (green dashed lines above)
    add_line(slide, bx, inner_y - Inches(0.05), bx4 + Inches(0.7), inner_y - Inches(0.05), C_FFN, Pt(1))

    # K,V arrow from encoder to decoder
    kv_x = lx + Inches(7.6)
    kv_y = enc_y + enc_h / 2
    add_text_box(slide, kv_x, kv_y - Inches(0.1), Inches(0.5), Inches(0.5),
                 "K, V", 9, C_EMBED, True, PP_ALIGN.CENTER)
    add_arrow_line(slide, lx + Inches(7.5), kv_y, kv_x + Inches(0.45), kv_y, C_EMBED, Pt(2))

    # Decoder x6 big box
    dec_y = enc_y + enc_h + Inches(0.3)
    dec_h = Inches(2.2)
    dec_outer = add_rect(slide, lx, dec_y, Inches(7.5), dec_h, None, C_DECODER, Pt(2.5))
    dec_outer.fill.background()
    add_text_box(slide, lx + Inches(0.1), dec_y + Inches(0.02), Inches(2), Inches(0.25),
                 "Decoder x 6", 10, C_DECODER, True)

    # Inside decoder: 3 sub-layers
    d_inner_y = dec_y + Inches(0.3)
    d_comp_w = Inches(1.4)
    d_comp_h = Inches(0.7)
    d_gap = Inches(0.08)
    dx = lx + Inches(0.3)

    # Masked Self-Attention
    s = add_rect(slide, dx, d_inner_y, d_comp_w, d_comp_h, C_BG_ATTENTION_FILL, C_ATTENTION, Pt(1))
    set_shape_text(s, "Masked\nSelf-Attention", 8, C_TEXT, True)
    dx2 = dx + d_comp_w + d_gap
    s = add_rect(slide, dx2, d_inner_y, Inches(0.6), d_comp_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 7, C_FFN, True)
    add_arrow_line(slide, dx + d_comp_w, d_inner_y + d_comp_h/2, dx2, d_inner_y + d_comp_h/2, C_GRAY, Pt(1))

    # Cross-Attention
    dx3 = dx2 + Inches(0.6) + d_gap
    s = add_rect(slide, dx3, d_inner_y, d_comp_w, d_comp_h, C_BG_EMBED_FILL, C_EMBED, Pt(1))
    set_shape_text(s, "Cross-Attention\n(Q from Dec\nK,V from Enc)", 7, C_TEXT, True)
    add_arrow_line(slide, dx2 + Inches(0.6), d_inner_y + d_comp_h/2, dx3, d_inner_y + d_comp_h/2, C_GRAY, Pt(1))

    # K,V arrow into cross-attention
    add_arrow_line(slide, kv_x + Inches(0.45), kv_y, dx3 + d_comp_w/2, d_inner_y, C_EMBED, Pt(1.5))

    # Add & Norm 2
    dx4 = dx3 + d_comp_w + d_gap
    s = add_rect(slide, dx4, d_inner_y, Inches(0.6), d_comp_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 7, C_FFN, True)
    add_arrow_line(slide, dx3 + d_comp_w, d_inner_y + d_comp_h/2, dx4, d_inner_y + d_comp_h/2, C_GRAY, Pt(1))

    # FFN
    dx5 = dx4 + Inches(0.6) + d_gap
    s = add_rect(slide, dx5, d_inner_y, d_comp_w, d_comp_h, C_BG_FFN_FILL, C_FFN, Pt(1))
    set_shape_text(s, "Feed-Forward\n(512->2048->512)", 7, C_TEXT, True)
    add_arrow_line(slide, dx4 + Inches(0.6), d_inner_y + d_comp_h/2, dx5, d_inner_y + d_comp_h/2, C_GRAY, Pt(1))

    # Add & Norm 3
    dx6 = dx5 + d_comp_w + d_gap
    s = add_rect(slide, dx6, d_inner_y, Inches(0.6), d_comp_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 7, C_FFN, True)
    add_arrow_line(slide, dx5 + d_comp_w, d_inner_y + d_comp_h/2, dx6, d_inner_y + d_comp_h/2, C_GRAY, Pt(1))

    # Output area
    out_y = dec_y + dec_h + Inches(0.15)
    add_rect(slide, lx, out_y, Inches(2.5), Inches(0.45), C_LIGHT_BLUE, C_ENCODER, Pt(1))
    set_shape_text(slide.shapes[-1], "Linear (512 x 37000)", 9, C_ENCODER, True)
    add_arrow_line(slide, lx + Inches(3.7), out_y + Inches(0.22), lx + Inches(4.1), out_y + Inches(0.22), C_GRAY, Pt(1))
    add_rect(slide, lx + Inches(4.1), out_y, Inches(1.5), Inches(0.45), C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(slide.shapes[-1], "Softmax", 9, C_FFN, True)
    add_arrow_line(slide, lx + Inches(5.6), out_y + Inches(0.22), lx + Inches(6.0), out_y + Inches(0.22), C_GRAY, Pt(1))
    add_rect(slide, lx + Inches(6.0), out_y, Inches(1.5), Inches(0.45), C_BG_DECODER_FILL, C_DECODER, Pt(1))
    set_shape_text(slide.shapes[-1], "Output Prob", 9, C_DECODER, True)

    # Dimension annotations
    dim_y = out_y + Inches(0.55)
    add_text_box(slide, lx, dim_y, Inches(7.5), Inches(0.4),
                 "d_model=512  |  8 heads x 64d  |  FFN d_ff=2048  |  总参数约6500万", 9, C_GRAY, False)

    # === RIGHT SIDE: Text ===
    rx = Inches(8.0)
    ry = Inches(0.8)
    rw = Inches(5.0)
    rh = Inches(6.3)

    paragraphs = [
        ("\u2460 编码器 = 编辑部:", 10, C_ENCODER, True, PP_ALIGN.LEFT, Pt(4)),
        ("6层堆叠,每层含多头自注意力+前馈网络,每个子层有残差连接和层归一化。数据流经编码器时,每个词的信息与其他所有词不断融合,6层后每个词的512维向量已蕴含整个句子的上下文。编码器参数约3450万。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 解码器 = 翻译部:", 10, C_DECODER, True, PP_ALIGN.LEFT, Pt(4)),
        ("同样6层,但多一个交叉注意力子层。掩码自注意力确保只能看之前的词(不能偷看答案),交叉注意力让解码器\"查阅\"编码器的输出(带着问题查资料)。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 数据流全景:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("分词(BPE) -> 嵌入(37000x512矩阵查表) -> 位置编码(sin/cos注入) -> 编码器6层(交流+提炼) -> 编码器输出 -> 解码器6层(多一步查资料) -> 线性变换(512->37000) -> Softmax -> 概率 -> 输出词。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 关键数字:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("编码器6层、解码器6层、512维、8头、FFN 2048维、总参数6500万、WMT 2014英德翻译、450万句对、8xP100 GPU、12小时、BLEU 28.4。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 编码器单层结构:", 10, C_ENCODER, True, PP_ALIGN.LEFT, Pt(4)),
        ("输入(512维) -> 多头自注意力(8头x64维) -> 残差+归一化(512维) -> FFN(512->2048->ReLU->512) -> 残差+归一化(512维) -> 输出。两条残差旁路(绿色Add标注)确保梯度稳定传播。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),

        ("\u2465 解码器单层结构:", 10, C_DECODER, True, PP_ALIGN.LEFT, Pt(4)),
        ("掩码自注意力(蓝) -> 残差+归一化 -> 交叉注意力(紫,Q来自解码器,KV来自编码器) -> 残差+归一化 -> FFN(橙) -> 残差+归一化。比编码器每层多一个交叉注意力子层,用于融合源语言信息。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# PAGE 4: Encoder Detail
# ============================================================
def make_page4():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "编码器详解:每层都在交流与提炼", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === LEFT: Encoder single layer data flow ===
    lx = Inches(0.3)
    ly = Inches(0.9)

    # Data flow: horizontal pipeline
    flow_y = ly + Inches(0.4)
    box_w = Inches(1.35)
    box_h = Inches(1.0)
    gap = Inches(0.2)
    sx = lx

    # Input
    s = add_rect(slide, sx, flow_y, Inches(0.8), box_h, C_LIGHT_BLUE, C_ENCODER, Pt(1))
    set_shape_text(s, "输入\n512d", 9, C_ENCODER, True)
    sx2 = sx + Inches(0.8) + gap
    add_arrow_line(slide, sx + Inches(0.8), flow_y + box_h/2, sx2, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Multi-Head Self-Attention
    s = add_rect(slide, sx2, flow_y, box_w, box_h, C_BG_ATTENTION_FILL, C_ATTENTION, Pt(1.5))
    set_shape_text(s, "多头自注意力\n8头 x 64d\nQ,K,V线性变换", 8, C_TEXT, True)
    sx3 = sx2 + box_w + gap
    add_arrow_line(slide, sx2 + box_w, flow_y + box_h/2, sx3, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Residual + LayerNorm
    s = add_rect(slide, sx3, flow_y, Inches(0.9), box_h, C_LIGHT_GREEN, C_FFN, Pt(1.5))
    set_shape_text(s, "残差\n连接\n+ LayerNorm", 8, C_FFN, True)
    sx4 = sx3 + Inches(0.9) + gap
    add_arrow_line(slide, sx3 + Inches(0.9), flow_y + box_h/2, sx4, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # FFN
    s = add_rect(slide, sx4, flow_y, box_w, box_h, C_BG_FFN_FILL, C_FFN, Pt(1.5))
    set_shape_text(s, "前馈网络\n512->2048\nReLU->512", 8, C_TEXT, True)
    sx5 = sx4 + box_w + gap
    add_arrow_line(slide, sx4 + box_w, flow_y + box_h/2, sx5, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Residual + LayerNorm 2
    s = add_rect(slide, sx5, flow_y, Inches(0.9), box_h, C_LIGHT_GREEN, C_FFN, Pt(1.5))
    set_shape_text(s, "残差\n连接\n+ LayerNorm", 8, C_FFN, True)
    sx6 = sx5 + Inches(0.9) + gap
    add_arrow_line(slide, sx5 + Inches(0.9), flow_y + box_h/2, sx6, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Output
    s = add_rect(slide, sx6, flow_y, Inches(0.8), box_h, C_LIGHT_BLUE, C_ENCODER, Pt(1))
    set_shape_text(s, "输出\n512d", 9, C_ENCODER, True)

    # Residual bypass lines (green)
    bypass_y = flow_y - Inches(0.25)
    add_line(slide, lx + Inches(0.8), bypass_y, sx3 + Inches(0.45), bypass_y, C_FFN, Pt(1.5))
    add_line(slide, lx + Inches(0.8), bypass_y, lx + Inches(0.8), flow_y, C_FFN, Pt(1))
    add_line(slide, sx3 + Inches(0.45), bypass_y, sx3 + Inches(0.45), flow_y, C_FFN, Pt(1))

    bypass_y2 = flow_y + box_h + Inches(0.2)
    add_line(slide, sx3 + Inches(0.9), bypass_y2, sx5 + Inches(0.45), bypass_y2, C_FFN, Pt(1.5))
    add_line(slide, sx3 + Inches(0.9), flow_y + box_h, sx3 + Inches(0.9), bypass_y2, C_FFN, Pt(1))
    add_line(slide, sx5 + Inches(0.45), flow_y + box_h, sx5 + Inches(0.45), bypass_y2, C_FFN, Pt(1))

    # Residual labels
    add_text_box(slide, lx + Inches(0.1), bypass_y - Inches(0.3), Inches(3), Inches(0.25),
                 "残差旁路1: output = input + SubLayer(input)", 8, C_FFN, True)
    add_text_box(slide, sx3, bypass_y2 + Inches(0.05), Inches(3), Inches(0.25),
                 "残差旁路2: output = input + FFN(input)", 8, C_FFN, True)

    # === Multi-head attention detail diagram ===
    mh_y = flow_y + box_h + Inches(0.7)
    add_text_box(slide, lx, mh_y, Inches(5), Inches(0.3),
                 "多头注意力内部结构(8头并行):", 10, C_TITLE, True)

    mh_inner_y = mh_y + Inches(0.35)
    # 8 small head boxes
    head_w = Inches(0.75)
    head_h = Inches(0.55)
    head_gap = Inches(0.08)
    head_labels = ["Head 1\n语法", "Head 2\n语义", "Head 3\n位置", "Head 4\n指代",
                   "Head 5\n...", "Head 6\n...", "Head 7\n...", "Head 8\n..."]
    for i, label in enumerate(head_labels):
        hx = lx + i * (head_w + head_gap)
        s = add_rect(slide, hx, mh_inner_y, head_w, head_h, C_BG_ATTENTION_FILL, C_ATTENTION, Pt(0.5))
        set_shape_text(s, label, 7, C_TEXT, False)

    # Concat arrow
    concat_y = mh_inner_y + head_h + Inches(0.1)
    add_text_box(slide, lx + Inches(1.5), concat_y, Inches(4), Inches(0.3),
                 "Concat(8 x 64 = 512) -> Linear -> Output", 9, C_TEXT, False, PP_ALIGN.CENTER)

    # LayerNorm detail
    ln_y = concat_y + Inches(0.45)
    add_rect(slide, lx, ln_y, Inches(7), Inches(0.55), C_LIGHT_GREEN, C_FFN, Pt(1))
    s = slide.shapes[-1]
    set_shape_text(s, "LayerNorm: y = gamma * (x - mu) / sigma + beta  |  gamma, beta 各512维, 每层1024个可学习参数  |  稳定训练,防止梯度爆炸/消失", 8, C_TEXT, False)

    # === RIGHT SIDE: Text ===
    rx = Inches(7.8)
    ry = Inches(0.9)
    rw = Inches(5.2)
    rh = Inches(6.2)

    paragraphs = [
        ("\u2460 多头自注意力 = 8个讨论组:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("单头只能从一个角度分析,多头用8个独立头学习不同关系。头1专注语法(主谓关系),头2关注语义(同义词),头3关注位置(相邻词),头4关注指代(代词-名词)。每个头独立计算QKV(各64维),8个结果拼接(8x64=512)后线性变换融合。多头让模型同时关注多种关系,大幅提升表达能力。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 残差连接 = 高速公路匝道:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("output = input + SubLayer(input)。从梯度看:dLoss/dx = dLoss/dOutput x (dSubLayer/dx + 1),加号后面的\"1\"确保梯度永远不会消失。这是Transformer能堆叠多层的关键。没有残差连接,深层网络会因为梯度消失而无法训练。每层有两个残差连接(注意力后+FFN后),形成信息高速通道。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 层归一化 = 标准化体检:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("y = gamma x (x - mu) / sigma + beta,对512维向量独立做标准化。没有归一化时,数值范围可能从[-0.01,0.01]突变到[-100,100],导致梯度爆炸/消失,训练无法收敛。gamma和beta各512维,每层1024个可学习参数。归一化不改变信息,只调整尺度,让每层输入分布稳定。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 前馈网络 = 独立思考空间:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("FFN(x) = ReLU(xW1 + b1)W2 + b2,维度512->2048->ReLU->512。参数量约210万(512x2048x2),比自注意力(约100万)大2倍,是计算量最大的组件。注意力负责词间交流,FFN负责每个词自身的深层特征提取。两层线性变换中间的非线性ReLU让模型能学习复杂函数。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 6层堆叠的威力:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("单层只能学习局部模式(相邻词关系),6层堆叠后逐步抽象:第1层学习局部语法,第2-3层学习短语级语义,第4-5层学习句子级逻辑,第6层学习全局上下文。越深层捕捉越抽象的关系。6层不是随意选择,实验发现6层时性能趋于饱和,更多层收益递减。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# PAGE 5: Decoder Detail
# ============================================================
def make_page5():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "解码器详解:比编码器多了一座翻译桥梁", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === LEFT: Decoder single layer structure ===
    lx = Inches(0.3)
    ly = Inches(0.9)

    # Horizontal flow
    flow_y = ly + Inches(0.4)
    box_w = Inches(1.3)
    box_h = Inches(0.9)
    gap = Inches(0.15)
    sx = lx

    # Input
    s = add_rect(slide, sx, flow_y, Inches(0.7), box_h, C_LIGHT_BLUE, C_ENCODER, Pt(1))
    set_shape_text(s, "输入\n(已生成词)", 8, C_ENCODER, True)
    sx2 = sx + Inches(0.7) + gap
    add_arrow_line(slide, sx + Inches(0.7), flow_y + box_h/2, sx2, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Masked Self-Attention
    s = add_rect(slide, sx2, flow_y, box_w, box_h, C_BG_ATTENTION_FILL, C_ATTENTION, Pt(1.5))
    set_shape_text(s, "掩码自注意力\n(Masked\nSelf-Attn)", 8, C_TEXT, True)
    sx3 = sx2 + box_w + gap
    add_arrow_line(slide, sx2 + box_w, flow_y + box_h/2, sx3, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Add & Norm
    s = add_rect(slide, sx3, flow_y, Inches(0.65), box_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 8, C_FFN, True)
    sx4 = sx3 + Inches(0.65) + gap
    add_arrow_line(slide, sx3 + Inches(0.65), flow_y + box_h/2, sx4, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Cross-Attention (highlighted)
    s = add_rect(slide, sx4, flow_y, box_w, box_h, C_BG_EMBED_FILL, C_EMBED, Pt(2))
    set_shape_text(s, "交叉注意力\nCross-Attn\n(Q Dec, KV Enc)", 7, C_TEXT, True)
    sx5 = sx4 + box_w + gap
    add_arrow_line(slide, sx4 + box_w, flow_y + box_h/2, sx5, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Add & Norm 2
    s = add_rect(slide, sx5, flow_y, Inches(0.65), box_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 8, C_FFN, True)
    sx6 = sx5 + Inches(0.65) + gap
    add_arrow_line(slide, sx5 + Inches(0.65), flow_y + box_h/2, sx6, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # FFN
    s = add_rect(slide, sx6, flow_y, box_w, box_h, C_BG_FFN_FILL, C_FFN, Pt(1.5))
    set_shape_text(s, "FFN\n512->2048\n->512", 8, C_TEXT, True)
    sx7 = sx6 + box_w + gap
    add_arrow_line(slide, sx6 + box_w, flow_y + box_h/2, sx7, flow_y + box_h/2, C_GRAY, Pt(1.5))

    # Add & Norm 3
    s = add_rect(slide, sx7, flow_y, Inches(0.65), box_h, C_LIGHT_GREEN, C_FFN, Pt(1))
    set_shape_text(s, "Add\n&\nNorm", 8, C_FFN, True)

    # Encoder K,V arrow into Cross-Attention
    enc_kv_y = flow_y - Inches(0.6)
    enc_s = add_rect(slide, sx4 - Inches(0.3), enc_kv_y, Inches(2.0), Inches(0.4), C_BG_ENCODER_FILL, C_ENCODER, Pt(1.5))
    set_shape_text(enc_s, "编码器输出 (K, V)", 9, C_ENCODER, True)
    add_arrow_line(slide, sx4 + box_w/2, enc_kv_y + Inches(0.4), sx4 + box_w/2, flow_y, C_EMBED, Pt(2))

    # Causal mask matrix visualization
    mask_y = flow_y + box_h + Inches(0.5)
    add_text_box(slide, lx, mask_y, Inches(4), Inches(0.3),
                 "因果掩码矩阵(Causal Mask):", 10, C_TITLE, True)

    # 4x4 matrix
    mat_y = mask_y + Inches(0.35)
    cell_size = Inches(0.6)
    for r in range(4):
        for c in range(4):
            cx = lx + c * cell_size
            cy = mat_y + r * cell_size
            if c <= r:
                fill = C_BG_ATTENTION_FILL
                lc = C_ATTENTION
                label = "1"
            else:
                fill = C_BG_DECODER_FILL
                lc = C_DECODER
                label = "-inf"
            s = add_rect(slide, cx, cy, cell_size - Inches(0.02), cell_size - Inches(0.02), fill, lc, Pt(0.5))
            set_shape_text(s, label, 8, lc, True)

    # Labels for mask
    add_text_box(slide, lx + Inches(0.1), mat_y + Inches(0.1), Inches(0.5), Inches(0.3),
                 "词1", 7, C_TEXT, False, PP_ALIGN.CENTER)
    add_text_box(slide, lx + cell_size + Inches(0.1), mat_y + Inches(0.1), Inches(0.5), Inches(0.3),
                 "词2", 7, C_TEXT, False, PP_ALIGN.CENTER)
    add_text_box(slide, lx + 2*cell_size + Inches(0.1), mat_y + Inches(0.1), Inches(0.5), Inches(0.3),
                 "词3", 7, C_TEXT, False, PP_ALIGN.CENTER)
    add_text_box(slide, lx + 3*cell_size + Inches(0.1), mat_y + Inches(0.1), Inches(0.5), Inches(0.3),
                 "词4", 7, C_TEXT, False, PP_ALIGN.CENTER)

    # Explanation below matrix
    add_text_box(slide, lx, mat_y + 4 * cell_size + Inches(0.1), Inches(4), Inches(0.6),
                 "黄色=可见(下三角), 红色=-inf(上三角)\nSoftmax后-inf变0, 确保每个词只能看自己和之前的词",
                 8, C_TEXT, False)

    # Cross-attention QKV explanation
    ca_y = mat_y + 4 * cell_size + Inches(0.8)
    add_rect(slide, lx, ca_y, Inches(7), Inches(0.7), C_BG_EMBED_FILL, C_EMBED, Pt(1))
    set_shape_text(slide.shapes[-1], "交叉注意力: Q来自解码器(需要找什么) x K^T来自编码器(有什么特征) -> Softmax -> x V(编码器的实际内容)\nQ和K的序列长度可以不同(源语言和目标语言长度不同), 不像自注意力是T x T方阵", 8, C_TEXT, False)

    # === RIGHT SIDE: Text ===
    rx = Inches(7.8)
    ry = Inches(0.9)
    rw = Inches(5.2)
    rh = Inches(6.2)

    paragraphs = [
        ("\u2460 掩码自注意力 = 考试不能偷看后面的题:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("因果掩码将注意力矩阵上三角设为负无穷,Softmax后变为0。处理\"爱\"时只能看到[START,我,爱],看不到\"你\"。训练和推理都必须严格执行。训练时如果偷看答案,推理时没有答案就会崩溃。这保证了训练和推理行为一致,避免了\"暴露偏差\"问题。掩码矩阵是一个下三角全1、上三角全负无穷的方阵,加到注意力分数上即可。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 交叉注意力 = 带着问题去查资料:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("Query来自解码器(\"我在翻译,需要找原文信息\"),Key和Value来自编码器(\"编辑对原文的理解\")。Q x K^T -> 缩放 -> Softmax -> x V。生成\"爱\"时,Query\"找情感动词\"与\"love\"的Key匹配度最高(0.7),大量提取Value信息辅助生成。注意:Q和K的序列长度可以不同(源语言和目标语言长度不同),不像自注意力是T x T方阵。交叉注意力是编码器和解码器之间唯一的桥梁。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 完整层结构:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("掩码自注意力 + 残差归一化 -> 交叉注意力 + 残差归一化 -> FFN + 残差归一化。解码器6层,每层3个子层 + 3个残差归一化,比编码器每层多一个交叉注意力子层。三个子层各司其职:掩码自注意力处理已生成序列,交叉注意力融合源语言信息,FFN做特征提取。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 与编码器的对比:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("编码器:自注意力+FFN,2个子层/层,可以看到完整输入,双向理解。解码器:掩码自注意力+交叉注意力+FFN,3个子层/层,只能看到已生成部分,单向生成。编码器像\"通读全文做笔记\",解码器像\"根据笔记逐字翻译\"。两者共享相似的子层结构(都有残差和归一化),但信息流方向不同。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 为什么需要掩码:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("自注意力天然会看到所有位置,包括未来位置。如果不加掩码,训练时模型能看到目标词(答案),导致训练和推理不一致。训练时\"作弊\"看答案学到的模式,推理时没有答案就无法使用。掩码让训练和推理使用相同的信息可见性,这是自回归生成模型的核心设计原则。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# PAGE 6: Training Overview
# ============================================================
def make_page6():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "训练阶段:模型如何从数据中学会语言?", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === TOP: Training data flow diagram ===
    lx = Inches(0.3)
    ly = Inches(0.8)

    # Input sentence
    s = add_rect(slide, lx, ly, Inches(2.0), Inches(0.5), C_BG_ENCODER_FILL, C_ENCODER, Pt(1.5))
    set_shape_text(s, "Input: \"I love you\"", 9, C_ENCODER, True)

    add_arrow_line(slide, lx + Inches(2.0), ly + Inches(0.25), lx + Inches(2.3), ly + Inches(0.25), C_GRAY, Pt(1.5))

    # Encoder
    s = add_rect(slide, lx + Inches(2.3), ly, Inches(2.5), Inches(0.5), C_BG_ENCODER_FILL, C_ENCODER, Pt(2))
    set_shape_text(s, "Encoder x 6", 10, C_ENCODER, True)

    add_arrow_line(slide, lx + Inches(4.8), ly + Inches(0.25), lx + Inches(5.1), ly + Inches(0.25), C_EMBED, Pt(1.5))

    # Encoder output with K,V
    s = add_rect(slide, lx + Inches(5.1), ly, Inches(1.8), Inches(0.5), C_BG_EMBED_FILL, C_EMBED, Pt(1.5))
    set_shape_text(s, "编码器输出(K,V)", 9, C_EMBED, True)

    # Decoder input
    dec_y = ly + Inches(0.7)
    s = add_rect(slide, lx, dec_y, Inches(2.0), Inches(0.5), C_LIGHT_GREEN, C_FFN, Pt(1.5))
    set_shape_text(s, "Decoder Input: [START,我,爱]", 8, C_FFN, True)

    add_arrow_line(slide, lx + Inches(2.0), dec_y + Inches(0.25), lx + Inches(2.3), dec_y + Inches(0.25), C_GRAY, Pt(1.5))

    # Decoder
    s = add_rect(slide, lx + Inches(2.3), dec_y, Inches(2.5), Inches(0.5), C_BG_DECODER_FILL, C_DECODER, Pt(2))
    set_shape_text(s, "Decoder x 6", 10, C_DECODER, True)

    add_arrow_line(slide, lx + Inches(4.8), dec_y + Inches(0.25), lx + Inches(5.1), dec_y + Inches(0.25), C_GRAY, Pt(1.5))

    # Linear + Softmax
    s = add_rect(slide, lx + Inches(5.1), dec_y, Inches(1.8), Inches(0.5), C_BG_ATTENTION_FILL, C_ATTENTION, Pt(1.5))
    set_shape_text(s, "Linear + Softmax", 9, C_ATTENTION, True)

    add_arrow_line(slide, lx + Inches(6.9), dec_y + Inches(0.25), lx + Inches(7.2), dec_y + Inches(0.25), C_GRAY, Pt(1.5))

    # Prediction
    s = add_rect(slide, lx + Inches(7.2), dec_y, Inches(1.5), Inches(0.5), C_BG_FFN_FILL, C_FFN, Pt(1.5))
    set_shape_text(s, "预测值", 9, C_FFN, True)

    # Target (truth)
    s = add_rect(slide, lx + Inches(9.0), dec_y, Inches(2.2), Inches(0.5), C_BG_DECODER_FILL, C_DECODER, Pt(1.5))
    set_shape_text(s, "真实目标: [我,爱,你,END]", 8, C_DECODER, True)

    # Loss arrow between prediction and target
    loss_y = dec_y + Inches(0.6)
    add_text_box(slide, lx + Inches(7.5), loss_y, Inches(3), Inches(0.35),
                 "<-- 交叉熵损失(Cross-Entropy Loss) -->", 9, C_DECODER, True, PP_ALIGN.CENTER)

    # Note about parallel computation
    note_y = dec_y + Inches(1.0)
    add_rect(slide, lx, note_y, Inches(12), Inches(0.35), C_LIGHT_YELLOW, C_ATTENTION, Pt(0.5))
    set_shape_text(slide.shapes[-1], "训练核心: 所有位置同时并行计算(Teacher Forcing), 解码器一次性看到完整目标序列, GPU充分利用并行能力", 9, C_ATTENTION, False)

    # Embedding matrix detail
    emb_y = note_y + Inches(0.5)
    add_rect(slide, lx, emb_y, Inches(5.8), Inches(0.5), C_BG_EMBED_FILL, C_EMBED, Pt(1))
    set_shape_text(slide.shapes[-1], "嵌入矩阵: 37000 x 512 = 约1890万参数(占模型30%) + 位置编码: PE(pos,2i) = sin(pos/10000^(2i/d))", 8, C_EMBED, False)

    # Training loop detail
    loop_y = emb_y + Inches(0.6)
    add_rect(slide, lx, loop_y, Inches(5.8), Inches(0.5), C_LIGHT_BLUE, C_ENCODER, Pt(1))
    set_shape_text(slide.shapes[-1], "训练循环: 前向传播 -> 交叉熵Loss -> 反向传播(自动微分) -> Adam更新6500万参数 -> 重复30万步", 8, C_ENCODER, False)

    # === RIGHT SIDE: Text ===
    rx = Inches(7.8)
    ry = Inches(0.8)
    rw = Inches(5.2)
    rh = Inches(6.3)

    paragraphs = [
        ("\u2460 训练 = 开卷考试:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("核心特点是所有位置并行计算。解码器同时看到完整输入[START,我,爱,你](Teacher Forcing),并行预测每个位置的下一个词。GPU同时处理所有位置,充分利用并行能力。这和推理的逐词串行形成鲜明对比。训练时我们有标准答案,可以同时计算所有位置的损失。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 嵌入矩阵 = 超级字典:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("BPE子词分词 -> 37000x512嵌入矩阵查表(约1890万参数,占模型30%)。每个词映射为512维向量,编码语义信息。位置编码用sin/cos函数:PE(pos,2i)=sin(pos/10000^(2i/d)),给每个词贴\"座位号\",不增加可训练参数。嵌入+位置编码=最终输入向量。嵌入矩阵是模型最大的单个参数组件。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 训练目标 = 学会预测下一个词:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("根据[START]预测\"我\",根据[START,我]预测\"爱\",根据[START,我,爱]预测\"你\",根据[START,我,爱,你]预测\"<END>\"。每个位置的交叉熵损失取平均,反向传播更新6500万参数。训练初期Loss约等于log(37000)约等于10.5(随机猜测),充分训练后降到1.5-3.0(预测较准)。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 编码器和解码器的分工:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("编码器只处理输入,每批训练只运行一次,输出(K,V)缓存供解码器使用。解码器处理目标序列,每层通过交叉注意力查阅编码器输出。编码器专注\"理解\",解码器专注\"生成\",各司其职。这种分工让两个模块各自优化自己的任务。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 训练配置与优化:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("WMT 2014英德翻译,450万句对,批量25000 token,Adam优化器(beta1=0.9,beta2=0.98),学习率峰值0.0003(前4000步线性增温后余弦衰减),标签平滑0.1,dropout 0.3,8xP100训练12小时(约30万步),BLEU 28.4。学习率增温策略是训练成功的关键:太大会发散,太小会收敛极慢。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# PAGE 7: Self-Attention 5 Steps
# ============================================================
def make_page7():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55),
                 "自注意力的5步计算:从输入到输出", 16, C_TITLE, True, PP_ALIGN.LEFT)

    # === LEFT: 5-step vertical flow ===
    lx = Inches(0.3)
    ly = Inches(0.85)

    step_h = Inches(1.05)
    step_gap = Inches(0.15)
    box_w = Inches(7.0)

    steps = [
        ("Step 1: 生成QKV", "X(512d) -> W_Q/W_K/W_V(各512x64) -> Q,K,V(各64d)", C_BG_ATTENTION_FILL, C_ATTENTION, "所有词同时计算,完全并行"),
        ("Step 2: 计算分数矩阵", "Q x K^T -> T x T 分数矩阵 (Score[i,j]=词i对词j的关注度)", C_LIGHT_BLUE, C_ENCODER, "本质是高维空间中的余弦相似度"),
        ("Step 3: 缩放", "Score / sqrt(d_k) = Score / sqrt(64) = Score / 8", C_BG_ATTENTION_FILL, C_ATTENTION, "防止点积过大导致Softmax饱和(梯度=0)"),
        ("Step 4: Softmax归一化", "exp(x_i) / Sum(exp(x_j)), 每行和=1", C_LIGHT_GREEN, C_FFN, "如\"it\"的分布: [0.02, 0.60, 0.05, 0.01, 0.03, 0.29]"),
        ("Step 5: 加权求和输出", "Attention_weights x V -> 每个词的输出是所有Value的加权平均", C_BG_EMBED_FILL, C_EMBED, "8头各自计算,拼接后线性变换,最终输出512维"),
    ]

    for i, (title, desc, fill, lc, note) in enumerate(steps):
        sy = ly + i * (step_h + step_gap)

        # Step number circle-like box
        num_s = add_rect(slide, lx, sy, Inches(0.5), step_h, lc, None, Pt(0))
        set_shape_text(num_s, str(i+1), 16, C_WHITE, True)

        # Main box
        s = add_rect(slide, lx + Inches(0.6), sy, box_w - Inches(0.6), step_h, fill, lc, Pt(1.5))
        tf = s.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = lc
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_CN
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = C_TEXT
        p2.font.name = FONT_CN
        p2.space_before = Pt(3)
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.size = Pt(8)
        p3.font.color.rgb = C_GRAY
        p3.font.name = FONT_CN
        p3.space_before = Pt(2)

        # Arrow to next step
        if i < len(steps) - 1:
            arrow_y = sy + step_h
            add_arrow_line(slide, lx + Inches(0.85), arrow_y, lx + Inches(0.85), arrow_y + step_gap, C_GRAY, Pt(1.5))

    # Small 3x3 attention matrix illustration
    mat_x = lx + Inches(5.0)
    mat_y = ly + Inches(1.2)
    add_text_box(slide, mat_x, mat_y - Inches(0.3), Inches(2), Inches(0.25),
                 "3x3注意力矩阵示意:", 8, C_TITLE, True)
    cell = Inches(0.45)
    mat_data = [
        ["0.5", "0.3", "0.2"],
        ["0.1", "0.7", "0.2"],
        ["0.1", "0.1", "0.8"],
    ]
    for r in range(3):
        for c in range(3):
            cx = mat_x + c * cell
            cy = mat_y + r * cell
            val = float(mat_data[r][c])
            if val >= 0.5:
                fill = C_BG_ATTENTION_FILL
            elif val >= 0.2:
                fill = C_LIGHT_YELLOW
            else:
                fill = C_LIGHT_GRAY
            s = add_rect(slide, cx, cy, cell - Inches(0.02), cell - Inches(0.02), fill, C_GRAY, Pt(0.3))
            set_shape_text(s, mat_data[r][c], 8, C_TEXT, False)

    # === RIGHT SIDE: Text ===
    rx = Inches(7.8)
    ry = Inches(0.85)
    rw = Inches(5.2)
    rh = Inches(6.3)

    paragraphs = [
        ("\u2460 步骤1 - 生成QKV:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("每个词的512维向量X分别乘以W_Q、W_K、W_V(各512x64),得到64维的Q、K、V。Query=\"我在找什么信息\"(如\"it\"的Q编码了\"需要找名词\"),Key=\"我有什么特征\"(如\"cat\"的K编码了\"名词、动物、单数\"),Value=\"我的实际内容\"(被提取传递的信息)。所有词同时计算,完全并行,单头参数量=512x64x3=98304。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2461 步骤2 - 计算分数矩阵:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("Q x K^T得到T x T矩阵。Score[i,j]表示第i个词对第j个词的关注度。本质是高维空间中的余弦相似度。矩阵中每行代表一个词对所有词的关注分数。对角线(自注意力)通常较大,但其他位置也可能很大(如\"it\"对\"cat\")。这个矩阵是注意力的核心产物。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2462 步骤3 - 缩放:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("分数除以sqrt(64)=8。当d_k较大时,点积方差为d_k,Softmax会进入饱和区(梯度约等于0)。除以sqrt(d_k)将方差归一化为1。像考试分数从百分制转十分制,保持相对大小但避免极端值。不缩放时,最大值和最小值差距可能达到数百,Softmax会变成one-hot(梯度消失)。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2463 步骤4 - Softmax归一化:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("exp(x_i)/Sum(exp(x_j)),每行和=1。例如\"it\"的注意力分布:[0.02, 0.60, 0.05, 0.01, 0.03, 0.29],\"cat\"最高(0.60)。Softmax放大最大值但不完全置零其他值,保留一定的信息多样性。温度参数可以调节分布的尖锐程度。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(6)),

        ("\u2464 步骤5 - 加权求和:", 10, C_TITLE, True, PP_ALIGN.LEFT, Pt(4)),
        ("Output = Attention_weights x V。每个词的输出是所有Value的加权平均。以\"it\"为例:0.02xV(\"The\")+0.60xV(\"cat\")+0.05xV(\"sat\")+... -> \"it\"的表示融合了\"cat\"的大量信息。8个头各自独立计算,拼接后通过线性变换(8x64=512->512),最终输出512维。多头的多样性让不同头关注不同的关系模式。", 10, C_TEXT, False, PP_ALIGN.LEFT, Pt(4)),
    ]
    add_rich_text_box(slide, rx, ry, rw, rh, paragraphs)


# ============================================================
# Generate all pages
# ============================================================
make_page1()
make_page2()
make_page3()
make_page4()
make_page5()
make_page6()
make_page7()

# Save
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v4_part1.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")

# Statistics
for i, slide in enumerate(prs.slides, 1):
    shape_count = len(slide.shapes)
    total_text = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                total_text += len(p.text)
    print(f"Page {i}: {shape_count} shapes, {total_text} chars of text")
