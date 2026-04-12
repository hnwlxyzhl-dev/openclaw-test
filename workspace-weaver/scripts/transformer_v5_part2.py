#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT v5 - Pages 8-13 (Chinese)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import math

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_R = Inches(13.2)
SAFE_B = Inches(7.3)

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
C_LIGHT_GREEN = RGBColor(0xE8, 0xF8, 0xF5)
C_LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xE7)
C_LIGHT_PURPLE = RGBColor(0xF4, 0xEC, 0xF7)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_DARK_GRAY = RGBColor(0x7F, 0x8C, 0x8D)

# Fonts
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank layout


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # MSO_SHAPE.ROUNDED_RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(11),
                 font_color=C_TEXT, bold=False, font_name=FONT_CN,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add a text box with text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    # Set line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=Pt(11),
                      default_color=C_TEXT, default_bold=False, default_font=FONT_CN,
                      alignment=PP_ALIGN.LEFT, line_spacing=1.15):
    """Add text box with multiple paragraphs. lines = [(text, size, color, bold, font), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold, font = line_data, default_size, default_color, default_bold, default_font
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else default_bold
            font = line_data[4] if len(line_data) > 4 else default_font
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = alignment
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        # line spacing
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox


def add_arrow_right(slide, left, top, width, height, color=C_TEXT):
    """Add right arrow."""
    shape = slide.shapes.add_shape(13, left, top, width, height)  # right arrow
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_page_number(slide, num):
    """Add page number at bottom right."""
    add_text_box(slide, Inches(12.4), Inches(7.0), Inches(0.6), Inches(0.3),
                 str(num), Pt(9), C_GRAY, False, FONT_EN, PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text):
    """Add title at top with accent line."""
    add_shape(slide, Inches(0), Inches(0), Inches(13.2), Inches(0.06), C_TITLE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
                 title_text, Pt(18), C_TITLE, True, FONT_CN)


def count_chars(text):
    """Count Chinese chars + other chars."""
    return len(text)


def draw_matrix(slide, left, top, cell_size, n, values, header=None):
    """Draw an n×n matrix. values = list of (row, col, color, text)"""
    # Draw grid
    total_w = cell_size * n
    total_h = cell_size * n
    # Background
    add_shape(slide, left - Inches(0.05), top - Inches(0.05),
              total_w + Inches(0.1), total_h + Inches(0.1), C_LIGHT_BG, C_DARK_GRAY, Pt(1))
    
    for r in range(n):
        for c in range(n):
            x = left + c * cell_size
            y = top + r * cell_size
            # Determine color from values
            bg = None
            txt = ""
            txt_color = C_TEXT
            for v in values:
                if v[0] == r and v[1] == c:
                    bg = v[2]
                    txt = v[3] if len(v) > 3 else ""
                    txt_color = v[4] if len(v) > 4 else C_TEXT
                    break
            if bg is None:
                bg = C_WHITE
            add_shape(slide, x, y, cell_size, cell_size, bg, C_DARK_GRAY, Pt(0.5))
            if txt:
                add_text_box(slide, x, y, cell_size, cell_size, txt,
                             Pt(10), txt_color, False, FONT_EN, PP_ALIGN.CENTER)


def draw_flow_box(slide, left, top, width, height, text, fill_color, text_color=C_WHITE, font_size=Pt(9)):
    """Draw a rounded box with centered text."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color, None)
    add_text_box(slide, left, top, width, height, text, font_size, text_color, True, FONT_CN, PP_ALIGN.CENTER)
    return shape


# ══════════════════════════════════════════════════════════════
# PAGE 8: Training Details
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_title_bar(slide8, "训练细节：安全网、掩码与损失")

# LEFT 60%: Causal mask matrix + Loss flow
# Causal mask 4x4
matrix_left = Inches(0.5)
matrix_top = Inches(1.2)
cell = Inches(0.7)
add_text_box(slide8, matrix_left, matrix_top - Inches(0.35), Inches(3), Inches(0.3),
             "Causal Mask (4x4)", Pt(11), C_TITLE, True, FONT_EN)

mask_values = [
    (0, 0, C_LIGHT_BLUE, "0"), (0, 1, C_LIGHT_RED, "-inf", C_DECODER), (0, 2, C_LIGHT_RED, "-inf", C_DECODER), (0, 3, C_LIGHT_RED, "-inf", C_DECODER),
    (1, 0, C_LIGHT_BLUE, "0"), (1, 1, C_LIGHT_BLUE, "0"), (1, 2, C_LIGHT_RED, "-inf", C_DECODER), (1, 3, C_LIGHT_RED, "-inf", C_DECODER),
    (2, 0, C_LIGHT_BLUE, "0"), (2, 1, C_LIGHT_BLUE, "0"), (2, 2, C_LIGHT_BLUE, "0"), (2, 3, C_LIGHT_RED, "-inf", C_DECODER),
    (3, 0, C_LIGHT_BLUE, "0"), (3, 1, C_LIGHT_BLUE, "0"), (3, 2, C_LIGHT_BLUE, "0"), (3, 3, C_LIGHT_BLUE, "0"),
]
draw_matrix(slide8, matrix_left, matrix_top, cell, 4, mask_values)

# After softmax annotation
add_text_box(slide8, matrix_left, matrix_top + cell * 4 + Inches(0.1), Inches(3), Inches(0.3),
             "Softmax -> 0 (red cells)", Pt(9), C_DECODER, False, FONT_EN)

# Loss flow diagram (bottom left)
flow_top = Inches(4.2)
add_text_box(slide8, matrix_left, flow_top - Inches(0.35), Inches(7), Inches(0.3),
             "Loss Flow", Pt(11), C_TITLE, True, FONT_EN)

box_w = Inches(1.1)
box_h = Inches(0.55)
gap = Inches(0.35)
flow_x = Inches(0.5)
flow_y = flow_top

boxes = [
    ("Decoder\n512d", C_ENCODER),
    ("Linear\n512x37K", C_EMBED),
    ("Softmax", C_ATTENTION),
    ("Loss", C_DECODER),
    ("Backprop", C_FFN),
]
for i, (txt, color) in enumerate(boxes):
    x = flow_x + i * (box_w + gap)
    draw_flow_box(slide8, x, flow_y, box_w, box_h, txt, color, C_WHITE, Pt(9))
    if i < len(boxes) - 1:
        ax = x + box_w
        ay = flow_y + box_h // 2 - Inches(0.1)
        add_arrow_right(slide8, ax, ay, gap - Inches(0.05), Inches(0.2), C_DARK_GRAY)

# "Update Parameters" label
add_text_box(slide8, flow_x + 4 * (box_w + gap) - Inches(0.1), flow_y + box_h + Inches(0.1),
             Inches(1.3), Inches(0.3), "Update 65M params", Pt(8), C_TEXT, False, FONT_EN, PP_ALIGN.CENTER)

# RIGHT 40%: 4 bullet points
right_x = Inches(7.8)
right_w = Inches(5.0)
bullet_top = Inches(1.0)

bullets8 = [
    ("Residual Connection", "output = x + SubLayer(x)", "Gradient has \"+1\" ensuring no vanishing. Key to stacking 6+ layers."),
    ("Causal Mask", "Upper triangle = -inf -> Softmax = 0", "Training/inference behave identically. Prevents exposure bias."),
    ("Cross-Entropy Loss", "Loss = -log(p(correct word))", "p=0.99 -> 0.01 (good); p=0.01 -> 4.6 (bad). Average over all positions."),
    ("Training Config", "Adam, lr=3e-4, dropout=0.3", "Label smoothing 0.1, 8xP100 GPUs, 12 hours, BLEU 28.4."),
]

for i, (title, formula, desc) in enumerate(bullets8):
    y = bullet_top + i * Inches(1.5)
    # Card background
    add_rounded_rect(slide8, right_x, y, right_w, Inches(1.35), C_LIGHT_BG, C_ENCODER, Pt(0.5))
    # Color bar
    add_shape(slide8, right_x, y, Inches(0.08), Inches(1.35), C_ENCODER)
    # Title
    add_text_box(slide8, right_x + Inches(0.2), y + Inches(0.05), right_w - Inches(0.3), Inches(0.25),
                 title, Pt(12), C_TITLE, True, FONT_CN)
    # Formula
    add_text_box(slide8, right_x + Inches(0.2), y + Inches(0.32), right_w - Inches(0.3), Inches(0.25),
                 formula, Pt(10), C_ATTENTION, True, FONT_EN)
    # Description
    add_text_box(slide8, right_x + Inches(0.2), y + Inches(0.6), right_w - Inches(0.3), Inches(0.7),
                 desc, Pt(10), C_TEXT, False, FONT_CN, line_spacing=1.2)

add_page_number(slide8, 8)

# Count chars page 8
p8_text = "训练细节安全网掩码与损失残差连接output等于x加SubLayer梯度中有加一确保永不消失因果掩码上三角负无穷Softmax等于零训练推理行为一致避免暴露偏差交叉熵损失Loss等于负log正确词概率点九九零点零一好零点零一四点六差训练配置Adam优化器学习率零点零零零三dropout零点三标签平滑零点一八块P100十二小时"
print(f"Page 8: ~{len(p8_text)} chars")


# ══════════════════════════════════════════════════════════════
# PAGE 9: Inference + KV Cache
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_title_bar(slide9, "推理阶段：闭卷考试与KV缓存")

# LEFT 60%: Autoregressive steps (top) + cache comparison (bottom)
auto_top = Inches(1.0)
add_text_box(slide9, Inches(0.5), auto_top, Inches(7), Inches(0.3),
             "Autoregressive Generation", Pt(11), C_TITLE, True, FONT_EN)

# Step boxes
steps = [
    ("[START]", C_GRAY),
    ("Step1\n[START]", C_LIGHT_BLUE),
    ("I love", C_ENCODER),
    ("Step2\n[START, I]", C_LIGHT_BLUE),
    ("you", C_FFN),
    ("Step3\n[START, I, love]", C_LIGHT_BLUE),
    ("[END]", C_ATTENTION),
]
step_w = Inches(1.0)
step_h = Inches(0.7)
step_gap = Inches(0.25)
start_x = Inches(0.3)
step_y = auto_top + Inches(0.35)

for i, (txt, color) in enumerate(steps):
    x = start_x + i * (step_w + step_gap)
    if "Step" in txt:
        draw_flow_box(slide9, x, step_y, step_w, step_h, txt, color, C_TEXT, Pt(8))
    else:
        draw_flow_box(slide9, x, step_y, step_w, step_h, txt, color, C_WHITE, Pt(10))
    if i < len(steps) - 1 and "Step" not in steps[i+1][0]:
        ax = x + step_w
        ay = step_y + step_h // 2 - Inches(0.08)
        add_arrow_right(slide9, ax, ay, step_gap - Inches(0.05), Inches(0.16), C_DARK_GRAY)

# Encoder runs once label
add_text_box(slide9, Inches(0.3), step_y + step_h + Inches(0.05), Inches(7), Inches(0.25),
             "Encoder runs ONCE -> cache output for entire generation", Pt(9), C_ENCODER, True, FONT_EN)

# Bottom: No cache vs Cache comparison
comp_top = Inches(3.8)
add_text_box(slide9, Inches(0.5), comp_top, Inches(7), Inches(0.3),
             "Without Cache vs With Cache", Pt(11), C_TITLE, True, FONT_EN)

# Without cache
nocache_x = Inches(0.3)
nocache_y = comp_top + Inches(0.35)
add_rounded_rect(slide9, nocache_x, nocache_y, Inches(3.2), Inches(2.8), C_LIGHT_RED, C_DECODER, Pt(1))
add_text_box(slide9, nocache_x, nocache_y - Inches(0.3), Inches(3.2), Inches(0.25),
             "Without KV Cache", Pt(10), C_DECODER, True, FONT_EN, PP_ALIGN.CENTER)

# Step 2 without cache - recompute
nocache_steps = [
    ("Step2: Recompute ALL K,V", C_DECODER),
    ("[START] K,V  (duplicate!)", C_LIGHT_RED),
    ("[I] K,V  (duplicate!)", C_LIGHT_RED),
    ("[love] K,V  (new)", C_LIGHT_GREEN),
    ("O(T x d x T) per step", C_TEXT),
]
for j, (txt, color) in enumerate(nocache_steps):
    add_text_box(slide9, nocache_x + Inches(0.15), nocache_y + Inches(0.15) + j * Inches(0.5),
                 Inches(2.9), Inches(0.45), txt, Pt(9), color, j == 0 or j == 4, FONT_EN)

# With cache
cache_x = Inches(3.8)
cache_y = nocache_y
add_rounded_rect(slide9, cache_x, cache_y, Inches(3.2), Inches(2.8), C_LIGHT_GREEN, C_FFN, Pt(1))
add_text_box(slide9, cache_x, cache_y - Inches(0.3), Inches(3.2), Inches(0.25),
             "With KV Cache", Pt(10), C_FFN, True, FONT_EN, PP_ALIGN.CENTER)

cache_steps = [
    ("Step2: Only compute new K,V", C_FFN),
    ("[START] K,V  (cached)", C_ENCODER),
    ("[I] K,V  (cached)", C_ENCODER),
    ("[love] K,V  (new, append!)", C_LIGHT_GREEN),
    ("O(d x T) per step", C_TEXT),
]
for j, (txt, color) in enumerate(cache_steps):
    add_text_box(slide9, cache_x + Inches(0.15), cache_y + Inches(0.15) + j * Inches(0.5),
                 Inches(2.9), Inches(0.45), txt, Pt(9), color, j == 0 or j == 4, FONT_EN)

# RIGHT 40%: 4 bullets
right_x = Inches(7.8)
right_w = Inches(5.0)

bullets9 = [
    ("Inference = Closed-Book Exam", "Serial word-by-word generation. Each step depends on previous. First mistake cascades."),
    ("Encoder Runs Once", "Output cached, reused for entire generation. Cost independent of sequence length."),
    ("Serial Bottleneck", "100 words = 100 forward passes. Training does it in 1 parallel pass."),
    ("KV Cache = Meeting Minutes", "Only compute new word K,V, append to cache. O(d*T) vs O(T^2). Memory: ~375MB for 12-layer."),
]

for i, (title, desc) in enumerate(bullets9):
    y = Inches(1.0) + i * Inches(1.5)
    add_rounded_rect(slide9, right_x, y, right_w, Inches(1.35), C_LIGHT_BG, C_ENCODER, Pt(0.5))
    add_shape(slide9, right_x, y, Inches(0.08), Inches(1.35), C_ENCODER)
    add_text_box(slide9, right_x + Inches(0.2), y + Inches(0.05), right_w - Inches(0.3), Inches(0.25),
                 title, Pt(12), C_TITLE, True, FONT_CN)
    add_text_box(slide9, right_x + Inches(0.2), y + Inches(0.35), right_w - Inches(0.3), Inches(0.9),
                 desc, Pt(10), C_TEXT, False, FONT_CN, line_spacing=1.2)

add_page_number(slide9, 9)

p9_text = "推理阶段闭卷考试与KV缓存推理等于闭卷考试逐词串行生成每步依赖前一步第一步错则后续全错编码器只运行一次输出缓存整个生成过程反复使用串行瓶颈一百词等于一百次前向传播训练只需一次KV缓存等于会议纪要每步只算新词KV追加缓存平方复杂度线性复杂度"
print(f"Page 9: ~{len(p9_text)} chars")


# ══════════════════════════════════════════════════════════════
# PAGE 10: Word Selection Strategies
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
add_title_bar(slide10, "词选择：从贪心到创意")

# LEFT 55%: 5 strategy cards
card_x = Inches(0.3)
card_w = Inches(6.8)
card_h = Inches(1.05)
card_gap = Inches(0.15)

strategies = [
    ("1. Greedy", "Always pick highest probability. Deterministic but may loop. Local optimum != global.", C_ENCODER),
    ("2. Beam Search", "Keep k best paths (beam=4-10). Good for translation but tends toward bland output.", C_ATTENTION),
    ("3. Top-K Sampling", "Sample from top K words (K=50). Balances quality and diversity.", C_FFN),
    ("4. Top-P (Nucleus)", "Smallest set reaching cumulative P=0.9. More adaptive than fixed K.", C_EMBED),
    ("5. Temperature", "T<1 (0.3): sharp/deterministic. T=1: original. T>1 (1.5): flat/random.", C_DECODER),
]

for i, (title, desc, color) in enumerate(strategies):
    y = Inches(1.0) + i * (card_h + card_gap)
    add_rounded_rect(slide10, card_x, y, card_w, card_h, C_LIGHT_BG, color, Pt(0.5))
    add_shape(slide10, card_x, y, Inches(0.08), card_h, color)
    add_text_box(slide10, card_x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.3), Inches(0.25),
                 title, Pt(12), color, True, FONT_EN)
    add_text_box(slide10, card_x + Inches(0.2), y + Inches(0.38), card_w - Inches(0.3), Inches(0.6),
                 desc, Pt(10), C_TEXT, False, FONT_EN, line_spacing=1.15)

# RIGHT 45%: 3 curves + 2 bullets
right_x = Inches(7.3)
right_w = Inches(5.5)

# Draw 3 distribution curves as text description with colored bars
dist_top = Inches(1.0)
add_text_box(slide10, right_x, dist_top, right_w, Inches(0.3),
             "Probability Distributions", Pt(11), C_TITLE, True, FONT_EN)

# Low temperature - sharp
add_text_box(slide10, right_x, dist_top + Inches(0.4), right_w, Inches(0.25),
             "T=0.3 (Low) - Sharp peak, one word dominates", Pt(9), C_ENCODER, False, FONT_EN)
bar_y = dist_top + Inches(0.65)
bars_low = [0.05, 0.1, 0.15, 0.6, 0.05, 0.03, 0.02]
for j, v in enumerate(bars_low):
    bw = Inches(v * 4)
    add_shape(slide10, right_x + j * Inches(0.55), bar_y, bw, Inches(0.25), C_ENCODER)

# Medium temperature
add_text_box(slide10, right_x, bar_y + Inches(0.35), right_w, Inches(0.25),
             "T=0.7 (Medium) - Balanced", Pt(9), C_ATTENTION, False, FONT_EN)
bar_y2 = bar_y + Inches(0.6)
bars_med = [0.08, 0.12, 0.2, 0.35, 0.12, 0.08, 0.05]
for j, v in enumerate(bars_med):
    bw = Inches(v * 4)
    add_shape(slide10, right_x + j * Inches(0.55), bar_y2, bw, Inches(0.25), C_ATTENTION)

# High temperature
add_text_box(slide10, right_x, bar_y2 + Inches(0.35), right_w, Inches(0.25),
             "T=1.5 (High) - Flat, more random", Pt(9), C_DECODER, False, FONT_EN)
bar_y3 = bar_y2 + Inches(0.6)
bars_high = [0.12, 0.14, 0.15, 0.18, 0.15, 0.14, 0.12]
for j, v in enumerate(bars_high):
    bw = Inches(v * 4)
    add_shape(slide10, right_x + j * Inches(0.55), bar_y3, bw, Inches(0.25), C_DECODER)

# 2 key bullets
bullet10_top = Inches(4.8)
add_rounded_rect(slide10, right_x, bullet10_top, right_w, Inches(1.1), C_LIGHT_BG, C_ENCODER, Pt(0.5))
add_shape(slide10, right_x, bullet10_top, Inches(0.08), Inches(1.1), C_ENCODER)
add_text_box(slide10, right_x + Inches(0.2), bullet10_top + Inches(0.05), right_w - Inches(0.3), Inches(0.25),
             "ChatGPT Strategy", Pt(12), C_TITLE, True, FONT_EN)
add_text_box(slide10, right_x + Inches(0.2), bullet10_top + Inches(0.35), right_w - Inches(0.3), Inches(0.7),
             "Temperature=0.7 + Top-P=0.9 + Repetition penalty. Code: T=0.2; Creative: T=0.7-1.0.",
             Pt(10), C_TEXT, False, FONT_EN)

add_rounded_rect(slide10, right_x, bullet10_top + Inches(1.3), right_w, Inches(1.1), C_LIGHT_BG, C_DECODER, Pt(0.5))
add_shape(slide10, right_x, bullet10_top + Inches(1.3), Inches(0.08), Inches(1.1), C_DECODER)
add_text_box(slide10, right_x + Inches(0.2), bullet10_top + Inches(1.35), right_w - Inches(0.3), Inches(0.25),
             "Greedy Trap", Pt(12), C_TITLE, True, FONT_EN)
add_text_box(slide10, right_x + Inches(0.2), bullet10_top + Inches(0.65), right_w - Inches(0.3), Inches(0.7),
             "Local optimum != global optimum. Open-domain dialogue needs diversity. No single 'best' answer.",
             Pt(10), C_TEXT, False, FONT_EN)

add_page_number(slide10, 10)

p10_text = "词选择从贪心到创意贪心搜索选最高概率确定但可能循环束搜索保留条路径翻译好但偏平庸TopK前个随机选平衡质量和多样性TopP累计达最小词集比固定更灵活温度低于一确定高于一随机ChatGPT策略温度加TopP加重复惩罚贪心陷阱局部最优不等于全局最优开放域需要多样性"
print(f"Page 10: ~{len(p10_text)} chars")


# ══════════════════════════════════════════════════════════════
# PAGE 11: Training vs Inference Comparison
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank_layout)
add_title_bar(slide11, "训练 vs 推理")

# TOP 70%: Comparison table (7 rows x 3 cols)
table_top = Inches(0.9)
table_left = Inches(0.5)
col_w = [Inches(2.5), Inches(4.8), Inches(4.8)]
row_h = Inches(0.7)
header_h = Inches(0.55)

# Table header
headers = ["Dimension", "Training (Open-Book)", "Inference (Closed-Book)"]
header_colors = [C_TITLE, C_ENCODER, C_DECODER]

cx = table_left
for i, (hdr, color) in enumerate(zip(headers, header_colors)):
    add_shape(slide11, cx, table_top, col_w[i], header_h, color)
    add_text_box(slide11, cx, table_top, col_w[i], header_h, hdr,
                 Pt(12), C_WHITE, True, FONT_CN, PP_ALIGN.CENTER)
    cx += col_w[i]

# Table rows
rows = [
    ("Input", "Full target sequence (Teacher Forcing)", "Model generates word by word"),
    ("Computation", "All positions parallel (1 pass)", "Serial word-by-word (T passes)"),
    ("Loss", "Cross-entropy at each position", "No loss, only prediction"),
    ("Gradient", "Backprop updates 65M params", "None (params frozen)"),
    ("KV Cache", "Not needed", "Required"),
    ("GPU Utilization", "High (fully parallel)", "Low (serial wait)"),
    ("Goal", "Learn to predict", "Generate text"),
]

for r, (dim, train, infer) in enumerate(rows):
    ry = table_top + header_h + r * row_h
    bg = C_LIGHT_BG if r % 2 == 0 else C_WHITE
    cx = table_left
    # Dimension
    add_shape(slide11, cx, ry, col_w[0], row_h, bg, C_DARK_GRAY, Pt(0.5))
    add_text_box(slide11, cx + Inches(0.1), ry, col_w[0] - Inches(0.2), row_h, dim,
                 Pt(11), C_TITLE, True, FONT_CN, PP_ALIGN.LEFT)
    cx += col_w[0]
    # Training
    add_shape(slide11, cx, ry, col_w[1], row_h, C_LIGHT_BLUE, C_DARK_GRAY, Pt(0.5))
    add_text_box(slide11, cx + Inches(0.1), ry, col_w[1] - Inches(0.2), row_h, train,
                 Pt(10), C_TEXT, False, FONT_CN, PP_ALIGN.LEFT)
    cx += col_w[1]
    # Inference
    add_shape(slide11, cx, ry, col_w[2], row_h, C_LIGHT_GREEN, C_DARK_GRAY, Pt(0.5))
    add_text_box(slide11, cx + Inches(0.1), ry, col_w[2] - Inches(0.2), row_h, infer,
                 Pt(10), C_TEXT, False, FONT_CN, PP_ALIGN.LEFT)

# BOTTOM 30%: Metaphor
meta_top = Inches(6.2)
add_rounded_rect(slide11, Inches(0.5), meta_top, Inches(12.2), Inches(0.9), C_LIGHT_YELLOW, C_ATTENTION, Pt(1))
add_text_box(slide11, Inches(0.8), meta_top + Inches(0.08), Inches(11.6), Inches(0.7),
             "Training = new employee training (with answer key). Inference = on the job (on your own). "
             "Core contradiction: Exposure Bias.",
             Pt(11), C_TEXT, False, FONT_EN, PP_ALIGN.CENTER)

add_page_number(slide11, 11)

p11_text = "训练对比推理输入完整目标序列Teacher Forcing逐词生成计算所有位置并行一次逐词串行次损失交叉熵无梯度反向传播更新参数无冻结缓存不需要必须利用率高低目标学会预测生成文本训练等于新员工培训有答案参考推理等于正式上岗全靠自己核心矛盾暴露偏差"
print(f"Page 11: ~{len(p11_text)} chars")


# ══════════════════════════════════════════════════════════════
# PAGE 12: Applications
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(blank_layout)
add_title_bar(slide12, "从翻译工具到AI通用引擎")

# LEFT 60%: Tree diagram
tree_top = Inches(1.0)
# Root node
root_x = Inches(2.5)
root_y = tree_top
root_w = Inches(3.0)
root_h = Inches(0.6)
draw_flow_box(slide12, root_x, root_y, root_w, root_h, "Transformer 2017", C_TITLE, C_WHITE, Pt(13))

# Three branches
branch_y = root_y + root_h + Inches(0.6)
branch_w = Inches(2.8)
branch_h = Inches(0.5)
branch_gap = Inches(0.4)
branch_start_x = Inches(0.3)

branches = [
    ("Decoder Only (GPT)", C_ENCODER, ["GPT-1/2/3/4", "ChatGPT, GLM, Qwen", "Generation King"]),
    ("Full Arch (T5)", C_FFN, ["T5, BART", "Text-to-Text", "Conversion King"]),
    ("Encoder Only (BERT)", C_EMBED, ["BERT, RoBERTa", "DeBERTa", "Understanding King"]),
]

for i, (name, color, children) in enumerate(branches):
    bx = branch_start_x + i * (branch_w + Inches(0.3))
    draw_flow_box(slide12, bx, branch_y, branch_w, branch_h, name, color, C_WHITE, Pt(10))
    # Connector line
    mid_x = bx + branch_w // 2
    add_shape(slide12, mid_x, root_y + root_h, Inches(0.03), Inches(0.6), color)
    # Children
    for j, child in enumerate(children):
        cy = branch_y + branch_h + Inches(0.1) + j * Inches(0.4)
        child_color = C_TEXT if j < 2 else color
        child_bold = j == 2
        add_text_box(slide12, bx + Inches(0.1), cy, branch_w - Inches(0.2), Inches(0.35),
                     child, Pt(9), child_color, child_bold, FONT_EN, PP_ALIGN.CENTER)

# Timeline at bottom
timeline_y = Inches(4.5)
add_text_box(slide12, Inches(0.5), timeline_y, Inches(7.5), Inches(0.3),
             "Beyond Language: Sequence is the Universal Format", Pt(11), C_TITLE, True, FONT_EN)

timeline_items = [
    ("ViT\n2020", "Image", C_ENCODER),
    ("Whisper\n2022", "Speech", C_FFN),
    ("GPT-4V\n2023", "Multi-modal", C_EMBED),
    ("Sora\n2024", "Video", C_DECODER),
    ("AlphaFold\nProtein", "Biology", C_ATTENTION),
]

tl_x = Inches(0.3)
tl_gap = Inches(1.3)
for i, (name, domain, color) in enumerate(timeline_items):
    tx = tl_x + i * tl_gap
    draw_flow_box(slide12, tx, timeline_y + Inches(0.4), Inches(1.15), Inches(0.8), name, color, C_WHITE, Pt(8))
    add_text_box(slide12, tx, timeline_y + Inches(1.25), Inches(1.15), Inches(0.25),
                 domain, Pt(8), color, True, FONT_EN, PP_ALIGN.CENTER)

# Timeline arrow
add_shape(slide12, Inches(0.3), timeline_y + Inches(1.6), Inches(6.5), Inches(0.03), C_DARK_GRAY)

# RIGHT 40%: 4 bullets
right_x = Inches(7.8)
right_w = Inches(5.0)

bullets12 = [
    ("GPT (Generation)", "Predict next word. ChatGPT = Pretrain + Instruction tuning + RLHF. 175B -> 1.8T params.", C_ENCODER),
    ("BERT (Understanding)", "Masked Language Modeling (mask 15%). Good at classification, QA, search.", C_EMBED),
    ("T5 (Conversion)", "All NLP tasks unified as Text-to-Text. Translation, summary, classification.", C_FFN),
    ("Beyond Language", "ViT (image) -> Whisper (speech) -> AlphaFold2 (protein). Sequence = universal format.", C_ATTENTION),
]

for i, (title, desc, color) in enumerate(bullets12):
    y = Inches(1.0) + i * Inches(1.5)
    add_rounded_rect(slide12, right_x, y, right_w, Inches(1.35), C_LIGHT_BG, color, Pt(0.5))
    add_shape(slide12, right_x, y, Inches(0.08), Inches(1.35), color)
    add_text_box(slide12, right_x + Inches(0.2), y + Inches(0.05), right_w - Inches(0.3), Inches(0.25),
                 title, Pt(12), C_TITLE, True, FONT_EN)
    add_text_box(slide12, right_x + Inches(0.2), y + Inches(0.35), right_w - Inches(0.3), Inches(0.9),
                 desc, Pt(10), C_TEXT, False, FONT_EN, line_spacing=1.2)

add_page_number(slide12, 12)

p12_text = "从翻译工具到AI通用引擎GPT生成预测下一个词预训练指令微调RLHF从亿到万亿参数BERT理解掩码语言建模遮盖擅长分类问答搜索T5转换所有任务统一为文本到文本翻译摘要分类超越语言图像语音蛋白质序列等于通用格式"
print(f"Page 12: ~{len(p12_text)} chars")


# ══════════════════════════════════════════════════════════════
# PAGE 13: Summary
# ══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(blank_layout)
add_title_bar(slide13, "总结：Transformer的核心要点")

# 5 summary cards
cards = [
    ("1", "Information Flows Freely", "Self-attention lets any word talk to any word directly. Solves all 3 RNN problems.", C_ENCODER, "360deg"),
    ("2", "Encoder Understands\nDecoder Generates", "Connected via K, V. Division of labor: read first, then write.", C_DECODER, "gear"),
    ("3", "Train Open-Book\nInference Closed-Book", "Parallel (fast) vs serial (KV cache). Core tension: exposure bias.", C_ATTENTION, "clock"),
    ("4", "Three Families", "GPT=Generate, BERT=Understand, T5=Convert. Same DNA, 3 careers.", C_FFN, "tree"),
    ("5", "Universal AI Engine", "65M -> 1T params, architecture unchanged. The foundation of the AI era.", C_EMBED, "star"),
]

card_w = Inches(2.2)
card_h = Inches(2.8)
card_gap = Inches(0.24)
card_start_x = Inches(0.5)
card_y = Inches(1.1)

for i, (num, title, desc, color, icon) in enumerate(cards):
    cx = card_start_x + i * (card_w + card_gap)
    # Card background
    add_rounded_rect(slide13, cx, card_y, card_w, card_h, C_LIGHT_BG, color, Pt(1))
    # Top color band
    add_shape(slide13, cx, card_y, card_w, Inches(0.5), color)
    # Number circle
    circle_size = Inches(0.35)
    circle_x = cx + (card_w - circle_size) // 2
    add_shape(slide13, circle_x, card_y + Inches(0.07), circle_size, circle_size, C_WHITE)
    add_text_box(slide13, circle_x, card_y + Inches(0.07), circle_size, circle_size,
                 num, Pt(14), color, True, FONT_EN, PP_ALIGN.CENTER)
    # Title
    add_text_box(slide13, cx + Inches(0.1), card_y + Inches(0.6), card_w - Inches(0.2), Inches(0.7),
                 title, Pt(11), C_TITLE, True, FONT_EN, PP_ALIGN.CENTER)
    # Description
    add_text_box(slide13, cx + Inches(0.15), card_y + Inches(1.4), card_w - Inches(0.3), Inches(1.3),
                 desc, Pt(9.5), C_TEXT, False, FONT_EN, PP_ALIGN.CENTER, line_spacing=1.25)

# Bottom quote
quote_y = Inches(4.3)
add_rounded_rect(slide13, Inches(1.5), quote_y, Inches(10.3), Inches(0.7), C_LIGHT_PURPLE, C_EMBED, Pt(1))
add_text_box(slide13, Inches(1.8), quote_y + Inches(0.1), Inches(9.7), Inches(0.5),
             "Transformer is not a model. It is an era.",
             Pt(14), C_EMBED, True, FONT_EN, PP_ALIGN.CENTER)

# Parameter evolution bar
evo_y = Inches(5.3)
add_text_box(slide13, Inches(0.5), evo_y, Inches(12), Inches(0.3),
             "Parameter Evolution", Pt(10), C_TITLE, True, FONT_EN)

params = [
    ("2017\n65M", Inches(1.5), C_ENCODER),
    ("2019\n1.5B", Inches(3.5), C_FFN),
    ("2020\n175B", Inches(6.0), C_ATTENTION),
    ("2023\n1.8T", Inches(9.0), C_DECODER),
    ("2025\n???", Inches(11.0), C_EMBED),
]

for label, x_pos, color in params:
    add_shape(slide13, x_pos, evo_y + Inches(0.4), Inches(0.08), Inches(0.08), color)
    add_text_box(slide13, x_pos - Inches(0.5), evo_y + Inches(0.55), Inches(1.2), Inches(0.5),
                 label, Pt(8), color, True, FONT_EN, PP_ALIGN.CENTER)

# Arrow line
add_shape(slide13, Inches(1.6), evo_y + Inches(0.44), Inches(9.5), Inches(0.03), C_DARK_GRAY)

add_page_number(slide13, 13)

p13_text = "总结Transformer核心要点信息自由流动自注意力让任意词直接交流解决三大问题编码器理解解码器生成通过连接分工协作训练开卷推理闭卷并行高效串行缓存加速核心矛盾暴露偏差三大家族生成理解转换同一DNA三种职业通用引擎从万到万亿参数架构未变AI时代基石Transformer不是一个模型而是一个时代"
print(f"Page 13: ~{len(p13_text)} chars")


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v5_part2.pptx"
prs.save(output_path)
print(f"\nSaved to: {output_path}")
print("Done! 6 slides (pages 8-13) generated.")
