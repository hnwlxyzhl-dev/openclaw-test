#!/usr/bin/env python3
"""Transformer PPT v7 Part1: Pages 1-7 — SD v10 design quality"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = "output/transformer_v7_part1.pptx"

# Colors
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_ENC = RGBColor(0x34,0x98,0xDB)
C_DEC = RGBColor(0xE7,0x4C,0x3C)
C_ATT = RGBColor(0xF3,0x9C,0x12)
C_FFN = RGBColor(0x2E,0xCC,0x71)
C_EMB = RGBColor(0x9B,0x59,0xB6)
C_LT = RGBColor(0xEB,0xF5,0xFB)
C_W = RGBColor(0xFF,0xFF,0xFF)
C_D = RGBColor(0x2C,0x3E,0x50)
C_G = RGBColor(0x88,0x88,0x88)
C_BG_ENC = RGBColor(0xE8,0xF0,0xFE)
C_BG_DEC = RGBColor(0xFD,0xED,0xED)
C_BG_ATT = RGBColor(0xFE,0xF9,0xE7)
C_BG_FFN = RGBColor(0xE8,0xF8,0xF5)
C_BG_EMB = RGBColor(0xF4,0xEC,0xF7)
C_BG_GRN = RGBColor(0xE8,0xF8,0xF5)
C_BG_ORG = RGBColor(0xFD,0xF2,0xE9)

FC = "微软雅黑"
FE = "Arial"

def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc; sh.line.width=Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc; sh.line.width=Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    for r in p.runs: r.font.size=Pt(fs); r.font.color.rgb=fc; r.font.bold=b; r.font.name=fn
    p.font.size=Pt(fs); p.font.color.rgb=fc; p.font.bold=b; p.font.name=fn
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
        for r in p.runs: r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=bld; r.font.name=fn
        p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld; p.font.name=fn
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def line_h(s,l,t,w,c=C_ACC,lw=2):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Emu(int(lw*12700)))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Transformer 架构深度解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    ok = True
    for sh in s.shapes:
        if sh.left < Emu(-100) or sh.top < Emu(-100): ok = False
        if sh.left+sh.width > SLIDE_W+Emu(5000) or sh.top+sh.height > SLIDE_H+Emu(5000): ok = False
    print(f"  {'OK' if ok else 'FAIL'} {msg}")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ==================== P1 封面 ====================
print("P1: 封面")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_PRI)
rect(sl, Inches(0.8), Inches(2.8), Inches(3.0), Inches(0.06), fc=C_ACC)
tb(sl, Inches(0.8), Inches(3.0), Inches(11.0), Inches(1.20),
   "Transformer", fs=44, fc=C_W, b=True, fn="Arial", ls=1.0)
tb(sl, Inches(0.8), Inches(4.0), Inches(11.0), Inches(0.60),
   "架构深度解析 — 从直觉到原理", fs=24, fc=RGBColor(0xBB,0xDE,0xFB), ls=1.0)
rect(sl, Inches(0.8), Inches(4.6), Inches(1.5), Inches(0.04), fc=C_ACC)
tb(sl, Inches(0.8), Inches(5.0), Inches(8.0), Inches(0.40),
   "一篇看懂AI的核心引擎：从训练到推理的完整拆解", fs=13, fc=RGBColor(0x99,0xBB,0xDD))
tb(sl, Inches(0.8), Inches(5.5), Inches(8.0), Inches(0.40),
   "Attention Is All You Need — Vaswani et al., 2017 | Google Brain | 6500万参数", fs=12, fc=RGBColor(0x88,0xAA,0xCC))
rect(sl, Inches(10.0), Inches(0.0), Inches(3.33), Inches(7.50), fc=RGBColor(0x17,0x2E,0x4A))
tb(sl, Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.50),
   "NATURAL\nLANGUAGE\nPROCESSING", fs=18, fc=C_ACC, b=True, fn="Arial", a=PP_ALIGN.CENTER, ls=1.3)
chk(sl, "P1")

# ==================== P2 目录 ====================
print("P2: 目录")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "目录", "CONTENTS — 11 CHAPTERS")
toc = [
    ("01", "为什么需要Transformer", "RNN的三大痛点与注意力突破"),
    ("02", "架构总览", "编码器-解码器结构与数据流全景"),
    ("03", "编码器详解", "多头自注意力、残差连接与层归一化"),
    ("04", "解码器详解", "掩码注意力、交叉注意力与因果掩码"),
    ("05", "训练阶段", "教师强制、嵌入矩阵与位置编码"),
    ("06", "自注意力5步计算", "从QKV生成到加权求和的完整过程"),
    ("07", "训练细节", "残差数学、因果掩码与交叉熵损失"),
    ("08", "推理阶段与KV缓存", "闭卷考试、串行生成与缓存加速"),
    ("09", "词选择策略", "从贪心搜索到温度采样的5种方法"),
    ("10", "训练vs推理全面对比", "开卷考试与闭卷考试的8维对比"),
    ("11", "应用与总结", "三大家族与AI通用引擎的未来"),
]
row_h = 0.40
for i, (num, title, desc) in enumerate(toc):
    y = Inches(1.30 + i * row_h)
    if i % 2 == 0:
        rect(sl, Inches(0.40), y, Inches(12.50), Inches(row_h), fc=C_LT)
    tb(sl, Inches(0.60), y, Inches(0.50), Inches(row_h), num, fs=11, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(1.20), y, Inches(4.20), Inches(row_h), title, fs=11, fc=C_PRI, b=False)
    tb(sl, Inches(5.60), y, Inches(7.00), Inches(row_h), desc, fs=11, fc=C_G)
chk(sl, "P2")

# ==================== P3 为什么需要Transformer ====================
print("P3: 为什么需要Transformer")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "01  为什么需要Transformer", "RNN的三大痛点与注意力机制的突破")

# RNN痛点卡片
rrect(sl, Inches(0.30), Inches(1.20), Inches(6.10), Inches(1.60), fc=C_BG_DEC, lc=C_DEC, lw=1)
tb(sl, Inches(0.50), Inches(1.25), Inches(5.70), Inches(0.25),
   "RNN的三大痛点", fs=12, fc=C_DEC, b=True)
mtb(sl, Inches(0.50), Inches(1.55), Inches(5.70), Inches(1.15), [
    ("① 顺序处理瓶颈：", True, C_DEC, 10),
    ("RNN必须逐词处理，GPU并行能力浪费99.9%。1000词=1000步串行。", False, C_D, 9),
    ("② 长距离遗忘：", True, C_DEC, 10),
    ("信息经数百次矩阵乘法指数级衰减，100步后可能只剩10的负30次方。", False, C_D, 9),
    ("③ 信息瓶颈：", True, C_DEC, 10),
    ("整个序列压缩成512维向量，就像用100字摘要概括500页书。", False, C_D, 9),
], ls=1.15)

# 注意力突破卡片
rrect(sl, Inches(6.60), Inches(1.20), Inches(6.40), Inches(1.60), fc=C_BG_GRN, lc=C_FFN, lw=1)
tb(sl, Inches(6.80), Inches(1.25), Inches(6.00), Inches(0.25),
   "核心突破 — 自注意力机制", fs=12, fc=C_FFN, b=True)
mtb(sl, Inches(6.80), Inches(1.55), Inches(6.00), Inches(1.15), [
    ("从\"排队发言\"到\"圆桌会议\"", True, C_FFN, 10),
    ("每个词都能与序列中所有其他词直接交流。以\"it\"为例，自注意力让\"it\"与\"cat\"建立", False, C_D, 9),
    ("最强关联(权重0.6)，模型\"知道\"it指代的是cat。", False, C_D, 9),
    ("数学本质：Q×K^T（查询×键的转置），高维空间中的余弦相似度。", False, C_D, 9),
], ls=1.15)

# 对比表
tbl = sl.shapes.add_table(5, 3, Inches(0.30), Inches(3.00), Inches(12.70), Inches(2.00))
table = tbl.table
table.columns[0].width = Inches(2.50)
table.columns[1].width = Inches(5.10)
table.columns[2].width = Inches(5.10)
headers = ["对比维度", "RNN/LSTM", "Transformer"]
data = [
    ["计算方式", "逐词串行，无法并行", "所有词同时处理，完全并行"],
    ["长距离依赖", "信息指数衰减，100词后遗忘严重", "任意词对直接交互，无距离限制"],
    ["信息压缩", "整个序列压缩成固定长度向量", "每个词保留独立表示，无信息损失"],
    ["并行能力", "串行O(T)，GPU利用率<1%", "并行O(1)，GPU充分跑满"],
]
for j, h in enumerate(headers):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
            if j == 2: p.font.color.rgb = C_FFN; p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 底部总结
rrect(sl, Inches(0.30), Inches(5.15), Inches(12.70), Inches(1.70), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.20), Inches(12.30), Inches(0.25),
   "自注意力的数学直觉", fs=12, fc=C_ACC, b=True)
mtb(sl, Inches(0.50), Inches(5.50), Inches(12.30), Inches(1.25), [
    "自注意力的核心计算：Attention(Q,K,V) = Softmax(Q×K^T / √d_k) × V",
    ("Q（查询）= \"我在找什么信息\"  ", False, C_ENC, 10),
    ("K（键）= \"我有什么特征\"  ", False, C_ATT, 10),
    ("V（值）= \"我的实际内容\"  ", False, C_FFN, 10),
    "点积Q×K^T一次性计算所有词对的关联度，形成T×T注意力矩阵，完全并行。这个操作同时解决了RNN的三个根本问题。",
], ls=1.2)
chk(sl, "P3")

# ==================== P4 架构总览 ====================
print("P4: 架构总览")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "02  架构总览", "编码器-解码器结构 — 编码器理解输入，解码器生成输出")

# 信息卡片
info_cards = [
    ("编码器=编辑部", "6层堆叠，每层:\n自注意力+前馈网络\n每个子层有残差+归一化\n约3450万参数", C_ENC),
    ("解码器=翻译部", "6层，多一个交叉注意力:\n掩码自注意力\n→交叉注意力(查资料)\n→前馈网络", C_DEC),
    ("关键数字", "512维模型\n8个注意力头\nFFN 2048维\n总参数6500万", C_ATT),
    ("训练数据", "WMT 2014英德翻译\n450万句对\n8块P100 GPU\n训练12小时", C_FFN),
]
for i, (title, body, color) in enumerate(info_cards):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.30), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(1.50), Inches(2.80), Inches(0.90), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# 数据流全景 - 用卡片+箭头
flow_items = [
    ("分词(BPE)", "文本→子词序列", C_G),
    ("嵌入矩阵", "37000×512查表", C_EMB),
    ("位置编码", "sin/cos注入", C_EMB),
    ("编码器×6", "每层交流+提炼", C_ENC),
    ("解码器×6", "每步查资料+生成", C_DEC),
    ("线性变换", "512→37000", C_ATT),
    ("Softmax", "→概率分布", C_FFN),
]
for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.30 + i * 1.82)
    rrect(sl, x, Inches(2.70), Inches(1.60), Inches(1.10), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(2.70), Inches(1.60), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.02), Inches(2.72), Inches(1.56), Inches(0.22), title, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(2.98), Inches(1.50), Inches(0.70), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
    if i < len(flow_items) - 1:
        arrow_r(sl, x+Inches(1.60), Inches(3.10), Inches(0.22), Inches(0.16), color)

# K,V连线标注
rrect(sl, Inches(7.40), Inches(3.90), Inches(2.00), Inches(0.35), fc=C_BG_EMB, lc=C_EMB, lw=1)
tb(sl, Inches(7.40), Inches(3.92), Inches(2.00), Inches(0.30),
   "K,V向量传递到解码器", fs=8, fc=C_EMB, b=True, a=PP_ALIGN.CENTER)

# 组件对照表
rect(sl, Inches(0.15), Inches(4.40), Inches(12.80), Inches(0.28), fc=C_PRI)
tb(sl, Inches(0.30), Inches(4.41), Inches(12.50), Inches(0.26), "编码器 vs 解码器对照", fs=11, fc=C_W, b=True)

tbl = sl.shapes.add_table(4, 3, Inches(0.30), Inches(4.80), Inches(12.70), Inches(1.50))
table = tbl.table
table.columns[0].width = Inches(3.00)
table.columns[1].width = Inches(4.85)
table.columns[2].width = Inches(4.85)
h_data = [["对比维度", "编码器 Encoder", "解码器 Decoder"],
          ["自注意力类型", "标准自注意力（可看所有词）", "掩码自注意力（只能看之前的词）"],
          ["独有子层", "无", "交叉注意力（Q来自解码器，K/V来自编码器）"],
          ["核心任务", "深度理解输入文本的上下文", "逐步生成输出，每步参考编码器结果"]]
for j, h in enumerate(h_data[0]):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
for i, row_data in enumerate(h_data[1:]):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 底部总结
rrect(sl, Inches(0.30), Inches(6.45), Inches(12.70), Inches(0.55), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(6.48), Inches(12.30), Inches(0.45),
   "核心比喻：编码器像编辑部（深度理解原文），解码器像翻译部（逐词翻译输出）。编码器的输出(K,V)通过交叉注意力传递给解码器，就像编辑部把理解好的资料交给翻译部参考。",
   fs=9, fc=C_D, ls=1.15)
chk(sl, "P4")

# ==================== P5 编码器详解 ====================
print("P5: 编码器详解")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "03  编码器单层详解", "每层都在\"交流\"与\"提炼\" — 多头注意力 + 前馈网络 + 残差连接 + 层归一化")

# 单层数据流 - 水平流程卡片
flow = [
    ("输入 X", "512维\n向量", C_G),
    ("多头自注意力", "8头×64维\n词间交流", C_ENC),
    ("残差+归一化", "Add&Norm\n512维", C_FFN),
    ("前馈网络", "512→2048\n→ReLU→512", C_ATT),
    ("残差+归一化", "Add&Norm\n512维", C_FFN),
    ("输出", "512维\n更丰富的表示", C_G),
]
for i, (title, desc, color) in enumerate(flow):
    x = Inches(0.30 + i * 2.15)
    rrect(sl, x, Inches(1.20), Inches(1.90), Inches(1.20), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(1.90), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.02), Inches(1.22), Inches(1.86), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(1.48), Inches(1.80), Inches(0.80), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
    if i < len(flow) - 1:
        arrow_r(sl, x+Inches(1.90), Inches(1.68), Inches(0.25), Inches(0.16), color)

# 残差旁路标注
rrect(sl, Inches(0.30), Inches(2.55), Inches(10.50), Inches(0.30), fc=C_BG_GRN, lc=C_FFN, lw=1)
tb(sl, Inches(0.50), Inches(2.57), Inches(10.10), Inches(0.25),
   "残差连接：output = input + SubLayer(input)  |  梯度中+1确保永不消失  |  重复6次，每轮让表示更丰富", fs=9, fc=C_FFN, b=True)

# 4个详解卡片
cards = [
    ("多头自注意力 = 8个讨论组", C_ENC, C_BG_ENC,
     "单头只能从一个角度分析，多头用8个独立头学习不同关系模式：",
     "头1专注语法(主谓)、头2关注语义(同义词)、头3关注位置(相邻词)、头4关注指代(代词-名词)等。",
     "每个头独立计算Q,K,V(各64维)，8个结果拼接(8×64=512)后线性变换融合。"),
    ("残差连接 = 高速公路匝道", C_FFN, C_BG_GRN,
     "output = input + SubLayer(input)，为梯度提供直接路径。",
     "∂Loss/∂x = ∂Loss/∂Output × (∂SubLayer/∂x + 1)，加号后面的\"1\"确保梯度永不消失。",
     "这是Transformer能堆叠100+层的核心原因。没有残差连接，深层网络根本无法训练。"),
    ("层归一化 = 标准化体检", C_ATT, C_BG_ATT,
     "y = γ×(x-μ)/σ+β，对每个词的512维向量独立做标准化。",
     "没有归一化时，数值范围可能从[-0.01,0.01]突变到[-100,100]，导致训练不稳定甚至无法收敛。",
     "γ和β各512维，每层1024个可学习参数。确保每层输入在标准化范围内。"),
    ("前馈网络 = 独立思考空间", C_EMB, C_BG_EMB,
     "FFN(x) = ReLU(xW₁+b₁)W₂+b₂，512→2048→ReLU→512。",
     "参数量约210万(512×2048×2)，比自注意力(约100万)大2倍，是计算量最大的组件。",
     "注意力负责词间交流，FFN负责每个词自身的深层特征提取。两个组件互补。"),
]
for i, (title, color, bgcolor, line1, line2, line3) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.30 + col * 6.50)
    y = Inches(3.00 + row * 2.00)
    rrect(sl, x, y, Inches(6.20), Inches(1.80), fc=C_W, lc=color, lw=1)
    rect(sl, x, y, Inches(6.20), Inches(0.28), fc=color)
    tb(sl, x+Inches(0.05), y+Inches(0.02), Inches(6.10), Inches(0.24), title, fs=10, fc=C_W, b=True)
    mtb(sl, x+Inches(0.10), y+Inches(0.35), Inches(6.00), Inches(1.35), [
        (line1, False, C_D, 9),
        (line2, False, C_D, 9),
        (line3, False, C_D, 9),
    ], ls=1.15)
chk(sl, "P5")

# ==================== P6 解码器详解 ====================
print("P6: 解码器详解")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "04  解码器单层详解", "比编码器多一座\"翻译桥梁\" — 掩码注意力 + 交叉注意力 + 因果掩码")

# 解码器结构流程
flow = [
    ("掩码自注意力", "只能看之前的词\n因果掩码", C_ENC),
    ("残差+归一化", "Add&Norm", C_FFN),
    ("交叉注意力", "Q来自解码器\nK/V来自编码器", C_EMB),
    ("残差+归一化", "Add&Norm", C_FFN),
    ("前馈网络", "512→2048→512", C_ATT),
    ("残差+归一化", "Add&Norm", C_FFN),
]
for i, (title, desc, color) in enumerate(flow):
    x = Inches(0.30 + i * 2.15)
    rrect(sl, x, Inches(1.20), Inches(1.90), Inches(1.10), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(1.90), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.02), Inches(1.22), Inches(1.86), Inches(0.22), title, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(1.48), Inches(1.80), Inches(0.70), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
    if i < len(flow) - 1:
        arrow_r(sl, x+Inches(1.90), Inches(1.62), Inches(0.25), Inches(0.16), color)

# K,V标注
rrect(sl, Inches(4.50), Inches(2.40), Inches(3.50), Inches(0.30), fc=C_BG_EMB, lc=C_EMB, lw=1)
tb(sl, Inches(4.50), Inches(2.42), Inches(3.50), Inches(0.25),
   "编码器输出(K,V) → 交叉注意力的K和V", fs=8, fc=C_EMB, b=True, a=PP_ALIGN.CENTER)

# 因果掩码矩阵
rrect(sl, Inches(0.30), Inches(2.85), Inches(4.00), Inches(2.50), fc=C_W, lc=C_DEC, lw=1)
rect(sl, Inches(0.30), Inches(2.85), Inches(4.00), Inches(0.28), fc=C_DEC)
tb(sl, Inches(0.35), Inches(2.87), Inches(3.90), Inches(0.24), "因果掩码矩阵 Causal Mask", fs=10, fc=C_W, b=True)
# 4x4 matrix
labels = ["", "[START]", "我", "爱", "你"]
for r_idx in range(5):
    for c_idx in range(5):
        x = Inches(0.50 + c_idx * 0.75)
        y = Inches(3.20 + r_idx * 0.40)
        if r_idx == 0:
            tb(sl, x, y, Inches(0.75), Inches(0.35), labels[c_idx], fs=8, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
        elif c_idx == 0:
            tb(sl, x, y, Inches(0.75), Inches(0.35), labels[r_idx], fs=8, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
        else:
            if c_idx <= r_idx:
                rrect(sl, x, y, Inches(0.75), Inches(0.35), fc=C_BG_ENC, lc=C_ENC, lw=0.5)
                tb(sl, x, y, Inches(0.75), Inches(0.35), "可见", fs=7, fc=C_ENC, a=PP_ALIGN.CENTER)
            else:
                rrect(sl, x, y, Inches(0.75), Inches(0.35), fc=C_BG_DEC, lc=C_DEC, lw=0.5)
                tb(sl, x, y, Inches(0.75), Inches(0.35), "-inf", fs=7, fc=C_DEC, a=PP_ALIGN.CENTER)

# 3个详解卡片
cards = [
    ("掩码自注意力 = 考试不能偷看", C_ENC, C_BG_ENC,
     "因果掩码将注意力矩阵上三角设为-∞，Softmax后变0。处理\"爱\"时只能看到[START,我,爱]，看不到\"你\"。",
     "训练和推理都必须严格执行——训练时如果偷看答案，推理时没有答案可抄就会崩溃。这保证了训练-推理行为一致，避免暴露偏差。"),
    ("交叉注意力 = 带着问题去查资料", C_EMB, C_BG_EMB,
     "Q来自解码器(\"我在翻译，需要找原文信息\")，K和V来自编码器(\"编辑对原文的理解\")。",
     "计算：Q×K^T→缩放→Softmax→×V。生成\"爱\"时，Query\"找情感动词\"与\"love\"的Key匹配度最高(0.7)。Q和K的序列长度可以不同（源语言和目标语言长度不同）。"),
    ("与编码器的核心区别", C_DEC, C_BG_DEC,
     "解码器每层3个子层（掩码自注意力→交叉注意力→前馈网络），编码器每层只有2个子层。",
     "解码器比编码器每层多一个交叉注意力子层，这是两者最核心的区别。6层解码器共18个子层+18个残差归一化操作。"),
]
for i, (title, color, bgcolor, line1, line2) in enumerate(cards):
    x = Inches(4.50)
    y = Inches(2.85 + i * 1.45)
    rrect(sl, x, y, Inches(8.50), Inches(1.30), fc=C_W, lc=color, lw=1)
    rect(sl, x, y, Inches(8.50), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.05), y+Inches(0.02), Inches(8.40), Inches(0.21), title, fs=10, fc=C_W, b=True)
    mtb(sl, x+Inches(0.10), y+Inches(0.30), Inches(8.30), Inches(0.90), [
        (line1, False, C_D, 9),
        (line2, False, C_D, 9),
    ], ls=1.15)
chk(sl, "P6")

# ==================== P7 训练阶段 ====================
print("P7: 训练阶段")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "05  训练阶段", "开卷考试 — 所有位置并行计算，GPU充分利用")

# 信息卡片
info = [
    ("训练方式", "教师强制(Teacher Forcing)\n解码器输入=真实答案\n所有位置并行计算\nGPU充分并行处理", C_ENC),
    ("嵌入矩阵", "37000×512 = 1890万参数\n占模型总参数30%\n每个词→512维向量\n编码语义信息", C_EMB),
    ("位置编码", "PE(pos,2i)=sin(pos/10000^(2i/d))\n注入位置信息\n固定不学习\n不增加可训练参数", C_ATT),
    ("训练目标", "预测下一个词\nLoss=-log(p(正确词))\n初始≈10.5(随机)\n训练后降到1.5-3.0", C_FFN),
]
for i, (title, body, color) in enumerate(info):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.50), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(1.50), Inches(2.80), Inches(1.10), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# 训练流程 - 水平卡片
flow = [
    ("输入文本", "\"I love you\"", C_G),
    ("编码器×6", "深度理解输入", C_ENC),
    ("K,V向量", "编码器输出缓存", C_EMB),
    ("解码器×6", "[START,我,爱]→预测", C_DEC),
    ("预测值", "[我,爱,你,END]", C_ATT),
    ("↔ 交叉熵", "与真实目标对比", C_DEC),
]
for i, (title, desc, color) in enumerate(flow):
    x = Inches(0.30 + i * 2.15)
    rrect(sl, x, Inches(2.90), Inches(1.90), Inches(1.00), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(2.90), Inches(1.90), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.02), Inches(2.92), Inches(1.86), Inches(0.22), title, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(3.18), Inches(1.80), Inches(0.60), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
    if i < len(flow) - 1:
        arrow_r(sl, x+Inches(1.90), Inches(3.25), Inches(0.25), Inches(0.16), color)

# 编码器/解码器分工
rrect(sl, Inches(0.30), Inches(4.05), Inches(6.10), Inches(1.20), fc=C_BG_ENC, lc=C_ENC, lw=1)
tb(sl, Inches(0.50), Inches(4.10), Inches(5.70), Inches(0.25), "编码器的职责", fs=11, fc=C_ENC, b=True)
mtb(sl, Inches(0.50), Inches(4.40), Inches(5.70), Inches(0.75), [
    "编码器只处理输入，每批训练只运行一次，输出(K,V)缓存供解码器使用。",
    "专注\"理解\"——6层堆叠，每层注意力+FFN，让每个词的表示蕴含整个句子的上下文。",
    "编码器参数约3450万，占总参数53%。计算成本只发生一次，与生成长度无关。",
], ls=1.15)

rrect(sl, Inches(6.60), Inches(4.05), Inches(6.40), Inches(1.20), fc=C_BG_DEC, lc=C_DEC, lw=1)
tb(sl, Inches(6.80), Inches(4.10), Inches(6.00), Inches(0.25), "解码器的职责", fs=11, fc=C_DEC, b=True)
mtb(sl, Inches(6.80), Inches(4.40), Inches(6.00), Inches(0.75), [
    "解码器处理目标序列，每层通过交叉注意力查阅编码器输出——\"带着问题查资料\"。",
    "专注\"生成\"——根据编码器理解和已生成内容，预测下一个词。掩码注意力防止偷看。",
    "训练时用教师强制(并行)，推理时用自回归(串行)。两种模式下掩码都必须执行。",
], ls=1.15)

# 训练配置
rrect(sl, Inches(0.30), Inches(5.40), Inches(12.70), Inches(1.50), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.45), Inches(12.30), Inches(0.25), "训练配置与核心比喻", fs=11, fc=C_ACC, b=True)
mtb(sl, Inches(0.50), Inches(5.75), Inches(12.30), Inches(1.00), [
    ("配置：", True, C_PRI, 9),
    ("Adam优化器(β₁=0.9,β₂=0.98) | 学习率峰值0.0003(前4000步增温后衰减) | 标签平滑0.1 | dropout=0.3 | 8×P100训练12小时 | BLEU 28.4", False, C_D, 9),
    ("比喻：", True, C_PRI, 9),
    ("训练就像开卷考试——每道题都能看到前面所有题的标准答案(教师强制)，所有题目同时做(并行)，效率很高但需要大量资源(算力、数据、时间)。", False, C_D, 9),
], ls=1.15)
chk(sl, "P7")

prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
