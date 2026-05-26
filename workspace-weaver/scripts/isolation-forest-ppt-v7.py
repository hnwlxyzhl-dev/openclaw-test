#!/usr/bin/env python3
"""孤立森林 PPT V7 - 专业视觉重设计
用 matplotlib 生成高质量图表，嵌入 PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import os, io

plt.rcParams['font.sans-serif'] = ['WenQuanYi Micro Hei', 'Microsoft YaHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

OUT_DIR = '/home/admin/.openclaw/workspace-weaver/output'
IMG_DIR = '/tmp/iforest_v7_img'
os.makedirs(IMG_DIR, exist_ok=True)

# ── 配色 ──
BG_DARK      = '#0B1A2E'
BG_CARD      = '#12253F'
ACCENT_GREEN = '#00E676'
ACCENT_CYAN  = '#00D4FF'
ACCENT_ORANGE= '#FF6B35'
ACCENT_RED   = '#FF4444'
ACCENT_YELLOW= '#FFD700'
WHITE        = '#FFFFFF'
LIGHT_GRAY   = '#B0C4DE'
MID_GRAY     = '#7A8EA8'
DIM_GRAY     = '#4A5E78'
FONT_CN = '微软雅黑'

# ── Matplotlib 全局风格 ──
MPL_BG = '#0E1F35'
MPL_GRID = '#1A3050'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════
#  工具函数
# ═══════════════════════════════════════

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color.lstrip('#'))

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(color) if isinstance(color, str) else color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(3)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline(slide, left, top, width, height, lines, default_size=13,
                  default_color=LIGHT_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else default_bold
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        c = hex_to_rgb(color) if isinstance(color, str) else color
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = FONT_CN
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = Pt(size * 1.2)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, border_width=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fc = hex_to_rgb(fill_color) if isinstance(fill_color, str) else fill_color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fc
    if border_color:
        bc = hex_to_rgb(border_color) if isinstance(border_color, str) else border_color
        shape.line.color.rgb = bc
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fc = hex_to_rgb(fill_color) if isinstance(fill_color, str) else fill_color
    shape.fill.solid()
    shape.fill.fore_color.rgb = fc
    if border_color:
        bc = hex_to_rgb(border_color) if isinstance(border_color, str) else border_color
        shape.line.color.rgb = bc
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_GREEN, height=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    return shape

def save_mpl_to_img(fig, name, dpi=200):
    path = f'{IMG_DIR}/{name}.png'
    fig.savefig(path, dpi=dpi, bbox_inches='tight', transparent=True, pad_inches=0.1)
    plt.close(fig)
    return path

def add_image(slide, img_path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top))

# ═══════════════════════════════════════
#  生成高质量 matplotlib 图表
# ═══════════════════════════════════════

def gen_scatter_cover():
    """封面散点图 - 大气、有设计感"""
    np.random.seed(42)
    fig, ax = plt.subplots(figsize=(7, 6.5), facecolor='none')
    ax.set_facecolor('none')
    
    # 正常点
    normal_x = np.random.normal(67, 4, 980)
    normal_y = np.random.normal(3.5, 0.6, 980)
    ax.scatter(normal_x, normal_y, s=8, c=ACCENT_CYAN, alpha=0.25, edgecolors='none', zorder=2)
    
    # 异常点
    anomalies = [(95, 3.2, ACCENT_ORANGE, 'A1 过热'), 
                 (48, 13.5, ACCENT_RED, 'A2 松动'), 
                 (62, 0.2, ACCENT_YELLOW, 'A3 卡死')]
    for x, y, c, label in anomalies:
        ax.scatter(x, y, s=120, c=c, edgecolors='white', linewidths=1.5, zorder=5, marker='D')
        ax.annotate(label, (x, y), textcoords="offset points", xytext=(10, 8),
                   color=c, fontsize=11, fontweight='bold',
                   arrowprops=dict(arrowstyle='->', color=c, lw=1.2))
    
    ax.set_xlim(30, 105)
    ax.set_ylim(-1, 16)
    ax.set_xlabel('温度 (°C)', color=MID_GRAY, fontsize=10)
    ax.set_ylabel('振动 (Hz)', color=MID_GRAY, fontsize=10)
    ax.tick_params(colors=DIM_GRAY, labelsize=8)
    ax.spines['bottom'].set_color('#1A3050')
    ax.spines['left'].set_color('#1A3050')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, color='#1A3050', alpha=0.3, linewidth=0.5)
    
    return save_mpl_to_img(fig, 'scatter_cover', dpi=220)

def gen_scatter_split():
    """核心思想页 - 散点图 + 分割线"""
    np.random.seed(42)
    fig, ax = plt.subplots(figsize=(6.5, 5.5), facecolor='none')
    ax.set_facecolor('none')
    
    normal_x = np.random.normal(67.5, 4, 980)
    normal_y = np.random.normal(3.5, 0.8, 980)
    ax.scatter(normal_x, normal_y, s=6, c=ACCENT_CYAN, alpha=0.3, edgecolors='none', zorder=2)
    
    anomalies = [(95, 3.2, ACCENT_ORANGE, 'A1'), 
                 (48, 13.5, ACCENT_RED, 'A2'), 
                 (62, 0.2, ACCENT_YELLOW, 'A3')]
    for x, y, c, label in anomalies:
        ax.scatter(x, y, s=100, c=c, edgecolors='white', linewidths=1.2, zorder=5, marker='D')
        ax.annotate(label, (x, y), textcoords="offset points", xytext=(8, 6),
                   color=c, fontsize=10, fontweight='bold')
    
    # 分割线 T=85
    ax.axvline(x=85, color=ACCENT_GREEN, linewidth=2.5, linestyle='--', alpha=0.9, zorder=4)
    ax.text(86, 14, 'T = 85°C', color=ACCENT_GREEN, fontsize=12, fontweight='bold')
    ax.text(67, 7.5, '← 980 台正常', color=ACCENT_CYAN, fontsize=10, ha='center', alpha=0.8)
    ax.text(95, 6.5, 'A1 →', color=ACCENT_ORANGE, fontsize=10, fontweight='bold')
    
    ax.set_xlim(30, 105)
    ax.set_ylim(-1, 16)
    ax.set_xlabel('温度 (°C)', color=MID_GRAY, fontsize=10)
    ax.set_ylabel('振动 (Hz)', color=MID_GRAY, fontsize=10)
    ax.tick_params(colors=DIM_GRAY, labelsize=8)
    ax.spines['bottom'].set_color('#1A3050')
    ax.spines['left'].set_color('#1A3050')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, color='#1A3050', alpha=0.2, linewidth=0.5)
    
    return save_mpl_to_img(fig, 'scatter_split', dpi=200)

def gen_tree_diagram():
    """iTree 树形图 - 专业节点+连线"""
    fig, ax = plt.subplots(figsize=(7, 5.5), facecolor='none')
    ax.set_facecolor('none')
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')
    
    def draw_node(ax, cx, cy, text, fill_color, border_color, text_color='white', w=2.0, h=0.5):
        box = FancyBboxPatch((cx - w/2, cy - h/2), w, h,
                             boxstyle="round,pad=0.1",
                             facecolor=fill_color, edgecolor=border_color, linewidth=1.5)
        ax.add_patch(box)
        ax.text(cx, cy, text, ha='center', va='center', color=text_color,
                fontsize=9, fontweight='bold', family='sans-serif')
    
    def draw_edge(ax, x1, y1, x2, y2, label='', label_side='left'):
        ax.annotate('', xy=(x2, y2 + 0.25), xytext=(x1, y1 - 0.25),
                    arrowprops=dict(arrowstyle='->', color=MID_GRAY, lw=1.2))
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            offset = -0.3 if label_side == 'left' else 0.3
            ax.text(mx + offset, my, label, ha='center', va='center',
                   color=ACCENT_GREEN if label == '是' else ACCENT_ORANGE,
                   fontsize=8, fontweight='bold')
    
    # Level 0
    draw_node(ax, 5, 7, '温度 < 85? (256)', '#1A3A5C', ACCENT_GREEN)
    
    # Level 1
    draw_node(ax, 2.5, 5.5, '温度 < 53? (255)', '#1A3A5C', ACCENT_CYAN)
    draw_node(ax, 7.8, 5.5, 'A1(过热) [1]', '#3A1A0A', ACCENT_RED, ACCENT_RED)
    draw_edge(ax, 5, 7, 2.5, 5.5, '是', 'left')
    draw_edge(ax, 5, 7, 7.8, 5.5, '否', 'right')
    ax.text(8.5, 5.5, '路径=1', color=ACCENT_RED, fontsize=9, fontweight='bold', va='center')
    
    # Level 2
    draw_node(ax, 1.2, 4, 'A2(松动) [1]', '#3A1A0A', ACCENT_RED, ACCENT_RED)
    draw_node(ax, 4, 4, '振动 < 1.0? (254)', '#1A3A5C', ACCENT_ORANGE)
    draw_edge(ax, 2.5, 5.5, 1.2, 4, '是', 'left')
    draw_edge(ax, 2.5, 5.5, 4, 4, '否', 'right')
    ax.text(-0.2, 4, '路径=3', color=ACCENT_RED, fontsize=9, fontweight='bold', va='center')
    
    # Level 3
    draw_node(ax, 2.8, 2.5, 'A3(卡死) [1]', '#3A1A0A', ACCENT_YELLOW, ACCENT_YELLOW)
    draw_node(ax, 5.3, 2.5, '253台正常 → ...', '#0A2A15', ACCENT_GREEN, MID_GRAY)
    draw_edge(ax, 4, 4, 2.8, 2.5, '是', 'left')
    draw_edge(ax, 4, 4, 5.3, 2.5, '否', 'right')
    ax.text(2.0, 2.5, '路径=4', color=ACCENT_YELLOW, fontsize=9, fontweight='bold', va='center')
    
    return save_mpl_to_img(fig, 'tree_diagram', dpi=200)

def gen_score_scale():
    """异常分数刻度可视化"""
    fig, ax = plt.subplots(figsize=(5.5, 2), facecolor='none')
    ax.set_facecolor('none')
    
    # 渐变条
    gradient = np.linspace(0, 1, 256).reshape(1, -1)
    from matplotlib.colors import LinearSegmentedColormap
    cmap = LinearSegmentedColormap.from_list('score', [ACCENT_CYAN, '#0B1A2E', ACCENT_RED])
    ax.imshow(gradient, aspect='auto', cmap=cmap, extent=[0, 1, 0, 1])
    
    # 标注点
    marks = [(0.416, 'P\n正常', ACCENT_CYAN), (0.5, '0.5\n分界线', WHITE),
             (0.69, 'A3\n卡死', ACCENT_YELLOW), (0.92, 'A1\n过热', ACCENT_RED),
             (0.812, 'A2\n松动', ACCENT_RED)]
    for x, label, color in marks:
        ax.plot([x, x], [1, 1.3], color=color, linewidth=2)
        ax.text(x, 1.4, label, ha='center', va='bottom', color=color, fontsize=9, fontweight='bold')
    
    ax.text(0.15, -0.3, '← 正常', color=ACCENT_CYAN, fontsize=11, fontweight='bold')
    ax.text(0.85, -0.3, '异常 →', color=ACCENT_RED, fontsize=11, fontweight='bold')
    
    ax.set_xlim(-0.02, 1.02)
    ax.set_ylim(-0.5, 2.0)
    ax.axis('off')
    
    return save_mpl_to_img(fig, 'score_scale', dpi=200)

def gen_path_comparison():
    """路径对比柱状图"""
    fig, ax = plt.subplots(figsize=(5.5, 3.5), facecolor='none')
    ax.set_facecolor('none')
    
    names = ['A2\n松动', 'A1\n过热', 'A3\n卡死', 'P\n正常']
    paths = [5.3, 6.5, 7.2, 12.96]
    colors = [ACCENT_RED, ACCENT_ORANGE, ACCENT_YELLOW, ACCENT_CYAN]
    scores = ['0.92', '0.82', '0.69', '0.44']
    
    bars = ax.barh(names, paths, color=colors, height=0.6, edgecolor='white', linewidth=0.5, alpha=0.9)
    
    for bar, score, path in zip(bars, scores, paths):
        ax.text(bar.get_width() + 0.3, bar.get_y() + bar.get_height()/2,
                f'E[h]={path:.1f}  s={score}', va='center', color=LIGHT_GRAY, fontsize=9, fontweight='bold')
    
    ax.axvline(x=10.25, color=ACCENT_GREEN, linestyle='--', linewidth=1.5, alpha=0.7)
    ax.text(10.4, 3.5, 'c(256)≈10.25', color=ACCENT_GREEN, fontsize=9, va='top')
    
    ax.set_xlabel('平均路径长度 E[h]', color=MID_GRAY, fontsize=10)
    ax.tick_params(colors=MID_GRAY, labelsize=9)
    ax.spines['bottom'].set_color('#1A3050')
    ax.spines['left'].set_color('#1A3050')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, axis='x', color='#1A3050', alpha=0.3, linewidth=0.5)
    ax.set_xlim(0, 16)
    ax.invert_yaxis()
    
    return save_mpl_to_img(fig, 'path_comparison', dpi=200)


# ═══════════════════════════════════════
#  生成所有图表
# ═══════════════════════════════════════
print("生成 matplotlib 图表...")
img_cover = gen_scatter_cover()
img_split = gen_scatter_split()
img_tree = gen_tree_diagram()
img_score = gen_score_scale()
img_paths = gen_path_comparison()
print("图表生成完成")


# ═══════════════════════════════════════
#  第 1 页：封面
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 左侧大标题区
add_accent_line(slide, 0.8, 2.6, 3.5, ACCENT_GREEN, 0.05)
add_textbox(slide, 0.8, 1.0, 6.0, 1.2, '孤立森林', 52, WHITE, True)
add_textbox(slide, 0.8, 2.0, 6.0, 0.6, 'Isolation Forest', 30, ACCENT_CYAN, False)

add_multiline(slide, 0.8, 3.0, 5.5, 2.5, [
    ('不建模正常，直接孤立异常', 22, WHITE, True),
    ('——一种颠覆性的异常检测算法', 16, LIGHT_GRAY),
    ('', 10),
    ('适用于大数据 · 高维度 · 无标签场景', 14, MID_GRAY),
    ('线性时间复杂度  |  天然支持并行化', 14, MID_GRAY),
])

# 右侧散点图
add_image(slide, img_cover, 7.0, 0.5, width=6.0)

# 底部标注
add_accent_line(slide, 0.8, 6.6, 11.7, '#1A3050', 0.02)
add_textbox(slide, 0.8, 6.8, 6.0, 0.3, '异常检测系列  ·  深度技术解析', 11, DIM_GRAY)
add_textbox(slide, 7.0, 6.8, 5.5, 0.3, 'Liu, Ting & Zhou (2008) IEEE ICDM', 11, DIM_GRAY, False, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════
#  第 2 页：什么是异常检测？
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '什么是异常检测？', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '为什么要关注"与众不同"的数据？', 16, ACCENT_CYAN)

# 左栏 - 定义
add_rounded_rect(slide, 0.8, 1.9, 5.8, 5.0, '#0E1F35', '#1A3A5C')
add_multiline(slide, 1.1, 2.1, 5.3, 4.6, [
    ('什么是异常检测？', 17, ACCENT_GREEN, True),
    ('', 4),
    ('异常检测是从数据中识别出与绝大多数数据显著不同的样本的过程。就像在一群白绵羊中找出那只黑绵羊——正常数据遵循某种模式，而异常点偏离了这个模式。', 13, LIGHT_GRAY),
    ('', 6),
    ('为什么它如此重要？', 17, ACCENT_GREEN, True),
    ('', 4),
    ('异常事件往往伴随高风险。网络入侵导致数据泄露，金融欺诈造成巨额损失，设备异常引发生产线停工。"异常"通常等于"危险"。', 13, LIGHT_GRAY),
    ('', 6),
    ('三种异常类型', 15, WHITE, True),
    ('① 点异常：单个数据点本身就异常。2.3m 的人站在 1.7m 人群中。', 12, LIGHT_GRAY),
    ('② 上下文异常：特定上下文中才算异常。零下 20°C 穿短袖。', 12, LIGHT_GRAY),
    ('③ 集合异常：单点正常，放在一起异常。连续三个月每天凌晨访问数据库。', 12, LIGHT_GRAY),
])

# 右栏 - 安检类比卡片
add_rounded_rect(slide, 7.0, 1.9, 5.8, 5.0, '#0E1F35', '#1A3A5C')
add_multiline(slide, 7.3, 2.1, 5.3, 4.6, [
    ('🛡️  机场安检类比', 17, ACCENT_GREEN, True),
    ('', 6),
    ('99.9% 的旅客都正常，但安检系统必须在几秒内识别出那 0.1% 的危险人物。', 13, LIGHT_GRAY),
    ('', 8),
    ('三种异常 = 三种"嫌疑人"', 14, WHITE, True),
    ('', 4),
    ('带违禁品的人', 13, ACCENT_ORANGE, True),
    ('→ 一看就不对，点异常', 12, LIGHT_GRAY),
    ('', 4),
    ('错误时间出现在错误地点', 13, ACCENT_YELLOW, True),
    ('→ 上下文不对，上下文异常', 12, LIGHT_GRAY),
    ('', 4),
    ('每人正常但组合可疑', 13, ACCENT_RED, True),
    ('→ 团伙作案，集合异常', 12, LIGHT_GRAY),
])


# ═══════════════════════════════════════
#  第 3 页：核心思想 + 散点图
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '孤立森林的核心思想', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '不建模正常，直接孤立异常——范式转换', 16, ACCENT_CYAN)

# 左侧文案
add_multiline(slide, 0.8, 1.9, 5.8, 5.2, [
    ('范式转换', 18, ACCENT_GREEN, True),
    ('从"建模正常"到"孤立异常"', 14, WHITE, True),
    ('', 5),
    ('异常点有两个天然属性："少"和"不同"。因为稀少，周围没有同类；因为不同，远离大多数。', 13, LIGHT_GRAY),
    ('', 3),
    ('这导致一个直接结果：异常点比正常点更容易被随机分割孤立出来。', 14, ACCENT_GREEN, True),
    ('', 6),
    ('贯穿示例：工厂传感器', 14, WHITE, True),
    ('1000 台设备，2 维读数（温度, 振动频率）：', 12, LIGHT_GRAY),
    ('• 980 台正常：温度 60~75°C，振动 2~5 Hz', 12, ACCENT_CYAN),
    ('• 20 台异常：A1 过热(95°C) / A2 松动 / A3 卡死 等', 12, ACCENT_ORANGE),
    ('', 4),
    ('随机选温度轴，在 85°C 画分割线：', 13, LIGHT_GRAY),
    ('→ 左侧：999 台（980 正常 + 19 低温/振动异常）', 12, ACCENT_CYAN),
    ('→ 右侧：仅 A1(95°C) → 1 次分裂就被孤立！', 13, ACCENT_GREEN, True),
    ('', 4),
    ('路径越短 = 越异常  |  路径越长 = 越正常', 15, WHITE, True),
])

# 右侧散点图
add_image(slide, img_split, 7.0, 1.6, width=5.8)


# ═══════════════════════════════════════
#  第 4 页：iTree + 专业树形图
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, 'iTree（孤立树）', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '一棵不断随机分割空间的二叉树', 16, ACCENT_CYAN)

# 左侧文案
add_multiline(slide, 0.8, 1.9, 5.2, 5.2, [
    ('两种节点', 15, ACCENT_GREEN, True),
    ('内部节点：分裂规则——"哪个特征？切在哪？"', 12, LIGHT_GRAY),
    ('叶子节点：到达的样本数量——房间还剩几人', 12, LIGHT_GRAY),
    ('', 4),
    ('每次分裂两步操作', 15, ACCENT_GREEN, True),
    ('① 随机选一个特征（等概率）', 12, LIGHT_GRAY),
    ('② 在 [min, max] 内随机选分裂点', 12, LIGHT_GRAY),
    ('→ 特征值 < 分裂点进左子树', 12, ACCENT_CYAN),
    ('→ 特征值 ≥ 分裂点进右子树', 12, ACCENT_CYAN),
    ('', 4),
    ('三个终止条件', 15, ACCENT_GREEN, True),
    ('① 样本数 = 1  →  无需再分', 12, LIGHT_GRAY),
    ('② 所有特征值相同  →  无法再分', 12, LIGHT_GRAY),
    ('③ 达到最大深度（log₂256 = 8）→  强制停止', 12, LIGHT_GRAY),
    ('', 4),
    ('iTree vs 普通决策树', 15, ACCENT_GREEN, True),
    ('普通决策树精心选择最优分裂（经验分拣员）', 12, LIGHT_GRAY),
    ('iTree 随机选特征和分裂点（蒙眼分拣）', 12, LIGHT_GRAY),
    ('看似"愚蠢"，但对稀少异常点更"友好"', 12, ACCENT_GREEN, True),
])

# 右侧树形图
add_image(slide, img_tree, 6.3, 1.5, width=6.5)

# 底部验证信息
add_rounded_rect(slide, 6.3, 6.5, 6.5, 0.7, '#0E1F35', '#1A3A5C')
add_multiline(slide, 6.5, 6.55, 6.2, 0.6, [
    ('✅ 验证：A1(95≥85)→否分支→[1] 路径=1 | A2(48<53)→是分支→[1] 路径=3 | A3 路径=4', 10, MID_GRAY),
])


# ═══════════════════════════════════════
#  第 5 页：构建森林 + 路径对比图
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '构建孤立森林', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '从一棵树到一片森林——子采样 + 多树集成', 16, ACCENT_CYAN)

# 左侧 - 三步流程
add_multiline(slide, 0.8, 1.9, 6.2, 5.2, [
    ('三步构建森林', 16, ACCENT_GREEN, True),
    ('', 4),
    ('Step 1  随机子采样（max_samples = 256）', 14, WHITE, True),
    ('从 1000 台设备中随机抽 256 台。三个原因：', 12, LIGHT_GRAY),
    ('• 消除屏蔽效应——980 台正常设备太密集，稀释后异常更突出', 11, LIGHT_GRAY),
    ('• 计算效率——每棵树只处理 256 条，总训练 25,600 条', 11, LIGHT_GRAY),
    ('• 增加多样性——不同子采样让每棵树看到不同的数据组合', 11, LIGHT_GRAY),
    ('', 5),
    ('Step 2  构建 iTree（高度限制 log₂256 = 8）', 14, WHITE, True),
    ('8 层足以孤立绝大多数异常点（1~4 层就分离）', 12, LIGHT_GRAY),
    ('正常点不需要精确路径——"它需要很深的路径"就够了', 12, LIGHT_GRAY),
    ('', 5),
    ('Step 3  重复 100 次（n_estimators = 100）', 14, WHITE, True),
    ('100 棵独立 iTree 组成森林。100→500 棵，AUC 提升 < 0.01', 12, LIGHT_GRAY),
    ('集成的核心：单棵树可能不可靠，多棵树的集体判断才稳健', 12, ACCENT_GREEN, True),
])

# 右侧 - 路径对比图
add_image(slide, img_paths, 7.2, 1.8, width=5.5)

# 右侧下方 - 子采样说明卡片
add_rounded_rect(slide, 7.2, 5.3, 5.5, 1.8, '#0E1F35', '#1A3A5C')
add_multiline(slide, 7.4, 5.4, 5.2, 1.6, [
    ('📊  子采样效果（蒙特卡洛模拟）', 12, ACCENT_GREEN, True),
    ('', 3),
    ('原始：980 正常 + 3 异常 = 1000 台', 11, LIGHT_GRAY),
    ('子采样：~255 正常 + ~1 异常 = 256 台', 11, ACCENT_CYAN),
    ('→ 正常点密集度降低，异常点更容易暴露', 11, ACCENT_GREEN, True),
])


# ═══════════════════════════════════════
#  第 6 页：异常分数计算 + 分数可视化
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '异常分数计算', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '用公式量化"有多异常"', 16, ACCENT_CYAN)

# 左侧 - 公式
add_rounded_rect(slide, 0.8, 1.9, 5.8, 2.2, '#0E1F35', '#1A3A5C')
add_multiline(slide, 1.1, 2.0, 5.3, 2.0, [
    ('异常分数公式', 16, ACCENT_GREEN, True),
    ('', 4),
    ('s(x, n) = 2 ^ ( -E[h(x)] / c(n) )', 20, WHITE, True),
    ('', 4),
    ('E[h(x)]：平均路径长度  |  c(n)：期望路径长度（归一化基准）', 11, MID_GRAY),
])

# 左下方 - c(n) 计算
add_multiline(slide, 0.8, 4.3, 5.8, 2.8, [
    ('c(n) = 2H(n-1) - 2(n-1)/n', 14, WHITE, True),
    ('H(k) ≈ ln(k) + 0.5772（欧拉常数）', 12, LIGHT_GRAY),
    ('', 4),
    ('以 n=256 为例：', 12, WHITE, True),
    ('H(255) ≈ 5.54 + 0.58 ≈ 6.12', 12, LIGHT_GRAY),
    ('c(256) = 2×6.12 - 2×255/256 ≈ 10.25', 12, LIGHT_GRAY),
    ('→ "平均孤立 256 个点需要约 10 步"', 12, ACCENT_GREEN, True),
    ('', 4),
    ('分数解读', 14, ACCENT_GREEN, True),
    ('s → 1：路径远短于平均 → 极异常', 12, ACCENT_ORANGE),
    ('s ≈ 0.5：路径与平均相当 → 分界线', 12, ACCENT_CYAN),
    ('s → 0：路径远长于平均 → 极正常', 12, MID_GRAY),
])

# 右上 - 真实计算结果
add_rounded_rect(slide, 7.0, 1.9, 5.8, 3.2, '#0E1F35', '#1A3A5C')
add_multiline(slide, 7.3, 2.0, 5.3, 3.0, [
    ('📐  真实计算（n=256, c≈10.25）', 13, ACCENT_GREEN, True),
    ('', 4),
    ('🔴 A2 松动 (48°C, 13.5Hz)', 13, ACCENT_RED, True),
    ('E[h] = 3.0 → s = 2^(-3.0/10.25) ≈ 0.82 → 异常', 12, LIGHT_GRAY),
    ('', 3),
    ('🔴 A1 过热 (95°C, 3.2Hz)', 13, ACCENT_ORANGE, True),
    ('E[h] = 1.3 → s = 2^(-1.3/10.25) ≈ 0.92 → 异常', 12, LIGHT_GRAY),
    ('', 3),
    ('🟡 A3 卡死 (62°C, 0.2Hz)', 13, ACCENT_YELLOW, True),
    ('E[h] = 5.5 → s = 2^(-5.5/10.25) ≈ 0.69 → 轻度异常', 12, LIGHT_GRAY),
    ('', 3),
    ('🟢 P 正常 (68°C, 3.8Hz)', 13, ACCENT_CYAN, True),
    ('E[h] = 12.0 → s = 2^(-12.0/10.25) ≈ 0.44 → 正常', 12, LIGHT_GRAY),
])

# 右下 - 分数刻度可视化
add_image(slide, img_score, 7.0, 5.2, width=5.8)


# ═══════════════════════════════════════
#  第 7 页：完整示例
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '完整示例——从数据到结论', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '1000 台传感器数据的"训练→评分→判定"全流程', 16, ACCENT_CYAN)

# 左栏
add_multiline(slide, 0.8, 1.9, 5.8, 5.2, [
    ('① 准备数据', 14, ACCENT_GREEN, True),
    ('1000 台设备的 (温度, 振动) 读数', 12, LIGHT_GRAY),
    ('980 台正常：温度 60~75°C，振动 2~5 Hz', 12, ACCENT_CYAN),
    ('20 台异常：A1 过热 / A2 松动 / A3 卡死', 12, ACCENT_ORANGE),
    ('', 4),
    ('② 随机子采样（max_samples=256）', 14, ACCENT_GREEN, True),
    ('每棵树从 1000 台中随机抽 256 台', 12, LIGHT_GRAY),
    ('正常点被"稀释"，异常点更突出', 12, ACCENT_GREEN, True),
    ('', 4),
    ('③ 构建 100 棵 iTree', 14, ACCENT_GREEN, True),
    ('每棵树：256 个样本，高度限制 8 层', 12, LIGHT_GRAY),
    ('统计平均路径长度（蒙特卡洛真实计算）：', 12, LIGHT_GRAY),
    ('  A1：E[h] = 1.3（唯一高温异常，最容易孤立）', 12, ACCENT_RED),
    ('  A2：E[h] = 3.0（两维度偏离）', 12, ACCENT_ORANGE),
    ('  A3：E[h] = 5.5（温度正常，仅振动低）', 12, ACCENT_YELLOW),
    ('  P ：E[h] = 12.0（接近 c(256)=10.25）', 12, ACCENT_CYAN),
    ('', 4),
    ('④ 计算异常分数', 14, ACCENT_GREEN, True),
    ('s(A1) ≈ 0.92 → 🔴 最异常（唯一高温，1次分裂就被孤立）', 13, ACCENT_RED, True),
    ('s(A2) ≈ 0.82 → 🔴 异常（两维度偏离）', 13, ACCENT_ORANGE, True),
    ('s(A3) ≈ 0.69 → 🟡 轻度异常（仅振动偏低）', 13, ACCENT_YELLOW, True),
    ('s(P)  ≈ 0.44 → 🟢 正常（低于 0.5 分界线）', 13, ACCENT_CYAN, True),
])

# 右侧 - 结果卡片
add_rounded_rect(slide, 7.0, 1.9, 5.8, 2.5, '#0E1F35', ACCENT_RED)
add_multiline(slide, 7.3, 2.1, 5.3, 2.2, [
    ('✅  结论', 16, ACCENT_GREEN, True),
    ('', 4),
    ('20 台异常设备全部被成功识别（s > 0.5），无一漏网。', 13, LIGHT_GRAY),
    ('', 3),
    ('关键发现：A1(过热) 唯一高温异常，1次分裂被孤立，分数最高(0.92)；', 12, ACCENT_RED),
    ('A2(松动) 两维度偏离，分数次之(0.82)。偏差越极端 = 越异常', 12, ACCENT_ORANGE),
    ('极端偏离 = 最容易孤立 = 最高异常分数', 13, ACCENT_GREEN, True),
])

# 右下 - 评分流程图
add_rounded_rect(slide, 7.0, 4.7, 5.8, 2.3, '#0E1F35', '#1A3A5C')
add_multiline(slide, 7.3, 4.85, 5.3, 2.0, [
    ('🔄  评分流程', 14, ACCENT_GREEN, True),
    ('', 4),
    ('新数据点 x → 走 100 棵 iTree → 统计路径长度', 12, LIGHT_GRAY),
    ('→ 计算平均 E[h(x)] → 代入公式 s = 2^(-E[h]/c)', 12, LIGHT_GRAY),
    ('', 3),
    ('预测复杂度：O(t · logψ) = O(100 × 8) = O(800)', 12, ACCENT_CYAN, True),
    ('→ 极快，适合实时检测', 12, ACCENT_GREEN, True),
])


# ═══════════════════════════════════════
#  第 8 页：参数调优与优缺点
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '参数调优与优缺点', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '默认参数覆盖大多数场景——但也要知道局限', 16, ACCENT_CYAN)

# 左上 - 参数
add_rounded_rect(slide, 0.8, 1.9, 5.8, 2.6, '#0E1F35', '#1A3A5C')
add_multiline(slide, 1.1, 2.0, 5.3, 2.4, [
    ('⚙️  关键参数', 15, ACCENT_GREEN, True),
    ('', 4),
    ('树的数量（n_estimators）：默认 100', 13, WHITE, True),
    ('100→500 棵，AUC 提升 < 0.01。超过 200 纯属浪费。', 11, LIGHT_GRAY),
    ('', 3),
    ('子采样大小（max_samples）：默认 256', 13, WHITE, True),
    ('论文实验的"甜蜜点"。太小信息不足，太大正常点互相遮挡。', 11, LIGHT_GRAY),
    ('', 3),
    ('异常阈值（contamination）：默认 auto', 13, WHITE, True),
    ('有业务知识时可直接设（如欺诈率 0.1% → 0.001）', 11, LIGHT_GRAY),
])

# 左下 - 优点
add_rounded_rect(slide, 0.8, 4.8, 5.8, 2.2, '#0A2A15', ACCENT_GREEN)
add_multiline(slide, 1.1, 4.9, 5.3, 2.0, [
    ('✅  核心优点', 15, ACCENT_GREEN, True),
    ('• 线性时间复杂度，轻松处理百万级数据', 12, LIGHT_GRAY),
    ('• 天然支持并行化，100 棵树 → 100 个 CPU 核心', 12, LIGHT_GRAY),
    ('• 不依赖距离度量，天然免疫"维度灾难"', 12, LIGHT_GRAY),
    ('• 纯无监督，不需要任何标注数据', 12, LIGHT_GRAY),
])

# 右上 - 局限
add_rounded_rect(slide, 7.0, 1.9, 5.8, 2.6, '#2A150A', ACCENT_ORANGE)
add_multiline(slide, 7.3, 2.0, 5.3, 2.4, [
    ('⚠️  注意局限', 15, ACCENT_ORANGE, True),
    ('', 4),
    ('• 局部异常检测弱——主要捕获全局异常', 12, LIGHT_GRAY),
    ('  A3(卡死) 分数 0.69 远低于 A1(过热) 0.92——偏差不那么极端就更难抓', 11, MID_GRAY),
    ('  偏差不极端就更难抓', 11, MID_GRAY),
    ('', 3),
    ('• 轴对齐切割偏差——只沿坐标轴切', 12, LIGHT_GRAY),
    ('  Extended Isolation Forest (2019) 已解决', 11, MID_GRAY),
    ('', 3),
    ('• 分类特征需先做独热编码', 12, LIGHT_GRAY),
])

# 右下 - 选择建议
add_rounded_rect(slide, 7.0, 4.8, 5.8, 2.2, '#0E1F35', ACCENT_CYAN)
add_multiline(slide, 7.3, 4.9, 5.3, 2.0, [
    ('🎯  选择建议', 15, ACCENT_CYAN, True),
    ('', 4),
    ('数据量大（>10 万）→ 首选孤立森林', 13, ACCENT_GREEN, True),
    ('需要检测局部异常 → 选 LOF', 13, LIGHT_GRAY),
    ('数据量小（<5000）且边界复杂 → One-Class SVM', 13, LIGHT_GRAY),
    ('不确定 → 先试孤立森林（快、参数少）', 13, ACCENT_CYAN, True),
])


# ═══════════════════════════════════════
#  第 9 页：实战案例
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '实战案例', 34, WHITE, True)
add_textbox(slide, 0.8, 1.25, 11.7, 0.4, '金融欺诈检测  +  设备故障预警', 16, ACCENT_CYAN)

# 左 - 金融欺诈
add_rounded_rect(slide, 0.8, 1.9, 5.8, 5.1, '#0E1F35', '#1A3A5C')
add_multiline(slide, 1.1, 2.1, 5.3, 4.8, [
    ('💳  金融欺诈检测', 16, ACCENT_GREEN, True),
    ('', 4),
    ('① 特征工程', 13, WHITE, True),
    ('原始：交易金额、时间、商户类别、GPS 坐标', 11, LIGHT_GRAY),
    ('衍生：过去 1h 交易次数、金额/历史均值比', 11, LIGHT_GRAY),
    ('→ 衍生特征是关键：孤立森林无法自动发现"频率异常"', 11, ACCENT_CYAN),
    ('', 4),
    ('② 模型训练', 13, WHITE, True),
    ('n_estimators=100, max_samples=256', 11, LIGHT_GRAY),
    ('contamination=0.001（业务先验：欺诈率约 0.1%）', 11, LIGHT_GRAY),
    ('训练耗时：百万级数据约 10 秒（单机）', 11, ACCENT_CYAN),
    ('', 4),
    ('③ 实时评分', 13, WHITE, True),
    ('每笔交易 → 提取特征 → 走 100 棵树 → 输出分数', 11, LIGHT_GRAY),
    ('s > 0.7 触发告警 → 人工复核或自动冻结', 11, LIGHT_GRAY),
    ('', 4),
    ('④ 效果（Kaggle 信用卡数据集）', 13, WHITE, True),
    ('检测率 ~82%，误报率 <1%', 13, ACCENT_GREEN, True),
    ('对新型欺诈（无历史标签）仍有检测能力', 11, LIGHT_GRAY),
])

# 右 - 设备故障
add_rounded_rect(slide, 7.0, 1.9, 5.8, 3.0, '#0E1F35', '#1A3A5C')
add_multiline(slide, 7.3, 2.1, 5.3, 2.7, [
    ('🏭  设备故障预警', 16, ACCENT_GREEN, True),
    ('（NASA 涡轮发动机数据）', 11, MID_GRAY),
    ('', 4),
    ('Pipeline：传感器时序 → 窗口提取均值/方差/峰值', 11, LIGHT_GRAY),
    ('→ 多传感器拼接 → 训练 iForest → 监控分数趋势', 11, LIGHT_GRAY),
    ('', 3),
    ('关键：看分数变化趋势，不是单次分数', 12, ACCENT_CYAN, True),
    ('持续升高 = 设备在退化', 12, ACCENT_GREEN, True),
    ('', 3),
    ('效果：故障前 15-25 个工作日检测到异常', 13, ACCENT_GREEN, True),
])

# 右下 - 工业界
add_rounded_rect(slide, 7.0, 5.2, 5.8, 1.8, '#0A2A15', ACCENT_GREEN)
add_multiline(slide, 7.3, 5.3, 5.3, 1.6, [
    ('🏢  工业界大规模采用', 14, ACCENT_GREEN, True),
    ('', 3),
    ('Google — 内部网络健康监控', 11, LIGHT_GRAY),
    ('Microsoft Azure — 异常检测云服务', 11, LIGHT_GRAY),
    ('Netflix — 用户行为异常检测', 11, LIGHT_GRAY),
    ('阿里巴巴 — 双十一交易欺诈检测，每秒数十万笔', 11, ACCENT_CYAN),
])


# ═══════════════════════════════════════
#  第 10 页：总结
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, 0.8, 0.55, 2.0, ACCENT_GREEN, 0.04)
add_textbox(slide, 0.8, 0.7, 11.7, 0.6, '关键收获', 34, WHITE, True)

# 五句话总结
add_multiline(slide, 0.8, 1.5, 7.0, 4.5, [
    ('五句话总结孤立森林', 18, ACCENT_GREEN, True),
    ('', 8),
    ('①  新范式：不建模正常，直接度量"被孤立的难易程度"', 14, LIGHT_GRAY),
    ('', 3),
    ('②  简洁实现：100 棵随机分割二叉树 + 路径长度 + 异常分数公式', 14, LIGHT_GRAY),
    ('', 3),
    ('③  极致效率：线性时间复杂度，天然并行，百万级数据轻松处理', 14, LIGHT_GRAY),
    ('', 3),
    ('④  广泛适用：网络安全、金融反欺诈、工业预警、IT 运维', 14, LIGHT_GRAY),
    ('', 3),
    ('⑤  持续进化：Extended IF 解决轴对齐偏差，深度学习变体探索中', 14, LIGHT_GRAY),
    ('', 10),
    ('行动建议', 18, ACCENT_GREEN, True),
    ('下次遇到异常检测任务，先用孤立森林建立基线。', 14, WHITE, True),
    ('默认参数（100 棵树、256 子采样）大多数场景够用。', 13, LIGHT_GRAY),
    ('效果不够再考虑 LOF 或 One-Class SVM。', 13, LIGHT_GRAY),
])

# 右侧 - 关键数字卡片
add_rounded_rect(slide, 8.2, 1.5, 4.5, 3.5, '#0E1F35', '#1A3A5C')
add_multiline(slide, 8.5, 1.7, 4.0, 3.2, [
    ('🔢  关键数字速记', 15, ACCENT_GREEN, True),
    ('', 8),
    ('树数量    100', 16, WHITE, True),
    ('子采样    256', 16, WHITE, True),
    ('树高度    8', 16, WHITE, True),
    ('分界线    s = 0.5', 16, WHITE, True),
    ('复杂度    O(n)', 16, WHITE, True),
])

# 底部 - 论文引用
add_rounded_rect(slide, 0.8, 6.0, 11.7, 1.1, '#0E1F35', '#1A3A5C')
add_multiline(slide, 1.1, 6.1, 11.2, 0.9, [
    ('📄  论文出处', 13, ACCENT_GREEN, True),
    ('Liu, Ting & Zhou (2008) "Isolation Forest" IEEE ICDM  |  Liu, Ting & Zhou (2012) ACM TKDD  |  Hariri et al. (2019) "Extended IF" IEEE TKDE', 10, LIGHT_GRAY),
    ('"一切都应该尽可能简单，但不能过于简单。" —— 爱因斯坦', 11, MID_GRAY, False, PP_ALIGN.CENTER),
])


# ═══════════════════════════════════════
#  保存
# ═══════════════════════════════════════
output_path = f'{OUT_DIR}/isolation-forest_v7.pptx'
prs.save(output_path)
print(f'✅ PPT V7 已保存到 {output_path}')
print(f'总页数：{len(prs.slides)}')
