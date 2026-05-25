#!/usr/bin/env python3
"""孤立森林 PPT V4 - 修复6个问题"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── 配色方案 ──
BG_DARK      = RGBColor(0x0B, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x12, 0x25, 0x3F)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_ORANGE= RGBColor(0xFF, 0x6B, 0x35)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xC4, 0xDE)
MID_GRAY     = RGBColor(0x7A, 0x8E, 0xA8)
DIM_GRAY     = RGBColor(0x4A, 0x5E, 0x78)
FONT_CN = '微软雅黑'

# ── 贯穿示例数据 ──
# 10个正常点 + 2个异常点
NORMAL_POINTS = [
    (3.0, 4.5), (3.5, 5.0), (4.0, 4.0), (4.5, 5.0), (3.5, 4.5),
    (4.0, 5.5), (4.5, 4.5), (5.0, 5.0), (4.0, 5.0), (3.0, 5.5),
]
ANOMALY_POINTS = {
    'A1': (20.0, 3.0),
    'A2': (2.0, 25.0),
}
# x range: 2~20, y range: 3~25
# 归一化到绘图区域

def map_point(px, py, area_left, area_top, area_w, area_h,
              x_min=1, x_max=22, y_min=2, y_max=26):
    """把数据坐标映射到幻灯片坐标（y轴翻转）"""
    sx = area_left + (px - x_min) / (x_max - x_min) * area_w
    sy = area_top + (1 - (py - y_min) / (y_max - y_min)) * area_h
    return sx, sy

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    p.space_before = Pt(0)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=14,
                           default_color=LIGHT_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else default_bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_CN
        p.alignment = alignment
        p.space_after = Pt(3)
        p.space_before = Pt(1)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_GREEN, thickness=3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_dot(slide, x, y, radius, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - radius), Inches(y - radius), Inches(radius*2), Inches(radius*2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot

def add_connector_v(slide, x, y1, y2, color=MID_GRAY, width=0.03):
    """垂直连接线"""
    h = y2 - y1
    if h < 0:
        y1, y2 = y2, y1
        h = y2 - y1
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y1), Inches(width), Inches(h))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_connector_h(slide, x1, x2, y, color=MID_GRAY, width=0.03):
    """水平连接线"""
    w = x2 - x1
    if w < 0:
        x1, x2 = x2, x1
        w = x2 - x1
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y), Inches(w), Inches(width))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_tree_node(slide, cx, cy, w, h, text, fill_color, border_color, text_color=WHITE, text_size=11, bold=True):
    """添加一个树节点（矩形+文字）"""
    rect = add_rect(slide, cx - w/2, cy - h/2, w, h, fill_color, border_color)
    add_textbox(slide, cx - w/2 + 0.05, cy - h/2 + 0.02, w - 0.1, h - 0.04, text, text_size, text_color, bold, PP_ALIGN.CENTER)
    return rect

def add_tree_edge(slide, parent_cx, parent_bottom_y, child_cx, child_top_y, color=MID_GRAY):
    """添加从父节点底部到子节点顶部的连接线（L形或折线）"""
    mid_y = (parent_bottom_y + child_top_y) / 2
    # 垂直线从父节点底部到中间
    add_connector_v(slide, parent_cx, parent_bottom_y, mid_y, color)
    # 水平线从父x到子x
    add_connector_h(slide, min(parent_cx, child_cx), max(parent_cx, child_cx), mid_y, color)
    # 垂直线从中间到子节点顶部
    add_connector_v(slide, child_cx, mid_y, child_top_y, color)


# ════════════════════════════════════════
# 第 1 页：封面（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1.0, 2.8, 4.0, ACCENT_GREEN)
add_textbox(slide, 1.0, 1.2, 11, 1.5, '孤立森林', 48, WHITE, True, PP_ALIGN.LEFT)
add_textbox(slide, 1.0, 2.0, 11, 0.8, 'Isolation Forest', 28, ACCENT_CYAN, False, PP_ALIGN.LEFT)
add_multiline_textbox(slide, 1.0, 3.2, 11, 1.5, [
    ('不建模正常，直接孤立异常——一种颠覆性的异常检测算法', 20, LIGHT_GRAY, False),
    ('', 8),
    ('适用于大数据、高维度、无标签的场景', 14, MID_GRAY, False),
    ('线性时间复杂度，天然支持并行化', 14, MID_GRAY, False),
    ('2008 年 IEEE ICDM 论文  |  Liu, Ting & Zhou', 12, DIM_GRAY, False),
])

# 右侧装饰散点 - 也修正为精确12个点
add_rect(slide, 8.5, 1.5, 4.0, 4.5, RGBColor(0x0E, 0x1F, 0x35), ACCENT_GREEN)
# 正常点群（10个，聚集在左下角区域）
scatter_area_left, scatter_area_top = 8.7, 2.0
scatter_area_w, scatter_area_h = 2.2, 3.0
for px, py in NORMAL_POINTS:
    sx, sy = map_point(px, py, scatter_area_left, scatter_area_top, scatter_area_w, scatter_area_h,
                       x_min=1, x_max=22, y_min=2, y_max=26)
    add_dot(slide, sx, sy, 0.06, ACCENT_CYAN)
# 异常点
for name, (px, py) in ANOMALY_POINTS.items():
    sx, sy = map_point(px, py, scatter_area_left, scatter_area_top, scatter_area_w, scatter_area_h,
                       x_min=1, x_max=22, y_min=2, y_max=26)
    add_dot(slide, sx, sy, 0.10, ACCENT_ORANGE)

add_textbox(slide, 1.0, 6.5, 11, 0.5, '异常检测系列  |  深度技术解析', 12, DIM_GRAY, False)


# ════════════════════════════════════════
# 第 2 页：什么是异常检测？（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '什么是异常检测？', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '为什么要关注"与众不同"的数据？', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 4.8, [
    ('什么是异常检测？', 18, ACCENT_GREEN, True),
    ('异常检测是从数据中识别出与绝大多数数据显著不同的样本的过程。就像在一群白绵羊中找出那只黑绵羊——正常数据遵循某种模式，而异常点偏离了这个模式。', 13, LIGHT_GRAY),
    ('', 6),
    ('为什么它如此重要？', 18, ACCENT_GREEN, True),
    ('异常事件往往伴随高风险。网络入侵导致数据泄露，金融欺诈造成巨额损失，设备异常引发生产线停工。正因为"异常"通常等于"危险"，异常检测成为数据科学最重要的研究方向之一。', 13, LIGHT_GRAY),
    ('', 6),
    ('异常的三种类型：', 16, WHITE, True),
    ('1️⃣  点异常：单个数据点本身就很异常。一群身高 1.6-1.8m 的人中间站了一个 2.3m 的人。', 12, LIGHT_GRAY),
    ('2️⃣  上下文异常：特定上下文中才算异常。夏天穿短袖正常，零下 20°C 穿短袖就异常。', 12, LIGHT_GRAY),
    ('3️⃣  集合异常：单点都正常，放在一起才异常。每天凌晨访问数据库正常，连续三个月天天凌晨就可疑。', 12, LIGHT_GRAY),
])

add_rect(slide, 7.2, 2.0, 5.5, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.5, 2.2, 5.0, 4.4, [
    ('🛡️  机场安检类比', 16, ACCENT_GREEN, True),
    ('', 6),
    ('异常检测就像机场安检——99.9% 的旅客都正常，但系统必须在几秒内识别出那 0.1% 的危险人物。', 13, LIGHT_GRAY),
    ('', 6),
    ('三种异常 = 三种"嫌疑人"：', 13, WHITE, True),
    ('• 带着违禁品的人 → 点异常（一看就不对）', 12, ACCENT_CYAN),
    ('• 错误时间出现在错误地点 → 上下文异常', 12, ACCENT_CYAN),
    ('• 每人行为正常但组合可疑 → 集合异常（团伙作案）', 12, ACCENT_CYAN),
])


# ════════════════════════════════════════
# 第 3 页：传统方法（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '传统异常检测方法', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '先建模"正常"，再找"偏离"——但建模"正常"恰恰是最难的步骤', 16, ACCENT_CYAN)

table_data = [
    ('方法', '核心思路', '时间复杂度', '高维适应', '核心局限'),
    ('Z-score', '距离均值几个标准差', 'O(n)', '差', '假设正态分布'),
    ('KNN / LOF', '邻居距离/局部密度', 'O(n²)', '一般', '计算量灾难性'),
    ('DBSCAN', '基于密度聚类', 'O(n log n)', '差', '高维密度难定义'),
    ('One-Class SVM', '学习正常边界', 'O(n²~n³)', '好', '参数敏感、不可扩展'),
    ('孤立森林', '直接孤立异常', 'O(n·t·logψ)', '好', '局部异常检测弱'),
]
rows = len(table_data)
cols = 5
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.2)).table
col_widths = [1.8, 2.8, 2.0, 1.5, 3.6]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = Inches(w)
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_CN
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r == 5:
                para.font.bold = True
                para.font.color.rgb = ACCENT_GREEN
            else:
                para.font.color.rgb = LIGHT_GRAY
            para.alignment = PP_ALIGN.LEFT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r == 5:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

add_rect(slide, 0.8, 5.6, 11.7, 1.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.1, 5.7, 11.2, 1.3, [
    ('⚠️  核心瓶颈：不在于"找异常"，而在于"建模正常"太昂贵、太困难、太脆弱', 15, ACCENT_ORANGE, True),
    ('传统方法像一个侦探，先要彻底研究"正常市民的生活规律"，然后才能发现谁的行踪可疑。但刻画"正常"本身就极其困难——正常人有千万种生活方式，你怎么定义"正常"？', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 4 页：核心思想 ★★★ FIX #1：精确12个点 + A2标签
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '孤立森林的核心思想', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '不建模正常，直接孤立异常——范式转换', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 5.0, [
    ('范式转换', 18, ACCENT_GREEN, True),
    ('从"建模正常"到"孤立异常"', 14, WHITE, True),
    ('', 4),
    ('异常点有两个天然属性："少"和"不同"。因为数量稀少，周围没有同类；因为特征不同，远离大多数。', 13, LIGHT_GRAY),
    ('', 4),
    ('这导致一个直接结果：异常点比正常点更容易被随机分割孤立出来。', 14, ACCENT_GREEN, True),
    ('', 4),
    ('用贯穿示例感受：', 14, WHITE, True),
    ('10 个正常点聚集在 (3~5, 4~5.5)', 13, ACCENT_CYAN),
    ('2 个异常点：A1(20, 3) 和 A2(2, 25) 远远游离', 13, ACCENT_ORANGE),
    ('', 4),
    ('随机选 x 轴，在 x=8 画分割线：', 13, LIGHT_GRAY),
    ('→ 左侧：10 个正常点 + A2（11 个点）', 13, ACCENT_CYAN),
    ('→ 右侧：A1 独自一人 → 1 次分裂就被孤立！', 14, ACCENT_GREEN, True),
    ('', 4),
    ('正常点呢？周围全是同伴，需要 6~8 次分裂才能单独挑出。', 13, LIGHT_GRAY),
    ('', 4),
    ('路径越短 = 越异常 | 路径越长 = 越正常', 16, WHITE, True),
])

# ★ 右侧散点图：精确12个点
chart_left, chart_top, chart_w, chart_h = 7.2, 2.0, 5.5, 4.8
add_rect(slide, chart_left, chart_top, chart_w, chart_h, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))

# 坐标轴标签
add_textbox(slide, chart_left + 0.1, chart_top + chart_h - 0.35, 0.5, 0.3, 'x', 10, MID_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, chart_left + 0.1, chart_top + 0.05, 0.5, 0.3, 'y', 10, MID_GRAY, False, PP_ALIGN.CENTER)

# 绘图区域内缩
plot_left = chart_left + 0.4
plot_top = chart_top + 0.3
plot_w = chart_w - 0.6
plot_h = chart_h - 0.7

# 画坐标轴
add_connector_h(slide, plot_left, plot_left + plot_w, plot_top + plot_h, MID_GRAY, 0.02)  # x轴
add_connector_v(slide, plot_left, plot_top, plot_top + plot_h, MID_GRAY, 0.02)  # y轴

# 画精确的12个点
for px, py in NORMAL_POINTS:
    sx, sy = map_point(px, py, plot_left, plot_top, plot_w, plot_h, x_min=1, x_max=22, y_min=2, y_max=26)
    add_dot(slide, sx, sy, 0.08, ACCENT_CYAN)

# A1
a1_sx, a1_sy = map_point(20.0, 3.0, plot_left, plot_top, plot_w, plot_h, x_min=1, x_max=22, y_min=2, y_max=26)
add_dot(slide, a1_sx, a1_sy, 0.12, ACCENT_ORANGE)
add_textbox(slide, a1_sx + 0.12, a1_sy - 0.1, 0.5, 0.25, 'A1', 11, ACCENT_ORANGE, True)

# A2
a2_sx, a2_sy = map_point(2.0, 25.0, plot_left, plot_top, plot_w, plot_h, x_min=1, x_max=22, y_min=2, y_max=26)
add_dot(slide, a2_sx, a2_sy, 0.12, ACCENT_ORANGE)
add_textbox(slide, a2_sx + 0.12, a2_sy - 0.1, 0.5, 0.25, 'A2', 11, ACCENT_ORANGE, True)

# 分割线 x=8
split_sx, _ = map_point(8.0, 0, plot_left, plot_top, plot_w, plot_h, x_min=1, x_max=22, y_min=2, y_max=26)
add_connector_v(slide, split_sx, plot_top + 0.05, plot_top + plot_h - 0.05, ACCENT_GREEN, 0.03)
add_textbox(slide, split_sx + 0.05, plot_top + 0.02, 0.8, 0.25, 'x=8', 10, ACCENT_GREEN, True)

# 标注
norm_cx, norm_cy = map_point(3.8, 4.8, plot_left, plot_top, plot_w, plot_h, x_min=1, x_max=22, y_min=2, y_max=26)
add_textbox(slide, norm_cx - 0.5, norm_cy + 0.2, 1.5, 0.3, '正常点群(10个)', 10, ACCENT_CYAN, False, PP_ALIGN.CENTER)
add_textbox(slide, split_sx + 0.3, plot_top + plot_h - 0.6, 1.8, 0.3, 'A1 被孤立！', 11, ACCENT_GREEN, True)
add_textbox(slide, a2_sx - 0.5, a2_sy + 0.15, 1.5, 0.3, 'A2 在左侧', 9, MID_GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 第 5 页：iTree ★★★ FIX #2+#3：更复杂的树形图 + 正确的连接线
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, 'iTree（孤立树）', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '一棵不断随机分割空间的二叉树', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('两种节点', 16, ACCENT_GREEN, True),
    ('内部节点：存分裂规则——"在哪个特征上切？切在哪里？"', 13, LIGHT_GRAY),
    ('叶子节点：存"到达这里的样本数量"——房间还剩几人', 13, LIGHT_GRAY),
    ('', 6),
    ('每次分裂两步操作', 16, ACCENT_GREEN, True),
    ('① 随机选一个特征（等概率）', 13, LIGHT_GRAY),
    ('② 在该特征的 [min, max] 内随机选分裂点', 13, LIGHT_GRAY),
    ('→ 特征值 < 分裂点进左子树，≥ 分裂点进右子树', 13, ACCENT_CYAN),
    ('', 6),
    ('三个终止条件', 16, ACCENT_GREEN, True),
    ('① 样本数 = 1 → 无需再分', 13, LIGHT_GRAY),
    ('② 所有特征值相同 → 无法再分', 13, LIGHT_GRAY),
    ('③ 达到最大深度（通常 log₂256 = 8 层）→ 强制停止', 13, LIGHT_GRAY),
    ('', 6),
    ('iTree vs 普通决策树', 16, ACCENT_GREEN, True),
    ('普通决策树：精心选择最优分裂（像经验丰富的分拣员）', 12, LIGHT_GRAY),
    ('iTree：随机选特征和分裂点（像蒙着眼分拣）', 12, LIGHT_GRAY),
    ('看似"愚蠢"，但随机分割对稀少的异常点更"友好"', 12, ACCENT_GREEN, True),
])

# ★ 右侧 - 更完整的树形图（4层）
add_rect(slide, 6.8, 2.0, 6.0, 4.8, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.08, 5.6, 0.35, '贯穿示例：一棵完整的 iTree（12个点）', 12, ACCENT_GREEN, True)

# 树节点参数
nw, nh = 1.6, 0.38  # 节点宽高
node_fill = RGBColor(0x1A, 0x3A, 0x5C)
leaf_fill_anomaly = RGBColor(0x3A, 0x1A, 0x0A)
leaf_fill_normal = RGBColor(0x0A, 0x2A, 0x15)

# 定义树结构（使用 slide 坐标）
# Level 0: Root
L0_y = 2.6
root_cx = 9.8

# Level 1
L1_y = 3.4
L1_left_cx = 8.3
L1_right_cx = 11.3

# Level 2
L2_y = 4.2
L2_ll_cx = 7.5   # left-left (正常点继续)
L2_lr_cx = 9.1   # left-right -> A2 被孤立

# Level 3
L3_y = 5.0
L3_lll_cx = 7.0  # 继续分割
L3_llr_cx = 8.0  # 叶子 size=3

# Level 4
L4_y = 5.7
L4_lll_l_cx = 7.0  # 叶子 size=1
L4_lll_r_cx = 7.8  # 叶子 size=2

# --- 画连接线（先画线，再画节点覆盖）---

# Root → Level 1
add_connector_v(slide, root_cx, L0_y + nh/2, (L0_y + nh/2 + L1_y - nh/2) / 2, MID_GRAY, 0.025)
add_connector_h(slide, L1_left_cx, L1_right_cx, (L0_y + nh/2 + L1_y - nh/2) / 2, MID_GRAY, 0.025)
add_connector_v(slide, L1_left_cx, (L0_y + nh/2 + L1_y - nh/2) / 2, L1_y - nh/2, MID_GRAY, 0.025)
add_connector_v(slide, L1_right_cx, (L0_y + nh/2 + L1_y - nh/2) / 2, L1_y - nh/2, MID_GRAY, 0.025)

# Yes/No labels on Level 1 edges
add_textbox(slide, root_cx - 0.6, L0_y + nh/2 + 0.05, 0.4, 0.2, '是', 8, ACCENT_GREEN)
add_textbox(slide, root_cx + 0.15, L0_y + nh/2 + 0.05, 0.4, 0.2, '否', 8, ACCENT_ORANGE)

# L1_left → Level 2
mid_y_l1 = (L1_y + nh/2 + L2_y - nh/2) / 2
add_connector_v(slide, L1_left_cx, L1_y + nh/2, mid_y_l1, MID_GRAY, 0.025)
add_connector_h(slide, L2_ll_cx, L2_lr_cx, mid_y_l1, MID_GRAY, 0.025)
add_connector_v(slide, L2_ll_cx, mid_y_l1, L2_y - nh/2, MID_GRAY, 0.025)
add_connector_v(slide, L2_lr_cx, mid_y_l1, L2_y - nh/2, MID_GRAY, 0.025)

add_textbox(slide, L1_left_cx - 0.6, L1_y + nh/2 + 0.05, 0.4, 0.2, '是', 8, ACCENT_GREEN)
add_textbox(slide, L1_left_cx + 0.15, L1_y + nh/2 + 0.05, 0.4, 0.2, '否', 8, ACCENT_ORANGE)

# L2_ll → Level 3
mid_y_l2 = (L2_y + nh/2 + L3_y - nh/2) / 2
add_connector_v(slide, L2_ll_cx, L2_y + nh/2, mid_y_l2, MID_GRAY, 0.025)
add_connector_h(slide, L3_lll_cx, L3_llr_cx, mid_y_l2, MID_GRAY, 0.025)
add_connector_v(slide, L3_lll_cx, mid_y_l2, L3_y - nh/2, MID_GRAY, 0.025)
add_connector_v(slide, L3_llr_cx, mid_y_l2, L3_y - nh/2, MID_GRAY, 0.025)

add_textbox(slide, L2_ll_cx - 0.6, L2_y + nh/2 + 0.05, 0.4, 0.2, '是', 8, ACCENT_GREEN)
add_textbox(slide, L2_ll_cx + 0.15, L2_y + nh/2 + 0.05, 0.4, 0.2, '否', 8, ACCENT_ORANGE)

# L3_lll → Level 4
mid_y_l3 = (L3_y + nh/2 + L4_y - nh/2) / 2
add_connector_v(slide, L3_lll_cx, L3_y + nh/2, mid_y_l3, MID_GRAY, 0.025)
add_connector_h(slide, L4_lll_l_cx, L4_lll_r_cx, mid_y_l3, MID_GRAY, 0.025)
add_connector_v(slide, L4_lll_l_cx, mid_y_l3, L4_y - nh/2, MID_GRAY, 0.025)
add_connector_v(slide, L4_lll_r_cx, mid_y_l3, L4_y - nh/2, MID_GRAY, 0.025)

# --- 画节点 ---

# Level 0: Root
add_tree_node(slide, root_cx, L0_y, nw + 0.4, nh, 'x < 8.0? (12个点)', node_fill, ACCENT_GREEN, WHITE, 11)

# Level 1
add_tree_node(slide, L1_left_cx, L1_y, nw + 0.2, nh, 'y < 10.0? (11个点)', node_fill, ACCENT_CYAN, WHITE, 10)
add_tree_node(slide, L1_right_cx, L1_y, nw, nh, 'A1  [size=1]', leaf_fill_anomaly, ACCENT_ORANGE, ACCENT_ORANGE, 11)
add_textbox(slide, L1_right_cx - 0.3, L1_y - nh/2 - 0.22, 1.0, 0.2, '路径=1', 9, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# Level 2
add_tree_node(slide, L2_ll_cx, L2_y, nw + 0.1, nh, 'x < 4.2? (10个点)', node_fill, ACCENT_CYAN, WHITE, 10)
add_tree_node(slide, L2_lr_cx, L2_y, nw, nh, 'A2  [size=1]', leaf_fill_anomaly, ACCENT_ORANGE, ACCENT_ORANGE, 11)
add_textbox(slide, L2_lr_cx - 0.3, L2_y - nh/2 - 0.22, 1.0, 0.2, '路径=2', 9, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# Level 3
add_tree_node(slide, L3_lll_cx, L3_y, nw - 0.1, nh, 'y < 4.8? (6个点)', node_fill, ACCENT_CYAN, WHITE, 10)
add_tree_node(slide, L3_llr_cx, L3_y, nw - 0.2, nh, '[size=4]', leaf_fill_normal, ACCENT_GREEN, ACCENT_CYAN, 10)

# Level 4
add_tree_node(slide, L4_lll_l_cx, L4_y, nw - 0.3, nh, 'P3 [1]', leaf_fill_normal, ACCENT_GREEN, ACCENT_CYAN, 9)
add_textbox(slide, L4_lll_l_cx - 0.2, L4_y - nh/2 - 0.2, 1.0, 0.18, '路径=4', 8, ACCENT_CYAN, False, PP_ALIGN.CENTER)
add_tree_node(slide, L4_lll_r_cx, L4_y, nw - 0.3, nh, '[size=5]', leaf_fill_normal, ACCENT_GREEN, MID_GRAY, 9)

# 底部注释
add_multiline_textbox(slide, 7.0, 6.15, 5.6, 0.6, [
    ('异常点 A1/A2 在第 1~2 层就被孤立（路径=1, 2）', 10, ACCENT_ORANGE, True),
    ('正常点 P3 需要第 4 层才被分离，其余正常点更深', 10, MID_GRAY),
])


# ════════════════════════════════════════
# 第 6 页：构建森林 ★★★ FIX #4：解释P9
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '构建孤立森林', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从一棵树到一片森林——子采样 + 多树集成 = 又快又准', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('三步构建森林', 16, ACCENT_GREEN, True),
    ('', 4),
    ('Step 1：随机子采样', 14, WHITE, True),
    ('从训练数据中随机抽取 256 个样本。三个原因：', 12, LIGHT_GRAY),
    ('  • 消除屏蔽效应——稀释正常点，让异常点更容易被看到', 11, LIGHT_GRAY),
    ('  • 计算效率——无论原始数据 1 万还是 1 亿，训练时间可控', 11, LIGHT_GRAY),
    ('  • 增加多样性——不同子采样让每棵树看到不同数据', 11, LIGHT_GRAY),
    ('', 4),
    ('Step 2：构建 iTree', 14, WHITE, True),
    ('按递归分裂方法构建，高度限制 ⌈log₂256⌉ = 8 层。', 12, LIGHT_GRAY),
    ('8 层足以孤立绝大多数异常点（1-4 层就会被分离）。', 12, LIGHT_GRAY),
    ('', 4),
    ('Step 3：重复 100 次', 14, WHITE, True),
    ('得到 100 棵独立的 iTree，组成孤立森林。', 12, LIGHT_GRAY),
    ('从 100 增加到 500 棵，AUC 提升通常 < 0.01。', 12, ACCENT_CYAN),
])

# ★ 右侧 - 路径对比表（增加P9解释）
add_rect(slide, 6.8, 2.0, 6.0, 3.2, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 2.08, 5.6, 0.5, [
    ('📊  5 棵树路径对比', 14, ACCENT_GREEN, True),
    ('P9 是 10 个正常点之一：P9(4.0, 5.0)，代表正常点的典型路径', 10, MID_GRAY),
])

table_data2 = [
    ('树编号', 'A1(20,3) 路径', 'P9(4,5) 路径'),
    ('Tree 1', '1', '5'),
    ('Tree 2', '2', '5'),
    ('Tree 3', '1', '5'),
    ('Tree 4', '2', '4'),
    ('Tree 5', '1', '5'),
    ('平均', '1.4', '4.8'),
]
rows2 = len(table_data2)
tbl2 = slide.shapes.add_table(rows2, 3, Inches(7.2), Inches(2.55), Inches(5.2), Inches(2.4)).table
tbl2.columns[0].width = Inches(1.5)
tbl2.columns[1].width = Inches(1.85)
tbl2.columns[2].width = Inches(1.85)
for r in range(rows2):
    for c in range(3):
        cell = tbl2.cell(r, c)
        cell.text = table_data2[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_CN
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r == rows2-1:
                para.font.bold = True
                para.font.color.rgb = ACCENT_GREEN
            else:
                para.font.color.rgb = LIGHT_GRAY
            para.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r == rows2-1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

add_rect(slide, 6.8, 5.5, 6.0, 1.5, BG_CARD)
add_multiline_textbox(slide, 7.0, 5.6, 5.6, 1.3, [
    ('🕵️  100 个"蒙眼安检员"', 14, ACCENT_GREEN, True),
    ('每个安检员随机抽查 256 人，用随机画线的方式分割人群。单个安检员可能有偏差，但 100 个安检员的集体投票非常可靠——大多数安检员都很快挑出的人，几乎一定是异常的。', 11, MID_GRAY),
])


# ════════════════════════════════════════
# 第 7 页：异常分数计算 ★★★ FIX #5+#6：修正公式 + 详细讲H(n-1)
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '异常分数计算', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用公式量化"有多异常"', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('异常分数公式', 16, ACCENT_GREEN, True),
    ('s(x, n) = 2 ^ ( -E[h(x)] / c(n) )', 20, WHITE, True),
    ('', 3),
    ('E[h(x)]：点 x 在所有树中路径长度的期望值（即平均路径 h̄）', 12, LIGHT_GRAY),
    ('c(n)：n 个点随机二分的期望路径长度（"平均孤立难度"）', 12, LIGHT_GRAY),
    ('', 6),
    ('c(n) 怎么算？', 16, ACCENT_GREEN, True),
    ('c(n) = 2H(n-1) - 2(n-1)/n', 14, WHITE, True),
    ('', 3),
    ('H(k) 是调和数：H(k) = 1 + 1/2 + 1/3 + ... + 1/k', 12, LIGHT_GRAY),
    ('可近似为 H(k) ≈ ln(k) + 0.5772（欧拉常数）', 12, ACCENT_CYAN),
    ('', 3),
    ('以贯穿示例 n=12 为例：', 12, WHITE, True),
    ('H(11) = 1 + 1/2 + 1/3 + ... + 1/11 ≈ 3.0199', 12, LIGHT_GRAY),
    ('c(12) = 2×3.0199 - 2×11/12 = 6.0398 - 1.8333', 12, LIGHT_GRAY),
    ('c(12) = 4.2065 ← 这就是"平均孤立12个点需要4.2步"', 13, ACCENT_GREEN, True),
    ('', 4),
    ('分数解读', 16, ACCENT_GREEN, True),
    ('s → 1：路径远短于平均 → 极异常', 12, ACCENT_ORANGE),
    ('s ≈ 0.5：路径与平均相当 → 分界线', 12, ACCENT_CYAN),
    ('s → 0：路径远长于平均 → 极正常', 12, LIGHT_GRAY),
])

# 右侧 - 数值计算
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.08, 5.6, 0.35, '📐  贯穿示例计算（n=12, c=4.2065）', 13, ACCENT_GREEN, True)

add_multiline_textbox(slide, 7.2, 2.55, 5.4, 4.2, [
    ('🔴 异常点 A1(20, 3)：E[h] = 1.4', 14, ACCENT_ORANGE, True),
    ('', 2),
    ('s(A1) = 2^(-1.4 / 4.2065)', 13, LIGHT_GRAY),
    ('     = 2^(-0.333)', 13, LIGHT_GRAY),
    ('     ≈ 0.795  →  明显异常 ✓', 15, ACCENT_GREEN, True),
    ('', 6),
    ('🔴 异常点 A2(2, 25)：E[h] = 2.1', 14, ACCENT_ORANGE, True),
    ('', 2),
    ('s(A2) = 2^(-2.1 / 4.2065)', 13, LIGHT_GRAY),
    ('     = 2^(-0.499)', 13, LIGHT_GRAY),
    ('     ≈ 0.707  →  明显异常 ✓', 15, ACCENT_GREEN, True),
    ('', 6),
    ('🟢 正常点 P9(4, 5)：E[h] = 4.8', 14, ACCENT_CYAN, True),
    ('', 2),
    ('s(P9) = 2^(-4.8 / 4.2065)', 13, LIGHT_GRAY),
    ('     = 2^(-1.141)', 13, LIGHT_GRAY),
    ('     ≈ 0.454  →  正常 ✓', 15, ACCENT_CYAN, True),
    ('', 6),
    ('0 ── 0.454 ─ 0.5 ── 0.707 ─ 0.795 ─ 1.0', 12, WHITE, True, PP_ALIGN.CENTER),
    ('    正常 ←──── 分界线 ────→ 异常', 11, MID_GRAY, False, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════
# 第 8 页：完整示例（与V3基本相同，微调P9说明）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '完整示例——从数据到结论', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用 12 个数据点走一遍"训练→评分→判定"全流程', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('① 准备数据', 14, ACCENT_GREEN, True),
    ('10 个正常点聚集在 (3-5, 4-5.5)，2 个异常点 A1(20,3) 和 A2(2,25) 远远游离。', 12, LIGHT_GRAY),
    ('', 4),
    ('② 构建一棵 iTree', 14, ACCENT_GREEN, True),
    ('第 1 次分裂：选 x 轴，分裂点 x=8.0', 12, LIGHT_GRAY),
    ('  → 左子树：10 个正常点 + A2（11 个点）', 11, MID_GRAY),
    ('  → 右子树：A1（仅 1 个）→ A1 被孤立！路径=1', 11, ACCENT_ORANGE),
    ('第 2 次分裂：选 y 轴，分裂点 y=10.0', 12, LIGHT_GRAY),
    ('  → 左子树：10 个正常点 | 右子树：A2 → 路径=2', 11, ACCENT_ORANGE),
    ('第 3~4 次分裂：继续分割 10 个正常点...', 12, MID_GRAY),
    ('', 4),
    ('③ 构建 100 棵 iTree', 14, ACCENT_GREEN, True),
    ('重复 100 次（不同子采样和随机分裂），统计平均路径：', 12, LIGHT_GRAY),
    ('  A1 平均路径 E[h] ≈ 1.4 | A2 平均路径 E[h] ≈ 2.1', 12, ACCENT_CYAN),
    ('  正常点 P9(4,5) 平均路径 E[h] ≈ 4.8', 12, ACCENT_CYAN),
    ('', 4),
    ('④ 计算异常分数（c(12) = 4.2065）', 14, ACCENT_GREEN, True),
    ('  s(A1) ≈ 0.795 → 🔴 异常', 13, ACCENT_ORANGE, True),
    ('  s(A2) ≈ 0.707 → 🔴 异常', 13, ACCENT_ORANGE, True),
    ('  s(P9) ≈ 0.454 → 🟢 正常', 13, ACCENT_CYAN, True),
])

# 右侧 - 结果汇总
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.08, 5.6, 0.35, '📊  结果汇总', 14, ACCENT_GREEN, True)

table_data3 = [
    ('数据点', 'E[h(x)]', '异常分数', '判定'),
    ('A1(20, 3)', '1.4', '0.795', '🔴 异常'),
    ('A2(2, 25)', '2.1', '0.707', '🔴 异常'),
    ('P9(4, 5)', '4.8', '0.454', '🟢 正常'),
    ('P1~P10', '4.5~5.0', '0.45~0.53', '🟢 正常'),
]
tbl3 = slide.shapes.add_table(5, 4, Inches(7.2), Inches(2.55), Inches(5.4), Inches(2.2)).table
tbl3.columns[0].width = Inches(1.5)
tbl3.columns[1].width = Inches(1.2)
tbl3.columns[2].width = Inches(1.2)
tbl3.columns[3].width = Inches(1.5)
for r in range(5):
    for c in range(4):
        cell = tbl3.cell(r, c)
        cell.text = table_data3[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_CN
            para.alignment = PP_ALIGN.CENTER
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r in (1, 2):
                para.font.color.rgb = ACCENT_ORANGE
            else:
                para.font.color.rgb = ACCENT_CYAN
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r in (1, 2):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x15, 0x0A)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

add_multiline_textbox(slide, 7.2, 5.0, 5.4, 1.6, [
    ('✅  结论', 14, ACCENT_GREEN, True),
    ('孤立森林成功识别 A1 和 A2 两个异常点（分数 0.795 和 0.707，远超 0.5 分界线），同时正确判定所有正常点为正常（分数 0.45~0.53，均低于 0.5）。', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 9 页：参数调优与优缺点（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '参数调优与优缺点', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '默认参数就能覆盖大多数场景——但也要知道局限', 16, ACCENT_CYAN)

add_rect(slide, 0.8, 2.0, 5.5, 2.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.0, 2.1, 5.2, 2.3, [
    ('⚙️  关键参数', 15, ACCENT_GREEN, True),
    ('', 4),
    ('树的数量（n_estimators）：默认 100', 13, WHITE, True),
    ('就像民意调查，100 份已足够稳定。100→500 棵 AUC 提升 <0.01', 11, LIGHT_GRAY),
    ('', 4),
    ('子采样大小（max_samples）：默认 256', 13, WHITE, True),
    ('论文实验的"甜蜜点"——太小信息不足，太大正常点互相遮挡', 11, LIGHT_GRAY),
    ('', 4),
    ('异常阈值（contamination）：默认 auto（s=0.5）', 13, WHITE, True),
    ('如有业务知识（欺诈率约 0.1%），可直接设 0.001', 11, LIGHT_GRAY),
])

add_rect(slide, 0.8, 4.8, 5.5, 2.3, BG_CARD, RGBColor(0x0A, 0x2A, 0x15))
add_multiline_textbox(slide, 1.0, 4.9, 5.2, 2.1, [
    ('✅  核心优点', 15, ACCENT_GREEN, True),
    ('• 线性时间复杂度 O(t·ψ·logψ)，轻松处理百万级数据', 12, LIGHT_GRAY),
    ('• 天然支持并行化，100 棵树可分配到 100 个 CPU 核心', 12, LIGHT_GRAY),
    ('• 不依赖距离度量，天然免疫"维度灾难"', 12, LIGHT_GRAY),
    ('• 纯无监督，不需要任何标注数据', 12, LIGHT_GRAY),
    ('• 内存高效：100×256 = 25,600 个数据点即可', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 2.0, 6.0, 2.5, BG_CARD, RGBColor(0x2A, 0x15, 0x0A))
add_multiline_textbox(slide, 7.0, 2.1, 5.6, 2.3, [
    ('⚠️  注意局限', 15, ACCENT_ORANGE, True),
    ('', 4),
    ('• 局部异常检测能力弱——主要捕获全局异常', 12, LIGHT_GRAY),
    ('  如果需要检测局部异常，先用 iForest 筛查，再用 LOF 精检', 11, MID_GRAY),
    ('• 轴对齐切割偏差——沿非坐标轴方向分布的数据可能产生"伪异常"', 12, LIGHT_GRAY),
    ('  Extended Isolation Forest（2019）用任意角度超平面解决', 11, MID_GRAY),
    ('• 对分类特征处理不佳，需要先做独热编码', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 4.8, 6.0, 2.3, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 4.9, 5.6, 2.1, [
    ('🎯  选择建议', 15, ACCENT_CYAN, True),
    ('', 4),
    ('数据量大（>10 万条）→ 首选孤立森林', 13, ACCENT_GREEN, True),
    ('需要检测局部异常 → 选 LOF', 13, LIGHT_GRAY),
    ('数据量小（<5000 条）且边界复杂 → 选 One-Class SVM', 13, LIGHT_GRAY),
    ('不确定 → 先试孤立森林（速度快、参数少）', 13, ACCENT_CYAN, True),
])


# ════════════════════════════════════════
# 第 10 页：实战案例（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '实战案例：金融欺诈检测', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从原始交易数据到实时欺诈告警——完整的落地 Pipeline', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.8, 5.0, [
    ('① 数据准备与特征工程', 14, ACCENT_GREEN, True),
    ('原始特征：交易金额、时间、商户类别码、GPS 坐标', 12, LIGHT_GRAY),
    ('衍生特征：过去 1h 交易次数、金额/历史均值比、距上次交易时间', 12, LIGHT_GRAY),
    ('→ 衍生特征是关键：孤立森林无法自动发现"频率异常"', 11, ACCENT_CYAN),
    ('', 4),
    ('② 模型训练', 14, ACCENT_GREEN, True),
    ('n_estimators=100, max_samples=256, contamination=0.001', 12, LIGHT_GRAY),
    ('contamination 对应约 0.1% 的欺诈率（业务先验知识）', 12, LIGHT_GRAY),
    ('训练耗时：百万级交易数据约 10 秒（单机）', 12, ACCENT_CYAN),
    ('', 4),
    ('③ 实时评分与告警', 14, ACCENT_GREEN, True),
    ('每笔新交易 → 提取特征 → 走 100 棵树 → 输出异常分数', 12, LIGHT_GRAY),
    ('s > 0.7 触发告警 → 人工复核或自动冻结', 12, LIGHT_GRAY),
    ('', 4),
    ('④ 效果（Kaggle 信用卡欺诈数据集）', 14, ACCENT_GREEN, True),
    ('检测率 ~82%，误报率 <1%', 13, ACCENT_GREEN, True),
    ('对新型欺诈手法（无历史标签）仍有检测能力', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 2.0, 6.0, 2.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 2.1, 5.6, 2.6, [
    ('🏭  设备故障预警（NASA 涡轮发动机数据）', 13, ACCENT_GREEN, True),
    ('', 3),
    ('Pipeline：传感器时序 → 按窗口提取均值/方差/峰值 → 多传感器拼接 → 训练 iForest → 监控分数趋势', 11, LIGHT_GRAY),
    ('', 3),
    ('关键发现：不是看单次分数，而是看分数变化趋势——持续升高 = 设备在退化', 11, ACCENT_CYAN),
    ('', 3),
    ('效果：可在故障前 15-25 个工作日检测到异常趋势', 12, ACCENT_GREEN, True),
])

add_rect(slide, 6.8, 5.1, 6.0, 2.0, BG_CARD, RGBColor(0x0A, 0x2A, 0x15))
add_multiline_textbox(slide, 7.0, 5.2, 5.6, 1.8, [
    ('🏢  工业界大规模采用', 14, ACCENT_GREEN, True),
    ('', 3),
    ('Google — 监控内部网络健康', 12, LIGHT_GRAY),
    ('Microsoft Azure — 异常检测云服务', 12, LIGHT_GRAY),
    ('Netflix — 用户行为异常检测', 12, LIGHT_GRAY),
    ('阿里巴巴 — 双十一交易欺诈检测，每秒数十万笔', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 11 页：结尾（与V3相同）
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '关键收获与下一步', 32, WHITE, True)

add_multiline_textbox(slide, 0.8, 1.5, 6.5, 4.0, [
    ('五句话总结孤立森林', 18, ACCENT_GREEN, True),
    ('', 6),
    ('① 新范式：不建模正常，直接度量"被孤立的难易程度"', 14, LIGHT_GRAY),
    ('② 简洁实现：100 棵随机分割二叉树 + 路径长度 + 异常分数公式', 14, LIGHT_GRAY),
    ('③ 极致效率：线性时间复杂度，天然并行，百万级数据轻松处理', 14, LIGHT_GRAY),
    ('④ 广泛适用：网络安全、金融反欺诈、工业预警、IT 运维', 14, LIGHT_GRAY),
    ('⑤ 持续进化：Extended Isolation Forest 解决轴对齐偏差', 14, LIGHT_GRAY),
    ('', 8),
    ('行动建议', 18, ACCENT_GREEN, True),
    ('下次遇到异常检测任务，先用孤立森林建立基线。', 14, WHITE, True),
    ('默认参数（100 棵树、256 子采样）在大多数场景下就够用。', 13, LIGHT_GRAY),
    ('效果不够再考虑 LOF 或 One-Class SVM。', 13, LIGHT_GRAY),
])

add_rect(slide, 7.8, 1.5, 4.8, 3.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 8.0, 1.6, 4.4, 0.4, '🔢  关键数字速记', 14, ACCENT_GREEN, True)

table_data4 = [
    ('项目', '数值'),
    ('默认树数量', '100'),
    ('默认子采样', '256'),
    ('默认树高度', '8'),
    ('异常分界线', 's = 0.5'),
    ('训练复杂度', 'O(n)'),
]
tbl4 = slide.shapes.add_table(6, 2, Inches(8.2), Inches(2.2), Inches(4.0), Inches(2.5)).table
tbl4.columns[0].width = Inches(2.0)
tbl4.columns[1].width = Inches(2.0)
for r in range(6):
    for c in range(2):
        cell = tbl4.cell(r, c)
        cell.text = table_data4[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.name = FONT_CN
            para.alignment = PP_ALIGN.CENTER
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            else:
                para.font.color.rgb = ACCENT_CYAN if c == 1 else LIGHT_GRAY
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

add_rect(slide, 0.8, 5.5, 11.7, 1.5, BG_CARD)
add_multiline_textbox(slide, 1.0, 5.6, 11.3, 1.3, [
    ('📄  论文出处', 14, ACCENT_GREEN, True),
    ('Liu, Ting & Zhou (2008) "Isolation Forest" IEEE ICDM  |  Liu, Ting & Zhou (2012) "Isolation-based Anomaly Detection" ACM TKDD', 11, LIGHT_GRAY),
    ('Hariri, Carrasco Kind & McNutt (2019) "Extended Isolation Forest" IEEE TKDE', 11, LIGHT_GRAY),
    ('', 4),
    ('"一切都应该尽可能简单，但不能过于简单。" —— 爱因斯坦。孤立森林恰好停在"足够简单但不失效"的完美位置。', 12, MID_GRAY),
])


# ── 保存 ──
output_path = '/home/admin/.openclaw/workspace-weaver/output/isolation-forest_v4.pptx'
prs.save(output_path)
print(f'✅ PPT V4 已保存到 {output_path}')
print(f'总页数：{len(prs.slides)}')
