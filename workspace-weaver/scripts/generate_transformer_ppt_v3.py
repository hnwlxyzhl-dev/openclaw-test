#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer v3 PPT Generator - Main Script
Integrates architecture diagrams with text content to produce the final 22-page PPT.
Uses diagrams from transformer-v3-diagrams.pptx and content from the planning doc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ============================================================
# Configuration
# ============================================================
SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_R = 13.2
SAFE_B = 7.3

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xC7)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_DARK = RGBColor(0x34, 0x49, 0x5E)

FONT_ZH = "Microsoft YaHei"
FONT_EN = "Arial"

# ============================================================
# Utility Functions
# ============================================================
def chk(l, t, w, h, label=""):
    assert l + w <= SAFE_R + 0.05, f"Right overflow {label}: {l+w:.2f} > {SAFE_R}"
    assert t + h <= SAFE_B + 0.05, f"Bottom overflow {label}: {t+h:.2f} > {SAFE_B}"
    assert w > 0 and h > 0, f"Zero size {label}: {w}x{h}"

def add_textbox(slide, left, top, width, height, text, font_size=11,
                bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font_name=FONT_ZH,
                anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.0):
    chk(left, top, width, height, f"TB:{text[:20]}")
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tb

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=10,
                          color=C_TEXT, bold_first=False, line_spacing=1.05,
                          font_name=FONT_ZH, align=PP_ALIGN.LEFT, bullet=False):
    """Add a text box with multiple paragraphs."""
    chk(left, top, width, height, f"MTB")
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_first and i == 0:
            p.font.bold = True
        if bullet and not (line.startswith("**") or line.startswith("  ")):
            p.level = 0
    return tb

def add_shape_box(slide, left, top, width, height, text="",
                  fill_color=None, line_color=None, line_width=Pt(1),
                  text_color=C_WHITE, font_size=11, bold=True,
                  font_name=FONT_ZH, align=PP_ALIGN.CENTER,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    chk(left, top, width, height, f"SB:{text[:15]}")
    shape = slide.shapes.add_shape(shape_type,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.font.name = font_name
        p.alignment = align
    return shape

def add_page_title(slide, title, subtitle="", page_num=""):
    """Add a page title bar at the top."""
    # Title background bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  Inches(SLIDE_W), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TITLE
    bar.line.fill.background()
    # Title text
    add_textbox(slide, 0.4, 0.1, 11.5, 0.65, title,
                font_size=22, bold=True, color=C_WHITE, font_name=FONT_ZH)
    if page_num:
        add_textbox(slide, 12.2, 0.1, 0.8, 0.65, page_num,
                    font_size=14, color=C_WHITE, align=PP_ALIGN.RIGHT, font_name=FONT_EN)

def add_footer(slide, text="Transformer Architecture Deep Dive"):
    add_textbox(slide, 0.3, 7.05, 6, 0.3, text,
                font_size=8, color=C_GRAY, font_name=FONT_EN)

# ============================================================
# Content definitions for all 22 pages
# ============================================================
# Each page: (title, layout_type, text_content, diagram_page_index or None)
# layout_type: "cover", "toc", "left_right", "top_bottom", "summary"

PAGES = []  # Will be built below

# ============================================================
# Build the PPT
# ============================================================
def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # Load diagrams PPTX
    diag = Presentation("/home/admin/.openclaw/workspace-weaver/output/transformer-v3-diagrams.pptx")

    print(f"Diagrams PPTX has {len(diag.slides)} slides")
    print(f"Main PPTX: creating {22} slides...")

    # ============================================================
    # Page 1: Cover
    # ============================================================
    slide = prs.slides.add_slide(blank)
    # Background decoration - dark blue gradient bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  Inches(SLIDE_W), Inches(SLIDE_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    bar.line.fill.background()

    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(0),
                                      Inches(SLIDE_W), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_TITLE
    top_bar.line.fill.background()

    # Main title
    add_textbox(slide, 1.5, 1.8, 10, 1.2, "Transformer",
                font_size=48, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, 1.5, 2.8, 10, 0.8, "Architecture Deep Dive",
                font_size=32, bold=True, color=C_ENCODER, align=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, 1.5, 3.8, 10, 0.6, "From Intuition to Principles - A Complete Guide to AI's Core Engine",
                font_size=16, color=C_TEXT, align=PP_ALIGN.CENTER, font_name=FONT_EN, italic=True)

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.5), Inches(4.6),
                                   Inches(4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ENCODER
    line.line.fill.background()

    # Subtitle info
    add_textbox(slide, 1.5, 4.9, 10, 0.5,
                "2017 Google | ChatGPT, GPT-4 and all LLMs are built on this architecture",
                font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, 1.5, 5.6, 10, 0.5,
                "From Training to Inference - Complete Breakdown of How Transformer Works",
                font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER, font_name=FONT_EN)

    print("  Page 1: Cover - Done")

    # ============================================================
    # Page 2: Table of Contents
    # ============================================================
    slide = prs.slides.add_slide(blank)
    add_page_title(slide, "Today's Learning Roadmap", page_num="02/22")

    # Left side: TOC list
    toc_items = [
        ("01", "Why Transformer", "What problems did RNN/LSTM have?", C_ENCODER),
        ("02", "Architecture Overview", "Encoder, Decoder, and each module", C_ENCODER),
        ("03", "Training Phase (Detailed)", "How does the model learn language?", C_ATTENTION),
        ("04", "Inference Phase (Detailed)", "How does the trained model generate?", C_FFN),
        ("05", "Applications & Summary", "ChatGPT, BERT, T5 and beyond", C_EMBED),
    ]
    y = 1.3
    for num, title, desc, color in toc_items:
        # Number circle
        add_shape_box(slide, 0.5, y, 0.55, 0.45, num, fill_color=color,
                      font_size=14, bold=True, font_name=FONT_EN)
        # Title
        add_textbox(slide, 1.2, y, 5, 0.25, title,
                    font_size=14, bold=True, color=C_TEXT, font_name=FONT_ZH)
        # Description
        add_textbox(slide, 1.2, y + 0.25, 5, 0.25, desc,
                    font_size=10, color=C_GRAY, font_name=FONT_ZH)
        y += 0.75

    # Right side: Roadmap visual
    # Vertical line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(8.5), Inches(1.3),
                                   Inches(0.03), Inches(5.2))
    line.fill.solid()
    line.fill.fore_color.rgb = C_TITLE
    line.line.fill.background()

    roadmap_nodes = [
        (8.2, 1.5, "Module 1", "Understanding Problems", C_ENCODER, "Pages 3-4"),
        (8.2, 2.6, "Module 2", "Core Architecture", C_ENCODER, "Pages 5-7"),
        (8.2, 3.7, "Module 3", "Training Phase", C_ATTENTION, "Pages 8-15"),
        (8.2, 4.8, "Module 4", "Inference Phase", C_FFN, "Pages 16-19"),
        (8.2, 5.9, "Module 5", "Applications", C_EMBED, "Pages 20-22"),
    ]
    for rx, ry, label, sub, color, pages in roadmap_nodes:
        add_shape_box(slide, rx, ry, 1.8, 0.6, label, fill_color=color,
                      font_size=11, bold=True)
        add_textbox(slide, rx, ry + 0.6, 1.8, 0.2, sub,
                    font_size=8, color=C_GRAY, align=PP_ALIGN.CENTER, font_name=FONT_EN)
        add_textbox(slide, rx, ry + 0.78, 1.8, 0.2, pages,
                    font_size=8, color=C_GRAY, align=PP_ALIGN.CENTER, font_name=FONT_EN)

    # Arrow at bottom
    add_textbox(slide, 7.8, 6.5, 3.5, 0.3, "Goal: Fully understand Transformer",
                font_size=10, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_footer(slide)
    print("  Page 2: TOC - Done")

    # ============================================================
    # Pages 3-21: Copy diagrams + add text
    # ============================================================
    # Diagram slide index mapping:
    # diagrams[0] -> Slide 3 (RNN)
    # diagrams[1] -> Slide 4 (Attention intuition)
    # diagrams[2] -> Slide 5 (Overall architecture)
    # diagrams[3] -> Slide 6 (Encoder single layer)
    # diagrams[4] -> Slide 7 (Decoder single layer)
    # diagrams[5] -> Slide 8 (Training overview) -> actually page 9 (embedding)
    # diagrams[6] -> Slide 9 (Embedding + Position)
    # diagrams[7] -> Slide 10 (QKV generation)
    # diagrams[8] -> Slide 11 (Attention 5 steps)
    # diagrams[9] -> Slide 12 (Residual + FFN)
    # diagrams[10] -> Slide 13 (Masked + Cross attention)
    # diagrams[11] -> Slide 14 (Loss function)
    # diagrams[12] -> Slide 15 (Teacher Forcing)
    # diagrams[13] -> Slide 17 (Autoregressive + KV Cache)

    # Page content definitions (text to add on the right or bottom)
    page_contents = {
        3: {
            "title": "Why Transformer: Three Pain Points of RNN/LSTM",
            "page_num": "03/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Sequential Processing", "Like reading one word at a time, RNN must process text word by word. Modern GPUs have thousands of cores, but RNN can only use one. Transformer trained in 12 hours on 8 GPUs vs days for LSTM."),
                ("Long-range Forgetting", "Vanishing gradient problem - error signals gradually decay during backpropagation. When sentences exceed 50 words, standard RNN performance drops significantly."),
                ("Information Bottleneck", "RNN compresses all information into a fixed-length vector. Like summarizing a 500-page book into 100 words - massive information loss."),
            ]
        },
        4: {
            "title": "Core Breakthrough: Every Word Can Communicate Directly",
            "page_num": "04/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Round Table vs Queue", "If RNN is like 'queue to speak' (must wait for the previous person), Transformer is like a 'round table meeting' - everyone hears everyone simultaneously. This is Self-Attention."),
                ("Example: 'it' finds 'cat'", "In 'The cat sat on the mat, and it was happy', when processing 'it', the model checks all words simultaneously. 'it' has highest attention to 'cat' (grammatical relationship)."),
                ("One Solution, Three Problems", "Simultaneous attention solves all three RNN problems: parallel processing, no information decay, no compression bottleneck."),
            ]
        },
        5: {
            "title": "Architecture Overview: Left Understands, Right Generates",
            "page_num": "05/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Encoder - 'Editorial Team'", "Processes input sequence through 6 layers, outputs semantically rich vectors. Like a reading comprehension expert. 512-dim, ~65M parameters."),
                ("Decoder - 'Translation Dept'", "Receives encoder output, generates target sequence word by word. Like a translator. Also 6 layers, but with an extra Cross-Attention module."),
                ("Complete Data Flow", "Input -> Tokenization -> Embedding -> Positional Encoding -> Encoder -> Cross-Attention -> Decoder -> Linear -> Softmax -> Output probabilities."),
            ]
        },
        6: {
            "title": "Encoder Detail: Each Layer 'Communicates' and 'Refines'",
            "page_num": "06/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Multi-Head Self-Attention", "8 heads analyze word relationships from different angles (grammar, semantics, coreference). Same total computation as a single head, but more comprehensive perspective."),
                ("Residual Connection + Layer Norm", "Residual connection = highway bypass (information can skip layers). Layer normalization = 'standardized health check' keeping all vectors in similar ranges."),
                ("Feed-Forward Network", "Independent 'thinking time' for each word. Expands 512 -> 2048 -> ReLU -> 512. Actually ~2x more computation than attention layer, but essential for expressiveness."),
            ]
        },
        7: {
            "title": "Decoder Detail: One Extra 'Translation Bridge'",
            "page_num": "07/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Masked Self-Attention", "Same as encoder's self-attention but with causal mask - each word can only see previous words (including itself). Like an exam where you can't look ahead."),
                ("Cross-Attention", "The only connection between encoder and decoder. Q from decoder, K and V from encoder. Like a translator consulting editorial notes."),
                ("Feed-Forward Network", "Identical to encoder (512->2048->ReLU->512). Full data flow: Masked Attn -> Norm -> Cross Attn -> Norm -> FFN -> Norm -> Output."),
            ]
        },
        8: {
            "title": "Training Phase: How Does the Model 'Learn' Language?",
            "page_num": "08/22",
            "layout": "top_bottom",
            "text_x": 0.5, "text_y": 4.2, "text_w": 12, "text_h": 2.8,
            "lines": [
                ("Training = 'Open-book Exam'", "The model sees correct answers for all previous positions when predicting each word. This technique is called Teacher Forcing - enables efficient learning from massive data."),
                ("Training Objective", "Encoder processes 'I love you' and outputs encoding vectors. Decoder input is [START, I, love] (right-shifted), decoder predicts [I, love, you, END] (next word at each position)."),
                ("Division of Labor", "Encoder processes input once, output is cached. Decoder receives full target sequence but uses causal mask internally. Encoder = 'understand source', Decoder = 'learn to translate'."),
            ]
        },
        9: {
            "title": "Training - Encoder Step 1: Turning Words into 'Digital Portraits'",
            "page_num": "09/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Tokenization", "Split text into minimum units. Original Transformer uses BPE. 'I love you' -> ['I', 'love', 'you']. Vocabulary size: 37,000 tokens."),
                ("Embedding Matrix", "A 'super dictionary' of 37,000 x 512 = ~18.9M parameters (~30% of model). Words with similar meanings (e.g., 'happy' and 'joyful') are close in vector space."),
                ("Positional Encoding", "Self-attention is permutation-invariant - but 'dog bites man' != 'man bites dog'. Sinusoidal functions generate 512-dim position vectors, added directly to embeddings. Like adding an address to each person's ID card."),
            ]
        },
        10: {
            "title": "Training - Encoder Step 2: Preparing 'Search Tools' for Each Word",
            "page_num": "10/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Query - 'What am I looking for'", "Each word's 512-dim vector is multiplied by W^Q (512x64) to produce a 64-dim Query. For 'it', the Query might encode 'I need to find a noun to tell me what it refers to'."),
                ("Key - 'What features do I have'", "Each word multiplied by W^K (512x64) produces a 64-dim Key. 'cat's Key might encode 'noun, animal, singular, can be subject'. Different matrices capture different information."),
                ("Value - 'What is my actual content'", "Each word multiplied by W^V (512x64) produces a 64-dim Value. When 'it' matches 'cat's Key, it reads 'cat's Value. 8 heads = 8 independent sets of QKV, total = 512-dim."),
            ]
        },
        11: {
            "title": "Training - Encoder Step 3: 5-Step Attention Computation",
            "page_num": "11/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Step 1: Generate Q, K, V", "All words simultaneously produce their Q, K, V vectors through three weight matrices. Fully parallel."),
                ("Step 2: Q x K^T = Score Matrix", "Dot product of each Query with all Keys produces a T x T score matrix. Higher score = stronger association."),
                ("Step 3: Divide by sqrt(d_k)", "Divide by sqrt(64) = 8. Prevents Softmax 'saturation' where highest score dominates. Mathematically normalizes variance to 1."),
                ("Step 4: Softmax = Attention Weights", "Converts each row to probabilities summing to 1. Best match gets most attention, but related words also get some."),
                ("Step 5: Weights x V = Output", "Weighted sum of all Values. For 'it', 'cat' has highest weight, so 'it's output contains 'cat's semantics. Each word's output fuses context from entire sentence."),
            ]
        },
        12: {
            "title": "Training - Encoder Step 4: Safety Net + Independent Thinking",
            "page_num": "12/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Residual Connection", "Output = Input + Attention(Input). Like a highway bypass - if the main road has problems, information flows through the bypass. Solves vanishing gradients in deep networks."),
                ("Layer Normalization", "Standardizes each word's vector: subtract mean, divide by std, multiply by learnable scale. Like a 'health check' ensuring all vectors are in similar ranges. Stabilizes training."),
                ("Feed-Forward Network", "512 -> 2048 -> ReLU -> 512. Expands dimension for complex feature extraction, then compresses back. ~2x more computation than attention, but essential. Without FFN, model expressiveness drops significantly."),
                ("Complete Layer Data Flow", "Input -> Multi-Head Attention -> Residual + Norm -> FFN -> Residual + Norm -> Output. Repeated 6 times. Each round of 'communicate + refine' enriches word representations."),
            ]
        },
        13: {
            "title": "Training - Decoder: Masked Attention + Cross Attention",
            "page_num": "13/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Masked Self-Attention", "Same computation as encoder self-attention, but with causal mask: scores above the diagonal are set to -infinity. After Softmax, these become 0 - invisible. Processing 'love' can only see [START, I, love]."),
                ("Why Mask During Training?", "Without masking, the model would learn to 'cheat' - copy future answers instead of learning to predict. Causal mask ensures training and inference behavior is consistent."),
                ("Cross-Attention", "Q from decoder's masked attention output ('what information do I need?'), K and V from encoder's final output ('complete understanding of source'). Same computation as self-attention, just different Q/KV sources."),
            ]
        },
        14: {
            "title": "Training - Loss Function: Measuring How Well the Model Learns",
            "page_num": "14/22",
            "layout": "top_bottom",
            "text_x": 0.5, "text_y": 3.5, "text_w": 12, "text_h": 3.5,
            "lines": [
                ("Output Layer: Vector -> Vocabulary Probabilities", "Linear layer (512x37,000) maps 512-dim to vocabulary size, then Softmax produces probability distribution. In many implementations, this weight matrix is shared with the embedding matrix (weight tying)."),
                ("Cross-Entropy Loss: 'Multiple Choice Score'", "Measures distance between predicted distribution and true answer. Intuition: if correct answer is B and your confidence is 70%, loss = -log(0.7) = 0.36. Higher confidence = lower loss."),
                ("Loss Calculation", "Cross-entropy = -log(p(correct word)), averaged over all positions. Early training: random predictions (~1/37000), loss ~ 10.5. After training: p(correct) = 0.5-0.9, loss drops to 1-3."),
                ("Backpropagation", "Compute gradient for each parameter (how much each parameter contributes to loss), then update parameters using Adam optimizer to reduce loss. Repeated over entire dataset."),
            ]
        },
        15: {
            "title": "Training - Teacher Forcing: Efficiency and Hidden Risk",
            "page_num": "15/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("How It Works", "Decoder input is the TRUE target sequence, not model predictions. When predicting word 2, decoder sees [START, I] (true 'I', not model's possibly wrong prediction). Like open-book exam."),
                ("Advantage: Extremely Efficient", "Every position gets immediate correct feedback. Stable gradient signals. All position losses computed in parallel, fully utilizing GPU."),
                ("Disadvantage: Exposure Bias", "Training sees correct inputs, but inference sees model's own (possibly wrong) outputs. This distribution mismatch can degrade inference performance. Like a student who always practices with answer keys, then panics in a real exam."),
                ("Mitigation: Scheduled Sampling", "Randomly use model's own predictions as input with some probability during training, narrowing the gap between training and inference."),
            ]
        },
        16: {
            "title": "Inference Phase: How Does the Model Generate Answers Independently?",
            "page_num": "16/22",
            "layout": "top_bottom",
            "text_x": 0.5, "text_y": 4.2, "text_w": 12, "text_h": 2.8,
            "lines": [
                ("Inference = 'Closed-book Exam'", "No correct answers available. Model must rely entirely on its own previous predictions. If the first word is wrong, subsequent words may all be wrong, and there's no going back."),
                ("Encoder: Works Only Once", "Identical to training. Input sequence is processed through 6 encoder layers once, output is cached and reused throughout generation."),
                ("Decoder: Word-by-Word Serial Generation", "Unlike training where all positions are computed in parallel, inference must generate one word at a time: generate word 1, add to input, generate word 2, and so on. This serial process is the main inference speed bottleneck."),
            ]
        },
        17: {
            "title": "Inference - Autoregressive Generation + KV Cache",
            "page_num": "17/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Autoregressive Generation", "Each new word is added to the sequence before predicting the next. Like writing an essay - must finish one word before starting the next. Cannot go back to revise."),
                ("Step-by-Step Example", "Step 1: Input [START] -> output 'I'. Step 2: Input [START, I] -> output 'love'. Step 3: Input [START, I, love] -> output 'you'. Step 4: Output <END>, stop. Final: 'I love you'."),
                ("KV Cache: 'Meeting Minutes'", "Instead of recomputing all historical K, V at each step, cache previously computed K, V vectors. Each step only computes the new word's K, V and appends to cache. Like reviewing meeting minutes instead of re-discussing everything."),
            ]
        },
        18: {
            "title": "Inference - Word Selection Strategies",
            "page_num": "18/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("Greedy Search", "Always pick the highest probability word. Fast but prone to local optima - a globally suboptimal sequence."),
                ("Beam Search", "Keep top-k candidate sequences simultaneously (beam size 4-10). Good for translation, but outputs tend to be 'safe and mediocre'."),
                ("Top-K Sampling", "Random sample from top-k words (e.g., k=50). Excludes low-quality candidates while maintaining randomness."),
                ("Top-P (Nucleus) Sampling", "Sample from minimum set of words whose cumulative probability reaches p (e.g., 0.9). More adaptive than Top-K. ChatGPT typically uses Temperature=0.7 + Top-P=0.9."),
                ("Temperature", "Scale logits before Softmax. Low temp (0.3) = conservative output. High temp (1.5) = diverse output. Acts as a 'creativity valve'."),
            ]
        },
        19: {
            "title": "Training vs Inference: Open-book vs Closed-book Exam",
            "page_num": "19/22",
            "layout": "top_bottom",
            "text_x": 0.5, "text_y": 4.2, "text_w": 12, "text_h": 2.8,
            "lines": [
                ("Decoder Input: Correct Answers vs Own Output", "Training sees full target sequence (Teacher Forcing). Inference sees model's own previous predictions. This is the fundamental difference."),
                ("Computation: Parallel vs Serial", "Training computes all positions in parallel. Inference generates word-by-word serially. Training efficiency >> inference efficiency."),
                ("Loss & Gradients: Yes vs No", "Training needs cross-entropy loss and backpropagation. Inference has no loss, no gradients - parameters are fixed, forward pass only."),
                ("KV Cache: Not Needed vs Required", "Training doesn't need KV Cache (parallel computation). Inference must use KV Cache to avoid recomputing historical K, V. LLaMA-2 70B at 4096 tokens: KV Cache ~ 5.4GB."),
            ]
        },
        20: {
            "title": "Three Families: Same DNA, Three Career Paths",
            "page_num": "20/22",
            "layout": "left_right",
            "text_x": 7.5, "text_w": 5.2,
            "lines": [
                ("GPT (Decoder-Only) - 'Improvisational Speaker'", "Core ability: text generation. Given a prompt, continues word by word. ChatGPT: pre-training (predict next word) -> instruction fine-tuning -> RLHF. Models: GPT-3 (175B), GPT-4 (~1.8T), LLaMA, Qwen."),
                ("BERT (Encoder-Only) - 'Reading Comprehension Expert'", "Core ability: text understanding. Uses Masked Language Modeling (randomly mask 15% of tokens). Excels at classification, search, sentiment analysis. Google Search uses BERT since 2019. Models: BERT-base (110M), BERT-large (340M)."),
                ("T5 (Encoder-Decoder) - 'All-purpose Translator'", "Full encoder-decoder architecture. Unifies all NLP tasks as text-to-text. Google Translate uses T5 variants. Models: T5-Small (60M) to T5-11B (11B)."),
            ]
        },
        21: {
            "title": "From Translation Tool to AI's Universal Engine",
            "page_num": "21/22",
            "layout": "top_bottom",
            "text_x": 0.5, "text_y": 3.5, "text_w": 12, "text_h": 3.5,
            "lines": [
                ("Scale Explosion: 65M to Trillions", "2017: 65M parameters. Today: GPT-4 ~1.8T (30,000x growth). Same core architecture, just more layers (6->80), larger dimensions (512->8192), more data. Scaling Laws: more params + more data = stronger ability."),
                ("Beyond Language: Everything Can Be 'Transformed'", "ViT (2020): cut images into patches as 'words'. Whisper (2022): speech recognition. GPT-4V, Gemini: multimodal understanding. AlphaFold2: protein structure prediction."),
                ("Core Insight: Universal Information Processing Framework", "Transformer provides a universal framework - as long as you can represent data as a 'sequence', you can process it with Transformer. Text = sequence of words, Image = sequence of patches, Audio = sequence of samples. This is why it's called 'AI's universal engine'."),
            ]
        },
        22: {
            "title": "Summary: 5 Key Takeaways",
            "page_num": "22/22",
            "layout": "summary",
            "text_x": 0.5, "text_w": 12, "text_h": 5.5,
            "lines": [
                ("1. Transformer's Essence: Free Information Flow", "Core idea: every word can communicate directly with all other words. Solved RNN's three problems simultaneously. Paradigm shift from serial processing to parallel understanding."),
                ("2. Encoder = Understand, Decoder = Generate", "Encoder: self-attention + FFN fuses sentence semantics into each word. Decoder: masked attention + cross-attention + FFN generates target word by word."),
                ("3. Training = Open-book Exam, Inference = Closed-book Exam", "Training: Teacher Forcing provides correct answers, parallel loss computation. Inference: autoregressive word-by-word generation with KV Cache optimization. Key risk: exposure bias."),
                ("4. Three Families, Each with Strengths", "GPT (Decoder-Only) excels at generation -> ChatGPT. BERT (Encoder-Only) excels at understanding -> search. T5 (Encoder-Decoder) excels at transformation -> translation."),
                ("5. From Paper to Universal Engine", "From 65M parameter translation model (2017) to trillion-parameter multimodal systems today. Transformer is the foundational architecture of the entire AI era."),
            ]
        },
    }

    # ============================================================
    # Generate Pages 3-22 by copying diagrams + adding text
    # ============================================================

    # Mapping: target page number -> diagram slide index
    # Diagram slides: 0-13
    # Pages without diagrams: 8 (training overview), 16 (inference overview), 19 (comparison), 21 (impact), 22 (summary)
    diagram_map = {
        3: 0,   # RNN
        4: 1,   # Attention intuition
        5: 2,   # Overall architecture
        6: 3,   # Encoder single layer
        7: 4,   # Decoder single layer
        # Page 8: no diagram (training overview) - create from scratch
        9: 5,   # Embedding + Position
        10: 6,  # QKV generation
        11: 7,  # Attention 5 steps
        12: 8,  # Residual + FFN
        13: 9,  # Masked + Cross attention
        14: 10, # Loss function
        15: 11, # Teacher Forcing
        # Page 16: no diagram (inference overview)
        17: 12, # Autoregressive + KV Cache
        18: 13, # Word selection strategies
        # Page 19: no diagram (comparison)
        # Page 20: no diagram (three families)
        # Page 21: no diagram (impact)
        # Page 22: no diagram (summary)
    }

    for page_num in range(3, 23):
        content = page_contents[page_num]
        slide = prs.slides.add_slide(blank)

        # Add page title
        add_page_title(slide, content["title"], page_num=content["page_num"])

        # Check if this page has a diagram
        if page_num in diagram_map:
            diag_idx = diagram_map[page_num]
            diag_slide = diag.slides[diag_idx]
            # Copy all shapes from diagram slide
            for shape in diag_slide.shapes:
                el = shape.element
                new_el = copy.deepcopy(el)
                slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            print(f"  Page {page_num}: Copied diagram[{diag_idx}] + title")

        # Add text content
        if content["layout"] == "left_right":
            tx = content["text_x"]
            tw = content["text_w"]
            ty = 1.1
            n_items = len(content["lines"])
            avail_h = SAFE_B - ty - 0.3
            heading_h = 0.28
            item_gap = 0.08
            total_heading = n_items * (heading_h + item_gap)
            body_budget = avail_h - total_heading
            for i, (heading, body) in enumerate(content["lines"]):
                # Heading
                add_textbox(slide, tx, ty, tw, heading_h, heading,
                            font_size=11, bold=True, color=C_TITLE, font_name=FONT_EN)
                ty += heading_h + 0.02
                # Body - distribute remaining space
                body_h = max(0.4, body_budget / n_items - 0.05)
                body_h = min(body_h, 1.2)
                body_chars = len(body)
                if body_chars > 200:
                    fs = 8.5
                elif body_chars > 150:
                    fs = 9
                else:
                    fs = 9.5
                add_textbox(slide, tx, ty, tw, body_h, body,
                            font_size=fs, color=C_TEXT, font_name=FONT_EN,
                            line_spacing=0.95)
                ty += body_h + item_gap

        elif content["layout"] == "top_bottom":
            tx = content.get("text_x", 0.5)
            ty_start = content.get("text_y", 4.2)
            tw = content.get("text_w", 12)
            ty = ty_start
            n_lines = len(content["lines"])
            # Calculate available space and distribute
            avail_h = SAFE_B - ty_start - 0.2
            heading_h = 0.22
            gap = 0.03
            body_h_each = (avail_h - n_lines * (heading_h + gap)) / n_lines
            body_h_each = max(0.3, min(0.65, body_h_each))
            for i, (heading, body) in enumerate(content["lines"]):
                add_textbox(slide, tx, ty, tw, heading_h, heading,
                            font_size=10, bold=True, color=C_TITLE, font_name=FONT_EN)
                ty += heading_h + gap
                fs = 9
                add_textbox(slide, tx, ty, tw, body_h_each, body,
                            font_size=fs, color=C_TEXT, font_name=FONT_EN,
                            line_spacing=0.95)
                ty += body_h_each

        elif content["layout"] == "summary":
            tx = content.get("text_x", 0.5)
            ty = 1.1
            tw = content.get("text_w", 12)
            for i, (heading, body) in enumerate(content["lines"]):
                # Numbered card
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               Inches(tx), Inches(ty),
                                               Inches(tw), Inches(0.95))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA) if i % 2 == 0 else RGBColor(0xEB, 0xF5, 0xFB)
                card.line.fill.background()
                add_textbox(slide, tx + 0.2, ty + 0.05, tw - 0.4, 0.3, heading,
                            font_size=12, bold=True, color=C_TITLE, font_name=FONT_EN)
                add_textbox(slide, tx + 0.2, ty + 0.35, tw - 0.4, 0.55, body,
                            font_size=10, color=C_TEXT, font_name=FONT_EN,
                            line_spacing=1.0)
                ty += 1.05

            # Thank you footer
            add_textbox(slide, 1.5, 6.6, 10, 0.4,
                        "Thank you! Questions and discussion welcome.",
                        font_size=13, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
            add_textbox(slide, 1.5, 6.95, 10, 0.3,
                        "Transformer is not just a model - it's an era.",
                        font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

        add_footer(slide)

    # ============================================================
    # Save
    # ============================================================
    output_path = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3.pptx"
    prs.save(output_path)
    print(f"\nPPT saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Boundary check
    print("\n=== Boundary Check ===")
    errors = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            l = shape.left / 914400
            t = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            if l + w > SAFE_R + 0.1:
                print(f"  WARN Page {i+1}: right overflow {l+w:.2f} (shape: {shape.name})")
                errors += 1
            if t + h > SAFE_B + 0.1:
                print(f"  WARN Page {i+1}: bottom overflow {t+h:.2f} (shape: {shape.name})")
                errors += 1
            if h > 7:
                print(f"  WARN Page {i+1}: abnormal height {h:.2f} (shape: {shape.name})")
                errors += 1
    if errors == 0:
        print("  All shapes within safe boundaries!")
    else:
        print(f"  {errors} warnings found")

if __name__ == "__main__":
    build_ppt()
