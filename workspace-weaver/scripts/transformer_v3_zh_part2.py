#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT slides 12-22 - Chinese version"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# === Constants ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_W = Inches(13.2)
SAFE_H = Inches(7.3)

C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xD0)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_BG_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
C_BG_GREEN = RGBColor(0xEA, 0xFA, 0xEA)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

LQ = '\u300c'  # left corner bracket 「
RQ = '\u300d'  # right corner bracket 」

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def set_shape(shape, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=10,
                font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT_CN
    p.alignment = alignment
    for run in p.runs:
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_para(tf, text, font_size=10, font_color=C_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, space_after=Pt(4), space_before=Pt(0)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT_CN
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    for run in p.runs:
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
    return p


def add_rrect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    return set_shape(shape, left, top, width, height, fill_color, line_color, line_width)


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    return set_shape(shape, left, top, width, height, fill_color, line_color, line_width)


def add_circle(slide, left, top, size, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    return set_shape(shape, left, top, size, size, fill_color, line_color)


def add_line(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_arrow_line(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)
    return connector


def shape_text(shape, text, font_size=9, font_color=C_TEXT, bold=False, alignment=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT_CN
    for run in p.runs:
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
    txBody = shape._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    return tf


def add_title_bar(slide, title_text):
    bar = add_rect(slide, Inches(0), Inches(0), SAFE_W, Inches(0.85), C_TITLE)
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
    txBody = bar._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    bodyPr.set('lIns', str(Emu(Inches(0.3))))
    return bar


def add_page_num(slide, num):
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
                str(num), font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# Slide 12
# ============================================================
def slide_12():
    s = add_slide()
    add_title_bar(s, '训练\u00b7编码器第四步：安全网+独立思考')
    add_page_num(s, 12)

    lx = Inches(0.3)
    ly = Inches(1.1)
    bw = Inches(1.4)
    bh = Inches(0.55)
    gap = Inches(0.15)

    nodes = [
        ('输入\n(512维)', C_LIGHT_BLUE),
        ('多头注意力\n(8头x64维)', C_ATTENTION),
        ('残差+归一化\n(512维)', C_FFN),
        ('前馈网络\n(512->2048->512)', C_ENCODER),
        ('残差+归一化\n(512维)', C_FFN),
        ('输出\n(512维)', C_LIGHT_PURPLE),
    ]

    cy = ly + Inches(0.3)
    node_shapes = []
    for i, (txt, col) in enumerate(nodes):
        x = lx + i * (bw + gap)
        sh = add_rrect(s, x, cy, bw, bh, col, C_TEXT, Pt(0.5))
        shape_text(sh, txt, font_size=8, font_color=C_TEXT, bold=True)
        node_shapes.append(sh)

    for i in range(len(node_shapes) - 1):
        x1 = node_shapes[i].left + node_shapes[i].width
        y1 = node_shapes[i].top + node_shapes[i].height // 2
        x2 = node_shapes[i + 1].left
        add_arrow_line(s, x1, y1, x2, y1, C_TEXT, Pt(1.2))

    # Residual connections (bypass)
    bypass_y = cy + bh + Inches(0.25)
    x_start = node_shapes[0].left + node_shapes[0].width // 2
    x_end1 = node_shapes[2].left + node_shapes[2].width // 2
    add_line(s, x_start, cy + bh, x_start, bypass_y, C_ENCODER, Pt(1.5))
    add_line(s, x_start, bypass_y, x_end1, bypass_y, C_ENCODER, Pt(1.5))
    add_line(s, x_end1, bypass_y, x_end1, node_shapes[2].top + node_shapes[2].height, C_ENCODER, Pt(1.5))

    x_mid = node_shapes[2].left + node_shapes[2].width // 2
    x_end2 = node_shapes[4].left + node_shapes[4].width // 2
    bypass_y2 = bypass_y + Inches(0.25)
    add_line(s, x_mid, node_shapes[2].top + node_shapes[2].height, x_mid, bypass_y2, C_ENCODER, Pt(1.5))
    add_line(s, x_mid, bypass_y2, x_end2, bypass_y2, C_ENCODER, Pt(1.5))
    add_line(s, x_end2, bypass_y2, x_end2, node_shapes[4].top + node_shapes[4].height, C_ENCODER, Pt(1.5))

    add_textbox(s, lx + Inches(0.5), bypass_y + Inches(0.55), Inches(3), Inches(0.3),
                '旁路 = 高速公路匝道 (残差连接)', font_size=8, font_color=C_ENCODER, bold=True)

    rx = Inches(7.2)
    ry = Inches(1.1)
    _, tf = add_rich_textbox(s, rx, ry, Inches(5.8), Inches(6.0))
    t1 = '\u2460 残差连接：输出=输入+注意力(输入)。像高速公路匝道，主干道堵了信息可以从匝道通过。解决梯度消失，让梯度回传到任意一层'
    t2 = '\u2461 层归一化：对每个词向量做标准化（减均值、除标准差）。像' + LQ + '标准化体检' + RQ + '，确保数值范围合理，训练更稳定'
    t3 = '\u2462 前馈网络：512->2048（扩大思考空间）->ReLU（保留正值变零）->512（浓缩结果）。计算量比注意力大2倍'
    t4 = '\u2463 一层完整数据流：输入->注意力->安全网->FFN->安全网->输出。6层重复，每轮' + LQ + '交流+提炼' + RQ + '让表示更丰富'
    add_para(tf, t1, font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, t2, font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, t3, font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, t4, font_size=10, font_color=C_TEXT, space_after=Pt(10))


# ============================================================
# Slide 13
# ============================================================
def slide_13():
    s = add_slide()
    add_title_bar(s, '训练\u00b7解码器：掩码自注意力+交叉注意力')
    add_page_num(s, 13)

    mx = Inches(0.5)
    my = Inches(1.3)
    cell = Inches(0.85)
    labels = ['START', '\u6211', '\u7231']

    for j, lb in enumerate(labels):
        add_textbox(s, mx + j * cell + Inches(0.05), my - Inches(0.3), cell, Inches(0.25),
                    lb, font_size=9, font_color=C_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    for i, lb in enumerate(labels):
        add_textbox(s, mx - Inches(0.6), my + i * cell + Inches(0.15), Inches(0.55), Inches(0.3),
                    lb, font_size=9, font_color=C_TEXT, bold=True, alignment=PP_ALIGN.RIGHT)
        for j in range(3):
            x = mx + j * cell
            y = my + i * cell
            if j <= i:
                sh = add_rect(s, x, y, cell - Inches(0.02), cell - Inches(0.02), C_LIGHT_BLUE, C_ENCODER, Pt(0.5))
                shape_text(sh, '1.0', font_size=10, font_color=C_ENCODER, bold=True)
            else:
                sh = add_rect(s, x, y, cell - Inches(0.02), cell - Inches(0.02), C_LIGHT_RED, C_DECODER, Pt(0.5))
                shape_text(sh, '-inf', font_size=10, font_color=C_DECODER, bold=True)

    add_textbox(s, mx, my + 3 * cell + Inches(0.1), Inches(3), Inches(0.3),
                'Softmax后->0（看不见）', font_size=9, font_color=C_DECODER, bold=True,
                alignment=PP_ALIGN.CENTER)

    cx = Inches(4.5)
    cy = Inches(1.5)
    enc = add_rrect(s, cx, cy, Inches(2.2), Inches(0.8), C_LIGHT_PURPLE, C_EMBED, Pt(1.5))
    shape_text(enc, '编码器输出 (K, V)\n编辑对原文的理解', font_size=9, font_color=C_EMBED, bold=True)

    dec = add_rrect(s, cx, cy + Inches(1.5), Inches(2.2), Inches(0.7), C_LIGHT_BLUE, C_ENCODER, Pt(1.5))
    shape_text(dec, '解码器 Q\n我在翻译，需要找原文信息', font_size=9, font_color=C_ENCODER, bold=True)

    add_arrow_line(s, cx + Inches(1.1), cy + Inches(2.2), cx + Inches(1.1), cy + Inches(3.0), C_ENCODER, Pt(1.5))
    res = add_rrect(s, cx, cy + Inches(3.0), Inches(2.2), Inches(0.6), C_LIGHT_GREEN, C_FFN, Pt(1.5))
    shape_text(res, "辅助生成" + LQ + "爱" + RQ, font_size=10, font_color=C_FFN, bold=True)

    add_textbox(s, cx + Inches(2.3), cy + Inches(1.8), Inches(1.5), Inches(0.5),
                'Q x K^T -> 缩放\n-> Softmax -> x V', font_size=8, font_color=C_TEXT)

    rx = Inches(7.5)
    _, tf = add_rich_textbox(s, rx, Inches(1.1), Inches(5.5), Inches(6.0))
    t1 = '\u2460 掩码自注意力：输入[START,我,爱]。因果掩码将未来位置分数设为-inf，Softmax后变0。处理' + LQ + '爱' + RQ + '时只能看到[START,我,爱]'
    t2 = '\u2461 为什么训练也需要掩码？不让模型' + LQ + '作弊' + RQ + '抄后面的答案，否则推理时没有答案可抄就会崩溃。确保训练和推理行为一致'
    t3 = '\u2462 交叉注意力：Q来自解码器（' + LQ + '我在翻译，需要找原文信息' + RQ + '），K和V来自编码器（' + LQ + '编辑对原文的理解' + RQ + '）。Q x K^T->缩放->Softmax->xV。如生成' + LQ + '爱' + RQ + '时，Query' + LQ + '找感情动词' + RQ + '与' + LQ + 'love' + RQ + '的Key匹配度最高'
    add_para(tf, t1, font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, t2, font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, t3, font_size=10, font_color=C_TEXT, space_after=Pt(10))


# ============================================================
# Slide 14
# ============================================================
def slide_14():
    s = add_slide()
    add_title_bar(s, '训练\u00b7损失函数：怎么衡量模型学得好不好？')
    add_page_num(s, 14)

    ly = Inches(1.2)
    bw = Inches(1.8)
    bh = Inches(0.7)
    gap = Inches(0.3)
    lx = Inches(0.5)

    flow_nodes = [
        ('解码器输出\n(512维)', C_ENCODER),
        ('线性变换\n(512->37000)', C_FFN),
        ('Softmax', C_ATTENTION),
        ('预测概率\n分布', RGBColor(0xF3, 0x9C, 0x12)),
    ]

    shapes = []
    for i, (txt, col) in enumerate(flow_nodes):
        x = lx + i * (bw + gap)
        sh = add_rrect(s, x, ly, bw, bh, col, C_TEXT, Pt(0.8))
        fc = C_WHITE if col in [C_ENCODER, C_ATTENTION] else C_TEXT
        shape_text(sh, txt, font_size=9, font_color=fc, bold=True)
        shapes.append(sh)

    for i in range(len(shapes) - 1):
        x1 = shapes[i].left + shapes[i].width
        y1 = shapes[i].top + shapes[i].height // 2
        x2 = shapes[i + 1].left
        add_arrow_line(s, x1, y1, x2, y1, C_TEXT, Pt(1.2))

    truth_y = ly + Inches(1.2)
    truth = add_rrect(s, lx + Inches(2), truth_y, Inches(1.5), Inches(0.6), C_DECODER, C_DECODER, Pt(1.5))
    shape_text(truth, '真实答案', font_size=10, font_color=C_WHITE, bold=True)

    pred = add_rrect(s, lx + Inches(4.5), truth_y, Inches(1.5), Inches(0.6), C_ATTENTION, C_ATTENTION, Pt(1.5))
    shape_text(pred, '预测概率', font_size=10, font_color=C_WHITE, bold=True)

    add_line(s, truth.left + truth.width, truth.top + truth.height // 2,
             pred.left, pred.top + pred.height // 2, C_TEXT, Pt(1.5))
    add_textbox(s, lx + Inches(2.5), truth_y - Inches(0.25), Inches(2), Inches(0.25),
                '<-> 对比', font_size=9, font_color=C_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    loss_y = truth_y + Inches(1.0)
    loss = add_rrect(s, lx + Inches(2.5), loss_y, Inches(2), Inches(0.6), C_DECODER, C_DECODER, Pt(1.5))
    shape_text(loss, '交叉熵损失', font_size=10, font_color=C_WHITE, bold=True)
    add_arrow_line(s, lx + Inches(3.5), truth.top + truth.height, lx + Inches(3.5), loss.top, C_DECODER, Pt(1.2))

    bp = add_rrect(s, lx + Inches(5.5), loss_y, Inches(1.5), Inches(0.6), C_GRAY, C_TEXT, Pt(0.8))
    shape_text(bp, '更新参数', font_size=9, font_color=C_TEXT, bold=True)
    add_arrow_line(s, loss.left + loss.width, loss.top + loss.height // 2,
                   bp.left, bp.top + bp.height // 2, C_GRAY, Pt(1.2))

    add_textbox(s, lx + Inches(1), loss_y + Inches(0.7), Inches(5), Inches(0.3),
                'Loss = -log(p(正确词))，所有位置取平均', font_size=9, font_color=C_DECODER, bold=True,
                alignment=PP_ALIGN.CENTER)

    rx = Inches(0.5)
    ry = Inches(4.8)
    _, tf = add_rich_textbox(s, rx, ry, Inches(12.2), Inches(2.5))
    add_para(tf, '\u2460 输出层：线性变换(512x37000权重矩阵)+Softmax->词表概率分布。概率最高的词就是模型认为最应该出现的词',
             font_size=10, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, '\u2461 交叉熵损失=选择题得分。正确答案B，信心70%->损失-log(0.7)=0.36。信心越高损失越低',
             font_size=10, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, '\u2462 损失计算：每个位置算-log(p(正确词))，所有位置取平均。训练初期约10.5（随机），训练充分后降到1-3',
             font_size=10, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, '\u2463 反向传播：根据损失计算每个参数梯度，用Adam优化器按减小损失方向更新参数。反复训练直到损失不再下降',
             font_size=10, font_color=C_TEXT, space_after=Pt(6))


# ============================================================
# Slide 15
# ============================================================
def slide_15():
    s = add_slide()
    add_title_bar(s, '训练\u00b7Teacher Forcing：开卷考试的效率与隐患')
    add_page_num(s, 15)

    lx = Inches(0.5)
    bw = Inches(1.2)
    bh = Inches(0.6)

    ty = Inches(1.2)
    add_textbox(s, lx, ty, Inches(5), Inches(0.3),
                '使用Teacher Forcing', font_size=11, font_color=C_FFN, bold=True)

    tf_y = ty + Inches(0.35)
    tf_items = ['[START]', '\u6211 OK', '\u7231 OK']
    tf_colors = [C_LIGHT_BLUE, C_LIGHT_GREEN, C_LIGHT_GREEN]
    tf_shapes = []
    for i, (txt, col) in enumerate(zip(tf_items, tf_colors)):
        x = lx + i * Inches(1.5)
        sh = add_rrect(s, x, tf_y, bw, bh, col, C_TEXT, Pt(0.8))
        shape_text(sh, txt, font_size=10, font_color=C_TEXT, bold=True)
        tf_shapes.append(sh)
        if i > 0:
            add_arrow_line(s, tf_shapes[i-1].left + tf_shapes[i-1].width,
                           tf_y + bh // 2, x, tf_y + bh // 2, C_FFN, Pt(1.2))
        add_textbox(s, x, tf_y + bh + Inches(0.05), bw, Inches(0.2),
                    '|预测', font_size=7, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(s, lx + Inches(0.5), tf_y + Inches(0.85), Inches(4), Inches(0.25),
                'OK = 真实答案 (效率高)', font_size=8, font_color=C_FFN, bold=True)

    sep_y = tf_y + Inches(1.3)
    add_line(s, lx, sep_y, lx + Inches(4.5), sep_y, C_GRAY, Pt(1))
    add_textbox(s, lx + Inches(1.5), sep_y - Inches(0.15), Inches(2), Inches(0.25),
                '效率高 vs 错误累积', font_size=9, font_color=C_ATTENTION, bold=True,
                alignment=PP_ALIGN.CENTER)

    ntf_y = sep_y + Inches(0.3)
    add_textbox(s, lx, ntf_y, Inches(5), Inches(0.3),
                '不使用Teacher Forcing', font_size=11, font_color=C_DECODER, bold=True)

    ntf_y2 = ntf_y + Inches(0.35)
    ntf_items = ['[START]', '\u4ed6 ERR', '\u4e0d ERR']
    ntf_colors = [C_LIGHT_BLUE, C_LIGHT_RED, C_LIGHT_RED]
    ntf_shapes = []
    for i, (txt, col) in enumerate(zip(ntf_items, ntf_colors)):
        x = lx + i * Inches(1.5)
        sh = add_rrect(s, x, ntf_y2, bw, bh, col, C_TEXT, Pt(0.8))
        shape_text(sh, txt, font_size=10, font_color=C_TEXT, bold=True)
        ntf_shapes.append(sh)
        if i > 0:
            add_arrow_line(s, ntf_shapes[i-1].left + ntf_shapes[i-1].width,
                           ntf_y2 + bh // 2, x, ntf_y2 + bh // 2, C_DECODER, Pt(1.2))
        add_textbox(s, x, ntf_y2 + bh + Inches(0.05), bw, Inches(0.2),
                    '|预测', font_size=7, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(s, lx + Inches(0.2), ntf_y2 + Inches(0.85), Inches(4.5), Inches(0.25),
                'ERR = 模型预测(可能出错)->错误累积', font_size=8, font_color=C_DECODER, bold=True)

    rx = Inches(6.0)
    _, tf_box = add_rich_textbox(s, rx, Inches(1.1), Inches(6.8), Inches(6.0))
    add_para(tf_box, '\u2460 原理：解码器输入不是模型预测，而是真实目标序列。每一步基于正确前文学习，像开卷考试',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf_box, '\u2461 优点：效率极高。所有位置并行计算，梯度信号稳定，充分利用GPU并行能力',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf_box, '\u2462 缺点：暴露偏差。训练看正确输入，推理看自己的输出，分布不一致。像学生练习有答案，考试没有',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf_box, '\u2463 缓解：Scheduled Sampling以一定概率用模型预测替代正确答案，缩小差距',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf_box, '\u2464 原始训练：WMT 2014英德翻译，450万句对，批量25000 token，学习率0.0003，30万步，8块P100 GPU，12小时',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))


# ============================================================
# Slide 16
# ============================================================
def slide_16():
    s = add_slide()
    add_title_bar(s, '推理阶段：模型如何独立生成回答？')
    add_page_num(s, 16)

    enc = add_rrect(s, Inches(0.3), Inches(1.2), Inches(2.5), Inches(1.0), C_ENCODER, C_ENCODER, Pt(1.5))
    shape_text(enc, '编码器 (只运行一次->缓存)\n处理' + LQ + 'I love you' + RQ, font_size=9, font_color=C_WHITE, bold=True)

    steps = [
        ('Step 1', '[START]', LQ + '\u6211' + RQ),
        ('Step 2', '[START, \u6211]', LQ + '\u7231' + RQ),
        ('Step 3', '[START, \u6211, \u7231]', LQ + '\u4f60' + RQ),
        ('Step 4', '[START, \u6211, \u7231, \u4f60]', LQ + '<END>' + RQ),
    ]

    sx = Inches(0.3)
    sy = Inches(2.6)
    sh = Inches(0.75)
    sg = Inches(0.2)

    for i, (label, inp, out) in enumerate(steps):
        y = sy + i * (sh + sg)
        add_textbox(s, sx, y - Inches(0.02), Inches(0.6), sh, label, font_size=8,
                    font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
        inp_sh = add_rrect(s, sx + Inches(0.6), y, Inches(1.5), sh, C_LIGHT_BLUE, C_ENCODER, Pt(0.5))
        shape_text(inp_sh, inp, font_size=8, font_color=C_ENCODER)
        add_arrow_line(s, sx + Inches(2.1), y + sh // 2, sx + Inches(2.4), y + sh // 2, C_TEXT, Pt(1))
        add_textbox(s, sx + Inches(2.4), y + Inches(0.1), Inches(0.3), Inches(0.5),
                    '解\n码\n器', font_size=7, font_color=C_TEXT, alignment=PP_ALIGN.CENTER)
        add_arrow_line(s, sx + Inches(2.7), y + sh // 2, sx + Inches(3.0), y + sh // 2, C_TEXT, Pt(1))
        out_sh = add_rrect(s, sx + Inches(3.0), y, Inches(0.8), sh, C_ATTENTION, C_ATTENTION, Pt(0.8))
        shape_text(out_sh, out, font_size=9, font_color=C_WHITE, bold=True)
        if i < len(steps) - 1:
            ny = y + sh
            add_arrow_line(s, sx + Inches(1.35), ny, sx + Inches(1.35), ny + sg, C_DECODER, Pt(1))

    add_textbox(s, sx + Inches(3.8), sy + Inches(1.2), Inches(1), Inches(0.3),
                '逐词生成', font_size=9, font_color=C_DECODER, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, sx + Inches(3.8), sy + Inches(1.4), Inches(1), Inches(0.3),
                '| 串行', font_size=8, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

    final = add_rrect(s, sx + Inches(0.5), sy + 4 * (sh + sg) - Inches(0.1), Inches(3.0), Inches(0.5),
                      C_ATTENTION, C_ATTENTION, Pt(1.5))
    shape_text(final, '最终输出：\u6211\u7231\u4f60', font_size=12, font_color=C_WHITE, bold=True)

    rx = Inches(6.5)
    _, tf = add_rich_textbox(s, rx, Inches(1.2), Inches(6.5), Inches(5.5))
    t1 = '\u2460 推理像' + LQ + '闭卷考试' + RQ + '：没有正确答案参考，完全靠自己之前生成的词预测下一个。第一题错了后面可能全错'
    t2 = '\u2461 编码器只需工作一次：处理输入->6层编码->输出向量->缓存，整个生成过程反复使用'
    t3 = '\u2462 解码器逐词串行生成：先生成第1个词，加入输入，生成第2个词...严格串行，无法并行化，推理速度瓶颈'
    add_para(tf, t1, font_size=11, font_color=C_TEXT, space_after=Pt(12))
    add_para(tf, t2, font_size=11, font_color=C_TEXT, space_after=Pt(12))
    add_para(tf, t3, font_size=11, font_color=C_TEXT, space_after=Pt(12))


# ============================================================
# Slide 17
# ============================================================
def slide_17():
    s = add_slide()
    add_title_bar(s, '推理\u00b7自回归生成：一个词一个词地写出回答')
    add_page_num(s, 17)

    lx1 = Inches(0.3)
    ly1 = Inches(1.2)
    add_textbox(s, lx1, ly1, Inches(3.5), Inches(0.3),
                '无缓存（慢）', font_size=12, font_color=C_DECODER, bold=True)

    no_cache = [
        'Step1: 算[START]的K,V',
        'Step2: 重新算全部K,V',
        'Step3: 重新算全部K,V',
    ]
    for i, txt in enumerate(no_cache):
        y = ly1 + Inches(0.4) + i * Inches(0.7)
        sh = add_rrect(s, lx1, y, Inches(3.2), Inches(0.5), C_LIGHT_RED, C_DECODER, Pt(0.8))
        shape_text(sh, txt, font_size=9, font_color=C_TEXT, bold=True)
        if i < 2:
            add_arrow_line(s, lx1 + Inches(1.6), y + Inches(0.5), lx1 + Inches(1.6), y + Inches(0.7),
                           C_DECODER, Pt(1))

    add_textbox(s, lx1 + Inches(0.5), ly1 + Inches(2.7), Inches(2.5), Inches(0.3),
                '重复计算！', font_size=11, font_color=C_DECODER, bold=True, alignment=PP_ALIGN.CENTER)

    lx2 = Inches(4.0)
    add_textbox(s, lx2, ly1, Inches(3.5), Inches(0.3),
                '有缓存（快）', font_size=12, font_color=C_FFN, bold=True)

    cache_steps = [
        ('Step1: 算[START]的K,V', C_LIGHT_GREEN),
        ('存入Cache', C_LIGHT_BLUE),
        ('Step2: 只算[\u6211]的K,V', C_LIGHT_GREEN),
        ('Step3: 只算[\u7231]的K,V', C_LIGHT_GREEN),
    ]
    for i, (txt, col) in enumerate(cache_steps):
        y = ly1 + Inches(0.4) + i * Inches(0.6)
        sh = add_rrect(s, lx2, y, Inches(3.2), Inches(0.45), col, C_FFN, Pt(0.8))
        shape_text(sh, txt, font_size=9, font_color=C_TEXT, bold=True)
        if i < 3:
            add_arrow_line(s, lx2 + Inches(1.6), y + Inches(0.45), lx2 + Inches(1.6), y + Inches(0.6),
                           C_FFN, Pt(1))

    add_textbox(s, lx2 + Inches(0.5), ly1 + Inches(2.9), Inches(2.5), Inches(0.3),
                '每步只算新词！', font_size=11, font_color=C_FFN, bold=True, alignment=PP_ALIGN.CENTER)

    cache_y = ly1 + Inches(3.3)
    cache_box = add_rrect(s, lx2, cache_y, Inches(3.2), Inches(1.5), C_LIGHT_PURPLE, C_EMBED, Pt(1.5))
    tf = cache_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'KV Cache'
    p.font.size = Pt(10)
    p.font.color.rgb = C_EMBED
    p.font.bold = True
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = FONT_EN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
    bodyPr = cache_box._element.txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    add_para(tf, 'Step1后: K[START], V[START]', font_size=8, font_color=C_TEXT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))
    add_para(tf, 'Step2后: K[START,\u6211], V[START,\u6211]', font_size=8, font_color=C_TEXT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))
    add_para(tf, 'Step3后: K[START,\u6211,\u7231], V[START,\u6211,\u7231]', font_size=8, font_color=C_TEXT, alignment=PP_ALIGN.CENTER, space_after=Pt(2))

    rx = Inches(7.8)
    _, tf2 = add_rich_textbox(s, rx, Inches(1.1), Inches(5.2), Inches(6.0))
    add_para(tf2, '\u2460 自回归生成：每生成一个新词就加入之前序列再预测下一个。像写文章，必须先想好第一个字才能写第二个',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf2, '\u2461 生成步骤：编码器处理' + LQ + 'I love you' + RQ + '(只做一次，缓存)。解码器[START]->' + LQ + '\u6211' + RQ + '->加入->[START,\u6211]->' + LQ + '\u7231' + RQ + '->...' + LQ + '<END>' + RQ + '停止',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf2, '\u2462 与训练区别：训练一次性看[START,\u6211,\u7231]并行计算；推理逐步扩展输入，每步重算。训练效率远高于推理',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))
    add_para(tf2, '\u2463 KV Cache=会议纪要：每步只算新词的Q,K,V，将新K,V追加到缓存。像开会只需查之前纪要，只讨论新增内容',
             font_size=10, font_color=C_TEXT, space_after=Pt(8))


# ============================================================
# Slide 18
# ============================================================
def slide_18():
    s = add_slide()
    add_title_bar(s, '推理\u00b7词选择：从贪心到创意')
    add_page_num(s, 18)

    lx = Inches(0.3)
    ly = Inches(1.2)
    cw = Inches(3.8)
    ch = Inches(0.7)
    cg = Inches(0.15)

    strategies = [
        ('\u2460 贪心搜索', '每次选最高概率词。无随机性。像走迷宫每次选最近路', C_ENCODER),
        ('\u2461 Beam Search', '保留top-k条路径同时探索。beam size=4-10。翻译效果好但偏平庸', C_FFN),
        ('\u2462 Top-K采样', '前k个词随机选。k=50既排除低质量又保留随机性', C_ATTENTION),
        ('\u2463 Top-P采样(核采样)', '累计概率达p(如0.9)的最小词集中随机选。比Top-K更灵活', C_EMBED),
        ('\u2464 Temperature', 'Softmax前缩放分数。低温(0.3)确定保守，高温(1.5)随机多样', C_DECODER),
    ]

    for i, (title, desc, col) in enumerate(strategies):
        y = ly + i * (ch + cg)
        title_bar = add_rect(s, lx, y, Inches(1.8), ch, col)
        shape_text(title_bar, title, font_size=9, font_color=C_WHITE, bold=True)
        add_textbox(s, lx + Inches(1.85), y + Inches(0.08), Inches(2.0), ch - Inches(0.1),
                    desc, font_size=8, font_color=C_TEXT)

    rx = Inches(6.8)
    add_textbox(s, rx, Inches(1.2), Inches(2.5), Inches(0.3),
                '低温 (T=0.3) - 尖锐确定', font_size=10, font_color=C_ENCODER, bold=True)

    low_t_heights = [Inches(2.0), Inches(0.3), Inches(0.15), Inches(0.08), Inches(0.05)]
    bar_w = Inches(0.3)
    for j, h in enumerate(low_t_heights):
        x = rx + j * Inches(0.4)
        y = Inches(1.5) + Inches(2.0) - h
        add_rect(s, x, y, bar_w, h, C_ENCODER)

    add_textbox(s, rx, Inches(4.0), Inches(2.5), Inches(0.3),
                '高温 (T=1.5) - 平坦多样', font_size=10, font_color=C_DECODER, bold=True)

    high_t_heights = [Inches(0.7), Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.45)]
    for j, h in enumerate(high_t_heights):
        x = rx + j * Inches(0.4)
        y = Inches(4.3) + Inches(0.7) - h
        add_rect(s, x, y, bar_w, h, C_DECODER)

    rx2 = Inches(9.5)
    _, tf = add_rich_textbox(s, rx2, Inches(1.2), Inches(3.5), Inches(5.5))
    add_para(tf, '贪心搜索最简单但最无聊，每次都选最优，输出固定不变', font_size=9, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, 'Beam Search平衡了质量和效率，是翻译任务标配', font_size=9, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, 'Top-K和Top-P引入随机性，让输出更有创意', font_size=9, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, 'Temperature是最常用参数，直接控制随机程度', font_size=9, font_color=C_TEXT, space_after=Pt(12))
    add_para(tf, 'ChatGPT通常用Temperature=0.7 + Top-P=0.9', font_size=10, font_color=C_ATTENTION, bold=True)


# ============================================================
# Slide 19
# ============================================================
def slide_19():
    s = add_slide()
    add_title_bar(s, '训练vs推理：开卷考试与闭卷考试的全面对比')
    add_page_num(s, 19)

    tx = Inches(0.5)
    ty = Inches(1.2)
    col_w = [Inches(2.0), Inches(4.5), Inches(4.5)]
    row_h = Inches(0.5)

    headers = [('维度', C_TITLE), ('训练', C_ENCODER), ('推理', C_FFN)]
    header_bg = [C_TITLE, C_ENCODER, C_FFN]
    for j, (txt, col) in enumerate(headers):
        x = tx + sum(col_w[:j])
        sh = add_rect(s, x, ty, col_w[j], row_h, header_bg[j])
        shape_text(sh, txt, font_size=11, font_color=C_WHITE, bold=True)

    rows = [
        ('解码器输入', '完整正确序列(Teacher Forcing)', '模型自己生成的词'),
        ('计算方式', '所有位置并行', '逐词串行'),
        ('损失函数', '交叉熵(每个位置计算)', '无损失，只做预测'),
        ('梯度计算', '反向传播更新参数', '无梯度(参数固定)'),
        ('KV Cache', '不需要', '必须使用'),
        ('效率', '高(并行)', '低(串行)'),
        ('目标', '学会预测', '生成结果'),
    ]

    for i, (dim, train, infer) in enumerate(rows):
        y = ty + (i + 1) * row_h
        bg = C_LIGHT_BLUE if i % 2 == 0 else C_WHITE
        sh1 = add_rect(s, tx, y, col_w[0], row_h, bg, C_TEXT, Pt(0.5))
        shape_text(sh1, dim, font_size=10, font_color=C_TEXT, bold=True)
        sh2 = add_rect(s, tx + col_w[0], y, col_w[1], row_h, C_BG_BLUE, C_TEXT, Pt(0.3))
        shape_text(sh2, train, font_size=10, font_color=C_ENCODER)
        sh3 = add_rect(s, tx + col_w[0] + col_w[1], y, col_w[2], row_h, C_BG_GREEN, C_TEXT, Pt(0.3))
        shape_text(sh3, infer, font_size=10, font_color=C_FFN)

    by = ty + 8 * row_h + Inches(0.3)
    _, tf = add_rich_textbox(s, Inches(0.5), by, Inches(12), Inches(1.5))
    add_para(tf, '比喻总结：', font_size=11, font_color=C_TITLE, bold=True, space_after=Pt(4))
    add_para(tf, '训练 = 新员工培训期，师傅手把手教，效率高需大量资源。推理 = 正式上岗，全靠自己，每步参考之前结果，只需一台电脑',
             font_size=11, font_color=C_TEXT, space_after=Pt(4))


# ============================================================
# Slide 20
# ============================================================
def slide_20():
    s = add_slide()
    add_title_bar(s, '从一篇论文到三大家族：同一个DNA，三种职业方向')
    add_page_num(s, 20)

    root_x = Inches(3.0)
    root_y = Inches(1.2)
    root = add_circle(s, root_x, root_y, Inches(1.2), C_TITLE)
    shape_text(root, 'Transformer\n2017', font_size=10, font_color=C_WHITE, bold=True)

    branch_y = Inches(2.8)

    # Left: GPT
    l1x = Inches(0.3)
    add_arrow_line(s, root_x + Inches(0.3), root_y + Inches(1.2), l1x + Inches(1.2), branch_y, C_ENCODER, Pt(1.5))
    l1_label = add_rrect(s, l1x, branch_y, Inches(2.5), Inches(0.5), C_ENCODER)
    shape_text(l1_label, '仅用解码器', font_size=10, font_color=C_WHITE, bold=True)

    gpt_models = ['GPT-1', 'GPT-2', 'GPT-3', 'GPT-4', 'ChatGPT']
    for i, m in enumerate(gpt_models):
        x = l1x + i * Inches(0.5)
        y = branch_y + Inches(0.6)
        sh = add_rrect(s, x, y, Inches(0.45), Inches(0.35), C_LIGHT_BLUE, C_ENCODER, Pt(0.5))
        shape_text(sh, m, font_size=7, font_color=C_ENCODER, bold=True)

    add_textbox(s, l1x, branch_y + Inches(1.1), Inches(2.5), Inches(0.3),
                '即兴演讲者', font_size=9, font_color=C_ENCODER, bold=True, alignment=PP_ALIGN.CENTER)

    # Middle: T5
    m1x = Inches(3.5)
    add_arrow_line(s, root_x + Inches(0.6), root_y + Inches(1.2), m1x + Inches(1.2), branch_y, C_FFN, Pt(1.5))
    m1_label = add_rrect(s, m1x, branch_y, Inches(2.8), Inches(0.5), C_FFN)
    shape_text(m1_label, '完整编码器+解码器', font_size=10, font_color=C_WHITE, bold=True)

    t5_models = ['T5', 'BART', 'Whisper']
    for i, m in enumerate(t5_models):
        x = m1x + Inches(0.3) + i * Inches(0.8)
        y = branch_y + Inches(0.6)
        sh = add_rrect(s, x, y, Inches(0.7), Inches(0.35), C_LIGHT_GREEN, C_FFN, Pt(0.5))
        shape_text(sh, m, font_size=8, font_color=C_FFN, bold=True)

    add_textbox(s, m1x, branch_y + Inches(1.1), Inches(2.8), Inches(0.3),
                '全能翻译官', font_size=9, font_color=C_FFN, bold=True, alignment=PP_ALIGN.CENTER)

    # Right: BERT
    r1x = Inches(7.0)
    add_arrow_line(s, root_x + Inches(0.9), root_y + Inches(1.2), r1x + Inches(1.2), branch_y, C_EMBED, Pt(1.5))
    r1_label = add_rrect(s, r1x, branch_y, Inches(2.5), Inches(0.5), C_EMBED)
    shape_text(r1_label, '仅用编码器', font_size=10, font_color=C_WHITE, bold=True)

    bert_models = ['BERT', 'RoBERTa', 'DeBERTa']
    for i, m in enumerate(bert_models):
        x = r1x + Inches(0.2) + i * Inches(0.8)
        y = branch_y + Inches(0.6)
        sh = add_rrect(s, x, y, Inches(0.7), Inches(0.35), C_LIGHT_PURPLE, C_EMBED, Pt(0.5))
        shape_text(sh, m, font_size=8, font_color=C_EMBED, bold=True)

    add_textbox(s, r1x, branch_y + Inches(1.1), Inches(2.5), Inches(0.3),
                '阅读理解专家', font_size=9, font_color=C_EMBED, bold=True, alignment=PP_ALIGN.CENTER)

    # Parameter scale bar
    bar_y = Inches(5.0)
    params = ['65M', '110M', '1.5B', '175B', '1.8T']
    colors_bar = [C_LIGHT_BLUE, C_ENCODER, C_FFN, C_ATTENTION, C_DECODER]
    for i, (p, col) in enumerate(zip(params, colors_bar)):
        x = Inches(0.5) + i * Inches(1.0)
        w = Inches(0.3) + i * Inches(0.15)
        h = Inches(0.2) + i * Inches(0.12)
        add_rect(s, x, bar_y + Inches(0.4) - h, w, h, col)
        add_textbox(s, x - Inches(0.1), bar_y + Inches(0.45), Inches(0.8), Inches(0.25),
                    p, font_size=8, font_color=C_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.5), bar_y - Inches(0.3), Inches(5), Inches(0.25),
                '参数量递增 ->', font_size=9, font_color=C_TEXT, bold=True)

    rx = Inches(7.0)
    _, tf = add_rich_textbox(s, rx, Inches(5.0), Inches(6.0), Inches(2.3))
    add_para(tf, '\u2460 GPT系列(仅解码器)：核心能力' + LQ + '文本生成' + RQ + '。ChatGPT三步训练：预训练(预测下一个词)->指令微调(理解人类指令)->RLHF(人类反馈强化学习)。代表：GPT-3(1750亿)、GPT-4(1.8万亿)、LLaMA、Qwen',
             font_size=9, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, '\u2461 BERT系列(仅编码器)：核心能力' + LQ + '文本理解' + RQ + '。掩码语言建模(随机遮盖15%token)。Google搜索2019年开始使用。代表：BERT-base(1.1亿)、BERT-large(3.4亿)',
             font_size=9, font_color=C_TEXT, space_after=Pt(6))
    add_para(tf, '\u2462 T5系列(完整架构)：将所有NLP任务统一为文本到文本格式。Google Translate基于T5变体。代表：T5-Small(6000万)到T5-11B(110亿)',
             font_size=9, font_color=C_TEXT, space_after=Pt(6))


# ============================================================
# Slide 21
# ============================================================
def slide_21():
    s = add_slide()
    add_title_bar(s, '从翻译工具到AI通用引擎：Transformer改变了一切')
    add_page_num(s, 21)

    cx = Inches(3.5)
    cy = Inches(3.5)
    r = Inches(2.0)

    center = add_circle(s, cx + Inches(0.3), cy + Inches(0.3), Inches(1.2), C_TITLE)
    shape_text(center, 'Transformer', font_size=11, font_color=C_WHITE, bold=True)

    domains = [
        ('NLP\n翻译/对话/写作', 270, C_ENCODER),
        ('计算机视觉\nViT/图像生成', 330, C_FFN),
        ('语音\nWhisper/语音识别', 30, C_ATTENTION),
        ('多模态\nGPT-4V/Gemini', 90, C_DECODER),
        ('蛋白质\nAlphaFold2', 150, C_EMBED),
        ('时间序列\n预测/推荐', 210, C_LIGHT_BLUE),
    ]

    for txt, angle_deg, col in domains:
        angle_rad = math.radians(angle_deg)
        dx = int(cx + Inches(0.9) + r * math.cos(angle_rad))
        dy = int(cy + Inches(0.9) + r * math.sin(angle_rad))
        line_end_x = int(cx + Inches(0.9) + (r - Inches(0.4)) * math.cos(angle_rad))
        line_end_y = int(cy + Inches(0.9) + (r - Inches(0.4)) * math.sin(angle_rad))
        add_line(s, cx + Inches(0.9), cy + Inches(0.9), line_end_x, line_end_y, col, Pt(1.5))
        sh = add_rrect(s, dx, dy, Inches(1.8), Inches(0.7), col)
        shape_text(sh, txt, font_size=8, font_color=C_WHITE, bold=True)

    tl_y = Inches(6.5)
    tl_x = Inches(0.5)
    milestones = [
        ('2017', 'Transformer论文'),
        ('2018', 'BERT/GPT-1'),
        ('2020', 'GPT-3/ViT'),
        ('2022', 'ChatGPT/Whisper'),
        ('2023', 'GPT-4'),
        ('2024+', '多模态爆发'),
    ]
    for i, (year, event) in enumerate(milestones):
        x = tl_x + i * Inches(2.0)
        add_circle(s, x + Inches(0.1), tl_y, Inches(0.15), C_TITLE)
        add_textbox(s, x - Inches(0.1), tl_y - Inches(0.3), Inches(0.8), Inches(0.25),
                    year, font_size=8, font_color=C_TITLE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x - Inches(0.3), tl_y + Inches(0.2), Inches(1.2), Inches(0.3),
                    event, font_size=7, font_color=C_TEXT, alignment=PP_ALIGN.CENTER)
        if i < 5:
            add_line(s, x + Inches(0.3), tl_y + Inches(0.07), x + Inches(1.9), tl_y + Inches(0.07), C_GRAY, Pt(1))

    rx = Inches(7.5)
    _, tf = add_rich_textbox(s, rx, Inches(1.2), Inches(5.5), Inches(5.0))
    add_para(tf, '\u2460 规模爆发：2017年6500万参数，今天GPT-4约1.8万亿(增长3万倍)。架构没变，只是层数更多(6->80)、维度更大(512->8192)',
             font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2461 超越语言：ViT(2020)将图像切成小块当词处理；Whisper(2022)语音识别；GPT-4V多模态；AlphaFold2预测蛋白质结构',
             font_size=10, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2462 核心洞察：Transformer是通用信息处理框架。只要能把数据表示成序列，就能用Transformer处理',
             font_size=10, font_color=C_TEXT, space_after=Pt(10))


# ============================================================
# Slide 22
# ============================================================
def slide_22():
    s = add_slide()
    add_title_bar(s, '总结：带走今天的5个核心要点')
    add_page_num(s, 22)

    lx = Inches(0.5)
    ly = Inches(1.2)
    cw = Inches(3.0)
    ch = Inches(0.8)
    cg = Inches(0.2)

    cards = [
        ('\u2460 信息自由流动', '每个词与所有词直接交流\n解决RNN三大问题', C_ENCODER),
        ('\u2461 编码器理解+解码器生成', '编码器像编辑部\n解码器像翻译部', C_EMBED),
        ('\u2462 训练开卷+推理闭卷', 'Teacher Forcing vs 自回归\n最大隐患：暴露偏差', C_ATTENTION),
        ('\u2463 GPT/BERT/T5三大家族', 'GPT生成->ChatGPT\nBERT理解->搜索\nT5转换->翻译', C_FFN),
        ('\u2464 AI通用引擎', '6500万->万亿参数\nTransformer是AI时代的基石', C_DECODER),
    ]

    icons = ['\U0001f517', '\U0001f4da', '\U0001f4dd', '\U0001f500', '\U0001f680']

    for i, ((title, desc, col), icon) in enumerate(zip(cards, icons)):
        y = ly + i * (ch + cg)
        icon_bg = add_rect(s, lx, y, Inches(0.6), ch, col)
        tf_icon = icon_bg.text_frame
        tf_icon.word_wrap = True
        p = tf_icon.paragraphs[0]
        p.text = icon
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        bodyPr = icon_bg._element.txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')

        content = add_rrect(s, lx + Inches(0.65), y, cw - Inches(0.65), ch, C_WHITE, col, Pt(1.5))
        _, tf = add_rich_textbox(s, lx + Inches(0.75), y + Inches(0.05), cw - Inches(0.85), ch - Inches(0.1))
        add_para(tf, title, font_size=10, font_color=col, bold=True, space_after=Pt(2))
        add_para(tf, desc, font_size=8, font_color=C_TEXT, space_after=Pt(0))

    rx = Inches(4.5)
    _, tf = add_rich_textbox(s, rx, Inches(1.2), Inches(8.5), Inches(5.5))
    add_para(tf, '\u2460 Transformer的本质：让信息自由流动。每个词与所有词直接交流，解决RNN三大问题',
             font_size=11, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2461 编码器=理解，解码器=生成。编码器像编辑部，解码器像翻译部',
             font_size=11, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2462 训练=开卷考试，推理=闭卷考试。Teacher Forcing vs 自回归生成。最大隐患是暴露偏差',
             font_size=11, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2463 三大家族各有所长：GPT生成->ChatGPT，BERT理解->搜索，T5转换->翻译',
             font_size=11, font_color=C_TEXT, space_after=Pt(10))
    add_para(tf, '\u2464 从论文到通用引擎：6500万到万亿参数，Transformer是AI时代的基石',
             font_size=11, font_color=C_TEXT, space_after=Pt(16))

    add_textbox(s, Inches(4.5), Inches(6.5), Inches(8.5), Inches(0.4),
                'Transformer不是一个模型，而是一个时代。感谢聆听！',
                font_size=10, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)


# === Generate all slides ===
slide_12()
slide_13()
slide_14()
slide_15()
slide_16()
slide_17()
slide_18()
slide_19()
slide_20()
slide_21()
slide_22()

# === Save ===
output_path = '/home/admin/.openclaw/workspace-weaver/output/transformer_v3_zh_part2.pptx'
prs.save(output_path)
print(f'OK: {output_path}')

# === Self-check ===
print('\n=== Self-check ===')
print(f'Slide count: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width / 914400:.3f} x {prs.slide_height / 914400:.3f} inches')

issues = 0
for i, slide in enumerate(prs.slides):
    shapes = slide.shapes
    oob = []
    for sh in shapes:
        if hasattr(sh, 'left') and hasattr(sh, 'width'):
            r = sh.left + sh.width
            b = sh.top + sh.height
            if r > SAFE_W + Inches(0.15):
                oob.append(f'  over-right: {sh.name} right={r/914400:.2f}')
            if b > SAFE_H + Inches(0.15):
                oob.append(f'  over-bottom: {sh.name} bottom={b/914400:.2f}')

    slide_num = i + 12
    status = 'OK' if not oob else 'WARN'
    print(f'Slide {slide_num}: {len(shapes)} shapes [{status}]')
    for w in oob:
        print(w)
        issues += 1

print(f'\nTotal issues: {issues}')
print('Done!')
