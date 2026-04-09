#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer 架构深度解析 PPT 生成脚本 v2
15页，面向非AI背景受众，图文并茂

v2 修复清单：
- 排版修复：Slide 3/5/7/9/15 的形状重叠
- 术语解释：梯度消失、权重矩阵、512维、Softmax、饱和区、线性投影、因果掩码
- 内容加深：Q/K/V角色、具体数字例子、RoPE比喻
- 衔接优化：Slide 8/9 过渡句
- 内容增加：Slide 15 推荐阅读
- 年份修正：8年→9年
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ============================================================
# 配置
# ============================================================
SLIDE_W = 13.333  # 英寸
SLIDE_H = 7.5

# 颜色
C_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)    # 深海蓝
C_ACCENT = RGBColor(0x34, 0x98, 0xDB)     # 科技蓝
C_LIGHT_BG = RGBColor(0xE8, 0xF4, 0xFD)   # 浅蓝背景
C_TEXT = RGBColor(0x33, 0x33, 0x33)        # 深灰文字
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)

# 字体
FONT_ZH = "Microsoft YaHei"
FONT_EN = "Arial"

# ============================================================
# 工具函数
# ============================================================
def add_textbox(slide, left, top, width, height, text, font_size=12, 
                bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font_name=FONT_ZH,
                line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """添加文本框，自动边界检查"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                       Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    
    # 边界检查
    assert left + width <= SLIDE_W - 0.1, f"右边界超出: {left+width:.2f} > {SLIDE_W-0.1:.2f}"
    assert top + height <= SLIDE_H - 0.2, f"下边界超出: {top+height:.2f} > {SLIDE_H-0.2:.2f}"
    
    return txBox

def add_shape(slide, shape_type, left, top, width, height, 
              fill_color=None, line_color=None, line_width=Pt(1)):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    
    # 边界检查
    assert left + width <= SLIDE_W, f"形状右边界超出: {left+width:.2f}"
    assert top + height <= SLIDE_H, f"形状下边界超出: {top+height:.2f}"
    
    return shape

def add_multiline_textbox(slide, left, top, width, height, lines, 
                           font_size=11, color=C_TEXT, line_spacing=1.1,
                           font_name=FONT_ZH, bold_list=None, color_list=None):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 支持加粗标记: **text** 加粗
        if "**" in line:
            run = p.add_run()
            parts = line.split("**")
            for j, part in enumerate(parts):
                r = p.add_run() if j > 0 else run
                r.text = part
                r.font.size = Pt(font_size)
                r.font.name = font_name
                if j % 2 == 1:  # 奇数部分是加粗的
                    r.font.bold = True
                    r.font.color.rgb = C_PRIMARY
                else:
                    r.font.color.rgb = color if not color_list else (color_list[i] if i < len(color_list) else color)
            p.font.size = Pt(font_size)
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.name = font_name
            p.font.color.rgb = color if not color_list else (color_list[i] if i < len(color_list) else color)
        
        if bold_list and i < len(bold_list) and bold_list[i]:
            p.font.bold = True
        
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        if line_spacing != 1.0:
            p.line_spacing = Pt(font_size * line_spacing)
    
    assert left + width <= SLIDE_W - 0.1, f"右边界超出"
    assert top + height <= SLIDE_H - 0.2, f"下边界超出"
    
    return txBox

def add_connector_line(slide, start_left, start_top, end_left, end_top, 
                        color=C_GRAY, width=Pt(2)):
    """添加连接线"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def set_shape_text(shape, text, font_size=11, color=C_WHITE, bold=False, 
                    font_name=FONT_ZH, align=PP_ALIGN.CENTER):
    """设置形状内的文字"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    return shape

def add_rounded_rect(slide, left, top, width, height, text="", 
                      fill_color=C_ACCENT, text_color=C_WHITE, font_size=11,
                      line_color=None, line_width=Pt(1)):
    """添加圆角矩形"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                       fill_color=fill_color, line_color=line_color, line_width=line_width)
    if text:
        set_shape_text(shape, text, font_size=font_size, color=text_color, bold=True)
    return shape

# ============================================================
# 创建主PPT
# ============================================================
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
blank_layout = prs.slide_layouts[6]  # 空白布局

print(f"主PPT幻灯片尺寸: {SLIDE_W} x {SLIDE_H} 英寸")

# ============================================================
# 第1页：封面
# ============================================================
print("创建第1页：封面...")
slide = prs.slides.add_slide(blank_layout)

# 深色背景
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, 
          fill_color=C_DARK_BG)

# 装饰线条
for i in range(5):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.5 + i * 1.2, SLIDE_W, 0.01,
              fill_color=C_ACCENT)

# 主标题
add_textbox(slide, 1.5, 1.8, 10, 1.2, "Transformer 架构深度解析",
            font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide, 2, 3.2, 9, 0.8, "从直觉到原理，一篇看懂AI的核心引擎",
            font_size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 分隔线
add_shape(slide, MSO_SHAPE.RECTANGLE, 4.5, 4.2, 4, 0.02, fill_color=C_ACCENT)

# 说明
add_textbox(slide, 2, 4.6, 9, 0.5, "2017年Google提出，ChatGPT、GPT-4等所有大语言模型的基石",
            font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

# 底部标签
add_textbox(slide, 2, 6.5, 9, 0.4, "Weaver  |  2026",
            font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 第2页：目录
# ============================================================
print("创建第2页：目录...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 6, 0.6, "今天的内容路线",
            font_size=26, bold=True, color=C_PRIMARY)

# 左侧目录列表
toc_items = [
    ("01", "前世今生", "RNN/LSTM 的三大痛点", C_ACCENT),
    ("02", "核心直觉", "每个词都能和所有其他词直接交流", C_ACCENT),
    ("03", "自注意力（上）", "Q、K、V——图书馆里的搜索系统", C_ORANGE),
    ("04", "自注意力（下）", "5步计算，从输入到输出的完整旅程", C_ORANGE),
    ("05", "多头注意力", "8个专家从不同角度同时分析", C_ORANGE),
    ("06", "交叉注意力", "编码器和解码器之间的翻译桥梁", C_ORANGE),
    ("07", "位置编码", "告诉 Transformer 词的先后顺序", C_ORANGE),
    ("08", "整体架构", "从输入到输出，数据如何流动", C_GREEN),
    ("09", "编码器详解", "信息如何被理解和提炼", C_GREEN),
    ("10", "解码器详解", "信息如何被生成和表达", C_GREEN),
    ("11", "训练与推理", "模型是怎么学会说话的？", C_GREEN),
    ("12", "演进与应用", "BERT、GPT、T5 三大家族", C_GREEN),
    ("13", "总结回顾", "五个核心要点", C_GREEN),
]

y_start = 1.1
for i, (num, title, desc, color) in enumerate(toc_items):
    y = y_start + i * 0.40
    
    # 编号圆圈
    circle = add_shape(slide, MSO_SHAPE.OVAL, 0.5, y, 0.35, 0.35,
                        fill_color=color)
    set_shape_text(circle, num, font_size=9, color=C_WHITE, bold=True)
    
    # 标题
    add_textbox(slide, 1.0, y, 3.5, 0.35, title,
                font_size=12, bold=True, color=C_TEXT)
    
    # 描述
    add_textbox(slide, 4.5, y, 4.5, 0.35, desc,
                font_size=10, color=C_GRAY)

# 右侧阶段说明
stages = [
    ("理解问题", C_ACCENT, "第1-2节"),
    ("核心机制", C_ORANGE, "第3-7节"),
    ("架构与应用", C_GREEN, "第8-13节"),
]
for i, (label, color, scope) in enumerate(stages):
    y = 1.5 + i * 1.8
    add_rounded_rect(slide, 9.5, y, 3.2, 0.5, label, fill_color=color, font_size=13)
    add_textbox(slide, 9.5, y + 0.55, 3.2, 0.3, scope,
                font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 第3页：Transformer 之前的世界
# ============================================================
print("创建第3页：Transformer 之前的世界...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 10, 0.6, "前世今生：RNN/LSTM 的三大痛点",
            font_size=24, bold=True, color=C_PRIMARY)

# 左侧：RNN流程图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "RNN 处理文本的方式",
            font_size=14, bold=True, color=C_ACCENT)

# 时间线
words = ["The", "cat", "sat", "on", "the", "mat"]
x_start = 0.8
for i, word in enumerate(words):
    x = x_start + i * 1.3
    # 词框
    add_rounded_rect(slide, x, 1.5, 1.0, 0.5, word, fill_color=C_ACCENT, font_size=12)
    # RNN单元
    add_rounded_rect(slide, x + 0.1, 2.2, 0.8, 0.4, "RNN", fill_color=C_ORANGE, font_size=10)
    # 箭头
    if i < len(words) - 1:
        add_connector_line(slide, x + 1.0, 1.75, x + 1.3, 1.75, color=C_GRAY, width=Pt(2))

# 三个痛点标注
pain_points = [
    (0.8, 3.0, "1. 顺序处理", "必须等前一个词处理完\n才能处理下一个词", "必须一个字一个字看"),
    (3.0, 3.0, "2. 长距离遗忘", "处理后面的词时\n逐渐忘记前面的信息", "听到结尾忘了开头"),
    (5.2, 3.0, "3. 信息瓶颈", "所有信息压缩成\n一个固定长度向量", "500页书压成100字摘要"),
]

for x, y, title, desc, analogy in pain_points:
    # 痛点框
    add_rounded_rect(slide, x, y, 2.0, 0.5, title, fill_color=C_RED, font_size=11)
    # 描述
    add_textbox(slide, x, y + 0.55, 2.0, 0.7, desc,
                font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER)
    # 比喻
    add_textbox(slide, x, y + 1.3, 2.0, 0.4, analogy,
                font_size=9, color=C_ORANGE, align=PP_ALIGN.CENTER)

# 右侧：文字说明 —— v2修复：left从7.5改为8.5，避免覆盖mat/RNN节点
lines = [
    "在2017年Transformer出现之前，AI处理",
    "语言主要靠RNN和LSTM，但它们有三个",
    "根本性的弱点：",
    "",
    "1. **顺序处理** -- 像读书只能一个字一个",
    "   字看，无法利用GPU并行计算，处理极慢。",
    "",
    "2. **长距离遗忘** -- 处理后面的词时逐渐",
    "   忘记前面的信息（梯度消失，即信息在反",
    "   向传播过程中逐渐衰减，就像传话游戏中",
    "   传到最后信息面目全非），句子超过50个",
    "   词性能就明显下降。",
    "",
    "3. **信息瓶颈** -- 所有信息压缩成一个",
    "   固定向量，就像把500页的书压成100字。",
    "",
    "**核心问题：信息流动不够自由。**",
    "**解决方案：让每个词直接和所有词交流！**",
]
add_multiline_textbox(slide, 8.5, 1.0, 4.3, 6.0, lines,
                       font_size=9, line_spacing=1.1)

# ============================================================
# 第4页：Transformer 的核心直觉
# ============================================================
print("创建第4页：Transformer 的核心直觉...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "一句话理解 Transformer：每个词都能和所有其他词直接交流",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：注意力可视化
add_textbox(slide, 0.5, 1.0, 6, 0.4, "注意力可视化：it 如何找到 cat？",
            font_size=13, bold=True, color=C_ACCENT)

# 句子词
sentence_words = ["The", "cat", "sat", "on", "the", "mat", ",", "and", "it", "was", "happy"]
x_start = 0.3
y_words = 1.6
for i, word in enumerate(sentence_words):
    x = x_start + i * 0.58
    # it 高亮
    if word == "it":
        add_rounded_rect(slide, x, y_words, 0.5, 0.4, word, fill_color=C_ACCENT, font_size=9)
    elif word == "cat":
        add_rounded_rect(slide, x, y_words, 0.5, 0.4, word, fill_color=C_GREEN, font_size=9)
    else:
        add_rounded_rect(slide, x, y_words, 0.5, 0.4, word, fill_color=C_GRAY, font_size=9)

# 注意力连线标注
add_textbox(slide, 0.5, 2.3, 6, 0.4, "it 对每个词的注意力权重（数字越大 = 关联度越强）：",
            font_size=10, color=C_TEXT)

weights = [
    ("The", 0.05), ("cat", 0.45), ("sat", 0.10), ("on", 0.03),
    ("the", 0.02), ("mat", 0.08), (",", 0.01), ("and", 0.02),
    ("it", 0.12), ("was", 0.07), ("happy", 0.05),
]
x_bar_start = 0.5
for i, (word, weight) in enumerate(weights):
    y = 2.7 + i * 0.33
    # 词
    add_textbox(slide, x_bar_start, y, 0.8, 0.3, word, font_size=9, color=C_TEXT)
    # 权重条
    bar_width = weight * 8
    bar_color = C_GREEN if word == "cat" else (C_ACCENT if weight > 0.1 else C_GRAY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.5, y + 0.05, bar_width, 0.2, fill_color=bar_color)
    # 权重值
    add_textbox(slide, 1.5 + bar_width + 0.1, y, 0.6, 0.3, f"{weight:.0%}", font_size=8, color=C_TEXT)

# cat标注
add_textbox(slide, 3.0, 2.7, 3.5, 0.3, "cat 关联度最高！it 指的就是 cat",
            font_size=11, bold=True, color=C_GREEN)

# 右侧：文字说明
lines = [
    "**Transformer 的核心思想：**",
    "RNN 像排队发言 -- 每个人等前面说完才能说，",
    "说完就忘了。",
    "",
    "Transformer 像圆桌会议 -- 每个人可以同时",
    "听到所有人说话，并关注任何一个人。",
    "",
    "**具体例子：**",
    "句子：The cat sat on the mat, and it was happy",
    "",
    "当模型处理 it 时，它会同时查看所有词，",
    "计算与每个词的关联度，发现 it 和 cat 的",
    "关联度最高（都是主语，语法关系密切），",
    "于是把更多注意力放在 cat 上。",
    "",
    "**效率优势：**",
    "所有词可以同时计算，充分利用GPU并行能力。",
    "原始论文中，8块GPU训练12小时就达到当时最优，",
    "而之前最好的模型需要数天。这是Transformer",
    "能被扩展到数千亿参数的根本原因。",
]
add_multiline_textbox(slide, 7.5, 1.0, 5.3, 6.2, lines,
                       font_size=10, line_spacing=1.1)

# ============================================================
# 第5页：Self-Attention（上）Q/K/V
# ============================================================
print("创建第5页：Self-Attention（上）...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 10, 0.6, "自注意力（上）：Q、K、V -- 图书馆里的搜索系统",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：Q/K/V 生成过程图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "Q、K、V 的生成过程",
            font_size=13, bold=True, color=C_ACCENT)

# 输入框 —— v2修复：512维加类比解释
input_box = add_rounded_rect(slide, 2.0, 1.5, 3.0, 0.6, "输入词向量 (512维)", 
                               fill_color=C_PRIMARY, font_size=12)

# 三条分叉箭头
add_connector_line(slide, 3.5, 2.1, 1.0, 2.8, color=C_ACCENT, width=Pt(2))
add_connector_line(slide, 3.5, 2.1, 3.5, 2.8, color=C_ORANGE, width=Pt(2))
add_connector_line(slide, 3.5, 2.1, 6.0, 2.8, color=C_GREEN, width=Pt(2))

# W^Q, W^K, W^V 矩阵
add_rounded_rect(slide, 0.3, 2.8, 1.8, 0.5, "W^Q 矩阵", fill_color=C_ACCENT, font_size=11)
add_rounded_rect(slide, 2.8, 2.8, 1.8, 0.5, "W^K 矩阵", fill_color=C_ORANGE, font_size=11)
add_rounded_rect(slide, 5.3, 2.8, 1.8, 0.5, "W^V 矩阵", fill_color=C_GREEN, font_size=11)

# Q, K, V 输出
add_rounded_rect(slide, 0.3, 3.6, 1.8, 0.5, "Query (Q)", fill_color=C_ACCENT, font_size=12)
add_rounded_rect(slide, 2.8, 3.6, 1.8, 0.5, "Key (K)", fill_color=C_ORANGE, font_size=12)
add_rounded_rect(slide, 5.3, 3.6, 1.8, 0.5, "Value (V)", fill_color=C_GREEN, font_size=12)

# 箭头
add_connector_line(slide, 1.2, 3.3, 1.2, 3.6, color=C_ACCENT, width=Pt(2))
add_connector_line(slide, 3.7, 3.3, 3.7, 3.6, color=C_ORANGE, width=Pt(2))
add_connector_line(slide, 6.2, 3.3, 6.2, 3.6, color=C_GREEN, width=Pt(2))

# 图书馆比喻图
add_textbox(slide, 0.5, 4.4, 7, 0.4, "图书馆搜索类比",
            font_size=13, bold=True, color=C_ACCENT)

# 搜索框 → 标签 → 内容 —— v2修复：V框右移避免与右侧文字重叠
add_rounded_rect(slide, 0.5, 5.0, 1.8, 0.5, "搜索关键词 (Q)", fill_color=C_ACCENT, font_size=10)
add_textbox(slide, 2.5, 5.1, 0.5, 0.3, "->", font_size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
add_rounded_rect(slide, 3.0, 5.0, 1.8, 0.5, "书名标签 (K)", fill_color=C_ORANGE, font_size=10)
add_textbox(slide, 5.0, 5.1, 0.5, 0.3, "->", font_size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
add_rounded_rect(slide, 5.5, 5.0, 1.8, 0.5, "书的内容 (V)", fill_color=C_GREEN, font_size=10)

# 比喻说明
add_textbox(slide, 0.5, 5.7, 7, 0.8, 
            "你输入搜索关键词(Q)，系统将它与每本书的标签(K)做匹配，\n找到匹配度最高的书后，提取书的内容(V)。",
            font_size=10, color=C_TEXT)

# 右侧：文字说明 —— v2修复：增加权重矩阵解释、512维类比、为什么三个向量
lines = [
    "**1. Query（查询）-- 你在找什么**",
    "每个词生成一个Query向量，代表我在找什么信息。",
    "就像图书馆搜索时输入的关键词 -- 处理it时，",
    "Query就是'我需要找到一个名词告诉我it指谁'。",
    "",
    "**2. Key（键）-- 我有什么特征**",
    "每个词同时生成一个Key向量，代表我有什么特征",
    "可以被别人找到。就像图书馆里每本书的书名和",
    "标签 -- cat的Key包含名词、动物、可作主语。",
    "",
    "**3. Value（值）-- 我的实际内容**",
    "每个词还生成一个Value向量，代表实际内容。",
    "就像书的内容本身 -- it找到cat的Key匹配最高后，",
    "就读取cat的Value（一只猫、动物、宠物等）。",
    "",
    "**4. 工作流程：搜索->匹配->提取**",
    "Q搜索所有K -> 找到最匹配的 -> 提取对应的V。",
    "每个词同时是搜索者也是被搜索者！",
    "",
    "**5. 为什么需要三个不同的向量？**",
    "因为Q是主动搜索的一方，K是被搜索的特征，",
    "V是被提取的内容 -- 角色不同，所以需要不同的表示。",
]
add_multiline_textbox(slide, 7.8, 1.0, 5.2, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第6页：Self-Attention（下）计算过程
# ============================================================
print("创建第6页：Self-Attention（下）...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "自注意力（下）：5步计算 -- 从输入到输出的完整旅程",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：5步流程图
steps = [
    ("1", "输入 -> 生成 Q/K/V", "每个词经三个权重矩阵\n变换得到Q、K、V向量", C_ACCENT),
    ("2", "Q x K^T -> 计算关联度", "每个词的Query和所有词的\nKey做点积，得到分数矩阵", C_ORANGE),
    ("3", "分数 / sqrt(d_k) -> 缩放", "防止分数过大导致Softmax\n进入饱和区", C_RED),
    ("4", "Softmax -> 转成概率", "将分数转换为0-1的概率\n分布，每行之和等于1", C_PURPLE),
    ("5", "概率 x V -> 加权求和", "按概率权重对所有词的\nValue加权求和，得到输出", C_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = 1.0 + i * 1.25
    
    # 步骤框
    add_rounded_rect(slide, 0.5, y, 0.6, 0.5, num, fill_color=color, font_size=16)
    
    # 标题
    add_textbox(slide, 1.3, y, 5.5, 0.3, title,
                font_size=12, bold=True, color=C_TEXT)
    
    # 描述
    add_textbox(slide, 1.3, y + 0.3, 5.5, 0.5, desc,
                font_size=9, color=C_GRAY)
    
    # 箭头
    if i < len(steps) - 1:
        add_connector_line(slide, 0.8, y + 0.5, 0.8, y + 1.0, color=C_GRAY, width=Pt(2))

# 右侧：文字说明 —— v2修复：Softmax详细解释、饱和区解释、具体数字例子
lines = [
    "**步骤1: 生成Q、K、V**",
    "输入词向量(512维，即每个词用512个数字来描述",
    "它的各种特征)经过三个权重矩阵(就是一组可学",
    "习的参数，决定如何从原始词向量中提取不同信息)",
    "变换，每个词得到64维的Q、K、V。可对所有词同时进行。",
    "",
    "**步骤2: 计算关联度打分**",
    "Q和K做点积运算，点积值越大关联度越强。",
    "10个词的句子会产生10x10的分数矩阵。",
    "**具体例子：**it的Query与cat的Key点积=8.5",
    "（高关联），与the的点积=1.2（低关联）。",
    "",
    "**步骤3: 缩放（关键！）**",
    "除以sqrt(64)=8。因为维度高时点积值会很大，",
    "导致Softmax进入饱和区。饱和区就是当分数差",
    "异太大时，Softmax会让最高分占据几乎全部注意",
    "力(比如99%)，其他词被完全忽略。除以sqrt(d_k)",
    "就像给分数做标准化，让注意力分布更均匀。",
    "",
    "**步骤4: Softmax转概率**",
    "Softmax函数把任意一组数字转换成0到1之间的",
    "概率值，且总和等于100%。比如分数[8.5, 1.2,",
    "0.3]变成概率[87%, 12%, 1%]。就像决定在每本",
    "书上花多少时间 -- 匹配度最高的花最多时间。",
    "",
    "**步骤5: 加权求和**",
    "按概率权重对所有V加权求和。cat权重最高，",
    "所以it的输出包含大量cat的语义信息 --",
    "it知道它指的是cat。每个词都融合了整个句子的上下文。",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第7页：Multi-Head Attention
# ============================================================
print("创建第7页：Multi-Head Attention...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "多头注意力：8个专家从不同角度同时分析",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：多头架构图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "多头注意力架构",
            font_size=13, bold=True, color=C_ACCENT)

# 输入向量
add_rounded_rect(slide, 2.0, 1.5, 3.5, 0.5, "输入向量 (512维)", fill_color=C_PRIMARY, font_size=12)

# 8个头（用4个代表）
head_colors = [C_ACCENT, C_ORANGE, C_GREEN, C_PURPLE, C_RED, C_PRIMARY, C_GRAY, RGBColor(0xE9, 0x1E, 0x63)]
head_labels = ["Head 1\n语法", "Head 2\n语义", "Head 3\n指代", "Head 4\n长距离"]

for i, (label, color) in enumerate(zip(head_labels, head_colors)):
    y = 2.3 + i * 0.65
    add_connector_line(slide, 3.75, 2.0, 1.0 + i * 1.3, y, color=color, width=Pt(1.5))
    add_rounded_rect(slide, 0.5 + i * 1.4, y, 1.2, 0.5, label, fill_color=color, font_size=8)

# v2修复：Concat标签下移到5.9，避开连线区域
add_textbox(slide, 0.5, 5.9, 6, 0.35, "Concat (拼接) -> 线性投影(混合整理8个头的结果) -> 输出",
            font_size=10, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

# 拼接框
add_rounded_rect(slide, 1.5, 6.35, 4.5, 0.45, "Concat + Linear Projection", 
                  fill_color=C_PRIMARY, font_size=11)

# 输出
add_rounded_rect(slide, 2.5, 6.95, 2.5, 0.35, "输出 (512维)", fill_color=C_GREEN, font_size=11)

# 箭头 —— v2修复：指向新的concat位置
for i in range(4):
    y = 2.3 + i * 0.65 + 0.5
    add_connector_line(slide, 0.5 + i * 1.4 + 0.6, y, 2.0, 6.35, color=head_colors[i], width=Pt(1))

add_connector_line(slide, 3.75, 6.8, 3.75, 6.95, color=C_PRIMARY, width=Pt(2))

# 右侧：文字说明
lines = [
    "**为什么需要多个头？**",
    "单个注意力头只能学习一种关注模式。",
    "就像分析一幅画，一个评论家只关注色彩，",
    "另一个关注构图，还有一个关注情感表达 --",
    "每个人看到的是同一幅画的不同侧面。",
    "",
    "**8个头 = 8种独立的分析维度**",
    "原始Transformer使用8个注意力头，",
    "每个头64维（512/8=64），有独立的Q/K/V矩阵。",
    "实验发现不同头学到了不同的语言现象：",
    "有的关注相邻词关系，有的关注句法结构，",
    "有的关注指代消解。这种分工是自动学会的。",
    "",
    "**关键亮点：计算量不变！**",
    "512维拆成8个64维，每个头独立计算，",
    "最后拼接起来。总计算量和单个512维完全相同！",
    "就像把512页书分8人同时各读一本，",
    "最后汇总笔记 -- 总阅读量一样，",
    "但效率更高、视角更全面。",
]
add_multiline_textbox(slide, 7.5, 1.0, 5.3, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第8页：Cross-Attention
# ============================================================
print("创建第8页：Cross-Attention...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "交叉注意力：编码器和解码器之间的翻译桥梁",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：对比图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "Self-Attention vs Cross-Attention",
            font_size=13, bold=True, color=C_ACCENT)

# Self-Attention
add_rounded_rect(slide, 0.5, 1.5, 6.0, 2.5, "", fill_color=None, line_color=C_ACCENT)
add_textbox(slide, 0.7, 1.6, 5.5, 0.4, "Self-Attention（自注意力）",
            font_size=13, bold=True, color=C_ACCENT)

sa_words = ["我", "爱", "北", "京", "天安门"]
for i, word in enumerate(sa_words):
    x = 1.0 + i * 1.0
    add_rounded_rect(slide, x, 2.1, 0.8, 0.4, word, fill_color=C_ACCENT, font_size=11)
    # 互相连线（简化：只画几条）
    for j in range(len(sa_words)):
        if i != j and abs(i-j) <= 2:
            add_connector_line(slide, x + 0.4, 2.5, 1.0 + j * 1.0 + 0.4, 2.1, 
                              color=RGBColor(0xBD, 0xC3, 0xC7), width=Pt(1))

add_textbox(slide, 0.7, 3.0, 5.5, 0.8, "Q、K、V 来自**同一个序列**\n词和词之间互相交流（团队内部开会）",
            font_size=10, color=C_TEXT)

# Cross-Attention
add_rounded_rect(slide, 0.5, 4.3, 6.0, 2.8, "", fill_color=None, line_color=C_GREEN)
add_textbox(slide, 0.7, 4.4, 5.5, 0.4, "Cross-Attention（交叉注意力）",
            font_size=13, bold=True, color=C_GREEN)

# 编码器（中文）
add_textbox(slide, 0.7, 4.9, 2.5, 0.3, "编码器输出（中文）：", font_size=10, bold=True, color=C_ACCENT)
enc_words = ["我", "爱", "北京"]
for i, word in enumerate(enc_words):
    add_rounded_rect(slide, 0.7 + i * 0.9, 5.3, 0.7, 0.4, word, fill_color=C_ACCENT, font_size=10)

# 解码器（英文）
add_textbox(slide, 3.8, 4.9, 2.5, 0.3, "解码器（英文）：", font_size=10, bold=True, color=C_GREEN)
dec_words = ["I", "love", "Beijing"]
for i, word in enumerate(dec_words):
    add_rounded_rect(slide, 3.8 + i * 1.0, 5.3, 0.8, 0.4, word, fill_color=C_GREEN, font_size=10)

# 跨界连线
add_connector_line(slide, 4.2, 5.7, 1.4, 5.3, color=C_RED, width=Pt(2))
add_connector_line(slide, 5.2, 5.7, 2.3, 5.3, color=C_RED, width=Pt(2))

add_textbox(slide, 0.7, 5.9, 5.5, 0.8, "Q 来自**解码器**，K、V 来自**编码器**\n英文词向中文词提问（拿问题清单去图书馆找答案）",
            font_size=10, color=C_TEXT)

# 右侧：文字说明 —— v2修复：加过渡句
lines = [
    "**到目前为止，我们讲的注意力都是在一个**",
    "**序列内部进行的。但翻译需要两个序列之间**",
    "**交流 -- 这就需要交叉注意力。**",
    "",
    "**Self-Attention vs Cross-Attention：**",
    "自注意力：Q、K、V 全部来自同一个序列，",
    "就像团队内部开会讨论。",
    "",
    "交叉注意力：Q 来自一个序列（解码器），",
    "K、V 来自另一个序列（编码器），",
    "就像你拿着问题清单去图书馆找答案。",
    "",
    "**翻译场景举例：**",
    "编码器先处理中文我 爱 北京 天安门，",
    "生成包含丰富语义的编码输出。",
    "",
    "解码器在生成英文时，每生成一个词都通过",
    "交叉注意力回头看编码器的输出。",
    "生成 Beijing 时，Query（我在找地名）",
    "与编码器中每个中文词的Key匹配，",
    "发现北京匹配度最高，提取其Value。",
    "",
    "**架构位置：**",
    "编码器内部：只有自注意力（中文词互交流）",
    "解码器内部：掩码自注意力 + 交叉注意力",
    "先自己讨论，再向对方请教！",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.0, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第9页：Position Encoding
# ============================================================
print("创建第9页：Position Encoding...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "位置编码：告诉 Transformer \"谁是第几个词\"",
            font_size=22, bold=True, color=C_PRIMARY)

# 上方：打乱顺序示意图
add_rounded_rect(slide, 0.5, 1.0, 5.5, 1.8, "", fill_color=C_LIGHT_BG, line_color=C_ACCENT)

# v2修复：缩短"原句："和"打乱："标签宽度，避免与字节点重叠
add_textbox(slide, 0.7, 1.15, 0.6, 0.3, "原句：", font_size=11, bold=True, color=C_TEXT)
for i, word in enumerate(["狗", "咬", "人"]):
    add_rounded_rect(slide, 1.4 + i * 0.8, 1.1, 0.6, 0.4, word, fill_color=C_ACCENT, font_size=11)

add_textbox(slide, 0.7, 1.65, 0.6, 0.3, "打乱：", font_size=11, bold=True, color=C_TEXT)
for i, word in enumerate(["人", "咬", "狗"]):
    add_rounded_rect(slide, 1.4 + i * 0.8, 1.6, 0.6, 0.4, word, fill_color=C_ORANGE, font_size=11)

# 等号和警告
add_textbox(slide, 3.8, 1.3, 2.0, 0.6, "= ???\n没有位置编码时，\nTransformer 分不清！",
            font_size=10, bold=True, color=C_RED)

# 解决方案
add_textbox(slide, 4.0, 2.0, 2.0, 0.5, "+ 位置编码\n-> 输出不同!",
            font_size=10, bold=True, color=C_GREEN)

# 下方：三种方法对比 —— v2修复：缩窄表格列宽，给右侧文字留空间
add_textbox(slide, 0.5, 3.1, 6, 0.4, "三种主流位置编码方法对比",
            font_size=13, bold=True, color=C_ACCENT)

# 表格（用形状模拟）—— v2修复：缩小列宽避免与右侧文字重叠
table_data = [
    ("方法", "原理", "优点", "缺点", "使用者"),
    ("正弦余弦", "sin/cos公式", "无需学习", "表达有限", "原始Transformer"),
    ("可学习编码", "当参数训练", "简单直观", "无法外推", "BERT/GPT-3"),
    ("RoPE旋转编码", "旋转矩阵", "相对位置", "稍复杂", "LLaMA等现代"),
]

for row_i, row in enumerate(table_data):
    y = 3.6 + row_i * 0.55
    col_widths = [1.1, 1.1, 1.1, 1.1, 1.2]
    x = 0.5
    for col_i, (cell, w) in enumerate(zip(row, col_widths)):
        color = C_PRIMARY if row_i == 0 else (C_LIGHT_BG if row_i % 2 == 1 else None)
        text_color = C_WHITE if row_i == 0 else C_TEXT
        font_s = 9 if row_i == 0 else 8
        shape = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.5, 
                           fill_color=color, line_color=C_GRAY, line_width=Pt(0.5))
        set_shape_text(shape, cell, font_size=font_s, color=text_color, bold=(row_i==0))
        x += w

# 右侧：文字说明 —— v2修复：加过渡句、RoPE直觉比喻、left从7.3改为7.8避免表格重叠
lines = [
    "**注意力机制讲完了，但它有一个致命缺陷：**",
    "**它不知道词的先后顺序。**",
    "",
    "**为什么需要位置编码？**",
    "自注意力本质上是排列不变的 -- 打乱词的顺序，",
    "计算结果不会变。狗咬人和人咬狗包含",
    "完全相同的词，没有位置编码就无法区分。",
    "",
    "**三种方法：**",
    "1. 正弦余弦（2017原始论文）：用数学公式算，",
    "   无需学习，可外推，但表达能力有限。",
    "",
    "2. 可学习编码（BERT/GPT）：当参数训练，",
    "   简单直观，但无法处理超过训练长度的序列。",
    "",
    "3. RoPE旋转编码（现代主流）：通过旋转矩阵",
    "   将位置信息注入Q和K。",
    "   **RoPE直觉比喻：**想象把每个词向量当作时钟",
    "   上的指针，不同位置的词旋转不同角度，距离越",
    "   远的词旋转角度差异越大，这样注意力计算自然",
    "   包含了词之间的距离信息。",
    "",
    "**为什么RoPE成为主流？**",
    "两个核心优势：",
    "1. 自动编码相对位置 -- 注意力只取决于词间距",
    "2. 支持长序列外推 -- 可处理128K+上下文",
    "",
    "目前90%以上的现代大模型都用RoPE：",
    "LLaMA、Mistral、Qwen、DeepSeek...",
]
add_multiline_textbox(slide, 7.8, 1.0, 5.2, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第10页：Transformer 整体架构
# ============================================================
print("创建第10页：整体架构...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "完整架构：从输入到输出，数据如何流动",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：架构图（编码器-解码器）
add_textbox(slide, 0.3, 1.0, 7.5, 0.4, "Encoder-Decoder 架构",
            font_size=13, bold=True, color=C_ACCENT)

# 输入
add_rounded_rect(slide, 0.3, 1.5, 2.5, 0.5, "输入嵌入 + 位置编码", fill_color=C_PRIMARY, font_size=10)

# 编码器（2层代表6层）
add_textbox(slide, 0.3, 2.2, 2.5, 0.3, "编码器 (x6层)", font_size=11, bold=True, color=C_ACCENT)
for i in range(2):
    y = 2.5 + i * 1.1
    add_rounded_rect(slide, 0.3, y, 2.5, 0.4, "多头自注意力", fill_color=C_ACCENT, font_size=10)
    add_rounded_rect(slide, 0.3, y + 0.45, 2.5, 0.4, "FFN + Add & Norm", fill_color=C_GREEN, font_size=10)
    if i < 1:
        add_connector_line(slide, 1.55, y + 0.85, 1.55, y + 1.1, color=C_GRAY, width=Pt(1))
add_textbox(slide, 0.3, 4.8, 2.5, 0.3, "...（共6层）", font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# 交叉注意力连接线
add_connector_line(slide, 2.8, 3.0, 4.0, 3.0, color=C_RED, width=Pt(2))
add_textbox(slide, 3.0, 2.7, 1.0, 0.3, "交叉注意力", font_size=8, color=C_RED, align=PP_ALIGN.CENTER)

# 解码器（2层代表6层）
add_textbox(slide, 4.0, 2.2, 3.0, 0.3, "解码器 (x6层)", font_size=11, bold=True, color=C_GREEN)
for i in range(2):
    y = 2.5 + i * 1.5
    add_rounded_rect(slide, 4.0, y, 3.0, 0.35, "掩码自注意力", fill_color=C_GREEN, font_size=9)
    add_rounded_rect(slide, 4.0, y + 0.4, 3.0, 0.35, "交叉注意力", fill_color=C_RED, font_size=9)
    add_rounded_rect(slide, 4.0, y + 0.8, 3.0, 0.35, "FFN + Add & Norm", fill_color=C_ORANGE, font_size=9)
    if i < 1:
        add_connector_line(slide, 5.5, y + 1.15, 5.5, y + 1.5, color=C_GRAY, width=Pt(1))
add_textbox(slide, 4.0, 5.7, 3.0, 0.3, "...（共6层）", font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# 输出嵌入
add_rounded_rect(slide, 4.0, 6.1, 3.0, 0.4, "输出嵌入 + 位置编码", fill_color=C_PRIMARY, font_size=10)

# 最终输出
add_rounded_rect(slide, 4.0, 6.7, 1.2, 0.4, "Linear", fill_color=C_PURPLE, font_size=10)
add_rounded_rect(slide, 5.5, 6.7, 1.2, 0.4, "Softmax", fill_color=C_PURPLE, font_size=10)

# 输入输出标签 —— v2修复：输出标签左移避免与右侧文字重叠
add_rounded_rect(slide, 0.3, 6.5, 2.5, 0.5, "输入: 我爱北京天安门", fill_color=C_ACCENT, font_size=10)
add_textbox(slide, 3.2, 6.6, 0.5, 0.3, "->", font_size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
add_rounded_rect(slide, 4.0, 6.5, 3.5, 0.5, "输出: I love Beijing...", fill_color=C_GREEN, font_size=10)

# 右侧：文字说明
lines = [
    "**宏观结构：左边理解，右边生成**",
    "编码器（Encoder）：负责理解输入内容",
    "解码器（Decoder）：负责生成输出内容",
    "两者通过交叉注意力连接。",
    "",
    "**数据流动过程：**",
    "1. 输入文字 -> 分词 -> 词向量 + 位置编码",
    "2. 编码器6层：每层包含多头自注意力 + FFN",
    "3. 编码器输出 -> 交叉注意力 -> 解码器",
    "4. 解码器6层：掩码自注意力 + 交叉注意力 + FFN",
    "5. 线性层 + Softmax -> 输出概率 -> 选最高概率词",
    "",
    "**原始配置（2017年论文）：**",
    "6层编码器 + 6层解码器",
    "8个注意力头，512维",
    "前馈网络2048维",
    "总共约6500万个参数",
    "",
    "**今天呢？**",
    "核心架构没变，但规模爆发式增长：",
    "GPT-4据传1.8万亿参数 -- 增长了近3万倍！",
    "但本质上还是同一个Encoder-Decoder结构。",
]
add_multiline_textbox(slide, 8.0, 1.0, 5.0, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第11页：Encoder 编码器详解
# ============================================================
print("创建第11页：编码器详解...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "编码器详解：信息如何被理解和提炼",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：编码器单层结构图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "编码器单层内部结构",
            font_size=13, bold=True, color=C_ACCENT)

# 流程：输入 → 注意力 → Add&Norm → FFN → Add&Norm → 输出
components = [
    ("输入向量", C_PRIMARY, 1.5),
    ("多头自注意力", C_ACCENT, 2.3),
    ("残差连接 + 层归一化", C_GREEN, 3.1),
    ("前馈网络 (FFN)", C_ORANGE, 3.9),
    ("残差连接 + 层归一化", C_GREEN, 4.7),
    ("输出向量", C_PRIMARY, 5.5),
]

for label, color, y in components:
    add_rounded_rect(slide, 1.0, y, 3.0, 0.5, label, fill_color=color, font_size=11)
    if y < 5.5:
        add_connector_line(slide, 2.5, y + 0.5, 2.5, y + 0.6, color=C_GRAY, width=Pt(2))

# 残差连接旁路
add_connector_line(slide, 4.0, 1.75, 4.5, 2.55, color=C_RED, width=Pt(2))
add_textbox(slide, 4.5, 2.8, 1.5, 0.4, "高速公路匝道\n(残差连接)", font_size=8, color=C_RED)

add_connector_line(slide, 4.0, 3.35, 4.5, 4.15, color=C_RED, width=Pt(2))
add_textbox(slide, 4.5, 4.4, 1.5, 0.4, "标准化体检\n(层归一化)", font_size=8, color=C_RED)

add_textbox(slide, 4.5, 3.95, 1.5, 0.4, "独立思考\n(前馈网络)", font_size=8, color=C_ORANGE)

# 6层堆叠标注
add_textbox(slide, 1.0, 6.2, 3.0, 0.3, "上述结构重复 6 次",
            font_size=10, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

# 右侧：文字说明
lines = [
    "**编码器单层：两个子模块 + 两个安全网**",
    "输入经过 注意力->安全网->FFN->安全网 四步，",
    "传给下一层。6层后每个词都融合了整个句子的深层语义。",
    "",
    "**1. 多头自注意力**",
    "让每个词与所有其他词交流信息。",
    "",
    "**2. 残差连接 -- 高速公路匝道**",
    "公式：输出 = 输入 + 子模块(输入)",
    "子模块的输出会和原始输入相加，不是替换。",
    "就像高速公路旁修匝道 -- 主干道堵了（信息衰减），",
    "信息可以直接从匝道通过。解决了梯度消失问题，",
    "让训练几十层甚至上百层网络成为可能。",
    "",
    "**3. 层归一化 -- 标准化体检**",
    "对每个词向量做标准化：减均值、除标准差。",
    "确保数值范围一致，不会某些维度特别大。",
    "让训练更稳定、加速收敛。",
    "",
    "**4. 前馈网络 -- 独立思考时间**",
    "对每个词独立应用两次线性变换 + 激活函数。",
    "隐藏维度是模型的4倍（512->2048->512），",
    "计算量比注意力层还大。",
    "注意力收集信息，FFN加工提炼。",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第12页：Decoder 解码器详解
# ============================================================
print("创建第12页：解码器详解...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "解码器详解：信息如何被生成和表达",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：解码器单层结构图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "解码器单层内部结构（比编码器多一层）",
            font_size=13, bold=True, color=C_ACCENT)

dec_components = [
    ("输入向量", C_PRIMARY, 1.5),
    ("掩码多头自注意力", C_ACCENT, 2.2),
    ("残差连接 + 层归一化", C_GREEN, 2.9),
    ("交叉注意力 (查阅编码器)", C_RED, 3.6),
    ("残差连接 + 层归一化", C_GREEN, 4.3),
    ("前馈网络 (FFN)", C_ORANGE, 5.0),
    ("残差连接 + 层归一化", C_GREEN, 5.7),
    ("输出向量", C_PRIMARY, 6.4),
]

for label, color, y in dec_components:
    add_rounded_rect(slide, 0.5, y, 3.5, 0.45, label, fill_color=color, font_size=10)
    if y < 6.4:
        add_connector_line(slide, 2.25, y + 0.45, 2.25, y + 0.55, color=C_GRAY, width=Pt(1.5))

# v2修复：标注"因果掩码"加括号解释，回指第8页
add_textbox(slide, 4.2, 2.0, 2.8, 0.5, "考试不能偷看后面的题\n(因果掩码，也叫掩码自注意力，\n第8页提到过，确保只能看到之前的词)", font_size=8, color=C_ACCENT)
add_textbox(slide, 4.2, 3.4, 2.5, 0.5, "翻译时回头查看原文\n(Q来自解码器, K/V来自编码器)", font_size=9, color=C_RED)

# 自回归生成示意
add_textbox(slide, 0.5, 6.9, 7, 0.3, "自回归生成：[START]->I->[START,I]->love->[START,I,love]->you->...",
            font_size=9, bold=True, color=C_GREEN)

# 右侧：文字说明
lines = [
    "**解码器三层结构（编码器只有两层）：**",
    "1. 掩码自注意力（已生成的词之间交流）",
    "2. 交叉注意力（向编码器提问）",
    "3. 前馈网络（独立思考）",
    "每层后有残差连接和层归一化。",
    "",
    "**1. 掩码机制 -- 考试不能偷看后面的题**",
    "计算注意力时，将当前位置之后的所有分数",
    "设为负无穷大，Softmax后变为0。确保每个词",
    "只能关注之前的词，不能偷看未来。",
    "",
    "**2. 交叉注意力 -- 翻译时回头查看原文**",
    "Q来自解码器（我正在生成的词需要什么信息），",
    "K、V来自编码器（输入句子的完整理解）。",
    "生成Beijing时，Query与编码器中每个中文词",
    "的Key匹配，发现北京匹配最高，提取其Value。",
    "",
    "**3. 自回归生成 -- 一个词一个词地说**",
    "先生成第1个词，加入输入，生成第2个词......",
    "直到生成结束符号。就像说话必须先想好第一个词。",
    "这也意味着推理速度是主要瓶颈 -- 无法完全并行。",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第13页：训练与推理
# ============================================================
print("创建第13页：训练与推理...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "训练与推理：模型是怎么学会说话的？",
            font_size=22, bold=True, color=C_PRIMARY)

# 左侧：对比图
# 训练
add_rounded_rect(slide, 0.5, 1.0, 6.0, 2.8, "", fill_color=C_LIGHT_BG, line_color=C_ACCENT)
add_textbox(slide, 0.7, 1.1, 5.5, 0.4, "训练阶段 = 开卷考试（Teacher Forcing）",
            font_size=12, bold=True, color=C_ACCENT)

# 训练输入输出
train_input = ["[START]", "我", "爱", "北京"]
train_target = ["我", "爱", "北京", "天安门"]

for i, (inp, tgt) in enumerate(zip(train_input, train_target)):
    y = 1.7 + i * 0.45
    add_rounded_rect(slide, 0.8, y, 0.8, 0.35, inp, fill_color=C_ACCENT, font_size=9)
    add_textbox(slide, 1.8, y, 0.4, 0.35, "->", font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, 2.3, y, 0.8, 0.35, tgt, fill_color=C_GREEN, font_size=9)
    add_textbox(slide, 3.3, y, 0.4, 0.35, "OK", font_size=12, color=C_GREEN, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.7, 3.5, 5.5, 0.3, "模型直接看到正确答案，所有位置同时计算",
            font_size=9, color=C_TEXT)

# 推理
add_rounded_rect(slide, 0.5, 4.1, 6.0, 3.0, "", fill_color=None, line_color=C_RED)
add_textbox(slide, 0.7, 4.2, 5.5, 0.4, "推理阶段 = 闭卷考试（自回归生成）",
            font_size=12, bold=True, color=C_RED)

infer_steps = [
    ("[START]", "我"),
    ("[START]我", "爱"),
    ("[START]我爱", "北京"),
    ("[START]我爱北京", "天安门"),
]

for i, (inp, out) in enumerate(infer_steps):
    y = 4.8 + i * 0.5
    add_textbox(slide, 0.8, y, 1.8, 0.35, f"输入: {inp}", font_size=8, color=C_TEXT)
    add_textbox(slide, 2.8, y, 0.4, 0.35, "->", font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.3, y, 1.5, 0.35, f"输出: {out}", font_size=8, color=C_GREEN)

add_textbox(slide, 0.7, 6.8, 5.5, 0.3, "模型靠自己预测，一个词一个词生成，不能回头修改",
            font_size=9, color=C_TEXT)

# 右侧：文字说明
lines = [
    "**训练阶段 -- Teacher Forcing：**",
    "模型不需要自己预测，而是直接看到正确答案。",
    "训练数据的输入是完整目标序列（向左移一位），",
    "模型根据前面的词预测下一个词，和正确答案对比，",
    "计算误差并更新参数。",
    "",
    "就像做练习题时，每做完一道老师立刻告诉你答案 --",
    "可以快速纠正错误、加速学习。",
    "",
    "**推理阶段 -- 闭卷考试：**",
    "模型不再有正确答案参考。必须完全靠自己之前",
    "生成的词来预测下一个词。生成过程是自回归的：",
    "输入起始符号->预测第1个词->加入输入->预测第2个词。",
    "",
    "每步生成无法回头修改。这就是为什么有时AI会",
    "说错话 -- 已经在前面说了错词，后面只能将错就错。",
    "",
    "**关键差异：**",
    "训练：看到完整目标序列，并行计算（效率高）",
    "推理：只能看到已生成部分，串行逐词（效率低）",
    "训练成本固定（一次性），推理成本持续（每次使用）",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第14页：Transformer 的演进与应用
# ============================================================
print("创建第14页：演进与应用...")
slide = prs.slides.add_slide(blank_layout)

# v2修复：标题"8年"改为"9年"
add_textbox(slide, 0.5, 0.3, 12, 0.6, "9年演进：从一篇论文到改变世界",
            font_size=22, bold=True, color=C_PRIMARY)

# 三大家族图
add_textbox(slide, 0.5, 1.0, 6, 0.4, "三大家族：同一个DNA，三种职业方向",
            font_size=13, bold=True, color=C_ACCENT)

# 起始节点
add_rounded_rect(slide, 2.8, 1.5, 1.8, 0.5, "Transformer\n2017", fill_color=C_PRIMARY, font_size=10)

# Encoder-Only
add_connector_line(slide, 2.8, 1.75, 0.5, 2.5, color=C_ACCENT, width=Pt(2))
add_rounded_rect(slide, 0.3, 2.5, 2.0, 0.4, "Encoder-Only", fill_color=C_ACCENT, font_size=10)
add_textbox(slide, 0.3, 3.0, 2.0, 0.3, "BERT -> RoBERTa -> DeBERTa", font_size=8, color=C_TEXT)
add_textbox(slide, 0.3, 3.3, 2.0, 0.5, "擅长：理解\n搜索、分类、情感分析", font_size=8, color=C_GREEN)

# Encoder-Decoder
add_connector_line(slide, 3.7, 2.0, 3.7, 2.5, color=C_ORANGE, width=Pt(2))
add_rounded_rect(slide, 2.8, 2.5, 2.0, 0.4, "Encoder-Decoder", fill_color=C_ORANGE, font_size=10)
add_textbox(slide, 2.8, 3.0, 2.0, 0.3, "T5 -> BART -> Whisper", font_size=8, color=C_TEXT)
add_textbox(slide, 2.8, 3.3, 2.0, 0.5, "擅长：转换\n翻译、摘要、语音识别", font_size=8, color=C_GREEN)

# Decoder-Only
add_connector_line(slide, 4.6, 1.75, 5.5, 2.5, color=C_GREEN, width=Pt(2))
add_rounded_rect(slide, 5.0, 2.5, 2.0, 0.4, "Decoder-Only", fill_color=C_GREEN, font_size=10)
add_textbox(slide, 5.0, 3.0, 2.0, 0.3, "GPT-1->2->3->4", font_size=8, color=C_TEXT)
add_textbox(slide, 5.0, 3.3, 2.0, 0.5, "擅长：生成\n对话、写作、代码", font_size=8, color=C_GREEN)

# 参数量增长
add_textbox(slide, 0.5, 4.1, 6, 0.4, "参数量爆发式增长",
            font_size=13, bold=True, color=C_ACCENT)

growth = [
    ("2017", "原始Transformer", "65M"),
    ("2018", "BERT-Large", "340M"),
    ("2019", "GPT-2 (1.5B)", "1.5B"),
    ("2020", "GPT-3", "175B"),
    ("2023", "GPT-4 (据传)", "1.8T"),
]

for i, (year, model, params) in enumerate(growth):
    y = 4.6 + i * 0.45
    bar_w = 0.5 + i * 0.9  # 递增
    add_textbox(slide, 0.5, y, 0.6, 0.3, year, font_size=9, color=C_GRAY)
    add_textbox(slide, 1.2, y, 1.8, 0.3, model, font_size=9, color=C_TEXT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 3.0, y + 0.05, min(bar_w, 3.5), 0.2, fill_color=C_ACCENT)
    add_textbox(slide, 3.0 + min(bar_w, 3.5) + 0.1, y, 0.8, 0.3, params, font_size=9, bold=True, color=C_ACCENT)

# 右侧：文字说明
lines = [
    "**1. 三大家族**",
    "Transformer诞生后，根据不同任务需求发展出三个方向：",
    "- Encoder-Only (BERT)：擅长理解类任务",
    "- Decoder-Only (GPT)：擅长生成类任务",
    "- Encoder-Decoder (T5)：擅长转换类任务",
    "",
    "**2. 规模爆发**",
    "核心架构没变，但规模爆炸式增长。",
    "从6500万到1.8万亿参数，增长了近3万倍。",
    "每一代增长都带来能力飞跃 -- 缩放定律：",
    "更多参数 + 更多数据 = 更强能力。",
    "",
    "**3. 超越语言**",
    "Transformer早已超越NLP领域：",
    "- ViT (2020)：图像切成小块当词处理",
    "- Whisper：语音识别",
    "- GPT-4V/Gemini：多模态理解",
    "- AlphaFold2：蛋白质结构预测",
    "",
    "Transformer正在成为AI的通用引擎。",
]
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 6.2, lines,
                       font_size=9, line_spacing=1.05)

# ============================================================
# 第15页：总结
# ============================================================
print("创建第15页：总结...")
slide = prs.slides.add_slide(blank_layout)

# 标题
add_textbox(slide, 0.5, 0.3, 12, 0.6, "总结：带走今天的5个核心要点",
            font_size=24, bold=True, color=C_PRIMARY)

# 左侧：5个要点图标
summaries = [
    ("1", "信息自由流动", "每个信息单元都能与\n所有其他单元直接交流", C_ACCENT),
    ("2", "自注意力 = 图书馆搜索", "Q搜索、K匹配、V提取\n融合整个句子的上下文", C_ORANGE),
    ("3", "多头 = 多视角分析", "8个独立头学习不同\n语言关系，计算量不变", C_GREEN),
    ("4", "位置编码 = 词序信息", "RoPE旋转编码是主流\n支持相对位置和长序列", C_PURPLE),
    ("5", "从翻译工具到AI通用引擎", "BERT/GPT/T5三大家族\n参数从65M增长到1.8T", C_RED),
]

for i, (num, title, desc, color) in enumerate(summaries):
    y = 1.0 + i * 1.2
    
    # 数字圆圈
    circle = add_shape(slide, MSO_SHAPE.OVAL, 0.5, y, 0.6, 0.6, fill_color=color)
    set_shape_text(circle, num, font_size=20, color=C_WHITE, bold=True)
    
    # 标题
    add_textbox(slide, 1.3, y, 5.0, 0.4, title,
                font_size=14, bold=True, color=color)
    
    # 描述
    add_textbox(slide, 1.3, y + 0.4, 5.0, 0.6, desc,
                font_size=10, color=C_TEXT)

# 右侧：详细回顾
lines = [
    "**要点1：Transformer的本质**",
    "让每个信息单元都能与所有其他信息单元",
    "直接交流。彻底解决RNN的三大痛点。",
    "",
    "**要点2：自注意力**",
    "每个词生成Q(搜索)、K(标签)、V(内容)，",
    "通过QxK^T计算关联度，Softmax转概率，",
    "对V加权求和。对所有词同时进行。",
    "",
    "**要点3：多头注意力**",
    "拆成8个独立头，学习不同语言关系。",
    "计算量与单头相同，但视角更全面。",
    "",
    "**要点4：位置编码**",
    "自注意力不知道词序，需要注入位置信息。",
    "RoPE是主流，编码相对位置，支持长序列。",
    "",
    "**要点5：从论文到改变世界**",
    "2017年至今，BERT(理解)、GPT(生成)、",
    "T5(转换)三大家族。从NLP到视觉、语音、",
    "多模态，Transformer成为AI通用引擎。",
    "",
    "**Transformer不是一个模型，而是一个时代。**",
]
# v2修复：右侧详细回顾高度缩小，给推荐阅读留空间
add_multiline_textbox(slide, 7.3, 1.0, 5.7, 5.4, lines,
                       font_size=9, line_spacing=1.05)

# v2修复：缩短感谢语宽度，避免与右侧面板重叠
add_textbox(slide, 0.5, 6.9, 7.0, 0.3, "感谢聆听！  |  Weaver  |  2026",
            font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# v2修复：加推荐阅读 —— 放在右侧面板下方，不与面板重叠
add_textbox(slide, 7.3, 6.5, 5.7, 0.6, 
            "推荐阅读：\n原文论文《Attention Is All You Need》\n入门教程：Jay Alammar的《The Illustrated Transformer》",
            font_size=8, color=C_GRAY)

# ============================================================
# 保存
# ============================================================
output_path = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v2.pptx"
prs.save(output_path)
print(f"\nPPT 已保存到: {output_path}")
print(f"总页数: {len(prs.slides)}")

# ============================================================
# 自检：边界和重叠检查
# ============================================================
print("\n" + "="*60)
print("自检报告：边界和重叠检查")
print("="*60)

for slide_idx, slide in enumerate(prs.slides, 1):
    shapes = list(slide.shapes)
    issues = []
    
    for i, shape in enumerate(shapes):
        # 获取边界
        left = shape.left / 914400  # EMU to inches
        top = shape.top / 914400
        width = shape.width / 914400
        height = shape.height / 914400
        
        # 边界检查
        if left + width > 13.2:
            issues.append(f"  Shape {i} ({shape.shape_type}): 右边界超出 {left+width:.2f} > 13.2")
        if top + height > 7.3:
            issues.append(f"  Shape {i} ({shape.shape_type}): 下边界超出 {top+height:.2f} > 7.3")
    
    if issues:
        print(f"\nSlide {slide_idx}: 发现 {len(issues)} 个问题")
        for issue in issues:
            print(issue)
    else:
        print(f"Slide {slide_idx}: 边界检查通过")

# 重叠检测（仅检查TextBox和有文字的Shape之间的重叠）
print("\n" + "-"*60)
print("重叠检测（仅检查可能的文字重叠）")
print("-"*60)

overlap_count = 0
for slide_idx, slide in enumerate(prs.slides, 1):
    shapes = list(slide.shapes)
    slide_overlaps = 0
    
    for i in range(len(shapes)):
        for j in range(i+1, len(shapes)):
            s1 = shapes[i]
            s2 = shapes[j]
            
            # 跳过connector（连线）
            if s1.shape_type == 6 or s2.shape_type == 6:  # MSO_SHAPE_TYPE.CONNECTOR
                continue
            # 跳过装饰线
            if s1.height / 914400 < 0.02 or s2.height / 914400 < 0.02:
                continue
            
            l1, t1 = s1.left / 914400, s1.top / 914400
            r1, b1 = l1 + s1.width / 914400, t1 + s1.height / 914400
            l2, t2 = s2.left / 914400, s2.top / 914400
            r2, b2 = l2 + s2.width / 914400, t2 + s2.height / 914400
            
            # 检查重叠
            overlap_h = max(0, min(r1, r2) - max(l1, l2))
            overlap_v = max(0, min(b1, b2) - max(t1, t2))
            overlap_area = overlap_h * overlap_v
            area1 = (r1 - l1) * (b1 - t1)
            area2 = (r2 - l2) * (b2 - t2)
            
            if overlap_area > 0.1:  # 显著重叠
                ratio = overlap_area / min(area1, area2)
                if ratio > 0.3:
                    slide_overlaps += 1
                    overlap_count += 1
                    # 获取名称
                    name1 = s1.name if hasattr(s1, 'name') else f"Shape_{i}"
                    name2 = s2.name if hasattr(s2, 'name') else f"Shape_{j}"
                    print(f"  Slide {slide_idx}: {name1} <-> {name2} (重叠比 {ratio:.0%})")
    
    if slide_overlaps == 0:
        print(f"Slide {slide_idx}: 无显著文字重叠")

print(f"\n总计重叠: {overlap_count} 处")
print("="*60)
print("自检完成！")
