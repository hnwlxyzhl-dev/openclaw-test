#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT v6 Part1 - Pages 1-7"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_TITLE   = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN     = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED   = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT    = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY    = RGBColor(0x95, 0xA5, 0xA6)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xED, 0xCB)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
C_LIGHT_YELLOW = RGBColor(0xFE, 0xF9, 0xE7)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY  = "Microsoft YaHei"
FONT_EN    = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ── Helpers ──
def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color, line_color, line_width)

def add_textbox(slide, left, top, width, height, text, font_size=Pt(10), font_color=C_TEXT,
                bold=False, align=PP_ALIGN.LEFT, font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    # Set East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', FONT_TITLE)
    rPr.append(ea)
    return txBox

def add_rich_textbox(slide, left, top, width, height, paragraphs_data, anchor=MSO_ANCHOR.TOP):
    """paragraphs_data: list of (text, font_size, font_color, bold, align, font_name, space_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, pdata in enumerate(paragraphs_data):
        text, font_size, font_color, bold, align, font_name = pdata[0], pdata[1], pdata[2], pdata[3], pdata[4], pdata[5]
        space_after = pdata[6] if len(pdata) > 6 else Pt(4)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {})
        ea.set('typeface', FONT_TITLE)
        rPr.append(ea)
    return txBox

def add_arrow_line(slide, x1, y1, x2, y2, color=C_ENCODER, width=Pt(1.5)):
    """Add a line with arrow from (x1,y1) to (x2,y2)"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    ln = connector.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)
    return connector

def add_line(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def set_shape_text(shape, text, font_size=Pt(9), font_color=C_TEXT, bold=False, align=PP_ALIGN.CENTER, font_name=FONT_BODY):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', FONT_TITLE)
    rPr.append(ea)

def add_page_title(slide, text, left=Inches(0.5), top=Inches(0.3)):
    add_textbox(slide, left, top, Inches(12), Inches(0.5), text,
                font_size=Pt(16), font_color=C_TITLE, bold=True)
    # Underline
    add_line(slide, left, top + Inches(0.45), left + Inches(12), top + Inches(0.45),
             color=C_TITLE, width=Pt(1.5))


# ═══════════════════════════════════════════════════
# PAGE 1: Cover
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Background gradient feel - light blue top bar
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.5), fill_color=C_TITLE)

# Decorative lines
add_line(slide, Inches(2), Inches(3.2), Inches(11.333), Inches(3.2), color=C_ENCODER, width=Pt(3))
add_line(slide, Inches(2), Inches(3.4), Inches(11.333), Inches(3.4), color=C_EMBED, width=Pt(2))
add_line(slide, Inches(2), Inches(3.6), Inches(11.333), Inches(3.6), color=C_ATTENTION, width=Pt(1.5))

# Main title
add_textbox(slide, Inches(1), Inches(0.8), Inches(11.333), Inches(1.2),
            "Transformer 架构深度解析",
            font_size=Pt(26), font_color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(0.5),
            "从直觉到原理，一篇看懂AI的核心引擎",
            font_size=Pt(14), font_color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

# Info bar
add_textbox(slide, Inches(2), Inches(3.9), Inches(9.333), Inches(0.4),
            "Attention Is All You Need — Vaswani et al., 2017  |  Google Brain  |  6500万参数",
            font_size=Pt(10), font_color=C_GRAY, align=PP_ALIGN.CENTER)

# Bottom text
add_textbox(slide, Inches(2), Inches(6.5), Inches(9.333), Inches(0.4),
            "从训练到推理，完整拆解Transformer的工作原理",
            font_size=Pt(10), font_color=C_TEXT, align=PP_ALIGN.CENTER)

# Decorative small boxes
for i, (c, x) in enumerate([(C_ENCODER, 3.5), (C_EMBED, 6.0), (C_ATTENTION, 8.5)]):
    add_rounded_rect(slide, Inches(x), Inches(5.8), Inches(1.2), Inches(0.35), fill_color=c)
    labels = ["编码器 Encoder", "嵌入 Embedding", "注意力 Attention"]
    set_shape_text(slide.shapes[-1], labels[i], font_size=Pt(9), font_color=C_WHITE, bold=True)


# ═══════════════════════════════════════════════════
# PAGE 2: Why Transformer
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "为什么需要Transformer：RNN的痛点与注意力突破")

# ── Left: Diagrams (8 inches wide) ──
# RNN flow diagram (top half)
add_textbox(slide, Inches(0.3), Inches(0.9), Inches(4), Inches(0.3),
            "RNN顺序处理", font_size=Pt(11), font_color=C_DECODER, bold=True)

words_rnn = ["The", "cat", "sat", "on", "the", "mat"]
box_w = Inches(0.9)
box_h = Inches(0.45)
start_x = Inches(0.5)
y_rnn = Inches(1.25)
gap = Inches(0.35)

for i, word in enumerate(words_rnn):
    x = start_x + i * (box_w + gap)
    s = add_rounded_rect(slide, x, y_rnn, box_w, box_h, fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1))
    set_shape_text(s, word, font_size=Pt(10), font_color=C_ENCODER, bold=True)
    # Red arrow to next
    if i < len(words_rnn) - 1:
        ax1 = x + box_w
        ay = y_rnn + box_h // 2
        ax2 = x + box_w + gap
        add_arrow_line(slide, ax1, ay, ax2, ay, color=C_DECODER, width=Pt(1.5))

# h_t labels
for i in range(len(words_rnn)):
    x = start_x + i * (box_w + gap) + box_w // 2
    add_textbox(slide, x - Inches(0.2), y_rnn + box_h + Inches(0.05), Inches(0.5), Inches(0.25),
                "h_t", font_size=Pt(9), font_color=C_DECODER, bold=True, align=PP_ALIGN.CENTER)

# Red annotation
add_textbox(slide, Inches(6.5), Inches(1.3), Inches(1.8), Inches(0.4),
            "顺序处理瓶颈", font_size=Pt(10), font_color=C_DECODER, bold=True, align=PP_ALIGN.CENTER)

# Attention diagram (bottom half)
add_textbox(slide, Inches(0.3), Inches(2.3), Inches(4), Inches(0.3),
            "自注意力机制", font_size=Pt(11), font_color=C_ENCODER, bold=True)

words_att = ["The", "cat", "sat", "on", "the", "mat", ".", "it", "was", "."]
att_w = Inches(0.75)
att_h = Inches(0.35)
att_start_x = Inches(0.4)
y_att = Inches(3.0)
att_gap = Inches(0.18)

it_idx = words_att.index("it")
att_weights = {
    0: 0.05, 1: 0.60, 2: 0.05, 3: 0.02, 4: 0.03, 5: 0.20, 6: 0.01, 7: 1.0, 8: 0.04, 9: 0.01
}

for i, word in enumerate(words_att):
    x = att_start_x + i * (att_w + att_gap)
    if word == "it":
        fill = C_ENCODER
        fc = C_WHITE
    else:
        fill = C_LIGHT_BLUE
        fc = C_ENCODER
    s = add_rounded_rect(slide, x, y_att, att_w, att_h, fill_color=fill, line_color=C_ENCODER, line_width=Pt(1))
    set_shape_text(s, word, font_size=Pt(9), font_color=fc, bold=True)

# Draw attention lines from "it" to others
it_x = att_start_x + it_idx * (att_w + att_gap) + att_w // 2
it_y = y_att + att_h

for i, word in enumerate(words_att):
    if i == it_idx:
        continue
    target_x = att_start_x + i * (att_w + att_gap) + att_w // 2
    w = att_weights.get(i, 0.01)
    if w >= 0.1:
        lw = Pt(3)
        c = C_ENCODER
    elif w >= 0.05:
        lw = Pt(2)
        c = C_ATTENTION
    else:
        lw = Pt(1)
        c = C_GRAY
    add_line(slide, it_x, it_y, target_x, it_y + Inches(1.0), color=c, width=lw)
    if w >= 0.05:
        add_textbox(slide, target_x - Inches(0.2), it_y + Inches(1.05), Inches(0.5), Inches(0.2),
                    str(w), font_size=Pt(9), font_color=C_TEXT, align=PP_ALIGN.CENTER)

# "it" label
add_textbox(slide, it_x - Inches(0.2), it_y + Inches(1.05), Inches(0.5), Inches(0.2),
            "1.0", font_size=Pt(9), font_color=C_ENCODER, bold=True, align=PP_ALIGN.CENTER)

# Weight labels
add_textbox(slide, Inches(0.3), Inches(5.0), Inches(8), Inches(0.3),
            "连线粗细 = 注意力权重（\"it\"关注\"cat\"权重最高0.60）", font_size=Pt(9), font_color=C_GRAY)

# ── Right: 5 bullet points ──
right_x = Inches(8.5)
right_w = Inches(4.5)
y_start = Inches(0.85)

bullets_p2 = [
    ("1. 顺序处理的效率陷阱", "循环神经网络(RNN)必须按顺序处理，每个词依赖前一个隐藏状态h_t。即使GPU有数千个并行单元，RNN也只能串行使用。处理1000词需要1000步，像雇1000人搬砖却每次只能1个人搬。LSTM虽引入门控缓解梯度消失，但顺序限制未变。"),
    ("2. 长距离遗忘（梯度消失）", "处理到句末时，早期信息经数百次矩阵乘法被指数级稀释。每步h_t=tanh(W*h_{t-1}+U*x_t)，W特征值小于1时衰减极快，100步后可能只剩10^(-30)的权重，几乎完全丢失。"),
    ("3. 信息瓶颈", "整个序列被压缩成固定长度向量（256或512维），再生成输出。像用100字摘要概括500页书再回答细节问题——信息瓶颈直接限制了处理长文本和理解复杂语义的能力。"),
    ("4. 核心突破——自注意力", "Transformer让每个词与所有其他词直接交流，像\"圆桌会议\"而非\"排队发言\"。以\"it\"为例，自注意力让它与\"cat\"建立最强关联（权重0.6），同时解决：并行化解决顺序瓶颈，直接连接解决长距离遗忘，不压缩解决信息瓶颈。"),
    ("5. 数学本质", "核心是Q*K^T（查询*键转置），本质是高维空间余弦相似度。点积越大语义关联越强。一次性计算所有词对关联度，形成T*T注意力矩阵，完全并行，无需顺序处理。"),
]

y = y_start
for title, body in bullets_p2:
    add_textbox(slide, right_x, y, right_w, Inches(0.25), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.25)
    add_textbox(slide, right_x, y, right_w, Inches(1.0), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(0.9)


# ═══════════════════════════════════════════════════
# PAGE 3: Architecture Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "架构总览：编码器-解码器结构")

# ── Left: Architecture diagram ──
lx = Inches(0.3)
# Input embedding + positional encoding
add_rounded_rect(slide, lx, Inches(1.0), Inches(1.6), Inches(0.55), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(slide.shapes[-1], "输入嵌入\n37000x512", font_size=Pt(9), font_color=C_ENCODER, bold=True)

add_textbox(slide, lx + Inches(1.6), Inches(1.1), Inches(0.3), Inches(0.3),
            "+", font_size=Pt(14), font_color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)

add_rounded_rect(slide, lx + Inches(1.9), Inches(1.0), Inches(1.5), Inches(0.55), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1.5))
set_shape_text(slide.shapes[-1], "位置编码\nsin/cos", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43), bold=True)

# Encoder block
enc_x = lx
enc_y = Inches(1.8)
enc_w = Inches(3.8)
enc_h = Inches(4.8)

# Encoder outer frame
enc_frame = add_rounded_rect(slide, enc_x, enc_y, enc_w, enc_h, fill_color=None, line_color=C_ENCODER, line_width=Pt(2.5))

add_textbox(slide, enc_x + Inches(0.1), enc_y - Inches(0.3), Inches(2.5), Inches(0.3),
            "编码器 Encoder x6", font_size=Pt(11), font_color=C_ENCODER, bold=True)

# Inside encoder: one detailed layer
ly = enc_y + Inches(0.3)
lh = Inches(0.45)

# Multi-head self-attention
s = add_rounded_rect(slide, enc_x + Inches(0.3), ly, enc_w - Inches(0.6), lh, fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1))
set_shape_text(s, "多头自注意力 Multi-Head Attention (8头x64维)", font_size=Pt(9), font_color=C_ENCODER, bold=True)

# Add/Normalize
ly += Inches(0.55)
s = add_rounded_rect(slide, enc_x + Inches(0.3), ly, enc_w - Inches(0.6), Inches(0.3), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差连接 + 层归一化 (Add & Norm)", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

# FFN
ly += Inches(0.4)
s = add_rounded_rect(slide, enc_x + Inches(0.3), ly, enc_w - Inches(0.6), lh, fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(s, "前馈网络 FFN (512->2048->ReLU->512)", font_size=Pt(9), font_color=RGBColor(0xB7, 0x7D, 0x0B), bold=True)

# Add/Normalize
ly += Inches(0.55)
s = add_rounded_rect(slide, enc_x + Inches(0.3), ly, enc_w - Inches(0.6), Inches(0.3), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差连接 + 层归一化 (Add & Norm)", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

# Repeat indicator
ly += Inches(0.5)
add_textbox(slide, enc_x + Inches(0.5), ly, Inches(3), Inches(0.3),
            "... 重复6次 ...", font_size=Pt(10), font_color=C_ENCODER, bold=True, align=PP_ALIGN.CENTER)

# Dotted lines for repetition
for dy in [Inches(0.35), Inches(0.6)]:
    add_line(slide, enc_x + Inches(0.5), ly + dy, enc_x + enc_w - Inches(0.5), ly + dy,
             color=C_ENCODER, width=Pt(0.75))

# Dimension annotations
add_textbox(slide, enc_x - Inches(0.1), enc_y + Inches(0.4), Inches(0.5), Inches(0.2),
            "512维", font_size=Pt(9), font_color=C_GRAY, align=PP_ALIGN.RIGHT)

# ── Decoder block ──
dec_x = Inches(4.8)
dec_y = Inches(1.8)
dec_w = Inches(3.8)
dec_h = Inches(4.8)

dec_frame = add_rounded_rect(slide, dec_x, dec_y, dec_w, dec_h, fill_color=None, line_color=C_DECODER, line_width=Pt(2.5))

add_textbox(slide, dec_x + Inches(0.1), dec_y - Inches(0.3), Inches(2.5), Inches(0.3),
            "解码器 Decoder x6", font_size=Pt(11), font_color=C_DECODER, bold=True)

# Inside decoder
dly = dec_y + Inches(0.2)
# Masked self-attention
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.4), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1))
set_shape_text(s, "掩码自注意力 Masked Self-Attention", font_size=Pt(9), font_color=C_ENCODER, bold=True)

dly += Inches(0.5)
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.25), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

# Cross attention (purple)
dly += Inches(0.35)
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.4), fill_color=C_LIGHT_PURPLE, line_color=C_EMBED, line_width=Pt(1.5))
set_shape_text(s, "交叉注意力 Cross-Attention (K,V来自编码器)", font_size=Pt(9), font_color=C_EMBED, bold=True)

dly += Inches(0.5)
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.25), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

# FFN
dly += Inches(0.35)
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.4), fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(s, "前馈网络 FFN (512->2048->ReLU->512)", font_size=Pt(9), font_color=RGBColor(0xB7, 0x7D, 0x0B), bold=True)

dly += Inches(0.5)
s = add_rounded_rect(slide, dec_x + Inches(0.3), dly, dec_w - Inches(0.6), Inches(0.25), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

dly += Inches(0.35)
add_textbox(slide, dec_x + Inches(0.5), dly, Inches(3), Inches(0.3),
            "... 重复6次 ...", font_size=Pt(10), font_color=C_DECODER, bold=True, align=PP_ALIGN.CENTER)

for dy in [Inches(0.3), Inches(0.55)]:
    add_line(slide, dec_x + Inches(0.5), dly + dy, dec_x + dec_w - Inches(0.5), dly + dy,
             color=C_DECODER, width=Pt(0.75))

# K,V arrows from encoder to decoder cross-attention
enc_out_x = enc_x + enc_w
cross_y = dec_y + Inches(0.95)

add_arrow_line(slide, enc_out_x, enc_y + Inches(0.5), dec_x, cross_y, color=C_EMBED, width=Pt(2))
add_textbox(slide, enc_out_x + Inches(0.05), enc_y + Inches(0.25), Inches(0.5), Inches(0.25),
            "K, V", font_size=Pt(9), font_color=C_EMBED, bold=True)

# Output area
out_y = Inches(6.85)
add_rounded_rect(slide, dec_x + Inches(0.3), out_y, Inches(1.5), Inches(0.35), fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(slide.shapes[-1], "线性变换 512x37000", font_size=Pt(9), font_color=RGBColor(0xB7, 0x7D, 0x0B))

add_arrow_line(slide, dec_x + Inches(1.85), out_y + Inches(0.17), dec_x + Inches(2.0), out_y + Inches(0.17), color=C_ATTENTION, width=Pt(1))

add_rounded_rect(slide, dec_x + Inches(2.0), out_y, Inches(0.8), Inches(0.35), fill_color=C_LIGHT_RED, line_color=C_DECODER, line_width=Pt(1))
set_shape_text(slide.shapes[-1], "Softmax", font_size=Pt(9), font_color=C_DECODER, bold=True)

add_arrow_line(slide, dec_x + Inches(2.85), out_y + Inches(0.17), dec_x + Inches(3.0), out_y + Inches(0.17), color=C_DECODER, width=Pt(1))

add_rounded_rect(slide, dec_x + Inches(3.0), out_y, Inches(0.8), Inches(0.35), fill_color=C_LIGHT_YELLOW, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(slide.shapes[-1], "输出概率", font_size=Pt(9), font_color=C_TEXT, bold=True)

# ── Right: 4 bullet points ──
rx = Inches(8.9)
rw = Inches(4.1)

bullets_p3 = [
    ("1. 编码器=编辑部", "编码器的任务是深度理解输入文本，由6个相同的层堆叠而成。每层包含两个子层——多头自注意力和前馈网络，每个子层都有残差连接和层归一化。数据流经编码器时，每个词的信息会与其他所有词不断融合，6层后每个词的512维向量已蕴含整个句子的上下文信息。编码器参数约3450万。"),
    ("2. 解码器=翻译部", "解码器负责逐步生成输出，也有6层，但比编码器多一个交叉注意力(Cross-Attention)子层。掩码自注意力确保只能看到之前的词（因果掩码），交叉注意力让解码器能够\"查阅\"编码器的输出——带着问题去查资料。"),
    ("3. 数据流全景", "输入文本经过BPE子词分词 -> 嵌入矩阵查表（37000x512，将词映射为512维向量）-> 位置编码（sin/cos函数注入位置信息）-> 编码器6层处理（每层都在做\"交流+提炼\"）-> 编码器输出 -> 解码器6层处理（每层多一步\"查资料\"）-> 线性变换（512维->37000维）-> Softmax（转为概率分布）-> 取概率最高的词作为输出。"),
    ("4. 关键数字", "编码器6层、解码器6层、模型维度512、8个注意力头、前馈网络2048维、参数总量6500万、训练数据WMT 2014英德翻译450万句对、8块P100 GPU训练12小时、BLEU得分28.4（当时最高）。"),
]

y = Inches(0.9)
for title, body in bullets_p3:
    add_textbox(slide, rx, y, rw, Inches(0.22), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.22)
    add_textbox(slide, rx, y, rw, Inches(1.35), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(1.2)


# ═══════════════════════════════════════════════════
# PAGE 4: Encoder Detail
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "编码器单层：交流与提炼")

# ── Left: Data flow diagram ──
lx = Inches(0.3)
fw = Inches(2.8)
fh = Inches(0.5)
fy = Inches(1.0)
arrow_len = Inches(0.6)
fx = lx

# Input X
s = add_rounded_rect(slide, fx, fy, fw, fh, fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(s, "输入 X (512维)", font_size=Pt(10), font_color=C_ENCODER, bold=True)

# Arrow down
fy += fh + Inches(0.15)
add_arrow_line(slide, fx + fw // 2, fy - Inches(0.15), fx + fw // 2, fy + Inches(0.05), color=C_ENCODER, width=Pt(1.5))

# Multi-head attention
fy += Inches(0.05)
s = add_rounded_rect(slide, fx, fy, fw, fh, fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(s, "多头自注意力\n8头x64维", font_size=Pt(10), font_color=C_ENCODER, bold=True)

# Residual connection (bypass line) - draw curved path as line to the right
res_x = fx + fw + Inches(0.2)
add_line(slide, res_x, fy + fh // 2, res_x, fy + fh + Inches(0.7) + Inches(0.25) // 2, color=C_FFN, width=Pt(2))
add_textbox(slide, res_x + Inches(0.05), fy + Inches(0.15), Inches(0.8), Inches(0.3),
            "残差连接", font_size=Pt(9), font_color=C_FFN, bold=True)

# Arrow down
fy += fh + Inches(0.15)
add_arrow_line(slide, fx + fw // 2, fy - Inches(0.15), fx + fw // 2, fy + Inches(0.05), color=C_ENCODER, width=Pt(1.5))

# Add & Norm
fy += Inches(0.05)
s = add_rounded_rect(slide, fx, fy, fw, Inches(0.35), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1.5))
set_shape_text(s, "Add & Norm (512维)", font_size=Pt(10), font_color=RGBColor(0x1B, 0x7A, 0x43), bold=True)

# Arrow down
fy += Inches(0.35) + Inches(0.15)
add_arrow_line(slide, fx + fw // 2, fy - Inches(0.15), fx + fw // 2, fy + Inches(0.05), color=C_ENCODER, width=Pt(1.5))

# FFN
fy += Inches(0.05)
s = add_rounded_rect(slide, fx, fy, fw, fh, fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1.5))
set_shape_text(s, "前馈网络 FFN\n512->2048->ReLU->512", font_size=Pt(10), font_color=RGBColor(0xB7, 0x7D, 0x0B), bold=True)

# Residual connection 2
add_line(slide, res_x, fy + fh // 2, res_x, fy + fh + Inches(0.7) + Inches(0.25) // 2, color=C_FFN, width=Pt(2))
add_textbox(slide, res_x + Inches(0.05), fy + Inches(0.15), Inches(0.8), Inches(0.3),
            "残差连接", font_size=Pt(9), font_color=C_FFN, bold=True)

# Arrow down
fy += fh + Inches(0.15)
add_arrow_line(slide, fx + fw // 2, fy - Inches(0.15), fx + fw // 2, fy + Inches(0.05), color=C_ENCODER, width=Pt(1.5))

# Add & Norm
fy += Inches(0.05)
s = add_rounded_rect(slide, fx, fy, fw, Inches(0.35), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1.5))
set_shape_text(s, "Add & Norm (512维)", font_size=Pt(10), font_color=RGBColor(0x1B, 0x7A, 0x43), bold=True)

# Arrow down
fy += Inches(0.35) + Inches(0.15)
add_arrow_line(slide, fx + fw // 2, fy - Inches(0.15), fx + fw // 2, fy + Inches(0.05), color=C_ENCODER, width=Pt(1.5))

# Output
fy += Inches(0.05)
s = add_rounded_rect(slide, fx, fy, fw, Inches(0.4), fill_color=C_LIGHT_PURPLE, line_color=C_EMBED, line_width=Pt(1.5))
set_shape_text(s, "输出 (512维)", font_size=Pt(10), font_color=C_EMBED, bold=True)

# Plus signs at Add&Norm positions
plus_positions = [
    (lx + fw // 2, Inches(1.0) + fh + Inches(0.15) + Inches(0.05) + fh + Inches(0.15) + Inches(0.05) + Inches(0.17)),
    (lx + fw // 2, fy - Inches(0.35) - Inches(0.15) - Inches(0.05) - fh - Inches(0.15) - Inches(0.05) - Inches(0.17)),
]

# ── Right: 4 bullet points ──
rx = Inches(8.3)
rw = Inches(4.7)

bullets_p4 = [
    ("1. 多头自注意力=8个讨论组", "8个独立注意力头学习不同关系模式：头1专注语法（主谓）、头2关注语义相似、头3关注位置、头4关注指代消解。每头独立计算Q、K、V（各64维），得到8个注意力矩阵，最后拼接（8x64=512维）经线性变换融合。"),
    ("2. 残差连接=高速公路匝道", "深度网络中梯度经链式法则反向传播，层数多时指数级衰减。残差连接将输入直接加到输出（output=input+sublayer(input)），为梯度提供\"匝道\"——即使中间权重很小，梯度也能通过加法路径无损传播。加号后的\"1\"确保梯度永不消失。"),
    ("3. 层归一化=标准化体检", "对每个词512维向量做标准化：减均值、除标准差，再经可学习缩放和偏移。公式：y=gamma*(x-mu)/sigma+beta。没有归一化时数值范围可能从[-0.01,0.01]突变到[-100,100]，导致训练不稳定。gamma和beta各512维，每层1024个可学习参数。"),
    ("4. 前馈网络=独立思考空间", "两层全连接：FFN(x)=ReLU(xW1+b1)W2+b2。第一层512维扩展到2048维（4倍），第二层压缩回512维，ReLU激活函数引入非线性。参数约210万（512x2048x2），比自注意力大2倍，是计算量最大的组件。注意力负责词间交流，FFN负责词自身特征提取。"),
]

y = Inches(0.85)
for title, body in bullets_p4:
    add_textbox(slide, rx, y, rw, Inches(0.22), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.22)
    add_textbox(slide, rx, y, rw, Inches(1.35), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(1.3)


# ═══════════════════════════════════════════════════
# PAGE 5: Decoder Detail
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "解码器单层：多一座翻译桥梁")

# ── Left top: Decoder structure ──
lx = Inches(0.3)
fw = Inches(2.8)
fh = Inches(0.45)
fy = Inches(1.0)

# Masked self-attention
s = add_rounded_rect(slide, lx, fy, fw, fh, fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(s, "掩码自注意力\nMasked Self-Attention", font_size=Pt(10), font_color=C_ENCODER, bold=True)

fy += fh + Inches(0.1)
s = add_rounded_rect(slide, lx, fy, fw, Inches(0.3), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

fy += Inches(0.35)
# Cross attention (purple)
s = add_rounded_rect(slide, lx, fy, fw, fh, fill_color=C_LIGHT_PURPLE, line_color=C_EMBED, line_width=Pt(1.5))
set_shape_text(s, "交叉注意力\nCross-Attention", font_size=Pt(10), font_color=C_EMBED, bold=True)

# Encoder output arrow pointing to cross-attention
add_arrow_line(slide, Inches(-0.1), fy + fh // 2, lx, fy + fh // 2, color=C_EMBED, width=Pt(2))
add_textbox(slide, Inches(-0.1), fy - Inches(0.3), Inches(1.5), Inches(0.3),
            "编码器输出(K,V)", font_size=Pt(9), font_color=C_EMBED, bold=True)

fy += fh + Inches(0.1)
s = add_rounded_rect(slide, lx, fy, fw, Inches(0.3), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

fy += Inches(0.35)
s = add_rounded_rect(slide, lx, fy, fw, fh, fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1.5))
set_shape_text(s, "前馈网络 FFN\n512->2048->ReLU->512", font_size=Pt(10), font_color=RGBColor(0xB7, 0x7D, 0x0B), bold=True)

fy += fh + Inches(0.1)
s = add_rounded_rect(slide, lx, fy, fw, Inches(0.3), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(s, "残差 + 归一化", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

# Arrows between blocks
for ay in [Inches(1.45), Inches(1.95), Inches(2.55), Inches(3.05)]:
    add_arrow_line(slide, lx + fw // 2, ay, lx + fw // 2, ay + Inches(0.1), color=C_ENCODER, width=Pt(1.5))

# ── Left bottom: Causal mask matrix ──
mask_x = Inches(0.5)
mask_y = Inches(4.3)
mask_cell = Inches(0.6)
mask_labels = ["START", "我", "爱", "你"]

add_textbox(slide, mask_x, mask_y - Inches(0.3), Inches(3), Inches(0.3),
            "因果掩码矩阵 Causal Mask", font_size=Pt(11), font_color=C_DECODER, bold=True)

# Column headers
for j, label in enumerate(mask_labels):
    x = mask_x + Inches(0.7) + j * mask_cell
    add_textbox(slide, x, mask_y, mask_cell, Inches(0.25), label,
                font_size=Pt(9), font_color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)

# Row labels + cells
for i, label in enumerate(mask_labels):
    y = mask_y + Inches(0.25) + i * mask_cell
    add_textbox(slide, mask_x, y, Inches(0.7), mask_cell, label,
                font_size=Pt(9), font_color=C_TEXT, bold=True, align=PP_ALIGN.RIGHT)
    for j in range(len(mask_labels)):
        x = mask_x + Inches(0.7) + j * mask_cell
        if j <= i:
            fill = C_LIGHT_BLUE
            txt = "可见"
            fc = C_ENCODER
        else:
            fill = C_LIGHT_RED
            txt = "-inf"
            fc = C_DECODER
        s = add_rounded_rect(slide, x, y, mask_cell - Inches(0.02), mask_cell - Inches(0.02),
                             fill_color=fill, line_color=C_GRAY, line_width=Pt(0.5))
        set_shape_text(s, txt, font_size=Pt(9), font_color=fc, bold=True)

add_textbox(slide, mask_x, mask_y + Inches(0.25) + 4 * mask_cell + Inches(0.1), Inches(3.5), Inches(0.3),
            "下三角=可见，上三角=-inf(Softmax后=0)", font_size=Pt(9), font_color=C_GRAY)

# ── Right: 3 bullet points ──
rx = Inches(8.3)
rw = Inches(4.7)

bullets_p5 = [
    ("1. 掩码自注意力=考试不能偷看后面的题", "解码器使用因果掩码(Causal Mask)，将注意力分数上三角设为负无穷(-inf)，Softmax后变0。处理\"爱\"时只能看到[START, 我, 爱]，看不到\"你\"。训练和推理都严格执行——训练时若能\"偷看\"答案，推理时没有答案可抄就会崩溃，避免暴露偏差(Exposure Bias)。"),
    ("2. 交叉注意力=带着问题去查资料", "解码器独有，编码器和解码器间的唯一\"桥梁\"。查询Q来自解码器（\"需要从原文找信息\"），键K和值V来自编码器（\"对原文的深度理解\"）。计算：Q*K^T->除以sqrt(d_k)->Softmax->*V。如生成\"爱\"时，Query编码\"需要情感动词\"，发现\"love\"的Key匹配最高（权重0.7），大量提取其Value。Q和K序列长度可以不同。"),
    ("3. 完整层结构与编码器对比", "解码器每层3个子层（掩码自注意力->交叉注意力->前馈网络），每个后都有残差连接和层归一化。编码器每层仅2个子层。多一个交叉注意力是两者最核心的区别。6层解码器共18个子层+18个残差归一化操作。"),
]

y = Inches(0.85)
for title, body in bullets_p5:
    add_textbox(slide, rx, y, rw, Inches(0.22), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.22)
    add_textbox(slide, rx, y, rw, Inches(1.7), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(1.65)


# ═══════════════════════════════════════════════════
# PAGE 6: Training
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "训练阶段：开卷考试")

# ── Top: Training data flow ──
top_y = Inches(0.85)
# Input
s = add_rounded_rect(slide, Inches(0.3), top_y, Inches(1.3), Inches(0.4), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1))
set_shape_text(s, "I love you", font_size=Pt(10), font_color=C_ENCODER, bold=True)

add_arrow_line(slide, Inches(1.65), top_y + Inches(0.2), Inches(2.0), top_y + Inches(0.2), color=C_ENCODER, width=Pt(1.5))

# Encoder x6
s = add_rounded_rect(slide, Inches(2.0), top_y, Inches(1.5), Inches(0.4), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(s, "编码器 x6", font_size=Pt(10), font_color=C_ENCODER, bold=True)

add_arrow_line(slide, Inches(3.55), top_y + Inches(0.2), Inches(3.9), top_y + Inches(0.2), color=C_ENCODER, width=Pt(1.5))

# Encoder output
s = add_rounded_rect(slide, Inches(3.9), top_y, Inches(1.5), Inches(0.4), fill_color=C_LIGHT_PURPLE, line_color=C_EMBED, line_width=Pt(1.5))
set_shape_text(s, "编码器输出\n(K,V向量)", font_size=Pt(9), font_color=C_EMBED, bold=True)

# Decoder input
add_rounded_rect(slide, Inches(5.8), top_y, Inches(1.6), Inches(0.4), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1))
set_shape_text(slide.shapes[-1], "解码器输入\n[START, 我, 爱]", font_size=Pt(9), font_color=RGBColor(0x1B, 0x7A, 0x43))

add_arrow_line(slide, Inches(7.45), top_y + Inches(0.2), Inches(7.8), top_y + Inches(0.2), color=C_FFN, width=Pt(1.5))

# Decoder x6
s = add_rounded_rect(slide, Inches(7.8), top_y, Inches(1.5), Inches(0.4), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1.5))
set_shape_text(s, "解码器 x6", font_size=Pt(10), font_color=RGBColor(0x1B, 0x7A, 0x43), bold=True)

add_arrow_line(slide, Inches(9.35), top_y + Inches(0.2), Inches(9.7), top_y + Inches(0.2), color=C_FFN, width=Pt(1.5))

# Linear + Softmax
s = add_rounded_rect(slide, Inches(9.7), top_y, Inches(1.3), Inches(0.4), fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(s, "Linear+Softmax", font_size=Pt(9), font_color=RGBColor(0xB7, 0x7D, 0x0B))

add_arrow_line(slide, Inches(11.05), top_y + Inches(0.2), Inches(11.4), top_y + Inches(0.2), color=C_ATTENTION, width=Pt(1.5))

# Predicted
s = add_rounded_rect(slide, Inches(11.4), top_y, Inches(1.2), Inches(0.4), fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION, line_width=Pt(1.5))
set_shape_text(s, "预测值", font_size=Pt(10), font_color=RGBColor(0xB7, 0x7D, 0x0B), bold=True)

# True target (red box below predicted)
true_y = top_y + Inches(0.55)
s = add_rounded_rect(slide, Inches(11.4), true_y, Inches(1.2), Inches(0.4), fill_color=C_LIGHT_RED, line_color=C_DECODER, line_width=Pt(1.5))
set_shape_text(s, "真实目标\n[我,爱,你,END]", font_size=Pt(9), font_color=C_DECODER, bold=True)

# Loss label
add_textbox(slide, Inches(11.4), true_y + Inches(0.45), Inches(1.5), Inches(0.3),
            "<-> 交叉熵损失", font_size=Pt(9), font_color=C_DECODER, bold=True)

# Bottom bracket text
add_textbox(slide, Inches(5.5), top_y + Inches(0.55), Inches(3), Inches(0.3),
            "所有位置同时并行计算（教师强制策略）", font_size=Pt(9), font_color=C_FFN, bold=True, align=PP_ALIGN.CENTER)

# K,V arrow from encoder output to decoder
add_arrow_line(slide, Inches(4.65), top_y + Inches(0.4), Inches(6.6), top_y + Inches(0.4), color=C_EMBED, width=Pt(1.5))

# ── Bottom left: Embedding + Positional encoding diagram ──
embed_y = Inches(2.5)

add_textbox(slide, Inches(0.3), embed_y, Inches(3), Inches(0.3),
            "嵌入 + 位置编码", font_size=Pt(11), font_color=C_ENCODER, bold=True)

# Word tokens
words = ["我", "爱", "你"]
for i, w in enumerate(words):
    x = Inches(0.5) + i * Inches(1.2)
    s = add_rounded_rect(slide, x, embed_y + Inches(0.35), Inches(0.7), Inches(0.35), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1))
    set_shape_text(s, w, font_size=Pt(10), font_color=C_ENCODER, bold=True)

# Arrow down to embedding matrix
add_arrow_line(slide, Inches(2.0), embed_y + Inches(0.7), Inches(2.0), embed_y + Inches(0.9), color=C_ENCODER, width=Pt(1.5))

# Embedding matrix
s = add_rounded_rect(slide, Inches(0.5), embed_y + Inches(0.9), Inches(3.0), Inches(1.0), fill_color=C_LIGHT_BLUE, line_color=C_ENCODER, line_width=Pt(1.5))
set_shape_text(s, "嵌入矩阵 Embedding\n37000 x 512\n约1890万参数", font_size=Pt(10), font_color=C_ENCODER, bold=True)

# Arrow down
add_arrow_line(slide, Inches(2.0), embed_y + Inches(1.9), Inches(2.0), embed_y + Inches(2.1), color=C_ENCODER, width=Pt(1.5))

# Positional encoding
s = add_rounded_rect(slide, Inches(0.5), embed_y + Inches(2.1), Inches(3.0), Inches(0.8), fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(1.5))
set_shape_text(s, "位置编码 Positional Encoding\nPE(pos,2i) = sin(pos/10000^(2i/d))", font_size=Pt(10), font_color=RGBColor(0x1B, 0x7A, 0x43), bold=True)

# Plus sign and output
add_textbox(slide, Inches(3.6), embed_y + Inches(2.3), Inches(0.5), Inches(0.4),
            "⊕", font_size=Pt(18), font_color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)

add_arrow_line(slide, Inches(2.0), embed_y + Inches(2.9), Inches(2.0), embed_y + Inches(3.1), color=C_FFN, width=Pt(1.5))

s = add_rounded_rect(slide, Inches(1.0), embed_y + Inches(3.1), Inches(2.0), Inches(0.4), fill_color=C_LIGHT_YELLOW, line_color=C_ATTENTION, line_width=Pt(1))
set_shape_text(s, "最终输入向量 (512维)", font_size=Pt(9), font_color=C_TEXT, bold=True)

# ── Bottom right: 4 bullet points ──
rx = Inches(6.5)
rw = Inches(6.3)

bullets_p6 = [
    ("1. 教师强制(Teacher Forcing)", "训练时解码器输入是真实目标序列而非模型预测。对于[我,爱,你,END]，解码器同时看到[START,我,爱,你]，并行预测每位置下一个词。像开卷考试——每题都能看前面标准答案。所有位置完全并行计算，GPU充分并行处理。"),
    ("2. 嵌入矩阵=超级字典", "BPE子词分词将句子切成最小单元。嵌入矩阵E大小37000x512，约1890万参数，占总量30%。每个词对应512维向量，编码语义——\"猫\"和\"狗\"的向量接近（都是动物），\"国王\"和\"女王\"的差异方向与\"男人\"和\"女人\"类似，由大量数据训练自动学到。"),
    ("3. 位置编码=给每个词贴座位号", "自注意力本身无位置概念。位置编码用sin/cos注入：PE(pos,2i)=sin(pos/10000^(2i/d))，pos是位置索引，d=512。不同位置编码唯一且能表示相对位置，编码是固定的（不可学习），不增加可训练参数。"),
    ("4. 训练目标与损失", "目标是预测下一个词：根据[START]预测\"我\"，[START,我]预测\"爱\"。交叉熵Loss=-log(p(正确词))取平均，反向传播更新6500万参数。初期Loss约10.5，训练后降到1.5-3.0。配置：Adam优化器，学习率0.0003，dropout=0.3，标签平滑0.1，8块P100训练12小时。"),
]

y = Inches(2.3)
for title, body in bullets_p6:
    add_textbox(slide, rx, y, rw, Inches(0.22), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.22)
    add_textbox(slide, rx, y, rw, Inches(1.05), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(1.05)


# ═══════════════════════════════════════════════════
# PAGE 7: Self-Attention 5 Steps
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_page_title(slide, "自注意力的5步计算：从输入到输出")

# ── Left: 5-step vertical flow ──
lx = Inches(0.3)
fw = Inches(3.0)
fh = Inches(0.8)
fy = Inches(1.0)
step_gap = Inches(0.25)

steps = [
    ("Step 1: 生成Q,K,V", "输入X -> W_Q/W_K/W_V -> Q,K,V", C_LIGHT_BLUE, C_ENCODER,
     "所有词同时计算，完全并行"),
    ("Step 2: 计算分数矩阵", "Q x K^T -> Score矩阵 (TxT)", C_LIGHT_BLUE, C_ENCODER,
     "本质=高维余弦相似度"),
    ("Step 3: 缩放", "Score / sqrt(d_k) = Score / 8", C_LIGHT_YELLOW, RGBColor(0xB7, 0x7D, 0x0B),
     "sqrt(64)=8, 防止梯度饱和"),
    ("Step 4: Softmax归一化", "Softmax(Score) -> 概率矩阵", C_LIGHT_GREEN, RGBColor(0x1B, 0x7A, 0x43),
     "每行和=1"),
    ("Step 5: 加权求和输出", "概率 x V -> Output向量", C_LIGHT_PURPLE, C_EMBED,
     "融合所有词的信息"),
]

for i, (title, desc, fill, fc, note) in enumerate(steps):
    s = add_rounded_rect(slide, lx, fy, fw, fh, fill_color=fill, line_color=fc, line_width=Pt(1.5))
    tf = s.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(10)
    r.font.color.rgb = fc
    r.font.bold = True
    r.font.name = FONT_BODY
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', FONT_TITLE)
    rPr.append(ea)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_TEXT
    r2.font.name = FONT_BODY
    rPr2 = r2._r.get_or_add_rPr()
    ea2 = rPr2.makeelement(qn('a:ea'), {})
    ea2.set('typeface', FONT_TITLE)
    rPr2.append(ea2)

    # Note to the right
    add_textbox(slide, lx + fw + Inches(0.15), fy + Inches(0.15), Inches(3.5), Inches(0.5),
                note, font_size=Pt(9), font_color=C_GRAY)

    # Arrow to next step
    if i < len(steps) - 1:
        next_fy = fy + fh
        add_arrow_line(slide, lx + fw // 2, next_fy, lx + fw // 2, next_fy + step_gap,
                       color=fc, width=Pt(1.5))
        fy = next_fy + step_gap
    else:
        fy += fh

# Small 3x3 attention matrix illustration next to Step 2
mat_x = Inches(4.2)
mat_y = Inches(1.0) + fh + step_gap + Inches(0.5)
mat_cell = Inches(0.4)
mat_labels = ["it", "cat", "sat"]
mat_values = [
    [1.0, 0.60, 0.05],
    [0.10, 1.0, 0.08],
    [0.03, 0.15, 1.0],
]
# This would overlap with notes, let's skip the matrix to keep clean layout

# ── Right: 5 bullet points ──
rx = Inches(7.8)
rw = Inches(5.2)

bullets_p7 = [
    ("1. 生成查询Q、键K、值V", "每个词的512维向量X分别乘以W_Q、W_K、W_V（各512x64），得64维的Q、K、V。Q编码\"找什么信息\"（\"it\"的Q编码\"需找名词指代\"），K编码\"有什么特征\"（\"cat\"的K编码\"名词、动物\"），V编码\"实际内容\"。所有词同时计算，完全并行。"),
    ("2. 计算注意力分数矩阵", "Q与K转置相乘得TxT矩阵Score=QK^T。Score[i,j]表示词i对词j的关注度。Score[\"it\",\"cat\"]可能=12.5（高度关注），Score[\"it\",\"the\"]可能=2.1（低）。本质是高维空间余弦相似度计算。"),
    ("3. 缩放", "分数除以sqrt(d_k)=sqrt(64)=8。d_k较大时点积方差为d_k，Softmax进入饱和区（梯度近0）。除以sqrt(d_k)将方差归一化为1，确保梯度不消失。像百分制转十分制，保持相对大小避免极端值。"),
    ("4. Softmax归一化", "每行做Softmax：e^(x_i)/Sum(e^(x_j))，行和=1。\"it\"的分布可能[0.02,0.60,0.05,0.01,0.03,0.29]，\"cat\"最高(0.60)。Softmax放大最大值但不完全置零其他，使每个词都能获得一些信息。"),
    ("5. 加权求和输出", "Output=权重xV，每个词输出是所有Value的加权平均。\"it\"：输出=0.02xV(\"The\")+0.60xV(\"cat\")+0.05xV(\"sat\")+0.01xV(\"on\")+0.03xV(\"the\")+0.29xV(\"mat\")。\"it\"融合了\"cat\"大量信息。8头独立计算，拼接后线性变换输出512维。"),
]

y = Inches(0.85)
for title, body in bullets_p7:
    add_textbox(slide, rx, y, rw, Inches(0.22), title,
                font_size=Pt(11), font_color=C_TITLE, bold=True)
    y += Inches(0.22)
    add_textbox(slide, rx, y, rw, Inches(0.95), body,
                font_size=Pt(10), font_color=C_TEXT)
    y += Inches(0.95)


# ═══════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v6_part1.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")

# ── Word count per page ──
import re

def count_chinese(text):
    """Count Chinese chars + English words"""
    chinese = len(re.findall(r'[\u4e00-\u9fff]', text))
    # Count English words as tokens
    english_words = re.findall(r'[a-zA-Z]+', text)
    return chinese + len(english_words)

pages_bullets = {
    "Page 2": bullets_p2,
    "Page 3": bullets_p3,
    "Page 4": bullets_p4,
    "Page 5": bullets_p5,
    "Page 6": bullets_p6,
    "Page 7": bullets_p7,
}

print("\n=== Word Count per page (target: 300-400) ===")
for page_name, bullets in pages_bullets.items():
    total = 0
    for title, body in bullets:
        total += count_chinese(title) + count_chinese(body)
    print(f"{page_name}: {total} chars/words")
    status = "OK" if 300 <= total <= 400 else "WARN"
    if status == "WARN":
        print(f"  -> {status}: {'UNDER' if total < 300 else 'OVER'} target range")
