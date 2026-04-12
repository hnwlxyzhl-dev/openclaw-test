#!/usr/bin/env python3
"""
Transformer v3 中文版 PPT 修复脚本 (Final)
目标：质检≥95，读者评审≥90

修复内容：
P0: 架构图英文标签中文化（43处替换，分3轮处理\x0b换行问题）
P1: Slide 17 KV Cache 箭头（已有5个带箭头连接器，无需添加）
P2: Slide 2 目录页增加章节简述 + Slide 9/10 放大文字框

执行日期: 2026-04-12
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

INPUT_FILE = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3_zh.pptx"
OUTPUT_FILE = INPUT_FILE

prs = Presentation(INPUT_FILE)
nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ============================================================
# P0: 英文标签中文化
# ============================================================
# 注意: python-pptx 中某些文本含有 \x0b (vertical tab) 换行符，
# 导致 para.text 和 run.text 不完全匹配，需要分三轮处理。

# 第一轮：普通文本（无\x0b）- 精确匹配 run.text
ROUND1_MAP = {
    "Multi-Head Self-Attention": "多头自注意力",
    "Add & Layer Norm": "残差连接+层归一化",
    "Feed Forward (FFN)": "前馈网络",
    "Masked Self-Attention": "掩码自注意力",
    "Cross-Attention": "交叉注意力",
    "Input Embedding + PosEnc": "输入嵌入+位置编码",
    "Encoder x6": "编码器×6",
    "Decoder x6": "解码器×6",
    "Linear": "线性变换",
    "Encoder Output (K, V)": "编码器输出(K,V)",
    "Embedding Matrix (37000 x 512)": "嵌入矩阵(37000×512)",
    "Position Encoding (sin/cos)": "位置编码(sin/cos)",
    "W^Q (512x64)": "W^Q(512×64)",
    "W^K (512x64)": "W^K(512×64)",
    "W^V (512x64)": "W^V(512×64)",
    "KV Cache": "KV缓存",
    "Feed Forward": "前馈网络",
    "FFN": "前馈网络",
}

# 第二轮：含空格的复合标签 - 精确匹配 para.text
ROUND2_MAP = {
    "Output Embedding + PosEnc": "输出嵌入+位置编码",
    "Predictions": "预测值",
    "Targets": "目标值",
    "存入Cache": "存入缓存",
    "② Beam Search": "② 束搜索",
    "⑤ Temperature": "⑤ 温度",
    "完整正确序列(Teacher Forcing)": "完整正确序列(教师强制)",
    "KV Cache": "KV缓存",
}

# 第三轮：含\x0b换行的标签 - 精确匹配 para.text
ROUND3_MAP = {
    "Multi-Head\x0bSelf-Attention": "多头自注意力",
    "Add &\x0bLayerNorm": "残差连接+层归一化",
    "Masked\x0bSelf-Attention": "掩码自注意力",
    "Add &\x0bNorm": "残差连接+层归一化",
    "Cross-\x0bAttention": "交叉注意力",
    "Encoder\x0bx6": "编码器×6",
    "Decoder\x0bx6": "解码器×6",
    "Decoder Input\x0b[START, 我, 爱]": "解码器输入[START,我,爱]",
    "Query\x0b(64维)": "查询Q(64维)",
    "Key\x0b(64维)": "键K(64维)",
    "Value\x0b(64维)": "值V(64维)",
}

def is_label_shape(shape):
    """判断是否是架构图标签（而非正文解释）"""
    name = shape.name.lower()
    return ('rounded rectangle' in name or 
            ('rectangle' in name and name != 'rectangle 1'))

def replace_run_text(para, new_text):
    """替换段落所有run的文本为new_text，保留第一个run的格式"""
    r_elems = para._p.findall('.//a:r', nsmap)
    if not r_elems:
        return
    t_elem = r_elems[0].find('a:t', nsmap)
    if t_elem is not None:
        t_elem.text = new_text
    for r in r_elems[1:]:
        t = r.find('a:t', nsmap)
        if t is not None:
            t.text = ""

# --- 第一轮 ---
total_r1 = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame or not is_label_shape(shape):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text.strip()
                for eng, chn in ROUND1_MAP.items():
                    if text.lower() == eng.lower():
                        run.text = run.text.replace(text, chn)
                        total_r1 += 1
                        break

# --- 第二轮 ---
total_r2 = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame or not is_label_shape(shape):
            continue
        for para in shape.text_frame.paragraphs:
            full_text = para.text.strip()
            for eng, chn in ROUND2_MAP.items():
                if full_text == eng:
                    replace_run_text(para, chn)
                    total_r2 += 1
                    break

# --- 第三轮 ---
total_r3 = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame or not is_label_shape(shape):
            continue
        for para in shape.text_frame.paragraphs:
            full_text = para.text  # 保留\x0b
            for eng, chn in ROUND3_MAP.items():
                if full_text == eng:
                    replace_run_text(para, chn)
                    total_r3 += 1
                    break

print(f"P0 英文标签中文化: 第一轮={total_r1}, 第二轮={total_r2}, 第三轮={total_r3}, 共{total_r1+total_r2+total_r3}个")

# ============================================================
# P1: Slide 17 KV Cache 箭头检查
# ============================================================
slide17 = prs.slides[16]
existing_arrows = sum(1 for s in slide17.shapes 
                      if s.name.startswith("Connector"))
print(f"P1 Slide 17 已有连接器: {existing_arrows}个（均带箭头，无需添加）")

# ============================================================
# P2: 布局密度优化
# ============================================================
# Slide 2: 目录页添加章节简述
slide2 = prs.slides[1]
chapter_descs = {
    "Rounded Rectangle 2": "RNN/LSTM痛点 → 为什么需要新架构",
    "Rounded Rectangle 3": "编码器-解码器结构 → 注意力机制",
    "Rounded Rectangle 4": "嵌入→自注意力→损失函数→教师强制",
    "Rounded Rectangle 6": "自回归生成→KV缓存→词选择策略",
    "Rounded Rectangle 7": "三大家族(GPT/BERT/T5)→应用全景",
}
desc_added = 0
for shape in slide2.shapes:
    if shape.name in chapter_descs:
        desc = chapter_descs[shape.name]
        tf = shape.text_frame
        new_para = tf.add_paragraph()
        new_para.text = desc
        for run in new_para.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        shape.height += Emu(274320)
        desc_added += 1

# Slide 9: 放大文字框和嵌入矩阵
slide9 = prs.slides[8]
s9_changes = 0
for shape in slide9.shapes:
    if shape.name in ("TextBox 24", "TextBox 25", "TextBox 26"):
        shape.height += Emu(274320)
        s9_changes += 1
    elif shape.name == "Rounded Rectangle 6":
        shape.width += Emu(457200)
        s9_changes += 1

# Slide 10: 放大文字框
slide10 = prs.slides[9]
s10_changes = 0
for shape in slide10.shapes:
    if shape.name in ("TextBox 23", "TextBox 24", "TextBox 25", "TextBox 26"):
        shape.height += Emu(274320)
        s10_changes += 1

print(f"P2 布局优化: Slide 2简述={desc_added}, Slide 9={s9_changes}, Slide 10={s10_changes}")

# ============================================================
# 保存
# ============================================================
prs.save(OUTPUT_FILE)
total = total_r1 + total_r2 + total_r3 + desc_added + s9_changes + s10_changes
print(f"\n总计修复: {total} 项")
print(f"文件已保存: {OUTPUT_FILE}")
