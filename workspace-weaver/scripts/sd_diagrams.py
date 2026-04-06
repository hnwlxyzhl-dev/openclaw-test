#!/usr/bin/env python3
"""
Stable Diffusion 架构图 & 流程图 — 4页精美图表
用 python-pptx 的 MSO_SHAPE 绘制：
  A: SD 整体架构图（推理流程全景）
  B: U-Net 架构图（编码器-瓶颈-解码器 U 型结构）
  C: 训练过程流程图（6步训练管线）
  D: 推理过程流程图（去噪循环详解）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── 全局配置 ──────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

C_PRI = RGBColor(0x1E, 0x3A, 0x5F)
C_ACC = RGBColor(0x34, 0x98, 0xDB)
C_LT  = RGBColor(0xEB, 0xF5, 0xFB)
C_W   = RGBColor(0xFF, 0xFF, 0xFF)
C_D   = RGBColor(0x33, 0x33, 0x33)
C_G   = RGBColor(0x88, 0x88, 0x88)
C_GRN = RGBColor(0x27, 0xAE, 0x60)
C_ORG = RGBColor(0xE6, 0x7E, 0x22)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_PUR = RGBColor(0x8E, 0x44, 0xAD)
C_TEL = RGBColor(0x00, 0x96, 0x88)

FC = "微软雅黑"
FE = "Arial"

OUT = "/home/admin/.openclaw/workspace-weaver/output/sd_diagrams.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 工具函数 ──────────────────────────────────────────
def bg(s, c=C_W):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(s, l, t, w, h, fc=None, lc=None, lw=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb = lc
        if lw:
            sh.line.width = lw
    return sh

def rrect(s, l, t, w, h, fc=None, lc=None, lw=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb = lc
        if lw:
            sh.line.width = lw
    return sh

def tb(s, l, t, w, h, txt, fs=12, fc=C_D, b=False, a=PP_ALIGN.LEFT, fn=FC, ls=1.15):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(fs); p.font.color.rgb = fc
    p.font.bold = b; p.font.name = fn; p.alignment = a
    p.space_after = Pt(2)
    if ls != 1.0:
        p.line_spacing = ls
    return bx

def mtb(s, l, t, w, h, lines, fs=11, fc=C_D, fn=FC, a=PP_ALIGN.LEFT, ls=1.1):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        if isinstance(ld, str):
            txt, bld, c, sz = ld, False, fc, fs
        elif len(ld) == 2:
            txt, bld = ld; c, sz = fc, fs
        elif len(ld) == 3:
            txt, bld, c = ld; sz = fs
        else:
            txt, bld, c, sz = ld
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = bld; p.font.name = fn; p.alignment = a
        p.space_after = Pt(3); p.line_spacing = ls
    return bx

def arrow_r(s, l, t, w, h, c=C_ACC):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def arrow_d(s, l, t, w, h, c=C_ACC):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def arrow_l(s, l, t, w, h, c=C_ACC):
    sh = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def arrow_u(s, l, t, w, h, c=C_ACC):
    sh = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def chev(s, l, t, w, h, c=C_ACC):
    sh = s.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def oval(s, l, t, w, h, fc=None, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb = lc
    return sh

def pentagon(s, l, t, w, h, fc=None, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.PENTAGON, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb = lc
    return sh

def diamond(s, l, t, w, h, fc=None, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb = lc
    return sh

def tbl(s, l, t, w, rows, cols, data, cw=None, hc=C_PRI, hfc=C_W, fs=10, ac=C_LT):
    ts = s.shapes.add_table(rows, cols, l, t, w, Inches(0.32 * rows))
    table = ts.table
    if cw:
        for i, wi in enumerate(cw):
            table.columns[i].width = wi
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fs); p.font.name = FC; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(1)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hc
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = hfc; p.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ac
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_W
            tc = cell._tc; tcPr = tc.get_or_add_tcPr()
            for bn in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                b = tcPr.find(qn(bn))
                if b is None:
                    b = etree.SubElement(tcPr, qn(bn))
                b.set('w', '6350')
                sf = b.find(qn('a:solidFill'))
                if sf is None:
                    sf = etree.SubElement(b, qn('a:solidFill'))
                sc = sf.find(qn('a:srgbClr'))
                if sc is None:
                    sc = etree.SubElement(sf, qn('a:srgbClr'))
                sc.set('val', 'D0D0D0')
    return ts

def hdr(s, title, sub=""):
    rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fc=C_PRI)
    tb(s, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6), title, fs=28, fc=C_W, b=True)
    if sub:
        tb(s, Inches(0.8), Inches(0.65), Inches(11), Inches(0.35), sub, fs=13, fc=RGBColor(0xBB, 0xDE, 0xFB))
    rect(s, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), fc=C_PRI)
    tb(s, Inches(0.5), Inches(7.15), Inches(5), Inches(0.35), "Stable Diffusion 技术解析", fs=9, fc=RGBColor(0x88, 0xBB, 0xDD))


# ================================================================
#  页面 A：SD 整体架构图（推理流程全景）
# ================================================================
def page_a():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); hdr(sl, "SD 整体架构图", "完整推理流程：文字输入 → CLIP → U-Net(去噪循环) → VAE → 图像输出")

    # ── 主流程横向布局 ──
    # Step 1: 文字输入
    y_main = Inches(1.4)
    box_h = Inches(1.6)
    box_w = Inches(2.0)
    gap = Inches(0.35)

    # ① 文字输入（起始椭圆）
    x = Inches(0.4)
    oval(sl, x, y_main + Inches(0.3), Inches(1.6), Inches(1.0), fc=RGBColor(0xF0, 0xF4, 0xF8), lc=C_ACC, )
    tb(sl, x + Inches(0.05), y_main + Inches(0.42), Inches(1.5), Inches(0.75),
       "文字描述\n\"一只橘猫坐在\n沙发上\"", fs=11, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)

    # 箭头 ①→②
    arrow_r(sl, x + Inches(1.65), y_main + Inches(0.6), Inches(0.35), Inches(0.22), C_ACC)

    # ② CLIP 文本编码器（大框）
    x2 = x + Inches(2.05)
    rrect(sl, x2, y_main, box_w, box_h, fc=C_W, lc=C_ACC, lw=Pt(2))
    rect(sl, x2, y_main, box_w, Inches(0.32), fc=C_ACC)
    tb(sl, x2 + Inches(0.05), y_main + Inches(0.03), box_w - Inches(0.1), Inches(0.26),
       "CLIP 文本编码器", fs=12, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    mtb(sl, x2 + Inches(0.08), y_main + Inches(0.38), box_w - Inches(0.16), Inches(1.15), [
        ("参数量: 123M", True, C_ACC, 10),
        ("输入: 文字(≤77 token)", False, C_D, 9),
        ("输出: 77×768 向量", False, C_D, 9),
        ("12层 Transformer", False, C_G, 9),
        ("作用: 文字→AI数字语言", False, C_G, 9),
    ])

    # 箭头 ②→③
    arrow_r(sl, x2 + box_w + Inches(0.02), y_main + Inches(0.6), gap, Inches(0.22), C_ACC)

    # ③ U-Net 去噪网络（大框，最核心）
    x3 = x2 + box_w + Inches(0.4)
    unet_w = Inches(3.2)
    unet_h = Inches(2.6)
    rrect(sl, x3, y_main - Inches(0.5), unet_w, unet_h, fc=C_W, lc=C_RED, lw=Pt(3))
    rect(sl, x3, y_main - Inches(0.5), unet_w, Inches(0.32), fc=C_RED)
    tb(sl, x3 + Inches(0.05), y_main - Inches(0.47), unet_w - Inches(0.1), Inches(0.26),
       "U-Net 去噪网络（核心）", fs=13, fc=C_W, b=True, a=PP_ALIGN.CENTER)

    # U-Net 内部信息
    mtb(sl, x3 + Inches(0.1), y_main - Inches(0.12), unet_w - Inches(0.2), Inches(0.7), [
        ("参数量: 860M (占80%)", True, C_RED, 10),
        ("输入: 噪声潜空间(64×64×4) + 文本向量 + 时间步", False, C_D, 9),
        ("输出: 预测噪声(64×64×4)", False, C_D, 9),
    ])

    # 去噪循环箭头（U-Net内部）
    loop_y = y_main + Inches(0.7)
    loop_cx = x3 + unet_w / 2
    # 循环框
    rrect(sl, x3 + Inches(0.15), loop_y, unet_w - Inches(0.3), Inches(1.2), fc=RGBColor(0xFD, 0xED, 0xEC), lc=C_ORG, lw=Pt(1.5))
    tb(sl, x3 + Inches(0.2), loop_y + Inches(0.05), unet_w - Inches(0.4), Inches(0.2),
       "⟳ 去噪循环 (×20~50 步)", fs=11, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
    mtb(sl, x3 + Inches(0.25), loop_y + Inches(0.28), unet_w - Inches(0.5), Inches(0.85), [
        ("t=1000→950: 确定大致色块分布", False, C_D, 9),
        ("t=500→450: 模糊轮廓显现", False, C_D, 9),
        ("t=200→150: 细节逐渐清晰", False, C_D, 9),
        ("t=50→1: 精雕细琢完成", False, C_GRN, 9),
    ])

    # 循环弧形箭头（用上箭头 + 右箭头模拟）
    # 右侧上行箭头
    arrow_u(sl, x3 + unet_w - Inches(0.45), loop_y + Inches(0.15), Inches(0.18), Inches(0.25), C_ORG)
    # 顶部右行箭头
    arrow_r(sl, x3 + unet_w - Inches(0.45), loop_y - Inches(0.05), Inches(0.2), Inches(0.15), C_ORG)
    # 左侧下行箭头
    arrow_d(sl, x3 + Inches(0.2), loop_y - Inches(0.15), Inches(0.18), Inches(0.25), C_ORG)
    # 底部左行箭头
    arrow_l(sl, x3 + Inches(0.15), loop_y + Inches(0.95), Inches(0.2), Inches(0.15), C_ORG)

    # 箭头 ③→④
    arrow_r(sl, x3 + unet_w + Inches(0.02), y_main + Inches(0.6), gap, Inches(0.22), C_ACC)

    # ④ VAE 解码器
    x4 = x3 + unet_w + Inches(0.4)
    rrect(sl, x4, y_main, box_w, box_h, fc=C_W, lc=C_GRN, lw=Pt(2))
    rect(sl, x4, y_main, box_w, Inches(0.32), fc=C_GRN)
    tb(sl, x4 + Inches(0.05), y_main + Inches(0.03), box_w - Inches(0.1), Inches(0.26),
       "VAE 解码器", fs=12, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    mtb(sl, x4 + Inches(0.08), y_main + Inches(0.38), box_w - Inches(0.16), Inches(1.15), [
        ("参数量: 83M (解码~49M)", True, C_GRN, 10),
        ("输入: 64×64×4 潜空间", False, C_D, 9),
        ("输出: 512×512×3 RGB", False, C_D, 9),
        ("4层上采样(每层×2)", False, C_G, 9),
        ("作用: 潜空间→高清图像", False, C_G, 9),
    ])

    # 箭头 ④→⑤
    arrow_r(sl, x4 + box_w + Inches(0.02), y_main + Inches(0.6), gap, Inches(0.22), C_ACC)

    # ⑤ 图像输出（终止椭圆）
    x5 = x4 + box_w + Inches(0.4)
    oval(sl, x5, y_main + Inches(0.3), Inches(1.6), Inches(1.0), fc=RGBColor(0xE8, 0xF8, 0xF5), lc=C_GRN)
    tb(sl, x5 + Inches(0.05), y_main + Inches(0.42), Inches(1.5), Inches(0.75),
       "图像输出\n512×512×3\n高清 RGB", fs=11, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

    # ── 随机噪声输入（从下方进入U-Net） ──
    noise_x = x3 + Inches(0.3)
    noise_y = y_main + unet_h - Inches(0.4)
    rrect(sl, noise_x, noise_y, Inches(1.3), Inches(0.45), fc=RGBColor(0xFD, 0xF2, 0xE9), lc=C_ORG)
    tb(sl, noise_x + Inches(0.05), noise_y + Inches(0.03), Inches(1.2), Inches(0.38),
       "随机噪声\nN(0,I) 64×64×4", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
    arrow_u(sl, noise_x + Inches(0.55), noise_y - Inches(0.25), Inches(0.2), Inches(0.25), C_ORG)

    # ── Cross-Attention 连接线（CLIP → U-Net，虚线效果用浅色线表示） ──
    ca_x = x2 + box_w + Inches(0.02)
    ca_y = y_main + Inches(0.2)
    # 用一个小的弯曲路径模拟虚线：上方的连接线
    rect(sl, ca_x, ca_y, Inches(0.4), Inches(0.04), fc=C_PUR)
    tb(sl, ca_x - Inches(0.5), ca_y - Inches(0.18), Inches(1.4), Inches(0.2),
       "Cross-Attention", fs=8, fc=C_PUR, b=True, a=PP_ALIGN.CENTER)

    # ── 下方说明区域 ──
    desc_y = Inches(4.5)

    # 三列说明
    descs = [
        ("CLIP — 翻译官", C_ACC,
         "• OpenAI 2021年发布\n"
         "• 4亿张图文对训练\n"
         "• 理解文字与图像关系\n"
         "• 将文字翻译为77×768维\n"
         "  语义指纹向量\n"
         "• 通过Cross-Attention\n"
         "  注入U-Net每一步"),
        ("U-Net — 核心画师", C_RED,
         "• 860M参数，占80%算力\n"
         "• U型编码器-解码器结构\n"
         "• 3个关键机制:\n"
         "  - Cross-Attention(文本注入)\n"
         "  - Skip Connections(细节保留)\n"
         "  - Time Embedding(噪声感知)\n"
         "• 每步预测并去除噪声\n"
         "• 20-50步从噪声到清晰图像"),
        ("VAE — 压缩还原工具", C_GRN,
         "• 83M参数(编码34M+解码49M)\n"
         "• 48倍压缩: 786K→16K数值\n"
         "• 编码器: 512→64 (训练时用)\n"
         "• 解码器: 64→512 (推理时用)\n"
         "• 保留视觉重要信息\n"
         "• 丢弃人眼不敏感冗余\n"
         "• 确保潜空间连续平滑"),
    ]
    for i, (t, c, d) in enumerate(descs):
        dx = Inches(0.5) + Inches(i * 4.2)
        rrect(sl, dx, desc_y, Inches(3.9), Inches(2.4), fc=C_W, lc=c)
        rect(sl, dx, desc_y, Inches(3.9), Inches(0.3), fc=c)
        tb(sl, dx + Inches(0.08), desc_y + Inches(0.03), Inches(3.7), Inches(0.24),
           t, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
        tb(sl, dx + Inches(0.1), desc_y + Inches(0.38), Inches(3.7), Inches(1.95),
           d, fs=9.5, fc=C_D, ls=1.25)

    print("Page A: SD 整体架构图")


# ================================================================
#  页面 B：U-Net 架构图
# ================================================================
def page_b():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); hdr(sl, "U-Net 架构图", "编码器-瓶颈-解码器 U 型结构 | 860M 参数 | 3 个关键机制")

    # ── 左侧：U 型结构图 ──
    # 编码器（左侧，从上到下收缩）
    enc_x = Inches(0.6)
    enc_w = Inches(2.2)
    layer_h = Inches(0.65)
    layer_gap = Inches(0.15)

    # 编码器4层 + 瓶颈 + 解码器4层
    enc_layers = [
        ("64×64×320", "Res+Attn+Down", C_ACC, True),
        ("32×32×640", "Res+Attn+Down", C_ACC, True),
        ("16×16×1280", "Res+Attn+Down", C_ACC, True),
        ("8×8×1280", "Res+Down", C_ACC, False),
    ]
    bottleneck = ("4×4×1280", "Res+SelfAttn+Res", C_PUR)
    dec_layers = [
        ("8×8×1280", "Res+Up+Skip", C_GRN, True),
        ("16×16×640", "Res+Up+Skip", C_GRN, True),
        ("32×32×320", "Res+Up+Skip", C_GRN, True),
        ("64×64×4", "输出预测噪声", C_GRN, False),
    ]

    start_y = Inches(1.35)

    # 画编码器
    enc_positions = []
    for i, (dim, desc, c, has_ca) in enumerate(enc_layers):
        y = start_y + Inches(i * (layer_h + layer_gap))
        rrect(sl, enc_x, y, enc_w, layer_h, fc=C_W, lc=c, lw=Pt(2))
        tb(sl, enc_x + Inches(0.05), y + Inches(0.02), enc_w - Inches(0.1), Inches(0.22),
           dim, fs=10, fc=c, b=True, a=PP_ALIGN.CENTER, fn=FE)
        tb(sl, enc_x + Inches(0.05), y + Inches(0.25), enc_w - Inches(0.1), Inches(0.18),
           desc, fs=8, fc=C_G, a=PP_ALIGN.CENTER)
        tb(sl, enc_x + Inches(0.05), y + Inches(0.42), enc_w - Inches(0.1), Inches(0.18),
           "✦ Cross-Attention" if has_ca else "  (无注意力)", fs=8,
           fc=C_PUR if has_ca else C_G, a=PP_ALIGN.CENTER)
        enc_positions.append((enc_x, y, enc_w, layer_h))

        # 下箭头
        if i < len(enc_layers) - 1:
            arrow_d(sl, enc_x + enc_w / 2 - Inches(0.1), y + layer_h,
                    Inches(0.2), layer_gap, c)

    # 编码器到瓶颈的箭头
    last_enc_y = start_y + Inches((len(enc_layers) - 1) * (layer_h + layer_gap))
    arrow_d(sl, enc_x + enc_w / 2 - Inches(0.1), last_enc_y + layer_h,
            Inches(0.2), layer_gap, C_ACC)

    # 瓶颈层
    bot_y = last_enc_y + layer_h + layer_gap
    rrect(sl, enc_x - Inches(0.15), bot_y, enc_w + Inches(0.3), layer_h + Inches(0.1),
          fc=RGBColor(0xF5, 0xEE, 0xF8), lc=C_PUR, lw=Pt(3))
    tb(sl, enc_x - Inches(0.1), bot_y + Inches(0.02), enc_w + Inches(0.2), Inches(0.22),
       bottleneck[0], fs=11, fc=C_PUR, b=True, a=PP_ALIGN.CENTER, fn=FE)
    tb(sl, enc_x - Inches(0.1), bot_y + Inches(0.25), enc_w + Inches(0.2), Inches(0.18),
       "Bottleneck", fs=9, fc=C_PUR, b=True, a=PP_ALIGN.CENTER)
    tb(sl, enc_x - Inches(0.1), bot_y + Inches(0.42), enc_w + Inches(0.2), Inches(0.18),
       bottleneck[1], fs=8, fc=C_G, a=PP_ALIGN.CENTER)

    # 瓶颈到解码器的箭头
    arrow_d(sl, enc_x + enc_w / 2 - Inches(0.1), bot_y + layer_h + Inches(0.1),
            Inches(0.2), layer_gap, C_GRN)

    # 解码器
    dec_x = Inches(4.0)
    dec_w = Inches(2.2)
    dec_positions = []
    for i, (dim, desc, c, has_ca) in enumerate(dec_layers):
        y = bot_y + layer_h + Inches(0.1) + layer_gap + Inches(i * (layer_h + layer_gap))
        rrect(sl, dec_x, y, dec_w, layer_h, fc=C_W, lc=c, lw=Pt(2))
        tb(sl, dec_x + Inches(0.05), y + Inches(0.02), dec_w - Inches(0.1), Inches(0.22),
           dim, fs=10, fc=c, b=True, a=PP_ALIGN.CENTER, fn=FE)
        tb(sl, dec_x + Inches(0.05), y + Inches(0.25), dec_w - Inches(0.1), Inches(0.18),
           desc, fs=8, fc=C_G, a=PP_ALIGN.CENTER)
        tb(sl, dec_x + Inches(0.05), y + Inches(0.42), dec_w - Inches(0.1), Inches(0.18),
           "✦ Cross-Attention" if has_ca else "  (无注意力)", fs=8,
           fc=C_PUR if has_ca else C_G, a=PP_ALIGN.CENTER)
        dec_positions.append((dec_x, y, dec_w, layer_h))

        if i < len(dec_layers) - 1:
            arrow_d(sl, dec_x + dec_w / 2 - Inches(0.1), y + layer_h,
                    Inches(0.2), layer_gap, c)

    # Skip Connections（从编码器到解码器的水平线）
    skip_colors = [C_ORG, C_ORG, C_ORG, C_ORG]
    for i in range(min(len(enc_layers), len(dec_layers))):
        ex, ey, ew, eh = enc_positions[i]
        dx, dy, dw, dh = dec_positions[i]
        mid_y = ey + eh / 2
        # 水平线
        rect(sl, ex + ew, mid_y - Inches(0.02), dx - (ex + ew), Inches(0.04),
             fc=skip_colors[i])
        # Skip标签
        tb(sl, ex + ew + Inches(0.15), mid_y - Inches(0.2), Inches(1.0), Inches(0.18),
           "Skip", fs=8, fc=skip_colors[i], b=True, a=PP_ALIGN.CENTER)
        # 小箭头
        arrow_r(sl, dx - Inches(0.2), mid_y - Inches(0.08), Inches(0.18), Inches(0.14), skip_colors[i])

    # "编码器" 和 "解码器" 标签
    tb(sl, enc_x, start_y - Inches(0.25), enc_w, Inches(0.22),
       "编码器 Encoder ↓", fs=10, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
    tb(sl, dec_x, start_y - Inches(0.25), dec_w, Inches(0.22),
       "解码器 Decoder ↓", fs=10, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

    # Time Embedding 注入标注（左侧）
    te_x = Inches(0.05)
    for i in range(len(enc_layers)):
        y = start_y + Inches(i * (layer_h + layer_gap))
        arrow_r(sl, te_x, y + Inches(0.25), Inches(0.5), Inches(0.15), C_TEL)
    tb(sl, te_x, start_y + Inches(0.05), Inches(0.55), Inches(0.2),
       "Time", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)
    tb(sl, te_x, start_y + Inches(0.2), Inches(0.55), Inches(0.2),
       "Embed", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)

    # ── 右侧说明区域 ──
    info_x = Inches(6.8)
    info_w = Inches(6.0)

    # 三个关键机制
    mechanisms = [
        ("① Cross-Attention 交叉注意力", C_ACC,
         "文本条件注入的核心机制",
         "• 图像特征作为 Query，文本特征作为 Key-Value\n"
         "• 图像问：根据文字，我这里该怎样？\n"
         "• 文本答：你这里应该是橘猫的毛发\n"
         "• 出现位置：编码器1-3层 + 解码器1-3层\n"
         "• 共6个位置注入文本指导"),
        ("② Skip Connections 跳跃连接", C_ORG,
         "信息传递的桥梁",
         "• 编码器特征直接桥接到解码器对应层\n"
         "• 防止细节在压缩过程中丢失\n"
         "• 像建筑师随时参考草稿本\n"
         "• 没有跳跃连接，细节会大打折扣\n"
         "• 图中橙色线即Skip Connection"),
        ("③ Time Embedding 时间嵌入", C_TEL,
         "当前噪声水平的标尺",
         "• 告诉U-Net当前噪声水平(1-1000)\n"
         "• 高噪声(t≈900)：粗雕确定轮廓\n"
         "• 低噪声(t≈100)：精雕添加细节\n"
         "• 正弦编码 → 2层MLP → 1280维\n"
         "• 注入每一个ResBlock中"),
    ]

    for i, (title, c, subtitle, desc) in enumerate(mechanisms):
        y = Inches(1.3) + Inches(i * 1.9)
        rrect(sl, info_x, y, info_w, Inches(1.75), fc=C_W, lc=c)
        rect(sl, info_x, y, info_w, Inches(0.28), fc=c)
        tb(sl, info_x + Inches(0.1), y + Inches(0.02), info_w - Inches(0.2), Inches(0.24),
           title, fs=11, fc=C_W, b=True)
        tb(sl, info_x + Inches(0.1), y + Inches(0.3), info_w - Inches(0.2), Inches(0.2),
           subtitle, fs=10, fc=c, b=True)
        tb(sl, info_x + Inches(0.1), y + Inches(0.52), info_w - Inches(0.2), Inches(1.15),
           desc, fs=9.5, fc=C_D, ls=1.2)

    print("Page B: U-Net 架构图")


# ================================================================
#  页面 C：训练过程流程图
# ================================================================
def page_c():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); hdr(sl, "训练过程流程图", "6步训练管线 — 15万GPU小时，50-60万美元")

    # ── 6步流程图（横向两行，每行3步） ──
    steps = [
        ("Step 1", "采样图像+文本", C_ACC,
         "从 LAION-5B 随机抽取\n图像 x₀ + 文本描述",
         "6亿图文对"),
        ("Step 2", "VAE 编码", C_GRN,
         "512×512×3 → 64×64×4\n压缩到潜空间(48倍)",
         "z₀ = VAE(x₀)"),
        ("Step 3", "添加随机噪声", C_ORG,
         "随机选 t ∈ [1, 1000]\n按比例添加高斯噪声",
         "z_t = √ᾱ_t·z₀ + √(1-ᾱ_t)·ε"),
        ("Step 4", "U-Net 预测噪声", C_RED,
         "输入: z_t + t + 文本向量\n输出: 预测噪声 ε_pred",
         "ε_pred = U-Net(z_t, t, c)"),
        ("Step 5", "计算损失", C_PUR,
         "预测噪声 vs 真实噪声\n均方误差(MSE)",
         "L = ‖ε - ε_pred‖²"),
        ("Step 6", "反向传播", C_TEL,
         "更新 U-Net 的 860M 参数\nCLIP + VAE 保持冻结",
         "θ ← θ - η·∇L"),
    ]

    box_w = Inches(3.6)
    box_h = Inches(1.7)
    arrow_gap = Inches(0.3)
    row_gap_x = box_w + arrow_gap
    start_x = Inches(0.5)
    row1_y = Inches(1.3)
    row2_y = row1_y + box_h + Inches(0.55)

    for i, (step_num, title, c, desc, formula) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = start_x + Inches(col * (box_w + Inches(0.6)))
        y = row1_y if row == 0 else row2_y

        # 主框
        rrect(sl, x, y, box_w, box_h, fc=C_W, lc=c, lw=Pt(2))
        # 标题条
        rect(sl, x, y, box_w, Inches(0.3), fc=c)
        tb(sl, x + Inches(0.08), y + Inches(0.03), Inches(0.7), Inches(0.24),
           step_num, fs=10, fc=RGBColor(0xDD, 0xEE, 0xFF))
        tb(sl, x + Inches(0.8), y + Inches(0.03), box_w - Inches(0.9), Inches(0.24),
           title, fs=12, fc=C_W, b=True)
        # 描述
        tb(sl, x + Inches(0.1), y + Inches(0.38), box_w - Inches(0.2), Inches(0.85),
           desc, fs=10, fc=C_D, ls=1.2)
        # 公式
        rrect(sl, x + Inches(0.1), y + Inches(1.25), box_w - Inches(0.2), Inches(0.35),
              fc=RGBColor(0xF8, 0xF9, 0xFA), lc=c)
        tb(sl, x + Inches(0.15), y + Inches(1.28), box_w - Inches(0.3), Inches(0.3),
           formula, fs=10, fc=c, b=True, a=PP_ALIGN.CENTER, fn=FE)

        # 步骤间的箭头
        if col < 2:
            arrow_r(sl, x + box_w + Inches(0.05), y + Inches(0.6),
                    arrow_gap, Inches(0.2), C_G)

    # 行间连接（Step 3 → Step 4）
    arrow_d(sl, start_x + Inches(2 * (box_w + Inches(0.6))) + box_w / 2 - Inches(0.1),
            row1_y + box_h + Inches(0.05), Inches(0.2), Inches(0.4), C_G)
    # 回绕指示
    tb(sl, start_x + Inches(2 * (box_w + Inches(0.6))) + box_w / 2 - Inches(0.3),
       row1_y + box_h + Inches(0.1), Inches(0.6), Inches(0.2),
       "继续 →", fs=9, fc=C_G, b=True, a=PP_ALIGN.CENTER)

    # ── 训练循环指示 ──
    loop_x = start_x - Inches(0.15)
    loop_y = row2_y + box_h + Inches(0.1)
    # 循环指示箭头（从Step6回到Step1）
    rrect(sl, loop_x, loop_y, Inches(12.3), Inches(0.45), fc=RGBColor(0xFE, 0xF9, 0xE7), lc=C_ORG)
    arrow_l(sl, loop_x + Inches(0.1), loop_y + Inches(0.12), Inches(0.4), Inches(0.18), C_ORG)
    tb(sl, loop_x + Inches(0.55), loop_y + Inches(0.08), Inches(11.5), Inches(0.3),
       "⟳ 重复训练循环：重复 Step 1-6 数十万次，直到损失收敛 | 每次迭代随机采样不同的图像和噪声水平",
       fs=10, fc=C_ORG, b=True)

    # ── 底部：训练成本 ──
    cost_y = Inches(6.0)
    tbl(sl, Inches(0.5), cost_y, Inches(12.3), 3, 5, [
        ["训练数据", "训练硬件", "训练时长", "训练成本", "模型大小"],
        ["LAION-5B 子集\n~6亿图文对", "多张\nA100 80GB", "~15万\nGPU小时", "~50-60万\n美元", "~4GB\n(float16)"],
        ["冻结组件: CLIP(123M) + VAE(83M)  |  只训练: U-Net(860M)  |  损失函数: MSE ‖ε - ε_pred‖²", "", "", "", ""],
    ], cw=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.3)], fs=10)

    print("Page C: 训练过程流程图")


# ================================================================
#  页面 D：推理过程流程图
# ================================================================
def page_d():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl); hdr(sl, "推理过程流程图", "4步推理 + 去噪循环详解 — 从纯噪声到高清图像")

    # ── 上方：4步概览 ──
    overview_y = Inches(1.25)
    ov_steps = [
        ("Step 1", "文本编码", "文字 → CLIP → 77×768 向量\n只需执行一次", C_ACC),
        ("Step 2", "初始化噪声", "采样 64×64×4 随机张量\n= 纯噪声雪花屏", C_ORG),
        ("Step 3", "逐步去噪循环", "U-Net 反复预测噪声\n20-50 步逐步清晰", C_RED),
        ("Step 4", "VAE 解码", "64×64×4 → 512×512×3\n4层上采样还原", C_GRN),
    ]

    ov_box_w = Inches(2.9)
    for i, (sn, title, desc, c) in enumerate(ov_steps):
        x = Inches(0.5) + Inches(i * 3.15)
        rrect(sl, x, overview_y, ov_box_w, Inches(1.15), fc=C_W, lc=c, lw=Pt(2))
        rect(sl, x, overview_y, ov_box_w, Inches(0.26), fc=c)
        tb(sl, x + Inches(0.05), overview_y + Inches(0.02), Inches(0.7), Inches(0.22),
           sn, fs=9, fc=RGBColor(0xDD, 0xEE, 0xFF))
        tb(sl, x + Inches(0.75), overview_y + Inches(0.02), ov_box_w - Inches(0.85), Inches(0.22),
           title, fs=11, fc=C_W, b=True)
        tb(sl, x + Inches(0.08), overview_y + Inches(0.32), ov_box_w - Inches(0.16), Inches(0.8),
           desc, fs=10, fc=C_D, ls=1.2)
        if i < 3:
            arrow_r(sl, x + ov_box_w + Inches(0.03), overview_y + Inches(0.35),
                    Inches(0.25), Inches(0.18), c)

    # ── 下方：去噪循环详解（重点） ──
    loop_section_y = Inches(2.6)
    # 循环区域背景
    rrect(sl, Inches(0.4), loop_section_y, Inches(8.5), Inches(4.3),
          fc=RGBColor(0xFD, 0xED, 0xEC), lc=C_RED, lw=Pt(2))
    tb(sl, Inches(0.6), loop_section_y + Inches(0.05), Inches(5), Inches(0.25),
       "去噪循环详解 (Step 3 展开)", fs=13, fc=C_RED, b=True)

    # 时间步进度条（从上到下）
    ts_x = Inches(0.7)
    ts_y = loop_section_y + Inches(0.4)
    ts_w = Inches(1.2)
    ts_h = Inches(0.5)

    # 时间步标注
    timesteps = [
        ("t = 1000", "纯噪声\n随机像素", C_RED),
        ("t ≈ 750", "粗雕阶段\n大致色块", C_ORG),
        ("t ≈ 500", "中雕阶段\n模糊轮廓", C_ORG),
        ("t ≈ 200", "精雕阶段\n细节显现", C_GRN),
        ("t = 1", "完成\n清晰图像", C_GRN),
    ]

    for i, (t, desc, c) in enumerate(timesteps):
        y = ts_y + Inches(i * 0.75)
        # 时间步圆
        oval(sl, ts_x, y, ts_w, ts_h, fc=C_W, lc=c, )
        tb(sl, ts_x + Inches(0.05), y + Inches(0.02), ts_w - Inches(0.1), Inches(0.2),
           t, fs=11, fc=c, b=True, a=PP_ALIGN.CENTER, fn=FE)
        tb(sl, ts_x + Inches(0.05), y + Inches(0.22), ts_w - Inches(0.1), Inches(0.25),
           desc, fs=8, fc=C_D, a=PP_ALIGN.CENTER)

        # 下箭头
        if i < len(timesteps) - 1:
            arrow_d(sl, ts_x + ts_w / 2 - Inches(0.08), y + ts_h,
                    Inches(0.16), Inches(0.22), c)

    # 循环内部详细步骤（右侧）
    detail_x = Inches(2.2)
    detail_y = ts_y + Inches(0.1)

    # 循环体框
    rrect(sl, detail_x, detail_y, Inches(6.4), Inches(3.6), fc=C_W, lc=C_ACC, lw=Pt(1.5))
    tb(sl, detail_x + Inches(0.1), detail_y + Inches(0.05), Inches(6.2), Inches(0.22),
       "每步去噪操作（for t = T, T-Δ, ..., 1）", fs=11, fc=C_ACC, b=True)

    # 内部3步
    inner_steps = [
        ("① U-Net 预测噪声", C_RED,
         "ε_pred = U-Net(x_t, t, text_embedding)\n"
         "输入: 当前带噪声图像 x_t + 时间步 t + 文本向量\n"
         "输出: 预测的噪声 ε_pred (64×64×4)"),
        ("② 采样器计算去噪", C_PUR,
         "x_{t-1} = sampler_step(x_t, ε_pred, t)\n"
         "如 DPM++ 2M Karras: 智能调整去噪幅度\n"
         "噪声高→大步去噪 | 噪声低→精细去噪"),
        ("③ 更新图像", C_GRN,
         "x_t ← x_{t-1}\n"
         "用更干净的图像替换当前图像\n"
         "重复直到 t = 1，去噪完成"),
    ]

    for i, (title, c, desc) in enumerate(inner_steps):
        y = detail_y + Inches(0.35) + Inches(i * 1.05)
        rrect(sl, detail_x + Inches(0.15), y, Inches(6.1), Inches(0.9), fc=C_W, lc=c)
        tb(sl, detail_x + Inches(0.2), y + Inches(0.03), Inches(5.9), Inches(0.2),
           title, fs=10, fc=c, b=True)
        tb(sl, detail_x + Inches(0.2), y + Inches(0.25), Inches(5.9), Inches(0.6),
           desc, fs=9, fc=C_D, ls=1.15)
        if i < 2:
            arrow_d(sl, detail_x + Inches(3.15), y + Inches(0.9),
                    Inches(0.16), Inches(0.13), C_ACC)

    # 循环回箭头
    loop_bottom = detail_y + Inches(3.6) - Inches(0.15)
    arrow_l(sl, detail_x + Inches(2.5), loop_bottom, Inches(0.5), Inches(0.18), C_RED)
    tb(sl, detail_x + Inches(1.0), loop_bottom - Inches(0.02), Inches(1.5), Inches(0.2),
       "⟳ 循环", fs=9, fc=C_RED, b=True, a=PP_ALIGN.CENTER)

    # ── 右侧：推理时间参考 ──
    ref_x = Inches(9.2)
    ref_y = Inches(2.6)
    rrect(sl, ref_x, ref_y, Inches(3.6), Inches(4.3), fc=C_W, lc=C_PRI, lw=Pt(1.5))
    rect(sl, ref_x, ref_y, Inches(3.6), Inches(0.3), fc=C_PRI)
    tb(sl, ref_x + Inches(0.05), ref_y + Inches(0.03), Inches(3.5), Inches(0.24),
       "推理时间参考", fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)

    mtb(sl, ref_x + Inches(0.1), ref_y + Inches(0.4), Inches(3.4), Inches(3.7), [
        ("配置              步数    GPU       耗时", True, C_PRI, 9),
        ("", False, C_D, 4),
        ("快速预览          10-15   3060      2-4秒", False, C_D, 9),
        ("标准质量          20-30   3060      5-8秒", False, C_GRN, 9),
        ("高质量            40-50   4090      3-5秒", False, C_D, 9),
        ("最高质量          50-100  4090      5-10秒", False, C_D, 9),
        ("", False, C_D, 4),
        ("推荐配置:", True, C_ACC, 10),
        ("DPM++ 2M Karras", True, C_ACC, 10),
        ("25步 / CFG 7.5", False, C_ACC, 10),
        ("RTX 3060 / 512×512", False, C_G, 9),
        ("", False, C_D, 4),
        ("采样器选择:", True, C_PUR, 10),
        ("通用: DPM++ 2M Karras", False, C_D, 9),
        ("快速: DPM++ SDE (15步)", False, C_D, 9),
        ("极速: LCM (4-8步)", False, C_D, 9),
    ], fn=FE, a=PP_ALIGN.LEFT, ls=1.15)

    print("Page D: 推理过程流程图")


# ================================================================
#  生成
# ================================================================
page_a()
page_b()
page_c()
page_d()

prs.save(OUT)
print(f"\nDone! {OUT} | {len(prs.slides)} pages")
