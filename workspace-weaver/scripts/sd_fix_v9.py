#!/usr/bin/env python3
"""
SD PPT Fix v9: Fix 6 issues reported by user on v8
  1. P4 (Slide 4): 整体架构 — 排版布局重做，框清晰
  2. P5 (Slide 5): CLIP详解 — 文字重叠修复
  3. P6 (Slide 6): VAE详解 — 箭头方向修正 + 文字超框修复
  4. P7 (Slide 7): U-Net架构 — 框闭合 + 底部箭头指向明确
  5. P10 (Slide 10): CFG — 箭头修正 + 文字超框修复 + 移到推理页之后

Strategy: Read v8 PPTX, replace problematic slides with fixed versions.
Slide order change: CFG (currently P10) moves after 推理过程 (currently P11)

New slide order:
  P1=封面, P2=目录, P3=Diffusion基础, P4=整体架构(修), P5=CLIP(修),
  P6=VAE(修), P7=U-Net(修), P8=协作流程, P9=训练过程, P10=推理过程,
  P11=CFG(修,移至此), P12=采样器+参数, P13=Prompt+LoRA+ControlNet,
  P14=SD3+Flux, P15=趋势总结

Output: output/Stable_Diffusion_技术解析_v9.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os, sys

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = "/home/admin/.openclaw/workspace-weaver/output/Stable_Diffusion_技术解析_v9.pptx"
SRC = "/home/admin/.openclaw/workspace-weaver/output/Stable_Diffusion_技术解析_v8.pptx"

# Colors
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_LT  = RGBColor(0xEB,0xF5,0xFB)
C_W   = RGBColor(0xFF,0xFF,0xFF)
C_D   = RGBColor(0x33,0x33,0x33)
C_G   = RGBColor(0x88,0x88,0x88)
C_GRN = RGBColor(0x27,0xAE,0x60)
C_ORG = RGBColor(0xE6,0x7E,0x22)
C_RED = RGBColor(0xE7,0x4C,0x3C)
C_PUR = RGBColor(0x8E,0x44,0xAD)
C_TEL = RGBColor(0x00,0x96,0x88)
C_BG_GRN = RGBColor(0xE8,0xF8,0xF5)
C_BG_ORG = RGBColor(0xFD,0xF2,0xE9)
C_BG_RED = RGBColor(0xFD,0xED,0xED)
C_YEL = RGBColor(0xF3,0x9C,0x12)

FC = "微软雅黑"
FE = "Arial"

# ── Helpers ──
def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb=lc
        if lw: sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb=lc
        if lw: sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    # Set font at run level for reliability
    for r in p.runs:
        r.font.size=Pt(fs); r.font.color.rgb=fc; r.font.bold=b; r.font.name=fn
    # Fallback: also set at paragraph level
    p.font.size=Pt(fs); p.font.color.rgb=fc; p.font.bold=b; p.font.name=fn
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
        for r in p.runs:
            r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=bld; r.font.name=fn
        p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld; p.font.name=fn
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_u(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.UP_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def line_h(s,l,t,w,c=C_ACC,lw=2):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Emu(int(lw*12700)))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Stable Diffusion 技术解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    violations = []
    for sh in s.shapes:
        if sh.left < Emu(-100):
            violations.append(f"  NEGATIVE X: left={sh.left/914400:.3f}")
        if sh.top < Emu(-100):
            violations.append(f"  NEGATIVE Y: top={sh.top/914400:.3f}")
        r = sh.left + sh.width
        b = sh.top + sh.height
        if r > SLIDE_W + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW X: right={r/914400:.3f} text={txt}")
        if b > SLIDE_H + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW Y: bottom={b/914400:.3f} text={txt}")
    if violations:
        print(f"  {msg} — {len(violations)} violations:")
        for v in violations: print(v)
        return False
    else:
        print(f"  OK {msg} — all shapes within bounds")
        return True

def copy_slide(prs, src_slide):
    """Copy a slide from one presentation to another."""
    slide_layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy background
    bg_fill = src_slide.background.fill
    if bg_fill.type is not None:
        # Can't directly copy fill, set solid color if possible
        try:
            new_bg = new_slide.background.fill
            new_bg.solid()
            new_bg.fore_color.rgb = bg_fill.fore_color.rgb
        except:
            pass
    
    # Copy shapes
    for shape in src_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide


# ============================================================
# Load source PPTX
# ============================================================
print("Loading v8 source...")
src_prs = Presentation(SRC)
print(f"  Source: {len(src_prs.slides)} slides, {src_prs.slide_width/914400:.2f}x{src_prs.slide_height/914400:.2f} inches")

# Create new presentation
prs = Presentation()
prs.slide_width = src_prs.slide_width
prs.slide_height = src_prs.slide_height

# ============================================================
# SLIDE 1 (REBUILD): 封面 — 修复字体显示
# ============================================================
print("\nRebuilding P1: 封面...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_PRI)

# 装饰线
rect(sl, Inches(0.8), Inches(2.8), Inches(3.0), Inches(0.06), fc=C_ACC)

# 主标题 Stable Diffusion
tb(sl, Inches(0.8), Inches(3.0), Inches(11.0), Inches(1.20),
   "Stable Diffusion", fs=44, fc=C_W, b=True, fn="Arial", ls=1.0)

# 副标题
tb(sl, Inches(0.8), Inches(4.0), Inches(11.0), Inches(0.60),
   "技术解析 — 从原理到前沿", fs=24, fc=RGBColor(0xBB,0xDE,0xFB), b=False, fn=FC, ls=1.0)

# 装饰线
rect(sl, Inches(0.8), Inches(4.6), Inches(1.5), Inches(0.04), fc=C_ACC)

# 描述
tb(sl, Inches(0.8), Inches(5.0), Inches(8.0), Inches(0.40),
   "深入理解扩散模型、CLIP、VAE、U-Net 及最新进展", fs=13, fc=RGBColor(0x99,0xBB,0xDD), fn=FC)
tb(sl, Inches(0.8), Inches(5.5), Inches(8.0), Inches(0.40),
   "2026 年 4 月  |  19 章节完整版", fs=12, fc=RGBColor(0x88,0xAA,0xCC), fn=FC)

# 右侧装饰区域
rect(sl, Inches(10.0), Inches(0.0), Inches(3.33), Inches(7.50), fc=RGBColor(0x17,0x2E,0x4A))
tb(sl, Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.50),
   "AI IMAGE\nGENERATION", fs=18, fc=RGBColor(0x34,0x98,0xDB), b=True, fn="Arial", a=PP_ALIGN.CENTER, ls=1.3)

chk(sl, "P1 封面")

# ============================================================
# Copy unchanged slides: 2(目录), 3(Diffusion基础)
# ============================================================
print("\nCopying unchanged slides...")
for idx in [1, 2]:  # slides 2,3
    copy_slide(prs, src_prs.slides[idx])
    print(f"  Copied slide {idx+1}")


# ============================================================
# SLIDE 4 (FIX): 整体架构概览 — 重做排版
# ============================================================
print("\nBuilding fixed P4: 整体架构概览...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "03  整体架构概览", "CLIP + U-Net + VAE 三大组件协同工作 — 推理与训练流程完全分离")

# ── 推理流程区域 ──
# 区域外框
rrect(sl, Inches(0.15), Inches(1.15), Inches(12.80), Inches(2.05), fc=None, lc=C_ACC, lw=2)
rect(sl, Inches(0.15), Inches(1.15), Inches(12.80), Inches(0.30), fc=C_ACC)
tb(sl, Inches(0.30), Inches(1.17), Inches(8.0), Inches(0.28),
   "▶ 推理流程（从文字到图像）— 主要使用场景", fs=11, fc=C_W, b=True)

# Step boxes (horizontal flow)
# ① 文字输入
rrect(sl, Inches(0.35), Inches(1.60), Inches(1.40), Inches(0.85), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(0.35), Inches(1.62), Inches(1.40), Inches(0.22),
   "① 文字输入", fs=10, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(1.87), Inches(1.40), Inches(0.50),
   '"一只橘猫\n坐在沙发上"', fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(1.80), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ② CLIP
rrect(sl, Inches(2.15), Inches(1.55), Inches(1.80), Inches(0.95), fc=C_LT, lc=C_ACC, lw=1)
rect(sl, Inches(2.15), Inches(1.55), Inches(1.80), Inches(0.25), fc=C_ACC)
tb(sl, Inches(2.15), Inches(1.57), Inches(1.80), Inches(0.22),
   "② CLIP 文本编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(2.15), Inches(1.83), Inches(1.80), Inches(0.60),
   "123M参数\n12层Transformer\n输出: 77×768向量", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(4.00), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ③ U-Net
rrect(sl, Inches(4.35), Inches(1.50), Inches(2.00), Inches(1.05), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI, lw=2)
rect(sl, Inches(4.35), Inches(1.50), Inches(2.00), Inches(0.25), fc=C_PRI)
tb(sl, Inches(4.35), Inches(1.52), Inches(2.00), Inches(0.22),
   "③ U-Net 去噪（核心）", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(4.35), Inches(1.78), Inches(2.00), Inches(0.70),
   "860M参数(80%)\n⟳ 循环20-50步\n噪声→逐步去噪", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# CA label
rect(sl, Inches(3.95), Inches(1.78), Inches(0.35), Inches(0.20), fc=C_PUR)
tb(sl, Inches(3.95), Inches(1.78), Inches(0.35), Inches(0.20),
   "CA", fs=7, fc=C_W, b=True, a=PP_ALIGN.CENTER)

arrow_r(sl, Inches(6.40), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ④ VAE 解码器
rrect(sl, Inches(6.75), Inches(1.55), Inches(1.80), Inches(0.95), fc=C_BG_GRN, lc=C_GRN, lw=1)
rect(sl, Inches(6.75), Inches(1.55), Inches(1.80), Inches(0.25), fc=C_GRN)
tb(sl, Inches(6.75), Inches(1.57), Inches(1.80), Inches(0.22),
   "④ VAE 解码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(6.75), Inches(1.83), Inches(1.80), Inches(0.60),
   "~49M参数\n64×64→512×512\n4层上采样", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(8.60), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ⑤ 图像输出
rrect(sl, Inches(8.95), Inches(1.60), Inches(1.40), Inches(0.85), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(8.95), Inches(1.62), Inches(1.40), Inches(0.22),
   "⑤ 图像输出", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(8.95), Inches(1.87), Inches(1.40), Inches(0.50),
   "512×512×3\n高清RGB图像", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# 随机噪声 → U-Net
rrect(sl, Inches(4.65), Inches(2.65), Inches(1.40), Inches(0.45), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(4.65), Inches(2.68), Inches(1.40), Inches(0.40),
   "随机噪声 N(0,I)\n64×64×4 纯噪声", fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(5.28), Inches(2.55), Inches(0.18), Inches(0.15), C_RED)

# VAE编码器：推理时不需要
rrect(sl, Inches(10.60), Inches(1.60), Inches(2.20), Inches(0.85), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(10.60), Inches(1.65), Inches(2.20), Inches(0.75),
   "❌ VAE 编码器\n推理时不需要！\n（编码器只在训练时使用）", fs=9, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.1)

# ── 训练流程区域 ──
rrect(sl, Inches(0.15), Inches(3.35), Inches(12.80), Inches(1.80), fc=None, lc=C_ORG, lw=2)
rect(sl, Inches(0.15), Inches(3.35), Inches(12.80), Inches(0.30), fc=C_ORG)
tb(sl, Inches(0.30), Inches(3.37), Inches(8.0), Inches(0.28),
   "▶ 训练流程（从图像学习去噪能力）— 模型训练阶段", fs=11, fc=C_W, b=True)

# T1: Image+Text
rrect(sl, Inches(0.35), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_W, lc=C_ORG, lw=1)
tb(sl, Inches(0.35), Inches(3.82), Inches(1.60), Inches(0.22),
   "图像+文本对", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(4.07), Inches(1.60), Inches(0.40),
   "512×512×3 图像\n+ 文字描述", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(2.00), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

# T2: VAE编码器
rrect(sl, Inches(2.35), Inches(3.75), Inches(1.80), Inches(0.85), fc=C_BG_GRN, lc=C_GRN, lw=1)
rect(sl, Inches(2.35), Inches(3.75), Inches(1.80), Inches(0.25), fc=C_GRN)
tb(sl, Inches(2.35), Inches(3.77), Inches(1.80), Inches(0.22),
   "VAE 编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(2.35), Inches(4.03), Inches(1.80), Inches(0.50),
   "~34M参数\n512²→64² 压缩48倍", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(4.20), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

# T3: 添加噪声
rrect(sl, Inches(4.55), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(4.55), Inches(3.82), Inches(1.60), Inches(0.22),
   "添加噪声", fs=10, fc=C_RED, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(4.55), Inches(4.07), Inches(1.60), Inches(0.40),
   "z_t = √ᾱ·z₀ + √(1-ᾱ)·ε\n前向扩散过程", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(6.20), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

# T4: U-Net
rrect(sl, Inches(6.55), Inches(3.75), Inches(2.00), Inches(0.85), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI, lw=1)
rect(sl, Inches(6.55), Inches(3.75), Inches(2.00), Inches(0.25), fc=C_PRI)
tb(sl, Inches(6.55), Inches(3.77), Inches(2.00), Inches(0.22),
   "U-Net 预测噪声", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(6.55), Inches(4.03), Inches(2.00), Inches(0.50),
   "输入: 噪声潜空间+文本向量\n输出: ε_pred 预测噪声", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(8.60), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

# T5: 计算损失
rrect(sl, Inches(8.95), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(8.95), Inches(3.82), Inches(1.60), Inches(0.22),
   "计算损失+反向传播", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(8.95), Inches(4.07), Inches(1.60), Inches(0.40),
   "L = ||ε - ε_pred||²\n更新U-Net 860M参数", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# CLIP in training
rrect(sl, Inches(4.55), Inches(4.70), Inches(1.60), Inches(0.35), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(4.55), Inches(4.72), Inches(1.60), Inches(0.30),
   "CLIP编码文本(同时进行)", fs=8, fc=C_ACC, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(5.28), Inches(4.60), Inches(0.18), Inches(0.12), C_ACC)

# VAE解码器：训练时不需要
rrect(sl, Inches(10.60), Inches(3.80), Inches(2.20), Inches(0.75), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(10.60), Inches(3.85), Inches(2.20), Inches(0.65),
   "❌ VAE 解码器\n训练时不需要！\n（解码器只在推理时使用）", fs=9, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.1)

# ── 组件对照表 ──
rect(sl, Inches(0.15), Inches(5.30), Inches(12.80), Inches(0.28), fc=C_PRI)
tb(sl, Inches(0.30), Inches(5.31), Inches(12.50), Inches(0.26),
   "组件使用对照表", fs=11, fc=C_W, b=True)

# CLIP card
rrect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(1.20), fc=C_W, lc=C_ACC, lw=1)
rect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_ACC)
tb(sl, Inches(0.25), Inches(5.67), Inches(3.90), Inches(0.22),
   "CLIP 文本编码器 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(0.25), Inches(5.95), Inches(3.90), Inches(0.85),
   "• 翻译官 — 123M参数\n• 12层Transformer, 768维\n• 作用: 把人类语言翻译成AI数字语言", fs=9, fc=C_D, ls=1.15)

# U-Net card
rrect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(1.20), fc=C_W, lc=C_PRI, lw=1)
rect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_PRI)
tb(sl, Inches(4.50), Inches(5.67), Inches(3.90), Inches(0.22),
   "U-Net 去噪网络 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(4.50), Inches(5.95), Inches(3.90), Inches(0.85),
   "• 核心画师 — 860M参数(80%)\n• U型编码器-解码器结构\n• 作用: 一步步去掉噪声", fs=9, fc=C_D, ls=1.15)

# VAE card
rrect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(1.20), fc=C_W, lc=C_GRN, lw=1)
rect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(0.25), fc=C_GRN)
tb(sl, Inches(8.75), Inches(5.67), Inches(4.10), Inches(0.22),
   "VAE — 编码器(训练) 解码器(推理)", fs=10, fc=C_W, b=True)
tb(sl, Inches(8.75), Inches(5.95), Inches(4.10), Inches(0.85),
   "• 编码器~34M: 图像→潜空间 [训练]\n• 解码器~49M: 潜空间→图像 [推理]\n• 48倍压缩: 786K→16K数值", fs=9, fc=C_D, ls=1.15)

chk(sl, "P4 整体架构")


# ============================================================
# SLIDE 5 (FIX): CLIP 文本编码器详解 — 修复文字重叠
# ============================================================
print("\nBuilding fixed P5: CLIP详解...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "04  CLIP 文本编码器详解", "Contrastive Language-Image Pre-training — 从文字到语义向量的完整旅程")

# 顶部4个信息卡（间距拉大）
info_cards = [
    ("输入", "文本描述字符串\n分词器处理\n最多77个token\n约50-60个汉字", C_ACC),
    ("输出", "77×768维张量\n共59,136个浮点数\n文本的语义指纹\n包含所有含义信息", C_GRN),
    ("参数", "约1.23亿(123M)\n12层Transformer\n隐藏维度768\n12个注意力头", C_ORG),
    ("作用", "人类文字→数字向量\n通过Cross-Attention\n注入U-Net指导去噪\n没有它AI不知道画什么", C_PUR),
]
for i, (title, body, color) in enumerate(info_cards):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.25), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.10), Inches(1.50), Inches(2.80), Inches(0.90), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# CLIP 内部流程 — 4步骤（水平排列，每步独立框，间距充足）
steps = [
    ("① 分词", 
     "文本拆成最小语义单元(token)\n\n例: \"一只橘猫坐在沙发上\"\n→ [一只/橘/猫/坐在/沙发/上]\n\n每个token映射成768维向量",
     C_ACC),
    ("② 位置编码",
     "为什么需要位置编码？\n\"猫坐在沙发上\" ≠ \"沙发坐在猫上\"\n\n词一样，顺序不同，含义天差地别。\nTransformer本身不区分顺序，\n必须额外注入位置信息。",
     C_ORG),
    ("③ Transformer",
     "自注意力机制 (Self-Attention)\n每个词都能\"看到\"其他所有词\n\n\"橘猫\"→注意\"坐在\"+\"沙发\"\n\"坐在\"→连接\"橘猫\"+\"沙发\"\n\n12层: 浅层语法→深层语义",
     C_PUR),
    ("④ 输出向量",
     "输出: 77×768 语义向量\n\n每个token变成768维向量\n共77×768 = 59,136个浮点数\n\n不仅含每个词的独立含义，\n还包含词与词之间的关系信息。",
     C_GRN),
]

for i, (title, body, color) in enumerate(steps):
    x = Inches(0.30 + i * 3.20)
    # 步骤框
    rrect(sl, x, Inches(2.65), Inches(3.00), Inches(3.50), fc=C_W, lc=color, lw=1)
    # 标题栏
    rect(sl, x, Inches(2.65), Inches(3.00), Inches(0.28), fc=color)
    tb(sl, x + Inches(0.05), Inches(2.67), Inches(2.90), Inches(0.24), title, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    # 内容
    tb(sl, x + Inches(0.15), Inches(3.00), Inches(2.70), Inches(3.00), body, fs=9, fc=C_D, ls=1.2)
    
    # 箭头连接
    if i < len(steps) - 1:
        arrow_r(sl, x + Inches(3.00), Inches(4.20), Inches(0.20), Inches(0.18), color)

# 底部总结
rrect(sl, Inches(0.30), Inches(6.30), Inches(12.70), Inches(0.60), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(6.35), Inches(12.30), Inches(0.50),
   "总结: CLIP = 分词 → 位置编码 → 12层Transformer(自注意力) → 77×768语义向量。这个向量通过Cross-Attention注入U-Net，指导每一步去噪方向，确保生成的图像符合文字描述。",
   fs=10, fc=C_D, ls=1.2)

chk(sl, "P5 CLIP详解")


# ============================================================
# SLIDE 6 (FIX): VAE 编解码器详解 — 修复箭头+文字超框
# ============================================================
print("\nBuilding fixed P6: VAE详解...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "05  VAE 编解码器详解", "Variational Autoencoder — 压缩48倍但肉眼几乎看不出差别")

# 顶部4个信息卡
info_cards = [
    ("编码器输入", "原始图像\n512×512×3\n786,432个像素值\n值域 0-255", C_ACC),
    ("编码器输出", "潜空间\n64×64×4\n16,384个浮点数\n值域约[-5,5]", C_GRN),
    ("总参数量", "约8300万(83M)\n编码器~34M\n解码器~49M\n解码器更大", C_ORG),
    ("为什么用VAE?", "正则化确保\n潜空间连续平滑\n普通AE潜空间不规则\n扩散操作需要平滑空间", C_PUR),
]
for i, (title, body, color) in enumerate(info_cards):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.15), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.10), Inches(1.50), Inches(2.80), Inches(0.80), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# ── 编码器 (左侧) ──
rrect(sl, Inches(0.30), Inches(2.50), Inches(6.20), Inches(3.05), fc=C_W, lc=C_ACC, lw=1)
rect(sl, Inches(0.30), Inches(2.50), Inches(6.20), Inches(0.28), fc=C_ACC)
tb(sl, Inches(0.40), Inches(2.52), Inches(6.00), Inches(0.24),
   "▼ 编码器 (Encoder): 4层下采样 — 图像→潜空间（训练时使用）", fs=10, fc=C_W, b=True)

# 编码流程 — 每层一行，清晰的下行箭头
enc_steps = [
    ("512×512×3", "输入图像（RGB三通道）", C_ACC),
    ("256×256×128", "通道3→128，提取低级特征（边缘、色块）", C_ACC),
    ("128×128×256", "提取中级特征（纹理、局部结构）", C_ACC),
    ("64×64×512", "提取高级特征（物体部件、语义信息）", C_ACC),
    ("32×32×512", "最深层特征表示", C_ACC),
    ("64×64×4", "输出潜空间（48倍压缩！）", C_GRN),
]

for i, (dim, desc, color) in enumerate(enc_steps):
    y = Inches(2.85 + i * 0.42)
    # 尺寸标签
    rrect(sl, Inches(0.45), y, Inches(1.40), Inches(0.30), fc=C_LT, lc=color, lw=1)
    tb(sl, Inches(0.45), y + Inches(0.03), Inches(1.40), Inches(0.24), dim, fs=9, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    # 箭头
    if i < len(enc_steps) - 1:
        arrow_d(sl, Inches(1.10), y + Inches(0.30), Inches(0.15), Inches(0.12), C_ACC)
    else:
        arrow_d(sl, Inches(1.10), y + Inches(0.30), Inches(0.15), Inches(0.12), C_GRN)
    # 描述
    tb(sl, Inches(2.00), y + Inches(0.03), Inches(4.30), Inches(0.24), desc, fs=9, fc=C_D)

# 底部注释
tb(sl, Inches(0.45), Inches(5.45), Inches(5.90), Inches(0.20),
   "编码器参数量: ~34M | 宽高每次缩小一半，通道数逐步增加", fs=8, fc=C_G)

# ── 解码器 (右侧) — 从底到顶: 64×64→512×512 ──
rrect(sl, Inches(6.80), Inches(2.50), Inches(6.20), Inches(3.05), fc=C_W, lc=C_GRN, lw=1)
rect(sl, Inches(6.80), Inches(2.50), Inches(6.20), Inches(0.28), fc=C_GRN)
tb(sl, Inches(6.90), Inches(2.52), Inches(6.00), Inches(0.24),
   "▲ 解码器 (Decoder): 4层上采样 — 潜空间→图像（推理时使用）", fs=10, fc=C_W, b=True)

dec_steps = [
    ("64×64×4", "输入潜空间", C_GRN),
    ("128×128×512", "通道4→512，开始重建空间细节", C_GRN),
    ("256×256×256", "通道减半，分辨率翻倍", C_GRN),
    ("512×512×128", "继续增加分辨率，减少通道", C_GRN),
    ("512×512×64", "接近最终分辨率", C_GRN),
    ("512×512×3", "输出高清图像（肉眼几乎无损！）", C_ORG),
]

# Place decoder bottom-to-top: 64×64 at bottom, 512×512 at top
dec_y_start = 2.85 + (len(dec_steps) - 1) * 0.42  # bottom of last item
for i, (dim, desc, color) in enumerate(dec_steps):
    y = Inches(dec_y_start - i * 0.42)  # reverse: bottom to top
    rrect(sl, Inches(6.95), y, Inches(1.40), Inches(0.30), fc=C_BG_GRN, lc=color, lw=1)
    tb(sl, Inches(6.95), y + Inches(0.03), Inches(1.40), Inches(0.24), dim, fs=9, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    # 上箭头: from current row top → pointing to next row above
    if i < len(dec_steps) - 1:
        # Arrow between current row and next row above
        arrow_u(sl, Inches(7.60), y - Inches(0.12), Inches(0.15), Inches(0.12), C_GRN)
    # 描述
    tb(sl, Inches(8.50), y + Inches(0.03), Inches(4.30), Inches(0.24), desc, fs=9, fc=C_D)

tb(sl, Inches(6.95), Inches(5.45), Inches(5.90), Inches(0.20),
   "解码器参数量: ~49M（比编码器大，因为\"还原\"比\"压缩\"更复杂）", fs=8, fc=C_G)

# ── 底部对比 ──
rrect(sl, Inches(0.30), Inches(5.75), Inches(12.70), Inches(1.15), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.78), Inches(12.30), Inches(0.22),
   "💡 为什么用VAE而不是普通自编码器(AE)？", fs=11, fc=C_PRI, b=True)
tb(sl, Inches(0.50), Inches(6.05), Inches(12.30), Inches(0.80),
   "普通AE的问题: 潜空间\"不规则\" — 两个相近的潜空间点可能映射到完全不相关的图像，潜空间中存在大量\"空白区域\"。\n"
   "在 irregular 空间做扩散(加减噪声)会导致解码出无意义的图像。VAE 通过 KL 散度正则化强制潜空间服从标准正态分布，\n"
   "使其连续、平滑、无空洞 — 这是扩散模型能在潜空间上正常工作的前提条件。",
   fs=9, fc=C_D, ls=1.2)

chk(sl, "P6 VAE详解")


# ============================================================
# SLIDE 7 (FIX): U-Net 架构图 — 框闭合+箭头指向明确
# ============================================================
print("\nBuilding fixed P7: U-Net架构...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "06  U-Net 架构图", "编码器-瓶颈-解码器 U 型结构 | 860M 参数 | 3 个关键机制")

# ── 整体U-Net区域框 ──
rrect(sl, Inches(0.10), Inches(1.10), Inches(6.50), Inches(5.85), fc=None, lc=C_PRI, lw=2)

# ── 编码器区域框 ──
rrect(sl, Inches(0.20), Inches(1.20), Inches(2.50), Inches(4.50), fc=None, lc=C_ACC, lw=1)
tb(sl, Inches(0.30), Inches(1.22), Inches(2.30), Inches(0.22),
   "编码器 Encoder ↓", fs=11, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)

enc_layers = [
    ("64×64×320", "Res+Attn+Down", "Cross-Attention"),
    ("32×32×640", "Res+Attn+Down", "Cross-Attention"),
    ("16×16×1280", "Res+Attn+Down", "Cross-Attention"),
    ("8×8×1280", "Res+Down", "无注意力"),
]
enc_y = 1.50
for dim, ops, attn in enc_layers:
    rrect(sl, Inches(0.35), Inches(enc_y), Inches(2.20), Inches(0.70), fc=C_W, lc=C_ACC, lw=1)
    tb(sl, Inches(0.40), Inches(enc_y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.40), Inches(enc_y+0.28), Inches(2.10), Inches(0.20), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.40), Inches(enc_y+0.48), Inches(2.10), Inches(0.18), f"{'✦ ' if 'Cross' in attn else ''}{attn}", fs=8, fc=C_PUR if 'Cross' in attn else C_G, a=PP_ALIGN.CENTER)
    enc_y += 0.85
    if 'Cross' in attn:
        arrow_d(sl, Inches(1.40), Inches(enc_y-0.85+0.70), Inches(0.18), Inches(0.15), C_ACC)

# Down arrow after last encoder layer
arrow_d(sl, Inches(1.40), Inches(enc_y-0.85+0.70), Inches(0.18), Inches(0.15), C_ACC)

# ── Bottleneck ──
bottleneck_y = enc_y
rrect(sl, Inches(0.30), Inches(bottleneck_y), Inches(2.30), Inches(0.70), fc=C_BG_ORG, lc=C_ORG, lw=2)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.02), Inches(2.20), Inches(0.22), "4×4×1280", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.25), Inches(2.20), Inches(0.20), "Bottleneck", fs=9, fc=C_D, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.45), Inches(2.20), Inches(0.20), "Res+SelfAttn+Res", fs=8, fc=C_D, a=PP_ALIGN.CENTER)

# ── 解码器区域框 ──
rrect(sl, Inches(3.00), Inches(1.20), Inches(2.50), Inches(4.50), fc=None, lc=C_GRN, lw=1)
tb(sl, Inches(3.10), Inches(1.22), Inches(2.30), Inches(0.22),
   "解码器 Decoder ↑", fs=11, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

dec_layers = [
    ("8×8×1280", "Res+Up+Skip", "Cross-Attention"),
    ("16×16×640", "Res+Up+Skip", "Cross-Attention"),
    ("32×32×320", "Res+Up+Skip", "Cross-Attention"),
    ("64×64×4", "输出预测噪声", "无注意力"),
]

# Decoder placed bottom-up. Match with encoder:
# enc: 64(1.50), 32(2.35), 16(3.20), 8(4.05)  bottleneck(4.90)
# dec: 8(4.05), 16(3.20), 32(2.35), 64(1.50) — same Y as encoder!

dec_ys = [4.05, 3.20, 2.35, 1.50]  # Bottom to top, matching encoder

for i, (dim, ops, attn) in enumerate(dec_layers):
    y = dec_ys[i]
    rrect(sl, Inches(3.20), Inches(y), Inches(2.20), Inches(0.70), fc=C_BG_GRN, lc=C_GRN, lw=1)
    tb(sl, Inches(3.25), Inches(y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(3.25), Inches(y+0.28), Inches(2.10), Inches(0.20), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(3.25), Inches(y+0.48), Inches(2.10), Inches(0.18), f"{'✦ ' if 'Cross' in attn else ''}{attn}", fs=8, fc=C_PUR if 'Cross' in attn else C_G, a=PP_ALIGN.CENTER)
    
    # UP arrow between decoder layers
    if i < len(dec_layers) - 1:
        next_y = dec_ys[i+1]
        arrow_u(sl, Inches(4.25), Inches(y-0.15), Inches(0.18), Inches(0.15), C_GRN)
        tb(sl, Inches(4.50), Inches(y-0.18), Inches(0.90), Inches(0.18),
           "上采样↑", fs=8, fc=C_GRN, b=True)

# Bottleneck → first decoder layer (8×8) — clear arrow path
# Bottleneck is at ~y=4.90, decoder 8×8 is at y=4.05
# Draw arrow from bottleneck right edge → curving up to decoder 8×8
# Use a clear L-shaped path: right from bottleneck, then up to 8×8 level
# Step 1: horizontal line from bottleneck to gap area
line_h(sl, Inches(2.60), Inches(4.90 + 0.30), Inches(0.60), C_GRN, 3)
# Step 2: vertical line going up to decoder 8×8 level  
rect(sl, Inches(3.10), Inches(4.05 + 0.30), Emu(int(3*12700)), Inches(0.85), fc=C_GRN)
# Step 3: horizontal arrow into decoder 8×8 box
arrow_r(sl, Inches(3.10), Inches(4.05 + 0.25), Inches(0.15), Inches(0.12), C_GRN)
# Label
tb(sl, Inches(2.55), Inches(4.90 - 0.18), Inches(0.65), Inches(0.18),
   "过渡→", fs=7, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

# ── Skip Connections ──
# Connect encoder → decoder at same resolution levels
for i in range(4):
    enc_mid_y = Inches(1.50 + i * 0.85) + Inches(0.30)
    # Orange horizontal line
    line_h(sl, Inches(2.55), enc_mid_y, Inches(0.65), C_ORG, 3)
    # "Skip" label
    rrect(sl, Inches(2.60), enc_mid_y - Inches(0.18), Inches(0.55), Inches(0.18), fc=C_ORG)
    tb(sl, Inches(2.60), enc_mid_y - Inches(0.18), Inches(0.55), Inches(0.18),
       "Skip", fs=7, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    # Arrow into decoder
    arrow_r(sl, Inches(3.15), enc_mid_y - Inches(0.04), Inches(0.12), Inches(0.10), C_ORG)

# ── Time Embedding (left) ──
rrect(sl, Inches(5.70), Inches(1.20), Inches(0.80), Inches(5.30), fc=None, lc=C_TEL, lw=1)
tb(sl, Inches(5.72), Inches(1.25), Inches(0.76), Inches(0.22), "Time", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(5.72), Inches(1.45), Inches(0.76), Inches(0.22), "Embed", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)
for y_pos in [1.80, 2.65, 3.50, 4.35, 5.20]:
    arrow_l = arrow_r  # reuse right arrow
    line_h(sl, Inches(5.70), Inches(y_pos), Inches(0.80), C_TEL, 1)

# ── 右面板: 3个关键机制 ──
# ① Cross-Attention
rrect(sl, Inches(6.80), Inches(1.20), Inches(6.20), Inches(1.65), fc=C_W, lc=C_PUR, lw=1)
rect(sl, Inches(6.80), Inches(1.20), Inches(6.20), Inches(0.28), fc=C_PUR)
tb(sl, Inches(6.90), Inches(1.22), Inches(6.00), Inches(0.24),
   "① Cross-Attention 交叉注意力 — 文本条件注入的核心机制", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(1.55), Inches(6.00), Inches(1.20), [
    "• 图像特征作为 Query，文本特征作为 Key-Value",
    "• 图像问：\"根据文字，我这里该怎样？\" → 文本答：\"你这里应该是橘猫的毛发\"",
    "• 出现位置：编码器前3层 + 解码器前3层（共6处）",
    "• 注意力计算: Q·K^T/√d → Softmax → ×V",
], fs=9, fc=C_D, ls=1.15)

# ② Skip Connections
rrect(sl, Inches(6.80), Inches(3.00), Inches(6.20), Inches(1.65), fc=C_W, lc=C_ORG, lw=1)
rect(sl, Inches(6.80), Inches(3.00), Inches(6.20), Inches(0.28), fc=C_ORG)
tb(sl, Inches(6.90), Inches(3.02), Inches(6.00), Inches(0.24),
   "② Skip Connections 跳跃连接 — 信息传递的桥梁（图中橙色横线）", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(3.35), Inches(6.00), Inches(1.20), [
    "• 编码器特征直接桥接到解码器对应层（同分辨率级别）",
    "• 防止细节在压缩过程中丢失",
    "• 像建筑师随时参考草稿本，没有它解码器只能\"凭记忆\"还原",
    "• 连接: 64×64↔64×64, 32×32↔32×32, 16×16↔16×16, 8×8↔8×8",
], fs=9, fc=C_D, ls=1.15)

# ③ Time Embedding
rrect(sl, Inches(6.80), Inches(4.80), Inches(6.20), Inches(1.85), fc=C_W, lc=C_TEL, lw=1)
rect(sl, Inches(6.80), Inches(4.80), Inches(6.20), Inches(0.28), fc=C_TEL)
tb(sl, Inches(6.90), Inches(4.82), Inches(6.00), Inches(0.24),
   "③ Time Embedding 时间嵌入 — 当前噪声水平的标尺", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(5.15), Inches(6.00), Inches(1.40), [
    "• 告诉U-Net当前噪声水平（1-1000步）",
    "• 高噪声(t≈900)：粗雕确定轮廓",
    "• 低噪声(t≈100)：精雕添加细节",
    "• 正弦编码 → 2层MLP → 1280维",
    "• 注入每一个ResBlock，影响去噪策略",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P7 U-Net架构")


# ============================================================
# Copy slide 8 (协作流程) — unchanged
# ============================================================
print("\nCopying slide 8 (协作流程)...")
copy_slide(prs, src_prs.slides[7])
print(f"  Copied slide 8")

# ============================================================
# SLIDE 9 (REBUILD): 训练过程流程图 — 修复文字重叠
# ============================================================
print("\nRebuilding P9: 训练过程...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "训练过程流程图", "6步训练管线 — 15万GPU小时，50-60万美元")

steps_data = [
    ("Step 1", "采样图像+文本", "从LAION-5B随机抽取\n图像x₀ + 文本描述", "6亿图文对", C_ACC),
    ("Step 2", "VAE 编码", "512×512×3 → 64×64×4\n压缩到潜空间(48倍)", "z₀ = VAE(x₀)", C_GRN),
    ("Step 3", "添加随机噪声", "随机选t ∈ [1,1000]\n按比例添加高斯噪声", "z_t = √ᾱ·z₀ + √(1-ᾱ)·ε", C_RED),
    ("Step 4", "U-Net 预测噪声", "输入: z_t + t + 文本向量\n输出: 预测噪声 ε_pred", "ε_pred = U-Net(z_t, t, c)", C_PRI),
    ("Step 5", "计算损失", "预测噪声 vs 真实噪声\n均方误差(MSE)", "L = ‖ε - ε_pred‖²", C_ORG),
    ("Step 6", "反向传播", "更新U-Net的860M参数\nCLIP + VAE保持冻结", "θ ← θ - η·∇L", C_PUR),
]

for i, (step, title, desc, formula, color) in enumerate(steps_data):
    col = i % 3
    row = i // 3
    x = Inches(0.50 + col * 4.20)
    y = Inches(1.30 + row * 2.30)
    
    # Card background
    rrect(sl, x, y, Inches(3.80), Inches(2.00), fc=C_W, lc=color, lw=1)
    # Step badge
    rect(sl, x, y, Inches(0.70), Inches(0.30), fc=color)
    tb(sl, x + Inches(0.05), y + Inches(0.03), Inches(0.60), Inches(0.24), step, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    # Title
    tb(sl, x + Inches(0.80), y + Inches(0.03), Inches(2.90), Inches(0.24), title, fs=11, fc=color, b=True)
    # Description
    tb(sl, x + Inches(0.15), y + Inches(0.40), Inches(3.50), Inches(0.80), desc, fs=10, fc=C_D, ls=1.2)
    # Formula badge
    rrect(sl, x + Inches(0.15), y + Inches(1.35), Inches(3.50), Inches(0.45), fc=C_LT, lc=color, lw=1)
    tb(sl, x + Inches(0.25), y + Inches(1.40), Inches(3.30), Inches(0.35), formula, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER, fn="Consolas")
    
    # Arrow to next step (horizontal)
    if col < 2 and i < 5:
        arrow_r(sl, x + Inches(3.80), y + Inches(0.85), Inches(0.30), Inches(0.18), color)

# "继续→" arrow from Step 3 → Step 4
arrow_d(sl, Inches(10.50), y + Inches(0.15) if i == 2 else Inches(3.40), Inches(0.18), Inches(0.18), C_ACC)

# Loop indicator
rrect(sl, Inches(0.30), Inches(5.80), Inches(12.70), Inches(0.40), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.85), Inches(12.30), Inches(0.30),
   "⟳ 重复训练循环：重复 Step 1-6 数十万次，直到损失收敛 | 每次迭代随机采样不同的图像和噪声水平",
   fs=10, fc=C_D, b=False, a=PP_ALIGN.CENTER)

# Training params table
table_shape = sl.shapes.add_table(2, 5, Inches(0.50), Inches(6.35), Inches(12.30), Inches(0.60))
table = table_shape.table
headers_t = ["数据集", "GPU", "训练时间", "成本", "Batch Size"]
vals_t = ["LAION-5B (60亿图文对子集)", "NVIDIA A100 × 256", "~15万GPU小时", "$50-60万", "2048×2048"]
for j in range(5):
    table.columns[j].width = Inches(2.46)
    cell = table.cell(0, j); cell.text = headers_t[j]
    for p in cell.text_frame.paragraphs: p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=C_W; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
    cell = table.cell(1, j); cell.text = vals_t[j]
    for p in cell.text_frame.paragraphs: p.font.size=Pt(9); p.font.color.rgb=C_D; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_W

chk(sl, "P9 训练过程")


# ============================================================
# SLIDE 10 (was 11): 推理过程 — 保持原样（先于CFG）
# ============================================================
print("\nCopying slide 11 (推理过程) as new slide 10...")
copy_slide(prs, src_prs.slides[10])  # 0-indexed: slide 11


# ============================================================
# SLIDE 11 (FIX, was P10): CFG — 修复箭头+文字+移到推理后
# ============================================================
print("\nBuilding fixed P11: CFG详解 (moved after 推理过程)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "09  CFG（Classifier-Free Guidance）详解", "无分类器引导 — 让AI更听话地按描述画画")

# ── 引导文字 ──
rrect(sl, Inches(0.30), Inches(1.20), Inches(12.70), Inches(0.85), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(0.50), Inches(1.22), Inches(12.30), Inches(0.22),
   "🤔 为什么需要CFG？（紧接推理过程）", fs=12, fc=C_RED, b=True)
tb(sl, Inches(0.50), Inches(1.48), Inches(12.30), Inches(0.50),
   "在推理去噪循环中，U-Net每一步预测噪声时，需要知道\"该往哪个方向去噪\"。没有CFG时，U-Net只知道\"去掉噪声\"，"
   "生成的图像随机且不可控。CFG通过对比\"有文字引导\"和\"无文字引导\"两个预测，放大差异，让生成结果更忠实于文字描述。",
   fs=10, fc=C_D, ls=1.2)

# ── 核心公式 ──
rrect(sl, Inches(0.30), Inches(2.20), Inches(12.70), Inches(1.10), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(2.23), Inches(12.30), Inches(0.22),
   "核心思想：既知道什么是对的，也知道什么不是对的", fs=11, fc=C_PRI, b=True)
mtb(sl, Inches(0.50), Inches(2.50), Inches(12.30), Inches(0.70), [
    "公式: ε_guided = ε_uncond + s × (ε_cond - ε_uncond)",
    "ε_uncond = 没有文本条件时的预测（自由发挥）  |  ε_cond = 有文本条件时的预测（按指令画）  |  s = CFG Scale（通常7-12）",
], fs=10, fc=C_D, ls=1.2)

# ── 左侧: CFG Scale 参数表 ──
table_shape = sl.shapes.add_table(5, 3, Inches(0.30), Inches(3.45), Inches(6.50), Inches(1.70))
table = table_shape.table
table.columns[0].width = Inches(1.20)
table.columns[1].width = Inches(2.50)
table.columns[2].width = Inches(2.80)

headers = ["CFG Scale", "效果", "适用场景"]
data = [
    ["1-3", "几乎没有引导，图像\"自由发挥\"", "创意探索，想要意外惊喜"],
    ["5-7", "适中引导，平衡忠实度和创造力", "通用场景，推荐默认值 7"],
    ["7-12", "强引导，高度忠实于文字描述", "需要精确控制，如产品图"],
    ["15-20+", "过强引导，过度饱和出现伪影", "一般不推荐"],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# ── 右侧: CFG工作原理图 (修复箭头方向) ──
rrect(sl, Inches(7.10), Inches(3.45), Inches(5.90), Inches(2.80), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(7.30), Inches(3.50), Inches(5.50), Inches(0.25),
   "CFG 工作原理图", fs=11, fc=C_PRI, b=True)

# 正确的流程:
# ε_uncond 和 ε_cond → 计算差异 → × CFG Scale → + ε_uncond = ε_guided

# 无条件生成 (左上)
rrect(sl, Inches(7.30), Inches(3.90), Inches(2.20), Inches(0.50), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(7.30), Inches(3.95), Inches(2.20), Inches(0.40),
   "无条件生成（自由发挥）\nε_uncond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)

# 有条件生成 (右上)
rrect(sl, Inches(10.40), Inches(3.90), Inches(2.20), Inches(0.50), fc=C_BG_GRN, lc=C_GRN, lw=1)
tb(sl, Inches(10.40), Inches(3.95), Inches(2.20), Inches(0.40),
   "有条件生成（按指令画）\nε_cond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)

# 差异计算框 (中间)
# ε_cond - ε_uncond: 两个都指向中间
arrow_d(sl, Inches(8.35), Inches(4.40), Inches(0.15), Inches(0.15), C_ORG)
arrow_d(sl, Inches(11.45), Inches(4.40), Inches(0.15), Inches(0.15), C_ORG)

rrect(sl, Inches(9.30), Inches(4.60), Inches(1.60), Inches(0.40), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(9.30), Inches(4.63), Inches(1.60), Inches(0.35),
   "计算差异\nε_cond - ε_uncond", fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0, b=True)

# × CFG Scale
arrow_d(sl, Inches(10.05), Inches(5.00), Inches(0.15), Inches(0.12), C_ORG)
rrect(sl, Inches(9.30), Inches(5.15), Inches(1.60), Inches(0.35), fc=C_W, lc=C_ORG, lw=1)
tb(sl, Inches(9.30), Inches(5.18), Inches(1.60), Inches(0.30),
   "× CFG Scale (s)", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)

# → 结果
arrow_d(sl, Inches(10.05), Inches(5.50), Inches(0.15), Inches(0.12), C_GRN)
rrect(sl, Inches(9.00), Inches(5.65), Inches(2.20), Inches(0.45), fc=C_BG_GRN, lc=C_GRN, lw=2)
tb(sl, Inches(9.00), Inches(5.68), Inches(2.20), Inches(0.40),
   "= 引导后的高质量结果\nε_guided", fs=9, fc=C_D, b=True, a=PP_ALIGN.CENTER, ls=1.0)

# + ε_uncond 说明
tb(sl, Inches(7.30), Inches(5.75), Inches(1.60), Inches(0.30),
   "← + ε_uncond (加回\n     无条件结果)", fs=8, fc=C_G, a=PP_ALIGN.CENTER, ls=1.0)
arrow_r(sl, Inches(8.95), Inches(5.85), Inches(0.10), Inches(0.10), C_ACC)

# ── 底部比喻 ──
rrect(sl, Inches(0.30), Inches(5.55), Inches(6.50), Inches(1.35), fc=C_LT, lc=C_ACC, lw=1)
mtb(sl, Inches(0.50), Inches(5.58), Inches(6.10), Inches(1.25), [
    "比喻: CFG像严厉又公正的艺术指导。",
    "每画一笔问两个问题:",
    "  \"按客户要求画了吗？\"",
    "  \"不管要求你会怎么画？\"",
    "差异放大(×CFG Scale)，强制遵循要求。",
    "没有CFG→模糊偏离 | CFG过大→过度刻意",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P11 CFG详解")


# ============================================================
# Copy remaining slides: 12(采样器+参数), 13(Prompt+LoRA+ControlNet), 14(SD3+Flux), 15(趋势)
# ============================================================
print("\nCopying remaining slides...")
# Original indices: 11(采样器), 12(参数+Prompt), 13(应用+LoRA+ControlNet), 14(SD3+Flux), 15(趋势)
for idx in [11, 12, 13, 14]:
    copy_slide(prs, src_prs.slides[idx])
    print(f"  Copied slide {idx+1}")


# ============================================================
# Save
# ============================================================
print(f"\nSaving to: {OUT}")
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)

print(f"\n{'='*60}")
print(f"RESULT: {len(prs.slides)} slides")
print(f"{'='*60}")
print("Slide order:")
print("  P1  = 封面 (copy)")
print("  P2  = 目录 (copy)")
print("  P3  = Diffusion基础 (copy)")
print("  P4  = 整体架构概览 (FIXED)")
print("  P5  = CLIP详解 (FIXED)")
print("  P6  = VAE详解 (FIXED)")
print("  P7  = U-Net架构 (FIXED)")
print("  P8  = 三大组件协作 (copy)")
print("  P9  = 训练过程 (copy)")
print("  P10 = 推理过程 (moved before CFG)")
print("  P11 = CFG详解 (FIXED, moved after 推理)")
print("  P12 = 参数优化+Prompt (copy)")
print("  P13 = 应用+LoRA+ControlNet (copy)")
print("  P14 = SD3+Flux (copy)")
print("  P15 = 趋势+总结 (copy)")

# Final check
print(f"\n{'='*60}")
print("FINAL VERIFICATION:")
print(f"{'='*60}")
for i, slide in enumerate(prs.slides):
    chk(slide, f"Slide {i+1}")

print(f"\nDone! Output: {OUT}")
