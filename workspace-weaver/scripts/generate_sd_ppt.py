#!/usr/bin/env python3
"""
Stable Diffusion 技术解析 PPT 生成器
v2.0 — 基于子代理调研结果 + 血泪教训总结
输出: .pptx 格式，15页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 配置 ====================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 颜色
C_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)      # 深海蓝
C_ACCENT = RGBColor(0x34, 0x98, 0xDB)        # 科技蓝
C_LIGHT = RGBColor(0xEB, 0xF5, 0xFB)        # 浅蓝背景
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x33, 0x33, 0x33)         # 深灰文字
C_GRAY = RGBColor(0x88, 0x88, 0x88)         # 辅助灰
C_GREEN = RGBColor(0x27, 0xAE, 0x60)        # 绿色
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)       # 橙色
C_RED = RGBColor(0xE7, 0x4C, 0x3C)          # 红色
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)       # 紫色
C_TEAL = RGBColor(0x00, 0x96, 0x88)         # 青色

FONT_CN = "微软雅黑"
FONT_EN = "Arial"

OUTPUT_PATH = "/home/admin/.openclaw/workspace-weaver/output/Stable_Diffusion_技术解析.pptx"
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ==================== 工具函数 ====================

def add_bg(slide, color=C_WHITE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=12, font_color=C_DARK,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.15):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=11,
                          font_color=C_DARK, font_name=FONT_CN, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.1, bullet=False):
    """添加多行文本框，lines是[(text, bold, color, size)]列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, color, size = line_data, False, font_color, font_size
        elif len(line_data) == 2:
            text, bold = line_data
            color, size = font_color, font_size
        elif len(line_data) == 3:
            text, bold, color = line_data
            size = font_size
        else:
            text, bold, color, size = line_data
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet and not bold:
            p.text = "  • " + text
        else:
            p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(3)
        p.line_spacing = line_spacing
    return txBox

def add_arrow_right(slide, left, top, width, height, color=C_ACCENT):
    """添加右箭头"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width, height, color=C_ACCENT):
    """添加下箭头"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, color=C_ACCENT):
    """添加 chevron 箭头"""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_header(slide, title, subtitle=""):
    """添加页眉"""
    # 顶部色带
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=C_PRIMARY)
    # 标题
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
                title, font_size=28, font_color=C_WHITE, bold=True, font_name=FONT_CN)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.35),
                    subtitle, font_size=13, font_color=RGBColor(0xBB, 0xDE, 0xFB), font_name=FONT_CN)
    # 页码区域（底部）
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), fill_color=C_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(7.15), Inches(5), Inches(0.35),
                "Stable Diffusion 技术解析", font_size=9, font_color=RGBColor(0x88, 0xBB, 0xDD), font_name=FONT_CN)

def make_table(slide, left, top, width, rows, cols, data, col_widths=None,
               header_color=C_PRIMARY, header_font_color=C_WHITE, font_size=10,
               alt_color=C_LIGHT):
    """创建表格"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.35 * rows))
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_CN
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.space_after = Pt(1)
            
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = header_font_color
                    paragraph.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
            
            # 边框
            from pptx.oxml.ns import qn
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                border = tcPr.find(qn(border_name))
                if border is None:
                    from lxml import etree
                    border = etree.SubElement(tcPr, qn(border_name))
                border.set('w', '6350')
                solidFill = border.find(qn('a:solidFill'))
                if solidFill is None:
                    from lxml import etree
                    solidFill = etree.SubElement(border, qn('a:solidFill'))
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is None:
                    from lxml import etree
                    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', 'D0D0D0')
    
    return table_shape

# ==================== 第1页：封面 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_PRIMARY)

# 装饰线条
add_rect(slide, Inches(0.8), Inches(2.8), Inches(3), Inches(0.06), fill_color=C_ACCENT)

# 标题
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(1.2),
            "Stable Diffusion", font_size=52, font_color=C_WHITE, bold=True, font_name=FONT_EN)
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.8),
            "技术解析 — 从原理到应用", font_size=30, font_color=C_ACCENT, bold=False, font_name=FONT_CN)

# 装饰线条
add_rect(slide, Inches(0.8), Inches(5.0), Inches(1.5), Inches(0.04), fill_color=C_ACCENT)

# 副标题信息
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(8), Inches(0.4),
            "深入理解扩散模型、CLIP、VAE、U-Net 三大核心组件", font_size=14, font_color=RGBColor(0xAA, 0xCC, 0xEE), font_name=FONT_CN)
add_textbox(slide, Inches(0.8), Inches(5.7), Inches(8), Inches(0.4),
            "2026 年 4 月", font_size=13, font_color=C_GRAY, font_name=FONT_CN)

# 右下角装饰
add_rect(slide, Inches(10), Inches(0), Inches(3.333), SLIDE_H, fill_color=RGBColor(0x16, 0x2D, 0x4A))
add_textbox(slide, Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.5),
            "AI IMAGE\nGENERATION", font_size=16, font_color=RGBColor(0x44, 0x66, 0x88), bold=True, font_name=FONT_EN, alignment=PP_ALIGN.RIGHT)

print("✅ 第1页：封面")

# ==================== 第2页：目录 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "目录", "CONTENTS")

toc_items = [
    ("01", "Diffusion Model 基础原理", "墨水扩散与逆向去噪"),
    ("02", "Stable Diffusion vs 普通 Diffusion", "48倍压缩，让AI画画人人可用"),
    ("03", "整体架构概览", "CLIP + U-Net + VAE 三大组件"),
    ("04", "CLIP 文本编码器", "把文字翻译成AI能理解的数字"),
    ("05", "VAE 编解码器", "超级智能的ZIP压缩软件"),
    ("06", "U-Net 去噪网络", "核心画师，8.6亿参数"),
    ("07", "三大组件协作流程", "AI画图流水线"),
    ("08", "训练过程", "6亿图像-文本对的炼丹之旅"),
    ("09", "CFG 引导技术", "无分类器引导，让AI更听话"),
    ("10", "推理过程", "从噪声到图像的50步雕刻"),
    ("11", "采样器对比", "9种算法，速度与质量取舍"),
    ("12", "参数优化", "5个核心参数的最佳实践"),
    ("13", "应用场景", "5大领域的革命性影响"),
]

start_y = Inches(1.35)
for i, (num, title, desc) in enumerate(toc_items):
    y = start_y + Inches(i * 0.44)
    # 序号
    add_textbox(slide, Inches(0.8), y, Inches(0.6), Inches(0.4),
                num, font_size=14, font_color=C_ACCENT, bold=True, font_name=FONT_EN)
    # 标题
    add_textbox(slide, Inches(1.5), y, Inches(4), Inches(0.4),
                title, font_size=13, font_color=C_DARK, bold=True, font_name=FONT_CN)
    # 描述
    add_textbox(slide, Inches(5.8), y, Inches(6.5), Inches(0.4),
                desc, font_size=11, font_color=C_GRAY, font_name=FONT_CN)
    # 分隔线
    if i < len(toc_items) - 1:
        add_rect(slide, Inches(0.8), y + Inches(0.4), Inches(11.5), Inches(0.01), fill_color=RGBColor(0xE0, 0xE0, 0xE0))

print("✅ 第2页：目录")

# ==================== 第3页：Diffusion Model 基础原理 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "01  Diffusion Model 基础原理", "墨水扩散与逆向去噪")

# 左侧：前向扩散图示
add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.6), fill_color=RGBColor(0xFD, 0xF2, 0xE9))
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.35),
            "前向扩散：清晰照片 → 纯噪声（雪花屏）", font_size=13, font_color=C_ORANGE, bold=True)

# 流程图：清晰图 → 噪声图
add_rounded_rect(slide, Inches(0.9), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=C_WHITE, line_color=C_ORANGE)
add_textbox(slide, Inches(0.9), Inches(2.15), Inches(1.2), Inches(0.5),
            "清晰\n照片", font_size=11, font_color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(2.2), Inches(2.25), Inches(0.5), Inches(0.3), C_ORANGE)
add_textbox(slide, Inches(2.2), Inches(2.0), Inches(0.5), Inches(0.25),
            "+噪声", font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(2.8), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=C_WHITE, line_color=C_ORANGE)
add_textbox(slide, Inches(2.8), Inches(2.15), Inches(1.2), Inches(0.5),
            "轻微\n模糊", font_size=11, font_color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(4.1), Inches(2.25), Inches(0.5), Inches(0.3), C_ORANGE)
add_textbox(slide, Inches(4.1), Inches(2.0), Inches(0.5), Inches(0.25),
            "+噪声", font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(4.7), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=RGBColor(0xE0, 0xE0, 0xE0), line_color=C_ORANGE)
add_textbox(slide, Inches(4.7), Inches(2.15), Inches(1.2), Inches(0.5),
            "纯噪声\n(雪花屏)", font_size=11, font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# 说明
add_textbox(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(0.8),
            "就像把一滴蓝色墨水滴入清水中，墨水逐渐扩散，最终整杯水变成均匀的浅蓝色。\n"
            "这个过程需要 1000 步，图像从 786,432 个有意义的像素退化成完全随机的数值。",
            font_size=10, font_color=C_DARK, line_spacing=1.2)

# 右侧：反向去噪图示
add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.6), fill_color=RGBColor(0xE8, 0xF8, 0xF5))
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.35),
            "反向去噪：纯噪声 → 清晰照片（AI修复）", font_size=13, font_color=C_GREEN, bold=True)

add_rounded_rect(slide, Inches(7.1), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=RGBColor(0xE0, 0xE0, 0xE0), line_color=C_GREEN)
add_textbox(slide, Inches(7.1), Inches(2.15), Inches(1.2), Inches(0.5),
            "纯噪声\n(起点)", font_size=11, font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(8.4), Inches(2.25), Inches(0.5), Inches(0.3), C_GREEN)
add_textbox(slide, Inches(8.4), Inches(2.0), Inches(0.5), Inches(0.25),
            '去噪', font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(9.0), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=C_WHITE, line_color=C_GREEN)
add_textbox(slide, Inches(9.0), Inches(2.15), Inches(1.2), Inches(0.5),
            "模糊\n轮廓", font_size=11, font_color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(10.3), Inches(2.25), Inches(0.5), Inches(0.3), C_GREEN)
add_textbox(slide, Inches(10.3), Inches(2.0), Inches(0.5), Inches(0.25),
            '去噪', font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(10.9), Inches(2.0), Inches(1.2), Inches(0.8), fill_color=C_WHITE, line_color=C_GREEN)
add_textbox(slide, Inches(10.9), Inches(2.15), Inches(1.2), Inches(0.5),
            "清晰\n照片", font_size=11, font_color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(7.1), Inches(2.9), Inches(5.2), Inches(0.8),
            "AI修复师面对雪花屏，一步一步去掉噪声，每步恢复一点细节，最终还原清晰照片。\n"
            '这个去噪能力就是扩散模型通过海量训练数据学到的核心技能。',
            font_size=10, font_color=C_DARK, line_spacing=1.2)

# 底部：为什么比GAN更好
add_rounded_rect(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.6), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5), Inches(0.35),
            "为什么扩散模型比 GAN 更好？", font_size=14, font_color=C_PRIMARY, bold=True)

lines = [
    ("GAN 的工作方式：两个对手博弈", True, C_RED, 12),
    ("  • \"造假者\"（生成器）试图制造逼真假图", False, C_DARK, 10),
    ("  • \"鉴定师\"（判别器）试图分辨真假", False, C_DARK, 10),
    ("  • 问题：训练不稳定，容易陷入\"死循环\"，多样性受限", False, C_RED, 10),
    ("", False, C_DARK, 6),
    ("扩散模型的方式：一步一步走", True, C_GREEN, 12),
    ("  • 把图像生成变成一系列可预测的小步骤", False, C_DARK, 10),
    ("  • 就像走路和跳跃的区别：GAN是一下子跳到终点，扩散模型是一步一步走", False, C_DARK, 10),
    ("  • 最终图像质量更高、更稳定、更多样化", False, C_GREEN, 10),
    ("", False, C_DARK, 6),
    ("数学公式：x_t = √ᾱ_t · x₀ + √(1-ᾱ_t) · ε", True, C_PURPLE, 11),
    ("  • x₀ = 原始图像，ε = 随机噪声，ᾱ_t = 时间系数（1→0）", False, C_DARK, 10),
]
add_multiline_textbox(slide, Inches(0.8), Inches(4.85), Inches(11.5), Inches(2.0), lines)

print("✅ 第3页：Diffusion Model 基础原理")

# ==================== 第4页：SD vs 普通Diffusion ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "02  Stable Diffusion vs 普通 Diffusion", "核心创新：在潜在空间中工作，计算量降低48倍")

# 核心对比表
data = [
    ["对比维度", "普通 Diffusion（DDPM）", "Stable Diffusion"],
    ["操作空间", "像素空间（512×512×3）", "潜在空间（64×64×4）"],
    ["数据维度", "786,432 维", "16,384 维"],
    ["压缩比", "无压缩", "48 倍压缩"],
    ["计算量", "极高", "降低约 48 倍"],
    ["GPU 显存", "40GB+（高端 GPU）", "8-12GB（消费级 GPU）"],
    ["生成速度", "几十秒到几分钟", "几秒到十几秒"],
    ["硬件门槛", "高（普通人用不起）", "低（RTX 3060 即可）"],
]
make_table(slide, Inches(0.6), Inches(1.4), Inches(12), 8, 4, data,
           col_widths=[Inches(2.2), Inches(4.9), Inches(4.9)], font_size=11)

# 底部说明
add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.4), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.35),
            "核心区别：在哪里\"画画\"？", font_size=14, font_color=C_PRIMARY, bold=True)

lines = [
    ("普通 Diffusion 就像画家在一块超大画布上逐像素修改，每一笔都要顾及 80 万个像素点，工作量极其庞大。", False, C_DARK, 11),
    ("", False, C_DARK, 6),
    ("Stable Diffusion 先用 VAE 把图像压缩成缩略图（48倍压缩），在缩略图上去噪声（只需处理 1.6 万个数值），", False, C_DARK, 11),
    ("去噪完成后再用 VAE 解码器还原成高清图像。就像画家不是在画布上逐像素画画，而是在缩略图上画画，", False, C_DARK, 11),
    ("画完后再用放大镜还原成高清大图。", False, C_DARK, 11),
    ("", False, C_DARK, 6),
    ("正是因为这个创新，Stable Diffusion 把 AI 图像生成从实验室里的昂贵实验，变成了普通人也能在电脑上运行的工具。", True, C_GREEN, 11),
]
add_multiline_textbox(slide, Inches(0.8), Inches(5.05), Inches(11.5), Inches(1.8), lines)

print("✅ 第4页：SD vs 普通Diffusion")

# ==================== 第5页：SD整体架构概览 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "03  Stable Diffusion 整体架构", "CLIP + U-Net + VAE 三大组件协同工作")

# 架构图 - 三个组件框
# CLIP
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(2.2), fill_color=RGBColor(0xE8, 0xF0, 0xFE), line_color=C_ACCENT)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(3.1), Inches(0.3),
            "1️⃣ CLIP 文本编码器", font_size=14, font_color=C_ACCENT, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.95), Inches(3.1), Inches(0.25),
            "\"翻译官\" — 123M 参数", font_size=10, font_color=C_GRAY, bold=True)
clip_lines = [
    ("输入：文字描述（最多77个token）", False, C_DARK, 10),
    ("输出：77×768 维语义向量", False, C_DARK, 10),
    ("作用：把人类语言翻译成AI数字语言", False, C_DARK, 10),
    ("训练：4亿张图片+文字对", False, C_DARK, 10),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.3), Inches(3.1), Inches(1.3), clip_lines)

# U-Net
add_rounded_rect(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(2.2), fill_color=RGBColor(0xFD, 0xED, 0xEC), line_color=C_RED)
add_textbox(slide, Inches(5.0), Inches(1.6), Inches(3.4), Inches(0.3),
            "2️⃣ U-Net 去噪网络", font_size=14, font_color=C_RED, bold=True)
add_textbox(slide, Inches(5.0), Inches(1.95), Inches(3.4), Inches(0.25),
            "\"核心画师\" — 860M 参数（80%计算量）", font_size=10, font_color=C_GRAY, bold=True)
unet_lines = [
    ("输入：噪声潜空间 + 文本向量 + 时间步", False, C_DARK, 10),
    ("输出：预测噪声（64×64×4）", False, C_DARK, 10),
    ("作用：一步步去掉噪声，\"雕刻\"图像", False, C_DARK, 10),
    ("核心：Cross-Attention + Skip Connections", False, C_DARK, 10),
]
add_multiline_textbox(slide, Inches(5.0), Inches(2.3), Inches(3.4), Inches(1.3), unet_lines)

# VAE
add_rounded_rect(slide, Inches(9.3), Inches(1.5), Inches(3.4), Inches(2.2), fill_color=RGBColor(0xE8, 0xF8, 0xF5), line_color=C_GREEN)
add_textbox(slide, Inches(9.5), Inches(1.6), Inches(3.0), Inches(0.3),
            "3️⃣ VAE 编解码器", font_size=14, font_color=C_GREEN, bold=True)
add_textbox(slide, Inches(9.5), Inches(1.95), Inches(3.0), Inches(0.25),
            "\"压缩工具\" — 83M 参数", font_size=10, font_color=C_GRAY, bold=True)
vae_lines = [
    ("编码器：512×512×3 → 64×64×4（压缩48倍）", False, C_DARK, 10),
    ("解码器：64×64×4 → 512×512×3（还原高清）", False, C_DARK, 10),
    ("作用：训练时压缩，推理时还原", False, C_DARK, 10),
    ("比喻：超级智能的ZIP压缩软件", False, C_DARK, 10),
]
add_multiline_textbox(slide, Inches(9.5), Inches(2.3), Inches(3.0), Inches(1.3), vae_lines)

# 箭头连接
add_arrow_right(slide, Inches(4.15), Inches(2.4), Inches(0.6), Inches(0.3), C_ACCENT)
add_textbox(slide, Inches(4.15), Inches(2.15), Inches(0.6), Inches(0.25),
            "文本\n向量", font_size=8, font_color=C_ACCENT, alignment=PP_ALIGN.CENTER)

# 数据流
add_rounded_rect(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(2.9), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.3),
            "完整数据流", font_size=14, font_color=C_PRIMARY, bold=True)

# 推理流程
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.25),
            "推理时（生成图像）：", font_size=12, font_color=C_PRIMARY, bold=True)

flow_items = [
    ("Step 1", "文字描述 → CLIP → 77×768 文本向量", C_ACCENT),
    ("Step 2", "随机噪声（64×64×4）→ U-Net 执行 20-50 步去噪", C_RED),
    ("Step 3", "去噪后的潜空间（64×64×4）→ VAE 解码器", C_GREEN),
    ("Step 4", "输出：512×512×3 高清图像", C_DARK),
]
for i, (step, desc, color) in enumerate(flow_items):
    y = Inches(4.85) + Inches(i * 0.45)
    add_chevron(slide, Inches(0.9), y, Inches(0.8), Inches(0.3), color)
    add_textbox(slide, Inches(0.95), y + Inches(0.02), Inches(0.7), Inches(0.25),
                step, font_size=9, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.8), y, Inches(4.5), Inches(0.3),
                desc, font_size=11, font_color=C_DARK)

# 参数量分布表
add_textbox(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(0.25),
            "参数量分布（SD 1.5，总计约 10 亿）：", font_size=12, font_color=C_PRIMARY, bold=True)
param_data = [
    ["组件", "参数量", "占比", "显存"],
    ["CLIP", "123M", "~12%", "~0.5 GB"],
    ["VAE", "83M", "~8%", "~0.3 GB"],
    ["U-Net", "860M", "~80%", "~4-5 GB"],
    ["总计", "~1B", "100%", "6-8 GB"],
]
make_table(slide, Inches(7), Inches(4.85), Inches(5.5), 5, 4, param_data,
           col_widths=[Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.4)], font_size=10)

print("✅ 第5页：整体架构概览")

# ==================== 第6页：CLIP详解 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "04  CLIP 文本编码器详解", "Contrastive Language-Image Pre-training — 把文字翻译成AI能理解的数字")

# CLIP 4维度卡片
dims = [
    ("输入", "文本描述字符串\n经分词器处理后最多 77 个 token\n约 50-60 个汉字或 10-15 个英文单词\n超出部分会被截断", C_ACCENT),
    ("输出", "77×768 维张量\n共 59,136 个浮点数\n这就是文本的\"语义指纹\"\n包含文字描述的所有含义信息", C_GREEN),
    ("参数", "约 1.23 亿（123M）\n12 层 Transformer 结构\n隐藏维度 768\n12 个注意力头", C_ORANGE),
    ("作用", "将人类文字翻译成数字向量\n通过 Cross-Attention 注入 U-Net\n指导去噪方向\n没有它，AI不知道该画什么", C_PURPLE),
]

for i, (title, desc, color) in enumerate(dims):
    x = Inches(0.6 + i * 3.15)
    add_rounded_rect(slide, x, Inches(1.4), Inches(2.95), Inches(2.0), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.95), Inches(0.35), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.42), Inches(2.75), Inches(0.3),
                title, font_size=13, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(1.85), Inches(2.65), Inches(1.5),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.2)

# CLIP 内部工作流程
add_rounded_rect(slide, Inches(0.6), Inches(3.65), Inches(12), Inches(3.2), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(3.75), Inches(5), Inches(0.3),
            "CLIP 内部工作流程（4步）", font_size=14, font_color=C_PRIMARY, bold=True)

steps = [
    ("1️⃣ 分词", "文本被拆分成最小语义单元（token）\n\"一只橘猫坐在沙发上\" → \"一只\"/\"橘\"/\"猫\"/\"坐在\"/\"沙发\"/\"上\"\n每个 token 映射成一个 768 维向量，最多处理 77 个", C_ACCENT),
    ("2️⃣ 位置编码", "每个词加上\"位置向量\"\n\"猫坐在沙发上\" ≠ \"沙发坐在猫上\" — 顺序不同含义天差地别\n位置编码让模型理解词语的顺序关系", C_GREEN),
    ("3️⃣ Transformer", "12 层 Transformer 网络处理\n每层包含\"自注意力机制\"，每个词都能\"看到\"其他所有词\n\"橘猫\"会注意到\"坐在\"和\"沙发\"，理解是猫在沙发上的场景", C_ORANGE),
    ("4️⃣ 输出向量", "每个 token 变成 768 维向量，共 77×768 = 59,136 个数值\n不仅包含每个词的独立含义，还包含词与词之间的关系\n是对整段文本的完整\"语义编码\"", C_PURPLE),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 3.05)
    add_chevron(slide, x, Inches(4.15), Inches(2.85), Inches(0.35), color)
    add_textbox(slide, x + Inches(0.2), Inches(4.17), Inches(2.45), Inches(0.3),
                title, font_size=11, font_color=C_WHITE, bold=True)
    add_textbox(slide, x + Inches(0.05), Inches(4.6), Inches(2.75), Inches(2.1),
                desc, font_size=9, font_color=C_DARK, line_spacing=1.15)

print("✅ 第6页：CLIP详解")

# ==================== 第7页：VAE详解 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "05  VAE 编解码器详解", "Variational Autoencoder — 超级智能的ZIP压缩软件，压缩48倍")

# VAE 4维度
dims = [
    ("编码器输入", "原始图像 512×512×3\n共 786,432 个像素值\n每个值范围 0-255", C_ACCENT),
    ("编码器输出", "潜空间 64×64×4\n共 16,384 个浮点数\n值域约 [-5, 5]", C_GREEN),
    ("参数量", "约 8300 万（83M）\n编码器 ~34M，解码器 ~49M\n解码器更大（还原比压缩更复杂）", C_ORANGE),
    ("核心作用", "训练时压缩（降低扩散计算成本）\n推理时还原（潜空间→高清图像）\n保留人眼重要的信息，丢弃冗余", C_PURPLE),
]

for i, (title, desc, color) in enumerate(dims):
    x = Inches(0.6 + i * 3.15)
    add_rounded_rect(slide, x, Inches(1.4), Inches(2.95), Inches(1.6), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.95), Inches(0.32), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.42), Inches(2.75), Inches(0.28),
                title, font_size=12, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(1.8), Inches(2.65), Inches(1.1),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.15)

# 编码器结构
add_rounded_rect(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.7), fill_color=RGBColor(0xE8, 0xF0, 0xFE))
add_textbox(slide, Inches(0.8), Inches(3.3), Inches(5.4), Inches(0.3),
            "编码器：4层下采样（每层宽高缩小一半）", font_size=13, font_color=C_ACCENT, bold=True)

enc_steps = [
    ("512×512×3", "输入图像（RGB）", C_ACCENT),
    ("256×256×128", "第1层卷积 + 下采样", C_ACCENT),
    ("128×128×256", "第2层卷积 + 下采样", C_ACCENT),
    ("64×64×512", "第3层卷积 + 下采样", C_ACCENT),
    ("32×32×512", "第4层卷积 + 下采样", C_ACCENT),
    ("64×64×4", "输出潜空间（48倍压缩！）", C_GREEN),
]

for i, (size, desc, color) in enumerate(enc_steps):
    y = Inches(3.7) + Inches(i * 0.48)
    add_rounded_rect(slide, Inches(0.9), y, Inches(1.4), Inches(0.38), fill_color=color)
    add_textbox(slide, Inches(0.95), y + Inches(0.04), Inches(1.3), Inches(0.3),
                size, font_size=10, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(2.5), y + Inches(0.04), Inches(3.5), Inches(0.3),
                desc, font_size=10, font_color=C_DARK)
    if i < len(enc_steps) - 1:
        add_arrow_down(slide, Inches(1.4), y + Inches(0.38), Inches(0.3), Inches(0.1), color)

# 解码器结构
add_rounded_rect(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.7), fill_color=RGBColor(0xE8, 0xF8, 0xF5))
add_textbox(slide, Inches(7.0), Inches(3.3), Inches(5.4), Inches(0.3),
            "解码器：4层上采样（每层宽高翻倍）", font_size=13, font_color=C_GREEN, bold=True)

dec_steps = [
    ("64×64×4", "输入潜空间", C_GREEN),
    ("128×128×512", "第1层卷积 + 上采样", C_GREEN),
    ("256×256×256", "第2层卷积 + 上采样", C_GREEN),
    ("512×512×128", "第3层卷积 + 上采样", C_GREEN),
    ("512×512×64", "第4层卷积 + 上采样", C_GREEN),
    ("512×512×3", "输出高清图像（RGB）", C_ACCENT),
]

for i, (size, desc, color) in enumerate(dec_steps):
    y = Inches(3.7) + Inches(i * 0.48)
    add_rounded_rect(slide, Inches(7.1), y, Inches(1.4), Inches(0.38), fill_color=color)
    add_textbox(slide, Inches(7.15), y + Inches(0.04), Inches(1.3), Inches(0.3),
                size, font_size=10, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(8.7), y + Inches(0.04), Inches(3.5), Inches(0.3),
                desc, font_size=10, font_color=C_DARK)
    if i < len(dec_steps) - 1:
        add_arrow_down(slide, Inches(7.6), y + Inches(0.38), Inches(0.3), Inches(0.1), color)

print("✅ 第7页：VAE详解")

# ==================== 第8页：U-Net详解 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "06  U-Net 噪声预测网络详解", "核心画师 — 8.6亿参数，占据80%以上计算量")

# U-Net 4维度
dims = [
    ("输入", "① 带噪声潜空间（64×64×4）\n② 文本向量（77×768）\n③ 时间步 t（1-1000）\n④ 可选条件图像", C_ACCENT),
    ("输出", "预测噪声张量（64×64×4）\n形状与输入相同\n用于从带噪声图像中\"减去\"噪声", C_GREEN),
    ("参数", "约 8.6 亿（860M）\n三个组件中最大\n推理时占用 4-5 GB 显存", C_ORANGE),
    ("作用", "核心去噪引擎\n每步预测应去除的噪声\n从噪声逐步\"雕刻\"出清晰图像\n比喻：经验丰富的雕塑家", C_PURPLE),
]

for i, (title, desc, color) in enumerate(dims):
    x = Inches(0.6 + i * 3.15)
    add_rounded_rect(slide, x, Inches(1.4), Inches(2.95), Inches(1.55), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.95), Inches(0.32), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.42), Inches(2.75), Inches(0.28),
                title, font_size=12, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(1.8), Inches(2.65), Inches(1.05),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.15)

# U-Net 结构图
add_rounded_rect(slide, Inches(0.6), Inches(3.15), Inches(12), Inches(3.7), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(3.25), Inches(5), Inches(0.3),
            "U-Net 结构与三个关键机制", font_size=14, font_color=C_PRIMARY, bold=True)

# 编码器路径
add_textbox(slide, Inches(0.8), Inches(3.65), Inches(3), Inches(0.25),
            "编码器路径（逐步收缩）", font_size=11, font_color=C_RED, bold=True)
enc_layers = [
    ("64×64×320 → 32×32×640", "ResBlock + Attention"),
    ("32×32×640 → 16×16×1280", "ResBlock + Attention"),
    ("16×16×1280 → 8×8×1280", "ResBlock + Attention"),
    ("8×8×1280 → 4×4×1280", "ResBlock（最深层）"),
]
for i, (size, desc) in enumerate(enc_layers):
    y = Inches(3.95) + Inches(i * 0.38)
    add_textbox(slide, Inches(1.0), y, Inches(2.5), Inches(0.3),
                size, font_size=9, font_color=C_DARK, font_name=FONT_EN)
    add_textbox(slide, Inches(3.3), y, Inches(1.5), Inches(0.3),
                desc, font_size=8, font_color=C_GRAY)

# 瓶颈层
add_textbox(slide, Inches(5.2), Inches(3.65), Inches(2.5), Inches(0.25),
            "瓶颈层（最深层处理）", font_size=11, font_color=C_PURPLE, bold=True)
add_textbox(slide, Inches(5.2), Inches(3.95), Inches(2.5), Inches(0.8),
            "4×4×1280\n最抽象的语义信息\n空间最小，通道最大", font_size=10, font_color=C_DARK, line_spacing=1.2)

# 解码器路径
add_textbox(slide, Inches(8.2), Inches(3.65), Inches(4), Inches(0.25),
            "解码器路径（逐步扩张）", font_size=11, font_color=C_GREEN, bold=True)
dec_layers = [
    ("4×4 → 8×8×1280", "+ Skip Connection"),
    ("8×8 → 16×16×640", "+ Skip Connection"),
    ("16×16 → 32×32×320", "+ Skip Connection"),
    ("32×32 → 64×64×4", "最终输出（预测噪声）"),
]
for i, (size, desc) in enumerate(dec_layers):
    y = Inches(3.95) + Inches(i * 0.38)
    add_textbox(slide, Inches(8.4), y, Inches(2.5), Inches(0.3),
                size, font_size=9, font_color=C_DARK, font_name=FONT_EN)
    add_textbox(slide, Inches(10.7), y, Inches(1.8), Inches(0.3),
                desc, font_size=8, font_color=C_GREEN, bold=True)

# 三个关键机制
mechs = [
    ("Cross-Attention（交叉注意力）", "文本条件注入核心机制\n图像特征\"问\"：根据文字我该怎样？\n文本特征\"答\"：你这里应该是橘猫毛发\n出现在编码器1-3层 + 解码器1-3层（共6处）", C_ACCENT),
    ("Skip Connections（跳跃连接）", "编码器细节直接传递给解码器\n防止信息在压缩过程中丢失\n就像建筑师随时参考\"草稿本\"\n没有它，细节会大打折扣", C_ORANGE),
    ("Time Embedding（时间嵌入）", "告诉 U-Net 当前噪声水平（1-1000）\n高噪声时\"粗雕\"（确定轮廓）\n低噪声时\"精雕\"（添加细节）\n通过正弦编码 + 两层MLP注入每个ResBlock", C_PURPLE),
]

for i, (title, desc, color) in enumerate(mechs):
    x = Inches(0.8 + i * 4.1)
    add_rounded_rect(slide, x, Inches(5.55), Inches(3.8), Inches(1.2), fill_color=C_WHITE, line_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(5.6), Inches(3.6), Inches(0.25),
                title, font_size=10, font_color=color, bold=True)
    add_textbox(slide, x + Inches(0.1), Inches(5.85), Inches(3.6), Inches(0.85),
                desc, font_size=9, font_color=C_DARK, line_spacing=1.1)

print("✅ 第8页：U-Net详解")

# ==================== 第9页：三大组件协作 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "07  三大组件协作流程", "AI 画图流水线 — 各司其职，紧密配合")

# 三站流水线
stations = [
    ("第1站", "CLIP 翻译站", "用户输入：\"一只橘猫坐在木质沙发上，阳光从窗户照进来\"\n→ 通过 12 层 Transformer 处理\n→ 输出 77×768 维\"语义指纹\"向量\n→ 发送给下一站的 U-Net", C_ACCENT, "\"翻译官\" — 把人类语言翻译成AI数字语言"),
    ("第2站", "U-Net 画图站", "接收文本向量 + 64×64×4 纯随机噪声\n第 1 轮：只能\"猜\"出大致色块分布\n第 15 轮：能看到猫的轮廓、沙发形状\n第 30 轮：毛发纹理、阳光光线细节\n第 50 轮：去噪完成，输出干净潜空间", C_RED, "\"核心画师\" — 执行 20-50 轮去噪循环"),
    ("第3站", "VAE 还原站", "接收 U-Net 输出的干净潜空间\n通过 4 层上采样逐步还原\n64×64 → 128×128 → 256×256 → 512×512\n最终输出 512×512×3 高清图像", C_GREEN, "\"还原打印机\" — 潜空间→高清图像"),
]

for i, (station, title, desc, color, role) in enumerate(stations):
    x = Inches(0.6 + i * 4.15)
    add_rounded_rect(slide, x, Inches(1.4), Inches(3.9), Inches(3.8), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, Inches(1.4), Inches(3.9), Inches(0.55), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.42), Inches(3.7), Inches(0.25),
                station, font_size=11, font_color=RGBColor(0xDD, 0xEE, 0xFF), bold=True)
    add_textbox(slide, x + Inches(0.1), Inches(1.68), Inches(3.7), Inches(0.25),
                title, font_size=14, font_color=C_WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(2.1), Inches(3.6), Inches(2.2),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.2)
    # 底部角色说明
    add_rounded_rect(slide, x + Inches(0.1), Inches(4.65), Inches(3.7), Inches(0.4), fill_color=C_LIGHT)
    add_textbox(slide, x + Inches(0.2), Inches(4.68), Inches(3.5), Inches(0.35),
                role, font_size=10, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 箭头连接
    if i < 2:
        add_arrow_right(slide, x + Inches(3.95), Inches(2.8), Inches(0.2), Inches(0.25), color)

# 底部：为什么这样分工
add_rounded_rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.4), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.25),
            "为什么要这样分工？", font_size=13, font_color=C_PRIMARY, bold=True)
lines = [
    ("• CLIP 擅长理解语言，但不擅长生成图像 — 所以只让它做\"翻译\"", False, C_DARK, 10),
    ("• U-Net 擅长在低维空间做逐步变换，但不擅长理解人类语言 — 需要CLIP的翻译结果作为指导", False, C_DARK, 10),
    ("• VAE 擅长压缩和还原，但不擅长理解语义 — 只在开头和结尾使用", False, C_DARK, 10),
    ("就像专业厨房：有人负责理解客人的点单（CLIP），有人负责烹饪（U-Net），有人负责摆盘（VAE）", True, C_ACCENT, 10),
]
add_multiline_textbox(slide, Inches(0.8), Inches(5.85), Inches(11.5), Inches(1.0), lines)

print("✅ 第9页：三大组件协作")

# ==================== 第10页：训练过程 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "08  训练过程", "6亿图像-文本对的炼丹之旅 — 15万GPU小时")

# 6个关键步骤
train_steps = [
    ("1", "采样图像和文本对", "从 LAION 数据集随机抽取图像 x₀\n和对应文本描述", C_ACCENT),
    ("2", "VAE 编码", "512×512×3 → 64×64×4\n压缩到潜空间", C_GREEN),
    ("3", "添加随机噪声", "随机选 t ∈ [1,1000]\nz_t = √ᾱ_t·z₀ + √(1-ᾱ_t)·ε", C_ORANGE),
    ("4", "U-Net 预测噪声", "输入：z_t + t + 文本向量\n输出：预测噪声 ε_pred", C_RED),
    ("5", "计算损失函数", "L = ||ε - ε_pred||²\n预测噪声 vs 真实噪声的均方误差", C_PURPLE),
    ("6", "反向传播更新", "更新 U-Net 的 8.6 亿参数\nCLIP 和 VAE 保持冻结", C_TEAL),
]

for i, (num, title, desc, color) in enumerate(train_steps):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.4 + row * 1.8)
    
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(1.55), fill_color=C_WHITE, line_color=color)
    # 序号圆
    add_rounded_rect(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.4), fill_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.12), Inches(0.4), Inches(0.35),
                num, font_size=16, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_textbox(slide, x + Inches(0.6), y + Inches(0.12), Inches(3.1), Inches(0.3),
                title, font_size=13, font_color=color, bold=True)
    # 描述
    add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.6), Inches(0.9),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.2)

# 箭头连接（步骤间）
for i in range(5):
    row = i // 3
    col = i % 3
    if col < 2:
        x = Inches(0.6 + col * 4.15 + 3.9)
        y = Inches(1.4 + row * 1.8 + 0.6)
        add_arrow_right(slide, x, y, Inches(0.25), Inches(0.2), C_GRAY)

# 训练成本表
add_rounded_rect(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.7), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.25),
            "训练成本数据", font_size=13, font_color=C_PRIMARY, bold=True)
cost_data = [
    ["项目", "数据", "项目", "数据"],
    ["训练数据", "LAION-5B 子集，约 6 亿图文对", "训练硬件", "多张 A100 80GB GPU"],
    ["训练时长", "约 15 万 GPU 小时", "训练成本", "约 50-60 万美元"],
    ["模型大小", "约 4 GB（float16）", "冻结组件", "CLIP + VAE（只训练U-Net）"],
]
make_table(slide, Inches(0.8), Inches(5.6), Inches(11.5), 4, 4, cost_data,
           col_widths=[Inches(2), Inches(3.8), Inches(2), Inches(3.7)], font_size=10)

print("✅ 第10页：训练过程")

# ==================== 第11页：CFG详解 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "09  CFG（Classifier-Free Guidance）详解", "无分类器引导 — 让AI更听话地按你的描述画画")

# CFG 原理
add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(1.6), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.25),
            "核心思想：既知道\"什么是对的\"，也知道\"什么不是对的\"", font_size=13, font_color=C_PRIMARY, bold=True)

lines = [
    ("公式：ε_guided = ε_uncond + s × (ε_cond - ε_uncond)", True, C_PURPLE, 14),
    ("", False, C_DARK, 6),
    ("ε_uncond = 没有文本条件时的预测（\"自由发挥\"）    ε_cond = 有文本条件时的预测（\"按指令画\"）", False, C_DARK, 11),
    ("s = 引导系数（CFG Scale），通常 7-12    s 越大 → 越忠实于文字描述    s 过大 → 过度饱和、出现伪影", False, C_DARK, 11),
    ("比喻：严厉又公正的艺术指导，每画一笔都问\"你按客户要求画了吗？\"和\"不管要求你会怎么画？\"，放大差异强制画师更严格", True, C_ACCENT, 10),
]
add_multiline_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.1), lines)

# CFG Scale 参数影响表
add_textbox(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.3),
            "CFG Scale 参数影响", font_size=13, font_color=C_PRIMARY, bold=True)
cfg_data = [
    ["CFG Scale", "效果", "适用场景"],
    ["1-3", "几乎无引导，图像自由发挥", "创意探索，想要意外惊喜"],
    ["5-7", "适中引导，平衡忠实度和创造力", "通用场景（推荐默认值 7）"],
    ["7-12", "强引导，高度忠实文字描述", "需要精确控制，如产品图"],
    ["15-20", "过强引导，可能出现伪影", "一般不推荐"],
    ["20+", "严重过度饱和，质量急剧下降", "绝对不推荐"],
]
make_table(slide, Inches(0.6), Inches(3.55), Inches(6.5), 6, 3, cfg_data,
           col_widths=[Inches(1.5), Inches(2.8), Inches(2.2)], font_size=10)

# 图示：CFG工作原理
add_rounded_rect(slide, Inches(7.5), Inches(3.2), Inches(5.2), Inches(3.5), fill_color=C_WHITE, line_color=C_ACCENT)
add_textbox(slide, Inches(7.7), Inches(3.3), Inches(4.8), Inches(0.3),
            "CFG 工作原理图示", font_size=13, font_color=C_ACCENT, bold=True)

# 无条件
add_rounded_rect(slide, Inches(7.8), Inches(3.8), Inches(2), Inches(0.8), fill_color=RGBColor(0xE0, 0xE0, 0xE0))
add_textbox(slide, Inches(7.8), Inches(3.95), Inches(2), Inches(0.5),
            "无条件生成\n（自由发挥）", font_size=10, font_color=C_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# 有条件
add_rounded_rect(slide, Inches(10.3), Inches(3.8), Inches(2), Inches(0.8), fill_color=RGBColor(0xE8, 0xF0, 0xFE), line_color=C_ACCENT)
add_textbox(slide, Inches(10.3), Inches(3.95), Inches(2), Inches(0.5),
            "有条件生成\n（按指令画）", font_size=10, font_color=C_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# 差异箭头
add_textbox(slide, Inches(9.85), Inches(3.9), Inches(0.5), Inches(0.3),
            "→", font_size=20, font_color=C_PURPLE, bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
add_textbox(slide, Inches(9.8), Inches(4.15), Inches(0.6), Inches(0.25),
            "差异", font_size=8, font_color=C_PURPLE, alignment=PP_ALIGN.CENTER)

# 放大
add_arrow_down(slide, Inches(11.2), Inches(4.65), Inches(0.25), Inches(0.3), C_PURPLE)
add_textbox(slide, Inches(10.5), Inches(4.65), Inches(1.5), Inches(0.25),
            "× CFG Scale (s)", font_size=9, font_color=C_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# 最终结果
add_rounded_rect(slide, Inches(8.8), Inches(5.2), Inches(3.2), Inches(0.7), fill_color=C_GREEN)
add_textbox(slide, Inches(8.8), Inches(5.3), Inches(3.2), Inches(0.5),
            "= 引导后的高质量结果", font_size=12, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 加号和说明
add_textbox(slide, Inches(7.8), Inches(5.0), Inches(2), Inches(0.4),
            "+ 无条件结果", font_size=9, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

# 总结
add_rounded_rect(slide, Inches(0.6), Inches(6.3), Inches(6.5), Inches(0.6), fill_color=RGBColor(0xFE, 0xF9, 0xE7))
add_textbox(slide, Inches(0.8), Inches(6.35), Inches(6.1), Inches(0.5),
            "总结：没有 CFG，生成图像模糊、偏离描述；CFG 让图像既高质量又忠实于文字描述。", font_size=11, font_color=C_ORANGE, bold=True)

print("✅ 第11页：CFG详解")

# ==================== 第12页：推理过程 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "10  推理过程（去噪循环）", "从噪声到图像的雕刻过程 — 4个步骤")

# 4个步骤
infer_steps = [
    ("Step 1", "文本编码", "用户输入文字描述\n→ CLIP 文本编码器\n→ 77×768 维文本向量\n只需执行一次，后续复用", C_ACCENT),
    ("Step 2", "初始化噪声", "从标准正态分布采样\n64×64×4 随机张量\n= 纯噪声\"雪花屏\"\n种子(seed)决定随机数", C_GREEN),
    ("Step 3", "逐步去噪", "循环执行 20-50 步\n每步：U-Net预测噪声 → 减去噪声\n从 t=1000 递减到 t=1\n逐步\"雕刻\"出清晰图像", C_ORANGE),
    ("Step 4", "VAE 解码", "64×64×4 潜空间\n→ VAE 解码器（4层上采样）\n→ 512×512×3 RGB 图像\n输出给用户", C_PURPLE),
]

for i, (step, title, desc, color) in enumerate(infer_steps):
    x = Inches(0.6 + i * 3.15)
    add_rounded_rect(slide, x, Inches(1.4), Inches(2.95), Inches(2.6), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, Inches(1.4), Inches(2.95), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(0.1), Inches(1.42), Inches(2.75), Inches(0.2),
                step, font_size=10, font_color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(1.62), Inches(2.75), Inches(0.25),
                title, font_size=14, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(2.05), Inches(2.65), Inches(1.8),
                desc, font_size=10, font_color=C_DARK, line_spacing=1.25)
    if i < 3:
        add_arrow_right(slide, x + Inches(2.95), Inches(2.4), Inches(0.2), Inches(0.2), color)

# 去噪循环详解
add_rounded_rect(slide, Inches(0.6), Inches(4.2), Inches(5.5), Inches(2.7), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.25),
            "去噪循环详解（Step 3 的核心）", font_size=13, font_color=C_PRIMARY, bold=True)

denoise_lines = [
    ("对于 t 从 T 递减到 1（如 1000→950→900→...→1）：", True, C_DARK, 10),
    ("", False, C_DARK, 4),
    ("  ① U-Net 预测当前噪声：", False, C_DARK, 10),
    ("     ε_pred = U-Net(x_t, t, text_embedding)", False, C_PURPLE, 10),
    ("", False, C_DARK, 4),
    ("  ② 采样器计算去噪结果：", False, C_DARK, 10),
    ("     x_{t-1} = sampler_step(x_t, ε_pred, t)", False, C_PURPLE, 10),
    ("", False, C_DARK, 4),
    ("比喻：从一块毛石中雕刻出雕像", True, C_ACCENT, 10),
    ("  第1-5步：粗雕（确定轮廓和构图）", False, C_DARK, 10),
    ("  第5-20步：中雕（显现物体形状）", False, C_DARK, 10),
    ("  第20-50步：精雕（毛发纹理、光影细节）", False, C_DARK, 10),
]
add_multiline_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.1), Inches(2.2), denoise_lines)

# 推理时间参考表
add_textbox(slide, Inches(6.5), Inches(4.2), Inches(6), Inches(0.25),
            "推理时间参考", font_size=13, font_color=C_PRIMARY, bold=True)
time_data = [
    ["配置", "步数", "采样器", "GPU", "单张耗时"],
    ["快速预览", "10-15", "DPM++ SDE", "RTX 3060", "2-4 秒"],
    ["标准质量", "20-30", "DPM++ 2M", "RTX 3060", "5-8 秒"],
    ["高质量", "40-50", "DPM++ 2M Karras", "RTX 4090", "3-5 秒"],
    ["最高质量", "50-100", "DDIM", "RTX 4090", "5-10 秒"],
]
make_table(slide, Inches(6.5), Inches(4.5), Inches(6.2), 5, 5, time_data,
           col_widths=[Inches(1.2), Inches(0.8), Inches(1.8), Inches(1.2), Inches(1.2)], font_size=10)

print("✅ 第12页：推理过程")

# ==================== 第13页：采样器对比 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "11  采样器对比", "9种算法 — 从噪声到图像的不同路径")

# 采样器对比表
sampler_data = [
    ["采样器", "步数", "速度", "质量", "稳定性", "适用场景"],
    ["Euler", "20-30", "快", "中", "高", "快速预览、初学者"],
    ["Euler a", "20-30", "快", "中高", "高", "通用，带随机性"],
    ["DDIM", "20-50", "中", "高", "高", "确定性生成"],
    ["DPM++ 2M", "20-30", "快", "高", "高", "⭐ 最推荐通用"],
    ["DPM++ 2M Karras", "20-30", "快", "很高", "高", "⭐ 高质量首选"],
    ["DPM++ SDE Karras", "20-30", "中", "很高", "中", "艺术感强的图像"],
    ["DPM++ SDE", "10-15", "很快", "中高", "中", "快速出图"],
    ["UniPC", "10-20", "很快", "高", "高", "最新高效采样器"],
    ["LCM", "4-8", "极快", "中", "高", "实时生成"],
]
make_table(slide, Inches(0.6), Inches(1.4), Inches(12), 10, 6, sampler_data,
           col_widths=[Inches(2.5), Inches(1), Inches(1), Inches(1), Inches(1), Inches(5.5)], font_size=10)

# 推荐选择
add_rounded_rect(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.7), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.25),
            "推荐选择", font_size=13, font_color=C_PRIMARY, bold=True)

recs = [
    ("⭐ 通用推荐", "DPM++ 2M Karras，20-30 步", "公认最佳性价比，接近最优质量", C_GREEN),
    ("⚡ 追求速度", "DPM++ SDE（10-15步）或 LCM（4-8步）", "适合快速迭代、概念验证", C_ORANGE),
    ("🏆 追求质量", "DPM++ 2M Karras，40-50 步", "适合最终交付、打印输出", C_ACCENT),
]
for i, (tag, tool, desc, color) in enumerate(recs):
    y = Inches(5.65) + Inches(i * 0.4)
    add_rounded_rect(slide, Inches(0.8), y, Inches(1.5), Inches(0.32), fill_color=color)
    add_textbox(slide, Inches(0.85), y + Inches(0.02), Inches(1.4), Inches(0.28),
                tag, font_size=10, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), y + Inches(0.02), Inches(3.5), Inches(0.28),
                tool, font_size=10, font_color=C_DARK, bold=True)
    add_textbox(slide, Inches(6.2), y + Inches(0.02), Inches(6), Inches(0.28),
                desc, font_size=10, font_color=C_GRAY)

print("✅ 第13页：采样器对比")

# ==================== 第14页：参数优化 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "12  参数优化", "5个核心参数的最佳实践")

# 5个参数卡片
params = [
    ("Steps（去噪步数）", "U-Net执行去噪的次数", "推荐：20-50步\n经验值：DPM++ 2M 用25步\n超过30步后提升越来越小\n比喻：打磨家具，20遍已很好", C_ACCENT),
    ("CFG Scale（引导系数）", "控制图像对文本的忠实度", "推荐：5-12\n通用：7-8，精确控制：10-12\n太低→偏离描述，太高→过度饱和\n比喻：老师批改作业的严格程度", C_GREEN),
    ("Seed（随机种子）", "决定初始随机噪声的值", "范围：任意整数\n相同种子+相同参数=相同图像\n不同种子=完全不同的图像\n比喻：掷骰子，记住数字可复现", C_ORANGE),
    ("Sampler（采样器）", "去噪的数学算法", "推荐：DPM++ 2M Karras\n快速：DPM++ SDE（15步）\n极速：LCM（4-8步）\n不同算法在速度和质量间取舍", C_PURPLE),
    ("图像尺寸", "输出图像的分辨率", "SD1.5：512×512\nSDXL：1024×1024\n尺寸应为64的倍数\n更大的尺寸→更多细节但更慢", C_TEAL),
]

for i, (title, desc, detail, color) in enumerate(params):
    if i < 3:
        x = Inches(0.6 + i * 4.15)
        y = Inches(1.4)
    else:
        x = Inches(0.6 + (i-3) * 6.2)
        y = Inches(3.8)
    
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.1), fill_color=C_WHITE, line_color=color)
    add_rect(slide, x, y, Inches(3.9), Inches(0.35), fill_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.03), Inches(3.7), Inches(0.3),
                title, font_size=12, font_color=C_WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.6), Inches(0.25),
                desc, font_size=10, font_color=C_GRAY, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.7), Inches(3.6), Inches(1.3),
                detail, font_size=10, font_color=C_DARK, line_spacing=1.2)

# 最佳实践表
add_textbox(slide, Inches(0.6), Inches(6.1), Inches(5), Inches(0.25),
            "最佳实践参数组合", font_size=12, font_color=C_PRIMARY, bold=True)
best_data = [
    ["场景", "Steps", "CFG", "Sampler"],
    ["快速预览", "15", "7", "DPM++ SDE"],
    ["标准生成", "25", "7.5", "DPM++ 2M Karras"],
    ["高质量", "40", "8", "DPM++ 2M Karras"],
    ["精细肖像", "30", "7", "DPM++ 2M Karras"],
]
make_table(slide, Inches(0.6), Inches(6.35), Inches(12), 5, 4, best_data,
           col_widths=[Inches(2.5), Inches(2), Inches(2), Inches(5.5)], font_size=10)

print("✅ 第14页：参数优化")

# ==================== 第15页：应用场景 ====================

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "13  应用场景", "5大领域的革命性影响")

scenes = [
    ("🎨", "概念艺术与插画设计", "几秒生成数十种设计方案\n游戏/电影/广告公司大量使用\n设计师先用SD快速生成灵感图\n比喻：给设计师配备\"灵感无限的外脑\"", C_ACCENT),
    ("🎮", "游戏与影视资产制作", "角色设计、场景概念图、纹理素材\n结合ControlNet控制姿态+LoRA统一风格\n批量生成风格一致的资产\n许多独立游戏团队的核心美术工具", C_GREEN),
    ("🛒", "电商与广告创意", "产品图合成到各种场景中\n一双鞋→海滩/雪山/城市街道\n成本几乎为零\n广告公司快速生成创意初稿", C_ORANGE),
    ("🏗️", "建筑与室内设计", "快速生成建筑外观和室内渲染图\n\"现代极简客厅，落地窗，暖色木质\"\n几秒得到多张设计效果图\n用于早期方案沟通和客户展示", C_PURPLE),
    ("✨", "个人创意与社交媒体", "独一无二的头像、小说插图\n社交媒体内容创作\n完全免费开源，可本地运行\n不受任何使用限制", C_TEAL),
]

for i, (emoji, title, desc, color) in enumerate(scenes):
    x = Inches(0.6 + i * 2.5)
    add_rounded_rect(slide, x, Inches(1.4), Inches(2.3), Inches(4.5), fill_color=C_WHITE, line_color=color)
    
    # emoji
    add_textbox(slide, x, Inches(1.5), Inches(2.3), Inches(0.5),
                emoji, font_size=28, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_rect(slide, x, Inches(2.0), Inches(2.3), Inches(0.35), fill_color=color)
    add_textbox(slide, x + Inches(0.05), Inches(2.02), Inches(2.2), Inches(0.3),
                title, font_size=11, font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_textbox(slide, x + Inches(0.1), Inches(2.5), Inches(2.1), Inches(3.2),
                desc, font_size=9, font_color=C_DARK, line_spacing=1.25)

# 底部总结
add_rounded_rect(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.8), fill_color=C_LIGHT)
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
            "Stable Diffusion 正在彻底改变创意产业的工作流程 — 从概念设计到最终交付，AI 图像生成已经成为设计师、开发者、创作者不可或缺的工具。",
            font_size=12, font_color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

print("✅ 第15页：应用场景")

# ==================== 保存 ====================

prs.save(OUTPUT_PATH)
print(f"\n🎉 PPT 生成完成！")
print(f"   文件路径：{OUTPUT_PATH}")
print(f"   总页数：{len(prs.slides)} 页")
