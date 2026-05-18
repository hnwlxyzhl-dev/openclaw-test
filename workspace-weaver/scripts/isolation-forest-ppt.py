#!/usr/bin/env python3
"""孤立森林 PPT 生成脚本 - 深色科技风"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── 配色方案 ──
BG_DARK      = RGBColor(0x0B, 0x1A, 0x2E)  # 深蓝底
BG_CARD      = RGBColor(0x12, 0x25, 0x3F)  # 卡片背景
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)  # 荧光绿
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)  # 荧光青
ACCENT_ORANGE= RGBColor(0xFF, 0x6B, 0x35)  # 亮橙
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xC4, 0xDE)
MID_GRAY     = RGBColor(0x7A, 0x8E, 0xA8)
DIM_GRAY     = RGBColor(0x4A, 0x5E, 0x78)

FONT_CN = '微软雅黑'
FONT_EN = 'Calibri'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    p.space_before = Pt(0)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=14,
                           default_color=LIGHT_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, size, color, bold) or just text string"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)

    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            color = line[2] if len(line) > 2 else default_color
            bold = line[3] if len(line) > 3 else default_bold

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_CN
        p.alignment = alignment
        p.space_after = Pt(3)
        p.space_before = Pt(1)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_GREEN, thickness=3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ════════════════════════════════════════
# 第 1 页：封面
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 装饰线
add_accent_line(slide, 1.0, 2.8, 4.0, ACCENT_GREEN)

add_textbox(slide, 1.0, 1.2, 11, 1.5, '孤立森林', 48, WHITE, True, PP_ALIGN.LEFT)
add_textbox(slide, 1.0, 2.0, 11, 0.8, 'Isolation Forest', 28, ACCENT_CYAN, False, PP_ALIGN.LEFT)

add_multiline_textbox(slide, 1.0, 3.2, 11, 1.5, [
    ('不建模正常，直接孤立异常——一种颠覆性的异常检测算法', 20, LIGHT_GRAY, False),
    ('', 8),
    ('适用于大数据、高维度、无标签的场景', 14, MID_GRAY, False),
    ('线性时间复杂度，天然支持并行化', 14, MID_GRAY, False),
    ('2008 年 IEEE ICDM 论文  |  Liu, Ting & Zhou', 12, DIM_GRAY, False),
])

# 右侧装饰 - 散点图示意
add_rect(slide, 8.5, 1.5, 4.0, 4.5, RGBColor(0x0E, 0x1F, 0x35), ACCENT_GREEN)
# 模拟密集点
import random
random.seed(42)
for _ in range(20):
    x = 8.8 + random.random() * 1.5
    y = 2.0 + random.random() * 2.5
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_CYAN
    dot.line.fill.background()
# 模拟异常点
for ax, ay in [(11.2, 3.0), (9.5, 5.5)]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax), Inches(ay), Inches(0.14), Inches(0.14))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_ORANGE
    dot.line.fill.background()

add_textbox(slide, 1.0, 6.5, 11, 0.5, '异常检测系列  |  深度技术解析', 12, DIM_GRAY, False)


# ════════════════════════════════════════
# 第 2 页：什么是异常检测？
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '什么是异常检测？', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '为什么要关注"与众不同"的数据？', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 4.8, [
    ('什么是异常检测？', 18, ACCENT_GREEN, True),
    ('异常检测是从数据中识别出与绝大多数数据显著不同的样本的过程。就像在一群白绵羊中找出那只黑绵羊——正常数据遵循某种模式，而异常点偏离了这个模式。', 13, LIGHT_GRAY),
    ('', 6),
    ('为什么它如此重要？', 18, ACCENT_GREEN, True),
    ('异常事件往往伴随高风险。网络入侵导致数据泄露，金融欺诈造成巨额损失，设备异常引发生产线停工。正因为"异常"通常等于"危险"，异常检测成为数据科学最重要的研究方向之一。', 13, LIGHT_GRAY),
    ('', 6),
    ('异常的三种类型：', 16, WHITE, True),
    ('1️⃣  点异常：单个数据点本身就很异常。一群身高 1.6-1.8m 的人中间站了一个 2.3m 的人。', 12, LIGHT_GRAY),
    ('2️⃣  上下文异常：特定上下文中才算异常。夏天穿短袖正常，零下 20°C 穿短袖就异常。', 12, LIGHT_GRAY),
    ('3️⃣  集合异常：单点都正常，放在一起才异常。每天凌晨访问数据库正常，连续三个月天天凌晨就可疑。', 12, LIGHT_GRAY),
])

# 右侧 - 安检类比卡片
add_rect(slide, 7.2, 2.0, 5.5, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.5, 2.2, 5.0, 4.4, [
    ('🛡️  机场安检类比', 16, ACCENT_GREEN, True),
    ('', 6),
    ('异常检测就像机场安检——99.9% 的旅客都正常，但系统必须在几秒内识别出那 0.1% 的危险人物。', 13, LIGHT_GRAY),
    ('', 6),
    ('三种异常 = 三种"嫌疑人"：', 13, WHITE, True),
    ('• 带着违禁品的人 → 点异常（一看就不对）', 12, ACCENT_CYAN),
    ('• 错误时间出现在错误地点 → 上下文异常', 12, ACCENT_CYAN),
    ('• 每人行为正常但组合可疑 → 集合异常（团伙作案）', 12, ACCENT_CYAN),
])


# ════════════════════════════════════════
# 第 3 页：传统方法
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '传统异常检测方法', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '先建模"正常"，再找"偏离"——但建模"正常"恰恰是最难的步骤', 16, ACCENT_CYAN)

# 方法对比表
table_data = [
    ('方法', '核心思路', '时间复杂度', '高维适应', '核心局限'),
    ('Z-score', '距离均值几个标准差', 'O(n)', '差', '假设正态分布'),
    ('KNN / LOF', '邻居距离/局部密度', 'O(n²)', '一般', '计算量灾难性'),
    ('DBSCAN', '基于密度聚类', 'O(n log n)', '差', '高维密度难定义'),
    ('One-Class SVM', '学习正常边界', 'O(n²~n³)', '好', '参数敏感、不可扩展'),
    ('孤立森林', '直接孤立异常', 'O(n·t·logψ)', '好', '局部异常检测弱'),
]

rows = len(table_data)
cols = 5
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.2)).table

col_widths = [1.8, 2.8, 2.0, 1.5, 3.6]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = Inches(w)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_CN
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r == 5:
                para.font.bold = True
                para.font.color.rgb = ACCENT_GREEN
            else:
                para.font.color.rgb = LIGHT_GRAY
            para.alignment = PP_ALIGN.LEFT

        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r == 5:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

# 核心瓶颈
add_rect(slide, 0.8, 5.6, 11.7, 1.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.1, 5.7, 11.2, 1.3, [
    ('⚠️  核心瓶颈：不在于"找异常"，而在于"建模正常"太昂贵、太困难、太脆弱', 15, ACCENT_ORANGE, True),
    ('传统方法像一个侦探，先要彻底研究"正常市民的生活规律"，然后才能发现谁的行踪可疑。但刻画"正常"本身就极其困难——正常人有千万种生活方式，你怎么定义"正常"？', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 4 页：核心思想
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '孤立森林的核心思想', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '不建模正常，直接孤立异常——范式转换', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 5.0, [
    ('范式转换', 18, ACCENT_GREEN, True),
    ('从"建模正常"到"孤立异常"', 14, WHITE, True),
    ('', 4),
    ('异常点有两个天然属性："少"和"不同"。因为数量稀少，周围没有同类；因为特征不同，远离大多数。', 13, LIGHT_GRAY),
    ('', 4),
    ('这导致一个直接结果：异常点比正常点更容易被随机分割孤立出来。', 14, ACCENT_GREEN, True),
    ('', 4),
    ('用一个二维例子感受：10 个正常点聚集在 (3-5, 4-5.5)，2 个异常点 A1(20,3) 和 A2(2,25) 远远游离在外。', 13, LIGHT_GRAY),
    ('', 4),
    ('随机选 x 轴，在 x=8.0 画一条线：', 13, LIGHT_GRAY),
    ('→ 所有正常点落在左边，A1 独自落在右边', 13, ACCENT_CYAN),
    ('→ 只需 1 次分裂，A1 就被孤立了！', 14, ACCENT_GREEN, True),
    ('', 4),
    ('正常点呢？周围全是同伴，需要 6-8 次甚至更多分裂才能单独挑出来。', 13, LIGHT_GRAY),
    ('', 4),
    ('路径越短 = 越异常 | 路径越长 = 越正常', 16, WHITE, True),
])

# 右侧 - 散点图示意
add_rect(slide, 7.2, 2.0, 5.5, 4.8, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))
random.seed(10)
for _ in range(25):
    x = 7.8 + random.random() * 2.0
    y = 2.5 + random.random() * 2.5
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_CYAN
    dot.line.fill.background()

# A1
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(3.5), Inches(0.18), Inches(0.18))
dot.fill.solid()
dot.fill.fore_color.rgb = ACCENT_ORANGE
dot.line.fill.background()
add_textbox(slide, 11.2, 3.2, 1.2, 0.4, 'A1', 12, ACCENT_ORANGE, True)

# 分割线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(2.2), Inches(0.02), Inches(4.2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_GREEN
line.line.fill.background()
add_textbox(slide, 10.1, 2.2, 0.6, 0.3, 'x=8', 10, ACCENT_GREEN, True)

# 左侧标注
add_textbox(slide, 8.0, 5.5, 2.0, 0.4, '正常点群体', 11, ACCENT_CYAN)
add_textbox(slide, 10.3, 5.5, 2.2, 0.4, 'A1 被孤立！', 11, ACCENT_GREEN, True)

# 苹果类比
add_rect(slide, 7.2, 6.0, 5.5, 0.8, BG_CARD)
add_textbox(slide, 7.4, 6.05, 5.2, 0.7, '🍎  传统方法先给好苹果拍照建模，再逐一比对。孤立森林闭着眼随机拿，腐烂的因为少且不同，几步就挑出来了。', 10, MID_GRAY)


# ════════════════════════════════════════
# 第 5 页：iTree
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, 'iTree（孤立树）', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '一棵不断随机分割空间的二叉树', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('两种节点', 16, ACCENT_GREEN, True),
    ('内部节点：存分裂规则——"在哪个特征上切？切在哪里？"', 13, LIGHT_GRAY),
    ('叶子节点：存"到达这里的样本数量"——房间还剩几人', 13, LIGHT_GRAY),
    ('', 6),
    ('每次分裂两步操作', 16, ACCENT_GREEN, True),
    ('① 随机选一个特征（等概率）', 13, LIGHT_GRAY),
    ('② 在该特征的 [min, max] 内随机选分裂点', 13, LIGHT_GRAY),
    ('→ 特征值 < 分裂点进左子树，≥ 分裂点进右子树', 13, ACCENT_CYAN),
    ('', 6),
    ('三个终止条件', 16, ACCENT_GREEN, True),
    ('① 样本数 = 1 → 无需再分', 13, LIGHT_GRAY),
    ('② 所有特征值相同 → 无法再分', 13, LIGHT_GRAY),
    ('③ 达到最大深度（通常 log₂256 = 8 层）→ 强制停止', 13, LIGHT_GRAY),
    ('', 6),
    ('iTree vs 普通决策树', 16, ACCENT_GREEN, True),
    ('普通决策树：精心选择最优分裂（像经验丰富的分拣员）', 12, LIGHT_GRAY),
    ('iTree：随机选特征和分裂点（像蒙着眼分拣）', 12, LIGHT_GRAY),
    ('看似"愚蠢"，但随机分割对稀少的异常点更"友好"', 12, ACCENT_GREEN, True),
])

# 右侧 - 树形结构示意
add_rect(slide, 6.8, 2.0, 6.0, 4.8, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.1, 5.6, 0.4, '贯穿示例：iTree 分裂过程', 14, ACCENT_GREEN, True)

# 简化树形图
# 根节点
add_rect(slide, 8.8, 2.6, 2.5, 0.5, RGBColor(0x1A, 0x3A, 0x5C), ACCENT_GREEN)
add_textbox(slide, 8.9, 2.63, 2.3, 0.4, 'x < 8.0?', 13, WHITE, True, PP_ALIGN.CENTER)

# 左子节点
add_rect(slide, 7.2, 3.5, 2.2, 0.5, RGBColor(0x1A, 0x3A, 0x5C), ACCENT_CYAN)
add_textbox(slide, 7.3, 3.53, 2.0, 0.4, 'y < 10.0?', 12, WHITE, True, PP_ALIGN.CENTER)

# 右子节点 - A1被孤立
add_rect(slide, 10.0, 3.5, 2.2, 0.5, RGBColor(0x2A, 0x1A, 0x0A), ACCENT_ORANGE)
add_textbox(slide, 10.1, 3.53, 2.0, 0.4, 'A1  size=1', 12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
add_textbox(slide, 10.1, 3.2, 0.6, 0.3, '路径=1', 10, ACCENT_GREEN, True)

# A2被孤立
add_rect(slide, 10.0, 4.4, 2.2, 0.5, RGBColor(0x2A, 0x1A, 0x0A), ACCENT_ORANGE)
add_textbox(slide, 10.1, 4.43, 2.0, 0.4, 'A2  size=1', 12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
add_textbox(slide, 10.1, 4.1, 0.6, 0.3, '路径=2', 10, ACCENT_GREEN, True)

# 正常点继续
add_rect(slide, 7.2, 4.4, 2.2, 0.5, RGBColor(0x1A, 0x3A, 0x5C), ACCENT_CYAN)
add_textbox(slide, 7.3, 4.43, 2.0, 0.4, '继续分割...', 12, MID_GRAY, False, PP_ALIGN.CENTER)

add_textbox(slide, 7.3, 5.1, 5.3, 0.4, '10个正常点需6-8次分裂才能逐一分离', 11, MID_GRAY)

# 房间类比
add_rect(slide, 7.0, 5.7, 5.6, 1.0, BG_CARD)
add_textbox(slide, 7.2, 5.75, 5.2, 0.9, '🏠  iTree 就像把大房间不断随机竖墙隔开。独自站在角落的人（异常点）一两堵墙就隔出来了；挤在人群中的人需要很多堵墙。', 10, MID_GRAY)


# ════════════════════════════════════════
# 第 6 页：构建森林
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '构建孤立森林', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从一棵树到一片森林——子采样 + 多树集成 = 又快又准', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('三步构建森林', 16, ACCENT_GREEN, True),
    ('', 4),
    ('Step 1：随机子采样', 14, WHITE, True),
    ('从训练数据中随机抽取 256 个样本。三个原因：', 12, LIGHT_GRAY),
    ('  • 消除屏蔽效应——稀释正常点，让异常点更容易被看到', 11, LIGHT_GRAY),
    ('  • 计算效率——无论原始数据 1 万还是 1 亿，训练时间可控', 11, LIGHT_GRAY),
    ('  • 增加多样性——不同子采样让每棵树看到不同数据', 11, LIGHT_GRAY),
    ('', 4),
    ('Step 2：构建 iTree', 14, WHITE, True),
    ('按递归分裂方法构建，高度限制 ⌈log₂256⌉ = 8 层。', 12, LIGHT_GRAY),
    ('8 层足以孤立绝大多数异常点（1-4 层就会被分离）。', 12, LIGHT_GRAY),
    ('', 4),
    ('Step 3：重复 100 次', 14, WHITE, True),
    ('得到 100 棵独立的 iTree，组成孤立森林。', 12, LIGHT_GRAY),
    ('从 100 增加到 500 棵，AUC 提升通常 < 0.01。', 12, ACCENT_CYAN),
])

# 右侧 - 路径对比表
add_rect(slide, 6.8, 2.0, 6.0, 3.2, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.1, 5.6, 0.4, '📊  5 棵树路径对比', 14, ACCENT_GREEN, True)

table_data2 = [
    ('树编号', 'A1 路径', 'P9 路径'),
    ('Tree 1', '1', '5.0'),
    ('Tree 2', '2', '4.5'),
    ('Tree 3', '1', '5.5'),
    ('Tree 4', '2', '4.0'),
    ('Tree 5', '1', '5.0'),
    ('平均', '1.4', '4.8'),
]
rows2 = len(table_data2)
tbl2 = slide.shapes.add_table(rows2, 3, Inches(7.2), Inches(2.6), Inches(5.2), Inches(2.4)).table
tbl2.columns[0].width = Inches(1.5)
tbl2.columns[1].width = Inches(1.85)
tbl2.columns[2].width = Inches(1.85)

for r in range(rows2):
    for c in range(3):
        cell = tbl2.cell(r, c)
        cell.text = table_data2[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.name = FONT_CN
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r == rows2-1:
                para.font.bold = True
                para.font.color.rgb = ACCENT_GREEN
            else:
                para.font.color.rgb = LIGHT_GRAY
            para.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r == rows2-1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

# 安检员类比
add_rect(slide, 6.8, 5.5, 6.0, 1.5, BG_CARD)
add_multiline_textbox(slide, 7.0, 5.6, 5.6, 1.3, [
    ('🕵️  100 个"蒙眼安检员"', 14, ACCENT_GREEN, True),
    ('每个安检员随机抽查 256 人，用随机画线的方式分割人群。单个安检员可能有偏差，但 100 个安检员的集体投票非常可靠——大多数安检员都很快挑出的人，几乎一定是异常的。', 11, MID_GRAY),
])


# ════════════════════════════════════════
# 第 7 页：异常分数计算
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '异常分数计算', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用公式量化"有多异常"', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('异常分数公式', 16, ACCENT_GREEN, True),
    ('s(x, n) = 2^(-h̄(x) / c(n))', 20, WHITE, True),
    ('', 4),
    ('h̄(x)：数据点 x 在所有树中的平均路径长度', 13, LIGHT_GRAY),
    ('c(n)：n 个点的期望路径长度 = 2H(n-1) - 2(n-1)/n', 13, LIGHT_GRAY),
    ('H：调和数 ≈ ln(k) + 0.5772（欧拉常数）', 12, MID_GRAY),
    ('', 6),
    ('分数解读', 16, ACCENT_GREEN, True),
    ('s → 1：路径极短，几乎在根部就被孤立 → 极异常', 13, ACCENT_ORANGE),
    ('s ≈ 0.5：路径与平均相当 → 正常（分界线）', 13, ACCENT_CYAN),
    ('s → 0：路径极长，深深隐藏在"人群"中 → 极正常', 13, LIGHT_GRAY),
    ('', 6),
    ('直觉：c(n) 是"平均孤立难度"，当 h̄ = c(n) 时', 12, MID_GRAY),
    ('s = 2^(-1) = 0.5，即"既不异常也不特别正常"', 12, MID_GRAY),
])

# 右侧 - 数值计算
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.1, 5.6, 0.4, '📐  贯穿示例计算', 14, ACCENT_GREEN, True)

add_multiline_textbox(slide, 7.2, 2.6, 5.4, 4.0, [
    ('c(12) = 2 × 3.0199 - 2 × 11/12 = 4.2065', 13, WHITE, True),
    ('', 8),
    ('🔴 异常点 A1：h̄ = 1.4', 14, ACCENT_ORANGE, True),
    ('s(A1) = 2^(-1.4 / 4.2065) = 2^(-0.333)', 13, LIGHT_GRAY),
    ('s(A1) ≈ 0.795 → 明显异常 ✓', 15, ACCENT_GREEN, True),
    ('', 8),
    ('🟢 正常点 P9：h̄ = 4.8', 14, ACCENT_CYAN, True),
    ('s(P9) = 2^(-4.8 / 4.2065) = 2^(-1.141)', 13, LIGHT_GRAY),
    ('s(P9) ≈ 0.454 → 正常 ✓', 15, ACCENT_CYAN, True),
    ('', 8),
    ('0 ─── 0.454 ─ 0.5 ─── 0.795 ─── 1.0', 13, WHITE, True, PP_ALIGN.CENTER),
    ('  正常 ←── 分界线 ──→ 异常', 12, MID_GRAY, False, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════
# 第 8 页：完整示例
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '完整示例——从数据到结论', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用 12 个数据点走一遍"训练→评分→判定"全流程', 16, ACCENT_CYAN)

# 左侧 - 流程
add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('① 准备数据', 14, ACCENT_GREEN, True),
    ('10 个正常点聚集在 (3-5, 4-5.5)，2 个异常点 A1(20,3) 和 A2(2,25) 远远游离。', 12, LIGHT_GRAY),
    ('', 4),
    ('② 构建一棵 iTree', 14, ACCENT_GREEN, True),
    ('第 1 次分裂：选 x 轴，分裂点 x=8.0', 12, LIGHT_GRAY),
    ('  → 左子树：10 个正常点 + A2（11 个点）', 11, MID_GRAY),
    ('  → 右子树：A1（仅 1 个）→ A1 被孤立！路径=1', 11, ACCENT_ORANGE),
    ('第 2 次分裂：选 y 轴，分裂点 y=10.0', 12, LIGHT_GRAY),
    ('  → 左子树：10 个正常点 | 右子树：A2 → 路径=2', 11, ACCENT_ORANGE),
    ('', 4),
    ('③ 构建 100 棵 iTree', 14, ACCENT_GREEN, True),
    ('重复 100 次（不同子采样和随机分裂），统计平均路径：', 12, LIGHT_GRAY),
    ('  A1 平均路径 h̄ ≈ 1.4 | A2 平均路径 h̄ ≈ 2.1', 12, ACCENT_CYAN),
    ('  正常点平均路径 h̄ ≈ 4.5-5.0', 12, LIGHT_GRAY),
    ('', 4),
    ('④ 计算异常分数（c(12) = 4.2065）', 14, ACCENT_GREEN, True),
    ('  s(A1) ≈ 0.795 → 🔴 异常', 13, ACCENT_ORANGE, True),
    ('  s(A2) ≈ 0.707 → 🔴 异常', 13, ACCENT_ORANGE, True),
    ('  s(P9) ≈ 0.454 → 🟢 正常', 13, ACCENT_CYAN, True),
])

# 右侧 - 结果汇总
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.1, 5.6, 0.4, '📊  结果汇总', 14, ACCENT_GREEN, True)

table_data3 = [
    ('数据点', '平均路径', '异常分数', '判定'),
    ('A1(20, 3)', '1.4', '0.795', '🔴 异常'),
    ('A2(2, 25)', '2.1', '0.707', '🔴 异常'),
    ('P1~P10', '4.5~5.0', '0.45~0.53', '🟢 正常'),
]
tbl3 = slide.shapes.add_table(4, 4, Inches(7.2), Inches(2.6), Inches(5.4), Inches(2.0)).table
tbl3.columns[0].width = Inches(1.5)
tbl3.columns[1].width = Inches(1.2)
tbl3.columns[2].width = Inches(1.2)
tbl3.columns[3].width = Inches(1.5)

for r in range(4):
    for c in range(4):
        cell = tbl3.cell(r, c)
        cell.text = table_data3[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.name = FONT_CN
            para.alignment = PP_ALIGN.CENTER
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            elif r in (1, 2):
                para.font.color.rgb = ACCENT_ORANGE
            else:
                para.font.color.rgb = ACCENT_CYAN
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        elif r in (1, 2):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x15, 0x0A)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

add_multiline_textbox(slide, 7.2, 5.0, 5.4, 1.6, [
    ('✅  结论', 14, ACCENT_GREEN, True),
    ('孤立森林成功识别 A1 和 A2 两个异常点（分数 0.795 和 0.707，远超 0.5 分界线），同时正确判定所有正常点为正常（分数 0.45-0.53）。', 12, LIGHT_GRAY),
    ('', 4),
    ('筛子类比：100 张随机网格去筛数据，大的孤立颗粒（异常点）在几乎所有网格中都是第一批被筛出来的。', 10, MID_GRAY),
])


# ════════════════════════════════════════
# 第 9 页：参数调优与优缺点
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '参数调优与优缺点', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '默认参数就能覆盖大多数场景——但也要知道局限', 16, ACCENT_CYAN)

# 左侧 - 参数
add_rect(slide, 0.8, 2.0, 5.5, 2.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.0, 2.1, 5.2, 2.3, [
    ('⚙️  关键参数', 15, ACCENT_GREEN, True),
    ('', 4),
    ('树的数量（n_estimators）：默认 100', 13, WHITE, True),
    ('就像民意调查，100 份已足够稳定。100→500 棵 AUC 提升 <0.01', 11, LIGHT_GRAY),
    ('', 4),
    ('子采样大小（max_samples）：默认 256', 13, WHITE, True),
    ('论文实验的"甜蜜点"——太小信息不足，太大正常点互相遮挡', 11, LIGHT_GRAY),
    ('', 4),
    ('异常阈值（contamination）：默认 auto（s=0.5）', 13, WHITE, True),
    ('如有业务知识（欺诈率约 0.1%），可直接设 0.001', 11, LIGHT_GRAY),
])

# 左下 - 优点
add_rect(slide, 0.8, 4.8, 5.5, 2.3, BG_CARD, RGBColor(0x0A, 0x2A, 0x15))
add_multiline_textbox(slide, 1.0, 4.9, 5.2, 2.1, [
    ('✅  核心优点', 15, ACCENT_GREEN, True),
    ('• 线性时间复杂度 O(t·ψ·logψ)，轻松处理百万级数据', 12, LIGHT_GRAY),
    ('• 天然支持并行化，100 棵树可分配到 100 个 CPU 核心', 12, LIGHT_GRAY),
    ('• 不依赖距离度量，天然免疫"维度灾难"', 12, LIGHT_GRAY),
    ('• 纯无监督，不需要任何标注数据', 12, LIGHT_GRAY),
    ('• 内存高效：100×256 = 25,600 个数据点即可', 12, LIGHT_GRAY),
])

# 右侧 - 局限
add_rect(slide, 6.8, 2.0, 6.0, 2.5, BG_CARD, RGBColor(0x2A, 0x15, 0x0A))
add_multiline_textbox(slide, 7.0, 2.1, 5.6, 2.3, [
    ('⚠️  注意局限', 15, ACCENT_ORANGE, True),
    ('', 4),
    ('• 局部异常检测能力弱——主要捕获全局异常', 12, LIGHT_GRAY),
    ('  如果需要检测局部异常，先用 iForest 筛查，再用 LOF 精检', 11, MID_GRAY),
    ('• 轴对齐切割偏差——沿非坐标轴方向分布的数据可能产生"伪异常"', 12, LIGHT_GRAY),
    ('  Extended Isolation Forest（2019）用任意角度超平面解决', 11, MID_GRAY),
    ('• 对分类特征处理不佳，需要先做独热编码', 12, LIGHT_GRAY),
])

# 右下 - 选择建议
add_rect(slide, 6.8, 4.8, 6.0, 2.3, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 4.9, 5.6, 2.1, [
    ('🎯  选择建议', 15, ACCENT_CYAN, True),
    ('', 4),
    ('数据量大（>10 万条）→ 首选孤立森林', 13, ACCENT_GREEN, True),
    ('需要检测局部异常 → 选 LOF', 13, LIGHT_GRAY),
    ('数据量小（<5000 条）且边界复杂 → 选 One-Class SVM', 13, LIGHT_GRAY),
    ('不确定 → 先试孤立森林（速度快、参数少）', 13, ACCENT_CYAN, True),
    ('', 4),
    ('孤立森林像瑞士军刀——不是每种场景都最佳，但足够通用、方便、快速', 11, MID_GRAY),
])


# ════════════════════════════════════════
# 第 10 页：实际应用
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '实际应用场景', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从理论到落地——异常检测的"第一道防线"', 16, ACCENT_CYAN)

# 四个场景卡片
cards = [
    ('🌐 网络入侵检测', '实时识别恶意连接', '将日志转为数值特征，训练后实时评分。KDD Cup 数据集检测准确率 >95%。', ACCENT_GREEN),
    ('💰 金融欺诈检测', '监控异常交易', '增加衍生特征（1h 交易次数等），contamination=0.001。Kaggle 数据集可检测 ~80-85% 欺诈。', ACCENT_GREEN),
    ('🏭 设备故障预警', '传感器数据监控', '时间序列分段提取统计特征，多传感器融合。NASA 数据集可提前 15-25 天预警。', ACCENT_CYAN),
    ('📋 日志异常检测', '服务器集群监控', '按时间窗口聚合，训练模型实时检测。比单维度阈值方法更有效。', ACCENT_CYAN),
]

for i, (title, subtitle, desc, color) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.0 + row * 2.6
    add_rect(slide, x, y, 5.8, 2.3, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
    add_multiline_textbox(slide, x + 0.2, y + 0.1, 5.4, 2.1, [
        (title, 14, color, True),
        (subtitle, 12, WHITE, True),
        ('', 2),
        (desc, 11, LIGHT_GRAY),
    ])

# 底部 - 工业界采用
add_rect(slide, 0.8, 5.6, 11.7, 1.5, BG_CARD)
add_multiline_textbox(slide, 1.0, 5.7, 11.3, 1.3, [
    ('🏢  工业界采用', 14, ACCENT_GREEN, True),
    ('Google：监控内部网络健康  |  Microsoft Azure：异常检测服务  |  Netflix：用户行为异常检测', 12, LIGHT_GRAY),
    ('阿里巴巴：双十一大促交易欺诈检测，每秒处理数十万笔交易', 12, LIGHT_GRAY),
    ('', 2),
    ('孤立森林就像"烟雾报警器"——不一定告诉你火灾原因，但能在最早阶段发出警报，为精准排查争取时间。', 11, MID_GRAY),
])


# ════════════════════════════════════════
# 第 11 页：结尾
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '关键收获与下一步', 32, WHITE, True)

add_multiline_textbox(slide, 0.8, 1.5, 6.5, 4.0, [
    ('五句话总结孤立森林', 18, ACCENT_GREEN, True),
    ('', 6),
    ('① 新范式：不建模正常，直接度量"被孤立的难易程度"', 14, LIGHT_GRAY),
    ('② 简洁实现：100 棵随机分割二叉树 + 路径长度 + 异常分数公式', 14, LIGHT_GRAY),
    ('③ 极致效率：线性时间复杂度，天然并行，百万级数据轻松处理', 14, LIGHT_GRAY),
    ('④ 广泛适用：网络安全、金融反欺诈、工业预警、IT 运维', 14, LIGHT_GRAY),
    ('⑤ 持续进化：Extended Isolation Forest 解决轴对齐偏差', 14, LIGHT_GRAY),
    ('', 8),
    ('行动建议', 18, ACCENT_GREEN, True),
    ('下次遇到异常检测任务，先用孤立森林建立基线。', 14, WHITE, True),
    ('默认参数（100 棵树、256 子采样）在大多数场景下就够用。', 13, LIGHT_GRAY),
    ('效果不够再考虑 LOF 或 One-Class SVM。', 13, LIGHT_GRAY),
])

# 右侧 - 关键数字
add_rect(slide, 7.8, 1.5, 4.8, 3.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 8.0, 1.6, 4.4, 0.4, '🔢  关键数字速记', 14, ACCENT_GREEN, True)

table_data4 = [
    ('项目', '数值'),
    ('默认树数量', '100'),
    ('默认子采样', '256'),
    ('默认树高度', '8'),
    ('异常分界线', 's = 0.5'),
    ('训练复杂度', 'O(n)'),
]
tbl4 = slide.shapes.add_table(6, 2, Inches(8.2), Inches(2.2), Inches(4.0), Inches(2.5)).table
tbl4.columns[0].width = Inches(2.0)
tbl4.columns[1].width = Inches(2.0)

for r in range(6):
    for c in range(2):
        cell = tbl4.cell(r, c)
        cell.text = table_data4[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.name = FONT_CN
            para.alignment = PP_ALIGN.CENTER
            if r == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            else:
                para.font.color.rgb = ACCENT_CYAN if c == 1 else LIGHT_GRAY
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

# 论文引用
add_rect(slide, 0.8, 5.5, 11.7, 1.5, BG_CARD)
add_multiline_textbox(slide, 1.0, 5.6, 11.3, 1.3, [
    ('📄  论文出处', 14, ACCENT_GREEN, True),
    ('Liu, Ting & Zhou (2008) "Isolation Forest" IEEE ICDM  |  Liu, Ting & Zhou (2012) "Isolation-based Anomaly Detection" ACM TKDD', 11, LIGHT_GRAY),
    ('Hariri, Carrasco Kind & McNutt (2019) "Extended Isolation Forest" IEEE TKDE', 11, LIGHT_GRAY),
    ('', 4),
    ('"一切都应该尽可能简单，但不能过于简单。" —— 爱因斯坦。孤立森林恰好停在"足够简单但不失效"的完美位置。', 12, MID_GRAY),
])

# ── 保存 ──
output_path = '/home/admin/.openclaw/workspace-weaver/output/isolation-forest_v1.pptx'
prs.save(output_path)
print(f'✅ PPT 已保存到 {output_path}')
print(f'总页数：{len(prs.slides)}')
