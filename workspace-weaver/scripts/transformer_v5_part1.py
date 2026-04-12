#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Transformer PPT v5 Part1 - Pages 1-7"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# === Constants ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_R = Inches(13.2)
SAFE_B = Inches(7.3)

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTN = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
C_BG_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
C_BG_RED = RGBColor(0xFD, 0xED, 0xEC)
C_BG_YELLOW = RGBColor(0xFE, 0xF9, 0xE7)
C_BG_GREEN = RGBColor(0xEA, 0xFA, 0xEA)
C_BG_PURPLE = RGBColor(0xF4, 0xEC, 0xF7)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=Pt(11), font_color=C_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    # line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Return (textbox, text_frame) for multi-paragraph use."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_paragraph(tf, text, font_size=Pt(11), font_color=C_TEXT, bold=False,
                  font_name=FONT_CN, alignment=PP_ALIGN.LEFT, space_after=Pt(4),
                  space_before=Pt(0), line_spacing=1.15, first_indent=None):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    # line spacing
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    if first_indent:
        indent = pPr.makeelement(qn('a:indent'), {'first_line': str(first_indent)})
        pPr.append(indent)
    return p


def add_rounded_rect(slide, left, top, width, height, text="", fill_color=None,
                     line_color=None, line_width=Pt(1.5), font_size=Pt(10),
                     font_color=C_TEXT, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = FONT_CN
        tf.paragraphs[0].space_after = Pt(0)
        tf.paragraphs[0].space_before = Pt(0)
        shape.text_frame.auto_size = None
    return shape


def add_arrow_line(slide, start_x, start_y, end_x, end_y, color=C_TEXT, width=Pt(1.5)):
    """Add a line with arrowhead."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # add end arrow
    ln = connector.line._ln
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)
    return connector


def add_line(slide, start_x, start_y, end_x, end_y, color=C_TEXT, width=Pt(1)):
    connector = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def set_shape_gradient(shape, color1, color2):
    """Set linear gradient fill."""
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = color1
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[1].color.rgb = color2
    shape.fill.gradient_stops[1].position = 1.0


def count_text(slide):
    """Count Chinese chars + words on a slide."""
    total = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text
                # count Chinese chars
                cn = sum(1 for c in t if '\u4e00' <= c <= '\u9fff')
                # count English words
                en = len([w for w in t.split() if w.strip()])
                total += cn + en
    return total


# ============================================================
# PAGE 1: Cover
# ============================================================
def make_page1():
    slide = prs.slides.add_slide(blank_layout)

    # Background gradient bar at top
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(2.8))
    set_shape_gradient(bar, RGBColor(0x1E, 0x3A, 0x5F), RGBColor(0x2C, 0x3E, 0x50))

    # Decorative lines
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.9), Inches(10.3), Pt(2.5),
              fill_color=C_ENCODER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(3.05), Inches(9.3), Pt(2),
              fill_color=C_ATTN)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(3.18), Inches(8.3), Pt(1.5),
              fill_color=C_FFN)

    # Main title
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.2),
                "Transformer \u67b6\u6784\u6df1\u5ea6\u89e3\u6790",
                font_size=Pt(36), font_color=C_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(0.6),
                "\u4ece\u76f4\u89c9\u5230\u539f\u7406\uff0c\u4e00\u7bc7\u770b\u61c2AI\u7684\u6838\u5fc3\u5f15\u64ce",
                font_size=Pt(16), font_color=RGBColor(0xBD, 0xC3, 0xC7),
                alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    # Paper info
    add_textbox(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.4),
                "Attention Is All You Need \u2014 Vaswani et al., 2017 | Google Brain",
                font_size=Pt(11), font_color=C_GRAY,
                alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

    # Bottom info box
    add_rounded_rect(slide, Inches(3.5), Inches(4.2), Inches(6.3), Inches(0.45),
                     "\u7f16\u7801\u5668 - \u89e3\u7801\u5668 \u67b6\u6784 | 6500\u4e07\u53c2\u6570 | 8\u5934\u6ce8\u610f\u529b | WMT BLEU 28.4",
                     fill_color=C_BG_LIGHT, line_color=C_TITLE, font_size=Pt(10),
                     font_color=C_TITLE, bold=True)

    # Architecture mini icons at bottom
    colors = [C_ENCODER, C_ATTN, C_FFN, C_DECODER, C_EMBED]
    labels = ["\u7f16\u7801\u5668", "\u6ce8\u610f\u529b", "FFN", "\u89e3\u7801\u5668", "\u5d4c\u5165"]
    for i, (c, lb) in enumerate(zip(colors, labels)):
        x = Inches(2.8) + Inches(1.7) * i
        add_rounded_rect(slide, x, Inches(5.0), Inches(1.4), Inches(0.5),
                         text=lb, fill_color=c, font_size=Pt(10),
                         font_color=C_WHITE, bold=True)

    # Bottom bar
    bar2 = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4))
    set_shape_gradient(bar2, RGBColor(0x2C, 0x3E, 0x50), RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(slide, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35),
                "\u81ea\u6ce8\u610f\u529b | \u4f4d\u7f6e\u7f16\u7801 | \u6b8b\u5dee\u8fde\u63a5 | \u5c42\u5f52\u4e00\u5316 | \u591a\u5934\u673a\u5236 | \u56e0\u679c\u63a9\u7801",
                font_size=Pt(9), font_color=RGBColor(0xBD, 0xC3, 0xC7),
                alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    return slide


# ============================================================
# PAGE 2: Why Transformer
# ============================================================
def make_page2():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(8), Inches(0.45),
                "\u4e3a\u4ec0\u4e48\u9700\u8981Transformer",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    # Title underline
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_ENCODER)

    # --- LEFT: Diagrams ---
    # RNN Flow (top-left)
    add_textbox(slide, Inches(0.3), Inches(0.75), Inches(2.5), Inches(0.3),
                "RNN \u987a\u5e8f\u5904\u7406\uff08\u6392\u961f\u53d1\u8a00\uff09",
                font_size=Pt(10), font_color=C_DECODER, bold=True)

    words_rnn = ["The", "cat", "sat", "on", "the", "mat"]
    for i, w in enumerate(words_rnn):
        x = Inches(0.4) + Inches(1.1) * i
        y = Inches(1.1)
        add_rounded_rect(slide, x, y, Inches(0.9), Inches(0.4),
                         text=w, fill_color=C_BG_RED, line_color=C_DECODER,
                         font_size=Pt(9), font_color=C_DECODER, bold=True)
        if i > 0:
            add_arrow_line(slide, x - Inches(0.05), y + Inches(0.2),
                           x + Inches(0.05), y + Inches(0.2), C_DECODER, Pt(1.5))
    # Time annotation
    add_textbox(slide, Inches(0.4), Inches(1.55), Inches(6.5), Inches(0.25),
                "-> \u9010\u6b65\u5904\u7406\uff0c1000\u8bcd = 1000\u6b65\uff0cGPU\u5e76\u884c\u5b8c\u5168\u6d6a\u8d39",
                font_size=Pt(8), font_color=C_GRAY)

    # Attention visualization (bottom-left)
    add_textbox(slide, Inches(0.3), Inches(1.95), Inches(3), Inches(0.3),
                "\u81ea\u6ce8\u610f\u529b\u53ef\u89c6\u5316\uff08\u5706\u684c\u4f1a\u8bae\uff09",
                font_size=Pt(10), font_color=C_ENCODER, bold=True)

    attn_words = ["The", "cat", "sat", "on", "it", "..."]
    attn_weights = {
        4: [(0, 0.05, Pt(0.5)), (1, 0.60, Pt(3)), (2, 0.10, Pt(1)), (3, 0.05, Pt(0.5)), (5, 0.20, Pt(1.2))]
    }
    for i, w in enumerate(attn_words):
        x = Inches(0.4) + Inches(1.1) * i
        y = Inches(2.35)
        if i == 4:  # "it"
            add_rounded_rect(slide, x, y, Inches(0.9), Inches(0.4),
                             text=w, fill_color=C_ENCODER, line_color=C_ENCODER,
                             font_size=Pt(9), font_color=C_WHITE, bold=True)
        else:
            add_rounded_rect(slide, x, y, Inches(0.9), Inches(0.4),
                             text=w, fill_color=C_BG_BLUE, line_color=C_ENCODER,
                             font_size=Pt(9), font_color=C_ENCODER)

    # Draw attention lines from "it" (index 4)
    it_x = Inches(0.4) + Inches(1.1) * 4 + Inches(0.45)
    it_y = Inches(2.35)
    for j, weight, lw in attn_weights[4]:
        target_x = Inches(0.4) + Inches(1.1) * j + Inches(0.45)
        add_arrow_line(slide, it_x, it_y + Inches(0.4), target_x, it_y + Inches(0.65),
                       C_ENCODER, lw)
        if weight >= 0.10:
            add_textbox(slide, target_x - Inches(0.25), it_y + Inches(0.65), Inches(0.5), Inches(0.2),
                        f"{weight:.0%}", font_size=Pt(7), font_color=C_ENCODER, bold=True,
                        alignment=PP_ALIGN.CENTER)

    # Bottom annotation
    add_textbox(slide, Inches(0.4), Inches(3.15), Inches(6.5), Inches(0.25),
                "\u201cit\u201d\u4e0e\u201ccat\u201d\u5efa\u7acb\u6700\u5f3a\u5173\u8054(0.6)\uff0c\u6240\u6709\u8bcd\u540c\u65f6\u4ea4\u6d41",
                font_size=Pt(8), font_color=C_GRAY)

    # Comparison box
    add_rounded_rect(slide, Inches(0.3), Inches(3.55), Inches(3.2), Inches(0.8),
                     fill_color=C_BG_YELLOW, line_color=C_ATTN, font_size=Pt(9),
                     font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT)
    # Fill text manually
    _, tf = add_rich_textbox(slide, Inches(0.4), Inches(3.58), Inches(3.0), Inches(0.75))
    add_paragraph(tf, "[!] RNN: \u6392\u961f\u53d1\u8a00 -> \u4e32\u884c\uff0c\u6162", font_size=Pt(9),
                  font_color=C_DECODER, bold=True, space_after=Pt(2))
    add_paragraph(tf, "V Transformer: \u5706\u684c\u4f1a\u8bae -> \u5e76\u884c\uff0c\u5feb", font_size=Pt(9),
                  font_color=C_FFN, bold=True, space_after=Pt(2))
    add_paragraph(tf, "V \u4efb\u610f\u8bcd\u5bf9\u76f4\u63a5\u4ea4\u4e92\uff0c\u65e0\u8ddd\u79bb\u9650\u5236", font_size=Pt(9),
                  font_color=C_FFN, bold=True)

    # --- RIGHT: Key Points ---
    rx = Inches(7.8)
    ry = Inches(0.75)
    rw = Inches(5.2)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(3.7))
    points = [
        ("(1) \u987a\u5e8f\u5904\u7406\u74f6\u9888", "RNN\u5fc5\u987b\u9010\u8bcd\u5904\u7406\uff0cGPU\u5e76\u884c\u80fd\u529b\u5b8c\u5168\u6d6a\u8d39\u30021000\u8bcd->1000\u6b65\u3002"),
        ("(2) \u957f\u8ddd\u79bb\u9057\u5fd8", "\u4fe1\u606f\u7ecf\u6570\u767e\u6b21\u77e9\u9635\u4e58\u6cd5\u6307\u6570\u7ea7\u8870\u51cf\uff0ch_t=tanh(W\u00b7h+\u00b7U\u00b7x_t)\u3002"),
        ("(3) \u4fe1\u606f\u74f6\u9888", "\u6574\u4e2a\u5e8f\u5217\u538b\u7f29\u6210512\u7ef4\u5411\u91cf\uff0c\u4fe1\u606f\u4e25\u91cd\u4e22\u5931\u3002"),
        ("(4) \u6838\u5fc3\u7a81\u7834", "\u81ea\u6ce8\u610f\u529b\u8ba9\u6bcf\u4e2a\u8bcd\u4e0e\u6240\u6709\u8bcd\u76f4\u63a5\u4ea4\u6d41\uff0c\u5706\u684c\u4f1a\u8bae vs \u6392\u961f\u53d1\u8a00\u3002"),
        ("(5) \u6570\u5b66\u672c\u8d28", "QxK^T\u8ba1\u7b97\u6240\u6709\u8bcd\u5bf9\u5173\u8054\u5ea6\uff0cTxT\u77e9\u9635\u5b8c\u5168\u5e76\u884c\u3002"),
    ]
    for title, desc in points:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(4))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(3), line_spacing=1.1)

    # Bottom divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.55), Inches(12.6), Pt(0.5),
              fill_color=C_BG_LIGHT)

    # Bottom section: key takeaway
    add_rounded_rect(slide, Inches(0.3), Inches(4.7), Inches(12.6), Inches(0.55),
                     fill_color=RGBColor(0xEB, 0xF5, 0xFB), line_color=C_ENCODER)
    add_textbox(slide, Inches(0.5), Inches(4.75), Inches(12.2), Inches(0.45),
                "[!] \u6838\u5fc3\u601d\u60f3\uff1a\u7528\u201c\u6ce8\u610f\u529b\u201d\u66ff\u4ee3\u201c\u5faa\u73af\u201d\uff0c\u4e00\u6b65\u5230\u4f4d\u89e3\u51b3\u4e09\u5927\u75db\u70b9\uff1a\u5e76\u884c\u5316 + \u5168\u5c40\u4ea4\u4e92 + \u65e0\u538b\u7f29",
                font_size=Pt(11), font_color=C_ENCODER, bold=True)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "1/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)

    return slide


# ============================================================
# PAGE 3: Architecture Overview
# ============================================================
def make_page3():
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.45),
                "\u67b6\u6784\u603b\u89c8\uff1a\u7f16\u7801\u5668-\u89e3\u7801\u5668",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_ENCODER)

    # --- LEFT: Architecture diagram ---
    lx = Inches(0.3)
    ly = Inches(0.8)

    # Input embedding + pos enc
    add_rounded_rect(slide, lx, ly, Inches(3.2), Inches(0.45),
                     "\u8f93\u5165: Embedding + \u4f4d\u7f6e\u7f16\u7801",
                     fill_color=C_BG_PURPLE, line_color=C_EMBED, font_size=Pt(9),
                     font_color=C_EMBED, bold=True)

    # Encoder block (big blue box)
    enc_x = lx
    enc_y = ly + Inches(0.55)
    enc_w = Inches(5.8)
    enc_h = Inches(2.6)
    add_rounded_rect(slide, enc_x, enc_y, enc_w, enc_h,
                     "", fill_color=None, line_color=C_ENCODER, line_width=Pt(2.5))
    add_textbox(slide, enc_x + Inches(0.1), enc_y + Inches(0.05), Inches(2.5), Inches(0.25),
                "\u7f16\u7801\u5668 x 6", font_size=Pt(11), font_color=C_ENCODER, bold=True)

    # Inside encoder: one layer detail
    inner_x = enc_x + Inches(0.15)
    inner_y = enc_y + Inches(0.35)
    # Self-Attention
    add_rounded_rect(slide, inner_x, inner_y, Inches(2.2), Inches(0.4),
                     "Multi-Head Self-Attention", fill_color=C_BG_BLUE, line_color=C_ENCODER,
                     font_size=Pt(8), font_color=C_ENCODER, bold=True)
    add_arrow_line(slide, inner_x + Inches(2.2), inner_y + Inches(0.2),
                   inner_x + Inches(2.4), inner_y + Inches(0.2), C_ENCODER, Pt(1))
    add_rounded_rect(slide, inner_x + Inches(2.4), inner_y, Inches(1.5), Inches(0.4),
                     "Add & LayerNorm", fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(8), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)
    # FFN
    add_rounded_rect(slide, inner_x, inner_y + Inches(0.5), Inches(2.2), Inches(0.4),
                     "Feed-Forward Network", fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(8), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)
    add_arrow_line(slide, inner_x + Inches(2.2), inner_y + Inches(0.7),
                   inner_x + Inches(2.4), inner_y + Inches(0.7), C_ENCODER, Pt(1))
    add_rounded_rect(slide, inner_x + Inches(2.4), inner_y + Inches(0.5), Inches(1.5), Inches(0.4),
                     "Add & LayerNorm", fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(8), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)

    # Vertical arrows inside encoder (6 layers indication)
    for yy in [inner_y + Inches(1.05), inner_y + Inches(1.5), inner_y + Inches(1.95)]:
        add_textbox(slide, inner_x + Inches(1.0), yy, Inches(2.5), Inches(0.3),
                    "... \u91cd\u590d 6 \u5c42 ...", font_size=Pt(8), font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # K,V output arrows from encoder
    kv_y = enc_y + enc_h
    add_textbox(slide, enc_x + Inches(0.2), kv_y, Inches(1.5), Inches(0.25),
                "K, V v", font_size=Pt(9), font_color=C_ENCODER, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Decoder block (big red box)
    dec_x = enc_x
    dec_y = kv_y + Inches(0.3)
    dec_w = enc_w
    dec_h = Inches(2.6)
    add_rounded_rect(slide, dec_x, dec_y, dec_w, dec_h,
                     "", fill_color=None, line_color=C_DECODER, line_width=Pt(2.5))
    add_textbox(slide, dec_x + Inches(0.1), dec_y + Inches(0.05), Inches(2.5), Inches(0.25),
                "\u89e3\u7801\u5668 x 6", font_size=Pt(11), font_color=C_DECODER, bold=True)

    # Inside decoder
    d_inner_x = dec_x + Inches(0.15)
    d_inner_y = dec_y + Inches(0.35)
    # Masked Self-Attention
    add_rounded_rect(slide, d_inner_x, d_inner_y, Inches(2.0), Inches(0.38),
                     "Masked Self-Attn", fill_color=C_BG_RED, line_color=C_DECODER,
                     font_size=Pt(8), font_color=C_DECODER, bold=True)
    # Cross-Attention
    add_rounded_rect(slide, d_inner_x, d_inner_y + Inches(0.45), Inches(2.0), Inches(0.38),
                     "Cross-Attention", fill_color=C_BG_PURPLE, line_color=C_EMBED,
                     font_size=Pt(8), font_color=C_EMBED, bold=True)
    # FFN
    add_rounded_rect(slide, d_inner_x, d_inner_y + Inches(0.9), Inches(2.0), Inches(0.38),
                     "Feed-Forward", fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(8), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)

    # Arrows between decoder sub-layers
    add_arrow_line(slide, d_inner_x + Inches(1.0), d_inner_y + Inches(0.38),
                   d_inner_x + Inches(1.0), d_inner_y + Inches(0.45), C_DECODER, Pt(1))
    add_arrow_line(slide, d_inner_x + Inches(1.0), d_inner_y + Inches(0.83),
                   d_inner_x + Inches(1.0), d_inner_y + Inches(0.9), C_DECODER, Pt(1))

    # K,V arrow from encoder to decoder cross-attention
    kv_arrow_x = enc_x + enc_w + Inches(0.15)
    add_arrow_line(slide, kv_arrow_x, enc_y + enc_h,
                   kv_arrow_x, d_inner_y + Inches(0.64), C_EMBED, Pt(2))
    add_textbox(slide, kv_arrow_x + Inches(0.05), Inches(4.1), Inches(0.8), Inches(0.25),
                "K, V", font_size=Pt(9), font_color=C_EMBED, bold=True)

    # Output
    out_y = dec_y + dec_h + Inches(0.15)
    add_rounded_rect(slide, dec_x, out_y, Inches(1.5), Inches(0.4),
                     "Linear", fill_color=C_BG_LIGHT, line_color=C_TEXT,
                     font_size=Pt(9), font_color=C_TEXT, bold=True)
    add_arrow_line(slide, dec_x + Inches(1.5), out_y + Inches(0.2),
                   dec_x + Inches(1.7), out_y + Inches(0.2), C_TEXT, Pt(1))
    add_rounded_rect(slide, dec_x + Inches(1.7), out_y, Inches(1.5), Inches(0.4),
                     "Softmax", fill_color=C_BG_LIGHT, line_color=C_TEXT,
                     font_size=Pt(9), font_color=C_TEXT, bold=True)
    add_arrow_line(slide, dec_x + Inches(3.2), out_y + Inches(0.2),
                   dec_x + Inches(3.4), out_y + Inches(0.2), C_TEXT, Pt(1))
    add_rounded_rect(slide, dec_x + Inches(3.4), out_y, Inches(1.5), Inches(0.4),
                     "\u8f93\u51fa\u6982\u7387", fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)

    # Right side output labels
    add_textbox(slide, enc_x + Inches(4.0), enc_y + Inches(1.2), Inches(1.5), Inches(0.5),
                "\u7f16\u7801\u8f93\u51fa\n(512\u7ef4)", font_size=Pt(8), font_color=C_ENCODER,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, dec_x + Inches(4.0), dec_y + Inches(1.2), Inches(1.5), Inches(0.5),
                "\u89e3\u7801\u8f93\u51fa\n(\u6982\u7387)", font_size=Pt(8), font_color=C_DECODER,
                alignment=PP_ALIGN.CENTER)

    # --- RIGHT: Key Points ---
    rx = Inches(7.8)
    ry = Inches(0.8)
    rw = Inches(5.2)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(5.5))
    points3 = [
        ("\u7f16\u7801\u5668 = \u7f16\u8f91\u90e8",
         "6\u5c42\u5806\u53e0\uff0c\u6bcf\u5c42\u6ce8\u610f\u529b+FFN\u3002\u6bcf\u4e2a\u8bcd\u7684\u4fe1\u606f\u4e0e\u5176\u4ed6\u6240\u6709\u8bcd\u4e0d\u65ad\u878d\u5408\uff0c6\u5c42\u540e\u6bcf\u4e2a\u8bcd\u7684512\u7ef4\u5411\u91cf\u5df2\u8574\u542b\u6574\u4e2a\u53e5\u5b50\u7684\u4e0a\u4e0b\u6587\u3002\u53c2\u6570\u7ea63450\u4e07\u3002"),
        ("\u89e3\u7801\u5668 = \u7ffb\u8bd1\u90e8",
         "6\u5c42\uff0c\u591a\u4e00\u4e2a\u4ea4\u53c9\u6ce8\u610f\u529b\u5b50\u5c42\u3002\u63a9\u7801\u81ea\u6ce8\u610f\u529b\u786e\u4fdd\u53ea\u80fd\u770b\u4e4b\u524d\u7684\u8bcd\uff0c\u4ea4\u53c9\u6ce8\u610f\u529b\u8ba9\u89e3\u7801\u5668\u201c\u67e5\u9605\u201d\u7f16\u7801\u5668\u8f93\u51fa\u3002"),
        ("\u6570\u636e\u6d41",
         "\u5206\u8bcd(BPE) -> \u5d4c\u5165(\u67e5\u8868) -> \u4f4d\u7f6e\u7f16\u7801(sin/cos) -> \u7f16\u7801\u56686\u5c42 -> \u89e3\u7801\u56686\u5c42 -> Softmax -> \u8f93\u51fa\u3002"),
        ("\u5173\u952e\u6570\u5b57",
         "512\u7ef4 | 8\u5934 | FFN 2048\u7ef4 | \u603b\u53c2\u65706500\u4e07 | WMT BLEU 28.4"),
    ]
    for title, desc in points3:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(6))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(4), line_spacing=1.1)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "2/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)
    return slide


# ============================================================
# PAGE 4: Encoder Detail
# ============================================================
def make_page4():
    slide = prs.slides.add_slide(blank_layout)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.45),
                "\u7f16\u7801\u5668\u5355\u5c42\uff1a\u4ea4\u6d41\u4e0e\u63d0\u70bc",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_ENCODER)

    # --- LEFT: Data flow diagram ---
    lx = Inches(0.3)
    ly = Inches(0.85)
    box_w = Inches(2.5)
    box_h = Inches(0.5)

    # Input
    add_rounded_rect(slide, lx, ly, box_w, box_h,
                     "X (512\u7ef4)", fill_color=C_BG_PURPLE, line_color=C_EMBED,
                     font_size=Pt(10), font_color=C_EMBED, bold=True)

    # Arrow
    ay1 = ly + box_h
    add_arrow_line(slide, lx + box_w / 2, ay1, lx + box_w / 2, ay1 + Inches(0.3), C_ENCODER, Pt(1.5))

    # Multi-Head Attention
    mha_y = ay1 + Inches(0.3)
    add_rounded_rect(slide, lx, mha_y, box_w, Inches(0.7),
                     "Multi-Head Attention\n(8\u5934 x 64\u7ef4)",
                     fill_color=C_BG_BLUE, line_color=C_ENCODER,
                     font_size=Pt(10), font_color=C_ENCODER, bold=True)

    # Residual arrow 1 (bypass)
    res_x = lx + box_w + Inches(0.2)
    add_line(slide, lx + box_w, ly + box_h / 2, res_x + Inches(0.8), ly + box_h / 2, C_FFN, Pt(1))
    add_line(slide, res_x + Inches(0.8), ly + box_h / 2, res_x + Inches(0.8), mha_y + Inches(1.4), C_FFN, Pt(1))
    add_textbox(slide, res_x + Inches(0.85), ly + Inches(0.5), Inches(0.8), Inches(0.3),
                "\u6b8b\u5dee\u2191", font_size=Pt(7), font_color=C_FFN, bold=True)

    # Arrow
    ay2 = mha_y + Inches(0.7)
    add_arrow_line(slide, lx + box_w / 2, ay2, lx + box_w / 2, ay2 + Inches(0.3), C_ENCODER, Pt(1.5))

    # Add & Norm 1
    an1_y = ay2 + Inches(0.3)
    add_rounded_rect(slide, lx, an1_y, box_w, Inches(0.45),
                     "Add & LayerNorm (512\u7ef4)",
                     fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)

    # Arrow
    ay3 = an1_y + Inches(0.45)
    add_arrow_line(slide, lx + box_w / 2, ay3, lx + box_w / 2, ay3 + Inches(0.3), C_ENCODER, Pt(1.5))

    # FFN
    ffn_y = ay3 + Inches(0.3)
    add_rounded_rect(slide, lx, ffn_y, box_w, Inches(0.7),
                     "Feed-Forward Network\n512->2048->ReLU->512",
                     fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(10), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)

    # Residual arrow 2
    add_line(slide, lx + box_w, an1_y + Inches(0.22), res_x + Inches(0.8), an1_y + Inches(0.22), C_FFN, Pt(1))
    add_line(slide, res_x + Inches(0.8), an1_y + Inches(0.22), res_x + Inches(0.8), ffn_y + Inches(1.05), C_FFN, Pt(1))
    add_textbox(slide, res_x + Inches(0.85), an1_y + Inches(0.3), Inches(0.8), Inches(0.3),
                "\u6b8b\u5dee\u2191", font_size=Pt(7), font_color=C_FFN, bold=True)

    # Arrow
    ay4 = ffn_y + Inches(0.7)
    add_arrow_line(slide, lx + box_w / 2, ay4, lx + box_w / 2, ay4 + Inches(0.3), C_ENCODER, Pt(1.5))

    # Add & Norm 2
    an2_y = ay4 + Inches(0.3)
    add_rounded_rect(slide, lx, an2_y, box_w, Inches(0.45),
                     "Add & LayerNorm (512\u7ef4)",
                     fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)

    # Output arrow
    add_arrow_line(slide, lx + box_w / 2, an2_y + Inches(0.45),
                   lx + box_w / 2, an2_y + Inches(0.7), C_ENCODER, Pt(1.5))
    add_textbox(slide, lx + Inches(0.3), an2_y + Inches(0.7), Inches(2), Inches(0.25),
                "v \u8f93\u51fa\u81f3\u4e0b\u4e00\u5c42", font_size=Pt(8), font_color=C_ENCODER, bold=True)

    # Dimension annotations on the right side of diagram
    dim_x = lx + box_w + Inches(1.5)
    for yy, label in [(ly, "512"), (mha_y + Inches(0.1), "8x64"), (an1_y, "512"),
                       (ffn_y + Inches(0.1), "2048"), (an2_y, "512")]:
        add_textbox(slide, dim_x, yy, Inches(0.8), Inches(0.25),
                    label, font_size=Pt(8), font_color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # Large encoder stack indicator
    add_rounded_rect(slide, lx - Inches(0.1), ly - Inches(0.15), box_w + Inches(0.2), an2_y + Inches(0.6) - ly + Inches(0.3),
                     "", fill_color=None, line_color=C_ENCODER, line_width=Pt(2))

    # Repeat x6 annotation
    add_textbox(slide, lx + Inches(0.5), an2_y + Inches(0.75), Inches(2), Inches(0.3),
                "\u2191 \u4e0a\u8ff0\u7ed3\u6784\u91cd\u590d 6 \u5c42", font_size=Pt(9), font_color=C_ENCODER, bold=True,
                alignment=PP_ALIGN.CENTER)

    # --- RIGHT: Key Points ---
    rx = Inches(7.5)
    ry = Inches(0.85)
    rw = Inches(5.5)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(6.0))
    points4 = [
        ("\u591a\u5934\u6ce8\u610f\u529b = 8\u4e2a\u8ba8\u8bba\u7ec4",
         "\u5355\u5934\u53ea\u80fd\u4ece\u4e00\u4e2a\u89d2\u5ea6\u5206\u6790\uff0c\u591a\u5934\u75288\u4e2a\u72ec\u7acb\u5934\u5b66\u4e60\u4e0d\u540c\u5173\u7cfb\u2014\u2014\u8bed\u6cd5(\u4e3b\u8c13)\u3001\u8bed\u4e49(\u540c\u4e49)\u3001\u4f4d\u7f6e(\u76f8\u90bb)\u3001\u6307\u4ee3(\u4ee3\u8bcd-\u540d\u8bcd)\u30028x64=512\u3002"),
        ("\u6b8b\u5dee\u8fde\u63a5 = \u9ad8\u901f\u516c\u8def",
         "output = x + SubLayer(x)\u3002\u68af\u5ea6\u4e2d\u6c38\u8fdc\u6709\u4e00\u6761\u201c1\u201d\u7684\u76f4\u901a\u8def\uff0c\u786e\u4fdd\u68af\u5ea6\u6c38\u4e0d\u6d88\u5931\u3002\u8fd9\u662fTransformer\u80fd\u5806\u53e06\u5c42\u7684\u5173\u952e\u3002"),
        ("\u5c42\u5f52\u4e00\u5316 = \u6807\u51c6\u5316\u4f53\u68c0",
         "y = gamma(x-mu)/sigma+beta\uff0c\u5bf9512\u7ef4\u5411\u91cf\u72ec\u7acb\u505a\u6807\u51c6\u5316\u3002\u9632\u6b62\u6570\u503c\u7a81\u53d8\u5bfc\u81f4\u68af\u5ea6\u7206\u70b8/\u6d88\u5931\u3002gamma\u548cbeta\u5404\u596e512\u7ef4\u3002"),
        ("FFN = \u72ec\u7acb\u601d\u8003",
         "512->2048->ReLU->512\uff0c\u53c2\u6570210\u4e07\uff0c\u6bd4\u6ce8\u610f\u529b\u59272\u500d\u3002\u6bcf\u4e2a\u4f4d\u7f6e\u72ec\u7acb\u5904\u7406\uff0c\u201c\u6d88\u5316\u201d\u6ce8\u610f\u529b\u6536\u96c6\u7684\u4fe1\u606f\u3002"),
    ]
    for title, desc in points4:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(6))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(4), line_spacing=1.1)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "3/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)
    return slide


# ============================================================
# PAGE 5: Decoder Detail
# ============================================================
def make_page5():
    slide = prs.slides.add_slide(blank_layout)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.45),
                "\u89e3\u7801\u5668\u5355\u5c42\uff1a\u591a\u4e00\u5ea7\u7ffb\u8bd1\u6865\u6881",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_DECODER)

    # --- LEFT: Decoder structure ---
    lx = Inches(0.3)
    ly = Inches(0.85)
    box_w = Inches(2.5)

    # Outer decoder frame
    add_rounded_rect(slide, lx - Inches(0.1), ly - Inches(0.15), Inches(5.5), Inches(3.5),
                     "", fill_color=None, line_color=C_DECODER, line_width=Pt(2))
    add_textbox(slide, lx, ly - Inches(0.12), Inches(3), Inches(0.25),
                "\u89e3\u7801\u5668\u5355\u5c42", font_size=Pt(10), font_color=C_DECODER, bold=True)

    # Masked Self-Attention
    msa_y = ly + Inches(0.2)
    add_rounded_rect(slide, lx, msa_y, box_w, Inches(0.55),
                     "\u63a9\u7801\u81ea\u6ce8\u610f\u529b",
                     fill_color=C_BG_RED, line_color=C_DECODER,
                     font_size=Pt(10), font_color=C_DECODER, bold=True)
    add_arrow_line(slide, lx + box_w / 2, msa_y + Inches(0.55),
                   lx + box_w / 2, msa_y + Inches(0.75), C_DECODER, Pt(1))

    # Add & Norm 1
    an1_y = msa_y + Inches(0.75)
    add_rounded_rect(slide, lx, an1_y, box_w, Inches(0.4),
                     "Add & LayerNorm", fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)
    add_arrow_line(slide, lx + box_w / 2, an1_y + Inches(0.4),
                   lx + box_w / 2, an1_y + Inches(0.6), C_DECODER, Pt(1))

    # Cross-Attention (highlighted!)
    ca_y = an1_y + Inches(0.6)
    add_rounded_rect(slide, lx, ca_y, box_w, Inches(0.55),
                     "\u4ea4\u53c9\u6ce8\u610f\u529b",
                     fill_color=C_BG_PURPLE, line_color=C_EMBED,
                     font_size=Pt(10), font_color=C_EMBED, bold=True)

    # K,V from encoder annotation
    add_arrow_line(slide, lx + box_w + Inches(0.1), ca_y + Inches(0.27),
                   lx + box_w, ca_y + Inches(0.27), C_EMBED, Pt(2))
    add_textbox(slide, lx + box_w + Inches(0.15), ca_y, Inches(2), Inches(0.25),
                "K, V <- \u7f16\u7801\u5668", font_size=Pt(9), font_color=C_EMBED, bold=True)

    add_arrow_line(slide, lx + box_w / 2, ca_y + Inches(0.55),
                   lx + box_w / 2, ca_y + Inches(0.75), C_DECODER, Pt(1))

    # Add & Norm 2
    an2_y = ca_y + Inches(0.75)
    add_rounded_rect(slide, lx, an2_y, box_w, Inches(0.4),
                     "Add & LayerNorm", fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)
    add_arrow_line(slide, lx + box_w / 2, an2_y + Inches(0.4),
                   lx + box_w / 2, an2_y + Inches(0.6), C_DECODER, Pt(1))

    # FFN
    ffn_y = an2_y + Inches(0.6)
    add_rounded_rect(slide, lx, ffn_y, box_w, Inches(0.45),
                     "FFN (512->2048->512)",
                     fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(9), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)

    # Difference highlight box
    add_rounded_rect(slide, lx - Inches(0.05), ca_y - Inches(0.05), box_w + Inches(0.1), Inches(0.65),
                     "", fill_color=None, line_color=C_EMBED, line_width=Pt(2))
    add_textbox(slide, lx + box_w + Inches(0.15), ca_y + Inches(0.3), Inches(2), Inches(0.3),
                "* \u89e3\u7801\u5668\u72ec\u6709", font_size=Pt(9), font_color=C_EMBED, bold=True)

    # --- Causal Mask Matrix (bottom-left) ---
    mask_x = Inches(0.3)
    mask_y = Inches(4.3)
    add_textbox(slide, mask_x, mask_y - Inches(0.25), Inches(3), Inches(0.25),
                "\u56e0\u679c\u63a9\u7801\u77e9\u9635\uff084x4\u793a\u4f8b\uff09",
                font_size=Pt(10), font_color=C_DECODER, bold=True)

    # Headers
    headers = ["[START]", "\u6211", "\u7231", "\u4f60"]
    for j, h in enumerate(headers):
        add_textbox(slide, mask_x + Inches(0.7) + Inches(0.6) * j, mask_y,
                    Inches(0.55), Inches(0.22), h, font_size=Pt(7), font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER, bold=True)

    matrix_vals = [
        [1, 0, 0, 0],
        [1, 1, 0, 0],
        [1, 1, 1, 0],
        [1, 1, 1, 1],
    ]
    for i, row in enumerate(matrix_vals):
        add_textbox(slide, mask_x, mask_y + Inches(0.22) + Inches(0.35) * i,
                    Inches(0.65), Inches(0.3), headers[i], font_size=Pt(7), font_color=C_TEXT,
                    alignment=PP_ALIGN.RIGHT, bold=True)
        for j, v in enumerate(row):
            cx = mask_x + Inches(0.7) + Inches(0.6) * j
            cy = mask_y + Inches(0.22) + Inches(0.35) * i
            if v == 1:
                add_rounded_rect(slide, cx, cy, Inches(0.5), Inches(0.28),
                                 text="V", fill_color=C_BG_GREEN, line_color=C_FFN,
                                 font_size=Pt(8), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)
            else:
                add_rounded_rect(slide, cx, cy, Inches(0.5), Inches(0.28),
                                 text="X", fill_color=C_BG_RED, line_color=C_DECODER,
                                 font_size=Pt(8), font_color=C_DECODER)

    add_textbox(slide, mask_x, mask_y + Inches(1.75), Inches(3.5), Inches(0.25),
                "\u4e0b\u4e09\u89d2=\u53ef\u89c1\uff0c\u4e0a\u4e09\u89d2=\u5c4f\u853d(\u8bbe-inf)",
                font_size=Pt(8), font_color=C_GRAY)

    # --- RIGHT: Key Points ---
    rx = Inches(7.5)
    ry = Inches(0.85)
    rw = Inches(5.5)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(6.0))
    points5 = [
        ("\u63a9\u7801\u81ea\u6ce8\u610f\u529b",
         "\u4e0a\u4e09\u89d2\u8bbe-inf -> Softmax\u540e\u4e3a0\u3002\u5904\u7406\u201c\u7231\u201d\u65f6\u53ea\u80fd\u770b[START, \u6211, \u7231]\uff0c\u4e0d\u80fd\u5077\u770b\u201c\u4f60\u201d\u3002\u786e\u4fdd\u81ea\u56de\u5f52\u7279\u6027\uff0c\u9632\u6b62\u201c\u6cc4\u9732\u7b54\u6848\u201d\u3002"),
        ("\u4ea4\u53c9\u6ce8\u610f\u529b",
         "Q\u6765\u81ea\u89e3\u7801\u5668\uff08\u201c\u6211\u8981\u627e\u4ec0\u4e48\u201d\uff09\uff0cK/V\u6765\u81ea\u7f16\u7801\u5668\uff08\u201c\u6e90\u8bed\u8a00\u6709\u4ec0\u4e48\u201d\uff09\u3002\u5b9e\u73b0\u6e90\u8bed\u8a00->\u76ee\u6807\u8bed\u8a00\u7684\u4fe1\u606f\u4f20\u9012\uff0c\u8fd9\u662f\u7ffb\u8bd1\u7684\u6838\u5fc3\u3002"),
        ("\u4e0e\u7f16\u7801\u5668\u7684\u533a\u522b",
         "(1) \u591a\u4e00\u4e2aCross-Attention\u5b50\u5c42\uff0c\u63a5\u6536\u7f16\u7801\u5668\u7684K,V\uff1b(2) \u81ea\u6ce8\u610f\u529b\u7528\u56e0\u679c\u63a9\u7801\u800c\u975e\u666e\u901a\u63a9\u7801\uff1b(3) \u8f93\u51fa\u662f\u4e0b\u4e00\u4e2a\u8bcd\u7684\u6982\u7387\u5206\u5e03\u3002"),
    ]
    for title, desc in points5:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(8))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(4), line_spacing=1.1)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "4/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)
    return slide


# ============================================================
# PAGE 6: Training Phase
# ============================================================
def make_page6():
    slide = prs.slides.add_slide(blank_layout)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.45),
                "\u8bad\u7ec3\u9636\u6bb5\uff1a\u5f00\u5377\u8003\u8bd5",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_ENCODER)

    # --- TOP: Training data flow ---
    # Input side
    add_textbox(slide, Inches(0.3), Inches(0.75), Inches(2), Inches(0.25),
                "\u8f93\u5165\u4fa7\uff08\u6e90\u8bed\u8a00\uff09",
                font_size=Pt(9), font_color=C_ENCODER, bold=True)
    add_rounded_rect(slide, Inches(0.3), Inches(1.0), Inches(1.2), Inches(0.4),
                     "I love NLP", fill_color=C_BG_BLUE, line_color=C_ENCODER,
                     font_size=Pt(8), font_color=C_ENCODER, bold=True)
    add_arrow_line(slide, Inches(1.5), Inches(1.2), Inches(1.8), Inches(1.2), C_ENCODER, Pt(1))
    add_rounded_rect(slide, Inches(1.8), Inches(1.0), Inches(1.8), Inches(0.4),
                     "\u7f16\u7801\u5668 x 6", fill_color=C_BG_BLUE, line_color=C_ENCODER,
                     font_size=Pt(9), font_color=C_ENCODER, bold=True)
    add_arrow_line(slide, Inches(3.6), Inches(1.2), Inches(4.0), Inches(1.2), C_EMBED, Pt(1.5))
    add_rounded_rect(slide, Inches(4.0), Inches(1.0), Inches(0.7), Inches(0.4),
                     "K, V", fill_color=C_BG_PURPLE, line_color=C_EMBED,
                     font_size=Pt(8), font_color=C_EMBED, bold=True)

    # Decoder side
    add_textbox(slide, Inches(0.3), Inches(1.55), Inches(2), Inches(0.25),
                "\u89e3\u7801\u4fa7\uff08\u76ee\u6807\u8bed\u8a00\uff09",
                font_size=Pt(9), font_color=C_DECODER, bold=True)
    add_rounded_rect(slide, Inches(0.3), Inches(1.8), Inches(2.0), Inches(0.4),
                     "[START, \u6211, \u7231, NLP]", fill_color=C_BG_RED, line_color=C_DECODER,
                     font_size=Pt(8), font_color=C_DECODER, bold=True)
    add_arrow_line(slide, Inches(2.3), Inches(2.0), Inches(2.6), Inches(2.0), C_DECODER, Pt(1))
    add_rounded_rect(slide, Inches(2.6), Inches(1.8), Inches(1.8), Inches(0.4),
                     "\u89e3\u7801\u5668 x 6", fill_color=C_BG_RED, line_color=C_DECODER,
                     font_size=Pt(9), font_color=C_DECODER, bold=True)
    add_arrow_line(slide, Inches(4.4), Inches(2.0), Inches(4.8), Inches(2.0), C_DECODER, Pt(1))
    add_rounded_rect(slide, Inches(4.8), Inches(1.8), Inches(0.8), Inches(0.4),
                     "\u9884\u6d4b", fill_color=C_BG_LIGHT, line_color=C_TEXT,
                     font_size=Pt(8), font_color=C_TEXT, bold=True)

    # Target + Loss
    add_textbox(slide, Inches(5.8), Inches(1.85), Inches(0.5), Inches(0.25),
                "<->", font_size=Pt(14), font_color=C_DECODER, bold=True)
    add_rounded_rect(slide, Inches(6.2), Inches(1.8), Inches(1.8), Inches(0.4),
                     "\u76ee\u6807: \u6211\u7231 NLP [END]", fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(8), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)
    add_arrow_line(slide, Inches(8.0), Inches(2.0), Inches(8.3), Inches(2.0), C_DECODER, Pt(1))
    add_rounded_rect(slide, Inches(8.3), Inches(1.8), Inches(1.0), Inches(0.4),
                     "Loss", fill_color=C_BG_RED, line_color=C_DECODER,
                     font_size=Pt(10), font_color=C_DECODER, bold=True)

    # Parallel annotation
    add_rounded_rect(slide, Inches(0.3), Inches(2.35), Inches(9.0), Inches(0.35),
                     "V \u6240\u6709\u4f4d\u7f6e\u5e76\u884c\u8ba1\u7b97\uff0cGPU\u5145\u5206\u5229\u7528",
                     fill_color=C_BG_GREEN, line_color=C_FFN,
                     font_size=Pt(9), font_color=RGBColor(0x27, 0xAE, 0x60), bold=True)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(2.85), Inches(12.6), Pt(0.5),
              fill_color=C_BG_LIGHT)

    # --- BOTTOM LEFT: Embedding + PosEnc flow ---
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(3), Inches(0.25),
                "\u5d4c\u5165 + \u4f4d\u7f6e\u7f16\u7801\u6d41\u7a0b",
                font_size=Pt(10), font_color=C_EMBED, bold=True)

    flow_items = [
        ("\u8bcd\u6c47\u8868", C_BG_LIGHT, C_TEXT, Inches(3.3)),
        ("-> Embedding Matrix", C_BG_PURPLE, C_EMBED, Inches(3.8)),
        ("-> + \u4f4d\u7f6e\u7f16\u7801", C_BG_YELLOW, C_ATTN, Inches(4.3)),
        ("-> \u8f93\u5165\u5411\u91cf(512\u7ef4)", C_BG_GREEN, RGBColor(0x27, 0xAE, 0x60), Inches(4.8)),
    ]
    for label, bg, fc, fy in flow_items:
        add_rounded_rect(slide, Inches(0.3), fy, Inches(2.5), Inches(0.38),
                         label, fill_color=bg, line_color=fc,
                         font_size=Pt(8), font_color=fc, bold=True)

    # Embedding matrix visual
    add_textbox(slide, Inches(3.0), Inches(3.1), Inches(1.5), Inches(0.25),
                "37000 x 512", font_size=Pt(9), font_color=C_EMBED, bold=True)
    # Mini matrix grid
    for r in range(4):
        for c in range(4):
            add_rounded_rect(slide, Inches(3.0) + Inches(0.35) * c, Inches(3.4) + Inches(0.3) * r,
                             Inches(0.3), Inches(0.25),
                             "", fill_color=C_BG_PURPLE, line_color=C_EMBED, line_width=Pt(0.5))

    # Pos encoding formula
    add_rounded_rect(slide, Inches(0.3), Inches(5.4), Inches(4.5), Inches(0.8),
                     "", fill_color=C_BG_YELLOW, line_color=C_ATTN)
    _, ptf = add_rich_textbox(slide, Inches(0.4), Inches(5.45), Inches(4.3), Inches(0.7))
    add_paragraph(ptf, "PE(pos,2i) = sin(pos/10000^(2i/d))", font_size=Pt(9),
                  font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True, font_name=FONT_EN)
    add_paragraph(ptf, "PE(pos,2i+1) = cos(pos/10000^(2i/d))", font_size=Pt(9),
                  font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True, font_name=FONT_EN)
    add_paragraph(ptf, "V \u4e0d\u589e\u52a0\u53ef\u8bad\u7ec3\u53c2\u6570\uff0c\u53ef\u63a8\u7406\u672a\u89c1\u8fc7\u7684\u957f\u5ea6", font_size=Pt(8),
                  font_color=C_TEXT)

    # --- BOTTOM RIGHT: Key Points ---
    rx = Inches(7.5)
    ry = Inches(3.0)
    rw = Inches(5.5)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(3.5))
    points6 = [
        ("Teacher Forcing",
         "\u89e3\u7801\u5668\u8f93\u5165\u662f\u771f\u5b9e\u7b54\u6848\uff08\u800c\u975e\u81ea\u5df1\u7684\u9884\u6d4b\uff09\uff0c\u6240\u6709\u4f4d\u7f6e\u5e76\u884c\u8ba1\u7b97\uff0cGPU\u5145\u5206\u5229\u7528\u3002"),
        ("\u5d4c\u5165\u77e9\u9635",
         "37000x512\uff0c1890\u4e07\u53c2\u6570\uff0c\u5360\u6a21\u578b30%\u3002\u6bcf\u4e2a\u8bcd\u67e5\u8868\u5f97\u5230\u72ec\u7acb\u7684512\u7ef4\u5411\u91cf\u3002"),
        ("\u4f4d\u7f6e\u7f16\u7801",
         "PE=sin/cos\u51fd\u6570\uff0c\u6ce8\u5165\u4f4d\u7f6e\u4fe1\u606f\u3002\u4e0d\u589e\u52a0\u53ef\u8bad\u7ec3\u53c2\u6570\uff0c\u652f\u6301\u63a8\u7406\u65f6\u7684\u672a\u89c1\u957f\u5ea6\u3002"),
        ("\u8bad\u7ec3\u76ee\u6807",
         "\u9884\u6d4b\u4e0b\u4e00\u4e2a\u8bcd\uff0cLoss=-log(p(\u6b63\u786e\u8bcd))\u3002\u521d\u59cb\u224810.5\uff08\u968f\u673a\uff09\uff0c\u8bad\u7ec3\u540e1.5-3.0\u3002"),
    ]
    for title, desc in points6:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(5))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(3), line_spacing=1.1)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "5/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)
    return slide


# ============================================================
# PAGE 7: Self-Attention 5 Steps
# ============================================================
def make_page7():
    slide = prs.slides.add_slide(blank_layout)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(10), Inches(0.45),
                "\u81ea\u6ce8\u610f\u529b\u76845\u6b65\u8ba1\u7b97",
                font_size=Pt(20), font_color=C_TITLE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.6), Inches(4.5), Pt(2.5),
              fill_color=C_ATTN)

    # --- LEFT: 5-step vertical flow ---
    lx = Inches(0.4)
    step_w = Inches(3.2)
    step_h = Inches(0.85)

    steps = [
        ("(1) \u751f\u6210 QKV", "X x W_Q / W_K / W_V", C_BG_BLUE, C_ENCODER),
        ("(3) \u5206\u6570\u77e9\u9635", "Q x K^T", C_BG_YELLOW, C_ATTN),
        ("(4) \u7f29\u653e", "/ sqrtd_k (sqrt64 = 8)", C_BG_LIGHT, C_TEXT),
        ("(5) Softmax", "\u8f6c\u4e3a\u6982\u7387\u5206\u5e03", C_BG_GREEN, C_FFN),
        ("(6) \u52a0\u6743\u6c42\u548c", "\u6743\u91cd x V = \u8f93\u51fa", C_BG_PURPLE, C_EMBED),
    ]

    for i, (title, formula, bg, fc) in enumerate(steps):
        y = Inches(0.8) + Inches(1.2) * i
        add_rounded_rect(slide, lx, y, step_w, step_h,
                         "", fill_color=bg, line_color=fc, line_width=Pt(1.5))
        add_textbox(slide, lx + Inches(0.1), y + Inches(0.05), step_w - Inches(0.2), Inches(0.3),
                    title, font_size=Pt(10), font_color=fc, bold=True)
        add_textbox(slide, lx + Inches(0.1), y + Inches(0.35), step_w - Inches(0.2), Inches(0.4),
                    formula, font_size=Pt(12), font_color=C_TEXT, bold=True,
                    font_name=FONT_EN, alignment=PP_ALIGN.CENTER)

        # Arrow to next step
        if i < len(steps) - 1:
            add_arrow_line(slide, lx + step_w / 2, y + step_h,
                           lx + step_w / 2, y + step_h + Inches(0.35), fc, Pt(1.5))

    # --- RIGHT: Key Points ---
    rx = Inches(4.5)
    ry = Inches(0.8)
    rw = Inches(8.5)
    _, tf = add_rich_textbox(slide, rx, ry, rw, Inches(6.0))

    points7 = [
        ("(1) \u751f\u6210QKV",
         "X\u4e58\u4e09\u7ec4\u6743\u91cd\u77e9\u9635\u5f97\u5230Q\u3001K\u3001V\u3002Q=\u201c\u6211\u8981\u627e\u4ec0\u4e48\u201d\uff0cK=\u201c\u6211\u6709\u4ec0\u4e48\u7279\u5f81\u201d\uff0cV=\u201c\u6211\u7684\u5b9e\u9645\u5185\u5bb9\u201d\u3002\u4e09\u4e2a\u77e9\u9635\u5747\u4e3adxd\u3002"),
        ("(3) \u5206\u6570\u77e9\u9635",
         "QxK^T\u5f97\u5230TxT\u7684\u6ce8\u610f\u529b\u5206\u6570\u77e9\u9635\u3002Score[i,j]=\u8bcdi\u5bf9\u8bcdj\u7684\u5173\u6ce8\u5ea6\uff0c\u503c\u8d8a\u5927\u8d8a\u91cd\u8981\u3002"),
        ("(4) \u7f29\u653e",
         "/sqrt64=8\uff0c\u9632\u6b62\u7ef4\u5ea6\u8f83\u9ad8\u65f6Softmax\u9971\u548c\uff08\u8f93\u51fa\u63a5\u8fd10\u62161\uff0c\u68af\u5ea6\u6d88\u5931\uff09\u3002"),
        ("(5) Softmax",
         "\u8f6c\u4e3a\u6982\u7387\u5206\u5e03\uff0c\u6bcf\u884c\u548c=1\u3002\u653e\u5927\u6700\u5927\u503c\u4f46\u4e0d\u7f6e\u96f6\u5176\u4ed6\u503c\uff0c\u4fdd\u7559\u6240\u6709\u8bcd\u7684\u4fe1\u606f\u3002"),
        ("(6) \u52a0\u6743\u6c42\u548c",
         "Output=\u6743\u91cd\u77e9\u9635xV\u3002\u201cit\u201d\u7684\u8f93\u51fa\u878d\u5408\u4e860.6xV(cat)+0.1xV(sat)+...\uff0c\u542b\u6709\u201ccat\u201d\u7684\u4fe1\u606f\uff0c\u6a21\u578b\u201c\u77e5\u9053\u201dit\u6307\u4ee3\u7684\u662fcat\u3002"),
    ]
    for title, desc in points7:
        add_paragraph(tf, title, font_size=Pt(11), font_color=C_TITLE, bold=True,
                      space_after=Pt(1), space_before=Pt(5))
        add_paragraph(tf, desc, font_size=Pt(10), font_color=C_TEXT,
                      space_after=Pt(3), line_spacing=1.1)

    # Bottom formula summary
    add_rounded_rect(slide, Inches(0.3), Inches(6.6), Inches(12.6), Inches(0.4),
                     "Attention(Q,K,V) = softmax(QK^T / sqrtd_k) V",
                     fill_color=C_BG_YELLOW, line_color=C_ATTN,
                     font_size=Pt(11), font_color=RGBColor(0xE6, 0x7E, 0x22), bold=True)

    # Page number
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                "6/7", font_size=Pt(9), font_color=C_GRAY, alignment=PP_ALIGN.RIGHT)
    return slide


# ============================================================
# BUILD
# ============================================================
if __name__ == "__main__":
    slides = []
    slides.append(make_page1())
    slides.append(make_page2())
    slides.append(make_page3())
    slides.append(make_page4())
    slides.append(make_page5())
    slides.append(make_page6())
    slides.append(make_page7())

    out_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v5_part1.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")
    print()
    print("=" * 50)
    print("Per-slide text count (target: 150-300 chars)")
    print("=" * 50)
    total = 0
    for i, sl in enumerate(slides, 1):
        cnt = count_text(sl)
        total += cnt
        status = "\u2705" if 50 <= cnt <= 400 else "[!]"
        print(f"  Page {i}: {cnt:4d} chars {status}")
    print(f"  Total: {total} chars")
