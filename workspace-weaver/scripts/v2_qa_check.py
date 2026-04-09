#!/usr/bin/env python3
"""
v2 PPT 质检脚本 - Transformer_架构深度解析_v2.pptx
严格检查：文字重叠、边界、字体大小、内容完整性
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json, sys

pptx_path = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v2.pptx"
prs = Presentation(pptx_path)

SLIDE_W = prs.slide_width  # EMU
SLIDE_H = prs.slide_height  # EMU
SLIDE_W_IN = SLIDE_W / 914400
SLIDE_H_IN = SLIDE_H / 914400

print(f"Slide dimensions: {SLIDE_W_IN:.3f}\" x {SLIDE_H_IN:.3f}\"")
print(f"Total slides: {len(prs.slides)}")
print("="*80)

# Container patterns to exclude (design intent)
# We'll classify shapes first, then check overlaps

def emu_to_inch(emu):
    return emu / 914400.0

def get_bbox(shape):
    """Get bounding box in inches: (left, top, right, bottom)"""
    l = emu_to_inch(shape.left)
    t = emu_to_inch(shape.top)
    r = l + emu_to_inch(shape.width)
    b = t + emu_to_inch(shape.height)
    return (l, t, r, b)

def bbox_area(bbox):
    l, t, r, b = bbox
    return max(0, (r - l)) * max(0, (b - t))

def overlap_area(a, b):
    l = max(a[0], b[0])
    t = max(a[1], b[1])
    r = min(a[2], b[2])
    b2 = min(a[3], b[3])
    if l >= r or t >= b2:
        return 0
    return (r - l) * (b2 - t)

def overlap_ratio(a, b):
    """How much of shape A is covered by shape B"""
    area_a = bbox_area(a)
    if area_a < 0.001:
        return 0
    return overlap_area(a, b) / area_a

def get_text_content(shape):
    """Extract text from shape"""
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""

def shape_type_name(shape):
    st = shape.shape_type
    type_map = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
        MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
        MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
        MSO_SHAPE_TYPE.FREEFORM: "Freeform",
        MSO_SHAPE_TYPE.GROUP: "Group",
        MSO_SHAPE_TYPE.PICTURE: "Picture",
        MSO_SHAPE_TYPE.TABLE: "Table",
        MSO_SHAPE_TYPE.CHART: "Chart",
    }
    return type_map.get(st, f"Unknown({st})")

def get_font_sizes(shape):
    """Get all font sizes in a shape"""
    sizes = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(run.font.size.pt)
    return sizes

# ============= ANALYSIS =============

all_issues = []
all_overlaps = []
font_issues = []
boundary_issues = []
content_issues = []

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    shapes = list(slide.shapes)
    n_shapes = len(shapes)
    
    print(f"\n{'='*60}")
    print(f"Slide {slide_num}: {n_shapes} shapes")
    print(f"{'='*60}")
    
    # --- Shape inventory ---
    for si, shape in enumerate(shapes):
        bbox = get_bbox(shape)
        text = get_text_content(shape)
        text_preview = text[:50] + "..." if len(text) > 50 else text
        print(f"  [{si:2d}] {shape_type_name(shape):12s} | ({bbox[0]:.2f},{bbox[1]:.2f})-({bbox[2]:.2f},{bbox[3]:.2f}) | {shape.width/914400:.2f}x{shape.height/914400:.2f}\" | \"{text_preview}\"")
    
    # --- 1. OVERLAP DETECTION ---
    # Classify shapes
    text_shapes = []  # shapes with meaningful text
    container_shapes = []  # large background/container shapes
    connector_shapes = []  # connectors
    other_shapes = []  # decorative, etc.
    
    for si, shape in enumerate(shapes):
        bbox = get_bbox(shape)
        text = get_text_content(shape)
        area = bbox_area(bbox)
        
        # Classify
        is_connector = shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
        
        # Large empty/near-empty shapes as containers (backgrounds, panels)
        is_container = False
        if area > 5.0 and len(text) < 3:
            is_container = True
        # Very large shapes even with some text
        if area > 10.0:
            is_container = True
            
        if is_connector:
            connector_shapes.append((si, shape))
        elif is_container:
            container_shapes.append((si, shape))
        elif text:
            text_shapes.append((si, shape))
        else:
            other_shapes.append((si, shape))
    
    # Check all shape pairs for overlap
    real_overlaps = []
    for i in range(n_shapes):
        for j in range(i+1, n_shapes):
            si, sj = shapes[i], shapes[j]
            bbox_i = get_bbox(si)
            bbox_j = get_bbox(sj)
            
            area_i = bbox_area(bbox_i)
            area_j = bbox_area(bbox_j)
            
            if area_i < 0.01 or area_j < 0.01:
                continue
            
            oa = overlap_area(bbox_i, bbox_j)
            if oa < 0.01:  # Less than 0.01 sq inch overlap
                continue
            
            ratio_i = oa / area_i * 100
            ratio_j = oa / area_j * 100
            
            text_i = get_text_content(si)
            text_j = get_text_content(sj)
            
            # Skip container-child patterns
            # If one shape is much larger and contains the other
            if ratio_i > 90 and ratio_j < 10:
                continue  # i contains j
            if ratio_j > 90 and ratio_i < 10:
                continue  # j contains i
            
            # Skip connector-line overlaps (design intent in attention matrices)
            if si.shape_type == MSO_SHAPE_TYPE.FREEFORM or sj.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                continue
            
            # Skip if both are empty/decorative
            if not text_i and not text_j:
                continue
            
            # For text on container: if the container is very large and the text is small, skip
            if area_i > 5.0 and area_j < 2.0 and not text_i:
                continue
            if area_j > 5.0 and area_i < 2.0 and not text_j:
                continue
            
            overlap_entry = {
                "slide": slide_num,
                "shape_a": f"[{i}] {shape_type_name(si)} \"{text_i[:30]}\"",
                "shape_b": f"[{j}] {shape_type_name(sj)} \"{text_j[:30]}\"",
                "bbox_a": f"({bbox_i[0]:.2f},{bbox_i[1]:.2f})-({bbox_i[2]:.2f},{bbox_i[3]:.2f})",
                "bbox_b": f"({bbox_j[0]:.2f},{bbox_j[1]:.2f})-({bbox_j[2]:.2f},{bbox_j[3]:.2f})",
                "overlap_area": oa,
                "ratio_a": ratio_i,
                "ratio_b": ratio_j,
            }
            
            # Only flag if significant overlap on the smaller shape
            min_ratio = min(ratio_i, ratio_j)
            
            if min_ratio >= 30:
                severity = "P0"
            elif min_ratio >= 15:
                severity = "P1"
            elif min_ratio >= 5:
                severity = "P2"
            else:
                severity = "P3"
            
            overlap_entry["severity"] = severity
            real_overlaps.append(overlap_entry)
            
            print(f"  ⚠️  OVERLAP [{severity}]: {overlap_entry['shape_a']}")
            print(f"       vs {overlap_entry['shape_b']}")
            print(f"       Area: {oa:.3f}sq\" | ratio A: {ratio_i:.1f}% | ratio B: {ratio_j:.1f}%")
    
    all_overlaps.extend(real_overlaps)
    
    # --- 2. BOUNDARY CHECK ---
    for si, shape in enumerate(shapes):
        bbox = get_bbox(shape)
        if bbox[2] > SLIDE_W_IN + 0.01:
            issue = f"Slide {slide_num} [{si}] extends right: {bbox[2]:.2f}\" > {SLIDE_W_IN:.2f}\""
            boundary_issues.append(issue)
            print(f"  🔴 BOUNDARY: {issue}")
        if bbox[3] > SLIDE_H_IN + 0.01:
            issue = f"Slide {slide_num} [{si}] extends bottom: {bbox[3]:.2f}\" > {SLIDE_H_IN:.2f}\""
            boundary_issues.append(issue)
            print(f"  🔴 BOUNDARY: {issue}")
        if bbox[0] < -0.01:
            issue = f"Slide {slide_num} [{si}] extends left: {bbox[0]:.2f}\""
            boundary_issues.append(issue)
            print(f"  🔴 BOUNDARY: {issue}")
        if bbox[1] < -0.01:
            issue = f"Slide {slide_num} [{si}] extends top: {bbox[1]:.2f}\""
            boundary_issues.append(issue)
            print(f"  🔴 BOUNDARY: {issue}")
    
    # --- 3. FONT SIZE CHECK ---
    for si, shape in enumerate(shapes):
        if not shape.has_text_frame:
            continue
        text = get_text_content(shape)
        if not text:
            continue
        sizes = get_font_sizes(shape)
        if not sizes:
            # Check paragraph default
            for para in shape.text_frame.paragraphs:
                if para.font.size:
                    sizes.append(para.font.size.pt)
        
        if sizes:
            min_size = min(sizes)
            bbox = get_bbox(shape)
            area = bbox_area(bbox)
            
            # For small shapes (labels in diagrams), require >= 10pt
            # For larger text blocks, require >= 9pt
            threshold = 10 if area < 2.0 else 9
            
            if min_size < threshold:
                text_preview = text[:30]
                issue = f"Slide {slide_num} [{si}] font {min_size:.1f}pt < {threshold}pt threshold (area={area:.2f}) \"{text_preview}\""
                font_issues.append(issue)
                print(f"  🔴 FONT: {issue}")
        else:
            # No explicit font size set - check if default is ok
            # python-pptx can't always get default, flag as warning
            pass
    
    # --- 4. CONTENT CHECK ---
    # Check for empty slides
    has_content = False
    for shape in shapes:
        if get_text_content(shape):
            has_content = True
            break
    if not has_content:
        content_issues.append(f"Slide {slide_num}: EMPTY SLIDE - no text content found")
        print(f"  🔴 CONTENT: Empty slide!")
    
    # Check for very short text fragments that look like leftover template text
    for si, shape in enumerate(shapes):
        text = get_text_content(shape)
        if not text:
            continue
        # Skip single characters, common diagram labels
        if len(text) <= 2 and text not in ['→', '←', '×', '+', '=', '∑', '√']:
            pass  # Likely diagram labels, acceptable

# ============= SUMMARY =============
print(f"\n{'='*80}")
print("SUMMARY")
print(f"{'='*80}")

p0_overlaps = [o for o in all_overlaps if o["severity"] == "P0"]
p1_overlaps = [o for o in all_overlaps if o["severity"] == "P1"]
p2_overlaps = [o for o in all_overlaps if o["severity"] == "P2"]
p3_overlaps = [o for o in all_overlaps if o["severity"] == "P3"]

print(f"\nOverlaps found: {len(all_overlaps)} total")
print(f"  P0 (critical, >=30%): {len(p0_overlaps)}")
print(f"  P1 (serious, >=15%): {len(p1_overlaps)}")
print(f"  P2 (minor, >=5%): {len(p2_overlaps)}")
print(f"  P3 (negligible, <5%): {len(p3_overlaps)}")

print(f"\nBoundary issues: {len(boundary_issues)}")
print(f"Font size issues: {len(font_issues)}")
print(f"Content issues: {len(content_issues)}")

# Scoring
overlap_score = 30
if p0_overlaps:
    overlap_score -= len(p0_overlaps) * 4
if p1_overlaps:
    overlap_score -= len(p1_overlaps) * 2
if p2_overlaps:
    overlap_score -= len(p2_overlaps) * 0.5
if p3_overlaps:
    overlap_score -= len(p3_overlaps) * 0.1
overlap_score = max(0, overlap_score)

boundary_score = 20 - len(boundary_issues) * 5
boundary_score = max(0, boundary_score)

font_score = 20 - len(font_issues) * 2
font_score = max(0, font_score)

content_score = 30 - len(content_issues) * 10
content_score = max(0, content_score)

total = overlap_score + boundary_score + font_score + content_score

print(f"\n{'='*80}")
print(f"SCORING")
print(f"{'='*80}")
print(f"Overlap:   {overlap_score:.1f} / 30")
print(f"Boundary:  {boundary_score:.1f} / 20")
print(f"Font:      {font_score:.1f} / 20")
print(f"Content:   {content_score:.1f} / 30")
print(f"TOTAL:     {total:.1f} / 100")

# Save detailed data
result = {
    "total_slides": len(prs.slides),
    "slide_dimensions": f"{SLIDE_W_IN:.3f}\" x {SLIDE_H_IN:.3f}\"",
    "overlaps": all_overlaps,
    "boundary_issues": boundary_issues,
    "font_issues": font_issues,
    "content_issues": content_issues,
    "scores": {
        "overlap": overlap_score,
        "boundary": boundary_score,
        "font": font_score,
        "content": content_score,
        "total": total,
    },
    "p0_count": len(p0_overlaps),
    "p1_count": len(p1_overlaps),
    "p2_count": len(p2_overlaps),
    "p3_count": len(p3_overlaps),
}

with open("/home/admin/.openclaw/workspace-weaver/memory/v2_qa_data.json", "w") as f:
    json.dump(result, f, indent=2, ensure_ascii=False)

print(f"\nDetailed data saved to v2_qa_data.json")
