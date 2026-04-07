#!/usr/bin/env python3
"""
SD PPT Fix v10: ALL 15 slides generated from scratch — zero copy_slide / deepcopy.
Fixes v9's rendering issues caused by XML deepcopy of 8 slides from v8.

Pages rebuilt from scratch (content preserved from v8, layout re-done):
  P2 目录, P3 Diffusion基础, P8 协作流程, P10 推理过程, P12-P15

Pages preserved from v9 (already from-scratch):
  P1 封面, P4 整体架构, P5 CLIP, P6 VAE, P7 U-Net, P9 训练, P11 CFG

Output: output/Stable_Diffusion_技术解析_v10.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = "/home/admin/.openclaw/workspace-weaver/output/Stable_Diffusion_技术解析_v10.pptx"

# Colors
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_LT  = RGBColor(0xEB,0xF5,0xFB)
C_W   = RGBColor(0xFF,0xFF,0xFF)
C_D   = RGBColor(0x33,0x33,0x33)
C_G   = RGBColor(0x88,0x88,0x88)
C_GRN = RGBColor(0x27,0xAE,0x60)
C_ORG = RGBColor(0xE6,0x7E,0x22)
C_RED = RGBColor(0xE7,0x4C,0x3C)
C_PUR = RGBColor(0x8E,0x44,0xAD)
C_TEL = RGBColor(0x00,0x96,0x88)
C_BG_GRN = RGBColor(0xE8,0xF8,0xF5)
C_BG_ORG = RGBColor(0xFD,0xF2,0xE9)
C_BG_RED = RGBColor(0xFD,0xED,0xED)
C_YEL = RGBColor(0xF3,0x9C,0x12)

FC = "微软雅黑"
FE = "Arial"

# ── Helpers ──
def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb=lc
        if lw: sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc:
        sh.line.fill.solid(); sh.line.color.rgb=lc
        if lw: sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    for r in p.runs:
        r.font.size=Pt(fs); r.font.color.rgb=fc; r.font.bold=b; r.font.name=fn
    p.font.size=Pt(fs); p.font.color.rgb=fc; p.font.bold=b; p.font.name=fn
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
        for r in p.runs:
            r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=bld; r.font.name=fn
        p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld; p.font.name=fn
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_u(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.UP_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def line_h(s,l,t,w,c=C_ACC,lw=2):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Emu(int(lw*12700)))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Stable Diffusion 技术解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    violations = []
    for sh in s.shapes:
        if sh.left < Emu(-100):
            violations.append(f"  NEGATIVE X: left={sh.left/914400:.3f}")
        if sh.top < Emu(-100):
            violations.append(f"  NEGATIVE Y: top={sh.top/914400:.3f}")
        r = sh.left + sh.width
        b = sh.top + sh.height
        if r > SLIDE_W + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW X: right={r/914400:.3f} text={txt}")
        if b > SLIDE_H + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW Y: bottom={b/914400:.3f} text={txt}")
    if violations:
        print(f"  {msg} — {len(violations)} violations:")
        for v in violations: print(v)
        return False
    else:
        print(f"  OK {msg} — all shapes within bounds")
        return True


# ============================================================
# Create new presentation
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ============================================================
# P1 封面 (from v9)
# ============================================================
print("\nBuilding P1: 封面...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_PRI)
rect(sl, Inches(0.8), Inches(2.8), Inches(3.0), Inches(0.06), fc=C_ACC)
tb(sl, Inches(0.8), Inches(3.0), Inches(11.0), Inches(1.20),
   "Stable Diffusion", fs=44, fc=C_W, b=True, fn="Arial", ls=1.0)
tb(sl, Inches(0.8), Inches(4.0), Inches(11.0), Inches(0.60),
   "技术解析 — 从原理到前沿", fs=24, fc=RGBColor(0xBB,0xDE,0xFB), b=False, fn=FC, ls=1.0)
rect(sl, Inches(0.8), Inches(4.6), Inches(1.5), Inches(0.04), fc=C_ACC)
tb(sl, Inches(0.8), Inches(5.0), Inches(8.0), Inches(0.40),
   "深入理解扩散模型、CLIP、VAE、U-Net 及最新进展", fs=13, fc=RGBColor(0x99,0xBB,0xDD), fn=FC)
tb(sl, Inches(0.8), Inches(5.5), Inches(8.0), Inches(0.40),
   "2026 年 4 月  |  19 章节完整版", fs=12, fc=RGBColor(0x88,0xAA,0xCC), fn=FC)
rect(sl, Inches(10.0), Inches(0.0), Inches(3.33), Inches(7.50), fc=RGBColor(0x17,0x2E,0x4A))
tb(sl, Inches(10.5), Inches(5.5), Inches(2.5), Inches(1.50),
   "AI IMAGE\nGENERATION", fs=18, fc=RGBColor(0x34,0x98,0xDB), b=True, fn="Arial", a=PP_ALIGN.CENTER, ls=1.3)
chk(sl, "P1 封面")


# ============================================================
# P2 目录 (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P2: 目录 (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "目录", "CONTENTS — 19 CHAPTERS")

toc_items = [
    ("01", "Diffusion Model 基础原理", "墨水扩散与逆向去噪"),
    ("02", "Stable Diffusion vs 普通 Diffusion", "48倍压缩的核心创新"),
    ("03", "整体架构概览", "CLIP + U-Net + VAE 三大组件"),
    ("04", "CLIP 文本编码器", "把文字翻译成AI能理解的数字"),
    ("05", "VAE 编解码器", "超级智能的ZIP压缩软件"),
    ("06", "U-Net 去噪网络", "核心画师，8.6亿参数"),
    ("07", "三大组件协作流程", "AI画图流水线"),
    ("08", "训练过程", "6亿图像-文本对的炼丹之旅"),
    ("09", "CFG 引导技术", "无分类器引导，让AI更听话"),
    ("10", "推理过程", "从噪声到图像的雕刻过程"),
    ("11", "采样器对比", "9种算法，速度与质量取舍"),
    ("12", "参数优化", "5个核心参数的最佳实践"),
    ("13", "Prompt 技巧", "给AI画师下订单的艺术"),
    ("14", "LoRA 原理", "给画家添加专项风格手册"),
    ("15", "ControlNet 原理", "精确控制构图和姿态"),
    ("16", "应用场景", "5大领域的革命性影响"),
    ("17", "SD3 / SDXL 最新进展", "MMDiT架构与Flow Matching"),
    ("18", "Flux 等新模型对比", "2026年开源图像生成格局"),
    ("19", "2026年AI图像生成动态", "7大趋势与未来展望"),
]

row_h = 0.30
y_start = 1.30
for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(y_start + i * row_h)
    # Alternating row background
    if i % 2 == 0:
        rect(sl, Inches(0.40), y, Inches(12.50), Inches(row_h), fc=C_LT)
    # Number
    tb(sl, Inches(0.60), y, Inches(0.50), Inches(row_h), num, fs=11, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
    # Title
    tb(sl, Inches(1.20), y, Inches(4.20), Inches(row_h), title, fs=11, fc=C_PRI, b=False)
    # Description
    tb(sl, Inches(5.60), y, Inches(7.00), Inches(row_h), desc, fs=11, fc=C_G)

chk(sl, "P2 目录")


# ============================================================
# P3 Diffusion基础 (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P3: Diffusion基础 (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "01-02  Diffusion Model 基础 + SD vs 普通 Diffusion", "核心创新：在潜在空间中工作，计算量降低48倍")

# 前向扩散
rrect(sl, Inches(0.50), Inches(1.30), Inches(6.00), Inches(1.70), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(0.70), Inches(1.35), Inches(5.50), Inches(0.30),
   "前向扩散：清晰照片 → 纯噪声（1000步）", fs=11, fc=C_ORG, b=True)
mtb(sl, Inches(0.70), Inches(1.70), Inches(5.60), Inches(1.20), [
    "就像把一滴蓝色墨水滴入清水中，墨水逐渐扩散，最终整杯水变成均匀的浅蓝色。",
    "系统在1000步中，每步往图像上添加随机噪声，最终变成纯噪声的雪花屏。",
    "图像从 512x512x3 = 786,432 个有意义像素退化成完全随机数值。这个过程不需要学习。",
], fs=9, fc=C_D, ls=1.15)

# 反向去噪
rrect(sl, Inches(6.80), Inches(1.30), Inches(6.00), Inches(1.70), fc=C_BG_GRN, lc=C_GRN, lw=1)
tb(sl, Inches(7.00), Inches(1.35), Inches(5.50), Inches(0.30),
   "反向去噪：纯噪声 → 清晰照片（AI学到的能力）", fs=11, fc=C_GRN, b=True)
mtb(sl, Inches(7.00), Inches(1.70), Inches(5.60), Inches(1.20), [
    "AI修复师面对雪花屏，一步一步去掉噪声，每步恢复一点细节，最终还原清晰照片。",
    "这就是扩散模型通过海量训练数据学到的核心技能。",
    "比GAN更好：GAN像一下子跳到终点（不稳定），扩散模型是一步一步走（稳定高质量）。",
], fs=9, fc=C_D, ls=1.15)

# 对比表
table_shape = sl.shapes.add_table(8, 4, Inches(0.50), Inches(3.20), Inches(12.30), Inches(2.56))
table = table_shape.table
table.columns[0].width = Inches(2.00)
table.columns[1].width = Inches(3.60)
table.columns[2].width = Inches(3.60)
table.columns[3].width = Inches(3.10)

headers = ["对比维度", "普通 Diffusion（DDPM）", "Stable Diffusion", "核心区别"]
data = [
    ["操作空间", "像素空间（512x512x3）", "潜在空间（64x64x4）", "SD不在像素空间工作"],
    ["数据维度", "786,432 维", "16,384 维", "维度降低48倍"],
    ["计算量", "极高", "降低约 48 倍", "48倍压缩是核心创新"],
    ["GPU显存", "40GB+（高端GPU）", "8-12GB（消费级GPU）", "普通人也能运行"],
    ["生成速度", "几十秒到几分钟", "几秒到十几秒", "实用性大幅提升"],
    ["代表模型", "DDPM, DDIM", "SD 1.5, SDXL, SD3", "SD成为最成功开源模型"],
    ["硬件门槛", "高（普通人用不起）", "低（RTX 3060即可）", "让AI画画人人可用"],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
            if j == 3: p.font.color.rgb = C_ORG
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 底部总结
rrect(sl, Inches(0.50), Inches(5.90), Inches(12.30), Inches(1.00), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.70), Inches(5.95), Inches(11.80), Inches(0.90),
   "核心比喻：普通Diffusion像在4K画布上逐像素清理噪声（80万个数值）；SD先用VAE压缩成缩略图（1.6万个数值），在缩略图上去噪声后再还原。\n"
   "就像画家不是在画布上逐像素画画，而是在缩略图上画画，画完再用放大镜还原成高清大图。这正是SD成为2022年以来最成功开源AI图像生成模型的根本原因。",
   fs=9, fc=C_D, ls=1.2)

chk(sl, "P3 Diffusion基础")


# ============================================================
# P4 整体架构 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P4: 整体架构概览...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "03  整体架构概览", "CLIP + U-Net + VAE 三大组件协同工作 — 推理与训练流程完全分离")

# ── 推理流程区域 ──
rrect(sl, Inches(0.15), Inches(1.15), Inches(12.80), Inches(2.05), fc=None, lc=C_ACC, lw=2)
rect(sl, Inches(0.15), Inches(1.15), Inches(12.80), Inches(0.30), fc=C_ACC)
tb(sl, Inches(0.30), Inches(1.17), Inches(8.0), Inches(0.28),
   "▶ 推理流程（从文字到图像）— 主要使用场景", fs=11, fc=C_W, b=True)

# ① 文字输入
rrect(sl, Inches(0.35), Inches(1.60), Inches(1.40), Inches(0.85), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(0.35), Inches(1.62), Inches(1.40), Inches(0.22),
   "① 文字输入", fs=10, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(1.87), Inches(1.40), Inches(0.50),
   '"一只橘猫\n坐在沙发上"', fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(1.80), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ② CLIP
rrect(sl, Inches(2.15), Inches(1.55), Inches(1.80), Inches(0.95), fc=C_LT, lc=C_ACC, lw=1)
rect(sl, Inches(2.15), Inches(1.55), Inches(1.80), Inches(0.25), fc=C_ACC)
tb(sl, Inches(2.15), Inches(1.57), Inches(1.80), Inches(0.22),
   "② CLIP 文本编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(2.15), Inches(1.83), Inches(1.80), Inches(0.60),
   "123M参数\n12层Transformer\n输出: 77×768向量", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(4.00), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ③ U-Net
rrect(sl, Inches(4.35), Inches(1.50), Inches(2.00), Inches(1.05), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI, lw=2)
rect(sl, Inches(4.35), Inches(1.50), Inches(2.00), Inches(0.25), fc=C_PRI)
tb(sl, Inches(4.35), Inches(1.52), Inches(2.00), Inches(0.22),
   "③ U-Net 去噪（核心）", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(4.35), Inches(1.78), Inches(2.00), Inches(0.70),
   "860M参数(80%)\n⟳ 循环20-50步\n噪声→逐步去噪", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
rect(sl, Inches(3.95), Inches(1.78), Inches(0.35), Inches(0.20), fc=C_PUR)
tb(sl, Inches(3.95), Inches(1.78), Inches(0.35), Inches(0.20), "CA", fs=8, fc=C_W, b=True, a=PP_ALIGN.CENTER)
arrow_r(sl, Inches(6.40), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ④ VAE 解码器
rrect(sl, Inches(6.75), Inches(1.55), Inches(1.80), Inches(0.95), fc=C_BG_GRN, lc=C_GRN, lw=1)
rect(sl, Inches(6.75), Inches(1.55), Inches(1.80), Inches(0.25), fc=C_GRN)
tb(sl, Inches(6.75), Inches(1.57), Inches(1.80), Inches(0.22),
   "④ VAE 解码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(6.75), Inches(1.83), Inches(1.80), Inches(0.60),
   "~49M参数\n64×64→512×512\n4层上采样", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(8.60), Inches(1.88), Inches(0.30), Inches(0.20), C_ACC)

# ⑤ 图像输出
rrect(sl, Inches(8.95), Inches(1.60), Inches(1.40), Inches(0.85), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(8.95), Inches(1.62), Inches(1.40), Inches(0.22),
   "⑤ 图像输出", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(8.95), Inches(1.87), Inches(1.40), Inches(0.50),
   "512×512×3\n高清RGB图像", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# 随机噪声
rrect(sl, Inches(4.65), Inches(2.65), Inches(1.40), Inches(0.45), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(4.65), Inches(2.68), Inches(1.40), Inches(0.40),
   "随机噪声 N(0,I)\n64×64×4 纯噪声", fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(5.28), Inches(2.55), Inches(0.18), Inches(0.15), C_RED)

# VAE编码器不需要
rrect(sl, Inches(10.60), Inches(1.60), Inches(2.20), Inches(0.85), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(10.60), Inches(1.65), Inches(2.20), Inches(0.75),
   "❌ VAE 编码器\n推理时不需要！\n（编码器只在训练时使用）", fs=9, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.1)

# ── 训练流程区域 ──
rrect(sl, Inches(0.15), Inches(3.35), Inches(12.80), Inches(1.80), fc=None, lc=C_ORG, lw=2)
rect(sl, Inches(0.15), Inches(3.35), Inches(12.80), Inches(0.30), fc=C_ORG)
tb(sl, Inches(0.30), Inches(3.37), Inches(8.0), Inches(0.28),
   "▶ 训练流程（从图像学习去噪能力）— 模型训练阶段", fs=11, fc=C_W, b=True)

rrect(sl, Inches(0.35), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_W, lc=C_ORG, lw=1)
tb(sl, Inches(0.35), Inches(3.82), Inches(1.60), Inches(0.22), "图像+文本对", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(4.07), Inches(1.60), Inches(0.40), "512×512×3 图像\n+ 文字描述", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(2.00), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

rrect(sl, Inches(2.35), Inches(3.75), Inches(1.80), Inches(0.85), fc=C_BG_GRN, lc=C_GRN, lw=1)
rect(sl, Inches(2.35), Inches(3.75), Inches(1.80), Inches(0.25), fc=C_GRN)
tb(sl, Inches(2.35), Inches(3.77), Inches(1.80), Inches(0.22), "VAE 编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(2.35), Inches(4.03), Inches(1.80), Inches(0.50), "~34M参数\n512²→64² 压缩48倍", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(4.20), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

rrect(sl, Inches(4.55), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(4.55), Inches(3.82), Inches(1.60), Inches(0.22), "添加噪声", fs=10, fc=C_RED, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(4.55), Inches(4.07), Inches(1.60), Inches(0.40), "z_t = √ᾱ·z₀ + √(1-ᾱ)·ε\n前向扩散过程", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(6.20), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

rrect(sl, Inches(6.55), Inches(3.75), Inches(2.00), Inches(0.85), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI, lw=1)
rect(sl, Inches(6.55), Inches(3.75), Inches(2.00), Inches(0.25), fc=C_PRI)
tb(sl, Inches(6.55), Inches(3.77), Inches(2.00), Inches(0.22), "U-Net 预测噪声", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(6.55), Inches(4.03), Inches(2.00), Inches(0.50), "输入: 噪声潜空间+文本向量\n输出: ε_pred 预测噪声", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)
arrow_r(sl, Inches(8.60), Inches(4.05), Inches(0.30), Inches(0.20), C_ORG)

rrect(sl, Inches(8.95), Inches(3.80), Inches(1.60), Inches(0.75), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(8.95), Inches(3.82), Inches(1.60), Inches(0.22), "计算损失+反向传播", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(8.95), Inches(4.07), Inches(1.60), Inches(0.40), "L = ||ε - ε_pred||²\n更新U-Net 860M参数", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

rrect(sl, Inches(4.55), Inches(4.70), Inches(1.60), Inches(0.35), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(4.55), Inches(4.72), Inches(1.60), Inches(0.30), "CLIP编码文本(同时进行)", fs=8, fc=C_ACC, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(5.28), Inches(4.60), Inches(0.18), Inches(0.12), C_ACC)

rrect(sl, Inches(10.60), Inches(3.80), Inches(2.20), Inches(0.75), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(10.60), Inches(3.85), Inches(2.20), Inches(0.65),
   "❌ VAE 解码器\n训练时不需要！\n（解码器只在推理时使用）", fs=9, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.1)

# ── 组件对照表 ──
rect(sl, Inches(0.15), Inches(5.30), Inches(12.80), Inches(0.28), fc=C_PRI)
tb(sl, Inches(0.30), Inches(5.31), Inches(12.50), Inches(0.26), "组件使用对照表", fs=11, fc=C_W, b=True)

rrect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(1.20), fc=C_W, lc=C_ACC, lw=1)
rect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_ACC)
tb(sl, Inches(0.25), Inches(5.67), Inches(3.90), Inches(0.22), "CLIP 文本编码器 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(0.25), Inches(5.95), Inches(3.90), Inches(0.85),
   "• 翻译官 — 123M参数\n• 12层Transformer, 768维\n• 作用: 把人类语言翻译成AI数字语言", fs=9, fc=C_D, ls=1.15)

rrect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(1.20), fc=C_W, lc=C_PRI, lw=1)
rect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_PRI)
tb(sl, Inches(4.50), Inches(5.67), Inches(3.90), Inches(0.22), "U-Net 去噪网络 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(4.50), Inches(5.95), Inches(3.90), Inches(0.85),
   "• 核心画师 — 860M参数(80%)\n• U型编码器-解码器结构\n• 作用: 一步步去掉噪声", fs=9, fc=C_D, ls=1.15)

rrect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(1.20), fc=C_W, lc=C_GRN, lw=1)
rect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(0.25), fc=C_GRN)
tb(sl, Inches(8.75), Inches(5.67), Inches(4.10), Inches(0.22), "VAE — 编码器(训练) 解码器(推理)", fs=10, fc=C_W, b=True)
tb(sl, Inches(8.75), Inches(5.95), Inches(4.10), Inches(0.85),
   "• 编码器~34M: 图像→潜空间 [训练]\n• 解码器~49M: 潜空间→图像 [推理]\n• 48倍压缩: 786K→16K数值", fs=9, fc=C_D, ls=1.15)

chk(sl, "P4 整体架构")


# ============================================================
# P5 CLIP详解 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P5: CLIP详解...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "04  CLIP 文本编码器详解", "Contrastive Language-Image Pre-training — 从文字到语义向量的完整旅程")

info_cards = [
    ("输入", "文本描述字符串\n分词器处理\n最多77个token\n约50-60个汉字", C_ACC),
    ("输出", "77×768维张量\n共59,136个浮点数\n文本的语义指纹\n包含所有含义信息", C_GRN),
    ("参数", "约1.23亿(123M)\n12层Transformer\n隐藏维度768\n12个注意力头", C_ORG),
    ("作用", "人类文字→数字向量\n通过Cross-Attention\n注入U-Net指导去噪\n没有它AI不知道画什么", C_PUR),
]
for i, (title, body, color) in enumerate(info_cards):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.25), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.10), Inches(1.50), Inches(2.80), Inches(0.90), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

steps = [
    ("① 分词",
     "文本拆成最小语义单元(token)\n\n例: \"一只橘猫坐在沙发上\"\n→ [一只/橘/猫/坐在/沙发/上]\n\n每个token映射成768维向量",
     C_ACC),
    ("② 位置编码",
     "为什么需要位置编码？\n\"猫坐在沙发上\" ≠ \"沙发坐在猫上\"\n\n词一样，顺序不同，含义天差地别。\nTransformer本身不区分顺序，\n必须额外注入位置信息。",
     C_ORG),
    ("③ Transformer",
     "自注意力机制 (Self-Attention)\n每个词都能\"看到\"其他所有词\n\n\"橘猫\"→注意\"坐在\"+\"沙发\"\n\"坐在\"→连接\"橘猫\"+\"沙发\"\n\n12层: 浅层语法→深层语义",
     C_PUR),
    ("④ 输出向量",
     "输出: 77×768 语义向量\n\n每个token变成768维向量\n共77×768 = 59,136个浮点数\n\n不仅含每个词的独立含义，\n还包含词与词之间的关系信息。",
     C_GRN),
]

for i, (title, body, color) in enumerate(steps):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(2.65), Inches(3.00), Inches(3.50), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(2.65), Inches(3.00), Inches(0.28), fc=color)
    tb(sl, x + Inches(0.05), Inches(2.67), Inches(2.90), Inches(0.24), title, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.15), Inches(3.00), Inches(2.70), Inches(3.00), body, fs=9, fc=C_D, ls=1.2)
    if i < len(steps) - 1:
        arrow_r(sl, x + Inches(3.00), Inches(4.20), Inches(0.20), Inches(0.18), color)

rrect(sl, Inches(0.30), Inches(6.30), Inches(12.70), Inches(0.60), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(6.35), Inches(12.30), Inches(0.50),
   "总结: CLIP = 分词 → 位置编码 → 12层Transformer(自注意力) → 77×768语义向量。这个向量通过Cross-Attention注入U-Net，指导每一步去噪方向，确保生成的图像符合文字描述。",
   fs=10, fc=C_D, ls=1.2)

chk(sl, "P5 CLIP详解")


# ============================================================
# P6 VAE详解 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P6: VAE详解...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "05  VAE 编解码器详解", "Variational Autoencoder — 压缩48倍但肉眼几乎看不出差别")

info_cards = [
    ("编码器输入", "原始图像\n512×512×3\n786,432个像素值\n值域 0-255", C_ACC),
    ("编码器输出", "潜空间\n64×64×4\n16,384个浮点数\n值域约[-5,5]", C_GRN),
    ("总参数量", "约8300万(83M)\n编码器~34M\n解码器~49M\n解码器更大", C_ORG),
    ("为什么用VAE?", "正则化确保\n潜空间连续平滑\n普通AE潜空间不规则\n扩散操作需要平滑空间", C_PUR),
]
for i, (title, body, color) in enumerate(info_cards):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.15), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.10), Inches(1.50), Inches(2.80), Inches(0.80), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# 编码器
rrect(sl, Inches(0.30), Inches(2.50), Inches(6.20), Inches(3.05), fc=C_W, lc=C_ACC, lw=1)
rect(sl, Inches(0.30), Inches(2.50), Inches(6.20), Inches(0.28), fc=C_ACC)
tb(sl, Inches(0.40), Inches(2.52), Inches(6.00), Inches(0.24),
   "▼ 编码器 (Encoder): 4层下采样 — 图像→潜空间（训练时使用）", fs=10, fc=C_W, b=True)

enc_steps = [
    ("512×512×3", "输入图像（RGB三通道）", C_ACC),
    ("256×256×128", "通道3→128，提取低级特征（边缘、色块）", C_ACC),
    ("128×128×256", "提取中级特征（纹理、局部结构）", C_ACC),
    ("64×64×512", "提取高级特征（物体部件、语义信息）", C_ACC),
    ("32×32×512", "最深层特征表示", C_ACC),
    ("64×64×4", "输出潜空间（48倍压缩！）", C_GRN),
]
for i, (dim, desc, color) in enumerate(enc_steps):
    y = Inches(2.85 + i * 0.42)
    rrect(sl, Inches(0.45), y, Inches(1.40), Inches(0.30), fc=C_LT, lc=color, lw=1)
    tb(sl, Inches(0.45), y + Inches(0.03), Inches(1.40), Inches(0.24), dim, fs=9, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    if i < len(enc_steps) - 1:
        arrow_d(sl, Inches(1.10), y + Inches(0.30), Inches(0.15), Inches(0.12), C_ACC)
    else:
        arrow_d(sl, Inches(1.10), y + Inches(0.30), Inches(0.15), Inches(0.12), C_GRN)
    tb(sl, Inches(2.00), y + Inches(0.03), Inches(4.30), Inches(0.24), desc, fs=9, fc=C_D)

tb(sl, Inches(0.45), Inches(5.45), Inches(5.90), Inches(0.20),
   "编码器参数量: ~34M | 宽高每次缩小一半，通道数逐步增加", fs=8, fc=C_G)

# 解码器
rrect(sl, Inches(6.80), Inches(2.50), Inches(6.20), Inches(3.05), fc=C_W, lc=C_GRN, lw=1)
rect(sl, Inches(6.80), Inches(2.50), Inches(6.20), Inches(0.28), fc=C_GRN)
tb(sl, Inches(6.90), Inches(2.52), Inches(6.00), Inches(0.24),
   "▲ 解码器 (Decoder): 4层上采样 — 潜空间→图像（推理时使用）", fs=10, fc=C_W, b=True)

dec_steps = [
    ("64×64×4", "输入潜空间", C_GRN),
    ("128×128×512", "通道4→512，开始重建空间细节", C_GRN),
    ("256×256×256", "通道减半，分辨率翻倍", C_GRN),
    ("512×512×128", "继续增加分辨率，减少通道", C_GRN),
    ("512×512×64", "接近最终分辨率", C_GRN),
    ("512×512×3", "输出高清图像（肉眼几乎无损！）", C_ORG),
]
dec_y_start = 2.85 + (len(dec_steps) - 1) * 0.42
for i, (dim, desc, color) in enumerate(dec_steps):
    y = Inches(dec_y_start - i * 0.42)
    rrect(sl, Inches(6.95), y, Inches(1.40), Inches(0.30), fc=C_BG_GRN, lc=color, lw=1)
    tb(sl, Inches(6.95), y + Inches(0.03), Inches(1.40), Inches(0.24), dim, fs=9, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    if i < len(dec_steps) - 1:
        arrow_u(sl, Inches(7.60), y - Inches(0.12), Inches(0.15), Inches(0.12), C_GRN)
    tb(sl, Inches(8.50), y + Inches(0.03), Inches(4.30), Inches(0.24), desc, fs=9, fc=C_D)

tb(sl, Inches(6.95), Inches(5.45), Inches(5.90), Inches(0.20),
   "解码器参数量: ~49M（比编码器大，因为\"还原\"比\"压缩\"更复杂）", fs=8, fc=C_G)

# 底部对比
rrect(sl, Inches(0.30), Inches(5.75), Inches(12.70), Inches(1.15), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.78), Inches(12.30), Inches(0.22),
   "💡 为什么用VAE而不是普通自编码器(AE)？", fs=11, fc=C_PRI, b=True)
tb(sl, Inches(0.50), Inches(6.05), Inches(12.30), Inches(0.80),
   "普通AE的问题: 潜空间\"不规则\" — 两个相近的潜空间点可能映射到完全不相关的图像，潜空间中存在大量\"空白区域\"。\n"
   "在 irregular 空间做扩散(加减噪声)会导致解码出无意义的图像。VAE 通过 KL 散度正则化强制潜空间服从标准正态分布，\n"
   "使其连续、平滑、无空洞 — 这是扩散模型能在潜空间上正常工作的前提条件。",
   fs=9, fc=C_D, ls=1.2)

chk(sl, "P6 VAE详解")


# ============================================================
# P7 U-Net架构 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P7: U-Net架构...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "06  U-Net 架构图", "编码器-瓶颈-解码器 U 型结构 | 860M 参数 | 3 个关键机制")

rrect(sl, Inches(0.10), Inches(1.10), Inches(6.50), Inches(5.85), fc=None, lc=C_PRI, lw=2)

rrect(sl, Inches(0.20), Inches(1.20), Inches(2.50), Inches(4.50), fc=None, lc=C_ACC, lw=1)
tb(sl, Inches(0.30), Inches(1.22), Inches(2.30), Inches(0.22),
   "编码器 Encoder ↓", fs=11, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)

enc_layers = [
    ("64×64×320", "Res+Attn+Down", "Cross-Attention"),
    ("32×32×640", "Res+Attn+Down", "Cross-Attention"),
    ("16×16×1280", "Res+Attn+Down", "Cross-Attention"),
    ("8×8×1280", "Res+Down", "无注意力"),
]
enc_y = 1.50
for dim, ops, attn in enc_layers:
    rrect(sl, Inches(0.35), Inches(enc_y), Inches(2.20), Inches(0.70), fc=C_W, lc=C_ACC, lw=1)
    tb(sl, Inches(0.40), Inches(enc_y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.40), Inches(enc_y+0.28), Inches(2.10), Inches(0.20), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.40), Inches(enc_y+0.48), Inches(2.10), Inches(0.18), f"{'✦ ' if 'Cross' in attn else ''}{attn}", fs=8, fc=C_PUR if 'Cross' in attn else C_G, a=PP_ALIGN.CENTER)
    enc_y += 0.85
    if 'Cross' in attn:
        arrow_d(sl, Inches(1.40), Inches(enc_y-0.85+0.70), Inches(0.18), Inches(0.15), C_ACC)

arrow_d(sl, Inches(1.40), Inches(enc_y-0.85+0.70), Inches(0.18), Inches(0.15), C_ACC)

bottleneck_y = enc_y
rrect(sl, Inches(0.30), Inches(bottleneck_y), Inches(2.30), Inches(0.70), fc=C_BG_ORG, lc=C_ORG, lw=2)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.02), Inches(2.20), Inches(0.22), "4×4×1280", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.25), Inches(2.20), Inches(0.20), "Bottleneck", fs=9, fc=C_D, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.35), Inches(bottleneck_y+0.45), Inches(2.20), Inches(0.20), "Res+SelfAttn+Res", fs=8, fc=C_D, a=PP_ALIGN.CENTER)

rrect(sl, Inches(3.00), Inches(1.20), Inches(2.50), Inches(4.50), fc=None, lc=C_GRN, lw=1)
tb(sl, Inches(3.10), Inches(1.22), Inches(2.30), Inches(0.22),
   "解码器 Decoder ↑", fs=11, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

dec_layers = [
    ("8×8×1280", "Res+Up+Skip", "Cross-Attention"),
    ("16×16×640", "Res+Up+Skip", "Cross-Attention"),
    ("32×32×320", "Res+Up+Skip", "Cross-Attention"),
    ("64×64×4", "输出预测噪声", "无注意力"),
]
dec_ys = [4.05, 3.20, 2.35, 1.50]

for i, (dim, ops, attn) in enumerate(dec_layers):
    y = dec_ys[i]
    rrect(sl, Inches(3.20), Inches(y), Inches(2.20), Inches(0.70), fc=C_BG_GRN, lc=C_GRN, lw=1)
    tb(sl, Inches(3.25), Inches(y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(3.25), Inches(y+0.28), Inches(2.10), Inches(0.20), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(3.25), Inches(y+0.48), Inches(2.10), Inches(0.18), f"{'✦ ' if 'Cross' in attn else ''}{attn}", fs=8, fc=C_PUR if 'Cross' in attn else C_G, a=PP_ALIGN.CENTER)
    if i < len(dec_layers) - 1:
        arrow_u(sl, Inches(4.25), Inches(y-0.15), Inches(0.18), Inches(0.15), C_GRN)
        tb(sl, Inches(4.50), Inches(y-0.18), Inches(0.90), Inches(0.18), "上采样↑", fs=8, fc=C_GRN, b=True)

line_h(sl, Inches(2.60), Inches(4.90 + 0.30), Inches(0.60), C_GRN, 3)
rect(sl, Inches(3.10), Inches(4.05 + 0.30), Emu(int(3*12700)), Inches(0.85), fc=C_GRN)
arrow_r(sl, Inches(3.10), Inches(4.05 + 0.25), Inches(0.15), Inches(0.12), C_GRN)
tb(sl, Inches(2.55), Inches(4.90 - 0.18), Inches(0.65), Inches(0.18), "过渡→", fs=8, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

for i in range(4):
    enc_mid_y = Inches(1.50 + i * 0.85) + Inches(0.30)
    line_h(sl, Inches(2.55), enc_mid_y, Inches(0.65), C_ORG, 3)
    rrect(sl, Inches(2.60), enc_mid_y - Inches(0.18), Inches(0.55), Inches(0.18), fc=C_ORG)
    tb(sl, Inches(2.60), enc_mid_y - Inches(0.18), Inches(0.55), Inches(0.18), "Skip", fs=8, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    arrow_r(sl, Inches(3.15), enc_mid_y - Inches(0.04), Inches(0.12), Inches(0.10), C_ORG)

rrect(sl, Inches(5.70), Inches(1.20), Inches(0.80), Inches(5.30), fc=None, lc=C_TEL, lw=1)
tb(sl, Inches(5.72), Inches(1.25), Inches(0.76), Inches(0.22), "Time", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(5.72), Inches(1.45), Inches(0.76), Inches(0.22), "Embed", fs=8, fc=C_TEL, b=True, a=PP_ALIGN.CENTER)
for y_pos in [1.80, 2.65, 3.50, 4.35, 5.20]:
    line_h(sl, Inches(5.70), Inches(y_pos), Inches(0.80), C_TEL, 1)

# 右面板
rrect(sl, Inches(6.80), Inches(1.20), Inches(6.20), Inches(1.65), fc=C_W, lc=C_PUR, lw=1)
rect(sl, Inches(6.80), Inches(1.20), Inches(6.20), Inches(0.28), fc=C_PUR)
tb(sl, Inches(6.90), Inches(1.22), Inches(6.00), Inches(0.24),
   "① Cross-Attention 交叉注意力 — 文本条件注入的核心机制", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(1.55), Inches(6.00), Inches(1.20), [
    "• 图像特征作为 Query，文本特征作为 Key-Value",
    "• 图像问：\"根据文字，我这里该怎样？\" → 文本答：\"你这里应该是橘猫的毛发\"",
    "• 出现位置：编码器前3层 + 解码器前3层（共6处）",
    "• 注意力计算: Q·K^T/√d → Softmax → ×V",
], fs=9, fc=C_D, ls=1.15)

rrect(sl, Inches(6.80), Inches(3.00), Inches(6.20), Inches(1.65), fc=C_W, lc=C_ORG, lw=1)
rect(sl, Inches(6.80), Inches(3.00), Inches(6.20), Inches(0.28), fc=C_ORG)
tb(sl, Inches(6.90), Inches(3.02), Inches(6.00), Inches(0.24),
   "② Skip Connections 跳跃连接 — 信息传递的桥梁（图中橙色横线）", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(3.35), Inches(6.00), Inches(1.20), [
    "• 编码器特征直接桥接到解码器对应层（同分辨率级别）",
    "• 防止细节在压缩过程中丢失",
    "• 像建筑师随时参考草稿本，没有它解码器只能\"凭记忆\"还原",
    "• 连接: 64×64↔64×64, 32×32↔32×32, 16×16↔16×16, 8×8↔8×8",
], fs=9, fc=C_D, ls=1.15)

rrect(sl, Inches(6.80), Inches(4.80), Inches(6.20), Inches(1.85), fc=C_W, lc=C_TEL, lw=1)
rect(sl, Inches(6.80), Inches(4.80), Inches(6.20), Inches(0.28), fc=C_TEL)
tb(sl, Inches(6.90), Inches(4.82), Inches(6.00), Inches(0.24),
   "③ Time Embedding 时间嵌入 — 当前噪声水平的标尺", fs=10, fc=C_W, b=True)
mtb(sl, Inches(6.90), Inches(5.15), Inches(6.00), Inches(1.40), [
    "• 告诉U-Net当前噪声水平（1-1000步）",
    "• 高噪声(t≈900)：粗雕确定轮廓",
    "• 低噪声(t≈100)：精雕添加细节",
    "• 正弦编码 → 2层MLP → 1280维",
    "• 注入每一个ResBlock，影响去噪策略",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P7 U-Net架构")


# ============================================================
# P8 协作流程 (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P8: 协作流程 (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "07  三大组件协作流程", "AI画图流水线 — 各司其职，紧密配合")

# 3站横向卡片
stations = [
    ("第1站", "CLIP 翻译站", C_ACC,
     "用户输入: 一只橘猫坐在木质沙发上\n\n通过12层Transformer处理\n\n输出77×768维语义指纹向量\n\n发送给下一站的U-Net",
     "翻译官 — 把人类语言翻译成AI数字语言"),
    ("第2站", "U-Net 画图站", C_RED,
     "接收文本向量+64×64×4纯随机噪声\n\n第1轮: 只能猜出大致色块分布\n第15轮: 能看到猫轮廓、沙发形状\n第30轮: 毛发纹理、阳光光线\n第50轮: 去噪完成",
     "核心画师 — 执行20-50轮去噪循环"),
    ("第3站", "VAE 还原站", C_GRN,
     "接收U-Net输出的干净潜空间\n\n通过4层上采样逐步还原\n64×64→128→256→512\n\n最终输出512×512×3高清图像",
     "还原打印机 — 潜空间→高清图像"),
]

for i, (station, name, color, content, footer) in enumerate(stations):
    x = Inches(0.50 + i * 4.20)
    # Card
    rrect(sl, x, Inches(1.30), Inches(3.90), Inches(3.50), fc=C_W, lc=color, lw=1)
    # Header
    rect(sl, x, Inches(1.30), Inches(3.90), Inches(0.50), fc=color)
    tb(sl, x + Inches(0.10), Inches(1.32), Inches(1.00), Inches(0.22), station, fs=10, fc=C_W, b=True)
    tb(sl, x + Inches(1.10), Inches(1.55), Inches(2.70), Inches(0.22), name, fs=12, fc=C_W, b=True)
    # Content
    tb(sl, x + Inches(0.15), Inches(1.90), Inches(3.60), Inches(2.30), content, fs=9, fc=C_D, ls=1.2)
    # Footer
    rrect(sl, x + Inches(0.10), Inches(4.30), Inches(3.70), Inches(0.35), fc=C_LT, lc=color, lw=1)
    tb(sl, x + Inches(0.15), Inches(4.32), Inches(3.60), Inches(0.30), footer, fs=8, fc=color, b=True, a=PP_ALIGN.CENTER)
    # Arrow between stations
    if i < 2:
        arrow_r(sl, x + Inches(3.90), Inches(2.50), Inches(0.25), Inches(0.20), color)

# 为什么要这样分工
rrect(sl, Inches(0.50), Inches(5.10), Inches(12.30), Inches(1.80), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.70), Inches(5.15), Inches(5.00), Inches(0.25),
   "为什么要这样分工？", fs=11, fc=C_PRI, b=True)
mtb(sl, Inches(0.70), Inches(5.45), Inches(11.80), Inches(1.40), [
    "• CLIP擅长理解语言但不擅长生成图像 — 所以只让它做翻译",
    "• U-Net擅长在低维空间做逐步变换但不擅长人类语言 — 需要CLIP的翻译结果作为指导",
    "• VAE擅长压缩和还原但不擅长语义 — 只在开头和结尾使用",
    "• 就像专业厨房：有人负责理解客人的点单(CLIP)，有人负责烹饪(U-Net)，有人负责摆盘(VAE)",
    "",
    "• 数据连贯性：CLIP输出的文本向量贯穿整个U-Net去噪过程(通过Cross-Attention)，确保每一步都\"记得\"用户要什么",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P8 协作流程")


# ============================================================
# P9 训练过程 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P9: 训练过程...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "训练过程流程图", "6步训练管线 — 15万GPU小时，50-60万美元")

steps_data = [
    ("Step 1", "采样图像+文本", "从LAION-5B随机抽取\n图像x₀ + 文本描述", "6亿图文对", C_ACC),
    ("Step 2", "VAE 编码", "512×512×3 → 64×64×4\n压缩到潜空间(48倍)", "z₀ = VAE(x₀)", C_GRN),
    ("Step 3", "添加随机噪声", "随机选t ∈ [1,1000]\n按比例添加高斯噪声", "z_t = √ᾱ·z₀ + √(1-ᾱ)·ε", C_RED),
    ("Step 4", "U-Net 预测噪声", "输入: z_t + t + 文本向量\n输出: 预测噪声 ε_pred", "ε_pred = U-Net(z_t, t, c)", C_PRI),
    ("Step 5", "计算损失", "预测噪声 vs 真实噪声\n均方误差(MSE)", "L = ‖ε - ε_pred‖²", C_ORG),
    ("Step 6", "反向传播", "更新U-Net的860M参数\nCLIP + VAE保持冻结", "θ ← θ - η·∇L", C_PUR),
]

for i, (step, title, desc, formula, color) in enumerate(steps_data):
    col = i % 3
    row = i // 3
    x = Inches(0.50 + col * 4.20)
    y = Inches(1.30 + row * 2.30)
    rrect(sl, x, y, Inches(3.80), Inches(2.00), fc=C_W, lc=color, lw=1)
    rect(sl, x, y, Inches(0.70), Inches(0.30), fc=color)
    tb(sl, x + Inches(0.05), y + Inches(0.03), Inches(0.60), Inches(0.24), step, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.80), y + Inches(0.03), Inches(2.90), Inches(0.24), title, fs=11, fc=color, b=True)
    tb(sl, x + Inches(0.15), y + Inches(0.40), Inches(3.50), Inches(0.80), desc, fs=10, fc=C_D, ls=1.2)
    rrect(sl, x + Inches(0.15), y + Inches(1.35), Inches(3.50), Inches(0.45), fc=C_LT, lc=color, lw=1)
    tb(sl, x + Inches(0.25), y + Inches(1.40), Inches(3.30), Inches(0.35), formula, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER, fn="Consolas")
    if col < 2 and i < 5:
        arrow_r(sl, x + Inches(3.80), y + Inches(0.85), Inches(0.30), Inches(0.18), color)

arrow_d(sl, Inches(10.50), Inches(3.40), Inches(0.18), Inches(0.18), C_ACC)

rrect(sl, Inches(0.30), Inches(5.80), Inches(12.70), Inches(0.40), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.85), Inches(12.30), Inches(0.30),
   "⟳ 重复训练循环：重复 Step 1-6 数十万次，直到损失收敛 | 每次迭代随机采样不同的图像和噪声水平",
   fs=10, fc=C_D, b=False, a=PP_ALIGN.CENTER)

table_shape = sl.shapes.add_table(2, 5, Inches(0.50), Inches(6.35), Inches(12.30), Inches(0.60))
table = table_shape.table
headers_t = ["数据集", "GPU", "训练时间", "成本", "Batch Size"]
vals_t = ["LAION-5B (60亿图文对子集)", "NVIDIA A100 × 256", "~15万GPU小时", "$50-60万", "2048×2048"]
for j in range(5):
    table.columns[j].width = Inches(2.46)
    cell = table.cell(0, j); cell.text = headers_t[j]
    for p in cell.text_frame.paragraphs: p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=C_W; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
    cell = table.cell(1, j); cell.text = vals_t[j]
    for p in cell.text_frame.paragraphs: p.font.size=Pt(9); p.font.color.rgb=C_D; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_W

chk(sl, "P9 训练过程")


# ============================================================
# P10 推理过程 (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P10: 推理过程 (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "推理过程流程图", "4步推理 + 去噪循环详解 — 从纯噪声到高清图像")

# 4步横向
infer_steps = [
    ("Step 1", "文本编码", "文字 → CLIP → 77×768 向量\n只需执行一次", C_ACC),
    ("Step 2", "初始化噪声", "采样 64×64×4 随机张量\n= 纯噪声雪花屏", C_ORG),
    ("Step 3", "逐步去噪循环", "U-Net 反复预测噪声\n20-50 步逐步清晰", C_RED),
    ("Step 4", "VAE 解码", "64×64×4 → 512×512×3\n4层上采样还原", C_GRN),
]

for i, (step, title, desc, color) in enumerate(infer_steps):
    x = Inches(0.50 + i * 3.20)
    rrect(sl, x, Inches(1.25), Inches(2.90), Inches(1.15), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.25), Inches(2.90), Inches(0.26), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.27), Inches(0.70), Inches(0.22), step, fs=9, fc=C_W, b=True)
    tb(sl, x + Inches(0.80), Inches(1.27), Inches(2.00), Inches(0.22), title, fs=10, fc=C_W, b=True)
    tb(sl, x + Inches(0.10), Inches(1.57), Inches(2.70), Inches(0.80), desc, fs=9, fc=C_D, ls=1.2)
    if i < 3:
        arrow_r(sl, x + Inches(2.90), Inches(1.60), Inches(0.25), Inches(0.18), color)

# 去噪循环详解区域
rrect(sl, Inches(0.40), Inches(2.60), Inches(8.50), Inches(4.30), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(0.60), Inches(2.65), Inches(5.00), Inches(0.25),
   "去噪循环详解 (Step 3 展开)", fs=11, fc=C_RED, b=True)

# 左侧: 噪声等级时间线
timeline = [
    ("t = 1000", "纯噪声\n随机像素", C_RED),
    ("t ≈ 750", "粗雕阶段\n大致色块", C_ORG),
    ("t ≈ 500", "中雕阶段\n模糊轮廓", C_ORG),
    ("t ≈ 200", "精雕阶段\n细节显现", C_GRN),
    ("t = 1", "完成\n清晰图像", C_GRN),
]

for i, (label, desc, color) in enumerate(timeline):
    y = Inches(3.00 + i * 0.68)
    rrect(sl, Inches(0.70), y, Inches(1.20), Inches(0.50), fc=C_W, lc=color, lw=1)
    tb(sl, Inches(0.75), y + Inches(0.02), Inches(1.10), Inches(0.20), label, fs=9, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.75), y + Inches(0.22), Inches(1.10), Inches(0.25), desc, fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)
    if i < len(timeline) - 1:
        arrow_d(sl, Inches(1.22), y + Inches(0.50), Inches(0.16), Inches(0.18), color)

# 右侧: 3步操作
rrect(sl, Inches(2.20), Inches(3.10), Inches(6.40), Inches(3.60), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(2.30), Inches(3.15), Inches(6.20), Inches(0.22),
   "每步去噪操作（for t = T, T-Δ, ..., 1）", fs=10, fc=C_PRI, b=True)

ops = [
    ("① U-Net 预测噪声",
     "ε_pred = U-Net(x_t, t, text_embedding)\n输入: 当前带噪声图像 x_t + 时间步 t + 文本向量\n输出: 预测的噪声 ε_pred (64×64×4)",
     C_ACC),
    ("② 采样器计算去噪",
     "x_{t-1} = sampler_step(x_t, ε_pred, t)\n如 DPM++ 2M Karras: 智能调整去噪幅度\n噪声高→大步去噪 | 噪声低→精细去噪",
     C_ORG),
    ("③ 更新图像",
     "x_t ← x_{t-1}\n用更干净的图像替换当前图像\n重复直到 t = 1，去噪完成",
     C_GRN),
]

for i, (title, desc, color) in enumerate(ops):
    y = Inches(3.45 + i * 1.05)
    rrect(sl, Inches(2.35), y, Inches(6.10), Inches(0.90), fc=C_W, lc=color, lw=1)
    tb(sl, Inches(2.40), y + Inches(0.03), Inches(5.90), Inches(0.20), title, fs=10, fc=color, b=True)
    tb(sl, Inches(2.40), y + Inches(0.25), Inches(5.90), Inches(0.60), desc, fs=9, fc=C_D, ls=1.15)
    if i < 2:
        arrow_d(sl, Inches(5.35), y + Inches(0.90), Inches(0.16), Inches(0.13), C_ACC)

# 循环箭头标注
tb(sl, Inches(3.20), Inches(6.53), Inches(1.50), Inches(0.20), "⟳ 循环", fs=10, fc=C_RED, b=True, a=PP_ALIGN.CENTER)

# 右侧面板: 推理时间参考
rrect(sl, Inches(9.20), Inches(2.60), Inches(3.60), Inches(4.30), fc=C_W, lc=C_PRI, lw=1)
rect(sl, Inches(9.20), Inches(2.60), Inches(3.60), Inches(0.30), fc=C_PRI)
tb(sl, Inches(9.25), Inches(2.63), Inches(3.50), Inches(0.24), "推理时间参考", fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)

# 时间参考表
table_shape = sl.shapes.add_table(5, 3, Inches(9.35), Inches(3.05), Inches(3.30), Inches(1.70))
table = table_shape.table
table.columns[0].width = Inches(1.10)
table.columns[1].width = Inches(0.80)
table.columns[2].width = Inches(1.40)

t_headers = ["配置", "步数", "耗时"]
t_data = [
    ["快速预览", "10-15", "2-4秒"],
    ["标准质量", "20-30", "5-8秒"],
    ["高质量", "40-50", "3-5秒"],
    ["最高质量", "50-100", "5-10秒"],
]

for j, h in enumerate(t_headers):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size=Pt(8); p.font.bold=True; p.font.color.rgb=C_W; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(t_data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(8); p.font.color.rgb=C_D; p.font.name=FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 推荐配置
mtb(sl, Inches(9.35), Inches(4.90), Inches(3.30), Inches(1.80), [
    ("推荐配置:", True, C_PRI, 9),
    ("DPM++ 2M Karras", False, C_D, 8),
    ("25步 / CFG 7.5", False, C_D, 8),
    ("RTX 3060 / 512×512", False, C_D, 8),
    ("", False, C_D, 8),
    ("采样器选择:", True, C_PRI, 9),
    ("通用: DPM++ 2M Karras", False, C_D, 8),
    ("快速: DPM++ SDE (15步)", False, C_D, 8),
    ("极速: LCM (4-8步)", False, C_D, 8),
], ls=1.1)

chk(sl, "P10 推理过程")


# ============================================================
# P11 CFG详解 (from v9 — already from-scratch)
# ============================================================
print("\nBuilding P11: CFG详解...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "09  CFG（Classifier-Free Guidance）详解", "无分类器引导 — 让AI更听话地按描述画画")

rrect(sl, Inches(0.30), Inches(1.20), Inches(12.70), Inches(0.85), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(0.50), Inches(1.22), Inches(12.30), Inches(0.22),
   "🤔 为什么需要CFG？（紧接推理过程）", fs=12, fc=C_RED, b=True)
tb(sl, Inches(0.50), Inches(1.48), Inches(12.30), Inches(0.50),
   "在推理去噪循环中，U-Net每一步预测噪声时，需要知道\"该往哪个方向去噪\"。没有CFG时，U-Net只知道\"去掉噪声\"，"
   "生成的图像随机且不可控。CFG通过对比\"有文字引导\"和\"无文字引导\"两个预测，放大差异，让生成结果更忠实于文字描述。",
   fs=10, fc=C_D, ls=1.2)

rrect(sl, Inches(0.30), Inches(2.20), Inches(12.70), Inches(1.10), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(2.23), Inches(12.30), Inches(0.22),
   "核心思想：既知道什么是对的，也知道什么不是对的", fs=11, fc=C_PRI, b=True)
mtb(sl, Inches(0.50), Inches(2.50), Inches(12.30), Inches(0.70), [
    "公式: ε_guided = ε_uncond + s × (ε_cond - ε_uncond)",
    "ε_uncond = 没有文本条件时的预测（自由发挥）  |  ε_cond = 有文本条件时的预测（按指令画）  |  s = CFG Scale（通常7-12）",
], fs=10, fc=C_D, ls=1.2)

# CFG Scale表
table_shape = sl.shapes.add_table(5, 3, Inches(0.30), Inches(3.45), Inches(6.50), Inches(1.70))
table = table_shape.table
table.columns[0].width = Inches(1.20)
table.columns[1].width = Inches(2.50)
table.columns[2].width = Inches(2.80)

headers = ["CFG Scale", "效果", "适用场景"]
data = [
    ["1-3", "几乎没有引导，图像\"自由发挥\"", "创意探索，想要意外惊喜"],
    ["5-7", "适中引导，平衡忠实度和创造力", "通用场景，推荐默认值 7"],
    ["7-12", "强引导，高度忠实于文字描述", "需要精确控制，如产品图"],
    ["15-20+", "过强引导，过度饱和出现伪影", "一般不推荐"],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# CFG工作原理图
rrect(sl, Inches(7.10), Inches(3.45), Inches(5.90), Inches(2.80), fc=C_W, lc=C_ACC, lw=1)
tb(sl, Inches(7.30), Inches(3.50), Inches(5.50), Inches(0.25), "CFG 工作原理图", fs=11, fc=C_PRI, b=True)

rrect(sl, Inches(7.30), Inches(3.90), Inches(2.20), Inches(0.50), fc=C_BG_RED, lc=C_RED, lw=1)
tb(sl, Inches(7.30), Inches(3.95), Inches(2.20), Inches(0.40),
   "无条件生成（自由发挥）\nε_uncond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)

rrect(sl, Inches(10.40), Inches(3.90), Inches(2.20), Inches(0.50), fc=C_BG_GRN, lc=C_GRN, lw=1)
tb(sl, Inches(10.40), Inches(3.95), Inches(2.20), Inches(0.40),
   "有条件生成（按指令画）\nε_cond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)

arrow_d(sl, Inches(8.35), Inches(4.40), Inches(0.15), Inches(0.15), C_ORG)
arrow_d(sl, Inches(11.45), Inches(4.40), Inches(0.15), Inches(0.15), C_ORG)

rrect(sl, Inches(9.30), Inches(4.60), Inches(1.60), Inches(0.40), fc=C_BG_ORG, lc=C_ORG, lw=1)
tb(sl, Inches(9.30), Inches(4.63), Inches(1.60), Inches(0.35),
   "计算差异\nε_cond - ε_uncond", fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0, b=True)

arrow_d(sl, Inches(10.05), Inches(5.00), Inches(0.15), Inches(0.12), C_ORG)
rrect(sl, Inches(9.30), Inches(5.15), Inches(1.60), Inches(0.35), fc=C_W, lc=C_ORG, lw=1)
tb(sl, Inches(9.30), Inches(5.18), Inches(1.60), Inches(0.30),
   "× CFG Scale (s)", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)

arrow_d(sl, Inches(10.05), Inches(5.50), Inches(0.15), Inches(0.12), C_GRN)
rrect(sl, Inches(9.00), Inches(5.65), Inches(2.20), Inches(0.45), fc=C_BG_GRN, lc=C_GRN, lw=2)
tb(sl, Inches(9.00), Inches(5.68), Inches(2.20), Inches(0.40),
   "= 引导后的高质量结果\nε_guided", fs=9, fc=C_D, b=True, a=PP_ALIGN.CENTER, ls=1.0)

tb(sl, Inches(7.30), Inches(5.75), Inches(1.60), Inches(0.30),
   "← + ε_uncond (加回\n     无条件结果)", fs=8, fc=C_G, a=PP_ALIGN.CENTER, ls=1.0)
arrow_r(sl, Inches(8.95), Inches(5.85), Inches(0.10), Inches(0.10), C_ACC)

# 底部比喻
rrect(sl, Inches(0.30), Inches(5.55), Inches(6.50), Inches(1.35), fc=C_LT, lc=C_ACC, lw=1)
mtb(sl, Inches(0.50), Inches(5.58), Inches(6.10), Inches(1.25), [
    "比喻: CFG像严厉又公正的艺术指导。",
    "每画一笔问两个问题:",
    "  \"按客户要求画了吗？\"",
    "  \"不管要求你会怎么画？\"",
    "差异放大(×CFG Scale)，强制遵循要求。",
    "没有CFG→模糊偏离 | CFG过大→过度刻意",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P11 CFG详解")


# ============================================================
# P12 参数优化+Prompt (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P12: 参数优化+Prompt (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "12-13  参数优化 + Prompt 技巧", "实战指南 — 让AI画出你想要的图像")

# 4个参数卡片
param_cards = [
    ("Steps(去噪步数)", C_ACC, "推荐: 20-50步\nDPM++ 2M用25步\n超过30步提升递减\n比喻: 打磨家具20遍已很好"),
    ("CFG Scale(引导系数)", C_GRN, "推荐: 5-12\n通用7-8 精确10-12\n太低偏离描述\n太高过度饱和\n比喻: 老师批改严格程度"),
    ("Seed(随机种子)", C_ORG, "任意整数\n相同种子+参数=相同图像\n不同种子=完全不同图像\n比喻: 掷骰子记住数字可复现"),
    ("Sampler(采样器)", C_PUR, "推荐: DPM++ 2M Karras\n快速: DPM++ SDE(15步)\n极速: LCM(4-8步)\n不同算法速度质量取舍"),
]

for i, (title, color, content) in enumerate(param_cards):
    x = Inches(0.50 + i * 3.20)
    rrect(sl, x, Inches(1.30), Inches(3.00), Inches(1.80), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.30), Inches(3.00), Inches(0.28), fc=color)
    tb(sl, x + Inches(0.05), Inches(1.32), Inches(2.90), Inches(0.24), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.10), Inches(1.62), Inches(2.80), Inches(1.40), content, fs=9, fc=C_D, ls=1.2)

# Prompt技巧区域
rrect(sl, Inches(0.50), Inches(3.30), Inches(12.30), Inches(3.60), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.70), Inches(3.35), Inches(5.00), Inches(0.25),
   "Prompt 技巧 — 给AI画师下订单的艺术", fs=11, fc=C_PRI, b=True)

tb(sl, Inches(0.70), Inches(3.65), Inches(6.00), Inches(0.20),
   "Prompt基本结构:", fs=10, fc=C_D, b=True)
tb(sl, Inches(0.70), Inches(3.85), Inches(12.00), Inches(0.20),
   "[主体描述] + [环境/背景] + [风格/艺术参考] + [画质/技术参数] + [负面提示词]",
   fs=9, fc=C_PRI, b=True, fn="Consolas")

prompt_tips = [
    ("1 主语前置", "A fluffy orange tabby cat sitting on a vintage sofa (具体)", "2 具体数值", "cinematic lighting, 85mm lens, f/1.8, shallow depth of field"),
    ("3 风格关键词", "photorealistic / oil painting / anime / concept art / dramatic", "4 权重调节", "(red dress:1.3)增加 / (blue sky:0.8)降低关键词权重"),
    ("5 负面提示词", "ugly, blurry, low quality, deformed, bad hands, watermark, text", "6 排列顺序", "最重要(主体) → 环境 → 风格 → 画质，越靠前影响越大"),
]

for ri, (l1_title, l1_desc, l2_title, l2_desc) in enumerate(prompt_tips):
    y = Inches(4.15 + ri * 0.80)
    # Left tip
    rrect(sl, Inches(0.70), y, Inches(6.00), Inches(0.65), fc=C_W, lc=C_ACC, lw=1)
    tb(sl, Inches(0.78), y + Inches(0.03), Inches(1.50), Inches(0.22), l1_title, fs=9, fc=C_ACC, b=True)
    tb(sl, Inches(2.30), y + Inches(0.03), Inches(4.30), Inches(0.58), l1_desc, fs=9, fc=C_D, fn="Consolas", ls=1.0)
    # Right tip
    rrect(sl, Inches(7.00), y, Inches(6.00), Inches(0.65), fc=C_W, lc=C_ACC, lw=1)
    tb(sl, Inches(7.08), y + Inches(0.03), Inches(1.50), Inches(0.22), l2_title, fs=9, fc=C_ACC, b=True)
    tb(sl, Inches(8.60), y + Inches(0.03), Inches(4.30), Inches(0.58), l2_desc, fs=9, fc=C_D, fn="Consolas", ls=1.0)

chk(sl, "P12 参数优化+Prompt")


# ============================================================
# P13 应用+LoRA+ControlNet (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P13: 应用+LoRA+ControlNet (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "14-16  应用场景 + LoRA + ControlNet", "从基础应用到高级控制")

# 5个应用场景卡片
apps = [
    ("概念艺术", "几秒生成数十种设计方案\n游戏/电影/广告大量使用\n给设计师配备灵感无限的外脑", C_ACC),
    ("游戏影视", "角色/场景/纹理素材\nControlNet控制姿态+LoRA统一风格\n批量生成风格一致资产", C_ORG),
    ("电商广告", "产品图合成到各种场景\n一双鞋→海滩/雪山/城市\n成本几乎为零", C_GRN),
    ("建筑室内", "快速生成设计效果图\n现代极简客厅落地窗\n用于早期方案沟通", C_PUR),
    ("个人创作", "头像/小说插图/社交媒体\n完全免费开源可本地运行\n不受任何使用限制", C_TEL),
]

for i, (title, desc, color) in enumerate(apps):
    x = Inches(0.50 + i * 2.55)
    rrect(sl, x, Inches(1.30), Inches(2.35), Inches(2.20), fc=C_W, lc=color, lw=1)
    tb(sl, x + Inches(0.05), Inches(1.35), Inches(2.20), Inches(0.22), title, fs=10, fc=color, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.08), Inches(1.60), Inches(2.10), Inches(1.80), desc, fs=9, fc=C_D, ls=1.2)

# LoRA
rrect(sl, Inches(0.50), Inches(3.70), Inches(6.00), Inches(3.20), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_ACC, lw=1)
tb(sl, Inches(0.70), Inches(3.75), Inches(5.50), Inches(0.25),
   "LoRA（Low-Rank Adaptation）— 给画家添加专项风格手册", fs=10, fc=C_PRI, b=True)
mtb(sl, Inches(0.70), Inches(4.05), Inches(5.60), Inches(2.70), [
    "核心思想: 不修改原始参数，在旁边挂载一组很小的额外参数",
    "就像给汽车加装外挂导航仪，不需要改装引擎就能获得新功能",
    "",
    "数学原理: W' = W + B × A (A: 1024×r, B: r×1024, r通常4-64)",
    "r=16时只需3.2万参数，相比原始100万减少97%",
    "",
    "训练数据: 10-50张图片 | 训练时间: 30分钟-2小时(RTX3060)",
    "文件大小: 10-200MB (主模型4GB的0.25%-5%)",
    "可以叠加多个LoRA: 风格LoRA + 角色LoRA + 场景LoRA",
    "LoRA系数0.5轻度/0.75中等/1.0完整风格影响",
], fs=9, fc=C_D, ls=1.15)

# ControlNet
rrect(sl, Inches(6.80), Inches(3.70), Inches(6.00), Inches(3.20), fc=C_BG_GRN, lc=C_GRN, lw=1)
tb(sl, Inches(7.00), Inches(3.75), Inches(5.50), Inches(0.25),
   "ControlNet — 精确控制构图和姿态", fs=10, fc=C_GRN, b=True)
mtb(sl, Inches(7.00), Inches(4.05), Inches(5.60), Inches(2.70), [
    "核心问题: 只用文字无法精确控制构图和结构",
    "解决方案: 通过额外条件图像(草图/边缘/深度/姿态)精确控制",
    "",
    "架构: 复制一份U-Net编码器作为控制分支，通过零卷积注入原始U-Net",
    "零卷积初始权重为0，训练开始时不影响原始模型，逐步学习注入方式",
    "",
    "8种条件类型:",
    "Canny Edge(边缘) | Depth(深度) | OpenPose(人体姿态) | Normal Map(法线)",
    "Scribble(简笔) | Segmentation(分割) | Lineart(线稿) | Tile(分块)",
    "",
    "比喻: 给画家配一个精确的构图模板，文字控制画什么，ControlNet控制画在哪里",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P13 应用+LoRA+ControlNet")


# ============================================================
# P14 SD3+Flux (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P14: SD3+Flux (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "17-18  SD3/SDXL 最新进展 + Flux 等新模型对比", "2026年AI图像生成格局")

# SD版本对比表
table_shape = sl.shapes.add_table(8, 4, Inches(0.50), Inches(1.30), Inches(6.00), Inches(2.56))
table = table_shape.table
table.columns[0].width = Inches(1.30)
table.columns[1].width = Inches(1.40)
table.columns[2].width = Inches(1.60)
table.columns[3].width = Inches(1.70)

sd_headers = ["特性", "SD 1.5", "SDXL", "SD3"]
sd_data = [
    ["架构", "U-Net", "U-Net(更大)", "MMDiT(全新)"],
    ["参数量", "~1B", "~3.5B", "800M~8B"],
    ["文本编码器", "CLIP L/14", "CLIP+OpenCLIP", "CLIP-L+CLIP-G+T5-XXL"],
    ["默认分辨率", "512x512", "1024x1024", "1024x1024"],
    ["训练方法", "DDPM噪声预测", "DDPM噪声预测", "Flow Matching"],
    ["VAE通道", "4通道", "4通道(更大)", "16通道"],
    ["文字渲染", "较差", "一般", "显著提升"],
]

for j, h in enumerate(sd_headers):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=C_W; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(sd_data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(8); p.font.color.rgb=C_D; p.font.name=FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# SD3 核心创新
rrect(sl, Inches(6.80), Inches(1.30), Inches(6.00), Inches(2.20), fc=RGBColor(0xF5,0xEE,0xF8), lc=C_PUR, lw=1)
tb(sl, Inches(7.00), Inches(1.35), Inches(5.50), Inches(0.25),
   "SD3 核心创新", fs=11, fc=C_PUR, b=True)
mtb(sl, Inches(7.00), Inches(1.65), Inches(5.60), Inches(1.80), [
    "• MMDiT: 纯Transformer架构，双流注意力(图像+文本双向交流)",
    "• 三编码器: CLIP-L(1.2亿)+CLIP-G(3.5亿)+T5-XXL(46亿)=50亿参数仅用于文本理解",
    "• Flow Matching: 从噪声分布到目标分布的直接路径，比DDPM更高效",
    "• U-Net是单向广播系统; MMDiT是双向对讲系统，画师能根据画面反问客户",
], fs=9, fc=C_D, ls=1.2)

# 模型对比表
tb(sl, Inches(0.50), Inches(3.70), Inches(5.00), Inches(0.25),
   "Flux + 2026年主要模型对比", fs=11, fc=C_PRI, b=True)

table_shape = sl.shapes.add_table(7, 6, Inches(0.50), Inches(4.00), Inches(12.30), Inches(2.24))
table = table_shape.table
table.columns[0].width = Inches(2.20)
table.columns[1].width = Inches(2.00)
table.columns[2].width = Inches(1.50)
table.columns[3].width = Inches(1.80)
table.columns[4].width = Inches(1.80)
table.columns[5].width = Inches(3.00)

m_headers = ["模型", "开发者", "参数量", "文字渲染", "开源", "ELO评分"]
m_data = [
    ["GPT Image 1.5", "OpenAI", "N/A", "最强", "闭源", "1264"],
    ["Gemini 3 Pro Image", "Google", "N/A", "强", "闭源", "1241"],
    ["Midjourney V7/V8", "Midjourney", "N/A", "中", "闭源", "N/A"],
    ["Flux 2 Dev", "Black Forest Labs", "~12B", "优秀", "开放", "~1200"],
    ["SD3.5 Large", "Stability AI", "8B", "较好", "部分", "~1170"],
    ["Hunyuan Image 3.0", "腾讯", "N/A", "中文强", "部分", "~1130"],
]

for j, h in enumerate(m_headers):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size=Pt(9); p.font.bold=True; p.font.color.rgb=C_W; p.font.name=FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI

for i, row_data in enumerate(m_data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(8); p.font.color.rgb=C_D; p.font.name=FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 底部总结
rrect(sl, Inches(0.50), Inches(6.30), Inches(12.30), Inches(0.60), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.70), Inches(6.33), Inches(11.80), Inches(0.50),
   "Flux由SD原始团队创建，文字渲染和手部生成显著改善。Flux 2 Dev与Max的ELO差距仅19分，免费开源已接近商业模型质量。\n"
   "2026年格局: 质量差距缩小，竞争转向专业化。多模型策略(根据用途选择不同模型)成为专业用户主流。",
   fs=9, fc=C_D, ls=1.2)

chk(sl, "P14 SD3+Flux")


# ============================================================
# P15 趋势+总结 (REBUILD — was copy_slide)
# ============================================================
print("\nBuilding P15: 趋势+总结 (from scratch)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "19  2026年AI图像生成动态 + 总结", "7大趋势与未来展望")

# 7大趋势卡片 (4+3布局)
trends = [
    ("1 质量差距缩小", "前9名ELO差距仅117分\n竞争从谁更好看转向谁在特定场景更好\n多模型策略成为主流工作方式", C_ACC),
    ("2 4K输出成标配", "主流模型普遍支持4K\n多档分辨率和多宽高比原生生成\n横幅广告等特殊尺寸不再需要裁剪", C_GRN),
    ("3 实时信息整合", "Grounding功能实时搜索互联网\n要求生成最新款iPhone会先搜索确认\n解决AI知识截止问题", C_ORG),
    ("4 角色一致性突破", "最多5个角色+14个物体保持一致\n可生成完整故事板\n漫画创作/品牌营销意义重大", C_PUR),
    ("5 Flash模型改变工作流", "极低价格接近Pro级质量\n迭代式AI图像工作流首次经济可行\n从一次生成变成快速迭代精雕细琢", C_RED),
    ("6 开源vs闭源博弈", "Flux 2 Dev与Max差距仅19分\n免费开源已接近顶级商业质量\n对数据敏感企业是重大利好", C_TEL),
    ("7 视频生成融合", "LTX 2.3(220亿参数,带同步音频4K视频)\nHelios(消费级GPU实时60秒视频)\n图像和视频技术边界快速模糊", C_TEL),
]

for i, (title, desc, color) in enumerate(trends):
    if i < 4:
        x = Inches(0.50 + i * 3.20)
        y = Inches(1.30)
        h = Inches(2.20)
    else:
        x = Inches(0.50 + (i - 4) * 4.20)
        y = Inches(3.70)
        h = Inches(2.00)

    rrect(sl, x, y, Inches(3.00) if i < 4 else Inches(3.00), h, fc=C_W, lc=color, lw=1)
    rect(sl, x, y, Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x + Inches(0.05), y + Inches(0.02), Inches(2.90), Inches(0.20), title, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.08), y + Inches(0.30), Inches(2.80), h - Inches(0.35), desc, fs=9, fc=C_D, ls=1.2)

# 底部总结
rrect(sl, Inches(0.50), Inches(5.90), Inches(12.30), Inches(1.00), fc=C_PRI, lc=C_PRI, lw=1)
tb(sl, Inches(0.70), Inches(5.95), Inches(11.80), Inches(0.90),
   "2026年AI图像生成就像从功能手机时代进入了智能手机时代。\n"
   "功能手机时代(2022-2023)比谁屏幕更大、待机更长。智能手机时代(2026)比生态——应用商店、多任务处理、云端整合。\n"
   "最终胜出的可能不是某一个模型，而是能够灵活调度多种模型的工作流平台。",
   fs=10, fc=C_W, ls=1.3)

chk(sl, "P15 趋势+总结")


# ============================================================
# Save
# ============================================================
print(f"\nSaving to: {OUT}")
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)

print(f"\n{'='*60}")
print(f"RESULT: {len(prs.slides)} slides — ALL from scratch, zero copy_slide")
print(f"{'='*60}")
print("Slide order:")
print("  P1  = 封面 (from scratch)")
print("  P2  = 目录 (REBUILT)")
print("  P3  = Diffusion基础 (REBUILT)")
print("  P4  = 整体架构概览 (from scratch)")
print("  P5  = CLIP详解 (from scratch)")
print("  P6  = VAE详解 (from scratch)")
print("  P7  = U-Net架构 (from scratch)")
print("  P8  = 三大组件协作 (REBUILT)")
print("  P9  = 训练过程 (from scratch)")
print("  P10 = 推理过程 (REBUILT)")
print("  P11 = CFG详解 (from scratch)")
print("  P12 = 参数优化+Prompt (REBUILT)")
print("  P13 = 应用+LoRA+ControlNet (REBUILT)")
print("  P14 = SD3+Flux (REBUILT)")
print("  P15 = 趋势+总结 (REBUILT)")

# Final verification
print(f"\n{'='*60}")
print("FINAL VERIFICATION:")
print(f"{'='*60}")
all_ok = True
for i, slide in enumerate(prs.slides):
    if not chk(slide, f"Slide {i+1}"):
        all_ok = False

if all_ok:
    print(f"\n✅ ALL {len(prs.slides)} slides passed verification!")
else:
    print(f"\n⚠️ Some slides have issues — review above")

print(f"\nDone! Output: {OUT}")
