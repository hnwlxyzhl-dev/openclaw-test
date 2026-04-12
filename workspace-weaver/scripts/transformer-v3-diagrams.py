#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer v3 Architecture Diagrams Generator
Generates 14 pages of architecture diagrams for the Transformer PPT.
Each page is a blank layout with diagrams on the left and space for text on the right.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# ============================================================
# Configuration
# ============================================================
SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_R = 13.2
SAFE_B = 7.3

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xC7)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
C_DARK = RGBColor(0x34, 0x49, 0x5E)
C_BG_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
C_BG_RED = RGBColor(0xFD, 0xED, 0xEC)
C_BG_GREEN = RGBColor(0xE8, 0xF8, 0xF5)
C_BG_ORANGE = RGBColor(0xFE, 0xF9, 0xE7)

FONT_ZH = "Microsoft YaHei"
FONT_EN = "Arial"
FONT_CODE = "Consolas"

# ============================================================
# Utility Functions
# ============================================================
def chk(l, t, w, h, label=""):
    """Assert safe boundaries."""
    assert l + w <= SAFE_R + 0.01, f"Right overflow {label}: {l+w:.2f} > {SAFE_R}"
    assert t + h <= SAFE_B + 0.01, f"Bottom overflow {label}: {t+h:.2f} > {SAFE_B}"
    assert w > 0 and h > 0, f"Zero size {label}: {w}x{h}"

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font_name=FONT_ZH,
                anchor=MSO_ANCHOR.TOP, italic=False):
    """Add a text box with boundary check."""
    chk(left, top, width, height, f"TB:{text[:15]}")
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tb

def add_shape_box(slide, left, top, width, height, text="",
                  fill_color=None, line_color=None, line_width=Pt(1),
                  text_color=C_WHITE, font_size=11, bold=True,
                  font_name=FONT_ZH, align=PP_ALIGN.CENTER,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, anchor=MSO_ANCHOR.MIDDLE):
    """Add a shape box with text."""
    chk(left, top, width, height, f"SB:{text[:15]}")
    shape = slide.shapes.add_shape(shape_type,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        try:
            tf.paragraphs[0].alignment = align
        except:
            pass
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    return shape

def add_connector(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(2)):
    """Add a straight connector line."""
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = width
    return c

def add_arrow_down(slide, cx, y1, y2, color=C_GRAY, width=Pt(2)):
    """Add a downward arrow."""
    return add_connector(slide, cx, y1, cx, y2, color, width)

def add_arrow_right(slide, x1, y, x2, color=C_GRAY, width=Pt(2)):
    """Add a rightward arrow."""
    return add_connector(slide, x1, y, x2, y, color, width)

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a plain rectangle (no text)."""
    chk(left, top, width, height, "rect")
    return add_shape_box(slide, left, top, width, height, "",
                         fill_color=fill_color, line_color=line_color)

def add_multiline_tb(slide, left, top, width, height, lines,
                     font_size=10, color=C_TEXT, line_spacing=1.1,
                     font_name=FONT_ZH, bold_first=False):
    """Add multi-line text box."""
    chk(left, top, width, height, "MTB")
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(1)
        p.space_before = Pt(0)
    return tb

def add_label(slide, left, top, text, font_size=9, color=C_TEXT, bold=False):
    """Add a small label text box."""
    return add_textbox(slide, left, top, 3.0, 0.25, text,
                       font_size=font_size, color=color, bold=bold)

def add_page_title(slide, title, page_num=""):
    """Add page title."""
    add_textbox(slide, 0.3, 0.15, 8.0, 0.5, title,
                font_size=22, bold=True, color=C_TITLE)
    if page_num:
        add_textbox(slide, 12.0, 0.15, 1.0, 0.5, page_num,
                    font_size=14, color=C_GRAY, align=PP_ALIGN.RIGHT)

def add_zone_box(slide, left, top, width, height, label,
                 fill_color, line_color, label_color=C_WHITE):
    """Add a zone/background box with a label at top."""
    chk(left, top, width, height, f"zone:{label}")
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = line_color
    box.line.width = Pt(2)
    # Label at top
    lbl = add_textbox(slide, left + 0.05, top + 0.02, width - 0.1, 0.3,
                      label, font_size=11, bold=True, color=label_color,
                      align=PP_ALIGN.CENTER, font_name=FONT_EN)
    return box

# ============================================================
# Create Presentation
# ============================================================
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
blank = prs.slide_layouts[6]

# Track all shapes for validation
all_shapes_log = []

def log_shape(label, left, top, width, height):
    all_shapes_log.append((label, left, top, width, height))

# ============================================================
# PAGE 3: RNN Processing Flow
# ============================================================
print("Creating Page 3: RNN Processing Flow...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "RNN/LSTM: Sequential Processing")

# Words along timeline
words = ["The", "cat", "sat", "on", "the", "mat"]
x_start = 0.5
y_words = 1.2
for i, w in enumerate(words):
    x = x_start + i * 1.05
    add_shape_box(slide, x, y_words, 0.85, 0.45, w,
                  fill_color=C_ENCODER, font_size=11)
    log_shape(f"word_{w}", x, y_words, 0.85, 0.45)

# RNN cells below words
for i, w in enumerate(words):
    x = x_start + i * 1.05
    add_shape_box(slide, x + 0.05, y_words + 0.6, 0.75, 0.4, "RNN",
                  fill_color=C_ATTENTION, font_size=9)
    log_shape(f"rnn_{i}", x + 0.05, y_words + 0.6, 0.75, 0.4)

# Arrows between RNN cells
for i in range(len(words) - 1):
    x1 = x_start + i * 1.05 + 0.85
    x2 = x_start + (i + 1) * 1.05
    y_mid = y_words + 0.8
    add_arrow_right(slide, x1, y_mid, x2, color=C_DARK, width=Pt(2))

# Arrows from words to RNN cells
for i in range(len(words)):
    x = x_start + i * 1.05 + 0.42
    add_arrow_down(slide, x, y_words + 0.45, y_words + 0.6, color=C_GRAY, width=Pt(1))

# Pain point 1: Sequential - clock icon area
y_pp = 2.5
add_shape_box(slide, 0.5, y_pp, 1.8, 0.4, "1 Sequential",
              fill_color=C_DECODER, font_size=9)
add_textbox(slide, 0.5, y_pp + 0.45, 1.8, 0.6,
            "Must wait for\nprevious word", font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER)
# Red dashed line bracket under first two words
add_connector(slide, 0.5, y_pp - 0.15, 2.4, y_pp - 0.15, color=C_DECODER, width=Pt(1))

# Pain point 2: Long distance forgetting
add_shape_box(slide, 2.8, y_pp, 1.8, 0.4, "2 Forget",
              fill_color=C_DECODER, font_size=9)
add_textbox(slide, 2.8, y_pp + 0.45, 1.8, 0.6,
            "Info fades over\ndistance (gradient\nvanishing)", font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER)
# Fading line from "The" to "mat"
for j in range(5):
    alpha = max(30, 255 - j * 50)
    c = RGBColor(0x34, 0x98, min(0xDB, 0x50 + j * 30))
    add_connector(slide, x_start + j * 1.05 + 0.42, y_words + 1.15,
                  x_start + (j+1) * 1.05 + 0.42, y_words + 1.15,
                  color=c, width=Pt(3 - j * 0.4))

# Pain point 3: Bottleneck
add_shape_box(slide, 5.1, y_pp, 1.8, 0.4, "3 Bottleneck",
              fill_color=C_DECODER, font_size=9)
add_textbox(slide, 5.1, y_pp + 0.45, 1.8, 0.6,
            "All info compressed\ninto one vector", font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER)
# Funnel shape approximation
add_shape_box(slide, 5.5, y_pp - 0.25, 1.0, 0.2, "", fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_shape_box(slide, 5.8, y_pp - 0.05, 0.4, 0.12, "", fill_color=C_ATTENTION)

# Bottom comparison
y_compare = 4.0
add_textbox(slide, 0.5, y_compare, 6.5, 0.35, "RNN vs Transformer Comparison",
            font_size=12, bold=True, color=C_TITLE)

# RNN box
add_shape_box(slide, 0.5, y_compare + 0.4, 3.0, 1.2, "",
              fill_color=C_LIGHT_RED, line_color=C_DECODER, line_width=Pt(2))
add_textbox(slide, 0.6, y_compare + 0.45, 2.8, 0.3, "RNN: Queue Speaking",
            font_size=10, bold=True, color=C_DECODER, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.6, y_compare + 0.8, 2.8, 0.7,
            "Person A speaks -> Person B speaks\n-> Person C speaks...\nEach forgets after speaking",
            font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER)

# Transformer box
add_shape_box(slide, 3.8, y_compare + 0.4, 3.0, 1.2, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN, line_width=Pt(2))
add_textbox(slide, 3.9, y_compare + 0.45, 2.8, 0.3, "Transformer: Round Table",
            font_size=10, bold=True, color=C_FFN, align=PP_ALIGN.CENTER)
add_textbox(slide, 3.9, y_compare + 0.8, 2.8, 0.7,
            "Everyone speaks simultaneously\nEveryone hears everyone\nNo one forgets",
            font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER)

# Key stats at bottom
y_stats = 5.8
stats = [
    ("8 GPUs", "Transformer\n12 hours", C_FFN),
    ("Many GPUs", "LSTM\nSeveral days", C_DECODER),
]
for i, (hw, time, color) in enumerate(stats):
    x = 1.0 + i * 2.5
    add_shape_box(slide, x, y_stats, 1.0, 0.6, hw, fill_color=color, font_size=10)
    add_textbox(slide, x + 1.1, y_stats + 0.05, 1.3, 0.5, time,
                font_size=9, color=C_TEXT)

add_textbox(slide, 0.5, y_stats + 0.7, 6.0, 0.3,
            "Speed difference: 10x+ (Transformer wins)",
            font_size=10, bold=True, color=C_TITLE)

# Right side: reserved for text (7.5 to 13.0)
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 3 done.")

# ============================================================
# PAGE 4: Attention Intuition Visualization
# ============================================================
print("Creating Page 4: Attention Intuition...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Attention Intuition: Round Table vs Queue")

# Sentence words with attention visualization
sentence = ["The", "cat", "sat", "on", "the", "mat", ",", "and", "it", "was", "happy"]
x_start = 0.3
y_sent = 1.0
box_w = 0.52
box_h = 0.4
gap = 0.08
it_idx = sentence.index("it")
cat_idx = sentence.index("cat")

for i, w in enumerate(sentence):
    x = x_start + i * (box_w + gap)
    if w == "it":
        fc = C_ENCODER
    elif w == "cat":
        fc = C_FFN
    else:
        fc = C_LIGHT_GRAY
    tc = C_WHITE if (w == "it" or w == "cat") else C_DARK
    add_shape_box(slide, x, y_sent, box_w, box_h, w, fill_color=fc, text_color=tc, font_size=9)

# Attention weight bars below
weights_map = {
    "The": 0.05, "cat": 0.42, "sat": 0.08, "on": 0.03,
    "the": 0.02, "mat": 0.06, ",": 0.01, "and": 0.03,
    "it": 0.10, "was": 0.08, "happy": 0.12
}
y_bar = 1.6
bar_max_w = 2.5
add_textbox(slide, 0.3, y_bar - 0.3, 6.0, 0.3,
            "Attention weights from \"it\" to each word:",
            font_size=10, bold=True, color=C_TITLE)

for i, w in enumerate(sentence):
    x = x_start + i * (box_w + gap)
    w_val = weights_map.get(w, 0.05)
    bar_w = w_val * bar_max_w / 0.45  # normalize to max
    bar_color = C_FFN if w == "cat" else (C_ENCODER if w_val > 0.08 else C_GRAY)
    add_shape_box(slide, x, y_bar + 0.1, box_w, 0.25, "", fill_color=bar_color)
    # Weight text
    add_textbox(slide, x, y_bar + 0.38, box_w, 0.2, f"{w_val:.0%}",
                font_size=7, color=C_TEXT, align=PP_ALIGN.CENTER)

# Highlight cat
add_textbox(slide, 0.3, y_bar + 0.6, 5.0, 0.3,
            "\"cat\" gets highest weight (42%) - model learned \"it\" refers to \"cat\"",
            font_size=10, bold=True, color=C_FFN)

# Comparison section
y_cmp = 3.2
# Left: RNN queue
add_zone_box(slide, 0.3, y_cmp, 3.2, 2.8, "RNN: Queue Speaking",
             C_LIGHT_RED, C_DECODER)
add_multiline_tb(slide, 0.5, y_cmp + 0.4, 2.8, 2.2, [
    "Person A: speaks...",
    "  (forgets everything)",
    "Person B: speaks...",
    "  (forgets everything)",
    "Person C: speaks...",
    "  (forgets everything)",
    "",
    "x Sequential, slow",
    "x Long-range forgotten",
], font_size=9, color=C_TEXT)

# Right: Transformer round table
add_zone_box(slide, 3.8, y_cmp, 3.2, 2.8, "Transformer: Round Table",
             C_LIGHT_GREEN, C_FFN)
add_multiline_tb(slide, 4.0, y_cmp + 0.4, 2.8, 2.2, [
    "Everyone at round table:",
    "  All speak simultaneously",
    "  All hear everyone",
    "  Focus on relevant speakers",
    "",
    "check Fully parallel (fast)",
    "check No distance limit",
    "check Rich information flow",
], font_size=9, color=C_TEXT)

# Three problems solved
y_solved = 6.3
add_textbox(slide, 0.3, y_solved, 6.5, 0.3,
            "One idea solves all three RNN problems:",
            font_size=10, bold=True, color=C_TITLE)
problems_solved = [
    ("Parallel", "Solves sequential"),
    ("Direct", "Solves forgetting"),
    ("Full info", "Solves bottleneck"),
]
for i, (key, desc) in enumerate(problems_solved):
    x = 0.3 + i * 2.2
    add_shape_box(slide, x, y_solved + 0.35, 2.0, 0.4, f"{key}: {desc}",
                  fill_color=C_FFN, font_size=8)

# Right side reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 4 done.")

# ============================================================
# PAGE 5: Transformer Overall Architecture (MOST IMPORTANT!)
# ============================================================
print("Creating Page 5: Transformer Overall Architecture...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Transformer Architecture: Encoder + Decoder")

# Input Processing Zone (left)
add_zone_box(slide, 0.2, 0.8, 1.6, 5.8, "Input", C_BG_BLUE, C_ENCODER, label_color=C_ENCODER)
add_shape_box(slide, 0.35, 1.3, 1.3, 0.45, "Input Embedding",
              fill_color=C_EMBED, font_size=9)
add_arrow_down(slide, 1.0, 1.75, 1.95, color=C_ENCODER)
add_shape_box(slide, 0.35, 1.95, 1.3, 0.45, "Pos Encoding",
              fill_color=C_FFN, font_size=9)
add_arrow_down(slide, 1.0, 2.4, 2.8, color=C_ENCODER)
# Input example
add_textbox(slide, 0.35, 2.6, 1.3, 0.4, "\"I love you\"",
            font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER, italic=True)

# Encoder Zone (middle-left)
enc_x = 2.1
enc_w = 2.3
enc_y = 0.8
enc_h = 5.8
add_zone_box(slide, enc_x, enc_y, enc_w, enc_h, "Encoder x6",
             C_BG_BLUE, C_ENCODER, label_color=C_ENCODER)

# Encoder internal: one layer detailed
ey = 1.4
add_shape_box(slide, enc_x + 0.15, ey, enc_w - 0.3, 0.5,
              "Multi-Head Self-Attention", fill_color=C_ENCODER, font_size=9)
add_arrow_down(slide, enc_x + enc_w/2, ey + 0.5, ey + 0.65, color=C_ENCODER)
add_shape_box(slide, enc_x + 0.15, ey + 0.65, enc_w - 0.3, 0.35,
              "Add & Layer Norm", fill_color=C_FFN, font_size=9)
add_arrow_down(slide, enc_x + enc_w/2, ey + 1.0, ey + 1.35, color=C_ENCODER)
add_shape_box(slide, enc_x + 0.15, ey + 1.35, enc_w - 0.3, 0.5,
              "Feed Forward", fill_color=C_ATTENTION, font_size=9)
add_arrow_down(slide, enc_x + enc_w/2, ey + 1.85, ey + 2.2, color=C_ENCODER)
add_shape_box(slide, enc_x + 0.15, ey + 2.2, enc_w - 0.3, 0.35,
              "Add & Layer Norm", fill_color=C_FFN, font_size=9)

# Repeat indicator
add_textbox(slide, enc_x + 0.2, ey + 2.7, enc_w - 0.4, 0.5,
            "... (repeat x6)", font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# Decoder Zone (middle-right)
dec_x = 4.8
dec_w = 2.6
dec_y = 0.8
dec_h = 5.8
add_zone_box(slide, dec_x, dec_y, dec_w, dec_h, "Decoder x6",
             C_BG_RED, C_DECODER, label_color=C_DECODER)

# Decoder internal
dy = 1.3
add_shape_box(slide, dec_x + 0.1, dy, dec_w - 0.2, 0.45,
              "Masked Self-Attention", fill_color=C_ENCODER, font_size=9)
add_arrow_down(slide, dec_x + dec_w/2, dy + 0.45, dy + 0.55, color=C_DECODER)
add_shape_box(slide, dec_x + 0.1, dy + 0.55, dec_w - 0.2, 0.3,
              "Add & Norm", fill_color=C_FFN, font_size=9)
add_arrow_down(slide, dec_x + dec_w/2, dy + 0.85, dy + 0.95, color=C_DECODER)
add_shape_box(slide, dec_x + 0.1, dy + 0.95, dec_w - 0.2, 0.45,
              "Cross-Attention", fill_color=C_EMBED, font_size=9)
add_arrow_down(slide, dec_x + dec_w/2, dy + 1.4, dy + 1.5, color=C_DECODER)
add_shape_box(slide, dec_x + 0.1, dy + 1.5, dec_w - 0.2, 0.3,
              "Add & Norm", fill_color=C_FFN, font_size=9)
add_arrow_down(slide, dec_x + dec_w/2, dy + 1.8, dy + 1.9, color=C_DECODER)
add_shape_box(slide, dec_x + 0.1, dy + 1.9, dec_w - 0.2, 0.45,
              "Feed Forward", fill_color=C_ATTENTION, font_size=9)
add_arrow_down(slide, dec_x + dec_w/2, dy + 2.35, dy + 2.45, color=C_DECODER)
add_shape_box(slide, dec_x + 0.1, dy + 2.45, dec_w - 0.2, 0.3,
              "Add & Norm", fill_color=C_FFN, font_size=9)
add_textbox(slide, dec_x + 0.2, dy + 2.85, dec_w - 0.4, 0.4,
            "... (repeat x6)", font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

# Output Processing Zone (right)
out_x = 7.8
add_zone_box(slide, out_x, 1.5, 1.4, 3.5, "Output", C_LIGHT_GRAY, C_GRAY, label_color=C_DARK)
add_shape_box(slide, out_x + 0.1, 2.2, 1.2, 0.45, "Linear",
              fill_color=C_DARK, font_size=10)
add_arrow_down(slide, out_x + 0.7, 2.65, 2.8, color=C_GRAY)
add_shape_box(slide, out_x + 0.1, 2.8, 1.2, 0.45, "Softmax",
              fill_color=C_DARK, font_size=10)
add_textbox(slide, out_x + 0.1, 3.4, 1.2, 0.4, "Output\nProbabilities",
            font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER)
add_textbox(slide, out_x + 0.1, 4.2, 1.2, 0.4, "\"wo ai ni\"",
            font_size=9, color=C_TEXT, align=PP_ALIGN.CENTER, italic=True)

# Decoder output embedding (top)
add_shape_box(slide, dec_x + 0.3, 6.0, dec_w - 0.6, 0.4,
              "Output Embedding + Pos Enc", fill_color=C_LIGHT_PURPLE,
              text_color=C_EMBED, font_size=8)
add_arrow_down(slide, dec_x + dec_w/2, 6.0, dec_y + 0.5, color=C_DECODER)

# Key connections
# Encoder output -> Cross-Attention (K, V)
enc_out_y = enc_y + enc_h - 0.6
cross_y = dy + 0.95 + 0.22
add_connector(slide, enc_x + enc_w, cross_y, dec_x + 0.1, cross_y,
              color=C_EMBED, width=Pt(3))
add_textbox(slide, enc_x + enc_w + 0.05, cross_y - 0.3, 1.0, 0.25,
            "K, V", font_size=9, bold=True, color=C_EMBED)

# Input -> Encoder
add_arrow_right(slide, 1.65, 2.17, enc_x + 0.05, color=C_ENCODER, width=Pt(2))

# Decoder -> Linear
add_arrow_right(slide, dec_x + dec_w, 2.42, out_x + 0.05, color=C_DECODER, width=Pt(2))

# Bottom annotations
add_textbox(slide, 0.2, 6.7, 3.0, 0.4, "Input: English sentence",
            font_size=9, color=C_ENCODER, bold=True)
add_textbox(slide, 4.8, 6.7, 3.0, 0.4, "Output: Chinese translation",
            font_size=9, color=C_DECODER, bold=True)

# Right side reserved
add_rect(slide, 9.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 9.8, 3.0, 3.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 5 done.")

# ============================================================
# PAGE 6: Encoder Single Layer Structure
# ============================================================
print("Creating Page 6: Encoder Single Layer...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Encoder Layer: Self-Attention + FFN")

# Input arrow
add_textbox(slide, 1.5, 0.7, 3.0, 0.3, "Layer Input (512-dim)",
            font_size=10, color=C_TITLE, align=PP_ALIGN.CENTER)
add_arrow_down(slide, 3.0, 1.0, 1.2, color=C_TITLE, width=Pt(2))

# Main flow boxes
y1 = 1.2
bw = 3.5
bx = 1.3

# Multi-Head Self-Attention
add_shape_box(slide, bx, y1, bw, 0.7, "Multi-Head Self-Attention",
              fill_color=C_ENCODER, font_size=13)
add_textbox(slide, bx + 0.1, y1 + 0.5, bw - 0.2, 0.2, "8 heads x 64-dim",
            font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)

# Residual connection label (bypass arrow)
add_connector(slide, bx - 0.4, 1.0, bx - 0.4, y1 + 0.35, color=C_FFN, width=Pt(2))
add_connector(slide, bx - 0.4, y1 + 0.35, bx, y1 + 0.35, color=C_FFN, width=Pt(2))
add_textbox(slide, bx - 0.8, y1 + 0.1, 0.5, 0.3, "+", font_size=14,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Layer Norm 1
y2 = y1 + 0.9
add_shape_box(slide, bx, y2, bw, 0.45, "Add & Layer Norm",
              fill_color=C_FFN, font_size=12)
add_textbox(slide, bx + 0.1, y2 + 0.3, bw - 0.2, 0.15, "output: 512-dim",
            font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)

add_arrow_down(slide, bx + bw/2, y2 + 0.45, y2 + 0.65, color=C_DARK, width=Pt(2))

# Feed Forward Network
y3 = y2 + 0.65
add_shape_box(slide, bx, y3, bw, 0.7, "Feed Forward Network",
              fill_color=C_ATTENTION, font_size=13)
add_textbox(slide, bx + 0.1, y3 + 0.5, bw - 0.2, 0.2, "512 -> 2048 -> ReLU -> 512",
            font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)

# Residual connection 2
add_connector(slide, bx - 0.4, y2 + 0.22, bx - 0.4, y3 + 0.35, color=C_FFN, width=Pt(2))
add_connector(slide, bx - 0.4, y3 + 0.35, bx, y3 + 0.35, color=C_FFN, width=Pt(2))
add_textbox(slide, bx - 0.8, y3 + 0.1, 0.5, 0.3, "+", font_size=14,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Layer Norm 2
y4 = y3 + 0.9
add_shape_box(slide, bx, y4, bw, 0.45, "Add & Layer Norm",
              fill_color=C_FFN, font_size=12)

# Output arrow
add_arrow_down(slide, bx + bw/2, y4 + 0.45, y4 + 0.7, color=C_TITLE, width=Pt(2))
add_textbox(slide, 1.5, y4 + 0.7, 3.0, 0.3, "Layer Output -> Next Layer",
            font_size=10, color=C_TITLE, align=PP_ALIGN.CENTER)

# Annotations on the right side of the diagram
ax = bx + bw + 0.4
# Self-Attention annotation
add_shape_box(slide, ax, y1 + 0.05, 2.0, 0.6, "",
              fill_color=C_LIGHT_BLUE, line_color=C_ENCODER)
add_multiline_tb(slide, ax + 0.05, y1 + 0.08, 1.9, 0.55, [
    "\"Round Table\"",
    "All words exchange",
    "information",
], font_size=8, color=C_ENCODER)

# Add&Norm annotation
add_shape_box(slide, ax, y2, 2.0, 0.4, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, ax + 0.05, y2 + 0.02, 1.9, 0.35, [
    "\"Safety Net\"",
    "Standardize + preserve",
], font_size=8, color=C_FFN)

# FFN annotation
add_shape_box(slide, ax, y3 + 0.05, 2.0, 0.6, "",
              fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_multiline_tb(slide, ax + 0.05, y3 + 0.08, 1.9, 0.55, [
    "\"Think Independently\"",
    "Non-linear transform",
    "per word",
], font_size=8, color=C_ATTENTION)

# Repeat label
add_textbox(slide, 1.0, y4 + 1.1, 4.0, 0.3,
            "x 6 layers (each layer refines word representations)",
            font_size=10, bold=True, color=C_TITLE)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 6 done.")

# ============================================================
# PAGE 7: Decoder Single Layer Structure
# ============================================================
print("Creating Page 7: Decoder Single Layer...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Decoder Layer: Masked + Cross + FFN")

# Encoder output reference (left side)
enc_ref_x = 0.3
enc_ref_y = 2.0
add_shape_box(slide, enc_ref_x, enc_ref_y, 2.2, 0.7,
              "Encoder Output\n(K, V from encoder)",
              fill_color=C_EMBED, font_size=9)
add_textbox(slide, enc_ref_x, enc_ref_y - 0.35, 2.2, 0.3,
            "From Encoder", font_size=9, color=C_EMBED, bold=True, align=PP_ALIGN.CENTER)

# Decoder main flow (right side)
dx = 3.0
dw = 3.8
dy = 0.9

# Input
add_textbox(slide, dx + 0.5, 0.7, 3.0, 0.3, "Decoder Layer Input",
            font_size=10, color=C_DECODER, align=PP_ALIGN.CENTER)
add_arrow_down(slide, dx + dw/2, 1.0, 1.1, color=C_DECODER, width=Pt(2))

# 1. Masked Self-Attention
add_shape_box(slide, dx, dy, dw, 0.6, "Masked Multi-Head Self-Attention",
              fill_color=C_ENCODER, font_size=11)
add_textbox(slide, dx + 0.1, dy + 0.42, dw - 0.2, 0.18,
            "Causal mask: can only see previous words",
            font_size=7, color=C_WHITE, align=PP_ALIGN.CENTER)

# Residual 1
add_connector(slide, dx - 0.4, 1.0, dx - 0.4, dy + 0.3, color=C_FFN, width=Pt(2))
add_connector(slide, dx - 0.4, dy + 0.3, dx, dy + 0.3, color=C_FFN, width=Pt(2))
add_textbox(slide, dx - 0.7, dy + 0.1, 0.4, 0.3, "+", font_size=14,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Norm 1
dy2 = dy + 0.75
add_shape_box(slide, dx, dy2, dw, 0.35, "Add & Norm", fill_color=C_FFN, font_size=10)
add_arrow_down(slide, dx + dw/2, dy2 + 0.35, dy2 + 0.55, color=C_DARK, width=Pt(2))

# 2. Cross-Attention
dy3 = dy2 + 0.55
add_shape_box(slide, dx, dy3, dw, 0.6, "Multi-Head Cross-Attention",
              fill_color=C_EMBED, font_size=11)
add_textbox(slide, dx + 0.1, dy3 + 0.42, dw - 0.2, 0.18,
            "Q from decoder, K,V from encoder",
            font_size=7, color=C_WHITE, align=PP_ALIGN.CENTER)

# Arrow from encoder output to cross-attention
add_connector(slide, enc_ref_x + 2.2, dy3 + 0.3, dx, dy3 + 0.3,
              color=C_EMBED, width=Pt(3))
add_textbox(slide, enc_ref_x + 2.3, dy3 + 0.05, 0.7, 0.25,
            "K, V", font_size=9, bold=True, color=C_EMBED)

# Residual 2
add_connector(slide, dx - 0.4, dy2 + 0.17, dx - 0.4, dy3 + 0.3, color=C_FFN, width=Pt(2))
add_connector(slide, dx - 0.4, dy3 + 0.3, dx, dy3 + 0.3, color=C_FFN, width=Pt(2))
add_textbox(slide, dx - 0.7, dy3 + 0.1, 0.4, 0.3, "+", font_size=14,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Norm 2
dy4 = dy3 + 0.75
add_shape_box(slide, dx, dy4, dw, 0.35, "Add & Norm", fill_color=C_FFN, font_size=10)
add_arrow_down(slide, dx + dw/2, dy4 + 0.35, dy4 + 0.55, color=C_DARK, width=Pt(2))

# 3. Feed Forward
dy5 = dy4 + 0.55
add_shape_box(slide, dx, dy5, dw, 0.6, "Feed Forward Network",
              fill_color=C_ATTENTION, font_size=11)

# Residual 3
add_connector(slide, dx - 0.4, dy4 + 0.17, dx - 0.4, dy5 + 0.3, color=C_FFN, width=Pt(2))
add_connector(slide, dx - 0.4, dy5 + 0.3, dx, dy5 + 0.3, color=C_FFN, width=Pt(2))
add_textbox(slide, dx - 0.7, dy5 + 0.1, 0.4, 0.3, "+", font_size=14,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Norm 3
dy6 = dy5 + 0.75
add_shape_box(slide, dx, dy6, dw, 0.35, "Add & Norm", fill_color=C_FFN, font_size=10)

# Output
add_arrow_down(slide, dx + dw/2, dy6 + 0.35, dy6 + 0.65, color=C_DECODER, width=Pt(2))
add_textbox(slide, dx + 0.5, dy6 + 0.65, 3.0, 0.3, "Layer Output",
            font_size=10, color=C_DECODER, align=PP_ALIGN.CENTER)

# Key difference annotations
ann_x = dx + dw + 0.3
add_shape_box(slide, ann_x, dy + 0.05, 1.8, 0.5, "",
              fill_color=C_LIGHT_BLUE, line_color=C_ENCODER)
add_textbox(slide, ann_x + 0.05, dy + 0.1, 1.7, 0.4,
            "Can only see\nprevious words", font_size=8, color=C_ENCODER, align=PP_ALIGN.CENTER)

add_shape_box(slide, ann_x, dy3 + 0.05, 1.8, 0.5, "",
              fill_color=C_LIGHT_PURPLE, line_color=C_EMBED)
add_textbox(slide, ann_x + 0.05, dy3 + 0.1, 1.7, 0.4,
            "Bridge to encoder\n(query encoder info)", font_size=8, color=C_EMBED, align=PP_ALIGN.CENTER)

add_shape_box(slide, ann_x, dy5 + 0.05, 1.8, 0.5, "",
              fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_textbox(slide, ann_x + 0.05, dy5 + 0.1, 1.7, 0.4,
            "Same as encoder\nFFN", font_size=8, color=C_ATTENTION, align=PP_ALIGN.CENTER)

# Repeat
add_textbox(slide, 3.0, dy6 + 1.1, 4.0, 0.3,
            "x 6 layers (3 sub-layers each, vs encoder's 2)",
            font_size=10, bold=True, color=C_DECODER)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 7 done.")

# ============================================================
# PAGE 9: Embedding Matrix + Positional Encoding
# ============================================================
print("Creating Page 9: Embedding + Positional Encoding...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Embedding + Positional Encoding")

# Token boxes at top
tokens = ["I", "love", "you"]
tx_start = 1.0
ty = 0.9
for i, tok in enumerate(tokens):
    x = tx_start + i * 1.5
    add_shape_box(slide, x, ty, 1.0, 0.45, f'"{tok}"',
                  fill_color=C_ENCODER, font_size=12)
    add_arrow_down(slide, x + 0.5, ty + 0.45, ty + 0.85, color=C_ENCODER)

# Embedding Matrix
em_y = ty + 0.85
add_shape_box(slide, 0.8, em_y, 4.0, 0.55, "Embedding Matrix (37000 x 512)",
              fill_color=C_EMBED, font_size=11)
add_textbox(slide, 0.8, em_y + 0.55, 4.0, 0.2,
            "Each row = one token's 512-dim vector",
            font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER)

# Arrows to word embedding vectors
we_y = em_y + 0.9
for i in range(3):
    x = tx_start + i * 1.5
    add_arrow_down(slide, x + 0.5, em_y + 0.55, we_y, color=C_EMBED)

# Word embedding vectors
for i, tok in enumerate(tokens):
    x = tx_start + i * 1.5
    add_shape_box(slide, x, we_y, 1.0, 0.4, f"{tok}: 512-d",
                  fill_color=C_LIGHT_PURPLE, text_color=C_EMBED, font_size=8)

# Positional Encoding (right side)
pe_y = em_y + 0.1
add_shape_box(slide, 5.2, pe_y, 2.2, 0.55, "Positional Encoding\n(sin/cos functions)",
              fill_color=C_FFN, font_size=9)

# Arrows from PE to position vectors
for i in range(3):
    x = tx_start + i * 1.5
    add_connector(slide, 5.2, pe_y + 0.27, x + 1.0, we_y + 0.2,
                  color=C_FFN, width=Pt(1))

# Plus symbol and final vectors
plus_y = we_y + 0.6
add_textbox(slide, 1.5, plus_y, 4.0, 0.3,
            "Word Embedding  +  Position Encoding  =  Final Input",
            font_size=10, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
add_textbox(slide, 2.0, plus_y + 0.2, 1.0, 0.4, "(+)", font_size=24,
            bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Final input vectors
fi_y = plus_y + 0.6
for i, tok in enumerate(tokens):
    x = tx_start + i * 1.5
    add_shape_box(slide, x, fi_y, 1.0, 0.45, f"512-d\nwith position",
                  fill_color=C_ENCODER, font_size=8)
    add_arrow_down(slide, x + 0.5, fi_y + 0.45, fi_y + 0.75, color=C_ENCODER)

# Arrow to encoder
add_textbox(slide, 1.0, fi_y + 0.75, 4.0, 0.3,
            "-> Feed into Encoder Layer 1",
            font_size=10, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)

# Position encoding detail (bottom)
pe_detail_y = 4.8
add_textbox(slide, 0.5, pe_detail_y, 6.5, 0.3,
            "Positional Encoding Detail (sinusoidal):",
            font_size=11, bold=True, color=C_TITLE)
add_multiline_tb(slide, 0.5, pe_detail_y + 0.35, 6.5, 1.5, [
    "PE(pos, 2i)   = sin(pos / 10000^(2i/d_model))",
    "PE(pos, 2i+1) = cos(pos / 10000^(2i/d_model))",
    "",
    "Each position gets a unique 512-dim vector using different frequencies.",
    "No learnable parameters needed. Can generalize to longer sequences.",
], font_size=9, color=C_TEXT, font_name=FONT_CODE)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 9 done.")

# ============================================================
# PAGE 10: QKV Matrix Generation
# ============================================================
print("Creating Page 10: QKV Generation...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Q, K, V Generation: Search Tools for Each Word")

# Input embedding at top
add_shape_box(slide, 1.5, 0.9, 3.5, 0.55, "Input Embedding (512-dim)",
              fill_color=C_ENCODER, font_size=12)

# Three branches
# Q branch
add_connector(slide, 2.5, 1.45, 1.0, 1.85, color=C_ENCODER, width=Pt(2))
add_shape_box(slide, 0.3, 1.85, 1.8, 0.5, "W^Q (512x64)",
              fill_color=C_ENCODER, font_size=10)
add_arrow_down(slide, 1.2, 2.35, 2.6, color=C_ENCODER)
add_shape_box(slide, 0.3, 2.6, 1.8, 0.45, "Query (64-dim)",
              fill_color=C_ENCODER, font_size=11)

# K branch
add_connector(slide, 3.25, 1.45, 3.0, 1.85, color=C_ATTENTION, width=Pt(2))
add_shape_box(slide, 2.3, 1.85, 1.8, 0.5, "W^K (512x64)",
              fill_color=C_ATTENTION, font_size=10)
add_arrow_down(slide, 3.2, 2.35, 2.6, color=C_ATTENTION)
add_shape_box(slide, 2.3, 2.6, 1.8, 0.45, "Key (64-dim)",
              fill_color=C_ATTENTION, font_size=11)

# V branch
add_connector(slide, 4.0, 1.45, 5.0, 1.85, color=C_FFN, width=Pt(2))
add_shape_box(slide, 4.3, 1.85, 1.8, 0.5, "W^V (512x64)",
              fill_color=C_FFN, font_size=10)
add_arrow_down(slide, 5.2, 2.35, 2.6, color=C_FFN)
add_shape_box(slide, 4.3, 2.6, 1.8, 0.45, "Value (64-dim)",
              fill_color=C_FFN, font_size=11)

# Library analogy section
ana_y = 3.5
add_textbox(slide, 0.3, ana_y, 6.5, 0.35,
            "Library Search Analogy:",
            font_size=12, bold=True, color=C_TITLE)

# Q analogy
add_shape_box(slide, 0.3, ana_y + 0.45, 1.8, 0.5, "",
              fill_color=C_LIGHT_BLUE, line_color=C_ENCODER)
add_multiline_tb(slide, 0.35, ana_y + 0.48, 1.7, 0.45, [
    "Query = Search",
    "keyword",
], font_size=9, color=C_ENCODER)

# Arrow
add_textbox(slide, 2.2, ana_y + 0.55, 0.5, 0.3, "->", font_size=18,
            color=C_GRAY, align=PP_ALIGN.CENTER)

# K analogy
add_shape_box(slide, 2.6, ana_y + 0.45, 1.8, 0.5, "",
              fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_multiline_tb(slide, 2.65, ana_y + 0.48, 1.7, 0.45, [
    "Key = Book title",
    "label/tag",
], font_size=9, color=C_ATTENTION)

# Arrow
add_textbox(slide, 4.5, ana_y + 0.55, 0.5, 0.3, "->", font_size=18,
            color=C_GRAY, align=PP_ALIGN.CENTER)

# V analogy
add_shape_box(slide, 4.9, ana_y + 0.45, 1.8, 0.5, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, 4.95, ana_y + 0.48, 1.7, 0.45, [
    "Value = Book",
    "content",
], font_size=9, color=C_FFN)

# Multi-head note
mh_y = ana_y + 1.2
add_shape_box(slide, 0.3, mh_y, 6.0, 0.8, "",
              fill_color=C_LIGHT_GRAY, line_color=C_DARK, line_width=Pt(1))
add_multiline_tb(slide, 0.5, mh_y + 0.05, 5.6, 0.7, [
    "Multi-Head: 8 independent sets of W^Q, W^K, W^V",
    "Each head: 64-dim Q, K, V (8 x 64 = 512 = d_model)",
    "Different heads learn different relationship types (syntax, semantics, etc.)",
], font_size=9, color=C_TEXT)

# Per-word illustration
pw_y = mh_y + 1.1
add_textbox(slide, 0.3, pw_y, 6.5, 0.3,
            "Every word generates its own Q, K, V simultaneously:",
            font_size=10, bold=True, color=C_TITLE)

# Example with "it" and "cat"
for i, (word, desc) in enumerate([("it", "Q: 'find a noun'"), ("cat", "K: 'noun, animal'")]):
    x = 0.5 + i * 3.2
    add_shape_box(slide, x, pw_y + 0.35, 2.8, 0.7, "",
                  fill_color=C_WHITE, line_color=C_ENCODER, line_width=Pt(1))
    add_textbox(slide, x + 0.05, pw_y + 0.38, 2.7, 0.25, f'Word: "{word}"',
                font_size=10, bold=True, color=C_ENCODER)
    add_textbox(slide, x + 0.05, pw_y + 0.65, 2.7, 0.35, desc,
                font_size=9, color=C_TEXT)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 10 done.")

# ============================================================
# PAGE 11: Attention 5-Step Computation
# ============================================================
print("Creating Page 11: Attention 5 Steps...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Self-Attention: 5-Step Computation")

# Vertical flow of 5 steps
steps = [
    ("Step 1: Generate Q, K, V", "Input x W^Q/W^K/W^V -> Q, K, V", C_ENCODER, "All words in parallel"),
    ("Step 2: Compute Scores", "Q x K^T -> Score Matrix (TxT)", C_LIGHT_BLUE, "Dot product = relevance"),
    ("Step 3: Scale Scores", "Score / sqrt(d_k)  [sqrt(64) = 8]", C_LIGHT_ORANGE, "Prevent softmax saturation"),
    ("Step 4: Softmax", "Softmax -> Probability Matrix", C_FFN, "Each row sums to 1"),
    ("Step 5: Weighted Sum", "Probability x V -> Output", C_EMBED, "Fused context info"),
]

sy = 0.85
sh = 0.55
sg = 0.35  # gap between steps
bw = 4.0
bx = 0.8

for i, (title, detail, color, note) in enumerate(steps):
    y = sy + i * (sh + sg)
    add_shape_box(slide, bx, y, bw, sh, title,
                  fill_color=color, font_size=11,
                  text_color=C_DARK if i in [1, 2] else C_WHITE)
    # Detail text
    add_textbox(slide, bx + 0.1, y + sh + 0.02, bw - 0.2, 0.22,
                detail, font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER,
                font_name=FONT_CODE)
    # Note on the right
    add_textbox(slide, bx + bw + 0.15, y + 0.1, 2.0, 0.35,
                note, font_size=8, color=C_DARK, italic=True)
    # Arrow to next step
    if i < len(steps) - 1:
        add_arrow_down(slide, bx + bw/2, y + sh, y + sh + sg * 0.3,
                       color=C_DARK, width=Pt(2))

# Small matrix visualizations on the far right
mx = 7.0

# Score matrix (3x3)
my1 = 1.0
add_textbox(slide, mx, my1 - 0.25, 1.5, 0.25, "Score Matrix:", font_size=8, bold=True, color=C_TITLE)
# Draw a simple 3x3 matrix with color-coded cells
matrix_vals = [
    [(0.9, C_ENCODER), (0.1, C_LIGHT_GRAY), (0.0, C_WHITE)],
    [(0.3, C_LIGHT_BLUE), (0.6, C_ENCODER), (0.1, C_LIGHT_GRAY)],
    [(0.1, C_LIGHT_GRAY), (0.2, C_LIGHT_BLUE), (0.7, C_ENCODER)],
]
cell_w = 0.35
cell_h = 0.3
for r, row in enumerate(matrix_vals):
    for c, (val, clr) in enumerate(row):
        cx = mx + c * cell_w
        cy = my1 + r * cell_h
        add_shape_box(slide, cx, cy, cell_w, cell_h, f"{val:.1f}",
                      fill_color=clr, font_size=7,
                      text_color=C_DARK if val > 0.3 else C_GRAY)

# Softmax result
my2 = 2.5
add_textbox(slide, mx, my2 - 0.25, 1.5, 0.25, "After Softmax:", font_size=8, bold=True, color=C_FFN)
softmax_vals = [
    [(0.55, C_ENCODER), (0.25, C_LIGHT_BLUE), (0.20, C_LIGHT_GRAY)],
    [(0.15, C_LIGHT_GRAY), (0.50, C_ENCODER), (0.35, C_LIGHT_BLUE)],
    [(0.10, C_LIGHT_GRAY), (0.20, C_LIGHT_BLUE), (0.70, C_ENCODER)],
]
for r, row in enumerate(softmax_vals):
    for c, (val, clr) in enumerate(row):
        cx = mx + c * cell_w
        cy = my2 + r * cell_h
        add_shape_box(slide, cx, cy, cell_w, cell_h, f"{val:.2f}",
                      fill_color=clr, font_size=7,
                      text_color=C_DARK if val > 0.3 else C_GRAY)

# Formula box
fy = 4.5
add_shape_box(slide, 0.5, fy, 6.5, 1.2, "",
              fill_color=C_LIGHT_GRAY, line_color=C_DARK, line_width=Pt(1))
add_multiline_tb(slide, 0.7, fy + 0.05, 6.1, 1.1, [
    "Attention(Q, K, V) = Softmax(QK^T / sqrt(d_k)) x V",
    "",
    "Where: d_k = 64 (dimension per head)",
    "Q, K, V are computed for ALL words simultaneously (fully parallel)",
    "Output shape: same as input (each word gets an enriched representation)",
], font_size=9, color=C_TEXT, font_name=FONT_CODE)

# Right reserved
add_rect(slide, 7.5, 3.5, 0.01, 3.5, C_LIGHT_GRAY)

print("  Page 11 done.")

# ============================================================
# PAGE 12: Residual Connection + LayerNorm + FFN
# ============================================================
print("Creating Page 12: Residual + LayerNorm + FFN...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Residual Connection + LayerNorm + FFN")

# Detailed data flow with dimensions
bx = 1.0
bw = 4.0
y = 0.9

# Input
add_shape_box(slide, bx, y, bw, 0.4, "Input (512-dim)",
              fill_color=C_ENCODER, font_size=11)
add_textbox(slide, bx, y - 0.25, bw, 0.25, "Layer Input",
            font_size=9, color=C_TITLE, align=PP_ALIGN.CENTER)

# Arrow down + residual bypass
add_arrow_down(slide, bx + bw/2, y + 0.4, y + 0.65, color=C_DARK, width=Pt(2))

# Self-Attention
y2 = y + 0.65
add_shape_box(slide, bx, y2, bw, 0.6, "Multi-Head Attention",
              fill_color=C_ENCODER, font_size=12)
add_textbox(slide, bx, y2 + 0.42, bw, 0.2, "8 heads x 64-dim = 512-dim output",
            font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)

# Residual arrow (bypass)
res_x = bx - 0.5
add_connector(slide, res_x, y + 0.2, res_x, y2 + 0.3, color=C_FFN, width=Pt(2.5))
add_connector(slide, res_x, y2 + 0.3, bx, y2 + 0.3, color=C_FFN, width=Pt(2.5))
add_textbox(slide, res_x - 0.35, y2 + 0.05, 0.4, 0.3, "+",
            font_size=16, bold=True, color=C_FFN, align=PP_ALIGN.CENTER)
add_label(slide, res_x - 0.6, y + 0.5, "Residual\n(skip)", font_size=7, color=C_FFN)

# Add & Norm 1
y3 = y2 + 0.8
add_shape_box(slide, bx, y3, bw, 0.45, "Add & Layer Norm",
              fill_color=C_FFN, font_size=11)
add_textbox(slide, bx, y3 + 0.3, bw, 0.15, "output: 512-dim (normalized)",
            font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)
add_arrow_down(slide, bx + bw/2, y3 + 0.45, y3 + 0.65, color=C_DARK, width=Pt(2))

# FFN
y4 = y3 + 0.65
add_shape_box(slide, bx, y4, bw, 0.7, "Feed Forward Network",
              fill_color=C_ATTENTION, font_size=12)
# FFN internal detail
add_textbox(slide, bx + 0.1, y4 + 0.45, bw - 0.2, 0.25,
            "512 -> 2048 (ReLU) -> 512",
            font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER, font_name=FONT_CODE)

# Residual 2
add_connector(slide, res_x, y3 + 0.22, res_x, y4 + 0.35, color=C_FFN, width=Pt(2.5))
add_connector(slide, res_x, y4 + 0.35, bx, y4 + 0.35, color=C_FFN, width=Pt(2.5))
add_textbox(slide, res_x - 0.35, y4 + 0.1, 0.4, 0.3, "+",
            font_size=16, bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Add & Norm 2
y5 = y4 + 0.9
add_shape_box(slide, bx, y5, bw, 0.45, "Add & Layer Norm",
              fill_color=C_FFN, font_size=11)

# Output
add_arrow_down(slide, bx + bw/2, y5 + 0.45, y5 + 0.65, color=C_TITLE, width=Pt(2))
add_shape_box(slide, bx, y5 + 0.65, bw, 0.4, "Output (512-dim) -> Next Layer",
              fill_color=C_ENCODER, font_size=10)

# Explanation boxes on the right
ex = bx + bw + 0.5

# Residual explanation
add_shape_box(slide, ex, y2, 2.0, 0.7, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, ex + 0.05, y2 + 0.03, 1.9, 0.65, [
    "Residual Connection:",
    "output = x + F(x)",
    "Like a highway ramp",
    "Preserves original info",
], font_size=8, color=C_FFN)

# LayerNorm explanation
add_shape_box(slide, ex, y3 + 0.05, 2.0, 0.5, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, ex + 0.05, y3 + 0.08, 1.9, 0.45, [
    "Layer Normalization:",
    "Standardize each vector",
    "Stable training",
], font_size=8, color=C_FFN)

# FFN explanation
add_shape_box(slide, ex, y4 + 0.05, 2.0, 0.7, "",
              fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_multiline_tb(slide, ex + 0.05, y4 + 0.08, 1.9, 0.65, [
    "Feed Forward:",
    "Expand: 512 -> 2048",
    "Activate: ReLU",
    "Compress: 2048 -> 512",
    "~2x compute of attention",
], font_size=8, color=C_ATTENTION)

# Bottom note
add_textbox(slide, 0.5, 6.5, 6.5, 0.4,
            "This pattern (Attention -> Add&Norm -> FFN -> Add&Norm) repeats 6x in encoder",
            font_size=10, bold=True, color=C_TITLE)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 12 done.")

# ============================================================
# PAGE 13: Decoder Training - Masked + Cross Attention
# ============================================================
print("Creating Page 13: Decoder Training...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Decoder Training: Masked + Cross Attention")

# Left: Causal mask matrix
add_textbox(slide, 0.3, 0.85, 3.0, 0.3, "Causal Mask Matrix:",
            font_size=11, bold=True, color=C_TITLE)

mask_labels = ["[START]", "wo", "ai"]
cell_s = 0.7
mx = 0.5
my = 1.25

# Column headers
for j, lbl in enumerate(mask_labels):
    add_textbox(slide, mx + (j + 1) * cell_s, my, cell_s, 0.25, lbl,
                font_size=8, color=C_TEXT, align=PP_ALIGN.CENTER, bold=True)
# Row headers and cells
for i, lbl in enumerate(mask_labels):
    add_textbox(slide, mx, my + (i + 1) * cell_s, cell_s, 0.25, lbl,
                font_size=8, color=C_TEXT, bold=True)
    for j in range(len(mask_labels)):
        cx = mx + (j + 1) * cell_s
        cy = my + (i + 1) * cell_s
        if j <= i:  # visible
            add_shape_box(slide, cx, cy, cell_s - 0.05, cell_s - 0.05, "score",
                          fill_color=C_ENCODER, font_size=7)
        else:  # masked
            add_shape_box(slide, cx, cy, cell_s - 0.05, cell_s - 0.05, "-inf",
                          fill_color=C_DECODER, font_size=7)

add_textbox(slide, 0.5, my + 4 * cell_s + 0.1, 3.0, 0.4,
            "Upper triangle = -inf\nSoftmax -> 0 (invisible)",
            font_size=8, color=C_DECODER)

# Right: Cross-attention diagram
ca_x = 4.0
ca_y = 0.85

add_textbox(slide, ca_x, ca_y, 3.0, 0.3, "Cross-Attention:",
            font_size=11, bold=True, color=C_TITLE)

# Encoder output
add_shape_box(slide, ca_x, ca_y + 0.4, 2.8, 1.0,
              "Encoder Output\n(K: love's Key\n V: love's Value)",
              fill_color=C_EMBED, font_size=9)

# Decoder Q
add_shape_box(slide, ca_x, ca_y + 1.7, 2.8, 0.5,
              "Decoder Query\n(\"find a verb for love\")",
              fill_color=C_ENCODER, font_size=9)

# Arrow from Q to Encoder
add_connector(slide, ca_x + 2.8, ca_y + 1.95, ca_x + 2.8, ca_y + 0.9,
              color=C_EMBED, width=Pt(2))
add_textbox(slide, ca_x + 2.85, ca_y + 1.2, 0.5, 0.3, "match",
            font_size=7, color=C_EMBED)

# Result
add_shape_box(slide, ca_x, ca_y + 2.5, 2.8, 0.5,
              "Result: helps generate \"ai\"",
              fill_color=C_FFN, font_size=9)
add_arrow_down(slide, ca_x + 1.4, ca_y + 2.2, ca_y + 2.5, color=C_FFN)

# Bottom: Full decoder training flow
ft_y = 4.5
add_textbox(slide, 0.3, ft_y, 7.0, 0.3,
            "Decoder Training Flow:",
            font_size=11, bold=True, color=C_TITLE)

# Flow boxes
flow_items = [
    ("Input:\n[START, wo, ai]", C_LIGHT_PURPLE, C_EMBED, 0.3),
    ("Masked\nSelf-Attn", C_ENCODER, C_WHITE, 2.0),
    ("Cross\nAttention", C_EMBED, C_WHITE, 3.5),
    ("FFN", C_ATTENTION, C_WHITE, 5.0),
    ("Linear +\nSoftmax", C_DARK, C_WHITE, 6.3),
]
for text, fill, tc, fx in flow_items:
    add_shape_box(slide, fx, ft_y + 0.4, 1.4, 0.7, text,
                  fill_color=fill, text_color=tc, font_size=9)

# Arrows between flow items
for i in range(len(flow_items) - 1):
    x1 = flow_items[i][3] + 1.4
    x2 = flow_items[i + 1][3]
    y_mid = ft_y + 0.75
    add_arrow_right(slide, x1, y_mid, x2, color=C_DARK, width=Pt(1.5))

# Encoder input from below
add_shape_box(slide, 3.5, ft_y + 1.3, 1.4, 0.4, "Encoder K,V",
              fill_color=C_LIGHT_PURPLE, text_color=C_EMBED, font_size=8)
add_arrow_down(slide, 4.2, ft_y + 1.1, ft_y + 1.1, color=C_EMBED, width=Pt(1.5))

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 13 done.")

# ============================================================
# PAGE 14: Loss Function
# ============================================================
print("Creating Page 14: Loss Function...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Loss Function: Cross-Entropy")

# Flow diagram
fx = 0.5
fy = 0.9

add_shape_box(slide, fx, fy, 2.5, 0.45, "Decoder Output (512-dim)",
              fill_color=C_ENCODER, font_size=10)
add_arrow_right(slide, fx + 2.5, fy + 0.22, fx + 2.8, color=C_DARK, width=Pt(2))

add_shape_box(slide, fx + 2.8, fy, 2.0, 0.45, "Linear (512->37000)",
              fill_color=C_DARK, font_size=9)
add_arrow_right(slide, fx + 4.8, fy + 0.22, fx + 5.1, color=C_DARK, width=Pt(2))

add_shape_box(slide, fx + 5.1, fy, 1.8, 0.45, "Softmax",
              fill_color=C_DARK, font_size=10)

fy2 = fy + 0.7
add_arrow_down(slide, fx + 6.0, fy + 0.45, fy2, color=C_DARK, width=Pt(2))

add_shape_box(slide, fx + 4.5, fy2, 3.0, 0.45, "Predictions (probabilities)",
              fill_color=C_ATTENTION, font_size=10)

# Targets
add_shape_box(slide, fx + 0.5, fy2, 2.5, 0.45, "Targets (ground truth)",
              fill_color=C_DECODER, font_size=10)

# Both arrows to loss
fy3 = fy2 + 0.7
add_arrow_down(slide, fx + 1.75, fy2 + 0.45, fy3, color=C_DECODER, width=Pt(2))
add_arrow_down(slide, fx + 6.0, fy2 + 0.45, fy3, color=C_ATTENTION, width=Pt(2))

add_shape_box(slide, fx + 2.5, fy3, 3.5, 0.5, "Cross-Entropy Loss",
              fill_color=C_DECODER, font_size=12)

fy4 = fy3 + 0.7
add_arrow_down(slide, fx + 4.25, fy3 + 0.5, fy4, color=C_DARK, width=Pt(2))

add_shape_box(slide, fx + 2.0, fy4, 4.5, 0.45, "Backpropagation -> Update Parameters (Adam)",
              fill_color=C_DARK, font_size=10)

# Loss formula
fy5 = fy4 + 0.7
add_shape_box(slide, fx + 1.0, fy5, 6.0, 0.8, "",
              fill_color=C_LIGHT_GRAY, line_color=C_DARK, line_width=Pt(1))
add_multiline_tb(slide, fx + 1.2, fy5 + 0.05, 5.6, 0.7, [
    "Loss = -log(p(correct word))",
    "Average over all positions in the sequence",
    "",
    "p=0.95 -> Loss=0.05 (good)  |  p=0.01 -> Loss=4.6 (bad)",
    "Training start: Loss ~10.5  |  After training: Loss ~1-3",
], font_size=9, color=C_TEXT, font_name=FONT_CODE)

# Cross-entropy intuition
fy6 = fy5 + 1.0
add_textbox(slide, 0.5, fy6, 7.0, 0.3,
            "Cross-Entropy Intuition: Multiple Choice Score",
            font_size=11, bold=True, color=C_TITLE)

# Example
add_shape_box(slide, 0.5, fy6 + 0.35, 6.5, 0.8, "",
              fill_color=C_WHITE, line_color=C_ENCODER, line_width=Pt(1))
add_multiline_tb(slide, 0.7, fy6 + 0.4, 6.1, 0.7, [
    "Question: What is the next word after \"I love\"?",
    "Correct answer: \"you\"  |  Model prediction: p(\"you\") = 0.85",
    "Loss = -log(0.85) = 0.16  (low = model is confident and correct)",
], font_size=9, color=C_TEXT)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 14 done.")

# ============================================================
# PAGE 15: Teacher Forcing
# ============================================================
print("Creating Page 15: Teacher Forcing...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Teacher Forcing: Open-Book Exam")

# Top: With Teacher Forcing
tf_y = 0.9
add_textbox(slide, 0.3, tf_y, 6.5, 0.35,
            "With Teacher Forcing (Training):",
            font_size=12, bold=True, color=C_FFN)

# Input sequence
input_words_tf = ["[START]", "wo (true)", "ai (true)"]
pred_words_tf = ["predict: wo", "predict: ai", "predict: ni"]
for i, (inp, pred) in enumerate(zip(input_words_tf, pred_words_tf)):
    x = 0.5 + i * 2.2
    # Input box (green = correct answer)
    add_shape_box(slide, x, tf_y + 0.45, 1.8, 0.4, inp,
                  fill_color=C_FFN, font_size=9)
    add_textbox(slide, x, tf_y + 0.85, 1.8, 0.25, "(ground truth)",
                font_size=7, color=C_FFN, align=PP_ALIGN.CENTER)
    # Arrow down
    add_arrow_down(slide, x + 0.9, tf_y + 1.1, tf_y + 1.25, color=C_DARK, width=Pt(1.5))
    # Prediction box
    add_shape_box(slide, x, tf_y + 1.25, 1.8, 0.4, pred,
                  fill_color=C_LIGHT_BLUE, text_color=C_ENCODER, font_size=9)
    # Arrow to next input
    if i < len(input_words_tf) - 1:
        add_arrow_right(slide, x + 1.8, tf_y + 0.65, x + 2.2, color=C_FFN, width=Pt(1.5))

# Checkmark labels
add_textbox(slide, 0.5, tf_y + 1.75, 6.0, 0.3,
            "Every step uses the CORRECT previous word as input",
            font_size=9, bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# Separator
sep_y = tf_y + 2.2
add_connector(slide, 0.3, sep_y, 7.0, sep_y, color=C_GRAY, width=Pt(1))
add_textbox(slide, 2.5, sep_y + 0.05, 3.0, 0.25, "VS",
            font_size=12, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

# Bottom: Without Teacher Forcing
nf_y = sep_y + 0.35
add_textbox(slide, 0.3, nf_y, 6.5, 0.35,
            "Without Teacher Forcing (What if):",
            font_size=12, bold=True, color=C_DECODER)

input_words_nf = ["[START]", "ta (wrong)", "bu (wrong)"]
pred_words_nf = ["predict: ta", "predict: bu", "predict: ..."]
for i, (inp, pred) in enumerate(zip(input_words_nf, pred_words_nf)):
    x = 0.5 + i * 2.2
    fc = C_DECODER if i > 0 else C_GRAY
    add_shape_box(slide, x, nf_y + 0.45, 1.8, 0.4, inp,
                  fill_color=fc, font_size=9)
    add_textbox(slide, x, nf_y + 0.85, 1.8, 0.25, "(model's guess)" if i > 0 else "(start token)",
                font_size=7, color=C_DECODER if i > 0 else C_GRAY, align=PP_ALIGN.CENTER)
    add_arrow_down(slide, x + 0.9, nf_y + 1.1, nf_y + 1.25, color=C_DARK, width=Pt(1.5))
    add_shape_box(slide, x, nf_y + 1.25, 1.8, 0.4, pred,
                  fill_color=C_LIGHT_RED, text_color=C_DECODER, font_size=9)
    if i < len(input_words_nf) - 1:
        add_arrow_right(slide, x + 1.8, nf_y + 0.65, x + 2.2, color=C_DECODER, width=Pt(1.5))

add_textbox(slide, 0.5, nf_y + 1.75, 6.0, 0.3,
            "Errors accumulate: one wrong word -> all following words may be wrong",
            font_size=9, bold=True, color=C_DECODER, align=PP_ALIGN.CENTER)

# Pros and Cons
pc_y = nf_y + 2.2
add_shape_box(slide, 0.3, pc_y, 3.2, 1.0, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, 0.5, pc_y + 0.05, 2.8, 0.9, [
    "Pros (Teacher Forcing):",
    "+ Fast, parallel training",
    "+ Stable gradient signal",
    "+ Efficient GPU utilization",
], font_size=9, color=C_FFN)

add_shape_box(slide, 3.8, pc_y, 3.2, 1.0, "",
              fill_color=C_LIGHT_RED, line_color=C_DECODER)
add_multiline_tb(slide, 4.0, pc_y + 0.05, 2.8, 0.9, [
    "Cons (Exposure Bias):",
    "- Train/test input mismatch",
    "- Model never sees errors",
    "- May fail at inference",
], font_size=9, color=C_DECODER)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 15 done.")

# ============================================================
# PAGE 17: Inference - Autoregressive Generation
# ============================================================
print("Creating Page 17: Autoregressive Generation...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Inference: Autoregressive Generation")

# Left: Without KV Cache (naive)
add_textbox(slide, 0.3, 0.85, 3.0, 0.3, "Without KV Cache:",
            font_size=11, bold=True, color=C_DECODER)

nc_steps = [
    ("Compute [START] K,V", 0.85),
    ("Recompute [START]+wo K,V", 1.55),
    ("Recompute [START]+wo+ai K,V", 2.25),
    ("Recompute all K,V", 2.95),
]
for text, y in nc_steps:
    add_shape_box(slide, 0.3, y, 3.2, 0.5, text,
                  fill_color=C_LIGHT_RED, text_color=C_DECODER, font_size=9)

add_textbox(slide, 0.3, 3.6, 3.2, 0.3, "x Redundant computation!",
            font_size=9, bold=True, color=C_DECODER, align=PP_ALIGN.CENTER)

# Right: With KV Cache
add_textbox(slide, 4.0, 0.85, 3.0, 0.3, "With KV Cache:",
            font_size=11, bold=True, color=C_FFN)

wc_steps = [
    ("Compute [START] K,V -> cache", 0.85),
    ("Only compute [wo] K,V + cache", 1.55),
    ("Only compute [ai] K,V + cache", 2.25),
    ("Only compute [ni] K,V + cache", 2.95),
]
for text, y in wc_steps:
    add_shape_box(slide, 4.0, y, 3.2, 0.5, text,
                  fill_color=C_LIGHT_GREEN, text_color=C_FFN, font_size=9)

add_textbox(slide, 4.0, 3.6, 3.2, 0.3, "check Only compute new word!",
            font_size=9, bold=True, color=C_FFN, align=PP_ALIGN.CENTER)

# KV Cache box
add_shape_box(slide, 4.0, 4.0, 3.2, 0.8, "",
              fill_color=C_LIGHT_GREEN, line_color=C_FFN)
add_multiline_tb(slide, 4.2, 4.05, 2.8, 0.7, [
    "KV Cache stores:",
    "K: [START], wo, ai, ni, ...",
    "V: [START], wo, ai, ni, ...",
    "Reuse across steps!",
], font_size=9, color=C_FFN)

# Bottom: Autoregressive steps illustration
ar_y = 5.0
add_textbox(slide, 0.3, ar_y, 7.0, 0.3,
            "Autoregressive Steps (translate \"I love you\"):",
            font_size=11, bold=True, color=C_TITLE)

ar_steps = [
    ("Step 1: [START] -> \"wo\"", C_LIGHT_BLUE),
    ("Step 2: [START, wo] -> \"ai\"", C_LIGHT_BLUE),
    ("Step 3: [START, wo, ai] -> \"ni\"", C_LIGHT_BLUE),
    ("Step 4: [START, wo, ai, ni] -> <END>", C_LIGHT_GREEN),
]
for i, (text, color) in enumerate(ar_steps):
    x = 0.3 + i * 1.75
    add_shape_box(slide, x, ar_y + 0.35, 1.6, 0.5, text,
                  fill_color=color, text_color=C_DARK, font_size=8)
    if i < len(ar_steps) - 1:
        add_arrow_right(slide, x + 1.6, ar_y + 0.6, x + 1.75, color=C_DARK, width=Pt(1.5))

# Final result
add_shape_box(slide, 1.5, ar_y + 1.0, 4.0, 0.45, "Final Output: wo ai ni",
              fill_color=C_FFN, font_size=12)

# Right reserved
add_rect(slide, 7.5, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 7.8, 3.0, 5.0, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 17 done.")

# ============================================================
# PAGE 18: Training vs Inference Comparison
# ============================================================
print("Creating Page 18: Training vs Inference Comparison...")
slide = prs.slides.add_slide(blank)
add_page_title(slide, "Training vs Inference: Full Comparison")

# Comparison table
table_x = 0.3
table_y = 0.85
col_w = [2.5, 3.0, 3.0]
row_h = 0.55
header_h = 0.45

# Headers
headers = ["Dimension", "Training (Open Book)", "Inference (Closed Book)"]
header_colors = [C_TITLE, C_ENCODER, C_FFN]
for j, (hdr, hc) in enumerate(zip(headers, header_colors)):
    x = table_x + sum(col_w[:j])
    add_shape_box(slide, x, table_y, col_w[j], header_h, hdr,
                  fill_color=hc, font_size=11)

# Rows
rows = [
    ("Decoder Input", "Full correct sequence\n(Teacher Forcing)", "Model's own predictions\n(autoregressive)"),
    ("Computation", "All positions PARALLEL", "One word at a time\n(SERIAL)"),
    ("Loss Function", "Cross-entropy loss\n(per position)", "No loss\n(prediction only)"),
    ("Gradients", "Backpropagation\n(update parameters)", "No gradients\n(parameters fixed)"),
    ("KV Cache", "Not needed\n(parallel compute)", "REQUIRED\n(avoid recomputation)"),
    ("Speed", "FAST\n(parallel GPU)", "SLOW\n(serial generation)"),
    ("Analogy", "Open-book exam:\nteacher gives answers", "Closed-book exam:\non your own"),
]

for i, (dim, train, infer) in enumerate(rows):
    y = table_y + header_h + i * row_h
    bg = C_WHITE if i % 2 == 0 else C_LIGHT_GRAY
    # Dimension
    add_shape_box(slide, table_x, y, col_w[0], row_h, dim,
                  fill_color=bg, text_color=C_DARK, font_size=9, bold=True)
    # Training
    add_shape_box(slide, table_x + col_w[0], y, col_w[1], row_h, train,
                  fill_color=C_BG_BLUE, text_color=C_ENCODER, font_size=8)
    # Inference
    add_shape_box(slide, table_x + col_w[0] + col_w[1], y, col_w[2], row_h, infer,
                  fill_color=C_BG_GREEN, text_color=C_FFN, font_size=8)

# Bottom summary
sum_y = table_y + header_h + len(rows) * row_h + 0.3
add_shape_box(slide, 0.3, sum_y, 8.5, 0.7, "",
              fill_color=C_LIGHT_ORANGE, line_color=C_ATTENTION)
add_multiline_tb(slide, 0.5, sum_y + 0.05, 8.1, 0.6, [
    "Key Insight: The fundamental difference is what the decoder sees.",
    "Training = correct answers (fast but may cause exposure bias)",
    "Inference = model's own output (slow but realistic)",
], font_size=9, color=C_ATTENTION)

# Right reserved
add_rect(slide, 9.2, 0.8, 0.01, 6.2, C_LIGHT_GRAY)
add_textbox(slide, 9.5, 3.0, 3.5, 1.5, "[Text Area]",
            font_size=12, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

print("  Page 18 done.")

# ============================================================
# Save and Self-Check
# ============================================================
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer-v3-diagrams.pptx"
prs.save(output_path)
print(f"\nSaved to: {output_path}")

# Self-check: verify all slides
print("\n" + "=" * 60)
print("SELF-CHECK REPORT")
print("=" * 60)
print(f"Total slides: {len(prs.slides)}")
print(f"Slide size: {SLIDE_W} x {SLIDE_H} inches")
print(f"Safe bounds: right={SAFE_R}, bottom={SAFE_B}")
print()

total_shapes = 0
overflow_count = 0
for i, slide in enumerate(prs.slides):
    shapes = list(slide.shapes)
    total_shapes += len(shapes)
    slide_overflows = 0
    for shape in shapes:
        l = shape.left / 914400  # EMU to inches
        t = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        if l + w > SAFE_R + 0.05:
            slide_overflows += 1
            print(f"  WARN: Slide {i+1}, right overflow: {l+w:.2f} (safe={SAFE_R})")
        if t + h > SAFE_B + 0.05:
            slide_overflows += 1
            print(f"  WARN: Slide {i+1}, bottom overflow: {t+h:.2f} (safe={SAFE_B})")
        if w > 5.0:
            slide_overflows += 1
            print(f"  WARN: Slide {i+1}, suspicious width: {w:.2f}")
        if h > 5.0:
            slide_overflows += 1
            print(f"  WARN: Slide {i+1}, suspicious height: {h:.2f}")
    overflow_count += slide_overflows
    print(f"Slide {i+1}: {len(shapes)} shapes, {slide_overflows} issues")

print()
print(f"Total shapes: {total_shapes}")
print(f"Total overflow issues: {overflow_count}")
if overflow_count == 0:
    print("ALL CHECKS PASSED!")
else:
    print(f"WARNING: {overflow_count} issues found. Please review.")
print("=" * 60)
