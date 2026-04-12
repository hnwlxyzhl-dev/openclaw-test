#!/usr/bin/env python3
"""Transformer v3 PPT Quality Assurance Script
检查7个维度，输出详细的质检数据。"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import re

PPTX_PATH = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3.pptx"
SLIDE_W = 13.333
SLIDE_H = 7.5

# Boundary tolerances from quality standards
MAX_RIGHT = 13.333  # title bar can equal this
MAX_BOTTOM = 7.5    # background rects can equal this
MIN_FONT_PT = 9     # minimum font size

prs = Presentation(PPTX_PATH)

results = {
    "slide_count": len(prs.slides),
    "slides": [],
    "summary": {
        "boundary_violations": [],
        "small_font_issues": [],
        "empty_slides": [],
        "text_overflow": [],
    }
}

total_shapes = 0
total_text_frames = 0

for slide_idx, slide in enumerate(prs.slides):
    slide_data = {
        "slide_num": slide_idx + 1,
        "shapes": [],
        "issues": [],
        "text_content": [],  # all text on slide for content review
    }
    
    has_title = False
    has_content = False
    shape_count = 0
    text_frame_count = 0
    
    for shape in slide.shapes:
        shape_count += 1
        total_shapes += 1
        
        # Get shape bounds
        left = shape.left / 914400.0 if shape.left else 0  # EMU to inches
        top = shape.top / 914400.0 if shape.top else 0
        width = shape.width / 914400.0 if shape.width else 0
        height = shape.height / 914400.0 if shape.height else 0
        right = left + width
        bottom = top + height
        
        shape_info = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": round(left, 3),
            "top": round(top, 3),
            "width": round(width, 3),
            "height": round(height, 3),
            "right": round(right, 3),
            "bottom": round(bottom, 3),
            "has_text": shape.has_text_frame if hasattr(shape, 'has_text_frame') else False,
        }
        
        # Check if it's a title (typically full width bar)
        is_title_bar = (width >= 13.0 and height <= 1.2)
        is_background = (width >= 13.0 and height >= 7.0)
        
        # Boundary checks
        boundary_ok = True
        if right > MAX_RIGHT + 0.01 and not is_title_bar and not is_background:
            issue = f"  ⚠️ OVERFLOW RIGHT: '{shape.name}' right={right:.3f} > {MAX_RIGHT} (left={left:.3f}, width={width:.3f})"
            slide_data["issues"].append(issue)
            results["summary"]["boundary_violations"].append(f"Slide {slide_idx+1}: {issue}")
            boundary_ok = False
        
        if bottom > MAX_BOTTOM + 0.01 and not is_background:
            issue = f"  ⚠️ OVERFLOW BOTTOM: '{shape.name}' bottom={bottom:.3f} > {MAX_BOTTOM} (top={top:.3f}, height={height:.3f})"
            slide_data["issues"].append(issue)
            results["summary"]["boundary_violations"].append(f"Slide {slide_idx+1}: {issue}")
            boundary_ok = False
        
        # Text analysis
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text_frame_count += 1
            total_text_frames += 1
            has_content = True
            
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                para_text = para.text.strip()
                if para_text:
                    slide_data["text_content"].append(para_text)
                    
                    # Check font sizes
                    for run in para.runs:
                        font_size = None
                        if run.font.size:
                            font_size = run.font.size / 12700.0  # EMU to pt
                        shape_info["font_sizes"] = shape_info.get("font_sizes", [])
                        if font_size:
                            shape_info["font_sizes"].append(round(font_size, 1))
                            if font_size < MIN_FONT_PT and not is_title_bar:
                                issue = f"  ⚠️ SMALL FONT: '{shape.name}' font={font_size:.1f}pt < {MIN_FONT_PT}pt, text='{para_text[:30]}...'"
                                slide_data["issues"].append(issue)
                                results["summary"]["small_font_issues"].append(f"Slide {slide_idx+1}: {issue}")
            
            if is_title_bar:
                has_title = True
        
        # Check for connectors/arrows
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if 'Arrow' in shape.name or 'arrow' in shape.name or 'Connector' in shape.name:
                shape_info["is_arrow"] = True
        
        slide_data["shapes"].append(shape_info)
    
    slide_data["has_title"] = has_title
    slide_data["shape_count"] = shape_count
    slide_data["text_frame_count"] = text_frame_count
    
    if not has_content and slide_idx > 0:  # skip first slide if it's cover
        results["summary"]["empty_slides"].append(f"Slide {slide_idx+1}")
    
    results["slides"].append(slide_data)

# Print summary
print("=" * 70)
print(f"TRANSFORMER v3 PPT QA REPORT")
print(f"Total slides: {len(prs.slides)}")
print(f"Total shapes: {total_shapes}")
print(f"Total text frames: {total_text_frames}")
print("=" * 70)

print(f"\n--- BOUNDARY VIOLATIONS ({len(results['summary']['boundary_violations'])}) ---")
for v in results["summary"]["boundary_violations"]:
    print(v)

print(f"\n--- SMALL FONT ISSUES ({len(results['summary']['small_font_issues'])}) ---")
for f in results["summary"]["small_font_issues"]:
    print(f)

print(f"\n--- EMPTY/LOW-CONTENT SLIDES ({len(results['summary']['empty_slides'])}) ---")
for s in results["summary"]["empty_slides"]:
    print(s)

print("\n\n--- PER-SLIDE DETAIL ---")
for sd in results["slides"]:
    print(f"\n{'='*60}")
    print(f"SLIDE {sd['slide_num']} ({sd['shape_count']} shapes, {sd['text_frame_count']} text frames)")
    if sd['issues']:
        print(f"ISSUES ({len(sd['issues'])}):")
        for issue in sd['issues']:
            print(issue)
    else:
        print("✅ No automated issues detected")
    
    print(f"TEXT CONTENT ({len(sd['text_content'])} paragraphs):")
    for t in sd['text_content']:
        print(f"  | {t[:80]}{'...' if len(t)>80 else ''}")

# Print shape details for boundary analysis
print("\n\n--- ALL SHAPE POSITIONS (for boundary check) ---")
for sd in results["slides"]:
    print(f"\nSlide {sd['slide_num']}:")
    for sh in sd["shapes"]:
        right_ok = "✅" if sh["right"] <= MAX_RIGHT + 0.01 else "❌"
        bottom_ok = "✅" if sh["bottom"] <= MAX_BOTTOM + 0.01 else "❌"
        fonts = sh.get("font_sizes", [])
        font_str = f"fonts={fonts}" if fonts else ""
        print(f"  {right_ok}{bottom_ok} {sh['name']:30s} L={sh['left']:6.3f} T={sh['top']:5.3f} W={sh['width']:6.3f} H={sh['height']:5.3f} R={sh['right']:6.3f} B={sh['bottom']:5.3f} {font_str}")

# Save structured results
with open("/home/admin/.openclaw/workspace-weaver/memory/transformer-v3-qa-data.json", "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)

print("\n\nStructured data saved to memory/transformer-v3-qa-data.json")
