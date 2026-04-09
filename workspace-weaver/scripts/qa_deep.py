#!/usr/bin/env python3
"""深入质检 - 检查字体大小、内容质量"""
from pptx import Presentation
from pptx.util import Pt, Emu

PPT_PATH = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v1.pptx"
prs = Presentation(PPT_PATH)

print("=" * 70)
print("🔍 深入质检 - 字体 & 内容分析")
print("=" * 70)

all_fonts = []
all_texts = []
slide_text_summary = []

for idx, slide in enumerate(prs.slides):
    si = idx + 1
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full_text = para.text.strip()
                if not full_text:
                    continue
                # Get font size from first run, or from paragraph level
                fs = None
                for run in para.runs:
                    if run.font.size:
                        fs = run.font.size / 12700
                        break
                if fs is None and para.font.size:
                    fs = para.font.size / 12700
                texts.append({"text": full_text, "font_size": fs, "para_level": para.level})
                all_texts.append({"slide": si, "text": full_text, "font_size": fs})
                if fs:
                    all_fonts.append(fs)
    
    total_chars = sum(len(t["text"]) for t in texts)
    slide_text_summary.append({
        "slide": si,
        "text_blocks": len(texts),
        "total_chars": total_chars,
        "texts": texts
    })

# Font size analysis
print(f"\n--- 字体大小统计 ---")
if all_fonts:
    print(f"  最小字体: {min(all_fonts):.1f}pt")
    print(f"  最大字体: {max(all_fonts):.1f}pt")
    print(f"  平均字体: {sum(all_fonts)/len(all_fonts):.1f}pt")
    
    # Count fonts below thresholds
    below_9 = sum(1 for f in all_fonts if f < 9)
    below_10 = sum(1 for f in all_fonts if f < 10)
    below_11 = sum(1 for f in all_fonts if f < 11)
    print(f"  < 9pt: {below_9}个")
    print(f"  < 10pt: {below_10}个")
    print(f"  < 11pt: {below_11}个")
else:
    print("  未检测到显式字体大小设置（可能使用母版默认值）")

# Content quality per slide
print(f"\n--- 各页文本量分析 ---")
for s in slide_text_summary:
    avg_len = s["total_chars"] / max(s["text_blocks"], 1)
    bar = "█" * min(int(s["total_chars"] / 50), 30)
    print(f"  第{s['slide']:2d}页: {s['text_blocks']:3d}个文本块, {s['total_chars']:5d}字, 均长{avg_len:.0f}字 {bar}")

# Check for specific content quality indicators
print(f"\n--- 内容质量指标 ---")
has_metaphor = 0  # 比喻
has_number = 0    # 数字
has_long_text = 0 # 详细解释

for t in all_texts:
    text = t["text"]
    # Metaphor indicators
    if any(w in text for w in ["比喻", "就像", "好比", "想象", "如同", "图书馆", "考试", "桥梁", "专家", "快递"]):
        has_metaphor += 1
    # Numbers
    if any(c.isdigit() for c in text):
        has_number += 1
    # Long text (>50 chars = detailed explanation)
    if len(text) > 50:
        has_long_text += 1

print(f"  含比喻的文本: {has_metaphor}/{len(all_texts)} ({has_metaphor/len(all_texts)*100:.0f}%)")
print(f"  含数字的文本: {has_number}/{len(all_texts)} ({has_number/len(all_texts)*100:.0f}%)")
print(f"  长文本(>50字): {has_long_text}/{len(all_texts)} ({has_long_text/len(all_texts)*100:.0f}%)")

# Slide 1 shape details (the problematic one)
print(f"\n--- 第1页 Shape 详情 ---")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    left = shape.left / 914400
    top = shape.top / 914400
    w = shape.width / 914400
    h = shape.height / 914400
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text[:40]
    print(f"  [{shape.shape_type}] name={shape.name[:30]:30s} L={left:.2f} T={top:.2f} W={w:.2f} H={h:.2f} R={left+w:.2f} B={top+h:.2f} | {text}")

# Check shapes that are backgrounds (full-slide shapes)
print(f"\n--- 背景形状检测 ---")
for idx, slide in enumerate(prs.slides):
    si = idx + 1
    for shape in slide.shapes:
        w = shape.width / 914400
        h = shape.height / 914400
        if w >= 13.0 and h >= 7.0:
            print(f"  第{si}页: 全屏背景 shape [{shape.name}] W={w:.2f} H={h:.2f}")
