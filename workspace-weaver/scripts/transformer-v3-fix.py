#!/usr/bin/env python3
"""
Transformer v3 PPT Fix Script
Fixes:
  P0: Slide 18-19 content swap (word selection vs training/inference comparison)
  P0: Add diagrams to slides 8, 16, 20, 21
  P1: Fix 6 missing-space typos
  P1: Remove 13 [Text Area] placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# ============================================================
# Constants
# ============================================================
SLIDE_W = 12191695  # ~13.333 inches
SLIDE_H = 6858000   # 7.5 inches
MAX_RIGHT = 12192000  # ~13.2 inches
MAX_BOTTOM = 6667500  # ~7.3 inches

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTN = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xD0)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
C_BG_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
C_DARK_BLUE = RGBColor(0x1A, 0x5C, 0x8A)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

INPUT_FILE = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3.pptx"
OUTPUT_FILE = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3.pptx"


def inch(v):
    """Convert inches to EMU"""
    return int(v * 914400)


def pt(v):
    """Convert points to EMU"""
    return int(v * 12700)


def add_rounded_rect(slide, left, top, width, height, fill_color, text, font_size=10, font_color=C_TEXT, bold=False, alignment=PP_ALIGN.CENTER):
    """Add a rounded rectangle with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_EN
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=10, font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_EN
    return txBox


def add_arrow(slide, left, top, width, height, fill_color=C_TEXT):
    """Add a right-pointing arrow"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_connector_line(slide, x1, y1, x2, y2, color=C_TEXT, width_pt=1.5):
    """Add a straight connector line"""
    from pptx.oxml.ns import qn
    cxnSp = slide.shapes._spTree.makeelement(qn("p:cxnSp"), {})
    nvSpPr = cxnSp.makeelement(qn("p:nvSpPr"), {})
    cNvPr = nvSpPr.makeelement(qn("p:cNvPr"), {"id": str(9999), "name": "conn"})
    nvSpPr.append(cNvPr)
    cNvCxnPr = nvSpPr.makeelement(qn("p:cNvCxnPr"), {})
    nvSpPr.append(cNvCxnPr)
    cxnSp.append(nvSpPr)
    spPr = cxnSp.makeelement(qn("p:spPr"), {})
    xfrm = spPr.makeelement(qn("a:xfrm"), {})
    xfrm.set("flipV", "0")
    xfrm.set("flipH", "0")
    spPr.append(xfrm)
    prstGeom = spPr.makeelement(qn("a:prstGeom"), {"prst": "line"})
    prstGeom.append(prstGeom.makeelement(qn("a:avLst"), {}))
    spPr.append(prstGeom)
    ln = spPr.makeelement(qn("a:ln"), {"w": str(int(width_pt * 12700))})
    solidFill = ln.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2]) if isinstance(color, (tuple, list)) else "%02X%02X%02X" % (color.red if hasattr(color, 'red') else 0, color.green if hasattr(color, 'green') else 0, color.blue if hasattr(color, 'blue') else 0)})
    solidFill.append(srgbClr)
    ln.append(solidFill)
    spPr.append(ln)
    cxnSp.append(spPr)
    stCxn = cxnSp.makeelement(qn("p:stCxn"), {"id": "0", "idx": "0"})
    cxnSp.append(stCxn)
    endCxn = cxnSp.makeelement(qn("p:endCxn"), {"id": "0", "idx": "0"})
    cxnSp.append(endCxn)
    slide.shapes._spTree.append(cxnSp)
    return cxnSp


def fix_missing_spaces(prs):
    """P1: Fix 6 missing-space typos"""
    fixes = {
        "forprevious": "for previous",
        "overdistance": "over distance",
        "gradientvanishing": "gradient vanishing",
        "compressedinto": "compressed into",
        "seeprevious": "see previous",
        "512-dwith": "512-d with",
        "-infSoftmax": "-inf Softmax",
        "-inf Softmax": "-inf Softmax",  # already ok
    }
    count = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        original = run.text
                        for bad, good in fixes.items():
                            if bad in original and bad != good:
                                run.text = original.replace(bad, good)
                                if run.text != original:
                                    count += 1
    return count


def remove_text_area_placeholders(prs):
    """P1: Remove 13 [Text Area] placeholders"""
    count = 0
    for slide_idx, slide in enumerate(prs.slides):
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text == "[Text Area]":
                    shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            count += 1
    return count


def fix_slide_18_word_selection(prs):
    """P0: Redesign slide 18 to be Word Selection Strategies only"""
    slide = prs.slides[17]  # 0-indexed
    
    # Remove all shapes except title bar, page number, and footer
    shapes_to_remove = []
    for shape in slide.shapes:
        left = shape.left
        top = shape.top
        # Keep title bar (top ~0.0-0.85), page number, footer (top ~7.05)
        if top < inch(0.85) or (top > inch(7.0) and top < inch(7.5)):
            continue
        shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Now build 5 strategy cards on the left side
    strategies = [
        ("Greedy Search", "Always pick the highest probability word. Fast but prone to local optima.", "No randomness", C_ENCODER),
        ("Beam Search", "Keep top-k candidate sequences simultaneously (beam size 4-10). Good for translation.", "Low randomness", C_ATTN),
        ("Top-K Sampling", "Random sample from top-k words (e.g., k=50). Balances quality and diversity.", "Medium randomness", C_FFN),
        ("Top-P (Nucleus)", "Sample from words whose cumulative probability reaches p (e.g., 0.9). Adapts dynamically.", "Adaptive randomness", C_EMBED),
        ("Temperature", "Scale logits before Softmax. Low (0.3) = conservative. High (1.5) = creative/diverse.", "Adjustable randomness", C_DECODER),
    ]
    
    card_w = inch(7.5)
    card_h = inch(0.95)
    card_x = inch(0.30)
    start_y = inch(1.0)
    gap = inch(0.12)
    
    for i, (title, desc, randomness, color) in enumerate(strategies):
        y = start_y + i * (card_h + gap)
        if y + card_h > MAX_BOTTOM:
            break
        
        # Card background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y, card_w, card_h
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_LIGHT_GRAY
        bg.line.color.rgb = color
        bg.line.width = Pt(2)
        bg.shadow.inherit = False
        
        # Color accent bar on the left
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_x, y, inch(0.08), card_h
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        accent.shadow.inherit = False
        
        # Number badge
        num_x = card_x + inch(0.20)
        num_y = y + inch(0.15)
        num_size = inch(0.30)
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, num_x, num_y, num_size, num_size
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = color
        num_shape.line.fill.background()
        num_shape.shadow.inherit = False
        tf = num_shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(12)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Strategy name
        title_box = add_textbox(slide, num_x + inch(0.40), y + inch(0.08), inch(3.0), inch(0.30), 
                                title, font_size=12, font_color=C_TITLE, bold=True)
        
        # Description
        desc_box = add_textbox(slide, num_x + inch(0.40), y + inch(0.35), inch(5.5), inch(0.55),
                               desc, font_size=9, font_color=C_TEXT)
        
        # Randomness tag (right side)
        tag_x = card_x + card_w - inch(2.0)
        tag_y = y + inch(0.30)
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, tag_x, tag_y, inch(1.8), inch(0.30)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tag.shadow.inherit = False
        tf = tag.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = randomness
        run.font.size = Pt(8)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
    
    # Right side: Temperature visual explanation
    right_x = inch(8.2)
    
    # Temperature section title
    add_textbox(slide, right_x, inch(1.0), inch(4.8), inch(0.35),
                "Temperature Effect", font_size=13, font_color=C_TITLE, bold=True)
    
    # Low temperature box
    add_textbox(slide, right_x, inch(1.5), inch(4.8), inch(0.25),
                "Low (0.3) - Conservative", font_size=10, font_color=C_ENCODER, bold=True)
    add_textbox(slide, right_x, inch(1.75), inch(4.8), inch(0.40),
                "Sharp distribution: [cat: 0.85, dog: 0.10, bird: 0.03, ...] Almost always picks the most likely word.", 
                font_size=9, font_color=C_TEXT)
    
    # High temperature box
    add_textbox(slide, right_x, inch(2.3), inch(4.8), inch(0.25),
                "High (1.5) - Creative", font_size=10, font_color=C_DECODER, bold=True)
    add_textbox(slide, right_x, inch(2.55), inch(4.8), inch(0.40),
                "Flat distribution: [cat: 0.25, dog: 0.22, bird: 0.20, ...] More unpredictable, diverse outputs.",
                font_size=9, font_color=C_TEXT)
    
    # ChatGPT default
    chatgpt_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, inch(3.1), inch(4.5), inch(0.45)
    )
    chatgpt_box.fill.solid()
    chatgpt_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    chatgpt_box.line.color.rgb = C_FFN
    chatgpt_box.line.width = Pt(1)
    chatgpt_box.shadow.inherit = False
    tf = chatgpt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ChatGPT default: Temperature=0.7 + Top-P=0.9"
    run.font.size = Pt(10)
    run.font.color.rgb = C_TEXT
    run.font.bold = True
    run.font.name = FONT_EN
    
    # Bottom key insight
    insight_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.30), inch(6.30), inch(12.60), inch(0.60)
    )
    insight_bg.fill.solid()
    insight_bg.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
    insight_bg.line.color.rgb = C_ATTN
    insight_bg.line.width = Pt(1.5)
    insight_bg.shadow.inherit = False
    tf = insight_bg.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Key Insight: Greedy = always safe but boring. Beam Search = good for translation. Top-P + Temperature = best for chat/creative writing. The choice depends on the task!"
    run.font.size = Pt(9)
    run.font.color.rgb = C_TEXT
    run.font.name = FONT_EN
    
    return "Slide 18 redesigned: Word Selection Strategies with 5 cards"


def add_training_diagram(slide):
    """Add training data flow panorama to slide 8"""
    # Layout: top area (y=1.0 to 4.0) for diagram, below for existing text
    
    y_base = inch(1.0)
    
    # Section label
    add_textbox(slide, inch(0.30), y_base, inch(3.0), inch(0.30),
                "Training Data Flow:", font_size=12, font_color=C_TITLE, bold=True)
    
    # === Left: Encoder Input ===
    enc_in_y = y_base + inch(0.40)
    add_rounded_rect(slide, inch(0.40), enc_in_y, inch(2.2), inch(0.45),
                     C_LIGHT_BLUE, "Input: I love you", font_size=10, font_color=C_TEXT, bold=True)
    add_arrow(slide, inch(1.30), enc_in_y + inch(0.50), inch(0.50), inch(0.25), C_ENCODER)
    
    enc_y = enc_in_y + inch(0.80)
    add_rounded_rect(slide, inch(0.40), enc_y, inch(2.2), inch(0.80),
                     C_ENCODER, "Encoder\n(6 layers)", font_size=11, font_color=C_WHITE, bold=True)
    
    # === Middle: Encoder Output (K, V) ===
    kv_y = enc_y + inch(0.10)
    add_rounded_rect(slide, inch(3.0), kv_y, inch(1.5), inch(1.0),
                     C_LIGHT_PURPLE, "Encoder\nOutput\n(K, V)", font_size=10, font_color=C_TEXT, bold=True)
    
    # Arrow from encoder to KV
    add_arrow(slide, inch(2.60), enc_y + inch(0.25), inch(0.40), inch(0.20), C_ENCODER)
    
    # === Right: Decoder Input ===
    dec_in_y = y_base + inch(0.40)
    add_rounded_rect(slide, inch(5.0), dec_in_y, inch(3.0), inch(0.45),
                     C_LIGHT_GREEN, "Decoder Input: [START, wo, ai]", font_size=10, font_color=C_TEXT, bold=True)
    add_arrow(slide, inch(6.20), dec_in_y + inch(0.50), inch(0.50), inch(0.25), C_FFN)
    
    dec_y = dec_in_y + inch(0.80)
    add_rounded_rect(slide, inch(5.0), dec_y, inch(3.0), inch(0.80),
                     C_FFN, "Decoder\n(6 layers)", font_size=11, font_color=C_WHITE, bold=True)
    
    # Arrow from KV to Decoder (cross-attention)
    add_arrow(slide, inch(4.50), kv_y + inch(0.40), inch(0.50), inch(0.20), C_EMBED)
    
    # === Right: Predictions ===
    pred_y = dec_y + inch(0.95)
    add_arrow(slide, inch(6.20), dec_y + inch(0.80), inch(0.50), inch(0.25), C_ATTN)
    add_rounded_rect(slide, inch(5.0), pred_y, inch(3.0), inch(0.45),
                     C_LIGHT_ORANGE, "Predictions: [wo, ai, ni, END]", font_size=10, font_color=C_TEXT, bold=True)
    
    # === Loss comparison ===
    loss_y = pred_y + inch(0.55)
    add_rounded_rect(slide, inch(5.0), loss_y, inch(3.0), inch(0.45),
                     C_LIGHT_RED, "Targets: [wo, ai, ni, END]", font_size=10, font_color=C_TEXT, bold=True)
    
    # Loss box
    loss_box_y = loss_y + inch(0.55)
    add_rounded_rect(slide, inch(5.5), loss_box_y, inch(2.0), inch(0.40),
                     C_DECODER, "Cross-Entropy Loss", font_size=10, font_color=C_WHITE, bold=True)
    
    # Backprop arrow
    bp_y = loss_box_y + inch(0.50)
    add_rounded_rect(slide, inch(5.0), bp_y, inch(3.0), inch(0.40),
                     C_LIGHT_GRAY, "Backpropagation -> Update Parameters", font_size=9, font_color=C_TEXT, bold=True)
    
    # === Bottom annotation ===
    note_y = bp_y + inch(0.55)
    note = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.30), note_y, inch(8.0), inch(0.35)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
    note.line.color.rgb = C_ATTN
    note.line.width = Pt(1)
    note.shadow.inherit = False
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "All positions computed in PARALLEL (Open-book Exam)"
    run.font.size = Pt(10)
    run.font.color.rgb = C_TEXT
    run.font.bold = True
    run.font.name = FONT_EN
    
    # Move existing text blocks down to make room
    # Existing text starts at top=4.20, we need to move them below the diagram
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top >= inch(4.0) and shape.top < inch(7.0):
            text_shapes.append(shape)
    
    # The diagram goes from y=1.0 to about y=4.0. Move text to start at 4.10
    new_text_start = inch(4.10)
    for shape in text_shapes:
        # These are the 3 section headers + 3 body texts
        orig_top = shape.top
        # Keep relative spacing
        shape.top = new_text_start + (orig_top - inch(4.20))
    
    return "Training diagram added to slide 8"


def add_inference_diagram(slide):
    """Add autoregressive generation steps diagram to slide 16"""
    y_base = inch(1.0)
    
    # Section label
    add_textbox(slide, inch(0.30), y_base, inch(3.0), inch(0.30),
                "Autoregressive Generation Steps:", font_size=12, font_color=C_TITLE, bold=True)
    
    # 4 steps side by side
    steps = [
        ("[START]", "wo", "Step 1"),
        ("[START, wo]", "ai", "Step 2"),
        ("[START, wo, ai]", "ni", "Step 3"),
        ("[START, wo, ai, ni]", "<END>", "Step 4"),
    ]
    
    step_w = inch(2.7)
    step_h = inch(2.2)
    start_x = inch(0.40)
    gap = inch(0.20)
    
    for i, (input_text, output_text, step_name) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        y = y_base + inch(0.40)
        
        # Step container
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h
        )
        container.fill.solid()
        container.fill.fore_color.rgb = C_LIGHT_GRAY
        container.line.color.rgb = C_ENCODER if i < 3 else C_FFN
        container.line.width = Pt(1.5)
        container.shadow.inherit = False
        
        # Step label
        step_label = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, inch(0.70), inch(0.30)
        )
        step_label.fill.solid()
        step_label.fill.fore_color.rgb = C_ENCODER if i < 3 else C_FFN
        step_label.line.fill.background()
        step_label.shadow.inherit = False
        tf = step_label.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_name
        run.font.size = Pt(8)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Decoder input box
        dec_in = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + inch(0.10), y + inch(0.40), step_w - inch(0.20), inch(0.50)
        )
        dec_in.fill.solid()
        dec_in.fill.fore_color.rgb = C_LIGHT_GREEN
        dec_in.line.color.rgb = C_FFN
        dec_in.line.width = Pt(1)
        dec_in.shadow.inherit = False
        tf = dec_in.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Input: " + input_text
        run.font.size = Pt(8)
        run.font.color.rgb = C_TEXT
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Decoder box
        dec = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + inch(0.30), y + inch(1.0), step_w - inch(0.60), inch(0.40)
        )
        dec.fill.solid()
        dec.fill.fore_color.rgb = C_FFN
        dec.line.fill.background()
        dec.shadow.inherit = False
        tf = dec.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Decoder"
        run.font.size = Pt(9)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Output box
        out = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + inch(0.10), y + inch(1.50), step_w - inch(0.20), inch(0.45)
        )
        out.fill.solid()
        out.fill.fore_color.rgb = C_LIGHT_ORANGE
        out.line.color.rgb = C_ATTN
        out.line.width = Pt(1)
        out.shadow.inherit = False
        tf = out.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Output: " + output_text
        run.font.size = Pt(9)
        run.font.color.rgb = C_TEXT
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Arrow between steps
        if i < 3:
            arrow_x = x + step_w
            arrow_y = y + step_h / 2
            add_arrow(slide, arrow_x, arrow_y - inch(0.08), gap - inch(0.05), inch(0.16), C_ENCODER)
    
    # Encoder note at top-right
    enc_note = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(9.0), y_base + inch(0.40), inch(3.8), inch(0.50)
    )
    enc_note.fill.solid()
    enc_note.fill.fore_color.rgb = C_LIGHT_BLUE
    enc_note.line.color.rgb = C_ENCODER
    enc_note.line.width = Pt(1)
    enc_note.shadow.inherit = False
    tf = enc_note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Encoder runs ONCE -> Output cached"
    run.font.size = Pt(10)
    run.font.color.rgb = C_TEXT
    run.font.bold = True
    run.font.name = FONT_EN
    
    # Serial generation annotation
    serial_note = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.30), y_base + inch(2.75), inch(11.8), inch(0.35)
    )
    serial_note.fill.solid()
    serial_note.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    serial_note.line.color.rgb = C_DECODER
    serial_note.line.width = Pt(1)
    serial_note.shadow.inherit = False
    tf = serial_note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SERIAL generation: must complete one step before starting the next (Closed-book Exam)"
    run.font.size = Pt(10)
    run.font.color.rgb = C_DECODER
    run.font.bold = True
    run.font.name = FONT_EN
    
    # Move existing text blocks down
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top >= inch(4.0) and shape.top < inch(7.0):
            text_shapes.append(shape)
    
    new_text_start = inch(4.20)
    for shape in text_shapes:
        orig_top = shape.top
        shape.top = new_text_start + (orig_top - inch(4.20))
    
    return "Inference diagram added to slide 16"


def add_family_tree(slide):
    """Add Transformer family tree diagram to slide 20"""
    # The left side (0.3 to 7.3) is empty, fill it with the tree
    
    # Central node: Transformer 2017
    center_x = inch(3.5)
    center_y = inch(1.2)
    
    central = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, center_x - inch(1.2), center_y, inch(2.4), inch(0.65)
    )
    central.fill.solid()
    central.fill.fore_color.rgb = C_TITLE
    central.line.fill.background()
    central.shadow.inherit = False
    tf = central.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Transformer 2017"
    run.font.size = Pt(12)
    run.font.color.rgb = C_WHITE
    run.font.bold = True
    run.font.name = FONT_EN
    
    # Three branches
    branches = [
        ("Decoder-Only", "GPT-1 / GPT-2 / GPT-3\nGPT-4 / ChatGPT\nLLaMA / Qwen", 
         "Improvisational Speaker", C_ENCODER, inch(0.3)),
        ("Encoder-Decoder", "T5 / BART\nWhisper / mBART\nFlan-T5",
         "All-purpose Translator", C_FFN, inch(2.6)),
        ("Encoder-Only", "BERT / RoBERTa\nDeBERTa / ALBERT\nELECTRA",
         "Reading Comprehension Expert", C_EMBED, inch(4.9)),
    ]
    
    for title, models, metaphor, color, bx in branches:
        # Type box
        type_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bx, inch(2.3), inch(2.0), inch(0.50)
        )
        type_box.fill.solid()
        type_box.fill.fore_color.rgb = color
        type_box.line.fill.background()
        type_box.shadow.inherit = False
        tf = type_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(10)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Models list
        models_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bx, inch(2.90), inch(2.0), inch(1.2)
        )
        models_box.fill.solid()
        models_box.fill.fore_color.rgb = C_LIGHT_GRAY
        models_box.line.color.rgb = color
        models_box.line.width = Pt(1)
        models_box.shadow.inherit = False
        tf = models_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = models
        run.font.size = Pt(9)
        run.font.color.rgb = C_TEXT
        run.font.name = FONT_EN
        
        # Metaphor tag
        meta_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bx, inch(4.20), inch(2.0), inch(0.35)
        )
        meta_box.fill.solid()
        meta_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
        meta_box.line.color.rgb = C_ATTN
        meta_box.line.width = Pt(1)
        meta_box.shadow.inherit = False
        tf = meta_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = metaphor
        run.font.size = Pt(8)
        run.font.color.rgb = C_TEXT
        run.font.bold = True
        run.font.name = FONT_EN
    
    # Scale comparison bar at bottom
    scale_y = inch(4.85)
    add_textbox(slide, inch(0.30), scale_y, inch(3.0), inch(0.25),
                "Parameter Scale Growth:", font_size=10, font_color=C_TITLE, bold=True)
    
    models_scale = [
        ("65M", "2017", 0.3),
        ("110M", "2018", 0.9),
        ("1.5B", "2019", 1.8),
        ("175B", "2020", 3.0),
        ("1.8T", "2023", 5.0),
    ]
    
    bar_y = scale_y + inch(0.30)
    for name, year, width_in in models_scale:
        bar_w = inch(max(0.4, width_in))
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.30), bar_y, bar_w, inch(0.30)
        )
        # Gradient from blue to dark blue based on size
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(
            min(255, 0x34 + int(width_in * 20)),
            max(0, 0x98 - int(width_in * 15)),
            max(0, 0xDB - int(width_in * 20))
        )
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{name} ({year})"
        run.font.size = Pt(8)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        bar_y += inch(0.35)
    
    return "Family tree added to slide 20"


def add_impact_diagram(slide):
    """Add domain expansion radial diagram to slide 21"""
    # Top area (y=1.0 to 3.2) for the radial diagram
    
    # Central circle
    cx = inch(3.8)
    cy = inch(1.5)
    
    center = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - inch(0.8), cy, inch(1.6), inch(0.70)
    )
    center.fill.solid()
    center.fill.fore_color.rgb = C_TITLE
    center.line.fill.background()
    center.shadow.inherit = False
    tf = center.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Transformer"
    run.font.size = Pt(12)
    run.font.color.rgb = C_WHITE
    run.font.bold = True
    run.font.name = FONT_EN
    
    # 6 domain nodes arranged in a semi-circle
    domains = [
        ("NLP", "Translation, Chat,\nWriting", C_ENCODER, inch(0.30), inch(1.0)),
        ("Computer Vision", "ViT, Image Gen,\nDALL-E", C_DECODER, inch(5.50), inch(1.0)),
        ("Speech", "Whisper,\nVoice AI", C_FFN, inch(7.00), inch(1.8)),
        ("Multimodal", "GPT-4V, Gemini,\nSora", C_EMBED, inch(6.00), inch(2.7)),
        ("Protein", "AlphaFold2,\nDrug Design", RGBColor(0x27, 0xAE, 0x60), inch(0.80), inch(2.7)),
        ("Time Series", "Forecasting,\nRecommendation", C_ATTN, inch(0.00), inch(1.8)),
    ]
    
    for name, desc, color, dx, dy in domains:
        # Domain node
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, dx, dy, inch(1.8), inch(0.65)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        node.shadow.inherit = False
        tf = node.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(9)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
    
    # Timeline bar
    timeline_y = inch(3.05)
    add_textbox(slide, inch(0.30), timeline_y, inch(8.0), inch(0.25),
                "Timeline:", font_size=10, font_color=C_TITLE, bold=True)
    
    milestones = [
        ("2017", "Transformer", C_ENCODER),
        ("2018", "GPT-1/BERT", C_ATTN),
        ("2020", "GPT-3/ViT", C_FFN),
        ("2022", "ChatGPT/Whisper", C_DECODER),
        ("2023", "GPT-4/Gemini", C_EMBED),
        ("2024+", "Sora/o1/o3", RGBColor(0x27, 0xAE, 0x60)),
    ]
    
    tl_y = timeline_y + inch(0.30)
    for i, (year, event, color) in enumerate(milestones):
        tx = inch(0.30) + i * inch(1.40)
        
        # Year badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, tx, tl_y, inch(0.50), inch(0.25)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = year
        run.font.size = Pt(8)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        run.font.name = FONT_EN
        
        # Event text
        add_textbox(slide, tx, tl_y + inch(0.28), inch(1.30), inch(0.25),
                    event, font_size=7, font_color=C_TEXT)
    
    return "Impact diagram added to slide 21"


def self_check(prs):
    """Validate all shapes are within bounds"""
    issues = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            r = shape.left + shape.width
            b = shape.top + shape.height
            # Skip title bar and background
            if shape.top < inch(0.9) and shape.height > inch(0.8):
                continue
            if r > MAX_RIGHT + inch(0.1):
                issues.append(f"Slide {idx+1}: right={r/914400:.2f} > 13.2 (shape at left={shape.left/914400:.2f})")
            if b > MAX_BOTTOM + inch(0.1):
                issues.append(f"Slide {idx+1}: bottom={b/914400:.2f} > 7.3 (shape at top={shape.top/914400:.2f})")
    return issues


def main():
    print("=" * 60)
    print("Transformer v3 PPT Fix Script")
    print("=" * 60)
    
    prs = Presentation(INPUT_FILE)
    
    fix_log = []
    
    # === P0: Fix Slide 18 (Word Selection Strategies) ===
    print("\n[P0] Fixing Slide 18 - Word Selection Strategies...")
    msg = fix_slide_18_word_selection(prs)
    fix_log.append(msg)
    print(f"  Done: {msg}")
    
    # === P0: Add diagram to Slide 8 (Training Overview) ===
    print("\n[P0] Adding training diagram to Slide 8...")
    msg = add_training_diagram(prs.slides[7])
    fix_log.append(msg)
    print(f"  Done: {msg}")
    
    # === P0: Add diagram to Slide 16 (Inference Overview) ===
    print("\n[P0] Adding inference diagram to Slide 16...")
    msg = add_inference_diagram(prs.slides[15])
    fix_log.append(msg)
    print(f"  Done: {msg}")
    
    # === P0: Add diagram to Slide 20 (Three Families) ===
    print("\n[P0] Adding family tree to Slide 20...")
    msg = add_family_tree(prs.slides[19])
    fix_log.append(msg)
    print(f"  Done: {msg}")
    
    # === P0: Add diagram to Slide 21 (Impact) ===
    print("\n[P0] Adding impact diagram to Slide 21...")
    msg = add_impact_diagram(prs.slides[20])
    fix_log.append(msg)
    print(f"  Done: {msg}")
    
    # === P1: Fix missing spaces ===
    print("\n[P1] Fixing missing-space typos...")
    count = fix_missing_spaces(prs)
    fix_log.append(f"Fixed {count} missing-space typos")
    print(f"  Done: Fixed {count} typos")
    
    # === P1: Remove [Text Area] placeholders ===
    print("\n[P1] Removing [Text Area] placeholders...")
    count = remove_text_area_placeholders(prs)
    fix_log.append(f"Removed {count} [Text Area] placeholders")
    print(f"  Done: Removed {count} placeholders")
    
    # === Self-check ===
    print("\n[CHECK] Running boundary self-check...")
    issues = self_check(prs)
    if issues:
        print(f"  WARNING: {len(issues)} boundary issues found:")
        for i in issues:
            print(f"    {i}")
    else:
        print("  All shapes within bounds!")
    
    # === Save ===
    print(f"\nSaving to: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("Saved successfully!")
    
    # === Summary ===
    print("\n" + "=" * 60)
    print("FIX SUMMARY")
    print("=" * 60)
    for i, msg in enumerate(fix_log, 1):
        print(f"  {i}. {msg}")
    print(f"\nTotal fixes applied: {len(fix_log)}")
    print(f"Boundary issues: {len(issues)}")
    print("=" * 60)


if __name__ == "__main__":
    main()
