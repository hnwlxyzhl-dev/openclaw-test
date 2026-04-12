#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer PPT Part1 (Slides 1-11) - Full Chinese Version
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_R = Inches(13.2)
SAFE_B = Inches(7.3)

# Colors
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTN = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LINE = RGBColor(0xBD, 0xC3, 0xC7)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xC8)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_YELLOW = RGBColor(0xFC, 0xF3, 0xCF)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── Helper functions ──
def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else C_WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, color=C_TEXT,
                 bold=False, align=PP_ALIGN.LEFT, font_name=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name if font_name else FONT_CN
    p.alignment = align
    return txBox


def add_arrow(slide, left, top, width, height, color=C_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=C_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=C_GRAY, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_multiline_text(slide, left, top, width, height, lines, font_size=11, color=C_TEXT,
                       bold_first=False, line_spacing=1.2):
    """lines: list of (text, color, bold) or just text strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, clr, bld = line, color, False
        else:
            txt = line[0]
            clr = line[1] if len(line) > 1 else color
            bld = line[2] if len(line) > 2 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = FONT_CN
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


def label_in_rect(slide, left, top, width, height, text, font_size=9, color=C_TEXT,
                  fill=None, border=None, bold=False):
    shape = add_rect(slide, left, top, width, height, fill_color=fill, border_color=border)
    shape.text_frame.word_wrap = True
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = bold
    shape.text_frame.paragraphs[0].font.name = FONT_CN
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# ═══════════════════════════════════════════════════════════════
# Slide 1: Cover
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

# Decorative horizontal lines
for y_offset in [1.5, 2.0, 5.5, 6.0]:
    add_rect(slide, Inches(0.5), Inches(y_offset), Inches(12.333), Pt(1), fill_color=C_BG_LINE)

# Main title
add_text_box(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(1.2),
             "Transformer 架构深度解析", 36, C_TITLE, True, PP_ALIGN.CENTER, FONT_CN)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.333), Inches(0.8),
             "从直觉到原理，一篇看懂AI的核心引擎", 18, C_TEXT, False, PP_ALIGN.CENTER)

# Small description
add_text_box(slide, Inches(2), Inches(4.4), Inches(9.333), Inches(0.6),
             "2017年Google提出，ChatGPT、GPT-4等所有大语言模型的基石", 13, C_GRAY, False, PP_ALIGN.CENTER)

# Bottom line
add_text_box(slide, Inches(2), Inches(6.5), Inches(9.333), Inches(0.5),
             "从训练到推理，完整拆解Transformer的工作原理", 11, C_GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Slide 2: Table of Contents
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.7),
             "今天的内容路线", 22, C_TITLE, True)

# Left: 5 modules
modules = [
    ("1. 为什么需要Transformer", C_ENCODER),
    ("2. Transformer架构概述", C_ENCODER),
    ("3. 训练阶段详解", C_ATTN, "重点"),
    ("4. 推理阶段详解", C_FFN),
    ("5. 典型应用与总结", C_EMBED),
]

for i, (text, color, *extra) in enumerate(modules):
    y = Inches(1.5) + Inches(i * 1.1)
    label_in_rect(slide, Inches(0.8), y, Inches(5), Inches(0.7),
                  text, 14, color, fill=C_WHITE, border=color, bold=True)
    if extra and extra[0]:
        add_text_box(slide, Inches(6.0), y + Pt(2), Inches(1), Inches(0.6),
                     extra[0], 10, C_RED, True)

# Right: vertical roadmap
road_x = Inches(9.5)
for i in range(5):
    y = Inches(1.5) + Inches(i * 1.1) + Inches(0.2)
    colors = [C_ENCODER, C_ENCODER, C_ATTN, C_FFN, C_EMBED]
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, road_x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors[i]
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(i + 1)
    circle.text_frame.paragraphs[0].font.size = Pt(12)
    circle.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Connecting line between nodes
    if i < 4:
        add_line(slide, road_x + Inches(0.25), y + Inches(0.5),
                 road_x + Inches(0.25), y + Inches(1.1) - Inches(0.05), C_GRAY, Pt(2))

# Node labels
labels_right = ["痛点分析", "整体架构", "训练详解", "推理详解", "应用总结"]
for i, lbl in enumerate(labels_right):
    y = Inches(1.5) + Inches(i * 1.1) + Inches(0.28)
    add_text_box(slide, Inches(10.2), y, Inches(2.5), Inches(0.4),
                 lbl, 10, C_GRAY, False)


# ═══════════════════════════════════════════════════════════════
# Slide 3: RNN Three Pain Points
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "为什么需要Transformer：RNN/LSTM的三大痛点", 20, C_TITLE, True)

# Left: RNN sequential flow
words = ["The", "cat", "sat", "on", "the", "mat"]
start_x = Inches(0.5)
y_word = Inches(1.4)
y_rnn = Inches(2.0)
box_w = Inches(0.9)
box_h = Inches(0.4)
rnn_h = Inches(0.5)

for i, w in enumerate(words):
    x = start_x + Inches(i * 1.15)
    # Word box
    label_in_rect(slide, x, y_word, box_w, Inches(0.4), w, 11, C_TEXT,
                  fill=C_LIGHT_BLUE, border=C_ENCODER)
    # RNN box
    label_in_rect(slide, x, y_rnn, box_w, rnn_h, "RNN", 9, C_WHITE,
                  fill=C_ENCODER, bold=True)
    # Arrow between RNN boxes
    if i < len(words) - 1:
        add_arrow(slide, x + box_w + Pt(2), y_rnn + Inches(0.15), Inches(0.2), Inches(0.2), C_GRAY)

# "Sequential" annotation
add_text_box(slide, Inches(0.5), Inches(2.7), Inches(7), Inches(0.4),
             ">>> 顺序处理，一次一个词 >>>", 10, C_RED, True, PP_ALIGN.CENTER)

# Right: Three pain points
pain_points = [
    [("① 顺序处理", C_RED, True),
     ("RNN必须一个词一个词处理，GPU并行能力完全浪费。", C_TEXT),
     ("就像雇1000个工人搬砖，却规定每次只能1个人搬。", C_GRAY)],
    [("② 长距离遗忘", C_RED, True),
     ("处理后面词时会逐渐忘记前面信息，叫梯度消失。", C_TEXT),
     ("就像听10分钟故事，结尾记不清开头。", C_GRAY)],
    [("③ 信息瓶颈", C_RED, True),
     ("所有信息压缩成一个固定长度向量。", C_TEXT),
     ("就像读500页书，只能用100字摘要做判断。", C_GRAY)],
]

for i, lines in enumerate(pain_points):
    y = Inches(3.3) + Inches(i * 1.35)
    add_multiline_text(slide, Inches(0.6), y, Inches(12), Inches(1.2),
                       lines, 11, line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════
# Slide 4: Core Breakthrough
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "核心突破：让每个词都能和所有其他词直接交流", 20, C_TITLE, True)

# Left: Attention visualization
attn_words = ["The", "cat", "sat", "on", "the", "mat", "and", "it", "was", "happy"]
it_idx = 7
start_x = Inches(0.3)
y_words = Inches(1.6)
word_w = Inches(0.9)

for i, w in enumerate(attn_words):
    x = start_x + Inches(i * 1.15)
    fill = C_ENCODER if i == it_idx else C_LIGHT_BLUE
    border = C_ENCODER if i == it_idx else C_GRAY
    label_in_rect(slide, x, y_words, word_w, Inches(0.4), w, 10,
                  C_WHITE if i == it_idx else C_TEXT, fill=fill, border=border,
                  bold=(i == it_idx))

# Draw attention lines from "it" to other words
it_cx = start_x + Inches(it_idx * 1.15) + word_w / 2
it_y = y_words + Inches(0.4)

# Strongest line: it -> cat
cat_cx = start_x + Inches(1 * 1.15) + word_w / 2
add_line(slide, it_cx, it_y, cat_cx, y_words, C_RED, Pt(4))
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(2), Inches(0.3),
             "关联度最高", 9, C_RED, True, PP_ALIGN.CENTER)

# Medium lines
for target_idx in [0, 2, 3, 5, 8]:
    target_cx = start_x + Inches(target_idx * 1.15) + word_w / 2
    add_line(slide, it_cx, it_y, target_cx, y_words, C_ENCODER, Pt(2))

# Weak lines
for target_idx in [4, 6, 9]:
    target_cx = start_x + Inches(target_idx * 1.15) + word_w / 2
    add_line(slide, it_cx, it_y, target_cx, y_words, C_GRAY, Pt(1))

# Right: Text explanation
explains = [
    [("① 圆桌会议 vs 排队发言", C_TITLE, True),
     ("RNN像排队发言，Transformer像圆桌会议", C_TEXT),
     ("-- 每个人可以同时听到所有人", C_TEXT)],
    [("② 例句演示", C_TITLE, True),
('“it”如何找到“cat”？注意力权重最高', C_ENCODER, True),     ('通过Q-K匹配，“it”自动发现“cat”是最相关的名词', C_TEXT)],
    [("③ 三个问题，一个方案", C_TITLE, True),
     ("并行化 --> 解决顺序处理", C_FFN),
     ("直接交流 --> 解决长距离遗忘", C_ENCODER),
     ("不压缩 --> 解决信息瓶颈", C_EMBED)],
]

for i, lines in enumerate(explains):
    y = Inches(3.0) + Inches(i * 1.5)
    add_multiline_text(slide, Inches(0.6), y, Inches(12), Inches(1.3),
                       lines, 11, line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════
# Slide 5: Architecture Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
             "架构总览：左边理解，右边生成", 20, C_TITLE, True)

# ── Left: Full Architecture Diagram ──
arch_x = Inches(0.3)
arch_y = Inches(1.0)

# Input Embedding + Positional Encoding
label_in_rect(slide, arch_x, arch_y, Inches(2.2), Inches(0.45),
              "Input Embedding + PosEnc", 9, C_TEXT, fill=C_LIGHT_PURPLE, border=C_EMBED)

# Encoder x6 block
enc_x = arch_x
enc_y = arch_y + Inches(0.6)
enc_w = Inches(2.8)
enc_h = Inches(3.5)
enc_box = add_rect(slide, enc_x, enc_y, enc_w, enc_h, fill_color=C_WHITE, border_color=C_ENCODER, border_width=Pt(2))
# Label
add_text_box(slide, enc_x + Pt(2), enc_y + Pt(2), Inches(2), Inches(0.3),
             "Encoder x6", 10, C_ENCODER, True)

# Encoder internal detail (one layer shown)
ey = enc_y + Inches(0.4)
label_in_rect(slide, enc_x + Inches(0.1), ey, Inches(2.6), Inches(0.4),
              "Multi-Head Self-Attention", 9, C_WHITE, fill=C_ENCODER, bold=True)
ey += Inches(0.5)
label_in_rect(slide, enc_x + Inches(0.1), ey, Inches(2.6), Inches(0.3),
              "Add & Layer Norm", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)
ey += Inches(0.4)
label_in_rect(slide, enc_x + Inches(0.1), ey, Inches(2.6), Inches(0.4),
              "Feed Forward (FFN)", 9, C_WHITE, fill=C_ATTN, bold=True)
ey += Inches(0.5)
label_in_rect(slide, enc_x + Inches(0.1), ey, Inches(2.6), Inches(0.3),
              "Add & Layer Norm", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)
ey += Inches(0.5)
add_text_box(slide, enc_x + Inches(0.2), ey, Inches(2.4), Inches(0.3),
             "... x6 layers ...", 9, C_GRAY, False, PP_ALIGN.CENTER)

# Arrow from input to encoder
add_down_arrow(slide, enc_x + Inches(1.1), arch_y + Inches(0.45), Inches(0.15), Inches(0.15), C_GRAY)

# Decoder x6 block
dec_x = enc_x + enc_w + Inches(0.8)
dec_y = enc_y
dec_w = Inches(3.5)
dec_h = Inches(3.5)
dec_box = add_rect(slide, dec_x, dec_y, dec_w, dec_h, fill_color=C_WHITE, border_color=C_DECODER, border_width=Pt(2))
add_text_box(slide, dec_x + Pt(2), dec_y + Pt(2), Inches(2), Inches(0.3),
             "Decoder x6", 10, C_DECODER, True)

# Decoder internal detail
dy = dec_y + Inches(0.4)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.35),
              "Masked Self-Attention", 9, C_WHITE, fill=C_ENCODER, bold=True)
dy += Inches(0.42)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.25),
              "Add & Norm", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)
dy += Inches(0.32)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.35),
              "Cross-Attention", 9, C_WHITE, fill=C_EMBED, bold=True)
dy += Inches(0.42)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.25),
              "Add & Norm", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)
dy += Inches(0.32)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.35),
              "Feed Forward (FFN)", 9, C_WHITE, fill=C_ATTN, bold=True)
dy += Inches(0.42)
label_in_rect(slide, dec_x + Inches(0.1), dy, Inches(3.3), Inches(0.25),
              "Add & Norm", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)

# Decoder Input
label_in_rect(slide, dec_x, arch_y, Inches(2.2), Inches(0.45),
              "Output Embedding + PosEnc", 9, C_TEXT, fill=C_LIGHT_RED, border=C_DECODER)
add_down_arrow(slide, dec_x + Inches(1.1), arch_y + Inches(0.45), Inches(0.15), Inches(0.15), C_GRAY)

# K, V arrow from Encoder to Decoder Cross-Attention
kv_y = dec_y + Inches(0.4 + 0.42 + 0.32)  # Cross-Attention y
add_line(slide, enc_x + enc_w, enc_y + Inches(0.5), dec_x, kv_y + Inches(0.17), C_EMBED, Pt(2))
add_text_box(slide, enc_x + enc_w + Pt(2), enc_y + Inches(0.3), Inches(0.6), Inches(0.3),
             "K, V", 9, C_EMBED, True)

# Output: Linear -> Softmax
out_y = dec_y + dec_h + Inches(0.2)
label_in_rect(slide, dec_x + Inches(0.5), out_y, Inches(1.2), Inches(0.35),
              "Linear", 9, C_TEXT, fill=C_LIGHT_ORANGE, border=C_ATTN)
label_in_rect(slide, dec_x + Inches(1.9), out_y, Inches(1.2), Inches(0.35),
              "Softmax", 9, C_WHITE, fill=C_DECODER, bold=True)
add_down_arrow(slide, dec_x + Inches(1.7), dec_y + dec_h, Inches(0.15), Inches(0.2), C_GRAY)
add_arrow(slide, dec_x + Inches(1.7), out_y + Inches(0.1), Inches(0.15), Inches(0.15), C_GRAY)

# Right: Text
texts_right = [
    [("① 编码器 = \"编辑部\"", C_ENCODER, True),
     ("理解输入，6层，512维，6500万参数", C_TEXT)],
    [("② 解码器 = \"翻译部\"", C_DECODER, True),
     ("生成输出，6层，比编码器多交叉注意力", C_TEXT)],
    [("③ 数据流", C_TITLE, True),
     ("分词 -> 嵌入 -> 位置编码", C_TEXT),
     ("-> 编码器 -> 解码器 -> Softmax -> 概率", C_TEXT)],
]

rx = Inches(8.2)
for i, lines in enumerate(texts_right):
    y = Inches(1.0) + Inches(i * 1.8)
    add_multiline_text(slide, rx, y, Inches(4.8), Inches(1.6),
                       lines, 11, line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════
# Slide 6: Encoder Detail
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "编码器详解：每层都在交流与提炼", 20, C_TITLE, True)

# Left: Single encoder layer flow
flow_x = Inches(0.5)
flow_y = Inches(1.3)

# Input
label_in_rect(slide, flow_x, flow_y, Inches(1.5), Inches(0.5),
              "输入", 11, C_TEXT, fill=C_LIGHT_BLUE, border=C_ENCODER, bold=True)

# Main flow boxes
components = [
    ("Multi-Head\nSelf-Attention", C_ENCODER, C_WHITE),
    ("Add &\nLayerNorm", C_FFN, C_TEXT),
    ("Feed Forward\n(FFN)", C_ATTN, C_WHITE),
    ("Add &\nLayerNorm", C_FFN, C_TEXT),
]

cx = flow_x + Inches(2.0)
for i, (txt, fill, fc) in enumerate(components):
    bx = cx + Inches(i * 2.0)
    label_in_rect(slide, bx, flow_y, Inches(1.6), Inches(0.7),
                  txt, 9, fc, fill=fill, bold=True)
    # Arrow
    if i < len(components) - 1:
        add_arrow(slide, bx + Inches(1.6) + Pt(2), flow_y + Inches(0.25),
                  Inches(0.35), Inches(0.2), C_GRAY)

# Input arrow
add_arrow(slide, flow_x + Inches(1.5) + Pt(2), flow_y + Inches(0.2),
          Inches(0.45), Inches(0.2), C_GRAY)

# Output
out_x = cx + Inches(4 * 2.0)
label_in_rect(slide, out_x, flow_y, Inches(1.2), Inches(0.5),
              "输出", 11, C_TEXT, fill=C_LIGHT_BLUE, border=C_ENCODER, bold=True)
add_arrow(slide, cx + Inches(3 * 2.0) + Inches(1.6) + Pt(2), flow_y + Inches(0.2),
          Inches(0.35), Inches(0.2), C_GRAY)

# Residual connections (bypass lines)
res_y1 = flow_y - Inches(0.3)
res_y2 = flow_y + Inches(0.7) + Inches(0.15)
# Bypass attention: from before attn to after Add&Norm
add_line(slide, cx - Pt(4), flow_y + Inches(0.7), cx + Inches(2.0) + Inches(1.6) + Pt(4), flow_y + Inches(0.7), C_GRAY, Pt(1))
# Bypass FFN
add_line(slide, cx + Inches(4.0) - Pt(4), flow_y + Inches(0.7), cx + Inches(6.0) + Inches(1.6) + Pt(4), flow_y + Inches(0.7), C_GRAY, Pt(1))

# Residual labels
add_text_box(slide, cx + Inches(0.3), res_y2 + Inches(0.05), Inches(2), Inches(0.3),
             "残差连接（绕过注意力）", 9, C_GRAY, True)
add_text_box(slide, cx + Inches(4.3), res_y2 + Inches(0.05), Inches(2), Inches(0.3),
             "残差连接（绕过FFN）", 9, C_GRAY, True)

# Right: Explanation
texts = [
    [("① 多头自注意力 = \"圆桌会议\"", C_ENCODER, True),
     ("8个头从不同角度分析（语法/语义/指代）", C_TEXT)],
    [("② 残差连接 = \"高速公路匝道\"", C_FFN, True),
     ("信息可以抄近道，解决梯度消失", C_TEXT)],
    [("③ 层归一化 = \"标准化体检\"", C_ATTN, True),
     ("确保数值范围合理", C_TEXT)],
    [("④ 前馈网络 = \"独立思考\"", C_TITLE, True),
     ("512 -> 2048 -> ReLU -> 512", C_TEXT),
     ("计算量比注意力大2倍", C_TEXT)],
]

for i, lines in enumerate(texts):
    y = Inches(2.8) + Inches(i * 1.15)
    add_multiline_text(slide, Inches(0.6), y, Inches(12), Inches(1.0),
                       lines, 11, line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════
# Slide 7: Decoder Detail
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "解码器详解：比编码器多了一座翻译桥梁", 20, C_TITLE, True)

# Left: Single decoder layer
flow_x = Inches(0.5)
flow_y = Inches(1.3)

dec_components = [
    ("Masked\nSelf-Attention", C_ENCODER, C_WHITE),
    ("Add &\nNorm", C_FFN, C_TEXT),
    ("Cross-\nAttention", C_EMBED, C_WHITE),
    ("Add &\nNorm", C_FFN, C_TEXT),
    ("FFN", C_ATTN, C_WHITE),
    ("Add &\nNorm", C_FFN, C_TEXT),
]

cx = flow_x + Inches(0.5)
for i, (txt, fill, fc) in enumerate(dec_components):
    bx = cx + Inches(i * 1.6)
    label_in_rect(slide, bx, flow_y, Inches(1.3), Inches(0.6),
                  txt, 9, fc, fill=fill, bold=True)
    if i < len(dec_components) - 1:
        add_arrow(slide, bx + Inches(1.3) + Pt(2), flow_y + Inches(0.2),
                  Inches(0.25), Inches(0.15), C_GRAY)

# "Encoder Output (K,V)" arrow pointing to Cross-Attention
kv_arrow_y = flow_y + Inches(0.6) + Inches(0.2)
cross_attn_x = cx + Inches(2 * 1.6)
label_in_rect(slide, flow_x, kv_arrow_y, Inches(2), Inches(0.4),
              "编码器输出 (K, V)", 9, C_EMBED, fill=C_LIGHT_PURPLE, border=C_EMBED, bold=True)
add_line(slide, flow_x + Inches(2), kv_arrow_y + Inches(0.2),
         cross_attn_x + Inches(0.65), flow_y + Inches(0.6), C_EMBED, Pt(2))

# Right: Text
texts = [
    [("① 掩码自注意力 = \"考试不能偷看后面的题\"", C_ENCODER, True),
     ("因果掩码确保只能看到之前的词", C_TEXT)],
    [("② 交叉注意力 = \"带着问题去查资料\"", C_EMBED, True),
     ("Q 来自解码器，K 和 V 来自编码器", C_TEXT)],
    [("③ 前馈网络", C_ATTN, True),
     ("和编码器相同，512 -> 2048 -> ReLU -> 512", C_TEXT)],
]

for i, lines in enumerate(texts):
    y = Inches(2.6) + Inches(i * 1.3)
    add_multiline_text(slide, Inches(0.6), y, Inches(12), Inches(1.2),
                       lines, 11, line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════
# Slide 8: Training Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             "训练阶段：模型如何从数据中学会语言？", 20, C_TITLE, True)

# Top: Training data flow
# Left: Input -> Encoder
label_in_rect(slide, Inches(0.5), Inches(1.3), Inches(1.2), Inches(0.5),
              "Input\n(原文)", 10, C_TEXT, fill=C_LIGHT_BLUE, border=C_ENCODER, bold=True)
add_arrow(slide, Inches(1.75), Inches(1.45), Inches(0.4), Inches(0.2), C_GRAY)
label_in_rect(slide, Inches(2.2), Inches(1.3), Inches(1.8), Inches(0.7),
              "Encoder\nx6", 11, C_WHITE, fill=C_ENCODER, bold=True)
add_arrow(slide, Inches(4.05), Inches(1.55), Inches(0.4), Inches(0.2), C_GRAY)
label_in_rect(slide, Inches(4.5), Inches(1.3), Inches(2.0), Inches(0.5),
              "Encoder Output (K, V)", 9, C_EMBED, fill=C_LIGHT_PURPLE, border=C_EMBED, bold=True)

# Right: Decoder Input -> Decoder -> Predictions
label_in_rect(slide, Inches(7.5), Inches(1.3), Inches(1.8), Inches(0.5),
              "Decoder Input\n[START, 我, 爱]", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN, bold=True)
add_arrow(slide, Inches(9.35), Inches(1.45), Inches(0.4), Inches(0.2), C_GRAY)
label_in_rect(slide, Inches(9.8), Inches(1.3), Inches(1.5), Inches(0.7),
              "Decoder\nx6", 11, C_WHITE, fill=C_FFN, bold=True)
add_arrow(slide, Inches(11.35), Inches(1.55), Inches(0.4), Inches(0.2), C_GRAY)
label_in_rect(slide, Inches(11.8), Inches(1.3), Inches(1.2), Inches(0.5),
              "Predictions", 9, C_TEXT, fill=C_LIGHT_ORANGE, border=C_ATTN, bold=True)

# K,V arrow from encoder to decoder
add_line(slide, Inches(6.5), Inches(1.55), Inches(9.8), Inches(1.55), C_EMBED, Pt(2))

# Loss calculation
label_in_rect(slide, Inches(11.8), Inches(2.0), Inches(1.2), Inches(0.4),
              "Targets", 9, C_WHITE, fill=C_DECODER, bold=True)
add_text_box(slide, Inches(11.2), Inches(1.85), Inches(1.5), Inches(0.3),
             "<=> 计算Loss", 9, C_RED, True, PP_ALIGN.CENTER)

# Bottom brace text
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.4),
             "==================== 所有位置同时并行计算（开卷考试）====================",
             10, C_FFN, True, PP_ALIGN.CENTER)

# Bottom text
texts = [
    [("① 训练像\"开卷考试\"", C_FFN, True),
     ("预测每个词时都能看到前面正确答案", C_TEXT)],
    [("② 训练目标：学会\"预测下一个词\"", C_ENCODER, True),
     ("编码器处理输入，解码器预测 [我, 爱, 你, END]", C_TEXT)],
    [("③ 分工", C_EMBED, True),
     ("编码器理解原文（只做一次，缓存），解码器学会翻译", C_TEXT)],
]

for i, lines in enumerate(texts):
    y = Inches(3.3) + Inches(i * 1.2)
    add_multiline_text(slide, Inches(0.6), y, Inches(12), Inches(1.0),
                       lines, 11, line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════
# Slide 9: Embedding + Positional Encoding
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
             "训练 - 编码器第一步：把文字变成数字画像", 20, C_TITLE, True)

# Left: Embedding + Positional Encoding flow
# Top: 3 word boxes
words9 = ["我", "爱", "北"]
wx = Inches(0.5)
wy = Inches(1.0)
for i, w in enumerate(words9):
    label_in_rect(slide, wx + Inches(i * 1.5), wy, Inches(1.0), Inches(0.5),
                  w, 14, C_TEXT, fill=C_LIGHT_BLUE, border=C_ENCODER, bold=True)

# Arrow down
add_down_arrow(slide, Inches(1.8), Inches(1.55), Inches(0.15), Inches(0.15), C_GRAY)

# Embedding Matrix
label_in_rect(slide, Inches(0.5), Inches(1.9), Inches(3.5), Inches(0.5),
              "Embedding Matrix (37000 x 512)", 9, C_WHITE, fill=C_ENCODER, bold=True)

# Arrow down
add_down_arrow(slide, Inches(1.8), Inches(2.45), Inches(0.15), Inches(0.15), C_GRAY)

# 3 embedding vectors
for i in range(3):
    label_in_rect(slide, Inches(0.5) + Inches(i * 1.2), Inches(2.8), Inches(1.0), Inches(0.4),
                  "512维", 9, C_TEXT, fill=C_LIGHT_BLUE, border=C_ENCODER)

# Position Encoding on the right side
label_in_rect(slide, Inches(4.5), Inches(1.0), Inches(2.5), Inches(0.5),
              "Position Encoding (sin/cos)", 9, C_WHITE, fill=C_FFN, bold=True)

# Arrow from PosEnc down to vectors
add_down_arrow(slide, Inches(5.5), Inches(1.55), Inches(0.15), Inches(0.15), C_FFN)
for i in range(3):
    label_in_rect(slide, Inches(4.5) + Inches(i * 0.8), Inches(2.8), Inches(0.7), Inches(0.4),
                  "pos", 9, C_TEXT, fill=C_LIGHT_GREEN, border=C_FFN)

# Big plus sign
label_in_rect(slide, Inches(2.5), Inches(3.5), Inches(0.6), Inches(0.5),
              "+", 20, C_ATTN, bold=True)

# Annotation
add_text_box(slide, Inches(1.5), Inches(4.1), Inches(4), Inches(0.3),
             "词嵌入 + 位置编码 = 最终输入", 9, C_TITLE, True, PP_ALIGN.CENTER)

# Arrow down to final vectors
add_down_arrow(slide, Inches(2.5), Inches(4.5), Inches(0.15), Inches(0.15), C_GRAY)

# Final vectors
for i in range(3):
    label_in_rect(slide, Inches(0.8) + Inches(i * 1.3), Inches(4.9), Inches(1.1), Inches(0.4),
                  "512维\n语义+位置", 9, C_TEXT, fill=C_LIGHT_PURPLE, border=C_EMBED)

# Arrow to encoder
add_down_arrow(slide, Inches(2.5), Inches(5.45), Inches(0.15), Inches(0.15), C_GRAY)
label_in_rect(slide, Inches(1.5), Inches(5.8), Inches(2.5), Inches(0.4),
              "送入编码器第1层", 10, C_ENCODER, fill=C_LIGHT_BLUE, border=C_ENCODER, bold=True)

# Right: Text
texts = [
    [("① 分词", C_ENCODER, True),
     ("把句子切成最小单元，BPE方法", C_TEXT)],
    [("② 嵌入矩阵", C_EMBED, True),
     ("37000 x 512，约1890万参数，占模型30%", C_TEXT),
     ("像\"超级字典\"", C_GRAY)],
    [("③ 位置编码", C_FFN, True),
     ("sin/cos函数，给每个词贴\"座位号\"", C_TEXT),
     ("就像身份证加住址", C_GRAY)],
]

for i, lines in enumerate(texts):
    y = Inches(1.0) + Inches(i * 1.6)
    add_multiline_text(slide, Inches(8.0), y, Inches(5), Inches(1.4),
                       lines, 10, line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════
# Slide 10: QKV Matrix Generation
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
             "训练 - 编码器第二步：为每个词准备搜索工具", 20, C_TITLE, True)

# Left: QKV generation diagram
# Input vector
label_in_rect(slide, Inches(0.5), Inches(1.5), Inches(2.0), Inches(0.6),
              "输入向量\n(512维)", 10, C_WHITE, fill=C_ENCODER, bold=True)

# Three branches
qkv_data = [
    ("W^Q (512x64)", C_ENCODER, "Query\n(64维)"),
    ("W^K (512x64)", C_ATTN, "Key\n(64维)"),
    ("W^V (512x64)", C_FFN, "Value\n(64维)"),
]

for i, (mat, color, out) in enumerate(qkv_data):
    y = Inches(1.2) + Inches(i * 1.2)
    add_line(slide, Inches(2.5), Inches(1.8), Inches(3.2), y + Inches(0.3), C_GRAY, Pt(1))
    label_in_rect(slide, Inches(3.2), y, Inches(2.0), Inches(0.5),
                  mat, 9, C_WHITE, fill=color, bold=True)
    add_arrow(slide, Inches(5.25), y + Inches(0.15), Inches(0.35), Inches(0.2), C_GRAY)
    label_in_rect(slide, Inches(5.7), y, Inches(1.5), Inches(0.5),
                  out, 9, C_TEXT, fill=C_WHITE, border=color)

# Library analogy
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(7), Inches(0.3),
             "--- 图书馆搜索类比 ---", 10, C_TITLE, True, PP_ALIGN.CENTER)
analogy_items = [
    ("搜索框", C_ENCODER, "Query"),
    ("书架标签", C_ATTN, "Key"),
    ("打开书本", C_FFN, "Value"),
]
for i, (cn, color, en) in enumerate(analogy_items):
    x = Inches(0.8) + Inches(i * 2.3)
    label_in_rect(slide, x, Inches(5.1), Inches(1.0), Inches(0.35), cn, 9, C_WHITE, fill=color, bold=True)
    add_text_box(slide, x, Inches(5.5), Inches(1.0), Inches(0.3), en, 9, color, True, PP_ALIGN.CENTER)

# x8 heads annotation
add_text_box(slide, Inches(4.5), Inches(5.9), Inches(4), Inches(0.3),
             "x 8个头（每个头独立权重矩阵）", 10, C_EMBED, True)

# Right: Text
texts = [
    [("① Query = \"我在找什么信息\"", C_ENCODER, True),
     ("如\"it\"的Q编码了\"需要找名词\"", C_TEXT)],
    [("② Key = \"我有什么特征\"", C_ATTN, True),
     ("如\"cat\"的K编码了\"名词、动物、单数\"", C_TEXT)],
    [("③ Value = \"我的实际内容\"", C_FFN, True),
     ("被提取传递的信息", C_TEXT)],
    [("④ 8个头 = 8套独立QKV", C_EMBED, True),
     ("各64维，加起来恰好512维", C_TEXT)],
]

for i, lines in enumerate(texts):
    y = Inches(1.0) + Inches(i * 1.35)
    add_multiline_text(slide, Inches(7.8), y, Inches(5.2), Inches(1.2),
                       lines, 10, line_spacing=1.15)


# ═══════════════════════════════════════════════════════════════
# Slide 11: Attention 5-Step Calculation
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
             "训练 - 编码器第三步：自注意力的5步计算", 20, C_TITLE, True)

# Left: 5-step vertical flow
steps = [
    ("Input -> W^Q/W^K/W^V -> Q, K, V", C_ENCODER, C_LIGHT_BLUE, "所有词同时计算"),
    ("Q x K^T -> 分数矩阵", C_ENCODER, C_LIGHT_BLUE, ""),
    ("分数 / sqrt(d_k)", C_ATTN, C_LIGHT_YELLOW, "sqrt(64)=8, 防止分数过热"),
    ("Softmax -> 概率矩阵", C_FFN, C_LIGHT_GREEN, "每行和=1"),
    ("概率 x V -> 输出", C_EMBED, C_LIGHT_PURPLE, "融合所有词的信息"),
]

step_x = Inches(0.5)
step_w = Inches(3.5)
step_h = Inches(0.65)
start_y = Inches(1.0)

for i, (txt, border_c, fill_c, note) in enumerate(steps):
    y = start_y + Inches(i * 1.2)
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, step_x, y + Inches(0.1), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = border_c
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = str(i + 1)
    circle.text_frame.paragraphs[0].font.size = Pt(12)
    circle.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Step box
    label_in_rect(slide, step_x + Inches(0.5), y, step_w, step_h,
                  txt, 9, C_TEXT, fill=fill_c, border=border_c)

    # Annotation
    if note:
        add_text_box(slide, step_x + step_w + Inches(0.2), y + Inches(0.15), Inches(2), Inches(0.3),
                     note, 9, C_GRAY, True)

    # Down arrow
    if i < len(steps) - 1:
        add_down_arrow(slide, step_x + Inches(2.0), y + step_h, Inches(0.12), Inches(0.12), C_GRAY)

# Right: Detailed text (9pt, dense)
right_texts = [
    [("步骤1：每个词生成Q, K, V，完全并行", C_ENCODER, True),
     ("", C_TEXT)],
    [("步骤2：Q x K^T 点积，得到 T x T 分数矩阵", C_ENCODER, True),
     ("像搜索关键词和书标签逐一比对", C_TEXT)],
    [("步骤3：除以 sqrt(64) = 8，防止Softmax饱和", C_ATTN, True),
     ("如果Q和K方差为1，点积方差为d_k", C_TEXT)],
    [("步骤4：Softmax转为概率", C_FFN, True),
     ("最高分放大但不会完全变0", C_TEXT),
     ("像决定在每本书上花多少时间", C_GRAY)],
    [("步骤5：概率 x V 加权求和", C_EMBED, True),
     ("\"it\"融合了\"cat\"的信息", C_TEXT),
     ("-- 知道了it指的是cat", C_ENCODER, True)],
]

for i, lines in enumerate(right_texts):
    y = Inches(0.9) + Inches(i * 1.2)
    add_multiline_text(slide, Inches(7.8), y, Inches(5.2), Inches(1.1),
                       lines, 9, line_spacing=1.05)


# ═══════════════════════════════════════════════════════════════
# Save & Validate
# ═══════════════════════════════════════════════════════════════
output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v3_zh_part1.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved: {output_path}")

# Self-check
print("\n=== Self-Check Report ===")
issues = []
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        try:
            left = int(shape.left) if shape.left is not None else 0
            top = int(shape.top) if shape.top is not None else 0
            width = int(shape.width) if shape.width is not None else 0
            height = int(shape.height) if shape.height is not None else 0
            right = left + width
            bottom = top + height
            if right > int(SAFE_R):
                issues.append(f"Slide {si+1}: shape '{shape.name}' right={right} > {SAFE_R}")
            if bottom > int(SAFE_B):
                issues.append(f"Slide {si+1}: shape '{shape.name}' bottom={bottom} > {SAFE_B}")
        except (ValueError, TypeError):
            pass  # Skip connectors with float coords
        # Check font sizes
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.font.size and p.font.size < Pt(9):
                    issues.append(f"Slide {si+1}: font size {p.font.size} < 9pt in '{p.text[:20]}'")

if issues:
    print(f"WARNING: {len(issues)} issues found:")
    for iss in issues[:20]:
        print(f"  - {iss}")
else:
    print("PASS: All shapes within safe boundaries, all fonts >= 9pt")

print(f"\nTotal slides: {len(prs.slides)}")
print("Done!")
