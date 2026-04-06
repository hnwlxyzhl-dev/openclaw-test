#!/usr/bin/env python3
"""
Stable Diffusion 技术解析 PPT 生成器
v3.0 — 19章节完整版，15页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_LT = RGBColor(0xEB,0xF5,0xFB)
C_W = RGBColor(0xFF,0xFF,0xFF)
C_D = RGBColor(0x33,0x33,0x33)
C_G = RGBColor(0x88,0x88,0x88)
C_GRN = RGBColor(0x27,0xAE,0x60)
C_ORG = RGBColor(0xE6,0x7E,0x22)
C_RED = RGBColor(0xE7,0x4C,0x3C)
C_PUR = RGBColor(0x8E,0x44,0xAD)
C_TEL = RGBColor(0x00,0x96,0x88)
FC = "微软雅黑"
FE = "Arial"
OUT = "/home/admin/.openclaw/workspace-weaver/output/Stable_Diffusion_技术解析.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c
def rect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh
def rrect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh
def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(fs); p.font.color.rgb=fc
    p.font.bold=b; p.font.name=fn; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    return bx
def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld
        p.font.name=fn; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
    return bx
def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def chev(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.CHEVRON,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Stable Diffusion 技术解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))
def tbl(s,l,t,w,rows,cols,data,cw=None,hc=C_PRI,hfc=C_W,fs=10,ac=C_LT):
    ts=s.shapes.add_table(rows,cols,l,t,w,Inches(0.32*rows)); table=ts.table
    if cw:
        for i,wi in enumerate(cw): table.columns[i].width=wi
    from pptx.oxml.ns import qn
    from lxml import etree
    for r in range(rows):
        for c in range(cols):
            cell=table.cell(r,c)
            cell.text=data[r][c] if r<len(data) and c<len(data[r]) else ""
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(fs); p.font.name=FC; p.alignment=PP_ALIGN.CENTER; p.space_after=Pt(1)
            if r==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=hc
                for p in cell.text_frame.paragraphs: p.font.color.rgb=hfc; p.font.bold=True
            elif r%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=ac
            else: cell.fill.solid(); cell.fill.fore_color.rgb=C_W
            tc=cell._tc; tcPr=tc.get_or_add_tcPr()
            for bn in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                b=tcPr.find(qn(bn))
                if b is None: b=etree.SubElement(tcPr,qn(bn))
                b.set('w','6350')
                sf=b.find(qn('a:solidFill'))
                if sf is None: sf=etree.SubElement(b,qn('a:solidFill'))
                sc=sf.find(qn('a:srgbClr'))
                if sc is None: sc=etree.SubElement(sf,qn('a:srgbClr'))
                sc.set('val','D0D0D0')
    return ts

# ========== 第1页: 封面 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,C_PRI)
rect(sl,Inches(0.8),Inches(2.8),Inches(3),Inches(0.06),fc=C_ACC)
tb(sl,Inches(0.8),Inches(3.0),Inches(11),Inches(1.2),"Stable Diffusion",fs=52,fc=C_W,b=True,fn=FE)
tb(sl,Inches(0.8),Inches(4.0),Inches(11),Inches(0.8),"技术解析 — 从原理到前沿",fs=30,fc=C_ACC)
rect(sl,Inches(0.8),Inches(5.0),Inches(1.5),Inches(0.04),fc=C_ACC)
tb(sl,Inches(0.8),Inches(5.2),Inches(8),Inches(0.4),"深入理解扩散模型、CLIP、VAE、U-Net 及最新进展",fs=14,fc=RGBColor(0xAA,0xCC,0xEE))
tb(sl,Inches(0.8),Inches(5.7),Inches(8),Inches(0.4),"2026 年 4 月  |  19 章节完整版",fs=13,fc=C_G)
rect(sl,Inches(10),Inches(0),Inches(3.333),SLIDE_H,fc=RGBColor(0x16,0x2D,0x4A))
tb(sl,Inches(10.5),Inches(5.5),Inches(2.5),Inches(1.5),"AI IMAGE\nGENERATION",fs=16,fc=RGBColor(0x44,0x66,0x88),b=True,fn=FE,a=PP_ALIGN.RIGHT)
print("P1: cover")

# ========== 第2页: 目录 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"目录","CONTENTS — 19 CHAPTERS")
toc=[
("01","Diffusion Model 基础原理","墨水扩散与逆向去噪"),
("02","Stable Diffusion vs 普通 Diffusion","48倍压缩的核心创新"),
("03","整体架构概览","CLIP + U-Net + VAE 三大组件"),
("04","CLIP 文本编码器","把文字翻译成AI能理解的数字"),
("05","VAE 编解码器","超级智能的ZIP压缩软件"),
("06","U-Net 去噪网络","核心画师，8.6亿参数"),
("07","三大组件协作流程","AI画图流水线"),
("08","训练过程","6亿图像-文本对的炼丹之旅"),
("09","CFG 引导技术","无分类器引导，让AI更听话"),
("10","推理过程","从噪声到图像的雕刻过程"),
("11","采样器对比","9种算法，速度与质量取舍"),
("12","参数优化","5个核心参数的最佳实践"),
("13","Prompt 技巧","给AI画师下订单的艺术"),
("14","LoRA 原理","给画家添加专项风格手册"),
("15","ControlNet 原理","精确控制构图和姿态"),
("16","应用场景","5大领域的革命性影响"),
("17","SD3 / SDXL 最新进展","MMDiT架构与Flow Matching"),
("18","Flux 等新模型对比","2026年开源图像生成格局"),
("19","2026年AI图像生成动态","7大趋势与未来展望"),
]
sy=Inches(1.3)
for i,(n,t,d) in enumerate(toc):
    y=sy+Inches(i*0.30)
    tb(sl,Inches(0.6),y,Inches(0.5),Inches(28),n,fs=11,fc=C_ACC,b=True,fn=FE)
    tb(sl,Inches(1.2),y,Inches(4.2),Inches(28),t,fs=10,fc=C_D,b=True)
    tb(sl,Inches(5.6),y,Inches(7),Inches(28),d,fs=9,fc=C_G)
print("P2: TOC")

# ========== 第3页: Diffusion基础 + SD vs 普通 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"01-02  Diffusion Model 基础 + SD vs 普通 Diffusion","核心创新：在潜在空间中工作，计算量降低48倍")
# 左上: 前向扩散
rrect(sl,Inches(0.5),Inches(1.3),Inches(6),Inches(1.7),fc=RGBColor(0xFD,0xF2,0xE9))
tb(sl,Inches(0.7),Inches(1.35),Inches(5.5),Inches(0.3),"前向扩散：清晰照片 → 纯噪声（1000步）",fs=12,fc=C_ORG,b=True)
mtb(sl,Inches(0.7),Inches(1.7),Inches(5.6),Inches(1.2),[
    ("就像把一滴蓝色墨水滴入清水中，墨水逐渐扩散，最终整杯水变成均匀的浅蓝色。",False,C_D,10),
    ("系统在1000步中，每步往图像上添加随机噪声，最终变成纯噪声的雪花屏。",False,C_D,10),
    ("图像从 512x512x3 = 786,432 个有意义像素退化成完全随机数值。这个过程不需要学习。",False,C_D,10),
])
# 右上: 反向去噪
rrect(sl,Inches(6.8),Inches(1.3),Inches(6),Inches(1.7),fc=RGBColor(0xE8,0xF8,0xF5))
tb(sl,Inches(7.0),Inches(1.35),Inches(5.5),Inches(0.3),"反向去噪：纯噪声 → 清晰照片（AI学到的能力）",fs=12,fc=C_GRN,b=True)
mtb(sl,Inches(7.0),Inches(1.7),Inches(5.6),Inches(1.2),[
    ("AI修复师面对雪花屏，一步一步去掉噪声，每步恢复一点细节，最终还原清晰照片。",False,C_D,10),
    ("这就是扩散模型通过海量训练数据学到的核心技能。",False,C_D,10),
    ("比GAN更好：GAN像一下子跳到终点（不稳定），扩散模型是一步一步走（稳定高质量）。",False,C_GRN,10),
])
# 对比表
tbl(sl,Inches(0.5),Inches(3.2),Inches(12.3),8,4,[
    ["对比维度","普通 Diffusion（DDPM）","Stable Diffusion","核心区别"],
    ["操作空间","像素空间（512x512x3）","潜在空间（64x64x4）","SD不在像素空间工作"],
    ["数据维度","786,432 维","16,384 维","维度降低48倍"],
    ["计算量","极高","降低约 48 倍","48倍压缩是核心创新"],
    ["GPU显存","40GB+（高端GPU）","8-12GB（消费级GPU）","普通人也能运行"],
    ["生成速度","几十秒到几分钟","几秒到十几秒","实用性大幅提升"],
    ["代表模型","DDPM, DDIM","SD 1.5, SDXL, SD3","SD成为最成功开源模型"],
    ["硬件门槛","高（普通人用不起）","低（RTX 3060即可）","让AI画画人人可用"],
],cw=[Inches(2),Inches(3.8),Inches(3.8),Inches(2.7)],fs=10)
# 底部比喻
rrect(sl,Inches(0.5),Inches(5.9),Inches(12.3),Inches(1.0),fc=C_LT)
mtb(sl,Inches(0.7),Inches(5.95),Inches(11.8),Inches(0.9),[
    ("核心比喻：普通Diffusion像在4K画布上逐像素清理噪声（80万个数值）；SD先用VAE压缩成缩略图（1.6万个数值），在缩略图上去噪声后再还原。",False,C_D,10),
    ("就像画家不是在画布上逐像素画画，而是在缩略图上画画，画完再用放大镜还原成高清大图。这正是SD成为2022年以来最成功开源AI图像生成模型的根本原因。",True,C_PRI,10),
])
print("P3: Diffusion+SDvsDiff")

# ========== 第4页: 整体架构 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"03  整体架构概览","CLIP + U-Net + VAE 三大组件协同工作")
# 三个组件卡片
dims=[("1 CLIP 文本编码器","翻译官 — 123M参数\n输入: 文字(最多77token)\n输出: 77x768维语义向量\n训练: 4亿张图片+文字对\n作用: 把人类语言翻译成AI数字语言",C_ACC),
("2 U-Net 去噪网络","核心画师 — 860M参数(80%)\n输入: 噪声潜空间+文本向量+时间步\n输出: 预测噪声(64x64x4)\n核心: Cross-Attention+Skip Connections\n作用: 一步步去掉噪声",C_RED),
("3 VAE 编解码器","压缩工具 — 83M参数\n编码: 512x512x3 → 64x64x4(48倍压缩)\n解码: 64x64x4 → 512x512x3(还原高清)\n训练时压缩，推理时还原\n比喻: 超级智能的ZIP压缩软件",C_GRN)]
for i,(t,d,c) in enumerate(dims):
    x=Inches(0.5+i*4.2)
    rrect(sl,x,Inches(1.3),Inches(3.9),Inches(2.0),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.9),Inches(0.32),fc=c)
    tb(sl,x+Inches(0.1),Inches(1.33),Inches(3.7),Inches(0.28),t,fs=11,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.1),Inches(1.7),Inches(3.7),Inches(1.5),d,fs=9.5,fc=C_D,ls=1.2)
    if i<2: arrow_r(sl,x+Inches(3.95),Inches(2.0),Inches(0.22),Inches(0.2),c)
# 推理流程
rrect(sl,Inches(0.5),Inches(3.5),Inches(12.3),Inches(1.8),fc=C_LT)
tb(sl,Inches(0.7),Inches(3.55),Inches(5),Inches(0.28),"推理时完整数据流",fs=13,fc=C_PRI,b=True)
steps=[("Step 1","文字描述 → CLIP → 77x768文本向量",C_ACC),("Step 2","随机噪声(64x64x4) → U-Net执行20-50步去噪",C_RED),
        ("Step 3","去噪后的潜空间(64x64x4) → VAE解码器",C_GRN),("Step 4","输出: 512x512x3 高清图像",C_D)]
for i,(st,desc,c) in enumerate(steps):
    y=Inches(3.9)+Inches(i*0.35)
    chev(sl,Inches(0.8),y,Inches(0.9),Inches(0.28),c)
    tb(sl,Inches(0.85),y+Inches(0.02),Inches(0.8),Inches(0.24),st,fs=9,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,Inches(1.9),y+Inches(0.02),Inches(5),Inches(0.24),desc,fs=10,fc=C_D)
# 参数表
tbl(sl,Inches(0.5),Inches(5.5),Inches(12.3),5,4,[
    ["组件","参数量","占比","显存占用"],
    ["CLIP 文本编码器","约 123M","~12%","~0.5 GB"],
    ["VAE 编解码器","约 83M","~8%","~0.3 GB"],
    ["U-Net 去噪网络","约 860M","~80%","~4-5 GB"],
    ["总计（SD 1.5）","约 10 亿（1B）","100%","6-8 GB"],
],cw=[Inches(3.5),Inches(3),Inches(2.5),Inches(3.3)],fs=10)
print("P4: architecture")

# ========== 第5页: CLIP ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"04  CLIP 文本编码器详解","Contrastive Language-Image Pre-training")
cd=[("输入","文本描述字符串\n经分词器最多77个token\n约50-60个汉字\n超出部分截断",C_ACC),
    ("输出","77x768维张量\n共59,136个浮点数\n文本的语义指纹\n包含所有含义信息",C_GRN),
    ("参数","约1.23亿(123M)\n12层Transformer\n隐藏维度768\n12个注意力头",C_ORG),
    ("作用","人类文字→数字向量\n通过Cross-Attention注入U-Net\n指导去噪方向\n没有它AI不知道画什么",C_PUR)]
for i,(t,d,c) in enumerate(cd):
    x=Inches(0.5+i*3.2)
    rrect(sl,x,Inches(1.3),Inches(3.0),Inches(1.5),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.0),Inches(0.3),fc=c)
    tb(sl,x+Inches(0.08),Inches(1.32),Inches(2.8),Inches(0.26),t,fs=12,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.12),Inches(1.65),Inches(2.7),Inches(1.1),d,fs=10,fc=C_D,ls=1.2)
# 工作流程
rrect(sl,Inches(0.5),Inches(3.0),Inches(12.3),Inches(3.9),fc=C_LT)
tb(sl,Inches(0.7),Inches(3.05),Inches(5),Inches(0.28),"CLIP 内部工作流程（4步）",fs=13,fc=C_PRI,b=True)
ws=[("1 分词","文本拆成最小语义单元(token)\n一只橘猫坐在沙发上 → 一只/橘/猫/坐在/沙发/上\n每个token映射成768维向量",C_ACC),
    ("2 位置编码","每个词加上位置向量\n猫坐在沙发上 ≠ 沙发坐在猫上\n顺序不同含义天差地别\n位置编码让模型理解词语顺序",C_GRN),
    ("3 Transformer","12层Transformer网络处理\n每层含自注意力机制\n每个词都能看到其他所有词\n橘猫会注意到坐在和沙发",C_ORG),
    ("4 输出向量","每个token变成768维向量\n共77x768=59,136个数值\n不仅含独立含义还含词间关系\n对整段文本的完整语义编码",C_PUR)]
for i,(t,d,c) in enumerate(ws):
    x=Inches(0.7+i*3.1)
    chev(sl,x,Inches(3.45),Inches(2.9),Inches(0.32),c)
    tb(sl,x+Inches(0.15),Inches(3.47),Inches(2.6),Inches(0.28),t,fs=11,fc=C_W,b=True)
    tb(sl,x+Inches(0.08),Inches(3.85),Inches(2.7),Inches(2.0),d,fs=9,fc=C_D,ls=1.15)
    if i<3: arrow_r(sl,x+Inches(2.95),Inches(3.5),Inches(0.15),Inches(0.18),c)
print("P5: CLIP")

# ========== 第6页: VAE ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"05  VAE 编解码器详解","Variational Autoencoder — 压缩48倍但肉眼几乎看不出差别")
vd=[("编码器输入","原始图像 512x512x3\n786,432个像素值\n值域 0-255",C_ACC),("编码器输出","潜空间 64x64x4\n16,384个浮点数\n值域约[-5,5]",C_GRN),
    ("参数量","约8300万(83M)\n编码器~34M 解码器~49M\n解码器更大(还原比压缩更复杂)",C_ORG),("核心作用","训练时压缩(降低计算成本)\n推理时还原(潜空间→高清)\n保留人眼重要信息，丢弃冗余",C_PUR)]
for i,(t,d,c) in enumerate(vd):
    x=Inches(0.5+i*3.2)
    rrect(sl,x,Inches(1.3),Inches(3.0),Inches(1.3),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.0),Inches(0.28),fc=c)
    tb(sl,x+Inches(0.08),Inches(1.32),Inches(2.8),Inches(0.24),t,fs=11,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.1),Inches(1.62),Inches(2.7),Inches(0.9),d,fs=9.5,fc=C_D,ls=1.15)
# 编码器
rrect(sl,Inches(0.5),Inches(2.8),Inches(6),Inches(3.0),fc=RGBColor(0xE8,0xF0,0xFE))
tb(sl,Inches(0.7),Inches(2.85),Inches(5.5),Inches(0.25),"编码器: 4层下采样(每层宽高缩小一半)",fs=12,fc=C_ACC,b=True)
es=[("512x512x3","输入图像(RGB)",C_ACC),("256x256x128","第1层卷积+下采样",C_ACC),("128x128x256","第2层卷积+下采样",C_ACC),
    ("64x64x512","第3层卷积+下采样",C_ACC),("32x32x512","第4层卷积+下采样",C_ACC),("64x64x4","输出潜空间(48倍压缩!)",C_GRN)]
for i,(s,d,c) in enumerate(es):
    y=Inches(3.2)+Inches(i*0.38)
    rrect(sl,Inches(0.8),y,Inches(1.5),Inches(0.3),fc=c)
    tb(sl,Inches(0.85),y+Inches(0.02),Inches(1.4),Inches(0.26),s,fs=9,fc=C_W,b=True,a=PP_ALIGN.CENTER,fn=FE)
    tb(sl,Inches(2.5),y+Inches(0.02),Inches(3.5),Inches(0.26),d,fs=9.5,fc=C_D)
    if i<5: arrow_d(sl,Inches(1.35),y+Inches(0.3),Inches(0.2),Inches(0.08),c)
# 解码器
rrect(sl,Inches(6.8),Inches(2.8),Inches(6),Inches(3.0),fc=RGBColor(0xE8,0xF8,0xF5))
tb(sl,Inches(7.0),Inches(2.85),Inches(5.5),Inches(0.25),"解码器: 4层上采样(每层宽高翻倍)",fs=12,fc=C_GRN,b=True)
ds=[("64x64x4","输入潜空间",C_GRN),("128x128x512","第1层卷积+上采样",C_GRN),("256x256x256","第2层卷积+上采样",C_GRN),
    ("512x512x128","第3层卷积+上采样",C_GRN),("512x512x64","第4层卷积+上采样",C_GRN),("512x512x3","输出高清图像(RGB)",C_ACC)]
for i,(s,d,c) in enumerate(ds):
    y=Inches(3.2)+Inches(i*0.38)
    rrect(sl,Inches(7.1),y,Inches(1.5),Inches(0.3),fc=c)
    tb(sl,Inches(7.15),y+Inches(0.02),Inches(1.4),Inches(0.26),s,fs=9,fc=C_W,b=True,a=PP_ALIGN.CENTER,fn=FE)
    tb(sl,Inches(8.8),y+Inches(0.02),Inches(3.5),Inches(0.26),d,fs=9.5,fc=C_D)
    if i<5: arrow_d(sl,Inches(7.65),y+Inches(0.3),Inches(0.2),Inches(0.08),c)
# 底部
rrect(sl,Inches(0.5),Inches(6.0),Inches(12.3),Inches(0.9),fc=RGBColor(0xFE,0xF9,0xE7))
mtb(sl,Inches(0.7),Inches(6.05),Inches(11.8),Inches(0.8),[
    ("为什么不用普通自编码器？普通AE的潜空间不规则，两个相近点可能对应完全不相关的图像，扩散模型无法工作。",False,C_D,10),
    ("VAE通过正则化约束确保潜空间连续平滑，使扩散操作（加噪声、去噪声）成为可能。这就是为什么选VAE而不是普通AE。",True,C_ORG,10),
])
print("P6: VAE")

# ========== 第7页: U-Net ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"06  U-Net 噪声预测网络详解","核心画师 — 8.6亿参数，占据80%以上计算量")
ud=[("输入","带噪声潜空间(64x64x4)\n文本向量(77x768)\n时间步t(1-1000)\n可选条件图像",C_ACC),
    ("输出","预测噪声(64x64x4)\n形状与输入相同\n用于从带噪声图像中减去噪声",C_GRN),
    ("参数","约8.6亿(860M)\n三个组件中最大\n推理占用4-5GB显存",C_ORG),
    ("作用","核心去噪引擎\n每步预测应去除的噪声\n从噪声逐步雕刻出清晰图像\n比喻: 经验丰富的雕塑家",C_PUR)]
for i,(t,d,c) in enumerate(ud):
    x=Inches(0.5+i*3.2)
    rrect(sl,x,Inches(1.3),Inches(3.0),Inches(1.4),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.0),Inches(0.28),fc=c)
    tb(sl,x+Inches(0.08),Inches(1.32),Inches(2.8),Inches(0.24),t,fs=11,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.1),Inches(1.62),Inches(2.7),Inches(1.0),d,fs=9.5,fc=C_D,ls=1.15)
# 结构+机制
rrect(sl,Inches(0.5),Inches(2.9),Inches(12.3),Inches(4.0),fc=C_LT)
tb(sl,Inches(0.7),Inches(2.95),Inches(5),Inches(0.25),"U-Net 结构与三个关键机制",fs=13,fc=C_PRI,b=True)
# 编码器
tb(sl,Inches(0.7),Inches(3.3),Inches(3),Inches(0.22),"编码器(逐步收缩)",fs=10,fc=C_RED,b=True)
el=[("64x64x320 → 32x32x640","ResBlock+Attention"),("32x32x640 → 16x16x1280","ResBlock+Attention"),
    ("16x16x1280 → 8x8x1280","ResBlock+Attention"),("8x8x1280 → 4x4x1280","ResBlock(最深层)")]
for i,(s,d) in enumerate(el):
    y=Inches(3.55)+Inches(i*0.3)
    tb(sl,Inches(0.8),y,Inches(2.5),Inches(0.25),s,fs=9,fc=C_D,fn=FE)
    tb(sl,Inches(3.3),y,Inches(1.3),Inches(0.25),d,fs=8,fc=C_G)
# 瓶颈
tb(sl,Inches(5.2),Inches(3.3),Inches(2.5),Inches(0.22),"瓶颈层(最深层)",fs=10,fc=C_PUR,b=True)
tb(sl,Inches(5.2),Inches(3.55),Inches(2.5),Inches(0.7),"4x4x1280\n最抽象语义信息\n空间最小通道最大",fs=9.5,fc=C_D,ls=1.15)
# 解码器
tb(sl,Inches(8.2),Inches(3.3),Inches(4),Inches(0.22),"解码器(逐步扩张)",fs=10,fc=C_GRN,b=True)
dl=[("4x4 → 8x8x1280","+Skip"),("8x8 → 16x16x640","+Skip"),("16x16 → 32x32x320","+Skip"),("32x32 → 64x64x4","预测噪声")]
for i,(s,d) in enumerate(dl):
    y=Inches(3.55)+Inches(i*0.3)
    tb(sl,Inches(8.4),y,Inches(2.3),Inches(0.25),s,fs=9,fc=C_D,fn=FE)
    tb(sl,Inches(10.7),y,Inches(1.8),Inches(0.25),d,fs=8,fc=C_GRN,b=True)
# 三个机制
ms=[("Cross-Attention(交叉注意力)","文本条件注入核心\n图像问: 根据文字我该怎样？\n文本答: 你这里应该是橘猫\n出现在编码器1-3层+解码器1-3层(共6处)",C_ACC),
    ("Skip Connections(跳跃连接)","编码器细节直接传给解码器\n防止信息在压缩中丢失\n就像建筑师随时参考草稿本\n没有它细节会大打折扣",C_ORG),
    ("Time Embedding(时间嵌入)","告诉U-Net当前噪声水平(1-1000)\n高噪声时粗雕(确定轮廓)\n低噪声时精雕(添加细节)\n通过正弦编码+两层MLP注入",C_PUR)]
for i,(t,d,c) in enumerate(ms):
    x=Inches(0.7+i*4.1)
    rrect(sl,x,Inches(4.85),Inches(3.8),Inches(1.9),fc=C_W,lc=c)
    tb(sl,x+Inches(0.1),Inches(4.9),Inches(3.6),Inches(0.22),t,fs=9.5,fc=c,b=True)
    tb(sl,x+Inches(0.1),Inches(5.15),Inches(3.6),Inches(1.5),d,fs=9,fc=C_D,ls=1.1)
print("P7: U-Net")

# ========== 第8页: 三大组件协作 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"07  三大组件协作流程","AI画图流水线 — 各司其职，紧密配合")
sts=[("第1站","CLIP 翻译站","用户输入: 一只橘猫坐在木质沙发上\n通过12层Transformer处理\n输出77x768维语义指纹向量\n发送给下一站的U-Net","翻译官 — 把人类语言翻译成AI数字语言",C_ACC),
    ("第2站","U-Net 画图站","接收文本向量+64x64x4纯随机噪声\n第1轮: 只能猜出大致色块分布\n第15轮: 能看到猫轮廓、沙发形状\n第30轮: 毛发纹理、阳光光线\n第50轮: 去噪完成","核心画师 — 执行20-50轮去噪循环",C_RED),
    ("第3站","VAE 还原站","接收U-Net输出的干净潜空间\n通过4层上采样逐步还原\n64x64→128→256→512\n最终输出512x512x3高清图像","还原打印机 — 潜空间→高清图像",C_GRN)]
for i,(st,t,d,role,c) in enumerate(sts):
    x=Inches(0.5+i*4.2)
    rrect(sl,x,Inches(1.3),Inches(3.9),Inches(3.5),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.9),Inches(0.5),fc=c)
    tb(sl,x+Inches(0.1),Inches(1.32),Inches(3.7),Inches(0.2),st,fs=10,fc=RGBColor(0xDD,0xEE,0xFF))
    tb(sl,x+Inches(0.1),Inches(1.55),Inches(3.7),Inches(0.22),t,fs=13,fc=C_W,b=True)
    tb(sl,x+Inches(0.12),Inches(1.9),Inches(3.6),Inches(2.0),d,fs=10,fc=C_D,ls=1.2)
    rrect(sl,x+Inches(0.1),Inches(4.3),Inches(3.7),Inches(0.35),fc=C_LT)
    tb(sl,x+Inches(0.2),Inches(4.32),Inches(3.5),Inches(0.3),role,fs=10,fc=c,b=True,a=PP_ALIGN.CENTER)
    if i<2: arrow_r(sl,x+Inches(3.95),Inches(2.5),Inches(0.22),Inches(0.2),c)
# 底部
rrect(sl,Inches(0.5),Inches(5.1),Inches(12.3),Inches(1.8),fc=C_LT)
tb(sl,Inches(0.7),Inches(5.15),Inches(5),Inches(0.25),"为什么要这样分工？",fs=13,fc=C_PRI,b=True)
mtb(sl,Inches(0.7),Inches(5.45),Inches(11.8),Inches(1.4),[
    ("CLIP擅长理解语言但不擅长生成图像 — 所以只让它做翻译",False,C_D,10),
    ("U-Net擅长在低维空间做逐步变换但不擅长人类语言 — 需要CLIP的翻译结果作为指导",False,C_D,10),
    ("VAE擅长压缩和还原但不擅长语义 — 只在开头和结尾使用",False,C_D,10),
    ("就像专业厨房：有人负责理解客人的点单(CLIP)，有人负责烹饪(U-Net)，有人负责摆盘(VAE)",True,C_ACC,10),
    ("",False,C_D,4),
    ("数据连贯性：CLIP输出的文本向量贯穿整个U-Net去噪过程(通过Cross-Attention)，确保每一步都\"记得\"用户要画什么。",False,C_D,10),
])
print("P8: collaboration")

# ========== 第9页: 训练过程 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"08  训练过程","6亿图像-文本对的炼丹之旅 — 15万GPU小时")
ts2=[("1","采样图像和文本","从LAION数据集随机抽取\n图像x0和对应文本描述",C_ACC),
     ("2","VAE编码","512x512x3 → 64x64x4\n压缩到潜空间",C_GRN),
     ("3","添加随机噪声","随机选t∈[1,1000]\nz_t = sqrt(a_t)*z0 + sqrt(1-a_t)*epsilon",C_ORG),
     ("4","U-Net预测噪声","输入: z_t + t + 文本向量\n输出: 预测噪声epsilon_pred",C_RED),
     ("5","计算损失函数","L = ||epsilon - epsilon_pred||^2\n预测噪声vs真实噪声的均方误差",C_PUR),
     ("6","反向传播更新","更新U-Net的8.6亿参数\nCLIP和VAE保持冻结(不训练)",C_TEL)]
for i,(n,t,d,c) in enumerate(ts2):
    row,col=i//3,i%3
    x=Inches(0.5+col*4.2); y=Inches(1.3+row*1.6)
    rrect(sl,x,y,Inches(3.9),Inches(1.35),fc=C_W,lc=c)
    rrect(sl,x+Inches(0.1),y+Inches(0.08),Inches(0.35),Inches(0.35),fc=c)
    tb(sl,x+Inches(0.12),y+Inches(0.1),Inches(0.3),Inches(0.3),n,fs=15,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.55),y+Inches(0.1),Inches(3.2),Inches(0.3),t,fs=12,fc=c,b=True)
    tb(sl,x+Inches(0.12),y+Inches(0.5),Inches(3.6),Inches(0.8),d,fs=10,fc=C_D,ls=1.2)
    if col<2: arrow_r(sl,x+Inches(3.95),y+Inches(0.5),Inches(0.22),Inches(0.18),C_G)
# 成本表
tbl(sl,Inches(0.5),Inches(4.7),Inches(12.3),4,4,[
    ["项目","数据","项目","数据"],
    ["训练数据","LAION-5B子集，约6亿图文对","训练硬件","多张A100 80GB GPU"],
    ["训练时长","约15万GPU小时","训练成本","约50-60万美元"],
    ["模型大小","约4GB(float16)","冻结组件","CLIP+VAE(只训练U-Net)"],
],cw=[Inches(2),Inches(3.8),Inches(2),Inches(4.5)],fs=10)
# 比喻
rrect(sl,Inches(0.5),Inches(6.1),Inches(12.3),Inches(0.8),fc=C_LT)
mtb(sl,Inches(0.7),Inches(6.15),Inches(11.8),Inches(0.7),[
    ("比喻：训练就像培养从零学画的学徒。导师先给学徒看参考画，然后把画擦到不同程度的模糊(添加噪声)，让学徒猜被擦掉的部分(预测噪声)。",False,C_D,10),
    ("猜对了奖励(损失降低)，猜错了改进(梯度更新)。经过几十万次练习，学徒就能从空白画布上根据文字描述画出精美图像。",True,C_ACC,10),
])
print("P9: training")

# ========== 第10页: CFG ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"09  CFG（Classifier-Free Guidance）详解","无分类器引导 — 让AI更听话地按描述画画")
rrect(sl,Inches(0.5),Inches(1.3),Inches(12.3),Inches(1.3),fc=C_LT)
mtb(sl,Inches(0.7),Inches(1.35),Inches(11.8),Inches(1.2),[
    ("核心思想：既知道什么是对的，也知道什么不是对的",True,C_PRI,13),
    ("",False,C_D,4),
    ("公式: epsilon_guided = epsilon_uncond + s x (epsilon_cond - epsilon_uncond)",True,C_PUR,13),
    ("epsilon_uncond = 没有文本条件时的预测(自由发挥)    epsilon_cond = 有文本条件时的预测(按指令画)",False,C_D,10),
    ("s = 引导系数(CFG Scale)，通常7-12    s越大越忠实于文字    s过大则过度饱和出现伪影",False,C_D,10),
])
tbl(sl,Inches(0.5),Inches(2.8),Inches(6.5),6,3,[
    ["CFG Scale","效果","适用场景"],
    ["1-3","几乎无引导，图像自由发挥","创意探索"],
    ["5-7","适中引导，平衡忠实度和创造力","通用场景(推荐默认7)"],
    ["7-12","强引导，高度忠实文字描述","精确控制(产品图)"],
    ["15-20","过强引导，可能伪影","不推荐"],
    ["20+","严重过度饱和","绝对不推荐"],
],cw=[Inches(1.3),Inches(2.8),Inches(2.4)],fs=10)
# 图示
rrect(sl,Inches(7.5),Inches(2.8),Inches(5.3),Inches(3.0),fc=C_W,lc=C_ACC)
tb(sl,Inches(7.7),Inches(2.85),Inches(4.8),Inches(0.25),"CFG工作原理",fs=12,fc=C_ACC,b=True)
rrect(sl,Inches(7.8),Inches(3.3),Inches(2.2),Inches(0.6),fc=RGBColor(0xE0,0xE0,0xE0))
tb(sl,Inches(7.8),Inches(3.38),Inches(2.2),Inches(0.4),"无条件生成\n(自由发挥)",fs=9,fc=C_G,b=True,a=PP_ALIGN.CENTER)
tb(sl,Inches(10.1),Inches(3.3),Inches(0.4),Inches(0.3),"→",fs=18,fc=C_PUR,b=True,a=PP_ALIGN.CENTER,fn=FE)
tb(sl,Inches(10.1),Inches(3.55),Inches(0.5),Inches(0.25),"差异",fs=8,fc=C_PUR,a=PP_ALIGN.CENTER)
rrect(sl,Inches(10.6),Inches(3.3),Inches(2.0),Inches(0.6),fc=RGBColor(0xE8,0xF0,0xFE),lc=C_ACC)
tb(sl,Inches(10.6),Inches(3.38),Inches(2.0),Inches(0.4),"有条件生成\n(按指令画)",fs=9,fc=C_ACC,b=True,a=PP_ALIGN.CENTER)
arrow_d(sl,Inches(11.5),Inches(3.95),Inches(0.2),Inches(0.2),C_PUR)
tb(sl,Inches(10.8),Inches(3.95),Inches(1.5),Inches(0.2),"x CFG Scale (s)",fs=8,fc=C_PUR,b=True,a=PP_ALIGN.CENTER)
rrect(sl,Inches(9.3),Inches(4.5),Inches(2.8),Inches(0.55),fc=C_GRN)
tb(sl,Inches(9.3),Inches(4.53),Inches(2.8),Inches(0.45),"= 引导后的高质量结果",fs=11,fc=C_W,b=True,a=PP_ALIGN.CENTER)
tb(sl,Inches(7.8),Inches(4.5),Inches(1.5),Inches(0.3),"+ 无条件结果",fs=9,fc=C_G,a=PP_ALIGN.CENTER)
# 总结
rrect(sl,Inches(0.5),Inches(6.1),Inches(12.3),Inches(0.8),fc=RGBColor(0xFE,0xF9,0xE7))
mtb(sl,Inches(0.7),Inches(6.15),Inches(11.8),Inches(0.7),[
    ("比喻: CFG就像严厉又公正的艺术指导。每画一笔都问两个问题: 你按客户要求画了吗？不管要求你会怎么画？",False,C_D,10),
    ("然后把两个答案的差异放大(乘以CFG Scale)，强制画师更严格遵循客户要求。没有CFG，图像模糊且偏离描述。",True,C_ORG,10),
])
print("P10: CFG")

# ========== 第11页: 推理+采样器 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"10-11  推理过程 + 采样器对比","从噪声到图像的雕刻 — 4步推理，9种采样器选择")
# 推理4步
rrect(sl,Inches(0.5),Inches(1.3),Inches(12.3),Inches(1.8),fc=C_LT)
tb(sl,Inches(0.7),Inches(1.35),Inches(5),Inches(0.25),"推理的4个步骤",fs=13,fc=C_PRI,b=True)
ist=[("Step 1","文本编码","文字→CLIP→77x768向量\n只需一次，后续复用",C_ACC),
     ("Step 2","初始化噪声","采样64x64x4随机张量\n=纯噪声雪花屏\nseed决定随机数",C_GRN),
     ("Step 3","逐步去噪","循环20-50步\n每步U-Net预测噪声→减去\n从t=1000递减到t=1",C_ORG),
     ("Step 4","VAE解码","64x64x4→解码器→512x512x3\n4层上采样还原\n输出高清图像",C_PUR)]
for i,(st,t,d,c) in enumerate(ist):
    x=Inches(0.7+i*3.1)
    rrect(sl,x,Inches(1.7),Inches(2.9),Inches(1.25),fc=C_W,lc=c)
    rect(sl,x,Inches(1.7),Inches(2.9),Inches(0.25),fc=c)
    tb(sl,x+Inches(0.08),Inches(1.72),Inches(2.7),Inches(0.2),st,fs=9,fc=C_W,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.08),Inches(1.98),Inches(2.7),Inches(0.2),t,fs=11,fc=c,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.1),Inches(2.25),Inches(2.7),Inches(0.6),d,fs=9,fc=C_D,ls=1.15)
    if i<3: arrow_r(sl,x+Inches(2.95),Inches(2.1),Inches(0.15),Inches(0.15),c)
# 采样器表
tb(sl,Inches(0.5),Inches(3.3),Inches(5),Inches(0.25),"采样器对比",fs=12,fc=C_PRI,b=True)
tbl(sl,Inches(0.5),Inches(3.6),Inches(12.3),10,5,[
    ["采样器","步数","速度","质量","适用场景"],
    ["Euler","20-30","快","中","快速预览"],
    ["DDIM","20-50","中","高","确定性生成"],
    ["DPM++ 2M","20-30","快","高","通用推荐"],
    ["DPM++ 2M Karras","20-30","快","很高","高质量首选"],
    ["DPM++ SDE Karras","20-30","中","很高","艺术感强"],
    ["DPM++ SDE","10-15","很快","中高","快速出图"],
    ["UniPC","10-20","很快","高","最新高效"],
    ["LCM","4-8","极快","中","实时生成"],
],cw=[Inches(2.5),Inches(1),Inches(1),Inches(1),Inches(6.8)],fs=9)
# 推荐
rrect(sl,Inches(0.5),Inches(6.4),Inches(12.3),Inches(0.5),fc=RGBColor(0xE8,0xF8,0xF5))
mtb(sl,Inches(0.7),Inches(6.42),Inches(11.8),Inches(0.4),[
    ("通用推荐: DPM++ 2M Karras 20-30步 | 追求速度: DPM++ SDE 10-15步或LCM 4-8步 | 追求质量: DPM++ 2M Karras 40-50步",True,C_GRN,10),
])
print("P11: inference+samplers")

# ========== 第12页: 参数+Prompt ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"12-13  参数优化 + Prompt 技巧","实战指南 — 让AI画出你想要的图像")
# 参数卡片
ps=[("Steps(去噪步数)","推荐: 20-50步\nDPM++ 2M用25步\n超过30步提升递减\n比喻: 打磨家具20遍已很好",C_ACC),
    ("CFG Scale(引导系数)","推荐: 5-12\n通用7-8 精确10-12\n太低偏离描述太高过度饱和\n比喻: 老师批改严格程度",C_GRN),
    ("Seed(随机种子)","任意整数\n相同种子+参数=相同图像\n不同种子=完全不同图像\n比喻: 掷骰子记住数字可复现",C_ORG),
    ("Sampler(采样器)","推荐: DPM++ 2M Karras\n快速: DPM++ SDE(15步)\n极速: LCM(4-8步)\n不同算法速度质量取舍",C_PUR)]
for i,(t,d,c) in enumerate(ps):
    x=Inches(0.5+i*3.2)
    rrect(sl,x,Inches(1.3),Inches(3.0),Inches(1.8),fc=C_W,lc=c)
    rect(sl,x,Inches(1.3),Inches(3.0),Inches(0.28),fc=c)
    tb(sl,x+Inches(0.08),Inches(1.32),Inches(2.8),Inches(0.24),t,fs=11,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.1),Inches(1.62),Inches(2.7),Inches(1.4),d,fs=9.5,fc=C_D,ls=1.15)
# Prompt技巧
rrect(sl,Inches(0.5),Inches(3.3),Inches(12.3),Inches(3.6),fc=C_LT)
tb(sl,Inches(0.7),Inches(3.35),Inches(5),Inches(0.25),"Prompt 技巧 — 给AI画师下订单的艺术",fs=13,fc=C_PRI,b=True)
# 结构
tb(sl,Inches(0.7),Inches(3.65),Inches(6),Inches(0.2),"Prompt基本结构:",fs=11,fc=C_PRI,b=True)
tb(sl,Inches(0.7),Inches(3.85),Inches(12),Inches(0.2),"[主体描述] + [环境/背景] + [风格/艺术参考] + [画质/技术参数] + [负面提示词]",fs=10,fc=C_D,b=True)
# 6个技巧
tips=[("1 主语前置","A fluffy orange tabby cat sitting on a vintage sofa (具体)",C_ACC),
      ("2 具体数值","cinematic lighting, 85mm lens, f/1.8, shallow depth of field",C_GRN),
      ("3 风格关键词","photorealistic / oil painting / anime / concept art / dramatic",C_ORG),
      ("4 权重调节","(red dress:1.3)增加 / (blue sky:0.8)降低关键词权重",C_PUR),
      ("5 负面提示词","ugly, blurry, low quality, deformed, bad hands, watermark, text",C_RED),
      ("6 排列顺序","最重要(主体) → 环境 → 风格 → 画质，越靠前影响越大",C_TEL)]
for i,(t,d,c) in enumerate(tips):
    row,col=i//2,i%2
    x=Inches(0.7+col*6.3); y=Inches(4.15+row*0.8)
    rrect(sl,x,y,Inches(6.0),Inches(0.65),fc=C_W,lc=c)
    tb(sl,x+Inches(0.08),y+Inches(0.03),Inches(1.5),Inches(0.22),t,fs=10,fc=c,b=True)
    tb(sl,x+Inches(1.6),y+Inches(0.03),Inches(4.3),Inches(0.58),d,fs=9,fc=C_D,ls=1.1)
print("P12: params+prompt")

# ========== 第13页: 应用+LoRA+ControlNet ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"14-16  应用场景 + LoRA + ControlNet","从基础应用到高级控制")
# 应用场景
sc=[("概念艺术","几秒生成数十种设计方案\n游戏/电影/广告大量使用\n给设计师配备灵感无限的外脑",C_ACC),
    ("游戏影视","角色/场景/纹理素材\nControlNet控制姿态+LoRA统一风格\n批量生成风格一致资产",C_GRN),
    ("电商广告","产品图合成到各种场景\n一双鞋→海滩/雪山/城市\n成本几乎为零",C_ORG),
    ("建筑室内","快速生成设计效果图\n现代极简客厅落地窗\n用于早期方案沟通",C_PUR),
    ("个人创作","头像/小说插图/社交媒体\n完全免费开源可本地运行\n不受任何使用限制",C_TEL)]
for i,(t,d,c) in enumerate(sc):
    x=Inches(0.5+i*2.55)
    rrect(sl,x,Inches(1.3),Inches(2.35),Inches(2.2),fc=C_W,lc=c)
    tb(sl,x+Inches(0.05),Inches(1.35),Inches(2.2),Inches(0.22),t,fs=11,fc=c,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.08),Inches(1.6),Inches(2.1),Inches(1.8),d,fs=9,fc=C_D,ls=1.2)
# LoRA
rrect(sl,Inches(0.5),Inches(3.7),Inches(6),Inches(3.2),fc=RGBColor(0xE8,0xF0,0xFE))
tb(sl,Inches(0.7),Inches(3.75),Inches(5.5),Inches(0.25),"LoRA（Low-Rank Adaptation）— 给画家添加专项风格手册",fs=12,fc=C_ACC,b=True)
mtb(sl,Inches(0.7),Inches(4.05),Inches(5.6),Inches(2.7),[
    ("核心思想: 不修改原始参数，在旁边挂载一组很小的额外参数",False,C_D,10),
    ("就像给汽车加装外挂导航仪，不需要改装引擎就能获得新功能",False,C_D,10),
    ("",False,C_D,4),
    ("数学原理: W' = W + B x A (A: 1024xr, B: rx1024, r通常4-64)",True,C_PUR,10),
    ("r=16时只需3.2万参数，相比原始100万减少97%",False,C_D,10),
    ("",False,C_D,4),
    ("训练数据: 10-50张图片 | 训练时间: 30分钟-2小时(RTX3060)",False,C_D,10),
    ("文件大小: 10-200MB (主模型4GB的0.25%-5%)",False,C_D,10),
    ("可以叠加多个LoRA: 风格LoRA + 角色LoRA + 场景LoRA",True,C_GRN,10),
    ("LoRA系数0.5轻度/0.75中等/1.0完整风格影响",False,C_D,10),
])
# ControlNet
rrect(sl,Inches(6.8),Inches(3.7),Inches(6),Inches(3.2),fc=RGBColor(0xE8,0xF8,0xF5))
tb(sl,Inches(7.0),Inches(3.75),Inches(5.5),Inches(0.25),"ControlNet — 精确控制构图和姿态",fs=12,fc=C_GRN,b=True)
mtb(sl,Inches(7.0),Inches(4.05),Inches(5.6),Inches(2.7),[
    ("核心问题: 只用文字无法精确控制构图和结构",False,C_D,10),
    ("解决方案: 通过额外条件图像(草图/边缘/深度/姿态)精确控制",False,C_D,10),
    ("",False,C_D,4),
    ("架构: 复制一份U-Net编码器作为控制分支，通过零卷积注入原始U-Net",True,C_PUR,10),
    ("零卷积初始权重为0，训练开始时不影响原始模型，逐步学习注入方式",False,C_D,10),
    ("",False,C_D,4),
    ("8种条件类型:",True,C_GRN,10),
    ("Canny Edge(边缘) | Depth(深度) | OpenPose(人体姿态) | Normal Map(法线)",False,C_D,9),
    ("Scribble(简笔) | Segmentation(分割) | Lineart(线稿) | Tile(分块)",False,C_D,9),
    ("",False,C_D,4),
    ("比喻: 给画家配一个精确的构图模板，文字控制画什么，ControlNet控制画在哪里",True,C_ORG,10),
])
print("P13: apps+LoRA+ControlNet")

# ========== 第14页: SD3/SDXL + Flux ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"17-18  SD3/SDXL 最新进展 + Flux 等新模型对比","2026年AI图像生成格局")
# SDXL vs SD3
tbl(sl,Inches(0.5),Inches(1.3),Inches(6),8,4,[
    ["特性","SD 1.5","SDXL","SD3"],
    ["架构","U-Net","U-Net(更大)","MMDiT(全新)"],
    ["参数量","~1B","~3.5B","800M~8B"],
    ["文本编码器","CLIP L/14","CLIP+OpenCLIP","CLIP-L+CLIP-G+T5-XXL"],
    ["默认分辨率","512x512","1024x1024","1024x1024"],
    ["训练方法","DDPM噪声预测","DDPM噪声预测","Flow Matching"],
    ["VAE通道","4通道","4通道(更大)","16通道"],
    ["文字渲染","较差","一般","显著提升"],
],cw=[Inches(1.3),Inches(1.3),Inches(1.5),Inches(1.9)],fs=9)
# SD3亮点
rrect(sl,Inches(6.8),Inches(1.3),Inches(6),Inches(2.2),fc=RGBColor(0xF5,0xEE,0xF8))
tb(sl,Inches(7.0),Inches(1.35),Inches(5.5),Inches(0.25),"SD3 核心创新",fs=12,fc=C_PUR,b=True)
mtb(sl,Inches(7.0),Inches(1.65),Inches(5.6),Inches(1.8),[
    ("MMDiT: 纯Transformer架构，双流注意力(图像+文本双向交流)",False,C_D,10),
    ("三编码器: CLIP-L(1.2亿)+CLIP-G(3.5亿)+T5-XXL(46亿)=50亿参数仅用于文本理解",False,C_D,10),
    ("Flow Matching: 从噪声分布到目标分布的直接路径，比DDPM更高效",False,C_D,10),
    ("U-Net是单向广播系统; MMDiT是双向对讲系统，画师能根据画面反问客户",True,C_ACC,10),
])
# Flux对比
tb(sl,Inches(0.5),Inches(3.7),Inches(5),Inches(0.25),"Flux + 2026年主要模型对比",fs=12,fc=C_PRI,b=True)
tbl(sl,Inches(0.5),Inches(4.0),Inches(12.3),7,6,[
    ["模型","开发者","参数量","文字渲染","开源","ELO评分"],
    ["GPT Image 1.5","OpenAI","N/A","最强","闭源","1264"],
    ["Gemini 3 Pro Image","Google","N/A","强","闭源","1241"],
    ["Midjourney V7/V8","Midjourney","N/A","中","闭源","N/A"],
    ["Flux 2 Dev","Black Forest Labs","~12B","优秀","开放","~1200"],
    ["SD3.5 Large","Stability AI","8B","较好","部分","~1170"],
    ["Hunyuan Image 3.0","腾讯","N/A","中文强","部分","~1130"],
],cw=[Inches(2.2),Inches(2.5),Inches(1.5),Inches(1.5),Inches(1.3),Inches(3.3)],fs=9)
# 关键洞察
rrect(sl,Inches(0.5),Inches(6.3),Inches(12.3),Inches(0.6),fc=C_LT)
mtb(sl,Inches(0.7),Inches(6.33),Inches(11.8),Inches(0.5),[
    ("Flux由SD原始团队创建，文字渲染和手部生成显著改善。Flux 2 Dev与Max的ELO差距仅19分，免费开源已接近商业模型质量。",False,C_D,10),
    ("2026年格局: 质量差距缩小，竞争转向专业化。多模型策略(根据用途选择不同模型)成为专业用户主流。",True,C_PRI,10),
])
print("P14: SD3+Flux")

# ========== 第15页: 2026动态+总结 ==========
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl); hdr(sl,"19  2026年AI图像生成动态 + 总结","7大趋势与未来展望")
ts3=[("1 质量差距缩小","前9名ELO差距仅117分\n竞争从谁更好看转向谁在特定场景更好\n多模型策略成为主流工作方式",C_ACC),
     ("2 4K输出成标配","主流模型普遍支持4K\n多档分辨率和多宽高比原生生成\n横幅广告等特殊尺寸不再需要裁剪",C_GRN),
     ("3 实时信息整合","Grounding功能实时搜索互联网\n要求生成最新款iPhone会先搜索确认\n解决AI知识截止问题",C_ORG),
     ("4 角色一致性突破","最多5个角色+14个物体保持一致\n可生成完整故事板\n漫画创作/品牌营销意义重大",C_PUR),
     ("5 Flash模型改变工作流","极低价格接近Pro级质量\n迭代式AI图像工作流首次经济可行\n从一次生成变成快速迭代精雕细琢",C_RED),
     ("6 开源vs闭源博弈","Flux 2 Dev与Max差距仅19分\n免费开源已接近顶级商业质量\n对数据敏感企业是重大利好",C_TEL),
     ("7 视频生成融合","LTX 2.3(220亿参数,带同步音频4K视频)\nHelios(消费级GPU实时60秒视频)\n图像和视频技术边界快速模糊",C_TEL)]
for i,(t,d,c) in enumerate(ts3):
    if i<4:
        x=Inches(0.5+i*3.2)
        y=Inches(1.3)
        h=Inches(2.2)
    else:
        x=Inches(0.5+(i-4)*4.2)
        y=Inches(3.7)
        h=Inches(2.0)
    rrect(sl,x,y,Inches(3.0),h,fc=C_W,lc=c)
    rect(sl,x,y,Inches(3.0),Inches(0.25),fc=c)
    tb(sl,x+Inches(0.08),y+Inches(0.02),Inches(2.8),Inches(0.2),t,fs=10,fc=C_W,b=True,a=PP_ALIGN.CENTER)
    tb(sl,x+Inches(0.08),y+Inches(0.3),Inches(2.7),Inches(h-Inches(0.4)),d,fs=9,fc=C_D,ls=1.15)
# 总结
rrect(sl,Inches(0.5),Inches(5.9),Inches(12.3),Inches(1.0),fc=C_PRI)
mtb(sl,Inches(0.7),Inches(5.95),Inches(11.8),Inches(0.9),[
    ("2026年AI图像生成就像从功能手机时代进入了智能手机时代。",True,C_W,14),
    ("功能手机时代(2022-2023)比谁屏幕更大、待机更长。智能手机时代(2026)比生态——应用商店、多任务处理、云端整合。",False,RGBColor(0xAA,0xCC,0xEE),10),
    ("最终胜出的可能不是某一个模型，而是能够灵活调度多种模型的工作流平台。",True,C_W,11),
])
print("P15: trends+summary")

# ========== 保存 ==========
prs.save(OUT)
print(f"\nDone! {OUT} | {len(prs.slides)} pages")
