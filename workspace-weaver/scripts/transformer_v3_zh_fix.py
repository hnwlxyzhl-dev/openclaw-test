#!/usr/bin/env python3
"""
Transformer PPT v3 中文版 - 质量修复脚本
修复内容:
  P0: 76处字体过小(<9pt) → 统一为9pt
  P1: Part1(1-11页)和Part2(12-22页)风格统一
  P2: 排版优化
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from collections import Counter
import copy

PPT_PATH = '/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3_zh.pptx'

# ===== 标准配色 =====
COLOR_TITLE = RGBColor(0x1E, 0x3A, 0x5F)      # #1E3A5F 标题
COLOR_BODY = RGBColor(0x2C, 0x3E, 0x50)         # #2C3E50 正文
COLOR_ENCODER = RGBColor(0x34, 0x98, 0xDB)      # #3498DB 编码器
COLOR_DECODER = RGBColor(0xE7, 0x4C, 0x3C)      # #E74C3C 解码器
COLOR_ATTENTION = RGBColor(0xF3, 0x9C, 0x12)    # #F39C12 注意力
COLOR_FFN = RGBColor(0x2E, 0xCC, 0x71)          # #2ECC71 FFN
COLOR_EMBED = RGBColor(0x9B, 0x59, 0xB6)        # #9B59B6 嵌入

SIZE_TITLE = Pt(20)
SIZE_BODY_MIN = Pt(9)
SIZE_BODY_NORMAL = Pt(10)
SIZE_BODY_MAX = Pt(11)

# Shape fill color mapping - keywords to colors
KEYWORD_COLOR_MAP = {
    '编码器': COLOR_ENCODER,
    '解码器': COLOR_DECODER,
    '注意力': COLOR_ATTENTION,
    '自注意力': COLOR_ATTENTION,
    '多头注意力': COLOR_ATTENTION,
    '交叉注意力': COLOR_ATTENTION,
    'ffn': COLOR_FFN,
    '前馈': COLOR_FFN,
    '嵌入': COLOR_EMBED,
    '词嵌入': COLOR_EMBED,
    '位置编码': COLOR_EMBED,
    '位置嵌入': COLOR_EMBED,
}


def analyze_font_distribution(prs):
    """分析修复前的字体分布"""
    size_counter = Counter()
    small_fonts = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.size is not None:
                        pt_val = para.font.size.pt
                        size_counter[pt_val] += 1
                        if pt_val < 9:
                            small_fonts.append((i+1, shape.name, pt_val, para.text[:40]))
    return size_counter, small_fonts


def fix_p0_small_fonts(prs):
    """P0: 修复所有<9pt的字体为9pt"""
    fixed_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.size is not None and para.font.size < Pt(9):
                        para.font.size = Pt(9)
                        fixed_count += 1
    return fixed_count


def is_title_shape(shape, slide_index):
    """判断是否是标题shape（每页第一个大字号文本或特定名称）"""
    name_lower = shape.name.lower() if shape.name else ''
    if 'title' in name_lower:
        return True
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.font.size is not None and para.font.size >= Pt(18):
                return True
    return False


def fix_p1_style_consistency(prs):
    """P1: 统一Part1和Part2的风格"""
    fixes = {
        'title_size_fixed': 0,
        'title_color_fixed': 0,
        'title_bold_fixed': 0,
        'body_color_fixed': 0,
        'body_size_fixed': 0,
    }

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = is_title_shape(shape, page_num)

            for para in shape.text_frame.paragraphs:
                pf = para.font

                # 标题统一
                if is_title:
                    if pf.size is not None and pf.size != SIZE_TITLE:
                        # 保留封面大标题(36pt)和副标题
                        if pf.size > Pt(25):
                            pass  # 保留封面大标题
                        else:
                            pf.size = SIZE_TITLE
                            fixes['title_size_fixed'] += 1

                    # 标题颜色统一
                    try:
                        if pf.color and pf.color.rgb and pf.color.rgb != COLOR_TITLE:
                            # 白色标题保留（在深色背景上）
                            if pf.color.rgb != RGBColor(0xFF, 0xFF, 0xFF):
                                pf.color.rgb = COLOR_TITLE
                                fixes['title_color_fixed'] += 1
                    except:
                        pass

                    # 标题加粗
                    if not pf.bold:
                        pf.bold = True
                        fixes['title_bold_fixed'] += 1

                else:
                    # 正文统一
                    # 颜色：非白色文字统一为COLOR_BODY
                    try:
                        if pf.color and pf.color.rgb:
                            current = pf.color.rgb
                            # 保留白色（深色背景上的文字）
                            if current == RGBColor(0xFF, 0xFF, 0xFF):
                                pass
                            # 保留组件颜色（编码器/解码器等）
                            elif current in (COLOR_ENCODER, COLOR_DECODER, COLOR_ATTENTION,
                                           COLOR_FFN, COLOR_EMBED, COLOR_TITLE):
                                pass
                            # 保留灰色说明文字
                            elif current in (RGBColor(0x95, 0xA5, 0xA6), RGBColor(0x7F, 0x8C, 0x8D)):
                                pass
                            # 其他统一为正文色
                            else:
                                pf.color.rgb = COLOR_BODY
                                fixes['body_color_fixed'] += 1
                    except:
                        pass

                    # 正文大小：9-11pt范围内，太小的已由P0修复
                    if pf.size is not None:
                        if pf.size < Pt(9):
                            pf.size = Pt(9)
                            fixes['body_size_fixed'] += 1
                        elif pf.size > Pt(14):
                            # 非标题的大字可能是副标题，保留
                            pass

    return fixes


def fix_p2_layout(prs):
    """P2: 排版优化 - 文字框边距"""
    margin_fixes = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                # 设置合理的边距
                # word_wrap
                tf.word_wrap = True

                # 边距检查 (单位: Emu, Inches(0.1) ≈ 91440 Emu)
                min_left = Inches(0.15)
                min_top = Inches(0.05)
                min_right = Inches(0.15)
                min_bottom = Inches(0.05)

                if tf.margin_left < min_left:
                    tf.margin_left = min_left
                    margin_fixes += 1
                if tf.margin_top < min_top:
                    tf.margin_top = min_top
                    margin_fixes += 1
                if tf.margin_right < min_right:
                    tf.margin_right = min_right
                    margin_fixes += 1
                if tf.margin_bottom < min_bottom:
                    tf.margin_bottom = min_bottom
                    margin_fixes += 1

    return margin_fixes


def compute_slide_coverage(slide, slide_w, slide_h):
    """计算shape覆盖率"""
    total_area = slide_w * slide_h
    covered = 0
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                w = shape.width
                h = shape.height
                if w > 0 and h > 0:
                    covered += w * h
        except (ValueError, AttributeError, TypeError):
            pass
    return covered / total_area if total_area > 0 else 0


def main():
    print("=" * 60)
    print("Transformer PPT v3 中文版 - 质量修复")
    print("=" * 60)

    # 加载PPT
    prs = Presentation(PPT_PATH)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ===== 修复前分析 =====
    print("\n📊 修复前分析:")
    print("-" * 40)
    before_sizes, before_small = analyze_font_distribution(prs)
    print(f"字体大小分布:")
    for size, count in sorted(before_sizes.items(), key=lambda x: float(x[0])):
        label = " ⚠️ <9pt" if float(size) < 9 else ""
        print(f"  {size}pt: {count}处{label}")
    print(f"\n小字体(<9pt)总数: {len(before_small)}处")

    # 滑动覆盖率
    print(f"\n各页shape覆盖率:")
    for i, slide in enumerate(prs.slides):
        cov = compute_slide_coverage(slide, slide_w, slide_h)
        flag = " ⚠️ 低覆盖" if cov < 0.5 else ""
        print(f"  第{i+1:2d}页: {cov:.1%}{flag}")

    # ===== P0: 修复小字体 =====
    print("\n🔧 P0: 修复小字体(<9pt → 9pt)")
    print("-" * 40)
    p0_fixed = fix_p0_small_fonts(prs)
    print(f"✅ 修复了 {p0_fixed} 处小字体")

    # ===== P1: 风格统一 =====
    print("\n🔧 P1: 统一Part1/Part2风格")
    print("-" * 40)
    p1_fixes = fix_p1_style_consistency(prs)
    for k, v in p1_fixes.items():
        if v > 0:
            print(f"  ✅ {k}: {v}处")

    # ===== P2: 排版优化 =====
    print("\n🔧 P2: 排版优化")
    print("-" * 40)
    p2_fixed = fix_p2_layout(prs)
    print(f"✅ 文字框边距修复: {p2_fixed}处")

    # ===== 修复后验证 =====
    print("\n📊 修复后验证:")
    print("-" * 40)
    after_sizes, after_small = analyze_font_distribution(prs)
    print(f"字体大小分布:")
    for size, count in sorted(after_sizes.items(), key=lambda x: float(x[0])):
        label = " ⚠️ <9pt" if float(size) < 9 else ""
        print(f"  {size}pt: {count}处{label}")
    print(f"\n小字体(<9pt)剩余: {len(after_small)}处 {'✅ 全部修复!' if len(after_small) == 0 else '❌ 仍有残留!'}")

    # 保存
    prs.save(PPT_PATH)
    print(f"\n💾 已保存到: {PPT_PATH}")

    # ===== 汇总 =====
    print("\n" + "=" * 60)
    print("📋 修复汇总")
    print("=" * 60)
    print(f"  P0 小字体修复: {p0_fixed}处")
    print(f"  P1 风格统一: {sum(p1_fixes.values())}处")
    print(f"  P2 排版优化: {p2_fixed}处")
    print(f"  修复前<9pt: {len(before_small)}处")
    print(f"  修复后<9pt: {len(after_small)}处")
    print("=" * 60)


if __name__ == '__main__':
    main()
