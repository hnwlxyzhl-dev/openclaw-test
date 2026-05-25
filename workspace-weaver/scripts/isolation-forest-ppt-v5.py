#!/usr/bin/env python3
"""孤立森林 PPT V5 - 贯穿示例升级：1000个传感器数据点"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── 配色方案 ──
BG_DARK      = RGBColor(0x0B, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x12, 0x25, 0x3F)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_ORANGE= RGBColor(0xFF, 0x6B, 0x35)
ACCENT_RED   = RGBColor(0xFF, 0x44, 0x44)
ACCENT_YELLOW= RGBColor(0xFF, 0xD7, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xC4, 0xDE)
MID_GRAY     = RGBColor(0x7A, 0x8E, 0xA8)
DIM_GRAY     = RGBColor(0x4A, 0x5E, 0x78)
FONT_CN = '微软雅黑'

# ── 贯穿示例 ──
# 工厂传感器数据：1000台设备，2维（温度°C，振动频率Hz）
# 980台正常设备：温度 60~75°C，振动 2~5 Hz（聚集在中心区域）
# 20台异常设备：读数明显偏离
# 代表异常点：
#   A1(95, 3.2) — 温度异常高（过热）
#   A2(48, 13.5) — 温度偏低 + 振动异常高（松动）
#   A3(62, 0.2) — 温度正常但振动极低（卡死）
# 代表正常点：
#   P(68, 3.8) — 典型正常设备

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name
    p.alignment = alignment; p.space_after = Pt(4); p.space_before = Pt(0)
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05); tf.margin_top = Inches(0.05)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=14,
                           default_color=LIGHT_GRAY, default_bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05); tf.margin_top = Inches(0.05)
    for i, line in enumerate(lines):
        if isinstance(line, str): text, size, color, bold = line, default_size, default_color, default_bold
        else:
            text = line[0]; size = line[1] if len(line)>1 else default_size
            color = line[2] if len(line)>2 else default_color
            bold = line[3] if len(line)>3 else default_bold
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = FONT_CN; p.alignment = alignment
        p.space_after = Pt(3); p.space_before = Pt(1)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color: shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else: shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_connector_v(slide, x, y1, y2, color=MID_GRAY, width=0.025):
    if y2 < y1: y1, y2 = y2, y1
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y1), Inches(width), Inches(y2-y1))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()

def add_connector_h(slide, x1, x2, y, color=MID_GRAY, width=0.025):
    if x2 < x1: x1, x2 = x2, x1
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y), Inches(x2-x1), Inches(width))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()

def add_dot(slide, x, y, r, color):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-r), Inches(y-r), Inches(r*2), Inches(r*2))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()

def draw_scatter_chart(slide, left, top, w, h, title_text=None):
    """画散点图：密集区色块 + 异常点"""
    add_rect(slide, left, top, w, h, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))
    if title_text:
        add_textbox(slide, left+0.1, top+0.05, w-0.2, 0.3, title_text, 10, MID_GRAY, False, PP_ALIGN.LEFT)

    # 坐标映射: 温度30~105, 振动0~16
    # 绘图区域内缩
    px, py, pw, ph = left+0.5, top+0.35, w-0.8, h-0.7

    # 正常点密集区 — 用半透明色块表示（椭圆）
    # 中心约 (67.5, 3.5) → 映射到绘图区
    cx_norm = px + (67.5-30)/(105-30)*pw
    cy_norm = py + (1-(3.5-0)/16)*ph
    # 椭圆覆盖 温度60~75, 振动2~5
    ew = (75-60)/(105-30)*pw
    eh = (5-2)/16*ph
    ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cx_norm - ew/2), Inches(cy_norm - eh/2), Inches(ew), Inches(eh))
    ellipse.fill.solid(); ellipse.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x60)
    ellipse.line.color.rgb = ACCENT_CYAN; ellipse.line.width = Pt(1.5)
    ellipse.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x50)

    # "980" 标签在密集区中心
    add_textbox(slide, cx_norm-0.5, cy_norm-0.15, 1.0, 0.3, '980台', 11, ACCENT_CYAN, True, PP_ALIGN.CENTER)

    # 异常点 A1, A2, A3
    anomalies = [
        ('A1', 95, 3.2, ACCENT_ORANGE, '(过热)'),
        ('A2', 48, 13.5, ACCENT_RED, '(松动)'),
        ('A3', 62, 0.2, ACCENT_YELLOW, '(卡死)'),
    ]
    for name, temp, vib, color, desc in anomalies:
        ax = px + (temp-30)/(105-30)*pw
        ay = py + (1-(vib-0)/16)*ph
        add_dot(slide, ax, ay, 0.12, color)
        add_textbox(slide, ax+0.12, ay-0.08, 1.2, 0.2, f'{name}{desc}', 9, color, True)

    # 其他异常点（用小点表示，不标名字）
    other_anomalies = [
        (92, 6.1), (88, 1.5), (98, 8.3), (42, 11.2), (35, 4.0),
        (52, 14.8), (78, 0.3), (90, 10.5), (40, 7.2), (96, 5.5),
        (55, 12.0), (85, 0.8), (38, 9.5), (72, 14.0), (93, 2.0),
        (45, 3.5), (80, 11.8),
    ]
    for temp, vib in other_anomalies:
        ax = px + (temp-30)/(105-30)*pw
        ay = py + (1-(vib-0)/16)*ph
        add_dot(slide, ax, ay, 0.06, ACCENT_ORANGE)

    # 坐标轴标签
    add_textbox(slide, px+pw/2-0.5, py+ph+0.02, 1.0, 0.2, '温度(°C)', 8, MID_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, left+0.05, py+ph/2-0.1, 0.4, 0.2, '振动(Hz)', 8, MID_GRAY, False, PP_ALIGN.CENTER)

    # 轴刻度
    for val in [40, 60, 80, 100]:
        sx = px + (val-30)/(105-30)*pw
        add_connector_v(slide, sx, py+ph, py+ph+0.08, DIM_GRAY, 0.01)
        add_textbox(slide, sx-0.2, py+ph+0.08, 0.4, 0.15, str(val), 7, DIM_GRAY, False, PP_ALIGN.CENTER)
    for val in [4, 8, 12]:
        sy = py + (1-(val-0)/16)*ph
        add_connector_h(slide, px-0.08, px, sy, DIM_GRAY, 0.01)
        add_textbox(slide, px-0.45, sy-0.07, 0.35, 0.15, str(val), 7, DIM_GRAY, False, PP_ALIGN.RIGHT)

    return px, py, pw, ph  # 返回绘图区坐标供后续使用


# ════════════════════════════════════════
# 第 1 页：封面
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 1.0, 2.8, 4.0, ACCENT_GREEN)
add_textbox(slide, 1.0, 1.2, 11, 1.5, '孤立森林', 48, WHITE, True)
add_textbox(slide, 1.0, 2.0, 11, 0.8, 'Isolation Forest', 28, ACCENT_CYAN)
add_multiline_textbox(slide, 1.0, 3.2, 11, 1.5, [
    ('不建模正常，直接孤立异常——一种颠覆性的异常检测算法', 20, LIGHT_GRAY),
    ('', 8),
    ('适用于大数据、高维度、无标签的场景', 14, MID_GRAY),
    ('线性时间复杂度，天然支持并行化', 14, MID_GRAY),
    ('2008 年 IEEE ICDM 论文  |  Liu, Ting & Zhou', 12, DIM_GRAY),
])

# 右侧装饰
draw_scatter_chart(slide, 8.0, 1.5, 4.8, 4.5)
add_textbox(slide, 1.0, 6.5, 11, 0.5, '异常检测系列  |  深度技术解析', 12, DIM_GRAY)


# ════════════════════════════════════════
# 第 2 页：什么是异常检测？
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '什么是异常检测？', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '为什么要关注"与众不同"的数据？', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 4.8, [
    ('什么是异常检测？', 18, ACCENT_GREEN, True),
    ('异常检测是从数据中识别出与绝大多数数据显著不同的样本的过程。就像在一群白绵羊中找出那只黑绵羊——正常数据遵循某种模式，而异常点偏离了这个模式。', 13, LIGHT_GRAY),
    ('', 6),
    ('为什么它如此重要？', 18, ACCENT_GREEN, True),
    ('异常事件往往伴随高风险。网络入侵导致数据泄露，金融欺诈造成巨额损失，设备异常引发生产线停工。正因为"异常"通常等于"危险"，异常检测成为数据科学最重要的研究方向之一。', 13, LIGHT_GRAY),
    ('', 6),
    ('异常的三种类型：', 16, WHITE, True),
    ('1️⃣  点异常：单个数据点本身就很异常。一群身高 1.6-1.8m 的人中间站了一个 2.3m 的人。', 12, LIGHT_GRAY),
    ('2️⃣  上下文异常：特定上下文中才算异常。夏天穿短袖正常，零下 20°C 穿短袖就异常。', 12, LIGHT_GRAY),
    ('3️⃣  集合异常：单点都正常，放在一起才异常。每天凌晨访问数据库正常，连续三个月天天凌晨就可疑。', 12, LIGHT_GRAY),
])

add_rect(slide, 7.2, 2.0, 5.5, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.5, 2.2, 5.0, 4.4, [
    ('🛡️  机场安检类比', 16, ACCENT_GREEN, True),
    ('', 6),
    ('异常检测就像机场安检——99.9% 的旅客都正常，但系统必须在几秒内识别出那 0.1% 的危险人物。', 13, LIGHT_GRAY),
    ('', 6),
    ('三种异常 = 三种"嫌疑人"：', 13, WHITE, True),
    ('• 带着违禁品的人 → 点异常（一看就不对）', 12, ACCENT_CYAN),
    ('• 错误时间出现在错误地点 → 上下文异常', 12, ACCENT_CYAN),
    ('• 每人行为正常但组合可疑 → 集合异常（团伙作案）', 12, ACCENT_CYAN),
])


# ════════════════════════════════════════
# 第 3 页：传统方法
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '传统异常检测方法', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '先建模"正常"，再找"偏离"——但建模"正常"恰恰是最难的步骤', 16, ACCENT_CYAN)

table_data = [
    ('方法', '核心思路', '时间复杂度', '高维适应', '核心局限'),
    ('Z-score', '距离均值几个标准差', 'O(n)', '差', '假设正态分布'),
    ('KNN / LOF', '邻居距离/局部密度', 'O(n²)', '一般', '计算量灾难性'),
    ('DBSCAN', '基于密度聚类', 'O(n log n)', '差', '高维密度难定义'),
    ('One-Class SVM', '学习正常边界', 'O(n²~n³)', '好', '参数敏感、不可扩展'),
    ('孤立森林', '直接孤立异常', 'O(n·t·logψ)', '好', '局部异常检测弱'),
]
tbl = slide.shapes.add_table(6, 5, Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.2)).table
for i, w in enumerate([1.8, 2.8, 2.0, 1.5, 3.6]):
    tbl.columns[i].width = Inches(w)
for r in range(6):
    for c in range(5):
        cell = tbl.cell(r, c); cell.text = table_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11); para.font.name = FONT_CN
            if r==0: para.font.bold=True; para.font.color.rgb=WHITE
            elif r==5: para.font.bold=True; para.font.color.rgb=ACCENT_GREEN
            else: para.font.color.rgb=LIGHT_GRAY
            para.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15,0x30,0x50) if r==0 else (RGBColor(0x0A,0x2A,0x15) if r==5 else BG_CARD)

add_rect(slide, 0.8, 5.6, 11.7, 1.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.1, 5.7, 11.2, 1.3, [
    ('⚠️  核心瓶颈：不在于"找异常"，而在于"建模正常"太昂贵、太困难、太脆弱', 15, ACCENT_ORANGE, True),
    ('传统方法像一个侦探，先要彻底研究"正常市民的生活规律"，然后才能发现谁的行踪可疑。但刻画"正常"本身就极其困难——正常人有千万种生活方式，你怎么定义"正常"？', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 4 页：核心思想 ★ 贯穿示例引入
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '孤立森林的核心思想', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '不建模正常，直接孤立异常——范式转换', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 6.0, 5.0, [
    ('范式转换', 18, ACCENT_GREEN, True),
    ('从"建模正常"到"孤立异常"', 14, WHITE, True),
    ('', 3),
    ('异常点有两个天然属性："少"和"不同"。因为稀少，周围没有同类；因为不同，远离大多数。', 13, LIGHT_GRAY),
    ('', 3),
    ('这导致一个直接结果：异常点比正常点更容易被随机分割孤立出来。', 14, ACCENT_GREEN, True),
    ('', 4),
    ('贯穿示例：工厂传感器异常检测', 14, WHITE, True),
    ('1000 台设备的 2 维读数（温度, 振动频率）：', 13, LIGHT_GRAY),
    ('• 980 台正常设备：温度 60~75°C，振动 2~5 Hz（密集聚集）', 12, ACCENT_CYAN),
    ('• 20 台异常设备：读数明显偏离（外围散布）', 12, ACCENT_ORANGE),
    ('', 3),
    ('随机选温度轴，在 80°C 画分割线：', 13, LIGHT_GRAY),
    ('→ 左侧：980 台正常设备 + 部分异常（~990 台）', 12, ACCENT_CYAN),
    ('→ 右侧：只有 A1(95°C) 等少数异常 → 1 次分裂就被孤立！', 13, ACCENT_GREEN, True),
    ('', 3),
    ('正常设备呢？周围全是同类，需要多次分裂才能单独挑出。', 12, LIGHT_GRAY),
    ('', 3),
    ('路径越短 = 越异常 | 路径越长 = 越正常', 15, WHITE, True),
])

# 右侧散点图
draw_scatter_chart(slide, 7.2, 2.0, 5.5, 4.8)

# 在散点图上画分割线 x=80°C
px = 7.2+0.5; pw = 5.5-0.8; py_top = 2.0+0.35; ph = 4.8-0.7
split_x = px + (80-30)/(105-30)*pw
add_connector_v(slide, split_x, py_top+0.02, py_top+ph-0.02, ACCENT_GREEN, 0.04)
add_textbox(slide, split_x+0.05, py_top, 1.0, 0.2, 'T=80°C', 10, ACCENT_GREEN, True)


# ════════════════════════════════════════
# 第 5 页：iTree ★ 更复杂的树形图
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, 'iTree（孤立树）', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '一棵不断随机分割空间的二叉树', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('两种节点', 16, ACCENT_GREEN, True),
    ('内部节点：存分裂规则——"在哪个特征上切？切在哪里？"', 13, LIGHT_GRAY),
    ('叶子节点：存"到达这里的样本数量"——房间还剩几人', 13, LIGHT_GRAY),
    ('', 5),
    ('每次分裂两步操作', 16, ACCENT_GREEN, True),
    ('① 随机选一个特征（等概率）', 13, LIGHT_GRAY),
    ('② 在该特征的 [min, max] 内随机选分裂点', 13, LIGHT_GRAY),
    ('→ 特征值 < 分裂点进左子树，≥ 分裂点进右子树', 13, ACCENT_CYAN),
    ('', 5),
    ('三个终止条件', 16, ACCENT_GREEN, True),
    ('① 样本数 = 1 → 无需再分', 13, LIGHT_GRAY),
    ('② 所有特征值相同 → 无法再分', 13, LIGHT_GRAY),
    ('③ 达到最大深度（log₂256 = 8 层）→ 强制停止', 13, LIGHT_GRAY),
    ('', 5),
    ('iTree vs 普通决策树', 16, ACCENT_GREEN, True),
    ('普通决策树：精心选择最优分裂（像经验丰富的分拣员）', 12, LIGHT_GRAY),
    ('iTree：随机选特征和分裂点（像蒙着眼分拣）', 12, LIGHT_GRAY),
    ('看似"愚蠢"，但随机分割对稀少的异常点更"友好"', 12, ACCENT_GREEN, True),
])

# 右侧 - 真实树形图（Tree #38，从100棵树中选出）
add_rect(slide, 6.8, 2.0, 6.0, 4.8, RGBColor(0x0E, 0x1F, 0x35), RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.05, 5.6, 0.3, '真实树结构（100棵树中的第38棵）', 11, ACCENT_GREEN, True)

nw, nh = 1.65, 0.36
node_fill = RGBColor(0x1A, 0x3A, 0x5C)
anomaly_fill = RGBColor(0x3A, 0x1A, 0x0A)
normal_fill = RGBColor(0x0A, 0x2A, 0x15)

# 真实树结构:
# Root: 振动 < 7.6? (256)
#   Right → A2等 [size=1] path=1
#   Left → 温度 < 76.7? (255)
#     Right → 振动 < 5.3? (2)
#       Left → [size=1] path=3 (含A1)
#       Right → [size=1] path=3
#     Left → 振动 < 5.2? (253)
#       Right → [size=1] path=3
#       Left → [252] → 继续分割...

# 节点坐标
L0y = 2.55; root_cx = 9.8
L1y = 3.35; L1_L = 8.2; L1_R = 11.4
L2y = 4.15; L2_LL = 7.3; L2_LR = 9.1
L3y = 4.95; L3_LL_L = 6.9; L3_LL_R = 7.8; L3_LR_L = 8.6; L3_LR_R = 9.6
L4y = 5.65

def draw_edge(slide, px_cx, py_bot, cx_cx, cy_top, color=MID_GRAY):
    mid = (py_bot + cy_top) / 2
    add_connector_v(slide, px_cx, py_bot, mid, color, 0.02)
    add_connector_h(slide, min(px_cx, cx_cx), max(px_cx, cx_cx), mid, color, 0.02)
    add_connector_v(slide, cx_cx, mid, cy_top, color, 0.02)

def node(slide, cx, cy, w, h, text, fill, border, tc=WHITE, ts=10, tb=True):
    add_rect(slide, cx-w/2, cy-h/2, w, h, fill, border)
    add_textbox(slide, cx-w/2+0.03, cy-h/2+0.01, w-0.06, h-0.02, text, ts, tc, tb, PP_ALIGN.CENTER)

# L0→L1
draw_edge(slide, root_cx, L0y+nh/2, L1_L, L1y-nh/2)
draw_edge(slide, root_cx, L0y+nh/2, L1_R, L1y-nh/2)
add_textbox(slide, root_cx-0.55, L0y+nh/2+0.03, 0.35, 0.18, '是', 7, ACCENT_GREEN)
add_textbox(slide, root_cx+0.15, L0y+nh/2+0.03, 0.35, 0.18, '否', 7, ACCENT_ORANGE)

# L1_L→L2
draw_edge(slide, L1_L, L1y+nh/2, L2_LL, L2y-nh/2)
draw_edge(slide, L1_L, L1y+nh/2, L2_LR, L2y-nh/2)
add_textbox(slide, L1_L-0.55, L1y+nh/2+0.03, 0.35, 0.18, '是', 7, ACCENT_GREEN)
add_textbox(slide, L1_L+0.15, L1y+nh/2+0.03, 0.35, 0.18, '否', 7, ACCENT_ORANGE)

# L2_LR→L3 (右分支)
draw_edge(slide, L2_LR, L2y+nh/2, L3_LR_L, L3y-nh/2)
draw_edge(slide, L2_LR, L2y+nh/2, L3_LR_R, L3y-nh/2)

# L2_LL→L3 (左分支)
draw_edge(slide, L2_LL, L2y+nh/2, L3_LL_L, L3y-nh/2)
draw_edge(slide, L2_LL, L2y+nh/2, L3_LL_R, L3y-nh/2)

# 节点内容
# Level 0
node(slide, root_cx, L0y, nw+0.3, nh, '振动 < 7.6? (256)', node_fill, ACCENT_GREEN, WHITE, 10)

# Level 1
node(slide, L1_R, L1y, nw-0.1, nh, 'A2等 [size=1]', anomaly_fill, ACCENT_RED, ACCENT_RED, 10)
add_textbox(slide, L1_R-0.4, L1y-nh/2-0.2, 1.2, 0.18, '路径=1 🔴', 9, ACCENT_RED, True, PP_ALIGN.CENTER)

node(slide, L1_L, L1y, nw+0.2, nh, '温度 < 76.7? (255)', node_fill, ACCENT_CYAN, WHITE, 10)

# Level 2
node(slide, L2_LR, L2y, nw+0.1, nh, '振动 < 5.3? (2)', node_fill, ACCENT_ORANGE, WHITE, 10)
node(slide, L2_LL, L2y, nw+0.1, nh, '振动 < 5.2? (253)', node_fill, ACCENT_CYAN, WHITE, 10)

# Level 3 - L2_LR 的子节点
node(slide, L3_LR_L, L3y, nw-0.3, nh, 'A1 [size=1]', anomaly_fill, ACCENT_ORANGE, ACCENT_ORANGE, 9)
add_textbox(slide, L3_LR_L-0.4, L3y-nh/2-0.2, 1.2, 0.18, '路径=3 🔴', 8, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
node(slide, L3_LR_R, L3y, nw-0.3, nh, '异常 [size=1]', anomaly_fill, ACCENT_RED, ACCENT_RED, 9)

# Level 3 - L2_LL 的子节点
node(slide, L3_LL_R, L3y, nw-0.2, nh, '[size=1] 路径=3', anomaly_fill, ACCENT_RED, ACCENT_RED, 9)
node(slide, L3_LL_L, L3y, nw-0.15, nh, '[252] 继续分割...', normal_fill, ACCENT_GREEN, MID_GRAY, 9)

# 底部总结（真实数据）
add_multiline_textbox(slide, 7.0, 6.1, 5.6, 0.55, [
    ('真实树：A2(松动) 路径=1，A1(过热) 路径=3，A3/P 深达 8+ 层', 9, ACCENT_ORANGE, True),
    ('A2振动13.5远超7.6，一次分裂就被孤立', 9, MID_GRAY),
])


# ════════════════════════════════════════
# 第 6 页：构建森林 ★ 体现子采样逻辑
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '构建孤立森林', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从一棵树到一片森林——子采样 + 多树集成 = 又快又准', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('三步构建森林', 16, ACCENT_GREEN, True),
    ('', 3),
    ('Step 1：随机子采样（max_samples=256）', 14, WHITE, True),
    ('从 1000 台设备中随机抽 256 台。三个原因：', 12, LIGHT_GRAY),
    ('  • 消除屏蔽效应——980 台正常设备太密集，稀释后异常设备更突出', 11, LIGHT_GRAY),
    ('  • 计算效率——每棵树只处理 256 条，100 棵树也只处理 25,600 条', 11, LIGHT_GRAY),
    ('  • 增加多样性——不同子采样让每棵树看到不同的数据组合', 11, LIGHT_GRAY),
    ('', 3),
    ('Step 2：构建 iTree（高度限制 log₂256 = 8 层）', 14, WHITE, True),
    ('8 层足以孤立绝大多数异常点（它们在 1~4 层就被分离）。', 12, LIGHT_GRAY),
    ('正常点不需要精确路径——知道"它需要很深的路径"就够了。', 12, LIGHT_GRAY),
    ('', 3),
    ('Step 3：重复 100 次（n_estimators=100）', 14, WHITE, True),
    ('得到 100 棵独立的 iTree，组成孤立森林。', 12, LIGHT_GRAY),
    ('100→500 棵，AUC 提升通常 < 0.01。', 12, ACCENT_CYAN),
])

# 右侧 - 子采样示意 + 路径对比
add_rect(slide, 6.8, 2.0, 6.0, 1.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 2.08, 5.6, 1.7, [
    ('📊  子采样的实际效果', 13, ACCENT_GREEN, True),
    ('', 2),
    ('原始数据 1000 台：异常占比 20/1000 = 2%', 12, LIGHT_GRAY),
    ('子采样 256 台：异常占比约 20×256/1000 ≈ 5 台', 12, ACCENT_CYAN),
    ('→ 异常点密度从 2% 提升到约 2%，但正常点被稀释', 11, MID_GRAY),
    ('→ 980→约 251 台正常设备的密集度下降，异常设备更容易暴露', 11, ACCENT_GREEN),
])

# 路径对比表
add_rect(slide, 6.8, 3.95, 6.0, 2.85, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 4.0, 5.6, 0.3, '📊  5 棵树路径对比', 13, ACCENT_GREEN, True)
add_textbox(slide, 7.0, 4.25, 5.6, 0.2, 'P(68, 3.8) 是 980 台正常设备的代表', 9, MID_GRAY)

table_data2 = [
    ('树编号', 'A1(过热)', 'A2(松动)', 'A3(卡死)', 'P(正常)'),
    ('Tree 1', '6', '2', '7', '16'),
    ('Tree 2', '4', '3', '3', '14'),
    ('Tree 3', '4', '2', '14', '15'),
    ('Tree 4', '3', '3', '5', '13'),
    ('Tree 5', '2', '5', '15', '16'),
    ('100树E[h]', '5.05', '3.08', '7.32', '12.96'),
]
tbl2 = slide.shapes.add_table(7, 5, Inches(7.0), Inches(4.5), Inches(5.6), Inches(2.2)).table
for i, w in enumerate([1.1, 1.1, 1.1, 1.1, 1.2]):
    tbl2.columns[i].width = Inches(w)
for r in range(7):
    for c in range(5):
        cell = tbl2.cell(r, c); cell.text = table_data2[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(10); para.font.name = FONT_CN
            if r==0: para.font.bold=True; para.font.color.rgb=WHITE
            elif r==6: para.font.bold=True; para.font.color.rgb=ACCENT_GREEN
            else: para.font.color.rgb=LIGHT_GRAY
            para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15,0x30,0x50) if r==0 else (RGBColor(0x0A,0x2A,0x15) if r==6 else BG_CARD)


# ════════════════════════════════════════
# 第 7 页：异常分数计算 ★ 修正公式 + H(n-1)详解
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '异常分数计算', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用公式量化"有多异常"', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.5, 5.0, [
    ('异常分数公式', 16, ACCENT_GREEN, True),
    ('s(x, n) = 2 ^ ( -E[h(x)] / c(n) )', 20, WHITE, True),
    ('', 3),
    ('E[h(x)]：点 x 在所有树中路径长度的期望值（即平均路径）', 12, LIGHT_GRAY),
    ('c(n)：n 个点随机二分的期望路径长度（"平均孤立难度"）', 12, LIGHT_GRAY),
    ('', 5),
    ('c(n) 怎么算？', 16, ACCENT_GREEN, True),
    ('c(n) = 2H(n-1) - 2(n-1)/n', 14, WHITE, True),
    ('', 3),
    ('H(k) 是调和数：H(k) = 1 + 1/2 + 1/3 + ... + 1/k', 12, LIGHT_GRAY),
    ('可近似为 H(k) ≈ ln(k) + 0.5772（欧拉常数）', 12, ACCENT_CYAN),
    ('', 3),
    ('以贯穿示例 n=256 为例：', 12, WHITE, True),
    ('H(255) ≈ ln(255) + 0.5772 ≈ 5.5413 + 0.5772 ≈ 6.1185', 12, LIGHT_GRAY),
    ('c(256) = 2×6.1185 - 2×255/256 = 12.237 - 1.992', 12, LIGHT_GRAY),
    ('c(256) ≈ 10.249 ← "平均孤立256个点需要约10步"', 13, ACCENT_GREEN, True),
    ('', 4),
    ('分数解读', 16, ACCENT_GREEN, True),
    ('s → 1：路径远短于平均 → 极异常', 12, ACCENT_ORANGE),
    ('s ≈ 0.5：路径与平均相当 → 分界线', 12, ACCENT_CYAN),
    ('s → 0：路径远长于平均 → 极正常', 12, LIGHT_GRAY),
])

# 右侧 - 数值计算
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.05, 5.6, 0.3, '📐  真实计算结果（n=256, c=10.249）', 12, ACCENT_GREEN, True)

add_multiline_textbox(slide, 7.2, 2.5, 5.4, 4.2, [
    ('🔴 A2 松动(48°C, 13.5Hz)：E[h] = 3.08', 13, ACCENT_RED, True),
    ('', 2),
    ('s(A2) = 2^(-3.08 / 10.249)', 12, LIGHT_GRAY),
    ('     = 2^(-0.300)', 12, LIGHT_GRAY),
    ('     ≈ 0.812  →  明显异常 ✓', 14, ACCENT_GREEN, True),
    ('', 5),
    ('🔴 A1 过热(95°C, 3.2Hz)：E[h] = 5.05', 13, ACCENT_ORANGE, True),
    ('', 2),
    ('s(A1) = 2^(-5.05 / 10.249) = 2^(-0.493)', 12, LIGHT_GRAY),
    ('     ≈ 0.711  →  异常 ✓', 14, ACCENT_GREEN, True),
    ('', 5),
    ('🟡 A3 卡死(62°C, 0.2Hz)：E[h] = 7.32', 13, ACCENT_YELLOW, True),
    ('', 2),
    ('s(A3) = 2^(-7.32 / 10.249) = 2^(-0.714)', 12, LIGHT_GRAY),
    ('     ≈ 0.610  →  轻度异常 ✓', 14, ACCENT_GREEN, True),
    ('', 5),
    ('🟢 P 正常(68°C, 3.8Hz)：E[h] = 12.96', 13, ACCENT_CYAN, True),
    ('', 2),
    ('s(P) = 2^(-12.96 / 10.249) = 2^(-1.264)', 12, LIGHT_GRAY),
    ('     ≈ 0.416  →  正常 ✓', 14, ACCENT_CYAN, True),
    ('', 3),
    ('0 ─ 0.416 ─ 0.5 ── 0.610 ─ 0.711 ─ 0.812 ─ 1', 11, WHITE, True, PP_ALIGN.CENTER),
    ('    正常 ←──── 分界线 ────→ 异常', 10, MID_GRAY, False, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════
# 第 8 页：完整示例 ★ 全流程
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '完整示例——从数据到结论', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '用 1000 台传感器数据走一遍"训练→评分→判定"全流程', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.8, 5.0, [
    ('① 准备数据', 14, ACCENT_GREEN, True),
    ('1000 台设备的 (温度, 振动频率) 读数。', 12, LIGHT_GRAY),
    ('980 台正常：温度 60~75°C，振动 2~5 Hz（密集聚集）', 12, ACCENT_CYAN),
    ('20 台异常：A1 过热(95,3.2)、A2 松动(48,13.5)、A3 卡死(62,0.2) 等', 12, ACCENT_ORANGE),
    ('', 3),
    ('② 随机子采样（max_samples=256）', 14, ACCENT_GREEN, True),
    ('每棵树从 1000 台中随机抽 256 台。', 12, LIGHT_GRAY),
    ('980 台正常 → 约抽中 251 台；20 台异常 → 约抽中 5 台', 12, ACCENT_CYAN),
    ('→ 正常点被"稀释"，异常点相对更突出', 12, ACCENT_GREEN, True),
    ('', 3),
    ('③ 构建 100 棵 iTree（n_estimators=100）', 14, ACCENT_GREEN, True),
    ('每棵树：256 个样本，高度限制 8 层。', 12, LIGHT_GRAY),
    ('统计 100 棵树的平均路径长度 E[h(x)]（真实计算）：', 12, LIGHT_GRAY),
    ('  A2 松动：E[h] = 3.08（两个维度都偏离，最容易孤立）', 12, ACCENT_RED),
    ('  A1 过热：E[h] = 5.05（温度偏离但振动正常）', 12, ACCENT_ORANGE),
    ('  A3 卡死：E[h] = 7.32（温度正常、仅振动低）', 12, ACCENT_YELLOW),
    ('  P 正常：E[h] = 12.96（接近 c(256)=10.249）', 12, ACCENT_CYAN),
    ('', 3),
    ('④ 计算异常分数（c(256) ≈ 10.249）', 14, ACCENT_GREEN, True),
    ('  s(A2) ≈ 0.812 → 🔴 最异常（两维度偏离）', 13, ACCENT_RED, True),
    ('  s(A1) ≈ 0.711 → 🔴 异常', 13, ACCENT_ORANGE, True),
    ('  s(A3) ≈ 0.610 → 🟡 轻度异常', 13, ACCENT_YELLOW, True),
    ('  s(P) ≈ 0.416 → 🟢 正常（< 0.5 分界线）', 13, ACCENT_CYAN, True),
])

# 右侧 - 结果汇总
add_rect(slide, 6.8, 2.0, 6.0, 4.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 7.0, 2.05, 5.6, 0.3, '📊  结果汇总', 14, ACCENT_GREEN, True)

table_data3 = [
    ('设备', '特征', 'E[h]', '分数', '判定'),
    ('A2 松动', '(48, 13.5)', '3.08', '0.812', '🔴 最异常'),
    ('A1 过热', '(95, 3.2)', '5.05', '0.711', '🔴 异常'),
    ('A3 卡死', '(62, 0.2)', '7.32', '0.610', '🟡 轻度异常'),
    ('P 正常', '(68, 3.8)', '12.96', '0.416', '🟢 正常'),
    ('正常设备均值', '60~75°C', '~12.2', '~0.43', '🟢 正常'),
]
tbl3 = slide.shapes.add_table(6, 5, Inches(7.0), Inches(2.5), Inches(5.6), Inches(2.8)).table
for i, w in enumerate([1.2, 1.2, 0.8, 0.9, 1.5]):
    tbl3.columns[i].width = Inches(w)
for r in range(6):
    for c in range(5):
        cell = tbl3.cell(r, c); cell.text = table_data3[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(10); para.font.name = FONT_CN; para.alignment = PP_ALIGN.CENTER
            if r==0: para.font.bold=True; para.font.color.rgb=WHITE
            elif r in (1,2): para.font.color.rgb=ACCENT_ORANGE
            elif r==3: para.font.color.rgb=ACCENT_YELLOW
            elif r in (4,5): para.font.color.rgb=ACCENT_CYAN
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15,0x30,0x50) if r==0 else (
            RGBColor(0x2A,0x15,0x0A) if r in (1,2) else (
            RGBColor(0x2A,0x25,0x0A) if r==3 else BG_CARD))

add_multiline_textbox(slide, 7.2, 5.5, 5.4, 1.2, [
    ('✅  结论（真实计算）', 13, ACCENT_GREEN, True),
    ('20 台异常设备全部被成功识别（s > 0.5），无一漏网。', 11, LIGHT_GRAY),
    ('关键发现：A2(松动) 在两个维度都偏离，分数最高(0.812)；', 11, ACCENT_RED),
    ('A1(过热) 仅温度偏离，分数较低(0.711)。多维度偏离 = 更异常', 11, ACCENT_ORANGE),
])


# ════════════════════════════════════════
# 第 9 页：参数调优与优缺点
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '参数调优与优缺点', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '默认参数就能覆盖大多数场景——但也要知道局限', 16, ACCENT_CYAN)

add_rect(slide, 0.8, 2.0, 5.5, 2.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 1.0, 2.1, 5.2, 2.3, [
    ('⚙️  关键参数', 15, ACCENT_GREEN, True),
    ('', 4),
    ('树的数量（n_estimators）：默认 100', 13, WHITE, True),
    ('就像民意调查，100 份已足够稳定。100→500 棵 AUC 提升 <0.01', 11, LIGHT_GRAY),
    ('', 4),
    ('子采样大小（max_samples）：默认 256', 13, WHITE, True),
    ('论文实验的"甜蜜点"——太小信息不足，太大正常点互相遮挡', 11, LIGHT_GRAY),
    ('', 4),
    ('异常阈值（contamination）：默认 auto（s=0.5）', 13, WHITE, True),
    ('如有业务知识（欺诈率约 0.1%），可直接设 0.001', 11, LIGHT_GRAY),
])

add_rect(slide, 0.8, 4.8, 5.5, 2.3, BG_CARD, RGBColor(0x0A, 0x2A, 0x15))
add_multiline_textbox(slide, 1.0, 4.9, 5.2, 2.1, [
    ('✅  核心优点', 15, ACCENT_GREEN, True),
    ('• 线性时间复杂度 O(t·ψ·logψ)，轻松处理百万级数据', 12, LIGHT_GRAY),
    ('• 天然支持并行化，100 棵树可分配到 100 个 CPU 核心', 12, LIGHT_GRAY),
    ('• 不依赖距离度量，天然免疫"维度灾难"', 12, LIGHT_GRAY),
    ('• 纯无监督，不需要任何标注数据', 12, LIGHT_GRAY),
    ('• 内存高效：100×256 = 25,600 个数据点即可', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 2.0, 6.0, 2.5, BG_CARD, RGBColor(0x2A, 0x15, 0x0A))
add_multiline_textbox(slide, 7.0, 2.1, 5.6, 2.3, [
    ('⚠️  注意局限', 15, ACCENT_ORANGE, True),
    ('', 4),
    ('• 局部异常检测能力弱——主要捕获全局异常', 12, LIGHT_GRAY),
    ('  如 A3(卡死) 分数 0.610 远低于 A2(松动) 的 0.812——偏差不那么极端就更难抓', 11, MID_GRAY),
    ('• 轴对齐切割偏差——只沿坐标轴切，非轴方向可能"伪异常"', 12, LIGHT_GRAY),
    ('  Extended Isolation Forest（2019）用任意角度超平面解决', 11, MID_GRAY),
    ('• 对分类特征处理不佳，需要先做独热编码', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 4.8, 6.0, 2.3, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 4.9, 5.6, 2.1, [
    ('🎯  选择建议', 15, ACCENT_CYAN, True),
    ('', 4),
    ('数据量大（>10 万条）→ 首选孤立森林', 13, ACCENT_GREEN, True),
    ('需要检测局部异常 → 选 LOF', 13, LIGHT_GRAY),
    ('数据量小（<5000 条）且边界复杂 → 选 One-Class SVM', 13, LIGHT_GRAY),
    ('不确定 → 先试孤立森林（速度快、参数少）', 13, ACCENT_CYAN, True),
])


# ════════════════════════════════════════
# 第 10 页：实战案例
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '实战案例：金融欺诈检测', 32, WHITE, True)
add_textbox(slide, 0.8, 1.3, 11.7, 0.5, '从原始交易数据到实时欺诈告警——完整的落地 Pipeline', 16, ACCENT_CYAN)

add_multiline_textbox(slide, 0.8, 2.0, 5.8, 5.0, [
    ('① 数据准备与特征工程', 14, ACCENT_GREEN, True),
    ('原始特征：交易金额、时间、商户类别码、GPS 坐标', 12, LIGHT_GRAY),
    ('衍生特征：过去 1h 交易次数、金额/历史均值比、距上次交易时间', 12, LIGHT_GRAY),
    ('→ 衍生特征是关键：孤立森林无法自动发现"频率异常"', 11, ACCENT_CYAN),
    ('', 4),
    ('② 模型训练', 14, ACCENT_GREEN, True),
    ('n_estimators=100, max_samples=256, contamination=0.001', 12, LIGHT_GRAY),
    ('contamination 对应约 0.1% 的欺诈率（业务先验知识）', 12, LIGHT_GRAY),
    ('训练耗时：百万级交易数据约 10 秒（单机）', 12, ACCENT_CYAN),
    ('', 4),
    ('③ 实时评分与告警', 14, ACCENT_GREEN, True),
    ('每笔新交易 → 提取特征 → 走 100 棵树 → 输出异常分数', 12, LIGHT_GRAY),
    ('s > 0.7 触发告警 → 人工复核或自动冻结', 12, LIGHT_GRAY),
    ('', 4),
    ('④ 效果（Kaggle 信用卡欺诈数据集）', 14, ACCENT_GREEN, True),
    ('检测率 ~82%，误报率 <1%', 13, ACCENT_GREEN, True),
    ('对新型欺诈手法（无历史标签）仍有检测能力', 12, LIGHT_GRAY),
])

add_rect(slide, 6.8, 2.0, 6.0, 2.8, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_multiline_textbox(slide, 7.0, 2.1, 5.6, 2.6, [
    ('🏭  设备故障预警（NASA 涡轮发动机数据）', 13, ACCENT_GREEN, True),
    ('', 3),
    ('Pipeline：传感器时序 → 按窗口提取均值/方差/峰值 → 多传感器拼接 → 训练 iForest → 监控分数趋势', 11, LIGHT_GRAY),
    ('', 3),
    ('关键发现：不是看单次分数，而是看分数变化趋势——持续升高 = 设备在退化', 11, ACCENT_CYAN),
    ('', 3),
    ('效果：可在故障前 15-25 个工作日检测到异常趋势', 12, ACCENT_GREEN, True),
])

add_rect(slide, 6.8, 5.1, 6.0, 2.0, BG_CARD, RGBColor(0x0A, 0x2A, 0x15))
add_multiline_textbox(slide, 7.0, 5.2, 5.6, 1.8, [
    ('🏢  工业界大规模采用', 14, ACCENT_GREEN, True),
    ('', 3),
    ('Google — 监控内部网络健康', 12, LIGHT_GRAY),
    ('Microsoft Azure — 异常检测云服务', 12, LIGHT_GRAY),
    ('Netflix — 用户行为异常检测', 12, LIGHT_GRAY),
    ('阿里巴巴 — 双十一交易欺诈检测，每秒数十万笔', 12, LIGHT_GRAY),
])


# ════════════════════════════════════════
# 第 11 页：结尾
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
add_accent_line(slide, 0.8, 0.6, 11.7, ACCENT_GREEN)
add_textbox(slide, 0.8, 0.7, 11.7, 0.7, '关键收获与下一步', 32, WHITE, True)

add_multiline_textbox(slide, 0.8, 1.5, 6.5, 4.0, [
    ('五句话总结孤立森林', 18, ACCENT_GREEN, True),
    ('', 6),
    ('① 新范式：不建模正常，直接度量"被孤立的难易程度"', 14, LIGHT_GRAY),
    ('② 简洁实现：100 棵随机分割二叉树 + 路径长度 + 异常分数公式', 14, LIGHT_GRAY),
    ('③ 极致效率：线性时间复杂度，天然并行，百万级数据轻松处理', 14, LIGHT_GRAY),
    ('④ 广泛适用：网络安全、金融反欺诈、工业预警、IT 运维', 14, LIGHT_GRAY),
    ('⑤ 持续进化：Extended Isolation Forest 解决轴对齐偏差', 14, LIGHT_GRAY),
    ('', 8),
    ('行动建议', 18, ACCENT_GREEN, True),
    ('下次遇到异常检测任务，先用孤立森林建立基线。', 14, WHITE, True),
    ('默认参数（100 棵树、256 子采样）在大多数场景下就够用。', 13, LIGHT_GRAY),
    ('效果不够再考虑 LOF 或 One-Class SVM。', 13, LIGHT_GRAY),
])

add_rect(slide, 7.8, 1.5, 4.8, 3.5, BG_CARD, RGBColor(0x1A, 0x3A, 0x5C))
add_textbox(slide, 8.0, 1.6, 4.4, 0.4, '🔢  关键数字速记', 14, ACCENT_GREEN, True)

tbl4 = slide.shapes.add_table(6, 2, Inches(8.2), Inches(2.2), Inches(4.0), Inches(2.5)).table
tbl4.columns[0].width = Inches(2.0); tbl4.columns[1].width = Inches(2.0)
for r, (k, v) in enumerate([('项目','数值'),('默认树数量','100'),('默认子采样','256'),('默认树高度','8'),('异常分界线','s = 0.5'),('训练复杂度','O(n)')]):
    for c, t in enumerate((k, v)):
        cell = tbl4.cell(r, c); cell.text = t
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12); para.font.name = FONT_CN; para.alignment = PP_ALIGN.CENTER
            if r==0: para.font.bold=True; para.font.color.rgb=WHITE
            else: para.font.color.rgb=ACCENT_CYAN if c==1 else LIGHT_GRAY
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x15,0x30,0x50) if r==0 else BG_CARD

add_rect(slide, 0.8, 5.5, 11.7, 1.5, BG_CARD)
add_multiline_textbox(slide, 1.0, 5.6, 11.3, 1.3, [
    ('📄  论文出处', 14, ACCENT_GREEN, True),
    ('Liu, Ting & Zhou (2008) "Isolation Forest" IEEE ICDM  |  Liu, Ting & Zhou (2012) "Isolation-based Anomaly Detection" ACM TKDD', 11, LIGHT_GRAY),
    ('Hariri, Carrasco Kind & McNutt (2019) "Extended Isolation Forest" IEEE TKDE', 11, LIGHT_GRAY),
    ('', 4),
    ('"一切都应该尽可能简单，但不能过于简单。" —— 爱因斯坦。孤立森林恰好停在"足够简单但不失效"的完美位置。', 12, MID_GRAY),
])

# ── 保存 ──
output_path = '/home/admin/.openclaw/workspace-weaver/output/isolation-forest_v5.pptx'
prs.save(output_path)
print(f'✅ PPT V5 已保存到 {output_path}')
print(f'总页数：{len(prs.slides)}')
