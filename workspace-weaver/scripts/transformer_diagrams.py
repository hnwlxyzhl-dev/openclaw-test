#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Transformer 架构图 PPT 生成脚本
使用 python-pptx 创建 5 页架构图/流程图
- 页面1: Self-Attention Q/K/V 概念图
- 页面2: Self-Attention 5步计算流程图
- 页面3: Multi-Head Attention 架构图
- 页面4: Transformer 整体架构图 (Encoder-Decoder)
- 页面5: 训练 vs 推理 对比图
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy

# ============================================================
# 全局配置
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MAX_RIGHT = 13.2   # 英寸，所有形状右边界限制
MAX_BOTTOM = 7.3   # 英寸，所有形状底边界限制

# 颜色方案
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)    # 蓝色系 - 编码器
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)    # 红色系 - 解码器
C_ATTENTION = RGBColor(0x34, 0x98, 0xDB)  # 蓝色 - 注意力层
C_FFN = RGBColor(0x27, 0xAE, 0x60)        # 绿色 - FFN层
C_NORM = RGBColor(0xF3, 0x9C, 0x12)       # 橙色 - 归一化层
C_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # 白色文字
C_TEXT_DARK = RGBColor(0x33, 0x33, 0x33)   # 深灰文字
C_ARROW = RGBColor(0x55, 0x55, 0x55)       # 箭头颜色
C_BG = RGBColor(0xFF, 0xFF, 0xFF)          # 白色背景
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8) # 浅蓝背景
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)  # 浅红背景
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3) # 浅绿背景
C_LIGHT_ORANGE = RGBColor(0xFA, 0xE5, 0xD3) # 浅橙背景
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)      # 紫色
C_TEAL = RGBColor(0x1A, 0xBC, 0x9C)       # 青色
C_DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)  # 深海蓝
C_GRAY = RGBColor(0xBD, 0xC3, 0xC7)       # 灰色

# 多头注意力 8 个头的颜色
HEAD_COLORS = [
    RGBColor(0xE7, 0x4C, 0x3C),  # 红
    RGBColor(0xE6, 0x7E, 0x22),  # 橙
    RGBColor(0xF1, 0xC4, 0x0F),  # 黄
    RGBColor(0x27, 0xAE, 0x60),  # 绿
    RGBColor(0x34, 0x98, 0xDB),  # 蓝
    RGBColor(0x8E, 0x44, 0xAD),  # 紫
    RGBColor(0x1A, 0xBC, 0x9C),  # 青
    RGBColor(0x2C, 0x3E, 0x50),  # 深灰蓝
]

# 字体配置
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

# ============================================================
# 工具函数
# ============================================================

def check_bounds(left, top, width, height, name=""):
    """验证形状是否在幻灯片边界内"""
    r = left + width
    b = top + height
    ok = True
    if r > MAX_RIGHT:
        print(f"  [WARN] {name}: right={r:.2f} > {MAX_RIGHT}")
        ok = False
    if b > MAX_BOTTOM:
        print(f"  [WARN] {name}: bottom={b:.2f} > {MAX_BOTTOM}")
        ok = False
    return ok


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None,
                     text="", font_size=12, font_color=C_TEXT_WHITE, bold=False,
                     name="", alignment=PP_ALIGN.CENTER):
    """添加圆角矩形，带文字"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    width_i = width if isinstance(width, float) else width
    height_i = height if isinstance(height, float) else height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_i), Inches(top_i),
        Inches(width_i), Inches(height_i)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    # 调整圆角弧度 (0.0 ~ 1.0)
    # 默认圆角较大，我们用较小弧度
    shape.adjustments[0] = 0.1

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = FONT_EN

    check_bounds(left_i, top_i, width_i, height_i, name or text[:20])
    return shape


def add_rect(slide, left, top, width, height, fill_color, line_color=None,
             text="", font_size=12, font_color=C_TEXT_WHITE, bold=False,
             name="", alignment=PP_ALIGN.CENTER):
    """添加矩形，带文字"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    width_i = width if isinstance(width, float) else width
    height_i = height if isinstance(height, float) else height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_i), Inches(top_i),
        Inches(width_i), Inches(height_i)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = FONT_EN

    check_bounds(left_i, top_i, width_i, height_i, name or text[:20])
    return shape


def add_arrow(slide, left, top, width, height, color=C_ARROW, name=""):
    """添加右箭头"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    width_i = width if isinstance(width, float) else width
    height_i = height if isinstance(height, float) else height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left_i), Inches(top_i),
        Inches(width_i), Inches(height_i)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    check_bounds(left_i, top_i, width_i, height_i, name or "arrow")
    return shape


def add_down_arrow(slide, left, top, width, height, color=C_ARROW, name=""):
    """添加下箭头"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    width_i = width if isinstance(width, float) else width
    height_i = height if isinstance(height, float) else height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left_i), Inches(top_i),
        Inches(width_i), Inches(height_i)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    check_bounds(left_i, top_i, width_i, height_i, name or "down-arrow")
    return shape


def add_connector(slide, start_left, start_top, end_left, end_top, color=C_ARROW, width_pt=1.5):
    """添加直线连接器"""
    from pptx.oxml.ns import qn
    # 使用简单的直线形状代替连接器
    cx = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top)
    )
    cx.line.color.rgb = color
    cx.line.width = Pt(width_pt)
    return cx


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 font_color=C_TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 name=""):
    """添加文本框"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    width_i = width if isinstance(width, float) else width
    height_i = height if isinstance(height, float) else height

    txBox = slide.shapes.add_textbox(
        Inches(left_i), Inches(top_i),
        Inches(width_i), Inches(height_i)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = FONT_CN

    check_bounds(left_i, top_i, width_i, height_i, name or text[:20])
    return txBox


def add_circle(slide, left, top, diameter, fill_color, text="", font_size=11,
               font_color=C_TEXT_WHITE, name=""):
    """添加圆形"""
    left_i = left if isinstance(left, float) else left
    top_i = top if isinstance(top, float) else top
    diameter_i = diameter if isinstance(diameter, float) else diameter

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left_i), Inches(top_i),
        Inches(diameter_i), Inches(diameter_i)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = True
        run.font.name = FONT_EN

    check_bounds(left_i, top_i, diameter_i, diameter_i, name or text[:20])
    return shape


def set_shape_text(shape, text, font_size=12, font_color=C_TEXT_WHITE,
                   bold=False, alignment=PP_ALIGN.CENTER, font_name=None):
    """设置已有形状的文字"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name or FONT_CN


def add_page_title(slide, text):
    """添加页面标题"""
    return add_text_box(slide, 0.4, 0.2, 12.5, 0.6, text,
                        font_size=22, font_color=C_DARK_BLUE, bold=True,
                        name="page-title")


def add_background_rect(slide, left, top, width, height, fill_color, name="bg"):
    """添加背景矩形"""
    return add_rect(slide, left, top, width, height, fill_color, name=name)


# ============================================================
# 页面1: Self-Attention Q/K/V 概念图
# ============================================================
def create_page1(prs):
    """展示 Self-Attention 中 Q/K/V 的概念，用句子连线展示注意力"""
    print("\n=== Page 1: Self-Attention Q/K/V 概念图 ===")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_page_title(slide, "Self-Attention: Q / K / V 概念图")

    # --- 左侧: 句子注意力可视化 ---
    # 标题
    add_text_box(slide, 0.4, 0.85, 6.0, 0.4,
                 "句子注意力可视化", font_size=14, font_color=C_DARK_BLUE, bold=True)

    # 句子: "The cat sat on the mat, and it was happy"
    words = ["The", "cat", "sat", "on", "the", "mat,", "and", "it", "was", "happy"]
    word_colors = [
        C_GRAY, C_ENCODER, C_GRAY, C_GRAY, C_GRAY, C_GRAY, C_GRAY,
        C_DECODER, C_GRAY, C_GRAY
    ]

    # 词汇方块排列 - 顶部一行
    word_w = 0.85
    word_h = 0.45
    start_x = 0.5
    start_y = 1.4
    gap = 0.15

    word_shapes = []
    for i, (w, c) in enumerate(zip(words, word_colors)):
        x = start_x + i * (word_w + gap)
        s = add_rounded_rect(slide, x, start_y, word_w, word_h, c,
                             text=w, font_size=13, font_color=C_TEXT_WHITE, bold=True,
                             name=f"word-{i}-{w}")
        word_shapes.append(s)

    # 注意力权重说明 - "it" 关注 "cat"
    # 在词汇下方画连线区域
    # 标注 "it" 的注意力分布
    add_text_box(slide, 0.5, 2.05, 8.0, 0.35,
                 '"it" 的注意力权重分布 (越粗 = 越关注):',
                 font_size=11, font_color=C_TEXT_DARK, bold=True)

    # 画注意力权重条形图 (水平)
    # "it" 对每个词的注意力分数
    attn_scores = [0.05, 0.45, 0.08, 0.03, 0.02, 0.15, 0.03, 0.10, 0.04, 0.05]
    bar_start_y = 2.5
    bar_h = 0.35
    max_bar_w = 1.2

    for i, (w, score) in enumerate(zip(words, attn_scores)):
        x = start_x + i * (word_w + gap)
        # 分数标签
        add_text_box(slide, x, bar_start_y, word_w, 0.25,
                     f"{score:.0%}", font_size=9, font_color=C_TEXT_DARK,
                     alignment=PP_ALIGN.CENTER)
        # 权重条
        bar_w = score * max_bar_w / 0.45  # 归一化到最大值
        if score > 0.01:
            bar_color = C_DECODER if w == "it" else (C_ENCODER if score >= 0.15 else C_GRAY)
            add_rect(slide, x + 0.1, bar_start_y + 0.25, bar_w, bar_h, bar_color,
                     name=f"bar-{i}")

    # 高亮标注: cat 得到最高注意力
    add_text_box(slide, 0.5, 3.3, 8.0, 0.3,
                 '>>> "it" 最关注 "cat" (45%)，因为它指代的就是 "cat"！',
                 font_size=12, font_color=C_DECODER, bold=True)

    # --- 右侧: Q/K/V 比喻说明 ---
    # 图书馆比喻区域
    add_text_box(slide, 7.5, 0.85, 5.5, 0.4,
                 "Q / K / V 比喻: 图书馆搜索", font_size=14, font_color=C_DARK_BLUE, bold=True)

    # 三个大卡片: Q, K, V
    card_w = 1.8
    card_h = 2.2
    card_y = 1.4

    # Q 卡片 - 搜索关键词
    q_x = 7.8
    add_rounded_rect(slide, q_x, card_y, card_w, card_h, C_ATTENTION,
                     name="Q-card")
    add_text_box(slide, q_x, card_y + 0.1, card_w, 0.35,
                 "Q (Query)", font_size=14, font_color=C_TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, q_x + 0.1, card_y + 0.5, card_w - 0.2, 1.6,
                 "搜索关键词\n\n'我在找什么？'\n\n如: '量子力学'",
                 font_size=11, font_color=C_TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    # K 卡片 - 书的标签
    k_x = 9.8
    add_rounded_rect(slide, k_x, card_y, card_w, card_h, C_FFN,
                     name="K-card")
    add_text_box(slide, k_x, card_y + 0.1, card_w, 0.35,
                 "K (Key)", font_size=14, font_color=C_TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, k_x + 0.1, card_y + 0.5, card_w - 0.2, 1.6,
                 "书的标签\n\n'我有什么特征？'\n\n如: '物理>量子'",
                 font_size=11, font_color=C_TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    # V 卡片 - 书的内容 (缩小宽度并左移)
    v_x = 11.3
    card_w_v = 1.6
    add_rounded_rect(slide, v_x, card_y, card_w_v, card_h, C_NORM,
                     name="V-card")
    add_text_box(slide, v_x, card_y + 0.1, card_w_v, 0.35,
                 "V (Value)", font_size=14, font_color=C_TEXT_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
    add_text_box(slide, v_x + 0.1, card_y + 0.5, card_w_v - 0.2, 1.6,
                 "书的实际内容\n\n'我的实际内容'\n\n如: 波函数方程...",
                 font_size=11, font_color=C_TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    # 连接箭头 Q -> K -> V
    add_arrow(slide, q_x + card_w + 0.05, card_y + 0.9, 0.12, 0.25,
              C_ARROW, name="q-k-arrow")
    add_arrow(slide, k_x + card_w + 0.05, card_y + 0.9, 0.12, 0.25,
              C_ARROW, name="k-v-arrow")

    # --- 底部: 总结公式 ---
    add_text_box(slide, 0.4, 3.8, 12.5, 0.35,
                 "生成过程: 输入 X 通过三个不同的权重矩阵变换得到 Q, K, V",
                 font_size=12, font_color=C_TEXT_DARK, bold=True)

    # 公式区域
    formula_y = 4.3
    add_rect(slide, 0.8, formula_y, 3.2, 0.6, C_ATTENTION,
             text="Q = X * W_Q", font_size=14, font_color=C_TEXT_WHITE, bold=True,
             name="formula-Q")
    add_rect(slide, 4.5, formula_y, 3.2, 0.6, C_FFN,
             text="K = X * W_K", font_size=14, font_color=C_TEXT_WHITE, bold=True,
             name="formula-K")
    add_rect(slide, 8.2, formula_y, 3.2, 0.6, C_NORM,
             text="V = X * W_V", font_size=14, font_color=C_TEXT_WHITE, bold=True,
             name="formula-V")

    # 箭头: 输入 -> 三个公式
    add_text_box(slide, 0.2, formula_y + 0.1, 0.5, 0.4,
                 "X", font_size=16, font_color=C_TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_arrow(slide, 0.6, formula_y + 0.2, 0.15, 0.15, C_ARROW, "x-arrow1")
    add_arrow(slide, 0.6, formula_y + 0.2, 0.15, 0.15, C_ARROW, "x-arrow2")

    # --- 底部: 直觉总结 ---
    summary_y = 5.3
    add_rect(slide, 0.4, summary_y, 12.5, 1.8, C_LIGHT_BLUE,
             name="summary-box")
    add_text_box(slide, 0.6, summary_y + 0.1, 12.0, 0.35,
                 "直觉理解:", font_size=14, font_color=C_DARK_BLUE, bold=True)
    add_text_box(slide, 0.6, summary_y + 0.5, 5.8, 1.1,
                 "1. Q (Query) = 每个词问: '我需要什么信息？'\n"
                 "2. K (Key) = 每个词说: '我有什么信息可以提供？'\n"
                 "3. V (Value) = 每个词的: '实际内容是什么？'",
                 font_size=12, font_color=C_TEXT_DARK)
    add_text_box(slide, 6.8, summary_y + 0.5, 5.8, 1.1,
                 "Q 和 K 匹配度高的词对，注意力权重就大。\n\n"
                 "就像 'it' 用 Q 去 '搜索'，'cat' 的 K 匹配度\n"
                 "最高，所以 'it' 从 'cat' 的 V 中获取了信息。",
                 font_size=12, font_color=C_TEXT_DARK)

    print("  Page 1 完成")


# ============================================================
# 页面2: Self-Attention 5步计算流程图
# ============================================================
def create_page2(prs):
    """5步计算流程图: 输入 -> Q,K,V -> QK^T -> /sqrt(dk) -> softmax -> *V"""
    print("\n=== Page 2: Self-Attention 5步计算流程图 ===")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_page_title(slide, "Self-Attention: 5步计算流程")

    # 流程步骤定义
    steps = [
        {
            "title": "Step 1",
            "main": "生成 Q, K, V",
            "formula": "Q=XW^Q, K=XW^K, V=XW^V",
            "desc": "输入通过3个权重矩阵\n变换得到Q/K/V",
            "color": C_ATTENTION,
            "bg_color": C_LIGHT_BLUE,
        },
        {
            "title": "Step 2",
            "main": "Q x K^T",
            "formula": "Scores = Q * K^T",
            "desc": "计算每对词之间的\n点积相似度",
            "color": C_FFN,
            "bg_color": C_LIGHT_GREEN,
        },
        {
            "title": "Step 3",
            "main": "/ sqrt(d_k)",
            "formula": "Scores / sqrt(d_k)",
            "desc": "缩放防止分数过大\n(温度计比喻)",
            "color": C_NORM,
            "bg_color": C_LIGHT_ORANGE,
        },
        {
            "title": "Step 4",
            "main": "Softmax",
            "formula": "Weights = softmax(...)",
            "desc": "转为概率分布\n每行之和 = 1",
            "color": C_PURPLE,
            "bg_color": RGBColor(0xE8, 0xDA, 0xEF),
        },
        {
            "title": "Step 5",
            "main": "x V",
            "formula": "Output = Weights * V",
            "desc": "加权求和\n注意力高的词贡献大",
            "color": C_TEAL,
            "bg_color": RGBColor(0xD1, 0xF2, 0xEB),
        },
    ]

    # 整体布局: 5个步骤横排
    step_w = 2.2
    step_h = 2.8
    arrow_w = 0.3
    gap = 0.15
    total_w = 5 * step_w + 4 * arrow_w + 4 * gap
    start_x = (13.333 - total_w) / 2
    step_y = 1.0

    for i, step in enumerate(steps):
        x = start_x + i * (step_w + arrow_w + gap)

        # 步骤编号背景
        add_rect(slide, x, step_y, step_w, 0.4, step["color"],
                 text=step["title"], font_size=13, font_color=C_TEXT_WHITE, bold=True,
                 name=f"step{i}-header")

        # 步骤主体背景
        add_rounded_rect(slide, x, step_y + 0.4, step_w, step_h - 0.4,
                         step["bg_color"], line_color=step["color"],
                         name=f"step{i}-bg")

        # 步骤主标题
        add_text_box(slide, x + 0.1, step_y + 0.55, step_w - 0.2, 0.5,
                     step["main"], font_size=16, font_color=step["color"], bold=True,
                     alignment=PP_ALIGN.CENTER, name=f"step{i}-main")

        # 公式
        add_text_box(slide, x + 0.1, step_y + 1.1, step_w - 0.2, 0.4,
                     step["formula"], font_size=11, font_color=C_TEXT_DARK,
                     alignment=PP_ALIGN.CENTER, name=f"step{i}-formula")

        # 说明文字
        add_text_box(slide, x + 0.15, step_y + 1.6, step_w - 0.3, 1.0,
                     step["desc"], font_size=11, font_color=C_TEXT_DARK,
                     alignment=PP_ALIGN.CENTER, name=f"step{i}-desc")

        # 步骤间箭头
        if i < 4:
            arrow_x = x + step_w + gap * 0.3
            add_arrow(slide, arrow_x, step_y + 1.2, arrow_w, 0.25,
                      C_ARROW, name=f"step-arrow-{i}")

    # --- 底部: 注意力分数变化过程可视化 ---
    viz_y = 4.1
    add_text_box(slide, 0.4, viz_y, 12.5, 0.35,
                 "数值变化过程 (以 'it' 和 'cat' 的关系为例):",
                 font_size=13, font_color=C_DARK_BLUE, bold=True)

    # 5个阶段的小矩阵/图形展示
    viz_h = 2.8
    viz_w = 2.2

    # 阶段1: 输入向量
    x1 = start_x
    add_rect(slide, x1, viz_y + 0.4, viz_w, viz_h, C_LIGHT_BLUE,
             line_color=C_ATTENTION, name="viz1")
    add_text_box(slide, x1 + 0.1, viz_y + 0.5, viz_w - 0.2, 0.3,
                 "输入向量", font_size=11, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x1 + 0.1, viz_y + 0.9, viz_w - 0.2, 2.0,
                 'X = [0.3, -0.1, 0.8, ...]\n   (512维嵌入)\n\n'
                 '输入: "it" 的词嵌入',
                 font_size=10, font_color=C_TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # 阶段2: 注意力分数矩阵
    x2 = start_x + (step_w + arrow_w + gap)
    add_rect(slide, x2, viz_y + 0.4, viz_w, viz_h, C_LIGHT_GREEN,
             line_color=C_FFN, name="viz2")
    add_text_box(slide, x2 + 0.1, viz_y + 0.5, viz_w - 0.2, 0.3,
                 "注意力分数", font_size=11, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x2 + 0.1, viz_y + 0.9, viz_w - 0.2, 2.0,
                 'Q("it") * K^T("cat")\n= [高!]\n\n'
                 '"it"与"cat"匹配度最高\n分数 = 3.2',
                 font_size=10, font_color=C_TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # 阶段3: 缩放后
    x3 = start_x + 2 * (step_w + arrow_w + gap)
    add_rect(slide, x3, viz_y + 0.4, viz_w, viz_h, C_LIGHT_ORANGE,
             line_color=C_NORM, name="viz3")
    add_text_box(slide, x3 + 0.1, viz_y + 0.5, viz_w - 0.2, 0.3,
                 "缩放后", font_size=11, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x3 + 0.1, viz_y + 0.9, viz_w - 0.2, 2.0,
                 '3.2 / sqrt(64)\n= 3.2 / 8\n= 0.4\n\n'
                 '防止Softmax饱和',
                 font_size=10, font_color=C_TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # 阶段4: Softmax概率
    x4 = start_x + 3 * (step_w + arrow_w + gap)
    add_rect(slide, x4, viz_y + 0.4, viz_w, viz_h, RGBColor(0xE8, 0xDA, 0xEF),
             line_color=C_PURPLE, name="viz4")
    add_text_box(slide, x4 + 0.1, viz_y + 0.5, viz_w - 0.2, 0.3,
                 "Softmax概率", font_size=11, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x4 + 0.1, viz_y + 0.9, viz_w - 0.2, 2.0,
                 'softmax([0.4, 0.1, ...])\n= [0.45, 0.08, ...]\n\n'
                 '"cat"权重 = 45%',
                 font_size=10, font_color=C_TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # 阶段5: 最终输出
    x5 = start_x + 4 * (step_w + arrow_w + gap)
    add_rect(slide, x5, viz_y + 0.4, viz_w, viz_h, RGBColor(0xD1, 0xF2, 0xEB),
             line_color=C_TEAL, name="viz5")
    add_text_box(slide, x5 + 0.1, viz_y + 0.5, viz_w - 0.2, 0.3,
                 "最终输出", font_size=11, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x5 + 0.1, viz_y + 0.9, viz_w - 0.2, 2.0,
                 '0.45*V("cat")\n+ 0.08*V("sat")\n+ ...\n\n'
                 '"it" 融合了"cat"的信息',
                 font_size=10, font_color=C_TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # 底部公式
    add_text_box(slide, 0.4, 7.0, 12.5, 0.3,
                 "Attention(Q,K,V) = softmax(QK^T / sqrt(d_k)) * V",
                 font_size=14, font_color=C_DARK_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    print("  Page 2 完成")


# ============================================================
# 页面3: Multi-Head Attention 架构图
# ============================================================
def create_page3(prs):
    """展示 Multi-Head Attention: 输入 -> 拆分8份 -> 8个注意力头 -> Concat -> 线性投影"""
    print("\n=== Page 3: Multi-Head Attention 架构图 ===")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_page_title(slide, "Multi-Head Attention: 多头注意力架构")

    # --- 主流程图 ---
    # 整体布局: 输入 -> 投影 -> 8头 -> Concat -> 线性 -> 输出
    # 严格控制总宽度 <= 13.2
    flow_y = 1.85   # 所有流程块共用Y坐标
    block_h = 1.3   # 流程块高度

    # 输入块
    input_x = 0.3
    input_w = 0.9
    add_rounded_rect(slide, input_x, flow_y, input_w, block_h, C_DARK_BLUE,
                     text="输入\nInput", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="mha-input")

    # 箭头: 输入 -> 投影
    add_arrow(slide, 1.25, flow_y + 0.45, 0.25, 0.18, C_ARROW, "input-split-arrow")

    # 投影块
    proj_x = 1.55
    proj_w = 0.9
    add_rounded_rect(slide, proj_x, flow_y, proj_w, block_h, C_NORM,
                     text="线性投影\n(8份)", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="mha-split")

    # 箭头: 投影 -> 头
    add_arrow(slide, 2.50, flow_y + 0.45, 0.2, 0.15, C_ARROW, "split-heads-arrow")

    # 8个注意力头
    head_x = 2.75
    head_w = 0.58
    head_gap = 0.05
    heads_end = head_x + 8 * (head_w + head_gap) - head_gap  # 8个头的右边界

    for i in range(8):
        x = head_x + i * (head_w + head_gap)
        color = HEAD_COLORS[i]
        add_rounded_rect(slide, x, flow_y, head_w, block_h, color,
                         text=f"Head\n{i+1}", font_size=10, font_color=C_TEXT_WHITE, bold=True,
                         name=f"head-{i}")

    # 箭头: 头 -> Concat
    add_arrow(slide, heads_end + 0.02, flow_y + 0.45, 0.2, 0.15, C_ARROW, "heads-concat-arrow")

    # Concat 块
    concat_x = heads_end + 0.27
    concat_w = 0.9
    add_rounded_rect(slide, concat_x, flow_y, concat_w, block_h, C_PURPLE,
                     text="Concat\n拼接", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="mha-concat")

    # 箭头: Concat -> 线性投影
    add_arrow(slide, concat_x + concat_w + 0.03, flow_y + 0.45, 0.2, 0.15, C_ARROW, "concat-linear-arrow")

    # 线性投影块
    linear_x = concat_x + concat_w + 0.28
    linear_w = 0.9
    add_rounded_rect(slide, linear_x, flow_y, linear_w, block_h, C_TEAL,
                     text="Linear\n线性投影", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="mha-linear")

    # 箭头: 线性投影 -> 输出
    add_arrow(slide, linear_x + linear_w + 0.03, flow_y + 0.45, 0.2, 0.15, C_ARROW, "linear-output-arrow")

    # 输出块
    output_x = linear_x + linear_w + 0.28
    output_w = 0.9
    add_rounded_rect(slide, output_x, flow_y, output_w, block_h, C_DARK_BLUE,
                     text="输出\nOutput", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="mha-output")

    # --- "8个不同的视角" 标注 ---
    add_text_box(slide, head_x, flow_y - 0.45, 8 * (head_w + head_gap), 0.35,
                 "8 个注意力头 = 8 个不同的 '视角'",
                 font_size=13, font_color=C_DECODER, bold=True,
                 alignment=PP_ALIGN.CENTER, name="8-perspectives")

    # --- 底部: 每个头的视角说明 ---
    desc_y = 4.0
    add_text_box(slide, 0.4, desc_y, 12.5, 0.35,
                 "不同注意力头可以学习不同的关系模式:",
                 font_size=13, font_color=C_DARK_BLUE, bold=True)

    perspectives = [
        ("Head 1", "语法关系", "主谓搭配", HEAD_COLORS[0]),
        ("Head 2", "共指关系", "代词指代", HEAD_COLORS[1]),
        ("Head 3", "语义相似", "同义词", HEAD_COLORS[2]),
        ("Head 4", "位置关系", "相邻词", HEAD_COLORS[3]),
        ("Head 5", "情感关联", "褒贬词", HEAD_COLORS[4]),
        ("Head 6", "逻辑关系", "因果/转折", HEAD_COLORS[5]),
        ("Head 7", "主题关联", "同一主题", HEAD_COLORS[6]),
        ("Head 8", "未知模式", "自动学习", HEAD_COLORS[7]),
    ]

    card_w = 1.45
    card_h = 1.5
    card_gap = 0.12
    card_start_x = (13.333 - 8 * card_w - 7 * card_gap) / 2

    for i, (head, relation, example, color) in enumerate(perspectives):
        x = card_start_x + i * (card_w + card_gap)
        y = desc_y + 0.4
        add_rounded_rect(slide, x, y, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF),
                         line_color=color, name=f"perspective-{i}")
        # 色条
        add_rect(slide, x, y, card_w, 0.3, color,
                 text=head, font_size=11, font_color=C_TEXT_WHITE, bold=True,
                 name=f"perspective-head-{i}")
        add_text_box(slide, x + 0.05, y + 0.4, card_w - 0.1, 0.3,
                     relation, font_size=11, font_color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.05, y + 0.75, card_w - 0.1, 0.6,
                     example, font_size=10, font_color=C_TEXT_DARK,
                     alignment=PP_ALIGN.CENTER)

    # --- 关键点说明 ---
    key_y = 6.2
    add_rect(slide, 0.4, key_y, 12.5, 0.9, C_LIGHT_GREEN, name="key-point-box")
    add_text_box(slide, 0.6, key_y + 0.05, 12.0, 0.3,
                 "关键设计: 计算量与单头注意力相同!",
                 font_size=13, font_color=C_FFN, bold=True)
    add_text_box(slide, 0.6, key_y + 0.4, 12.0, 0.4,
                 "将 d_model(512) 分成 8 个头, 每个头维度 d_k=64。8个头的计算量之和 = 1个大头的计算量, "
                 "但能学习更丰富的表示。",
                 font_size=11, font_color=C_TEXT_DARK)

    print("  Page 3 完成")


# ============================================================
# 页面4: Transformer 整体架构图 (Encoder-Decoder)
# ============================================================
def create_page4(prs):
    """左:编码器6层, 右:解码器6层, 中间:交叉注意力连接"""
    print("\n=== Page 4: Transformer 整体架构图 ===")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_page_title(slide, "Transformer: Encoder-Decoder 整体架构")

    # 整体布局
    # 编码器在左侧, 解码器在右侧
    # 用堆叠的方式展示6层 (展示2层 + "x6" 标记)

    enc_x = 1.5     # 编码器左边界
    dec_x = 7.5     # 解码器左边界
    comp_w = 2.2    # 组件宽度
    row_h = 0.55    # 每行高度
    comp_gap = 0.05 # 组件间间隙
    row_gap = 0.08  # 行间间隙

    # --- 顶部: 输出层 ---
    top_y = 0.85
    add_text_box(slide, enc_x + 0.2, top_y, 2.0, 0.3,
                 "编码器 Encoder", font_size=14, font_color=C_ENCODER, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, dec_x + 0.2, top_y, 2.0, 0.3,
                 "解码器 Decoder", font_size=14, font_color=C_DECODER, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # --- 底部: Embedding + Position Encoding ---
    bottom_y = 6.3
    add_rounded_rect(slide, enc_x, bottom_y, comp_w, 0.7, C_GRAY,
                     text="Embedding\n+ Position Encoding", font_size=11,
                     font_color=C_TEXT_WHITE, bold=True, name="enc-embed")
    add_rounded_rect(slide, dec_x, bottom_y, comp_w, 0.7, C_GRAY,
                     text="Embedding\n+ Position Encoding", font_size=11,
                     font_color=C_TEXT_WHITE, bold=True, name="dec-embed")

    # --- 编码器层 (展示2层) ---
    enc_layers = [
        {"y": 1.3, "label": "Encoder Layer 1"},
        {"y": 3.0, "label": "Encoder Layer 2"},
    ]

    for li, layer in enumerate(enc_layers):
        y = layer["y"]
        # 层背景
        add_rect(slide, enc_x - 0.1, y - 0.05, comp_w + 0.2, 1.55, C_LIGHT_BLUE,
                 line_color=C_ENCODER, name=f"enc-layer-bg-{li}")

        # Multi-Head Self-Attention
        add_rounded_rect(slide, enc_x, y, comp_w, row_h, C_ATTENTION,
                         text="Multi-Head\nSelf-Attention", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"enc-attn-{li}")

        # Add & Norm
        y1 = y + row_h + comp_gap
        add_rounded_rect(slide, enc_x, y1, comp_w, row_h * 0.75, C_NORM,
                         text="Add & Layer Norm", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"enc-norm1-{li}")

        # FFN
        y2 = y1 + row_h * 0.75 + comp_gap
        add_rounded_rect(slide, enc_x, y2, comp_w, row_h, C_FFN,
                         text="Feed-Forward\nNetwork (FFN)", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"enc-ffn-{li}")

        # Add & Norm 2
        y3 = y2 + row_h + comp_gap
        add_rounded_rect(slide, enc_x, y3, comp_w, row_h * 0.75, C_NORM,
                         text="Add & Layer Norm", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"enc-norm2-{li}")

    # "x6" 标记
    add_text_box(slide, enc_x + comp_w + 0.2, 2.0, 0.8, 0.6,
                 "x N\n(N=6)", font_size=13, font_color=C_ENCODER, bold=True,
                 alignment=PP_ALIGN.CENTER, name="enc-x6")

    # --- 解码器层 (展示2层, 3个子层) ---
    dec_layers = [
        {"y": 1.3, "label": "Decoder Layer 1"},
        {"y": 3.8, "label": "Decoder Layer 2"},
    ]

    for li, layer in enumerate(dec_layers):
        y = layer["y"]
        # 层背景
        add_rect(slide, dec_x - 0.1, y - 0.05, comp_w + 0.2, 2.35, C_LIGHT_RED,
                 line_color=C_DECODER, name=f"dec-layer-bg-{li}")

        # Masked Multi-Head Self-Attention
        add_rounded_rect(slide, dec_x, y, comp_w, row_h, RGBColor(0xC0, 0x39, 0x2B),
                         text="Masked\nSelf-Attention", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-masked-attn-{li}")

        # Add & Norm
        y1 = y + row_h + comp_gap
        add_rounded_rect(slide, dec_x, y1, comp_w, row_h * 0.65, C_NORM,
                         text="Add & Norm", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-norm1-{li}")

        # Cross-Attention (编码器-解码器注意力)
        y2 = y1 + row_h * 0.65 + comp_gap
        add_rounded_rect(slide, dec_x, y2, comp_w, row_h, C_PURPLE,
                         text="Cross-Attention\n(Encoder-Decoder)", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-cross-attn-{li}")

        # Add & Norm
        y3 = y2 + row_h + comp_gap
        add_rounded_rect(slide, dec_x, y3, comp_w, row_h * 0.65, C_NORM,
                         text="Add & Norm", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-norm2-{li}")

        # FFN
        y4 = y3 + row_h * 0.65 + comp_gap
        add_rounded_rect(slide, dec_x, y4, comp_w, row_h, C_FFN,
                         text="Feed-Forward\nNetwork (FFN)", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-ffn-{li}")

        # Add & Norm 3
        y5 = y4 + row_h + comp_gap
        add_rounded_rect(slide, dec_x, y5, comp_w, row_h * 0.65, C_NORM,
                         text="Add & Norm", font_size=11,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"dec-norm3-{li}")

    # "x6" 标记
    add_text_box(slide, dec_x + comp_w + 0.2, 2.2, 0.8, 0.6,
                 "x N\n(N=6)", font_size=13, font_color=C_DECODER, bold=True,
                 alignment=PP_ALIGN.CENTER, name="dec-x6")

    # --- 中间: 交叉注意力连接箭头 ---
    # 从编码器到解码器的交叉注意力连接
    arrow_y = 2.3
    add_arrow(slide, 4.0, arrow_y, 3.2, 0.2, C_PURPLE, "cross-attn-arrow")
    add_text_box(slide, 4.5, arrow_y - 0.3, 2.0, 0.25,
                 "K, V", font_size=11, font_color=C_PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    arrow_y2 = 4.8
    add_arrow(slide, 4.0, arrow_y2, 3.2, 0.2, C_PURPLE, "cross-attn-arrow2")
    add_text_box(slide, 4.5, arrow_y2 - 0.3, 2.0, 0.25,
                 "K, V", font_size=11, font_color=C_PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # --- 顶部: Linear + Softmax ---
    out_y = 0.85
    add_rounded_rect(slide, dec_x + comp_w + 0.2, 1.0, 1.8, 0.6, C_DARK_BLUE,
                     text="Linear + Softmax", font_size=12,
                     font_color=C_TEXT_WHITE, bold=True, name="output-layer")
    add_arrow(slide, dec_x + comp_w + 0.05, 1.1, 0.1, 0.15, C_ARROW, "dec-out-arrow")

    # 输出标签
    add_text_box(slide, dec_x + comp_w + 0.2, 1.7, 1.8, 0.3,
                 "Output Probabilities", font_size=11, font_color=C_TEXT_DARK,
                 alignment=PP_ALIGN.CENTER)

    # --- 输入/输出标签 ---
    add_text_box(slide, enc_x - 0.3, bottom_y + 0.7, comp_w + 0.6, 0.3,
                 "Input Sequence", font_size=12, font_color=C_ENCODER, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, dec_x - 0.3, bottom_y + 0.7, comp_w + 0.6, 0.3,
                 "Output Sequence (shifted right)", font_size=12, font_color=C_DECODER, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # --- 图例 ---
    legend_y = 6.2
    legend_items = [
        (C_ATTENTION, "Attention"),
        (C_FFN, "FFN"),
        (C_NORM, "Add & Norm"),
        (C_PURPLE, "Cross-Attention"),
    ]
    lx = 4.8
    for i, (color, label) in enumerate(legend_items):
        x = lx + i * 1.8
        add_rect(slide, x, legend_y, 0.3, 0.25, color, name=f"legend-{i}")
        add_text_box(slide, x + 0.35, legend_y, 1.3, 0.25,
                     label, font_size=10, font_color=C_TEXT_DARK)

    print("  Page 4 完成")


# ============================================================
# 页面5: 训练 vs 推理 对比图
# ============================================================
def create_page5(prs):
    """左: 训练阶段 (Teacher Forcing), 右: 推理阶段 (逐步生成)"""
    print("\n=== Page 5: 训练 vs 推理 对比图 ===")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_page_title(slide, "训练 vs 推理: Teacher Forcing 对比")

    # --- 左半: 训练阶段 ---
    train_x = 0.4
    train_w = 6.2

    # 训练区域背景
    add_rounded_rect(slide, train_x, 0.9, train_w, 6.2, C_LIGHT_BLUE,
                     line_color=C_ENCODER, name="train-bg")

    # 标题
    add_text_box(slide, train_x + 0.2, 1.0, train_w - 0.4, 0.4,
                 "训练阶段 Training (Teacher Forcing)",
                 font_size=15, font_color=C_ENCODER, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 比喻标签
    add_rounded_rect(slide, train_x + 0.3, 1.5, 2.5, 0.4, C_ENCODER,
                     text='比喻: "带答案做练习题"', font_size=11,
                     font_color=C_TEXT_WHITE, bold=True, name="train-metaphor")

    # 输入: 完整句子
    input_label_y = 2.1
    add_text_box(slide, train_x + 0.3, input_label_y, train_w - 0.6, 0.3,
                 "输入 (完整句子, 一次全部给出):",
                 font_size=12, font_color=C_TEXT_DARK, bold=True)

    # 输入序列
    words_train = ["<BOS>", "I", "love", "you"]
    wx = train_x + 0.5
    wy = input_label_y + 0.35
    ww = 1.1
    wh = 0.45
    for i, w in enumerate(words_train):
        add_rounded_rect(slide, wx + i * (ww + 0.1), wy, ww, wh, C_ATTENTION,
                         text=w, font_size=12, font_color=C_TEXT_WHITE, bold=True,
                         name=f"train-input-{i}")

    # 下箭头
    add_down_arrow(slide, train_x + train_w / 2 - 0.15, wy + wh + 0.1, 0.3, 0.35,
                   C_ENCODER, "train-down-arrow")

    # 解码器处理
    dec_y = wy + wh + 0.55
    add_rounded_rect(slide, train_x + 0.8, dec_y, train_w - 1.6, 0.7, C_ENCODER,
                     text="Decoder (看到所有输入 + 之前已知的正确输出)",
                     font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="train-decoder")

    # 下箭头
    add_down_arrow(slide, train_x + train_w / 2 - 0.15, dec_y + 0.75, 0.3, 0.35,
                   C_ENCODER, "train-down-arrow2")

    # 预测输出
    pred_y = dec_y + 1.2
    add_text_box(slide, train_x + 0.3, pred_y, train_w - 0.6, 0.3,
                 "预测输出:", font_size=12, font_color=C_TEXT_DARK, bold=True)

    words_pred = ["我", "爱", "你", "<EOS>"]
    for i, w in enumerate(words_pred):
        color = C_FFN if i < 3 else C_GRAY
        add_rounded_rect(slide, wx + i * (ww + 0.1), pred_y + 0.35, ww, wh, color,
                         text=w, font_size=12, font_color=C_TEXT_WHITE, bold=True,
                         name=f"train-pred-{i}")

    # 对比: 正确答案
    truth_y = pred_y + 1.0
    add_text_box(slide, train_x + 0.3, truth_y, train_w - 0.6, 0.3,
                 "正确答案 (Ground Truth):", font_size=12, font_color=C_TEXT_DARK, bold=True)

    words_truth = ["我", "爱", "你", "<EOS>"]
    for i, w in enumerate(words_truth):
        add_rounded_rect(slide, wx + i * (ww + 0.1), truth_y + 0.35, ww, wh, C_TEAL,
                         text=w, font_size=12, font_color=C_TEXT_WHITE, bold=True,
                         name=f"train-truth-{i}")

    # Loss 计算
    loss_y = truth_y + 1.0
    add_rounded_rect(slide, train_x + 1.0, loss_y, train_w - 2.0, 0.5, C_DECODER,
                     text="Loss = CrossEntropy(预测, 正确答案)",
                     font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="train-loss")

    # --- 右半: 推理阶段 ---
    infer_x = 6.9
    infer_w = 6.2

    # 推理区域背景
    add_rounded_rect(slide, infer_x, 0.9, infer_w, 6.2, C_LIGHT_RED,
                     line_color=C_DECODER, name="infer-bg")

    # 标题
    add_text_box(slide, infer_x + 0.2, 1.0, infer_w - 0.4, 0.4,
                 "推理阶段 Inference (Auto-Regressive)",
                 font_size=15, font_color=C_DECODER, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 比喻标签
    add_rounded_rect(slide, infer_x + 0.3, 1.5, 2.2, 0.4, C_DECODER,
                     text='比喻: "闭卷考试"', font_size=11,
                     font_color=C_TEXT_WHITE, bold=True, name="infer-metaphor")

    # 逐步生成过程
    step_label_y = 2.1
    add_text_box(slide, infer_x + 0.3, step_label_y, infer_w - 0.6, 0.3,
                 "一个词一个词生成 (每步只能看到之前生成的词):",
                 font_size=12, font_color=C_TEXT_DARK, bold=True)

    # Step 1
    s1_y = step_label_y + 0.4
    add_text_box(slide, infer_x + 0.3, s1_y, 1.2, 0.3,
                 "Step 1:", font_size=11, font_color=C_DECODER, bold=True)
    add_rounded_rect(slide, infer_x + 1.4, s1_y, 0.9, 0.35, C_ENCODER,
                     text="<BOS>", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s1-input")
    add_text_box(slide, infer_x + 2.4, s1_y, 0.3, 0.35,
                 "->", font_size=14, font_color=C_ARROW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, infer_x + 2.8, s1_y, 0.9, 0.35, C_FFN,
                     text="我", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s1-out")

    # Step 2
    s2_y = s1_y + 0.45
    add_text_box(slide, infer_x + 0.3, s2_y, 1.2, 0.3,
                 "Step 2:", font_size=11, font_color=C_DECODER, bold=True)
    add_rounded_rect(slide, infer_x + 1.4, s2_y, 0.9, 0.35, C_ENCODER,
                     text="<BOS>", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s2-in1")
    add_rounded_rect(slide, infer_x + 2.4, s2_y, 0.9, 0.35, C_FFN,
                     text="我", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s2-in2")
    add_text_box(slide, infer_x + 3.4, s2_y, 0.3, 0.35,
                 "->", font_size=14, font_color=C_ARROW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, infer_x + 3.8, s2_y, 0.9, 0.35, C_FFN,
                     text="爱", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s2-out")

    # Step 3
    s3_y = s2_y + 0.45
    add_text_box(slide, infer_x + 0.3, s3_y, 1.2, 0.3,
                 "Step 3:", font_size=11, font_color=C_DECODER, bold=True)
    add_rounded_rect(slide, infer_x + 1.4, s3_y, 0.9, 0.35, C_ENCODER,
                     text="<BOS>", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s3-in1")
    add_rounded_rect(slide, infer_x + 2.4, s3_y, 0.9, 0.35, C_FFN,
                     text="我", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s3-in2")
    add_rounded_rect(slide, infer_x + 3.4, s3_y, 0.9, 0.35, C_FFN,
                     text="爱", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s3-in3")
    add_text_box(slide, infer_x + 4.4, s3_y, 0.3, 0.35,
                 "->", font_size=14, font_color=C_ARROW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, infer_x + 4.8, s3_y, 0.9, 0.35, C_FFN,
                     text="你", font_size=12, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s3-out")

    # Step 4
    s4_y = s3_y + 0.45
    add_text_box(slide, infer_x + 0.3, s4_y, 1.2, 0.3,
                 "Step 4:", font_size=11, font_color=C_DECODER, bold=True)
    # 显示 <BOS> 我 爱 你 -> <EOS>
    tokens_s4 = ["<BOS>", "我", "爱", "你"]
    token_w = 0.7
    colors_s4 = [C_ENCODER, C_FFN, C_FFN, C_FFN]
    for j, (t, c) in enumerate(zip(tokens_s4, colors_s4)):
        add_rounded_rect(slide, infer_x + 1.4 + j * (token_w + 0.05), s4_y,
                         token_w, 0.35, c, text=t, font_size=10,
                         font_color=C_TEXT_WHITE, bold=True,
                         name=f"infer-s4-in{j}")
    add_text_box(slide, infer_x + 1.4 + 4 * (token_w + 0.05), s4_y, 0.3, 0.35,
                 "->", font_size=14, font_color=C_ARROW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    eos_x = infer_x + 1.4 + 4 * (token_w + 0.05) + 0.35
    add_rounded_rect(slide, eos_x, s4_y, 0.9, 0.35, C_GRAY,
                     text="<EOS>", font_size=11, font_color=C_TEXT_WHITE, bold=True,
                     name="infer-s4-out")

    # ... 省略号
    s5_y = s4_y + 0.45
    add_text_box(slide, infer_x + 0.3, s5_y, infer_w - 0.6, 0.3,
                 "... 重复直到生成 <EOS> (结束标记)",
                 font_size=11, font_color=C_TEXT_DARK)

    # 关键差异说明
    diff_y = s5_y + 0.6
    add_rounded_rect(slide, infer_x + 0.3, diff_y, infer_w - 0.6, 1.3,
                     RGBColor(0xFF, 0xFF, 0xFF), line_color=C_DECODER,
                     name="infer-note")
    add_text_box(slide, infer_x + 0.5, diff_y + 0.1, infer_w - 1.0, 0.3,
                 "关键差异:", font_size=12, font_color=C_DECODER, bold=True)
    add_text_box(slide, infer_x + 0.5, diff_y + 0.45, infer_w - 1.0, 0.8,
                 "1. 没有 '正确答案' 可参考, 只能用自己之前生成的词\n"
                 "2. 一旦某步生成了错误的词, 后续会跟着错 (误差累积)\n"
                 "3. 推理速度受限于序列长度 (无法并行)",
                 font_size=11, font_color=C_TEXT_DARK)

    # --- 中间: VS 分隔线 ---
    mid_x = 6.55
    add_text_box(slide, mid_x - 0.15, 3.5, 0.5, 0.5,
                 "VS", font_size=20, font_color=C_ARROW, bold=True,
                 alignment=PP_ALIGN.CENTER)

    print("  Page 5 完成")


# ============================================================
# 主函数
# ============================================================
def main():
    print("=" * 60)
    print("Transformer 架构图 PPT 生成器")
    print("=" * 60)

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 创建 5 页架构图
    create_page1(prs)  # Self-Attention Q/K/V 概念图
    create_page2(prs)  # Self-Attention 5步计算流程图
    create_page3(prs)  # Multi-Head Attention 架构图
    create_page4(prs)  # Transformer 整体架构图
    create_page5(prs)  # 训练 vs 推理 对比图

    # 保存
    output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_diagrams.pptx"
    prs.save(output_path)
    print(f"\n=== 文件已保存到: {output_path} ===")

    # 验证: 检查所有 shapes 的边界
    print("\n=== 验证所有 shapes 边界 ===")
    errors = 0
    for si, slide in enumerate(prs.slides):
        print(f"\nSlide {si + 1}: {len(slide.shapes)} shapes")
        for shape in slide.shapes:
            left = shape.left / 914400  # EMU to inches
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            right = left + width
            bottom = top + height
            if right > MAX_RIGHT or bottom > MAX_BOTTOM:
                errors += 1
                print(f"  [ERROR] '{shape.name}': left={left:.2f}, top={top:.2f}, "
                      f"w={width:.2f}, h={height:.2f}, right={right:.2f}, bottom={bottom:.2f}")

    if errors == 0:
        print(f"\n=== 验证通过! 所有 {sum(len(s.shapes) for s in prs.slides)} 个 shapes 都在边界内 ===")
    else:
        print(f"\n=== 发现 {errors} 个越界错误, 需要修复 ===")

    return errors


if __name__ == "__main__":
    errors = main()
    exit(errors)
