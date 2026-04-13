#!/usr/bin/env python3
"""Transformer v8: Fix 4 issues from QA round 3 (94→98+)
Issues:
1. Slide 1: Title text box too tall (1097280→640080 EMU)
2. Slide 8: Two bottom cards overlap (right card beyond slide edge)
3. Slide 8: Bottom bar width inconsistent (11612880→12191695)
4. Slide 10: Table overlaps bottom bar
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import sys

SRC = "output/Transformer_v7_zh.pptx"
DST = "output/Transformer_v8_zh.pptx"

prs = Presentation(SRC)

SLIDE_W = prs.slide_width  # 12192000 EMU = 13.333 inches
SLIDE_H = prs.slide_height  # 6858000 EMU = 7.5 inches
BOTTOM_BAR_TOP = Emu(6537960)  # Inches(5.12)
BOTTOM_BAR_FULL_W = Emu(12191695)

print(f"Slide W={SLIDE_W} H={SLIDE_H}")
print(f"Bottom bar top={BOTTOM_BAR_TOP}")
print(f"Total slides: {len(prs.slides)}")

# ============================================================
# Fix 1: Slide 1 (index 0) — Title text box too tall
# ============================================================
sl = prs.slides[0]
print(f"\n=== Fix 1: Slide 1 (封面) ===")
for sh in sl.shapes:
    tf = sh.text_frame if sh.has_text_frame else None
    if tf and tf.paragraphs:
        txt = tf.paragraphs[0].text.strip() if tf.paragraphs[0].text else ""
        if txt == "Transformer":
            old_h = sh.height
            old_top = sh.top
            old_bottom = old_top + old_h
            print(f"  Found 'Transformer' box: top={old_top} h={old_h} bottom={old_bottom}")
            # Shrink height from 1097280 to 640080 EMU (~50pt)
            new_h = Emu(640080)
            sh.height = new_h
            print(f"  Fixed: h={old_h} → {new_h}")
            break

# ============================================================
# Fix 2 & 3: Slide 8 (index 7) — Card overlap + bottom bar width
# ============================================================
sl = prs.slides[7]
print(f"\n=== Fix 2&3: Slide 8 (自注意力5步) ===")

# Find the two overlapping rounded rectangles at bottom
# RR38 (多头注意力): L=274320, top=4617720, W=11612880
# RR41 (为什么除以sqrt(d_k)): L=6035040, top=4617720, W=11612880 (exceeds slide!)
# There are also background rounded rects (no text) that need fixing
# Solution: Make them side-by-side, each ~half width

cards_to_fix = []
textboxes_to_fix = []

# Track shapes at the bottom card area (y ~4617720)
# Some are title rects, some are background rects, some are text boxes
bottom_shapes = []  # (shape, role_hint)

for sh in sl.shapes:
    if sh.shape_type is not None and hasattr(sh, 'left'):
        # Collect shapes in the bottom card area
        if abs(sh.top - Emu(4617720)) < Emu(50000):  # Near y=4617720
            tf = sh.text_frame if sh.has_text_frame else None
            txt = ""
            if tf and tf.paragraphs:
                txt = tf.paragraphs[0].text.strip() if tf.paragraphs[0].text else ""
            bottom_shapes.append((sh, txt))

print(f"  Found {len(bottom_shapes)} shapes in bottom card area")

for sh, txt in bottom_shapes:
    # Find "多头注意力" card title (left)
    if "多头注意力" in txt and "讨论组" in txt:
        print(f"  Found left card title '{txt[:25]}': L={sh.left} T={sh.top} W={sh.width} H={sh.height}")
        cards_to_fix.append(("left_title", sh))
    
    # Find "为什么除以sqrt(d_k)" card title (right, overflowing)
    if "sqrt" in txt or "为什么除以" in txt:
        print(f"  Found right card title '{txt[:25]}': L={sh.left} T={sh.top} W={sh.width} H={sh.height}")
        cards_to_fix.append(("right_title", sh))
    
    # Background rounded rect with no text (left card bg)
    if not txt and sh.width > Emu(5000000) and sh.left < Emu(3000000):
        print(f"  Found left card bg (no text): L={sh.left} T={sh.top} W={sh.width} H={sh.height}")
        cards_to_fix.append(("left_bg", sh))
    
    # Background rounded rect with no text (right card bg, overflowing)
    if not txt and sh.width > Emu(5000000) and sh.left >= Emu(3000000):
        print(f"  Found right card bg (no text): L={sh.left} T={sh.top} W={sh.width} H={sh.height}")
        cards_to_fix.append(("right_bg", sh))

# Fix the cards - make them side by side properly
# Available width: SLIDE_W - 2*margin = 12192000 - 2*274320 = 11643360
# Each card: (11643360 - gap) / 2
margin = Emu(274320)
gap = Emu(182880)  # ~0.2 inch gap
available_w = SLIDE_W - 2 * margin - gap
card_w = available_w // 2  # ~5721600 each

# Also need to find and move text boxes inside these cards
# First, let's find all textboxes associated with these cards
left_card_body = None
right_card_body = None

for sh in sl.shapes:
    if sh.has_text_frame and sh.text_frame.paragraphs:
        txt = sh.text_frame.paragraphs[0].text.strip() if sh.text_frame.paragraphs[0].text else ""
        # Left card body text: "不是1个注意力头，而是8个独立头同时工作"
        if "不是1个注意力头" in txt:
            left_card_body = sh
            print(f"  Found left card body: L={sh.left} T={sh.top} W={sh.width} H={sh.height}")
        # Right card body text: "当d_k=512时"
        if "当d_k=512时" in txt:
            right_card_body = sh
            print(f"  Found right card body: L={sh.left} T={sh.top} W={sh.width} H={sh.height}")

# Fix positions
for role, sh in cards_to_fix:
    if role.startswith("left"):
        sh.left = margin
        sh.width = card_w
        print(f"  Fixed {role}: L={sh.left} W={sh.width}")
    else:
        sh.left = margin + card_w + gap
        sh.width = card_w
        print(f"  Fixed {role}: L={sh.left} W={sh.width}")

# Fix body text boxes to fit within new card widths
if left_card_body:
    left_card_body.left = Emu(274320 + 182880)  # margin + 0.15in padding
    left_card_body.width = Emu(int(card_w) - 365760)  # card_w - 2*padding
    print(f"  Fixed left body: L={left_card_body.left} W={left_card_body.width}")

if right_card_body:
    right_card_body.left = Emu(int(margin + card_w + gap) + 182880)
    right_card_body.width = Emu(int(card_w) - 365760)
    print(f"  Fixed right body: L={right_card_body.left} W={right_card_body.width}")

# Fix 3: Bottom bar width on Slide 8
print("\n  Fixing bottom bar width...")
for sh in sl.shapes:
    tf = sh.text_frame if sh.has_text_frame else None
    txt = ""
    if tf and tf.paragraphs:
        txt = tf.paragraphs[0].text.strip() if tf.paragraphs[0].text else ""
    
    # Bottom bar rectangle (no text, dark blue, at y=7.15)
    if sh.top >= Emu(6500000) and not sh.has_text_frame:
        if sh.width < SLIDE_W:
            old_w = sh.width
            sh.width = BOTTOM_BAR_FULL_W
            sh.left = Emu(0)
            print(f"  Fixed bottom bar rect: W={old_w} → {BOTTOM_BAR_FULL_W}")
    
    # Bottom bar text ("Transformer 架构深度解析")
    if "Transformer 架构深度解析" in txt:
        old_w = sh.width
        sh.width = BOTTOM_BAR_FULL_W
        print(f"  Fixed bottom bar text: W={old_w} → {BOTTOM_BAR_FULL_W}")
    
    # Page number text (e.g. "6/12")
    if "/" in txt and len(txt) <= 5:
        # Page number is at right side of bottom bar - adjust if needed
        pass

# ============================================================
# Fix 4: Slide 10 (index 9) — Table overlaps bottom bar
# ============================================================
sl = prs.slides[9]
print(f"\n=== Fix 4: Slide 10 (推理阶段) ===")

for sh in sl.shapes:
    if sh.has_table:
        tbl = sh.table
        old_top = sh.top
        old_bottom = sh.top + sh.height
        print(f"  Found table: top={old_top} h={sh.height} bottom={old_bottom}")
        print(f"  Bottom bar top={BOTTOM_BAR_TOP}")
        
        # Need table bottom <= bottom_bar_top with some margin
        target_bottom = BOTTOM_BAR_TOP - Emu(91440)  # 0.1 inch margin
        new_height = target_bottom - old_top
        
        if sh.top + sh.height > BOTTOM_BAR_TOP:
            print(f"  Table overflows by {(sh.top + sh.height) - BOTTOM_BAR_TOP} EMU")
            # Option: reduce row heights proportionally
            # Move table up slightly and reduce height
            sh.height = new_height
            print(f"  Fixed: h={old_bottom - old_top} → {new_height}")
            
            # Also try to move it up a bit
            sh.top = Emu(int(old_top) - 45720)  # Move up 0.05 inch
            print(f"  Moved up: top={old_top} → {sh.top}")

# ============================================================
# Also fix page number position on Slide 8 (ensure it's visible)
# ============================================================
sl8 = prs.slides[7]
for sh in sl8.shapes:
    if sh.has_text_frame and sh.text_frame.paragraphs:
        txt = sh.text_frame.paragraphs[0].text.strip() if sh.text_frame.paragraphs[0].text else ""
        # Page number like "6/12"
        if "/" in txt and len(txt) <= 5 and txt[0].isdigit():
            # Make sure it's in the right position (right side of bottom bar)
            print(f"\n  Slide 8 page number '{txt}': L={sh.left} T={sh.top}")

# ============================================================
# Save
# ============================================================
prs.save(DST)
print(f"\n✅ Saved: {DST}")

# Verify: check all shapes are within bounds
print("\n=== Verification ===")
issues = 0
for i, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        # Check right boundary
        if sh.left + sh.width > SLIDE_W + Emu(5000):
            tf = sh.text_frame if sh.has_text_frame else None
            txt = ""
            if tf and tf.paragraphs:
                txt = tf.paragraphs[0].text[:30] if tf.paragraphs[0].text else ""
            print(f"  ⚠️ Slide {i+1}: shape exceeds right edge by {(sh.left + sh.width) - SLIDE_W} EMU '{txt}'")
            issues += 1
        # Check bottom boundary
        if sh.top + sh.height > SLIDE_H + Emu(5000):
            tf = sh.text_frame if sh.has_text_frame else None
            txt = ""
            if tf and tf.paragraphs:
                txt = tf.paragraphs[0].text[:30] if tf.paragraphs[0].text else ""
            print(f"  ⚠️ Slide {i+1}: shape exceeds bottom edge by {(sh.top + sh.height) - SLIDE_H} EMU '{txt}'")
            issues += 1

if issues == 0:
    print("  ✅ All shapes within bounds!")
else:
    print(f"  ❌ {issues} issues found")
