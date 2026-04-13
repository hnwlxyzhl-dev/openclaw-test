#!/usr/bin/env python3
"""Transformer PPT v7 Part2: Pages 8-13"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = "output/transformer_v7_part2.pptx"

C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_ENC = RGBColor(0x34,0x98,0xDB)
C_DEC = RGBColor(0xE7,0x4C,0x3C)
C_ATT = RGBColor(0xF3,0x9C,0x12)
C_FFN = RGBColor(0x2E,0xCC,0x71)
C_EMB = RGBColor(0x9B,0x59,0xB6)
C_LT = RGBColor(0xEB,0xF5,0xFB)
C_W = RGBColor(0xFF,0xFF,0xFF)
C_D = RGBColor(0x2C,0x3E,0x50)
C_G = RGBColor(0x88,0x88,0x88)
C_BG_ENC = RGBColor(0xE8,0xF0,0xFE)
C_BG_DEC = RGBColor(0xFD,0xED,0xED)
C_BG_ATT = RGBColor(0xFE,0xF9,0xE7)
C_BG_FFN = RGBColor(0xE8,0xF8,0xF5)
C_BG_EMB = RGBColor(0xF4,0xEC,0xF7)
C_BG_GRN = RGBColor(0xE8,0xF8,0xF5)
C_BG_ORG = RGBColor(0xFD,0xF2,0xE9)

FC = "微软雅黑"
FE = "Arial"

def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc; sh.line.width=Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc; sh.line.width=Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    for r in p.runs: r.font.size=Pt(fs); r.font.color.rgb=fc; r.font.bold=b; r.font.name=fn
    p.font.size=Pt(fs); p.font.color.rgb=fc; p.font.bold=b; p.font.name=fn
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
        for r in p.runs: r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=bld; r.font.name=fn
        p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld; p.font.name=fn
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def line_h(s,l,t,w,c=C_ACC,lw=2):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,Emu(int(lw*12700)))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Transformer 架构深度解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    ok = True
    for sh in s.shapes:
        if sh.left < Emu(-100) or sh.top < Emu(-100): ok = False
        if sh.left+sh.width > SLIDE_W+Emu(5000) or sh.top+sh.height > SLIDE_H+Emu(5000): ok = False
    print(f"  {'OK' if ok else 'FAIL'} {msg}")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ==================== P8 自注意力5步计算 ====================
print("P8: 自注意力5步计算")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "06  自注意力5步计算", "从QKV生成到加权求和 — 所有词同时并行计算，无串行依赖")

# 5步流程卡片
steps = [
    ("Step 1", "生成Q,K,V", "X×W_Q=Q, X×W_K=K, X×W_V=V\n每个512维X乘以3个矩阵\n输出各512维", C_ENC),
    ("Step 2", "计算分数矩阵", "Q×K^T → Score矩阵\nT×T方阵(每个词对)\n本质=高维余弦相似度", C_ATT),
    ("Step 3", "缩放", "Score / sqrt(d_k)\nsqrt(64)=8\n防止点积过大→梯度饱和", C_EMB),
    ("Step 4", "Softmax归一化", "每行Softmax→概率\n每行和=1\n概率=注意力权重", C_FFN),
    ("Step 5", "加权求和", "概率×V → Output\n融合所有词的信息\n每个词得到上下文表示", C_DEC),
]
for i, (step, title, desc, color) in enumerate(steps):
    x = Inches(0.30 + i * 2.58)
    rrect(sl, x, Inches(1.20), Inches(2.35), Inches(2.00), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(2.35), Inches(0.55), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(2.25), Inches(0.25), step, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(1.45), Inches(2.25), Inches(0.25), title, fs=11, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(1.82), Inches(2.15), Inches(1.25), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)
    if i < len(steps) - 1:
        arrow_r(sl, x+Inches(2.35), Inches(2.05), Inches(0.22), Inches(0.16), color)

# 具体数值示例
rrect(sl, Inches(0.30), Inches(3.40), Inches(12.70), Inches(1.50), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(3.45), Inches(12.30), Inches(0.25), "具体数值示例：\"我 爱 AI\" — 计算\"爱\"的自注意力输出", fs=11, fc=C_ACC, b=True)
mtb(sl, Inches(0.50), Inches(3.75), Inches(12.30), Inches(1.00), [
    ("输入X: ", True, C_PRI, 9),
    ("\"我\"=[0.5,0.2,...]  \"爱\"=[0.1,0.8,...]  \"AI\"=[0.3,0.4,...]  (512维向量)", False, C_D, 9),
    ("Step1 Q,K,V: ", True, C_PRI, 9),
    ("X分别乘W_Q/W_K/W_V(各512×512) → Q_爱=[0.2,0.9,...], K_我=[0.4,0.1,...], K_爱=[0.1,0.7,...], K_AI=[0.3,0.5,...]", False, C_D, 9),
    ("Step2-3 Score: ", True, C_PRI, 9),
    ("Q_爱·K_我=0.68, Q_爱·K_爱=1.21, Q_爱·K_AI=0.42 → 缩放: [0.085, 0.151, 0.053]", False, C_D, 9),
    ("Step4 Softmax: ", True, C_PRI, 9),
    ("[0.085, 0.151, 0.053] → [0.355, 0.418, 0.227]  → Step5 Output = 0.355×V_我 + 0.418×V_爱 + 0.227×V_AI", False, C_D, 9),
], ls=1.1)

# 多头说明
rrect(sl, Inches(0.30), Inches(5.05), Inches(6.10), Inches(1.85), fc=C_BG_ENC, lc=C_ENC, lw=1)
tb(sl, Inches(0.50), Inches(5.10), Inches(5.70), Inches(0.25), "多头注意力 = 8个独立讨论组", fs=11, fc=C_ENC, b=True)
mtb(sl, Inches(0.50), Inches(5.40), Inches(5.70), Inches(1.40), [
    "不是1个注意力头，而是8个独立头同时工作：",
    "头1: 专注语法(主谓宾结构) | 头2: 关注语义(同义词) | 头3: 关注位置(相邻词) | 头4: 关注指代(代词→名词)",
    "每个头独立计算Q,K,V(各64维=512/8)，8个输出拼接: Concat(h1,...,h8) = 512维",
    "最后通过线性变换W_O(512×512)融合8个头的信息。参数量: W_Q+W_K+W_V+W_O = 4×512×512 = 约100万",
], ls=1.15)

rrect(sl, Inches(6.60), Inches(5.05), Inches(6.40), Inches(1.85), fc=C_BG_ATT, lc=C_ATT, lw=1)
tb(sl, Inches(6.80), Inches(5.10), Inches(6.00), Inches(0.25), "为什么除以sqrt(d_k)？", fs=11, fc=C_ATT, b=True)
mtb(sl, Inches(6.80), Inches(5.40), Inches(6.00), Inches(1.40), [
    "当d_k=512时，点积Q·K的值域约为[-100,100]（方差=d_k=512）。",
    "Softmax对大输入极其敏感：输入差10倍→输出差22000倍。数值过大→梯度接近0→训练停滞(饱和)。",
    "除以sqrt(64)=8后，值域缩到[-12,12]，Softmax梯度正常(约0.01-0.1)。",
    "这个缩放因子是Transformer论文的\"隐藏功臣\"——没有它，模型根本无法训练。数学原理：Var(Q·K)=d_k→除以sqrt(d_k)→方差归一化为1。",
], ls=1.15)
chk(sl, "P8")

# ==================== P9 训练细节 ====================
print("P9: 训练细节")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "07  训练细节", "残差连接的数学原理、因果掩码与交叉熵损失函数")

# 3个核心概念卡片
cards = [
    ("残差连接的数学原理", C_ENC, C_BG_ENC, 6.10, 2.00,
     [("公式：", True, C_ENC, 10),
      ("output = LayerNorm(x + SubLayer(x))", False, C_D, 10),
      ("", False, C_D, 6),
      ("梯度分析（为什么梯度不消失）：", True, C_ENC, 10),
      ("Loss对输入x的梯度 = ∂Loss/∂output × (∂SubLayer/∂x + 1)", False, C_D, 9),
      ("关键：加号后面的 \"+1\" 确保即使SubLayer梯度很小，总梯度也有下界1×∂Loss/∂output。", False, C_D, 9),
      ("这就是为什么Transformer能堆叠12层甚至100层（GPT-4据说128层），而不用操心梯度消失。", False, C_D, 9),
      ("对比RNN：梯度要乘以每一步的权重矩阵，T步后≈W^T→指数衰减。", False, C_D, 9)]),
    ("因果掩码实现", C_DEC, C_BG_DEC, 6.40, 2.00,
     [("实现方式：", True, C_DEC, 10),
      ("attn_mask = torch.triu(torch.ones(T,T), diagonal=1) * (-1e9)", False, C_D, 9),
      ("上三角全部设为-10亿（≈-∞），Softmax后≈0。", False, C_D, 9),
      ("", False, C_D, 6),
      ("为什么用-1e9而不是-∞？", True, C_DEC, 10),
      ("浮点数没有真正的-∞，-1e9已经足够大（Softmax(−10⁹)≈0，但不会产生NaN）。", False, C_D, 9),
      ("掩码在训练和推理时都必须严格执行。训练时偷看答案→推理时没答案可抄→崩溃。", False, C_D, 9),
      ("这是保证训练-推理行为一致的关键设计，避免暴露偏差(exposure bias)。", False, C_D, 9)]),
    ("交叉熵损失函数", C_ATT, C_BG_ATT, 6.10, 2.00,
     [("公式：", True, C_ATT, 10),
      ("Loss = -log(p(正确词))", False, C_D, 10),
      ("", False, C_D, 6),
      ("含义：模型对正确词的预测概率越高，损失越小。", True, C_ATT, 10),
      ("例：p(正确词)=0.001 → Loss=6.9（很差，接近随机）", False, C_D, 9),
      ("例：p(正确词)=0.9 → Loss=0.1（很好，很有信心）", False, C_D, 9),
      ("例：p(正确词)=0.5 → Loss=0.7（一般，模棱两可）", False, C_D, 9),
      ("37000个词的输出概率分布中，只有正确词的log被计入损失。标签平滑(ε=0.1)让正确词只占90%，其余10%分给其他词，防止过拟合。", False, C_D, 9)]),
]

for i, (title, color, bgcolor, w, h, lines) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.30 + col * 6.50)
    y = Inches(1.20 + row * 2.10)
    rrect(sl, x, y, Inches(w), Inches(h), fc=C_W, lc=color, lw=1)
    rect(sl, x, y, Inches(w), Inches(0.28), fc=color)
    tb(sl, x+Inches(0.05), y+Inches(0.02), Inches(w-0.10), Inches(0.24), title, fs=10, fc=C_W, b=True)
    mtb(sl, x+Inches(0.10), y+Inches(0.35), Inches(w-0.20), Inches(h-0.45), lines, ls=1.15)

# 底部
rrect(sl, Inches(0.30), Inches(5.55), Inches(12.70), Inches(1.40), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(5.60), Inches(12.30), Inches(0.25), "完整训练循环", fs=11, fc=C_ACC, b=True)
mtb(sl, Inches(0.50), Inches(5.90), Inches(12.30), Inches(0.90), [
    ("前向传播：", True, C_PRI, 9),
    ("输入 - 嵌入+位置编码 - 编码器x6 K,V - 解码器x6 - 线性层 - Softmax - 预测概率分布", False, C_D, 9),
    ("损失计算：", True, C_PRI, 9),
    ("交叉熵 - 预测分布 vs 真实目标 - 标量损失值，对所有位置取平均。", False, C_D, 9),
    ("反向传播：", True, C_PRI, 9),
    ("Adam更新所有6500万参数。残差连接让梯度直达底层，层归一化保持数值稳定。", False, C_D, 9),
], ls=1.1)
chk(sl, "P9")

# ==================== P10 推理阶段 ====================
print("P10: 推理阶段")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "08  推理阶段与KV缓存", "闭卷考试 — 逐词串行生成，但KV缓存避免重复计算")

# 信息卡片
info = [
    ("推理方式", "自回归(Autoregressive)\n每步只生成1个词\n必须等上一步完成\n无法并行", C_DEC),
    ("KV缓存", "缓存编码器K,V\n缓存已生成的K,V\n避免重复计算\n加速5-10倍", C_ENC),
    ("生成策略", "贪心/Beam/采样\n温度参数控制\nTop-k/Top-p过滤\n平衡多样性与质量", C_ATT),
    ("效率瓶颈", "编码器只跑1次\n解码器串行N步\nN=输出长度\nGPU利用率低", C_FFN),
]
for i, (title, body, color) in enumerate(info):
    x = Inches(0.30 + i * 3.20)
    rrect(sl, x, Inches(1.20), Inches(3.00), Inches(1.50), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(3.00), Inches(0.25), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(2.90), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(1.50), Inches(2.80), Inches(1.10), body, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)

# 推理流程 - 垂直步骤
gen_steps = [
    ("Step 1", "[START] → 编码器(1次) → K,V缓存", C_ENC),
    ("Step 2", "[START] → 解码器 → 概率分布 → 选\"我\"", C_DEC),
    ("Step 3", "[START,我] → 解码器(+KV缓存) → 选\"爱\"", C_ATT),
    ("Step 4", "[START,我,爱] → 解码器(+KV缓存) → 选\"你\"", C_FFN),
    ("Step 5", "[START,我,爱,你] → 解码器 → 选\"[END]\" → 停止", C_EMB),
]
for i, (step, desc, color) in enumerate(gen_steps):
    y = Inches(2.85 + i * 0.55)
    rrect(sl, Inches(0.30), y, Inches(1.00), Inches(0.45), fc=color, lc=None)
    tb(sl, Inches(0.32), y+Inches(0.05), Inches(0.96), Inches(0.35), step, fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    rrect(sl, Inches(1.40), y, Inches(5.00), Inches(0.45), fc=C_W, lc=color, lw=1)
    tb(sl, Inches(1.50), y+Inches(0.05), Inches(4.80), Inches(0.35), desc, fs=9, fc=C_D)
    if i < len(gen_steps) - 1:
        arrow_d(sl, Inches(0.70), y+Inches(0.45), Inches(0.18), Inches(0.10), C_G)

# KV缓存详解
rrect(sl, Inches(6.60), Inches(2.85), Inches(6.40), Inches(2.00), fc=C_BG_ENC, lc=C_ENC, lw=1)
tb(sl, Inches(6.80), Inches(2.90), Inches(6.00), Inches(0.25), "KV缓存的数学原理", fs=11, fc=C_ENC, b=True)
mtb(sl, Inches(6.80), Inches(3.20), Inches(6.00), Inches(1.50), [
    ("问题：", True, C_ENC, 9),
    ("第3步解码时，如果重新计算\"我\"和\"爱\"的K,V，就浪费了前面2步的工作。序列越长，浪费越严重。", False, C_D, 9),
    ("解决：", True, C_ENC, 9),
    ("缓存所有已计算过的K和V向量。第3步只需计算新词\"爱\"的Q,K,V，然后将新K,V拼接到缓存上。", False, C_D, 9),
    ("效率：", True, C_ENC, 9),
    ("无缓存：每步O(N²)复杂度 | 有缓存：每步只需O(N)（N=已生成长度）| 实测加速5-10倍。", False, C_D, 9),
    ("代价：", True, C_ENC, 9),
    ("KV缓存占用显存：2×层数×序列长度×头数×维度×2字节。长序列时可能成为瓶颈（FlashAttention优化）。", False, C_D, 9),
], ls=1.1)

# 训练vs推理对比小表
tbl = sl.shapes.add_table(5, 3, Inches(0.30), Inches(5.10), Inches(12.70), Inches(1.70))
table = tbl.table
table.columns[0].width = Inches(2.50)
table.columns[1].width = Inches(5.10)
table.columns[2].width = Inches(5.10)
h_data = [["对比维度", "训练（开卷考试）", "推理（闭卷考试）"],
          ["并行性", "所有位置同时计算，GPU跑满", "逐词串行生成，GPU利用率低"],
          ["输入", "解码器输入=真实答案(教师强制)", "解码器输入=上一步预测(自回归)"],
          ["掩码", "因果掩码(防止偷看后面)", "因果掩码(只看已生成的词)"],
          ["效率", "每批并行处理，速度快", "串行N步，N=输出长度，慢但准确"]]
for j, h in enumerate(h_data[0]):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
for i, row_data in enumerate(h_data[1:]):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W
chk(sl, "P10")

# ==================== P11 词选择策略 ====================
print("P11: 词选择策略")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "09  词选择策略", "从贪心搜索到温度采样 — 平衡确定性与多样性的艺术")

# 5种策略卡片
strategies = [
    ("贪心搜索", "Greedy", "每步选概率最高的词\n确定性强，多样性差\n容易陷入重复循环", "p=0.4\np=0.3\np=0.2 → 选第1个", C_DEC),
    ("束搜索", "Beam Search", "保留top-k条候选路径\nk=4: 同时追踪4个序列\n找到全局最优路径", "4条路径并行\n最终选总概率最高的\n标准翻译方案", C_ENC),
    ("温度采样", "Temperature", "logits/T，T控制随机性\nT→0: 确定性(接近贪心)\nT→∞: 完全随机", "T=0.1: 保守精确\nT=0.7: 平衡(推荐)\nT=1.5: 创意发散", C_ATT),
    ("Top-k采样", "Top-k", "只从概率前k个词中选\nk=50: 从top50中随机\n排除低概率噪声", "1000词中选top50\n在这50个中按概率采样\n平衡多样性和合理性", C_FFN),
    ("Top-p采样", "Nucleus", "从概率累积≥p的词中选\np=0.9: 最少词达到90%概率\n动态词表大小", "p=0.9: 可能只需要3个词\np=0.95: 可能需要10个词\n比Top-k更智能", C_EMB),
]
for i, (title, en, desc, example, color) in enumerate(strategies):
    x = Inches(0.30 + i * 2.58)
    rrect(sl, x, Inches(1.20), Inches(2.35), Inches(3.00), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(2.35), Inches(0.50), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(2.25), Inches(0.22), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.05), Inches(1.42), Inches(2.25), Inches(0.22), en, fs=9, fc=RGBColor(0xFF,0xFF,0xDD), a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(1.80), Inches(2.15), Inches(1.30), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.15)
    rrect(sl, x+Inches(0.10), Inches(3.30), Inches(2.15), Inches(0.80), fc=C_LT, lc=color, lw=0.5)
    tb(sl, x+Inches(0.15), Inches(3.32), Inches(2.05), Inches(0.75), example, fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# 对比表
tbl = sl.shapes.add_table(6, 4, Inches(0.30), Inches(4.40), Inches(12.70), Inches(2.00))
table = tbl.table
table.columns[0].width = Inches(2.50)
table.columns[1].width = Inches(3.40)
table.columns[2].width = Inches(3.40)
table.columns[3].width = Inches(3.40)
h_data = [["策略", "确定性", "多样性", "适用场景"],
          ["贪心搜索", "★★★★★", "☆☆☆☆☆", "翻译、摘要（需要准确）"],
          ["束搜索(k=4)", "★★★★☆", "★★☆☆☆", "机器翻译（标准方案）"],
          ["温度采样(T=0.7)", "★★★☆☆", "★★★★☆", "对话、写作（平衡）"],
          ["Top-k(k=50)", "★★☆☆☆", "★★★★☆", "故事生成（需要创意）"],
          ["Top-p(p=0.9)", "★★★☆☆", "★★★★★", "ChatGPT默认方案"]]
for j, h in enumerate(h_data[0]):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
for i, row_data in enumerate(h_data[1:]):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 底部
rrect(sl, Inches(0.30), Inches(6.55), Inches(12.70), Inches(0.45), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(6.58), Inches(12.30), Inches(0.38),
   "实战建议：翻译任务用束搜索(k=4-6)；对话/创作用Top-p(p=0.9)+温度(T=0.7-1.0)；代码生成用低温度(T=0.2)。ChatGPT/ Claude均采用Top-p采样。",
   fs=9, fc=C_D, ls=1.1)
chk(sl, "P11")

# ==================== P12 训练vs推理全面对比 ====================
print("P12: 训练vs推理全面对比")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "10  训练 vs 推理全面对比", "开卷考试与闭卷考试 — 8个维度的深度对比")

# 大对比表
tbl = sl.shapes.add_table(9, 3, Inches(0.30), Inches(1.20), Inches(12.70), Inches(4.20))
table = tbl.table
table.columns[0].width = Inches(2.50)
table.columns[1].width = Inches(5.10)
table.columns[2].width = Inches(5.10)
h_data = [["对比维度", "训练（开卷考试）", "推理（闭卷考试）"],
          ["核心比喻", "开卷考试：能看到所有答案，同时做完所有题", "闭卷考试：每道题只能看到前面已做的题，逐题作答"],
          ["并行性", "所有位置同时计算，GPU跑满。时间复杂度O(1)", "逐词串行生成，N步。时间复杂度O(N)，N=输出长度"],
          ["解码器输入", "真实目标序列（教师强制）：[START,我,爱,你,END]", "上一步预测序列（自回归）：[START]→我→[START,我]→爱"],
          ["编码器", "只运行1次，输出(K,V)供所有解码器步骤使用", "同样只运行1次，输出(K,V)缓存供所有步骤复用"],
          ["损失函数", "交叉熵 Loss=-log(p(正确词))，反向传播更新参数", "无损失计算，参数冻结，只做前向传播"],
          ["输出方式", "一次输出整个序列的概率分布", "每步输出1个词的概率分布，逐步拼接"],
          ["效率瓶颈", "GPU显存（需要存所有中间结果做反向传播）", "生成速度（串行N步，每步依赖上一步）"],
          ["KV缓存", "不需要（每次全量计算，训练阶段没有\"已生成\"的概念）", "必须（缓存已计算的K,V，避免重复，加速5-10倍）"]]
for j, h in enumerate(h_data[0]):
    cell = table.cell(0, j); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_W; p.font.name = FC
    cell.fill.solid(); cell.fill.fore_color.rgb = C_PRI
for i, row_data in enumerate(h_data[1:]):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.color.rgb = C_D; p.font.name = FC
            if j == 1: p.font.color.rgb = C_ENC
            elif j == 2: p.font.color.rgb = C_DEC
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_LT if i % 2 == 0 else C_W

# 关键洞察
rrect(sl, Inches(0.30), Inches(5.55), Inches(6.10), Inches(1.40), fc=C_BG_ENC, lc=C_ENC, lw=1)
tb(sl, Inches(0.50), Inches(5.60), Inches(5.70), Inches(0.25), "训练的核心洞察", fs=11, fc=C_ENC, b=True)
mtb(sl, Inches(0.50), Inches(5.90), Inches(5.70), Inches(0.90), [
    "训练的本质：给模型看海量\"输入→输出\"样本，让它通过反向传播逐步调整6500万参数。",
    "教师强制让GPU并行跑满（所有位置同时计算），但代价是训练-推理行为不一致（暴露偏差）。",
    "解决方案：标签平滑+课程学习+Scheduled Sampling（逐步减少教师强制比例）。",
], ls=1.15)

rrect(sl, Inches(6.60), Inches(5.55), Inches(6.40), Inches(1.40), fc=C_BG_DEC, lc=C_DEC, lw=1)
tb(sl, Inches(6.80), Inches(5.60), Inches(6.00), Inches(0.25), "推理的核心洞察", fs=11, fc=C_DEC, b=True)
mtb(sl, Inches(6.80), Inches(5.90), Inches(6.00), Inches(0.90), [
    "推理的本质：模型\"闭卷考试\"，每步只看已生成的内容，预测下一个词，逐步拼出完整输出。",
    "KV缓存是推理效率的关键——缓存编码器和已生成词的K,V，避免重复计算，加速5-10倍。",
    "量化(16bit→8bit→4bit)和FlashAttention进一步优化，让大模型在消费级GPU上也能流畅运行。",
], ls=1.15)
chk(sl, "P12")

# ==================== P13 应用与总结 ====================
print("P13: 应用与总结")
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "11  应用与总结", "三大家族与AI通用引擎 — Transformer改变了整个AI领域")

# 三大家族
families = [
    ("编码器家族\nEncoder-only", "BERT, RoBERTa\nDeBERTa, ALBERT", C_ENC, C_BG_ENC,
     "双向注意力，理解上下文\n文本分类、命名实体识别\n阅读理解、情感分析\n预训练+微调范式"),
    ("解码器家族\nDecoder-only", "GPT系列, LLaMA\nClaude, Gemini, Qwen", C_DEC, C_BG_DEC,
     "单向注意力，文本生成\nChatGPT就是GPT+RLHF\n代码生成、创意写作\n自回归生成范式"),
    ("编码器-解码器\nEncoder-Decoder", "T5, BART, mBART\nTransformer原版", C_ATT, C_BG_ATT,
     "两种注意力兼备\n翻译、摘要、问答\nSeq2Seq任务首选\n编码理解+解码生成"),
]
for i, (title, models, color, bgcolor, desc) in enumerate(families):
    x = Inches(0.30 + i * 4.30)
    rrect(sl, x, Inches(1.20), Inches(4.00), Inches(2.50), fc=C_W, lc=color, lw=1)
    rect(sl, x, Inches(1.20), Inches(4.00), Inches(0.50), fc=color)
    tb(sl, x+Inches(0.05), Inches(1.22), Inches(3.90), Inches(0.45), title, fs=10, fc=C_W, b=True, a=PP_ALIGN.CENTER)
    rrect(sl, x+Inches(0.10), Inches(1.80), Inches(3.80), Inches(0.50), fc=bgcolor, lc=color, lw=0.5)
    tb(sl, x+Inches(0.15), Inches(1.82), Inches(3.70), Inches(0.45), models, fs=9, fc=color, b=True, a=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.10), Inches(2.40), Inches(3.80), Inches(1.20), desc, fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.2)

# 影响力数据
rrect(sl, Inches(0.30), Inches(3.90), Inches(12.70), Inches(1.30), fc=C_LT, lc=C_ACC, lw=1)
tb(sl, Inches(0.50), Inches(3.95), Inches(12.30), Inches(0.25), "Transformer的影响力", fs=11, fc=C_ACC, b=True)
impact = [
    ("2017", "Transformer论文发表，\"Attention Is All You Need\""),
    ("2018", "BERT(编码器)和GPT-1(解码器)诞生，预训练时代开启"),
    ("2019", "GPT-2(15亿参数)展现强大的零样本生成能力"),
    ("2020", "GPT-3(1750亿参数)\"少样本学习\"震惊世界"),
    ("2022", "ChatGPT(GPT-3.5+RLHF)引爆AI应用，用户破亿"),
    ("2024+", "GPT-4/Claude/Gemini/Qwen突破万亿参数，多模态融合"),
]
for i, (year, event) in enumerate(impact):
    col = i % 2
    row = i // 2
    x = Inches(0.50 + col * 6.30)
    y = Inches(4.25 + row * 0.30)
    tb(sl, x, y, Inches(0.50), Inches(0.25), year, fs=9, fc=C_ACC, b=True)
    tb(sl, x+Inches(0.60), y, Inches(5.60), Inches(0.25), event, fs=9, fc=C_D)

# 总结
rrect(sl, Inches(0.30), Inches(5.35), Inches(12.70), Inches(1.60), fc=C_W, lc=C_PRI, lw=2)
rect(sl, Inches(0.30), Inches(5.35), Inches(12.70), Inches(0.30), fc=C_PRI)
tb(sl, Inches(0.50), Inches(5.37), Inches(12.30), Inches(0.26), "核心总结", fs=12, fc=C_W, b=True)
mtb(sl, Inches(0.50), Inches(5.70), Inches(12.30), Inches(1.10), [
    ("1. 核心创新：", True, C_PRI, 10),
    ("自注意力机制(Q×K^T/√d_k×V)让每个词都能直接与所有词交互，完全并行，解决了RNN的三大根本问题。", False, C_D, 9),
    ("2. 架构精髓：", True, C_PRI, 10),
    ("残差连接保证梯度不消失，层归一化保持数值稳定，多头注意力捕捉多种关系，FFN提取深层特征。", False, C_D, 9),
    ("3. 训练-推理差异：", True, C_PRI, 10),
    ("训练=开卷考试(教师强制，并行高效)；推理=闭卷考试(自回归，串行但KV缓存加速)。", False, C_D, 9),
    ("4. 历史地位：", True, C_PRI, 10),
    ("从2017年一篇论文到2026年万亿参数大模型，Transformer不仅是NLP的基础，更扩展到CV、语音、多模态，成为AI的\"通用引擎\"。", False, C_D, 9),
], ls=1.1)
chk(sl, "P13")

prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
