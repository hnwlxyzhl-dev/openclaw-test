#!/usr/bin/env python3
"""Transformer PPT v5 - Pages 8-13 (全中文版 - 高中文占比)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import re

# === 常量 ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 颜色
C_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
C_ENCODER = RGBColor(0x34, 0x98, 0xDB)
C_DECODER = RGBColor(0xE7, 0x4C, 0x3C)
C_ATTN = RGBColor(0xF3, 0x9C, 0x12)
C_FFN = RGBColor(0x2E, 0xCC, 0x71)
C_EMBED = RGBColor(0x9B, 0x59, 0xB6)
C_TEXT = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY = RGBColor(0x95, 0xA5, 0xA6)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
C_LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
C_LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xD0)
C_LIGHT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
C_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_BLUE = RGBColor(0x34, 0x98, 0xDB)
C_DARK_BLUE = RGBColor(0x1A, 0x25, 0x2F)

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"
FONT_MONO = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


def add_shape(slide, left, top, width, height, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=11, font_name=FONT_BODY, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_title(slide, text, top=Inches(0.15)):
    add_text_box(slide, Inches(0.3), top, Inches(12.7), Inches(0.45), text, 18, FONT_TITLE, C_TITLE, True)


def add_arrow_line(slide, x1, y1, x2, y2, color=C_TEXT, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector._element.spPr.find(qn('a:ln'))
    if ln is not None:
        tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tailEnd)
    return connector


def add_line(slide, x1, y1, x2, y2, color=C_TEXT, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_bullet_text(slide, left, top, width, height, items, font_size=11, color=C_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(3)
        p.line_spacing = Pt(font_size * 1.4)
        
        if "：" in item and item.index("：") < 25:
            colon_pos = item.index("：")
            title_part = item[:colon_pos]
            body_part = item[colon_pos:]
            run1 = p.add_run()
            run1.text = "● " + title_part + "："
            run1.font.size = Pt(font_size)
            run1.font.name = FONT_BODY
            run1.font.color.rgb = C_TITLE
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = body_part
            run2.font.size = Pt(font_size)
            run2.font.name = FONT_BODY
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = "● " + item
            run.font.size = Pt(font_size)
            run.font.name = FONT_BODY
            run.font.color.rgb = color
    return txBox


# ===========================
# 第8页：训练细节
# ===========================
def create_slide_8():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "训练细节：安全网、掩码与损失")
    
    left_x = Inches(0.3)
    left_w = Inches(7.8)
    
    # --- 因果掩码矩阵 ---
    mask_top = Inches(0.75)
    add_text_box(slide, left_x, mask_top, left_w, Inches(0.3), "因果掩码矩阵：下三角可见，上三角屏蔽", 13, FONT_BODY, C_TITLE, True)
    
    labels = ["起始", "我", "爱", "你"]
    cell_size = Inches(0.75)
    grid_start_x = Inches(1.5)
    grid_start_y = mask_top + Inches(0.38)
    
    # 列标签
    for j, lbl in enumerate(labels):
        add_text_box(slide, grid_start_x + j * cell_size + Inches(0.05), grid_start_y, cell_size, Inches(0.25), lbl, 10, FONT_BODY, C_TITLE, True, PP_ALIGN.CENTER)
    
    # 矩阵
    for i in range(4):
        add_text_box(slide, grid_start_x - Inches(0.7), grid_start_y + Inches(0.25) + i * cell_size + Inches(0.08), Inches(0.65), Inches(0.25), labels[i], 10, FONT_BODY, C_TITLE, True, PP_ALIGN.RIGHT)
        for j in range(4):
            cx = grid_start_x + j * cell_size
            cy = grid_start_y + Inches(0.28) + i * cell_size
            if j <= i:
                add_shape(slide, cx, cy, cell_size - Inches(0.03), cell_size - Inches(0.03), MSO_SHAPE.ROUNDED_RECTANGLE, C_ENCODER, C_ENCODER, Pt(0.5))
                add_text_box(slide, cx, cy + Inches(0.15), cell_size, Inches(0.25), "可见", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
            else:
                add_shape(slide, cx, cy, cell_size - Inches(0.03), cell_size - Inches(0.03), MSO_SHAPE.ROUNDED_RECTANGLE, C_RED, C_RED, Pt(0.5))
                add_text_box(slide, cx, cy + Inches(0.1), cell_size, Inches(0.25), "屏蔽", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    # 图例
    leg_y = grid_start_y + Inches(0.28) + 4 * cell_size + Inches(0.08)
    add_shape(slide, grid_start_x, leg_y, Inches(0.25), Inches(0.2), MSO_SHAPE.ROUNDED_RECTANGLE, C_ENCODER)
    add_text_box(slide, grid_start_x + Inches(0.3), leg_y - Inches(0.02), Inches(1.5), Inches(0.25), "= 可以看到（下三角）", 9, FONT_BODY, C_TEXT)
    add_shape(slide, grid_start_x + Inches(2.5), leg_y, Inches(0.25), Inches(0.2), MSO_SHAPE.ROUNDED_RECTANGLE, C_RED)
    add_text_box(slide, grid_start_x + Inches(2.8), leg_y - Inches(0.02), Inches(2.5), Inches(0.25), "= 被屏蔽（上三角设为负无穷）", 9, FONT_BODY, C_TEXT)
    
    # --- 损失计算流程 ---
    flow_top = Inches(4.15)
    add_text_box(slide, left_x, flow_top, left_w, Inches(0.3), "损失计算流程", 13, FONT_BODY, C_TITLE, True)
    
    flow_y = flow_top + Inches(0.4)
    box_w = Inches(1.1)
    box_h = Inches(0.45)
    arrow_len = Inches(0.3)
    
    flow_items = [
        ("解码器输出", C_ENCODER),
        ("线性变换", C_EMBED),
        ("归一化", C_ATTN),
        ("预测概率", C_FFN),
    ]
    
    start_x = left_x + Inches(0.1)
    for idx, (txt, clr) in enumerate(flow_items):
        bx = start_x + idx * (box_w + arrow_len)
        add_shape(slide, bx, flow_y, box_w, box_h, MSO_SHAPE.ROUNDED_RECTANGLE, clr, clr, Pt(1))
        add_text_box(slide, bx, flow_y + Inches(0.08), box_w, box_h - Inches(0.1), txt, 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
        if idx < len(flow_items) - 1:
            add_arrow_line(slide, bx + box_w, flow_y + box_h / 2, bx + box_w + arrow_len, flow_y + box_h / 2, C_TEXT, Pt(2))
    
    # 第二行
    flow_y2 = flow_y + box_h + Inches(0.4)
    tgt_x = start_x + 2 * (box_w + arrow_len)
    add_shape(slide, tgt_x, flow_y2, box_w, box_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_ATTN, C_ATTN, Pt(1))
    add_text_box(slide, tgt_x, flow_y2 + Inches(0.08), box_w, box_h - Inches(0.1), "目标词", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    add_text_box(slide, tgt_x + box_w + Inches(0.02), flow_y2 + Inches(0.08), Inches(0.35), box_h, "↔", 12, FONT_MONO, C_RED, True, PP_ALIGN.CENTER)
    
    loss_x = start_x + 3 * (box_w + arrow_len) + Inches(0.35)
    add_shape(slide, loss_x, flow_y2, box_w + Inches(0.3), box_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_RED, C_RED, Pt(1))
    add_text_box(slide, loss_x, flow_y2 + Inches(0.02), box_w + Inches(0.3), box_h, "交叉熵损失", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    # 第三行
    flow_y3 = flow_y2 + box_h + Inches(0.35)
    bp_items = [("反向传播", C_DECODER), ("更新参数", C_FFN)]
    bp_start = start_x + 2 * (box_w + arrow_len)
    for idx, (txt, clr) in enumerate(bp_items):
        bx = bp_start + idx * (box_w + arrow_len)
        add_shape(slide, bx, flow_y3, box_w, box_h, MSO_SHAPE.ROUNDED_RECTANGLE, clr, clr, Pt(1))
        add_text_box(slide, bx, flow_y3 + Inches(0.08), box_w, box_h - Inches(0.1), txt, 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
        if idx < len(bp_items) - 1:
            add_arrow_line(slide, bx + box_w, flow_y3 + box_h / 2, bx + box_w + arrow_len, flow_y3 + box_h / 2, C_TEXT, Pt(2))
    
    # 连接
    add_arrow_line(slide, start_x + 3 * (box_w + arrow_len) + box_w / 2, flow_y + box_h, start_x + 3 * (box_w + arrow_len) + box_w / 2, flow_y2, C_TEXT, Pt(1.5))
    
    # === 右侧40% ===
    right_x = Inches(8.5)
    right_w = Inches(4.5)
    
    items = [
        "残差连接：输出等于输入加上子层处理后的输入。梯度流中永远有一个加一的通路，确保梯度永不消失，因此可以训练超过一百层的深度网络",
        "因果掩码：将注意力矩阵上三角设为负无穷大，经过归一化后变为零。这样训练和推理的行为完全一致，杜绝了信息泄露导致的暴露偏差问题",
        "交叉熵损失：损失值等于正确词预测概率的负对数。预测越准确损失越小，预测越差损失越大，是分类任务的标准损失函数",
        "训练配置：使用自适应矩估计优化器，学习率万分之三，随机丢弃比例百分之三十，标签平滑系数零点一，八块显卡训练十二小时",
    ]
    add_bullet_text(slide, right_x, Inches(0.8), right_w, Inches(6.2), items, 11, C_TEXT)


# ===========================
# 第9页：推理+KV缓存
# ===========================
def create_slide_9():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "推理阶段：闭卷考试与键值缓存")
    
    left_x = Inches(0.3)
    left_w = Inches(7.8)
    
    # --- 自回归步骤图 ---
    step_top = Inches(0.8)
    add_text_box(slide, left_x, step_top, left_w, Inches(0.3), "自回归生成：逐词预测，每步依赖前面所有词", 13, FONT_BODY, C_TITLE, True)
    
    # 编码器
    enc_x = left_x + Inches(0.1)
    enc_y = step_top + Inches(0.35)
    add_shape(slide, enc_x, enc_y, Inches(2.2), Inches(0.5), MSO_SHAPE.ROUNDED_RECTANGLE, C_ENCODER, C_ENCODER, Pt(1))
    add_text_box(slide, enc_x, enc_y + Inches(0.05), Inches(2.2), Inches(0.4), "编码器（只运行一次后缓存）", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    # 步骤
    steps = [
        ("第一步", "[起始]", "输出\"我\"", C_ENCODER),
        ("第二步", "[起始,我]", "输出\"爱\"", C_ATTN),
        ("第三步", "[起始,我,爱]", "输出\"你\"", C_FFN),
        ("第四步", "[起始,我,爱,你]", "输出\"结束\"", C_GREEN),
    ]
    
    step_y = enc_y + Inches(0.7)
    step_box_w = Inches(1.5)
    step_arrow_w = Inches(0.25)
    
    for idx, (step_name, inp, out, clr) in enumerate(steps):
        sx = left_x + Inches(0.1) + idx * (step_box_w + step_arrow_w)
        add_text_box(slide, sx, step_y - Inches(0.15), step_box_w, Inches(0.2), step_name, 8, FONT_BODY, C_TITLE, True, PP_ALIGN.CENTER)
        add_shape(slide, sx, step_y + Inches(0.05), step_box_w, Inches(0.35), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_BLUE, C_ENCODER, Pt(0.5))
        add_text_box(slide, sx, step_y + Inches(0.08), step_box_w, Inches(0.3), inp, 8, FONT_BODY, C_TEXT, False, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_arrow_line(slide, sx + step_box_w, step_y + Inches(0.22), sx + step_box_w + step_arrow_w, step_y + Inches(0.22), C_TEXT, Pt(1.5))
        add_shape(slide, sx, step_y + Inches(0.5), step_box_w, Inches(0.35), MSO_SHAPE.ROUNDED_RECTANGLE, clr, clr, Pt(1))
        add_text_box(slide, sx, step_y + Inches(0.53), step_box_w, Inches(0.3), out, 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    # 解码器标签
    dec_x = left_x + Inches(0.1)
    dec_y = step_y + Inches(0.95)
    add_shape(slide, dec_x, dec_y, Inches(4.2), Inches(0.35), MSO_SHAPE.ROUNDED_RECTANGLE, C_DECODER, C_DECODER, Pt(1))
    add_text_box(slide, dec_x, dec_y + Inches(0.03), Inches(4.2), Inches(0.3), "解码器（每个步骤运行一次）", 9, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    # --- 无缓存 vs 有缓存 ---
    compare_top = Inches(4.55)
    add_text_box(slide, left_x, compare_top, left_w, Inches(0.3), "键值缓存对比：节省重复计算", 13, FONT_BODY, C_TITLE, True)
    
    nocache_x = left_x + Inches(0.2)
    nocache_y = compare_top + Inches(0.35)
    nocache_w = Inches(3.5)
    add_shape(slide, nocache_x, nocache_y, nocache_w, Inches(2.2), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_RED, C_RED, Pt(1))
    add_text_box(slide, nocache_x, nocache_y + Inches(0.05), nocache_w, Inches(0.3), "无缓存（每步重新计算全部键值）", 10, FONT_BODY, C_RED, True, PP_ALIGN.CENTER)
    
    nc_steps = ["第一步：计算第一个词的键和值", "第二步：重新计算第一个词+计算第二个词", "第三步：重新计算前两个词+计算第三个词", "时间复杂度：平方级，浪费严重"]
    for si, st in enumerate(nc_steps):
        add_text_box(slide, nocache_x + Inches(0.15), nocache_y + Inches(0.4) + si * Inches(0.35), nocache_w - Inches(0.3), Inches(0.3), st, 8, FONT_BODY, C_TEXT, si == 3, PP_ALIGN.LEFT)
    
    add_text_box(slide, nocache_x + Inches(0.5), nocache_y + Inches(1.85), Inches(2.5), Inches(0.25), "重复计算，效率低下", 10, FONT_BODY, C_RED, True, PP_ALIGN.CENTER)
    
    cache_x = left_x + Inches(4.0)
    cache_y = compare_top + Inches(0.35)
    cache_w = Inches(3.5)
    add_shape(slide, cache_x, cache_y, cache_w, Inches(2.2), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_GREEN, C_GREEN, Pt(1))
    add_text_box(slide, cache_x, cache_y + Inches(0.05), cache_w, Inches(0.3), "有缓存（只计算新词的键值并追加）", 10, FONT_BODY, C_GREEN, True, PP_ALIGN.CENTER)
    
    ca_steps = ["第一步：计算第一个词的键和值，存入缓存", "第二步：只计算新词键值，追加到缓存", "第三步：同样只算新词，缓存自动累积", "时间复杂度：线性级，内存换速度"]
    for si, st in enumerate(ca_steps):
        add_text_box(slide, cache_x + Inches(0.15), cache_y + Inches(0.4) + si * Inches(0.35), cache_w - Inches(0.3), Inches(0.3), st, 8, FONT_BODY, C_TEXT, si == 3, PP_ALIGN.LEFT)
    
    add_text_box(slide, cache_x + Inches(0.3), cache_y + Inches(1.85), Inches(2.9), Inches(0.25), "只算新词，高效节省", 10, FONT_BODY, C_GREEN, True, PP_ALIGN.CENTER)
    
    # === 右侧40% ===
    right_x = Inches(8.5)
    right_w = Inches(4.5)
    
    items = [
        "推理就像闭卷考试：模型必须逐词串行生成，每一步都依赖前面已经生成的所有词。如果第一个词预测错误，后续所有词都可能跟着出错",
        "编码器只需运行一次：将输入编码后缓存结果，整个文本生成过程中反复使用这些编码，与最终生成的文本长度完全无关",
        "串行瓶颈：生成一百个词需要一百次前向传播，而训练阶段一次前向传播就能处理所有位置，因此推理速度远慢于训练速度",
        "键值缓存如同会议纪要：每步只计算新产生词的键和值并追加到缓存中，将复杂度从平方级降低到线性级，代价是占用更多显存",
    ]
    add_bullet_text(slide, right_x, Inches(0.8), right_w, Inches(6.2), items, 11, C_TEXT)


# ===========================
# 第10页：词选择策略
# ===========================
def create_slide_10():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "词选择：从贪心搜索到创意生成")
    
    left_x = Inches(0.3)
    left_w = Inches(7.2)
    card_top = Inches(0.75)
    card_h = Inches(1.05)
    card_gap = Inches(0.12)
    
    strategies = [
        ("1. 贪心搜索", "每一步都选择概率最高的词，结果确定但可能陷入循环重复", C_LIGHT_GRAY, C_TEXT),
        ("2. 束搜索", "同时保留多条候选路径，最终选总概率最高的，翻译质量好但偏向平庸", C_LIGHT_BLUE, C_ENCODER),
        ("3. 前K采样", "只从前K个高概率词中随机选择，兼顾生成质量和多样性", C_LIGHT_GREEN, C_FFN),
        ("4. 核心采样", "选取累计概率达到阈值的最小词集，比前K采样更灵活自适应", C_LIGHT_ORANGE, C_ATTN),
        ("5. 温度调节", "在归一化前对概率分布进行缩放，低温更确定，高温更随机", C_LIGHT_RED, C_RED),
    ]
    
    for idx, (title, desc, bg, accent) in enumerate(strategies):
        cy = card_top + idx * (card_h + card_gap)
        add_shape(slide, left_x, cy, left_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE, bg, accent, Pt(1.5))
        add_shape(slide, left_x, cy, Inches(0.08), card_h, MSO_SHAPE.RECTANGLE, accent)
        add_text_box(slide, left_x + Inches(0.2), cy + Inches(0.08), left_w - Inches(0.4), Inches(0.3), title, 11, FONT_BODY, accent, True)
        add_text_box(slide, left_x + Inches(0.2), cy + Inches(0.42), left_w - Inches(0.4), Inches(0.55), desc, 10, FONT_BODY, C_TEXT)
    
    # === 右侧45% ===
    right_x = Inches(7.8)
    right_w = Inches(5.2)
    
    items = [
        "主流组合策略：温度设为点七加上核心采样设为零点九并配合重复惩罚机制。写代码时温度设为零点二确保精确，创意写作时温度设为零点七到一点零增加多样性",
        "贪心的陷阱：局部最优不等于全局最优。例如我喜欢在公园里这个句子，贪心可能选择散步这个最常见的词，但从整体来看跑步加之后喝杯咖啡的联合概率可能更高",
    ]
    add_bullet_text(slide, right_x, Inches(0.8), right_w, Inches(2.8), items, 11, C_TEXT)
    
    # === 右下：温度效果示意 ===
    curve_top = Inches(3.8)
    add_text_box(slide, right_x, curve_top, right_w, Inches(0.3), "温度对概率分布的影响", 12, FONT_BODY, C_TITLE, True, PP_ALIGN.CENTER)
    
    dist_y = curve_top + Inches(0.4)
    dist_w = Inches(5.0)
    dist_h = Inches(2.8)
    add_shape(slide, right_x + Inches(0.1), dist_y, dist_w, dist_h, MSO_SHAPE.ROUNDED_RECTANGLE, RGBColor(0xFA, 0xFA, 0xFA), C_GRAY, Pt(0.5))
    
    # 低温
    add_text_box(slide, right_x + Inches(0.3), dist_y + Inches(0.15), Inches(1.3), Inches(0.25), "低温 尖锐集中", 9, FONT_BODY, C_BLUE, True)
    add_shape(slide, right_x + Inches(0.3), dist_y + Inches(0.5), Inches(0.1), Inches(0.6), MSO_SHAPE.RECTANGLE, C_BLUE)
    add_text_box(slide, right_x + Inches(0.55), dist_y + Inches(0.6), Inches(2.0), Inches(0.5), "高概率词概率接近一\n其他词概率接近零\n输出非常确定和重复", 9, FONT_BODY, C_TEXT)
    
    # 中温
    add_text_box(slide, right_x + Inches(0.3), dist_y + Inches(1.2), Inches(1.3), Inches(0.25), "中温 自然分布", 9, FONT_BODY, C_GREEN, True)
    add_shape(slide, right_x + Inches(0.3), dist_y + Inches(1.55), Inches(0.25), Inches(0.45), MSO_SHAPE.RECTANGLE, C_GREEN)
    add_text_box(slide, right_x + Inches(0.65), dist_y + Inches(1.6), Inches(2.0), Inches(0.5), "概率分布较为自然\n高概率词仍然优先\n质量和多样性兼顾", 9, FONT_BODY, C_TEXT)
    
    # 高温
    add_text_box(slide, right_x + Inches(0.3), dist_y + Inches(2.2), Inches(1.3), Inches(0.25), "高温 均匀平坦", 9, FONT_BODY, C_RED, True)
    add_shape(slide, right_x + Inches(0.3), dist_y + Inches(2.5), Inches(0.6), Inches(0.15), MSO_SHAPE.RECTANGLE, C_RED)
    add_text_box(slide, right_x + Inches(1.0), dist_y + Inches(2.35), Inches(2.0), Inches(0.5), "所有词概率趋于均匀\n低概率词也可能被选中\n输出极具创意但可能不通顺", 9, FONT_BODY, C_TEXT)


# ===========================
# 第11页：训练vs推理对比
# ===========================
def create_slide_11():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "训练与推理：开卷考试与闭卷考试")
    
    table_left = Inches(0.4)
    table_top = Inches(0.75)
    col_widths = [Inches(2.2), Inches(4.8), Inches(4.8)]
    row_height = Inches(0.62)
    
    rows_data = [
        ("对比维度", "训练阶段（开卷考试）", "推理阶段（闭卷考试）"),
        ("解码器输入", "完整的目标序列（教师强制输入）", "模型自己逐词生成并作为下一步输入"),
        ("计算方式", "所有位置同时并行计算（一次前向）", "逐词串行计算（生成多少词算多少次）"),
        ("损失函数", "每个位置都计算交叉熵损失", "不计算损失，只做概率预测"),
        ("梯度计算", "通过反向传播更新所有参数", "不计算梯度，参数全部冻结"),
        ("键值缓存", "不需要（并行计算无需缓存）", "必须使用（避免重复计算）"),
        ("显卡利用率", "高（矩阵运算可以充分并行）", "低（串行生成导致大量等待）"),
        ("核心目标", "学习参数（让模型学会预测）", "生成结果（产出最终文本）"),
    ]
    
    for ri, (dim, train, infer) in enumerate(rows_data):
        y = table_top + ri * row_height
        if ri == 0:
            for ci, (txt, w) in enumerate(zip([dim, train, infer], col_widths)):
                x = table_left + sum(cw for cw in [Inches(0)] + col_widths[:ci])
                add_shape(slide, x, y, w - Inches(0.04), row_height - Inches(0.04), MSO_SHAPE.RECTANGLE, C_TITLE)
                add_text_box(slide, x, y + Inches(0.1), w - Inches(0.04), row_height - Inches(0.15), txt, 12, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
        else:
            add_shape(slide, table_left, y, col_widths[0] - Inches(0.04), row_height - Inches(0.04), MSO_SHAPE.RECTANGLE, C_LIGHT_GRAY, C_GRAY, Pt(0.5))
            add_text_box(slide, table_left + Inches(0.1), y + Inches(0.1), col_widths[0] - Inches(0.2), row_height - Inches(0.15), dim, 11, FONT_BODY, C_TITLE, True, PP_ALIGN.CENTER)
            
            train_bg = C_LIGHT_BLUE if ri % 2 == 1 else C_WHITE
            add_shape(slide, table_left + col_widths[0], y, col_widths[1] - Inches(0.04), row_height - Inches(0.04), MSO_SHAPE.RECTANGLE, train_bg, C_ENCODER, Pt(0.3))
            add_text_box(slide, table_left + col_widths[0] + Inches(0.15), y + Inches(0.1), col_widths[1] - Inches(0.3), row_height - Inches(0.15), train, 11, FONT_BODY, C_TEXT, False, PP_ALIGN.CENTER)
            
            infer_bg = C_LIGHT_GREEN if ri % 2 == 1 else C_WHITE
            add_shape(slide, table_left + col_widths[0] + col_widths[1], y, col_widths[2] - Inches(0.04), row_height - Inches(0.04), MSO_SHAPE.RECTANGLE, infer_bg, C_FFN, Pt(0.3))
            add_text_box(slide, table_left + col_widths[0] + col_widths[1] + Inches(0.15), y + Inches(0.1), col_widths[2] - Inches(0.3), row_height - Inches(0.15), infer, 11, FONT_BODY, C_TEXT, False, PP_ALIGN.CENTER)
    
    # 底部比喻
    summary_y = table_top + len(rows_data) * row_height + Inches(0.2)
    add_shape(slide, Inches(1.5), summary_y, Inches(10.3), Inches(0.5), MSO_SHAPE.ROUNDED_RECTANGLE, RGBColor(0xEB, 0xF5, 0xFB), C_ENCODER, Pt(1))
    add_text_box(slide, Inches(1.7), summary_y + Inches(0.05), Inches(9.9), Inches(0.4),
                 "训练就像新员工培训（有标准答案可以参考），推理就像正式上岗（全靠自己）。核心矛盾是暴露偏差。",
                 12, FONT_BODY, C_TITLE, True, PP_ALIGN.CENTER)


# ===========================
# 第12页：应用
# ===========================
def create_slide_12():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "从翻译工具到人工智能通用引擎")
    
    left_x = Inches(0.3)
    left_w = Inches(7.8)
    
    # 顶部
    top_box_w = Inches(3.0)
    top_box_h = Inches(0.55)
    top_x = Inches(3.9)
    top_y = Inches(0.8)
    
    add_shape(slide, top_x, top_y, top_box_w, top_box_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_DARK_BLUE, C_DARK_BLUE, Pt(1.5))
    add_text_box(slide, top_x, top_y + Inches(0.08), top_box_w, top_box_h, "变换器模型（二零一七）", 12, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    
    branch_y = top_y + top_box_h + Inches(0.5)
    branch_h = Inches(0.4)
    
    # 左分支
    lb_x = Inches(0.3)
    add_shape(slide, lb_x, branch_y, Inches(2.5), branch_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_ENCODER, C_ENCODER, Pt(1))
    add_text_box(slide, lb_x, branch_y + Inches(0.05), Inches(2.5), branch_h, "仅解码器架构", 10, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    add_line(slide, top_x + top_box_w / 2, top_y + top_box_h, lb_x + Inches(1.25), branch_y, C_ENCODER, Pt(2))
    
    gpt_y = branch_y + branch_h + Inches(0.15)
    gpt_models = ["一代", "二代", "三代", "四代", "对话版"]
    gpt_prefix = ["生成预训练", "生成预训练", "生成预训练", "生成预训练", "生成预训练"]
    for gi, gm in enumerate(gpt_models):
        gx = lb_x + gi * Inches(0.55)
        add_shape(slide, gx, gpt_y, Inches(0.5), Inches(0.3), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_BLUE, C_ENCODER, Pt(0.5))
        add_text_box(slide, gx, gpt_y + Inches(0.02), Inches(0.5), Inches(0.26), gm, 7, FONT_BODY, C_TEXT, True, PP_ALIGN.CENTER)
        if gi < len(gpt_models) - 1:
            add_arrow_line(slide, gx + Inches(0.5), gpt_y + Inches(0.15), gx + Inches(0.55), gpt_y + Inches(0.15), C_ENCODER, Pt(1))
    add_text_box(slide, lb_x, gpt_y + Inches(0.35), Inches(2.5), Inches(0.25), "生成之王", 10, FONT_BODY, C_ENCODER, True, PP_ALIGN.CENTER)
    
    # 中分支
    mb_x = Inches(3.5)
    add_shape(slide, mb_x, branch_y, Inches(2.5), branch_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_FFN, C_FFN, Pt(1))
    add_text_box(slide, mb_x, branch_y + Inches(0.05), Inches(2.5), branch_h, "完整编码解码架构", 10, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    add_line(slide, top_x + top_box_w / 2, top_y + top_box_h, mb_x + Inches(1.25), branch_y, C_FFN, Pt(2))
    
    t5_models = ["文本转换", "序列转换", "语音识别"]
    for ti, tm in enumerate(t5_models):
        tx = mb_x + ti * Inches(0.85)
        add_shape(slide, tx, gpt_y, Inches(0.8), Inches(0.3), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_GREEN, C_FFN, Pt(0.5))
        add_text_box(slide, tx, gpt_y + Inches(0.02), Inches(0.8), Inches(0.26), tm, 7, FONT_BODY, C_TEXT, True, PP_ALIGN.CENTER)
        if ti < len(t5_models) - 1:
            add_arrow_line(slide, tx + Inches(0.8), gpt_y + Inches(0.15), tx + Inches(0.85), gpt_y + Inches(0.15), C_FFN, Pt(1))
    add_text_box(slide, mb_x, gpt_y + Inches(0.35), Inches(2.5), Inches(0.25), "转换之王", 10, FONT_BODY, C_FFN, True, PP_ALIGN.CENTER)
    
    # 右分支
    rb_x = Inches(6.7)
    add_shape(slide, rb_x, branch_y, Inches(2.5), branch_h, MSO_SHAPE.ROUNDED_RECTANGLE, C_EMBED, C_EMBED, Pt(1))
    add_text_box(slide, rb_x, branch_y + Inches(0.05), Inches(2.5), branch_h, "仅编码器架构", 10, FONT_BODY, C_WHITE, True, PP_ALIGN.CENTER)
    add_line(slide, top_x + top_box_w / 2, top_y + top_box_h, rb_x + Inches(1.25), branch_y, C_EMBED, Pt(2))
    
    bert_models = ["语言理解", "增强版", "解构版"]
    for bi, bm in enumerate(bert_models):
        bx = rb_x + bi * Inches(0.85)
        add_shape(slide, bx, gpt_y, Inches(0.8), Inches(0.3), MSO_SHAPE.ROUNDED_RECTANGLE, C_LIGHT_PURPLE, C_EMBED, Pt(0.5))
        add_text_box(slide, bx, gpt_y + Inches(0.02), Inches(0.8), Inches(0.26), bm, 7, FONT_BODY, C_TEXT, True, PP_ALIGN.CENTER)
        if bi < len(bert_models) - 1:
            add_arrow_line(slide, bx + Inches(0.8), gpt_y + Inches(0.15), bx + Inches(0.85), gpt_y + Inches(0.15), C_EMBED, Pt(1))
    add_text_box(slide, rb_x, gpt_y + Inches(0.35), Inches(2.5), Inches(0.25), "理解之王", 10, FONT_BODY, C_EMBED, True, PP_ALIGN.CENTER)
    
    # 时间线
    timeline_y = Inches(4.2)
    add_text_box(slide, left_x, timeline_y, left_w, Inches(0.3), "技术演进时间线", 12, FONT_BODY, C_TITLE, True)
    
    tl_y = timeline_y + Inches(0.4)
    tl_line_y = tl_y + Inches(0.3)
    add_line(slide, left_x + Inches(0.5), tl_line_y, left_x + Inches(7.3), tl_line_y, C_GRAY, Pt(2))
    
    timeline_items = [
        ("二零一七", "变换器\n论文发表", C_ENCODER),
        ("二零二零", "视觉变换器\n图像识别", C_FFN),
        ("二零二二", "语音识别\n多语言", C_ATTN),
        ("二零二三", "多模态\n图文理解", C_EMBED),
        ("二零二四", "视频生成\n视觉创作", C_RED),
    ]
    
    for ti, (year, desc, clr) in enumerate(timeline_items):
        tx = left_x + Inches(0.5) + ti * Inches(1.7)
        add_shape(slide, tx + Inches(0.15), tl_line_y - Inches(0.08), Inches(0.16), Inches(0.16), MSO_SHAPE.OVAL, clr)
        add_text_box(slide, tx - Inches(0.15), tl_line_y - Inches(0.35), Inches(0.9), Inches(0.25), year, 9, FONT_BODY, clr, True, PP_ALIGN.CENTER)
        add_text_box(slide, tx - Inches(0.3), tl_line_y + Inches(0.15), Inches(1.1), Inches(0.5), desc, 8, FONT_BODY, C_TEXT, False, PP_ALIGN.CENTER)
    
    # === 右侧40% ===
    right_x = Inches(8.5)
    right_w = Inches(4.5)
    
    items = [
        "生成系列：仅使用解码器，通过预测下一个词来生成文本。对话机器人经过三步训练：大规模预训练、指令微调、人类反馈强化学习",
        "理解系列：仅使用编码器，通过掩码语言建模训练（随机遮盖百分之十五的词），擅长文本分类和问答任务，被搜索引擎广泛使用",
        "转换系列：使用完整的编码器和解码器架构，将所有自然语言处理任务统一为文本到文本的格式，灵活通用",
        "超越语言领域：视觉变换器将图像切块当作词来处理，语音模型处理音频，蛋白质折叠模型处理氨基酸序列，序列数据是通用格式",
    ]
    add_bullet_text(slide, right_x, Inches(0.8), right_w, Inches(6.2), items, 11, C_TEXT)


# ===========================
# 第13页：总结
# ===========================
def create_slide_13():
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "总结：变换器模型的核心要点")
    
    cards = [
        ("1. 信息自由流动", "自注意力机制让序列中任意两个位置直接交流信息，彻底解决了循环神经网络的长距离依赖、串行计算和梯度消失三大难题", C_ENCODER, C_LIGHT_BLUE),
        ("2. 编码器理解加解码器生成", "编码器负责理解输入信息，解码器负责生成输出结果，两者通过键值对实现信息传递和分工协作", C_EMBED, C_LIGHT_PURPLE),
        ("3. 训练开卷加推理闭卷", "训练阶段并行高效地学习参数，推理阶段串行逐词生成文本，键值缓存技术有效加速了推理过程", C_ATTN, C_LIGHT_ORANGE),
        ("4. 三大家族各有所长", "生成系列擅长文本创作，理解系列擅长分析理解，转换系列擅长格式转换，同一套架构三种不同的职业方向", C_FFN, C_LIGHT_GREEN),
        ("5. 人工智能通用引擎", "从六千五百万参数到万亿参数，核心架构始终未变，已经成为整个人工智能时代的基石", C_RED, C_LIGHT_RED),
    ]
    
    card_w = Inches(10.0)
    card_h = Inches(0.85)
    card_gap = Inches(0.15)
    card_x = Inches(1.7)
    start_y = Inches(0.9)
    
    for idx, (title, desc, accent, bg) in enumerate(cards):
        cy = start_y + idx * (card_h + card_gap)
        add_shape(slide, card_x, cy, card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE, bg, accent, Pt(1.5))
        add_shape(slide, card_x, cy, Inches(0.1), card_h, MSO_SHAPE.RECTANGLE, accent)
        add_shape(slide, card_x + Inches(0.3), cy + Inches(0.2), Inches(0.45), Inches(0.45), MSO_SHAPE.OVAL, accent)
        add_text_box(slide, card_x + Inches(0.3), cy + Inches(0.25), Inches(0.45), Inches(0.35), str(idx + 1), 14, FONT_MONO, C_WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, card_x + Inches(0.95), cy + Inches(0.08), Inches(5.0), Inches(0.35), title, 12, FONT_BODY, accent, True)
        add_text_box(slide, card_x + Inches(0.95), cy + Inches(0.42), Inches(8.8), Inches(0.4), desc, 10, FONT_BODY, C_TEXT)
    
    # 底部结语
    footer_y = start_y + len(cards) * (card_h + card_gap) + Inches(0.3)
    add_text_box(slide, Inches(2.0), footer_y, Inches(9.3), Inches(0.4),
                 "变换器不是一个模型，而是一个时代。感谢聆听！",
                 12, FONT_BODY, C_GRAY, True, PP_ALIGN.CENTER)


# === 生成 ===
create_slide_8()
create_slide_9()
create_slide_10()
create_slide_11()
create_slide_12()
create_slide_13()

output_path = "/home/admin/.openclaw/workspace-weaver/output/transformer_v5_part2.pptx"
prs.save(output_path)
print(f"PPT已保存: {output_path}")

# 统计中英文比例
page_names = ["第8页：训练细节", "第9页：推理+KV缓存", "第10页：词选择", "第11页：训练vs推理", "第12页：应用", "第13页：总结"]

def count_chars(text):
    chinese = len(re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]', text))
    english = len(re.findall(r'[a-zA-Z]', text))
    digits = len(re.findall(r'[0-9]', text))
    total = len(text)
    return chinese, english, digits, total

total_chinese = 0
total_english = 0
total_all = 0

print("\n=== 中英文比例统计 ===")
for i, slide in enumerate(prs.slides):
    page_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                page_text += para.text
    
    ch, en, dg, tot = count_chars(page_text)
    total_chinese += ch
    total_english += en
    total_all += tot
    
    if tot > 0:
        ch_pct = ch / tot * 100
        en_pct = en / tot * 100
        print(f"{page_names[i]}: 中文={ch}({ch_pct:.1f}%) 英文={en}({en_pct:.1f}%) 总字符={tot}")

print(f"\n总计: 中文={total_chinese}({total_chinese/total_all*100:.1f}%) 英文={total_english}({total_english/total_all*100:.1f}%)")
print(f"中文占比 > 80%: {'是' if total_chinese/total_all > 0.8 else '否'}")
