#!/usr/bin/env python3
"""
SD PPT Fix v7: Fix 3 issues on slides 4, 7, 10
  - Problem 1 (P7): U-Net decoder arrows pointing wrong direction (should go UP)
  - Problem 2 (P10): CFG appears without context, add bridging text
  - Problem 3 (P4): Architecture diagram mixes inference & training, need clear separation

Output: sd_fix_v7.pptx (3 slides: P4, P7, P10)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MAX_X = 13.333
MAX_Y = 7.5

# Colors (same as v6)
C_PRI = RGBColor(0x1E,0x3A,0x5F)
C_ACC = RGBColor(0x34,0x98,0xDB)
C_LT  = RGBColor(0xEB,0xF5,0xFB)
C_W   = RGBColor(0xFF,0xFF,0xFF)
C_D   = RGBColor(0x33,0x33,0x33)
C_G   = RGBColor(0x88,0x88,0x88)
C_GRN = RGBColor(0x27,0xAE,0x60)
C_ORG = RGBColor(0xE6,0x7E,0x22)
C_RED = RGBColor(0xE7,0x4C,0x3C)
C_PUR = RGBColor(0x8E,0x44,0xAD)
C_TEL = RGBColor(0x00,0x96,0x88)
C_BG_GRN = RGBColor(0xE8,0xF8,0xF5)
C_BG_ORG = RGBColor(0xFD,0xF2,0xE9)
C_BG_RED = RGBColor(0xFD,0xED,0xED)

FC = "微软雅黑"
FE = "Arial"
OUT = "/home/admin/.openclaw/workspace-weaver/output/sd_fix_v7.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ── Helpers ──
def bg(s, c=C_W):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh

def rrect(s,l,t,w,h,fc=None,lc=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb=fc
    else: sh.fill.background()
    if lc: sh.line.fill.solid(); sh.line.color.rgb=lc
    return sh

def tb(s,l,t,w,h,txt,fs=12,fc=C_D,b=False,a=PP_ALIGN.LEFT,fn=FC,ls=1.15):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(fs); p.font.color.rgb=fc
    p.font.bold=b; p.font.name=fn; p.alignment=a; p.space_after=Pt(2)
    if ls!=1.0: p.line_spacing=ls
    return bx

def mtb(s,l,t,w,h,lines,fs=11,fc=C_D,fn=FC,a=PP_ALIGN.LEFT,ls=1.1):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,ld in enumerate(lines):
        if isinstance(ld,str): txt,bld,c,sz=ld,False,fc,fs
        elif len(ld)==2: txt,bld=ld; c,sz=fc,fs
        elif len(ld)==3: txt,bld,c=ld; sz=fs
        else: txt,bld,c,sz=ld
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=bld
        p.font.name=fn; p.alignment=a; p.space_after=Pt(3); p.line_spacing=ls
    return bx

def arrow_r(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_d(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def arrow_u(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.UP_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def chev(s,l,t,w,h,c=C_ACC):
    sh=s.shapes.add_shape(MSO_SHAPE.CHEVRON,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def hdr(s,title,sub=""):
    rect(s,Inches(0),Inches(0),SLIDE_W,Inches(1.1),fc=C_PRI)
    tb(s,Inches(0.8),Inches(0.15),Inches(11),Inches(0.6),title,fs=28,fc=C_W,b=True)
    if sub: tb(s,Inches(0.8),Inches(0.65),Inches(11),Inches(0.35),sub,fs=13,fc=RGBColor(0xBB,0xDE,0xFB))
    rect(s,Inches(0),Inches(7.15),SLIDE_W,Inches(0.35),fc=C_PRI)
    tb(s,Inches(0.5),Inches(7.15),Inches(5),Inches(0.35),"Stable Diffusion 技术解析",fs=9,fc=RGBColor(0x88,0xBB,0xDD))

def chk(s, msg=""):
    violations = []
    for sh in s.shapes:
        if sh.left < Emu(-100):
            violations.append(f"  NEGATIVE X: left={sh.left/914400:.3f}")
        if sh.top < Emu(-100):
            violations.append(f"  NEGATIVE Y: top={sh.top/914400:.3f}")
        r = sh.left + sh.width
        b = sh.top + sh.height
        if r > SLIDE_W + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW X: right={r/914400:.3f} text={txt}")
        if b > SLIDE_H + Emu(5000):
            txt = sh.text_frame.text[:30] if sh.has_text_frame else ''
            violations.append(f"  OVERFLOW Y: bottom={b/914400:.3f} text={txt}")
    if violations:
        print(f"⚠️  {msg} — {len(violations)} violations:")
        for v in violations: print(v)
    else:
        print(f"✅ {msg} — all shapes within bounds")

def info(s, msg):
    shape_count = len(s.shapes)
    print(f"📊 {msg}: {shape_count} shapes")

# ── Create presentation ──
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# SLIDE 1 (P4): 整体架构图 — 推理/训练分开展示
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "03  整体架构概览", "CLIP + U-Net + VAE 三大组件协同工作 — 推理与训练流程完全分离")

# ── 左半: 推理流程 ──
# Section header with blue bg
rect(sl, Inches(0.15), Inches(1.15), Inches(6.40), Inches(0.30), fc=C_ACC)
tb(sl, Inches(0.25), Inches(1.15), Inches(6.20), Inches(0.30),
   "▶ 推理流程（从文字到图像）— 主要使用场景", fs=11, fc=C_W, b=True)

# Step 1: 文字输入
rrect(sl, Inches(0.20), Inches(1.65), Inches(1.20), Inches(0.75), fc=C_W, lc=C_ACC)
tb(sl, Inches(0.22), Inches(1.67), Inches(1.16), Inches(0.18),
   "① 文字输入", fs=10, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.22), Inches(1.88), Inches(1.16), Inches(0.45),
   '"一只橘猫坐在沙发上"', fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# Arrow →
arrow_r(sl, Inches(1.45), Inches(1.88), Inches(0.25), Inches(0.18), C_ACC)

# Step 2: CLIP
rrect(sl, Inches(1.75), Inches(1.55), Inches(1.60), Inches(0.95), fc=C_LT, lc=C_ACC)
rect(sl, Inches(1.75), Inches(1.55), Inches(1.60), Inches(0.22), fc=C_ACC)
tb(sl, Inches(1.77), Inches(1.57), Inches(1.56), Inches(0.18),
   "② CLIP 文本编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(1.77), Inches(1.80), Inches(1.56), Inches(0.65),
   "123M参数\n12层Transformer\n输出: 77×768向量", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# Arrow →
arrow_r(sl, Inches(3.40), Inches(1.88), Inches(0.25), Inches(0.18), C_ACC)

# Step 3: U-Net (core, highlighted)
rrect(sl, Inches(3.70), Inches(1.45), Inches(1.80), Inches(1.15), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI)
rect(sl, Inches(3.70), Inches(1.45), Inches(1.80), Inches(0.22), fc=C_PRI)
tb(sl, Inches(3.72), Inches(1.47), Inches(1.76), Inches(0.18),
   "③ U-Net 去噪（核心）", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(3.72), Inches(1.70), Inches(1.76), Inches(0.85),
   "860M参数(80%)\n⟳ 循环20-50步\n噪声→逐步去噪", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# Cross-Attention label
rect(sl, Inches(3.40), Inches(1.72), Inches(0.30), Inches(0.18), fc=C_PUR)
tb(sl, Inches(3.40), Inches(1.72), Inches(0.30), Inches(0.18),
   "CA", fs=7, fc=C_W, b=True, a=PP_ALIGN.CENTER)

# Arrow →
arrow_r(sl, Inches(5.55), Inches(1.88), Inches(0.25), Inches(0.18), C_ACC)

# Step 4: VAE 解码器
rrect(sl, Inches(5.85), Inches(1.55), Inches(1.50), Inches(0.95), fc=C_BG_GRN, lc=C_GRN)
rect(sl, Inches(5.85), Inches(1.55), Inches(1.50), Inches(0.22), fc=C_GRN)
tb(sl, Inches(5.87), Inches(1.57), Inches(1.46), Inches(0.18),
   "④ VAE 解码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(5.87), Inches(1.80), Inches(1.46), Inches(0.65),
   "~49M参数\n64×64→512×512\n4层上采样", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# Arrow →
arrow_r(sl, Inches(7.40), Inches(1.88), Inches(0.25), Inches(0.18), C_ACC)

# Step 5: 图像输出
rrect(sl, Inches(7.70), Inches(1.65), Inches(1.20), Inches(0.75), fc=C_BG_ORG, lc=C_ORG)
tb(sl, Inches(7.72), Inches(1.67), Inches(1.16), Inches(0.18),
   "⑤ 图像输出", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(7.72), Inches(1.88), Inches(1.16), Inches(0.45),
   "512×512×3\n高清RGB图像", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# 随机噪声输入 → U-Net
rrect(sl, Inches(4.00), Inches(2.70), Inches(1.20), Inches(0.45), fc=C_BG_RED, lc=C_RED)
tb(sl, Inches(4.02), Inches(2.72), Inches(1.16), Inches(0.40),
   "随机噪声 N(0,I)\n64×64×4 纯噪声", fs=8, fc=C_D, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(4.50), Inches(2.62), Inches(0.18), Inches(0.12), C_RED)

# ── VAE编码器标记：推理时不需要 ──
# Ghost VAE编码器 with red X
rrect(sl, Inches(2.30), Inches(2.70), Inches(1.50), Inches(0.45), fc=C_BG_RED, lc=C_RED)
# Draw X effect with red border + text
rect(sl, Inches(2.30), Inches(2.70), Inches(1.50), Inches(0.45), lc=C_RED)
tb(sl, Inches(2.32), Inches(2.72), Inches(1.46), Inches(0.40),
   "❌ VAE 编码器\n推理时不需要！", fs=8, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.0)

# ── 右半: 训练流程 ──
rect(sl, Inches(0.15), Inches(3.35), Inches(6.40), Inches(0.30), fc=C_ORG)
tb(sl, Inches(0.25), Inches(3.35), Inches(6.20), Inches(0.30),
   "▶ 训练流程（从图像学习去噪能力）— 模型训练阶段", fs=11, fc=C_W, b=True)

# Training flow: Image+Text → VAE Encoder → Add Noise → U-Net → Loss → Backprop
# Step T1: Image + Text
rrect(sl, Inches(0.20), Inches(3.85), Inches(1.20), Inches(0.70), fc=C_W, lc=C_ORG)
tb(sl, Inches(0.22), Inches(3.87), Inches(1.16), Inches(0.18),
   "图像+文本对", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.22), Inches(4.08), Inches(1.16), Inches(0.40),
   "512×512×3 图像\n+ 文字描述", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(1.45), Inches(4.08), Inches(0.25), Inches(0.18), C_ORG)

# Step T2: VAE 编码器
rrect(sl, Inches(1.75), Inches(3.80), Inches(1.60), Inches(0.80), fc=C_BG_GRN, lc=C_GRN)
rect(sl, Inches(1.75), Inches(3.80), Inches(1.60), Inches(0.22), fc=C_GRN)
tb(sl, Inches(1.77), Inches(3.82), Inches(1.56), Inches(0.18),
   "VAE 编码器", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(1.77), Inches(4.05), Inches(1.56), Inches(0.50),
   "~34M参数\n图像→潜空间\n512²→64² 压缩48倍", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(3.40), Inches(4.08), Inches(0.25), Inches(0.18), C_ORG)

# Step T3: 添加噪声
rrect(sl, Inches(3.70), Inches(3.85), Inches(1.30), Inches(0.70), fc=C_BG_RED, lc=C_RED)
tb(sl, Inches(3.72), Inches(3.87), Inches(1.26), Inches(0.18),
   "添加噪声", fs=10, fc=C_RED, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(3.72), Inches(4.08), Inches(1.26), Inches(0.40),
   "前向扩散\nz_t = √ᾱ·z₀ + √(1-ᾱ)·ε", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(5.05), Inches(4.08), Inches(0.25), Inches(0.18), C_ORG)

# Step T4: U-Net 预测噪声
rrect(sl, Inches(5.35), Inches(3.80), Inches(1.80), Inches(0.80), fc=RGBColor(0xE8,0xF0,0xFE), lc=C_PRI)
rect(sl, Inches(5.35), Inches(3.80), Inches(1.80), Inches(0.22), fc=C_PRI)
tb(sl, Inches(5.37), Inches(3.82), Inches(1.76), Inches(0.18),
   "U-Net 预测噪声", fs=9, fc=C_W, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(5.37), Inches(4.05), Inches(1.76), Inches(0.50),
   "输入: 带噪声潜空间\n+ CLIP文本向量\n输出: ε_pred 预测噪声", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

arrow_r(sl, Inches(7.20), Inches(4.08), Inches(0.25), Inches(0.18), C_ORG)

# Step T5: 计算损失 → 反向传播
rrect(sl, Inches(7.50), Inches(3.85), Inches(1.40), Inches(0.70), fc=C_BG_ORG, lc=C_ORG)
tb(sl, Inches(7.52), Inches(3.87), Inches(1.36), Inches(0.18),
   "计算损失", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(7.52), Inches(4.08), Inches(1.36), Inches(0.40),
   "L = ||ε - ε_pred||²\n反向传播更新U-Net", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# CLIP also used in training
rrect(sl, Inches(3.70), Inches(4.75), Inches(1.30), Inches(0.40), fc=C_LT, lc=C_ACC)
tb(sl, Inches(3.72), Inches(4.77), Inches(1.26), Inches(0.35),
   "CLIP编码文本", fs=8, fc=C_ACC, a=PP_ALIGN.CENTER, ls=1.0)
arrow_u(sl, Inches(4.28), Inches(4.67), Inches(0.15), Inches(0.12), C_ACC)

# VAE解码器标记：训练时不需要
rrect(sl, Inches(5.35), Inches(4.75), Inches(1.80), Inches(0.40), fc=C_BG_RED, lc=C_RED)
tb(sl, Inches(5.37), Inches(4.77), Inches(1.76), Inches(0.35),
   "❌ VAE 解码器  训练时不需要！", fs=8, fc=C_RED, b=True, a=PP_ALIGN.CENTER, ls=1.0)

# ── Component usage summary ──
rect(sl, Inches(0.15), Inches(5.30), Inches(12.80), Inches(0.25), fc=C_PRI)
tb(sl, Inches(0.25), Inches(5.30), Inches(12.60), Inches(0.25),
   "组件使用对照表", fs=11, fc=C_W, b=True)

# Component cards
# CLIP
rrect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(1.25), fc=C_W, lc=C_ACC)
rect(sl, Inches(0.15), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_ACC)
tb(sl, Inches(0.25), Inches(5.67), Inches(3.90), Inches(0.22),
   "CLIP 文本编码器 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(0.25), Inches(5.95), Inches(3.90), Inches(0.90),
   "• 翻译官 — 123M参数\n• 训练: 4亿张图文对预训练\n• 12层Transformer, 768维\n• 作用: 把人类语言翻译成AI数字语言", fs=9, fc=C_D, ls=1.15)

# U-Net
rrect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(1.25), fc=C_W, lc=C_PRI)
rect(sl, Inches(4.40), Inches(5.65), Inches(4.10), Inches(0.25), fc=C_PRI)
tb(sl, Inches(4.50), Inches(5.67), Inches(3.90), Inches(0.22),
   "U-Net 去噪网络 — 训练/推理 都用", fs=10, fc=C_W, b=True)
tb(sl, Inches(4.50), Inches(5.95), Inches(3.90), Inches(0.90),
   "• 核心画师 — 860M参数(80%)\n• U型编码器-解码器结构\n• Cross-Attention + Skip Connections\n• 作用: 一步步去掉噪声", fs=9, fc=C_D, ls=1.15)

# VAE
rrect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(1.25), fc=C_W, lc=C_GRN)
rect(sl, Inches(8.65), Inches(5.65), Inches(4.30), Inches(0.25), fc=C_GRN)
tb(sl, Inches(8.75), Inches(5.67), Inches(4.10), Inches(0.22),
   "VAE 编解码器 — 编码器(训练) 解码器(推理)", fs=10, fc=C_W, b=True)
tb(sl, Inches(8.75), Inches(5.95), Inches(4.10), Inches(0.90),
   "• 编码器(~34M): 图像→潜空间 [训练用]\n• 解码器(~49M): 潜空间→图像 [推理用]\n• 48倍压缩: 786K→16K数值\n• ⚠️ 编码器和解码器不在同一个流程中！", fs=9, fc=C_D, ls=1.15)

chk(sl, "P4 整体架构")
info(sl, "P4 整体架构")


# ============================================================
# SLIDE 2 (P7): U-Net 架构图 — 修复解码器箭头方向
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "06  U-Net 架构图", "编码器-瓶颈-解码器 U 型结构 | 860M 参数 | 3 个关键机制")

# ── 左侧: 编码器 (从上到下) ──
tb(sl, Inches(0.60), Inches(1.10), Inches(2.20), Inches(0.22),
   "编码器 Encoder ↓", fs=11, fc=C_ACC, b=True, a=PP_ALIGN.CENTER)

enc_layers = [
    ("64×64×320", "Res+Attn+Down", "✦ Cross-Attention"),
    ("32×32×640", "Res+Attn+Down", "✦ Cross-Attention"),
    ("16×16×1280", "Res+Attn+Down", "✦ Cross-Attention"),
    ("8×8×1280", "Res+Down", "  (无注意力)"),
]
enc_y = 1.35
for dim, ops, attn in enc_layers:
    rrect(sl, Inches(0.60), Inches(enc_y), Inches(2.20), Inches(0.65), fc=C_W, lc=C_ACC)
    tb(sl, Inches(0.65), Inches(enc_y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.65), Inches(enc_y+0.25), Inches(2.10), Inches(0.18), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(0.65), Inches(enc_y+0.42), Inches(2.10), Inches(0.18), attn, fs=8, fc=C_PUR, a=PP_ALIGN.CENTER)
    enc_y += 0.80
    # Down arrow between layers (except after last)
    if attn != "  (无注意力)":
        arrow_d(sl, Inches(1.60), Inches(enc_y-0.80+0.65), Inches(0.20), Inches(0.15), C_ACC)

# Down arrow after 8×8 layer
arrow_d(sl, Inches(1.60), Inches(enc_y-0.80+0.65), Inches(0.20), Inches(0.15), C_ACC)

# ── Bottleneck ──
rrect(sl, Inches(0.45), Inches(enc_y), Inches(2.50), Inches(0.75), fc=C_BG_ORG, lc=C_ORG)
tb(sl, Inches(0.50), Inches(enc_y+0.02), Inches(2.40), Inches(0.22), "4×4×1280", fs=10, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.50), Inches(enc_y+0.25), Inches(2.40), Inches(0.18), "Bottleneck", fs=9, fc=C_D, a=PP_ALIGN.CENTER)
tb(sl, Inches(0.50), Inches(enc_y+0.42), Inches(2.40), Inches(0.18), "Res+SelfAttn+Res", fs=8, fc=C_D, a=PP_ALIGN.CENTER)

bottleneck_y = enc_y

# ── Right side: Decoder (from bottom to top!) ──
tb(sl, Inches(4.00), Inches(1.10), Inches(2.20), Inches(0.22),
   "解码器 Decoder ↑", fs=11, fc=C_GRN, b=True, a=PP_ALIGN.CENTER)

dec_layers = [
    ("8×8×1280", "Res+Up+Skip", "✦ Cross-Attention"),
    ("16×16×640", "Res+Up+Skip", "✦ Cross-Attention"),
    ("32×32×320", "Res+Up+Skip", "✦ Cross-Attention"),
    ("64×64×4", "输出预测噪声", "  (无注意力)"),
]

# Decoder starts at bottom, goes up
# Place decoder bottom at same y as bottleneck
dec_start_y = bottleneck_y  # bottom of decoder = bottleneck level
dec_spacing = 0.80
dec_layer_h = 0.65

# Calculate: decoder layers from bottom to top
# Bottleneck → 8×8 → 16×16 → 32×32 → 64×64
# We place them starting from the bottom going UP
# 8×8 at dec_start_y - spacing
# 16×16 at dec_start_y - 2*spacing
# etc.

# Horizontal arrow from bottleneck to first decoder layer (8×8)
# Bottleneck right edge → 8×8 left edge at same Y
arrow_r(sl, Inches(2.95), Inches(bottleneck_y+0.25), Inches(1.05), Inches(0.18), C_GRN)

dec_y = dec_start_y - dec_spacing  # First decoder layer (8×8) above bottleneck

for i, (dim, ops, attn) in enumerate(dec_layers):
    rrect(sl, Inches(4.00), Inches(dec_y), Inches(2.20), dec_layer_h, fc=C_BG_GRN, lc=C_GRN)
    tb(sl, Inches(4.05), Inches(dec_y+0.02), Inches(2.10), Inches(0.22), dim, fs=10, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
    tb(sl, Inches(4.05), Inches(dec_y+0.25), Inches(2.10), Inches(0.18), ops, fs=9, fc=C_D, a=PP_ALIGN.CENTER)
    tb(sl, Inches(4.05), Inches(dec_y+0.42), Inches(2.10), Inches(0.18), attn, fs=8, fc=C_PUR, a=PP_ALIGN.CENTER)
    
    # UP arrow between decoder layers (from bottom layer to upper layer)
    if i < len(dec_layers) - 1:
        # Arrow pointing UP from current layer to next upper layer
        arrow_u(sl, Inches(5.00), Inches(dec_y-0.15), Inches(0.20), Inches(0.15), C_GRN)
        # Add "上采样↑" label
        tb(sl, Inches(5.22), Inches(dec_y-0.18), Inches(0.90), Inches(0.18),
           "上采样↑", fs=8, fc=C_GRN, b=True)
    
    dec_y -= dec_spacing

# Skip Connections (horizontal, encoder → decoder at same level)
# Encoder layers Y positions: 1.35, 2.15, 2.95, 3.75
# Decoder layers Y positions (bottom to top): calculated above
# Need to match: enc[0]=64×64 ↔ dec[3]=64×64, enc[1]=32×32 ↔ dec[2]=32×32, etc.

enc_ys = [1.35, 2.15, 2.95, 3.75]  # Y positions of encoder layers
dec_ys_bottom_to_top = []
tmp_y = dec_start_y - dec_spacing
for i in range(len(dec_layers)):
    dec_ys_bottom_to_top.append(tmp_y)
    tmp_y -= dec_spacing
dec_ys = list(reversed(dec_ys_bottom_to_top))  # Now top to bottom: 64×64, 32×32, 16×16, 8×8

for i in range(4):
    enc_mid_y = enc_ys[i] + 0.30
    # Skip line
    rect(sl, Inches(2.80), Inches(enc_mid_y), Inches(1.20), Inches(0.04), fc=C_ORG)
    tb(sl, Inches(2.90), Inches(enc_mid_y-0.16), Inches(0.80), Inches(0.18),
       "Skip", fs=8, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)
    arrow_r(sl, Inches(3.85), Inches(enc_mid_y-0.04), Inches(0.15), Inches(0.12), C_ORG)

# ── Time Embedding (left side) ──
tb(sl, Inches(0.05), Inches(1.40), Inches(0.50), Inches(0.18), "Time", fs=8, fc=C_TEL, b=True)
tb(sl, Inches(0.05), Inches(1.55), Inches(0.50), Inches(0.18), "Embed", fs=8, fc=C_TEL, b=True)
for y_pos in [1.70, 1.85, 2.90, 3.70]:
    arrow_r(sl, Inches(0.05), Inches(y_pos), Inches(0.50), Inches(0.15), C_TEL)

# ── Right panel: 3 key mechanisms ──
# ① Cross-Attention
rrect(sl, Inches(6.80), Inches(1.20), Inches(6.00), Inches(1.75), fc=C_W, lc=C_PUR)
rect(sl, Inches(6.80), Inches(1.20), Inches(6.00), Inches(0.28), fc=C_PUR)
tb(sl, Inches(6.90), Inches(1.22), Inches(5.80), Inches(0.24),
   "① Cross-Attention 交叉注意力", fs=10, fc=C_W, b=True)
tb(sl, Inches(6.90), Inches(1.52), Inches(5.80), Inches(0.18),
   "文本条件注入的核心机制", fs=9, fc=C_PUR, b=True)
mtb(sl, Inches(6.90), Inches(1.75), Inches(5.80), Inches(1.10), [
    "• 图像特征作为 Query，文本特征作为 Key-Value",
    "• 图像问：根据文字，我这里该怎样？",
    "• 文本答：你这里应该是橘猫的毛发",
    "• 出现位置：编码器1-3层 + 解码器1-3层（共6处）",
    "• 注意力计算: Q·K^T/√d → Softmax → ×V",
], fs=9, fc=C_D, ls=1.15)

# ② Skip Connections
rrect(sl, Inches(6.80), Inches(3.10), Inches(6.00), Inches(1.75), fc=C_W, lc=C_ORG)
rect(sl, Inches(6.80), Inches(3.10), Inches(6.00), Inches(0.28), fc=C_ORG)
tb(sl, Inches(6.90), Inches(3.12), Inches(5.80), Inches(0.24),
   "② Skip Connections 跳跃连接", fs=10, fc=C_W, b=True)
tb(sl, Inches(6.90), Inches(3.42), Inches(5.80), Inches(0.18),
   "信息传递的桥梁（图中橙色横线）", fs=9, fc=C_ORG, b=True)
mtb(sl, Inches(6.90), Inches(3.65), Inches(5.80), Inches(1.10), [
    "• 编码器特征直接桥接到解码器对应层",
    "• 防止细节在压缩过程中丢失",
    "• 像建筑师随时参考草稿本",
    "• 没有跳跃连接，解码器只能\"凭记忆\"还原，效果大打折扣",
    "• 横向连接：64×64↔64×64, 32×32↔32×32, 16×16↔16×16, 8×8↔8×8",
], fs=9, fc=C_D, ls=1.15)

# ③ Time Embedding
rrect(sl, Inches(6.80), Inches(5.00), Inches(6.00), Inches(1.85), fc=C_W, lc=C_TEL)
rect(sl, Inches(6.80), Inches(5.00), Inches(6.00), Inches(0.28), fc=C_TEL)
tb(sl, Inches(6.90), Inches(5.02), Inches(5.80), Inches(0.24),
   "③ Time Embedding 时间嵌入", fs=10, fc=C_W, b=True)
tb(sl, Inches(6.90), Inches(5.32), Inches(5.80), Inches(0.18),
   "当前噪声水平的标尺", fs=9, fc=C_TEL, b=True)
mtb(sl, Inches(6.90), Inches(5.55), Inches(5.80), Inches(1.20), [
    "• 告诉U-Net当前噪声水平（1-1000）",
    "• 高噪声(t≈900)：粗雕确定轮廓",
    "• 低噪声(t≈100)：精雕添加细节",
    "• 正弦编码 → 2层MLP → 1280维",
    "• 注入每一个ResBlock，影响去噪策略",
], fs=9, fc=C_D, ls=1.15)

chk(sl, "P7 U-Net架构")
info(sl, "P7 U-Net架构")


# ============================================================
# SLIDE 3 (P10): CFG — 增加引导文字，逻辑过渡
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
hdr(sl, "09  CFG（Classifier-Free Guidance）详解", "无分类器引导 — 让AI更听话地按描述画画")

# ── 引导文字: 为什么需要CFG？──
rrect(sl, Inches(0.50), Inches(1.20), Inches(12.30), Inches(1.10), fc=C_BG_RED, lc=C_RED)
tb(sl, Inches(0.70), Inches(1.22), Inches(11.80), Inches(0.22),
   "🤔 为什么需要CFG？（从推理过程过渡）", fs=12, fc=C_RED, b=True)
mtb(sl, Inches(0.70), Inches(1.48), Inches(11.80), Inches(0.75), [
    "在前面的推理过程中，U-Net去噪时需要一个\"方向盘\"来控制生成方向。没有CFG时，U-Net只知道\"去掉噪声\"，",
    "生成的图像随机且不可控——可能和你的prompt完全无关。CFG的作用就是让U-Net\"听\"prompt的话，生成符合描述的图像。",
], fs=10, fc=C_D, ls=1.2)

# ── 核心公式与原理 ──
rrect(sl, Inches(0.50), Inches(2.45), Inches(12.30), Inches(1.20), fc=C_W, lc=C_ACC)
tb(sl, Inches(0.70), Inches(2.48), Inches(11.80), Inches(0.22),
   "核心思想：既知道什么是对的，也知道什么不是对的", fs=11, fc=C_PRI, b=True)
mtb(sl, Inches(0.70), Inches(2.75), Inches(11.80), Inches(0.85), [
    "公式: ε_guided = ε_uncond + s × (ε_cond - ε_uncond)",
    "ε_uncond = 没有文本条件时的预测（自由发挥）    ε_cond = 有文本条件时的预测（按指令画）",
    "s = 引导系数（CFG Scale），通常 7-12    s越大越忠实于文字    s过大则过度饱和出现伪影",
], fs=10, fc=C_D, ls=1.2)

# ── CFG Scale 参数表 ──
# Create table
from pptx.util import Inches, Pt
table_shape = sl.shapes.add_table(5, 3, Inches(0.50), Inches(3.80), Inches(6.50), Inches(1.80))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(1.20)
table.columns[1].width = Inches(2.50)
table.columns[2].width = Inches(2.80)

headers = ["CFG Scale", "效果", "适用场景"]
data = [
    ["1-3", "几乎没有引导，图像\"自由发挥\"", "创意探索，想要意外惊喜"],
    ["5-7", "适中引导，平衡忠实度和创造力", "通用场景，推荐默认值 7"],
    ["7-12", "强引导，高度忠实于文字描述", "需要精确控制，如产品图"],
    ["15-20+", "过强引导，过度饱和出现伪影", "❌ 一般不推荐"],
]

# Style header
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_W
        p.font.name = FC
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRI

# Style data rows
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = C_D
            p.font.name = FC
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_LT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_W

# ── CFG 工作原理图 ──
rrect(sl, Inches(7.50), Inches(3.80), Inches(5.30), Inches(2.80), fc=C_W, lc=C_ACC)
tb(sl, Inches(7.70), Inches(3.85), Inches(4.80), Inches(0.25),
   "CFG 工作原理", fs=11, fc=C_PRI, b=True)

# 无条件生成
rrect(sl, Inches(7.80), Inches(4.25), Inches(2.20), Inches(0.55), fc=C_BG_RED, lc=C_RED)
tb(sl, Inches(7.80), Inches(4.35), Inches(2.20), Inches(0.35),
   "无条件生成（自由发挥）\nε_uncond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# → 差异
tb(sl, Inches(10.10), Inches(4.30), Inches(0.40), Inches(0.25), "→", fs=14, fc=C_PRI, b=True, a=PP_ALIGN.CENTER)
tb(sl, Inches(10.10), Inches(4.52), Inches(0.50), Inches(0.20), "差异", fs=8, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)

# 有条件生成
rrect(sl, Inches(10.60), Inches(4.25), Inches(2.00), Inches(0.55), fc=C_BG_GRN, lc=C_GRN)
tb(sl, Inches(10.60), Inches(4.35), Inches(2.00), Inches(0.35),
   "有条件生成（按指令画）\nε_cond", fs=9, fc=C_D, a=PP_ALIGN.CENTER, ls=1.1)

# × CFG Scale
tb(sl, Inches(10.80), Inches(4.90), Inches(1.50), Inches(0.20), "× CFG Scale (s)", fs=9, fc=C_ORG, b=True, a=PP_ALIGN.CENTER)

# ↓ 
arrow_d(sl, Inches(11.50), Inches(5.10), Inches(0.20), Inches(0.18), C_ORG)

# = 结果
rrect(sl, Inches(9.30), Inches(5.40), Inches(2.80), Inches(0.50), fc=C_BG_GRN, lc=C_GRN)
tb(sl, Inches(9.30), Inches(5.48), Inches(2.80), Inches(0.35),
   "= 引导后的高质量结果\nε_guided", fs=9, fc=C_D, b=True, a=PP_ALIGN.CENTER, ls=1.1)

# + 无条件
tb(sl, Inches(7.80), Inches(5.48), Inches(1.50), Inches(0.30), "+ 无条件结果", fs=9, fc=C_D, a=PP_ALIGN.CENTER)

# Arrow from uncond to result
arrow_r(sl, Inches(9.30), Inches(5.55), Inches(0.10), Inches(0.10), C_ACC)

# ── 底部比喻 ──
rrect(sl, Inches(0.50), Inches(5.90), Inches(12.30), Inches(0.95), fc=C_LT, lc=C_ACC)
mtb(sl, Inches(0.70), Inches(5.95), Inches(11.80), Inches(0.85), [
    "💡 比喻: CFG 就像严厉又公正的艺术指导。每画一笔都问两个问题：",
    "\"你按客户要求画了吗？\" 和 \"不管要求你会怎么画？\"",
    "然后把两个答案的差异放大（乘以 CFG Scale），强制画师更严格遵循客户要求。",
    "没有 CFG，图像模糊且偏离描述。CFG Scale 太大，画师太紧张反而画得太过刻意。",
], fs=9, fc=C_D, ls=1.2)

chk(sl, "P10 CFG详解")
info(sl, "P10 CFG详解")


# ── Save ──
prs.save(OUT)
print(f"\n✅ Saved to: {OUT}")
print(f"📊 Total slides: {len(prs.slides)}")

# ── Final verification ──
print("\n📋 Final verification:")
for i, slide in enumerate(prs.slides):
    shapes = len(slide.shapes)
    print(f"  Slide {i+1}: {shapes} shapes")
    violations = []
    for sh in slide.shapes:
        if sh.left < Emu(-100):
            violations.append(f"    NEGATIVE X: {sh.name} left={sh.left/914400:.3f}")
        if sh.top < Emu(-100):
            violations.append(f"    NEGATIVE Y: {sh.name} top={sh.top/914400:.3f}")
        r = sh.left + sh.width
        b = sh.top + sh.height
        if r > SLIDE_W + Emu(5000):
            violations.append(f"    OVERFLOW X: {sh.name} right={r/914400:.3f}")
        if b > SLIDE_H + Emu(5000):
            violations.append(f"    OVERFLOW Y: {sh.name} bottom={b/914400:.3f}")
    if violations:
        for v in violations: print(v)
    else:
        print(f"    ✅ All within bounds")
