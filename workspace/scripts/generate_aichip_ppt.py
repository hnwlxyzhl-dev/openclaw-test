#!/usr/bin/env python3
"""
A股AI芯片行业深度研究报告 PPT 生成脚本
风格：深色科技金融报告
配色：深蓝黑底 + 科技蓝(#0EA5E9) + 金色(#F59E0B) + 白色文字
页数：18页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import math

# ============================================================
# 全局配置
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色方案 - 深色科技金融
BG_DARK = RGBColor(0x0B, 0x11, 0x1A)        # 深蓝黑背景
BG_CARD = RGBColor(0x11, 0x1A, 0x2B)        # 卡片背景
BG_CARD2 = RGBColor(0x16, 0x21, 0x35)       # 次级卡片
ACCENT_BLUE = RGBColor(0x0E, 0xA5, 0xE9)    # 科技蓝
ACCENT_BLUE_DIM = RGBColor(0x07, 0x59, 0x8B) # 暗蓝
ACCENT_GOLD = RGBColor(0xF5, 0x9E, 0x0B)    # 金色
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # 绿色(正增长)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)     # 红色(风险/亏损)
TEXT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)     # 主文字白
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)      # 辅助灰
TEXT_DIM = RGBColor(0x64, 0x74, 0x8B)       # 暗灰
BORDER_COLOR = RGBColor(0x1E, 0x29, 0x3B)   # 边框线
TABLE_HEADER_BG = RGBColor(0x0C, 0x4A, 0x6E) # 表头蓝
TABLE_ROW1 = RGBColor(0x0F, 0x17, 0x2A)    # 表格行1
TABLE_ROW2 = RGBColor(0x13, 0x1C, 0x31)    # 表格行2

FONT_CN = "微软雅黑"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# 辅助函数
# ============================================================

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=Pt(12), color=TEXT_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.15,
                 anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    # 设置中文字体
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), font_name)
    # 行间距
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    # 垂直对齐
    txBody = txBox._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set('anchor', 'ctr')
    elif anchor == MSO_ANCHOR.BOTTOM:
        bodyPr.set('anchor', 'b')
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=Pt(11),
                   default_color=TEXT_GRAY, font_name=FONT_CN, line_spacing=1.2):
    """
    添加多行文本框，lines格式:
    [(text, font_size, color, bold), ...]
    如果只需要(text, font_size, color)，bold默认False
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        if len(line_data) >= 4:
            text, size, color, bold = line_data
        elif len(line_data) == 3:
            text, size, color = line_data
            bold = False
        else:
            text, size = line_data
            color = default_color
            bold = False

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), font_name)

        # 行间距
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

        # 段后间距
        spcAft = pPr.makeelement(qn('a:spcAft'), {})
        spcPts = spcAft.makeelement(qn('a:spcPts'), {'val': '400'})
        spcAft.append(spcPts)
        pPr.append(spcAft)

    return txBox

def add_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE_DIM, width=Pt(1)):
    """添加直线"""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_page_number(slide, num, total=18):
    """添加页码"""
    add_text_box(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.35),
                 f"{num} / {total}", Pt(9), TEXT_DIM, alignment=PP_ALIGN.RIGHT)

def add_section_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5), color=ACCENT_BLUE):
    """添加章节竖线装饰"""
    add_rect(slide, left, top, width, height, fill_color=color)

def add_kpi_card(slide, left, top, width, height, title, value, subtitle="", accent=ACCENT_BLUE):
    """添加KPI卡片"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER_COLOR)
    # 顶部装饰线
    add_rect(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Pt(2), fill_color=accent)
    # 标题
    add_text_box(slide, left + Inches(0.15), top + Inches(0.25), width - Inches(0.3), Inches(0.3),
                 title, Pt(9), TEXT_DIM, font_name=FONT_EN)
    # 数值
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), Inches(0.5),
                 value, Pt(22), TEXT_WHITE, bold=True, font_name=FONT_EN)
    # 副标题
    if subtitle:
        add_text_box(slide, left + Inches(0.15), top + height - Inches(0.4), width - Inches(0.3), Inches(0.3),
                     subtitle, Pt(9), accent)
    return card

def add_table(slide, left, top, width, col_widths, headers, rows, header_bg=TABLE_HEADER_BG):
    """添加表格"""
    row_count = len(rows) + 1
    col_count = len(headers)
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, Inches(0.4 * row_count))
    table = table_shape.table

    # 设置列宽
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 表头
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.size = Pt(10)
        run.font.color.rgb = TEXT_WHITE
        run.font.bold = True
        run.font.name = FONT_CN
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)
        p.alignment = PP_ALIGN.CENTER
        # 表头背景
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(*header_bg)})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
        # 去边框
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = tcPr.makeelement(qn(border_name), {'w': '6350'})
            sf = ln.makeelement(qn('a:solidFill'), {})
            sc = sf.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(*BORDER_COLOR)})
            sf.append(sc)
            ln.append(sf)
            tcPr.append(ln)

    # 数据行
    for i, row_data in enumerate(rows):
        bg = TABLE_ROW1 if i % 2 == 0 else TABLE_ROW2
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(9)
            run.font.color.rgb = TEXT_GRAY
            run.font.name = FONT_CN
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:ea'), FONT_CN)
            p.alignment = PP_ALIGN.CENTER
            # 行背景
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(*bg)})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.makeelement(qn(border_name), {'w': '6350'})
                sf = ln.makeelement(qn('a:solidFill'), {})
                sc = sf.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(*BORDER_COLOR)})
                sf.append(sc)
                ln.append(sf)
                tcPr.append(ln)

    return table_shape

def add_footer_bar(slide):
    """添加底部装饰条"""
    add_rect(slide, Inches(0), Inches(7.2), Inches(13.333), Pt(2), fill_color=ACCENT_BLUE_DIM)

def add_top_accent(slide):
    """添加顶部装饰线"""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Pt(3), fill_color=ACCENT_BLUE)

def slide_title(slide, title, subtitle=""):
    """标准页面标题"""
    add_top_accent(slide)
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(10), Inches(0.55),
                 title, Pt(24), TEXT_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.9), Inches(10), Inches(0.3),
                     subtitle, Pt(11), TEXT_DIM)
    add_line(slide, Inches(0.6), Inches(1.2), Inches(12.5), Inches(1.2), ACCENT_BLUE_DIM, Pt(1))
    add_footer_bar(slide)


# ============================================================
# Slide 1: 封面
# ============================================================
def make_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 背景装饰 - 右侧光晕
    glow = add_rect(slide, Inches(8), Inches(0), Inches(5.333), Inches(7.5),
                    fill_color=RGBColor(0x07, 0x1A, 0x2E))

    # 左侧装饰线
    add_rect(slide, Inches(0.8), Inches(1.5), Pt(3), Inches(2.5), fill_color=ACCENT_BLUE)
    add_rect(slide, Inches(0.8), Inches(4.3), Pt(3), Inches(1.2), fill_color=ACCENT_GOLD)

    # 主标题
    add_text_box(slide, Inches(1.2), Inches(1.5), Inches(9), Inches(1.0),
                 "A股AI芯片行业", Pt(40), TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.3), Inches(9), Inches(0.8),
                 "深度研究报告", Pt(40), ACCENT_BLUE, bold=True)

    # 副标题
    add_text_box(slide, Inches(1.2), Inches(3.3), Inches(9), Inches(0.4),
                 "从概念炒作到业绩兑现 —— 国产AI芯片的盈利拐点之年", Pt(14), TEXT_GRAY)

    # 分隔线
    add_line(slide, Inches(1.2), Inches(4.0), Inches(5), Inches(4.0), ACCENT_BLUE_DIM, Pt(1))

    # 信息行
    add_text_box(slide, Inches(1.2), Inches(4.3), Inches(4), Inches(0.3),
                 "研究机构：Weaver Research", Pt(11), TEXT_DIM)
    add_text_box(slide, Inches(1.2), Inches(4.7), Inches(4), Inches(0.3),
                 "报告日期：2026年5月", Pt(11), TEXT_DIM)
    add_text_box(slide, Inches(1.2), Inches(5.1), Inches(4), Inches(0.3),
                 "数据截止：2026年Q1", Pt(11), TEXT_DIM)

    # 右侧装饰标签
    tags = ["寒武纪", "海光信息", "景嘉微", "工业富联"]
    for i, tag in enumerate(tags):
        y = Inches(5.5) + Inches(i * 0.4)
        add_rounded_rect(slide, Inches(8.5), y, Inches(1.8), Inches(0.3),
                         fill_color=BG_CARD2, border_color=ACCENT_BLUE_DIM, border_width=Pt(1))
        add_text_box(slide, Inches(8.55), y + Pt(2), Inches(1.7), Inches(0.25),
                     tag, Pt(10), ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # 免责声明
    add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
                 "免责声明：本报告仅供参考，不构成投资建议。投资有风险，入市需谨慎。",
                 Pt(9), TEXT_DIM, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 目录
# ============================================================
def make_toc():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_top_accent(slide)
    add_footer_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(5), Inches(0.55),
                 "目录", Pt(28), TEXT_WHITE, bold=True)
    add_line(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(1.0), ACCENT_BLUE, Pt(2))

    sections = [
        ("01", "行业概览", "市场空间、产业链结构与发展阶段", ACCENT_BLUE),
        ("02", "寒武纪（688256）", "AI芯片绝对龙头，首次年度盈利", ACCENT_GOLD),
        ("03", "海光信息（688041）", "CPU+DCU双轮驱动，中国版AMD", ACCENT_BLUE),
        ("04", "景嘉微（300474）", "GPU+边端AI SoC双线转型", ACCENT_GOLD),
        ("05", "工业富联（601138）", "AI服务器制造龙头", ACCENT_BLUE),
        ("06", "竞争格局", "云端/边端/通用计算全景图", ACCENT_GOLD),
        ("07", "趋势与风险", "行业发展趋势与核心风险提示", ACCENT_BLUE),
        ("08", "投资建议", "综合评估与投资策略", ACCENT_GOLD),
    ]

    for i, (num, title, desc, accent) in enumerate(sections):
        y = Inches(1.3) + Inches(i * 0.72)
        # 编号圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.35), Inches(0.35),
                     num, Pt(10), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        # 标题
        add_text_box(slide, Inches(1.4), y, Inches(4), Inches(0.35),
                     title, Pt(14), TEXT_WHITE, bold=True)
        # 描述
        add_text_box(slide, Inches(1.4), y + Inches(0.33), Inches(5), Inches(0.25),
                     desc, Pt(10), TEXT_DIM)
        # 连接线
        add_line(slide, Inches(0.97), y + Inches(0.4), Inches(0.97), y + Inches(0.6), BORDER_COLOR, Pt(1))

    add_page_number(slide, 2)


# ============================================================
# Slide 3: 行业概览 - 市场空间
# ============================================================
def make_market_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "行业概览：AI芯片黄金时代", "全球AI产业从训练军备竞赛向推理应用落地范式切换")

    # KPI 卡片行
    add_kpi_card(slide, Inches(0.6), Inches(1.5), Inches(2.8), Inches(1.4),
                 "2027全球AI芯片市场", "$1,194亿", "Gartner预测(美元)", ACCENT_BLUE)
    add_kpi_card(slide, Inches(3.7), Inches(1.5), Inches(2.8), Inches(1.4),
                 "2027中国AI总投资", "$400亿+", "IDC预测(美元)", ACCENT_GOLD)
    add_kpi_card(slide, Inches(6.8), Inches(1.5), Inches(2.8), Inches(1.4),
                 "2024加速服务器YoY", "+134%", "规模$221亿(美元)", ACCENT_GREEN)
    add_kpi_card(slide, Inches(9.9), Inches(1.5), Inches(2.8), Inches(1.4),
                 "2029端侧AI市场", "1.22万亿元", "CAGR 39.6%(人民币)", ACCENT_BLUE)

    # 右侧：北美云厂商资本开支
    add_rounded_rect(slide, Inches(7.2), Inches(3.2), Inches(5.5), Inches(3.6),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(7.5), Inches(3.35), Inches(5), Inches(0.35),
                 "北美四大云服务商2026资本开支", Pt(12), TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(7.5), Inches(3.7), Inches(5), Inches(0.25),
                 "合计突破 $7000亿，算力基建投资创历史新高（含AI基建全口径）", Pt(9), TEXT_DIM)

    capex_data = [
        ("谷歌 Google", "$750亿", "AI基础设施+TPU扩产", ACCENT_BLUE),
        ("微软 Microsoft", "$800亿", "Azure AI+Copilot算力", ACCENT_BLUE),
        ("Meta", "$600亿+", "开源大模型+数据中心", ACCENT_BLUE),
        ("亚马逊 AWS", "$1,000亿+", "Trainium/Inferentia自研", ACCENT_BLUE),
    ]
    for i, (name, amount, detail, color) in enumerate(capex_data):
        y = Inches(4.15) + Inches(i * 0.6)
        add_text_box(slide, Inches(7.5), y, Inches(2.2), Inches(0.3),
                     name, Pt(10), TEXT_GRAY)
        add_text_box(slide, Inches(9.8), y, Inches(1.5), Inches(0.3),
                     amount, Pt(13), ACCENT_GOLD, bold=True, font_name=FONT_EN)
        add_text_box(slide, Inches(11.3), y, Inches(1.5), Inches(0.3),
                     detail, Pt(9), TEXT_DIM)

    # 左侧：行业发展阶段
    add_rounded_rect(slide, Inches(0.6), Inches(3.2), Inches(6.3), Inches(3.6),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(3.35), Inches(5), Inches(0.35),
                 "A股AI芯片板块发展三阶段", Pt(12), TEXT_WHITE, bold=True)

    stages = [
        ("Phase 1", "概念炒作期", "AI芯片概念火热，估值泡沫化", "2023-2024初", TEXT_DIM),
        ("Phase 2", "订单博弈期", "国产替代订单落地，营收快速增长", "2024-2025", ACCENT_BLUE),
        ("Phase 3", "业绩兑现期", "头部公司盈利拐点确认", "2025至今", ACCENT_GREEN),
    ]
    for i, (phase, name, desc, period, color) in enumerate(stages):
        y = Inches(3.9) + Inches(i * 0.9)
        # 阶段标签
        add_rounded_rect(slide, Inches(0.9), y, Inches(1.1), Inches(0.3),
                         fill_color=color, border_color=color)
        add_text_box(slide, Inches(0.9), y, Inches(1.1), Inches(0.3),
                     phase, Pt(9), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        # 名称
        add_text_box(slide, Inches(2.2), y, Inches(2), Inches(0.3),
                     name, Pt(11), TEXT_WHITE, bold=True)
        # 时间
        add_text_box(slide, Inches(5.5), y, Inches(1.2), Inches(0.3),
                     period, Pt(9), TEXT_DIM, alignment=PP_ALIGN.RIGHT)
        # 描述
        add_text_box(slide, Inches(2.2), y + Inches(0.3), Inches(4.5), Inches(0.3),
                     desc, Pt(9), TEXT_GRAY)
        # 连接箭头
        if i < 2:
            add_rect(slide, Inches(1.4), y + Inches(0.65), Pt(2), Inches(0.2), fill_color=ACCENT_BLUE_DIM)

    add_page_number(slide, 3)


# ============================================================
# Slide 4: 产业链结构
# ============================================================
def make_industry_chain():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "产业链全景图", "上中下游协同，国产AI芯片生态逐步成型")

    # 三个层级卡片
    layers = [
        {
            "title": "上游 · 芯片设计",
            "color": ACCENT_BLUE,
            "items": [
                ("寒武纪", "云端/边端AI芯片，思元系列"),
                ("海光信息", "CPU+DCU，x86生态"),
                ("景嘉微", "GPU+边端AI SoC"),
                ("华为昇腾", "昇腾系列，云端训练推理"),
            ]
        },
        {
            "title": "中游 · 整机制造",
            "color": ACCENT_GOLD,
            "items": [
                ("工业富联", "AI服务器代工龙头"),
                ("中科曙光", "智算中心方案商"),
                ("浪潮信息", "AI服务器制造商"),
                ("紫光股份", "ICT基础设施"),
            ]
        },
        {
            "title": "下游 · 应用落地",
            "color": ACCENT_GREEN,
            "items": [
                ("智算中心", "政府/城市级算力基建"),
                ("算力租赁", "云服务商推理算力"),
                ("终端应用", "手机/PC/汽车/机器人"),
                ("行业方案", "金融/教育/医疗AI"),
            ]
        }
    ]

    for i, layer in enumerate(layers):
        left = Inches(0.6) + Inches(i * 4.2)
        # 卡片背景
        add_rounded_rect(slide, left, Inches(1.5), Inches(3.8), Inches(5.3),
                         fill_color=BG_CARD, border_color=BORDER_COLOR)
        # 顶部色条
        add_rect(slide, left, Inches(1.5), Inches(3.8), Pt(3), fill_color=layer["color"])
        # 标题
        add_text_box(slide, left + Inches(0.2), Inches(1.65), Inches(3.4), Inches(0.4),
                     layer["title"], Pt(13), layer["color"], bold=True)

        # 项目列表
        for j, (name, desc) in enumerate(layer["items"]):
            y = Inches(2.2) + Inches(j * 1.1)
            # 小圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), y + Inches(0.07), Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = layer["color"]
            dot.line.fill.background()
            # 名称
            add_text_box(slide, left + Inches(0.5), y, Inches(3), Inches(0.3),
                         name, Pt(12), TEXT_WHITE, bold=True)
            # 描述
            add_text_box(slide, left + Inches(0.5), y + Inches(0.32), Inches(3), Inches(0.3),
                         desc, Pt(9), TEXT_DIM)

        # 连接箭头
        if i < 2:
            arrow_left = left + Inches(3.9)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.8), Inches(0.25), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE_DIM
            arrow.line.fill.background()

    add_page_number(slide, 4)


# ============================================================
# Slide 5: 寒武纪 - 公司概况与业绩
# ============================================================
def make_cambricon_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "寒武纪（688256）—— AI芯片绝对龙头", "A股第一高价股，国产AI芯片标杆，2025年首次年度盈利")

    # 左侧：核心标签
    add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(4), Inches(2.2),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(1.6), Inches(3.4), Inches(0.3),
                 "核心标签", Pt(11), ACCENT_GOLD, bold=True)
    tags_info = [
        "A股第一高价股 · 市值突破7000亿",
        "2025年首次年度盈利，扭亏为盈",
        "云边端全品类覆盖，全球少数",
        "适配365款主流大模型",
    ]
    for i, tag in enumerate(tags_info):
        add_text_box(slide, Inches(0.9), Inches(2.0) + Inches(i * 0.35), Inches(3.4), Inches(0.3),
                     f"▸ {tag}", Pt(10), TEXT_GRAY)

    # 右侧：2025年年报核心数据
    add_rounded_rect(slide, Inches(5.0), Inches(1.5), Inches(3.9), Inches(2.2),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(5.3), Inches(1.6), Inches(3.4), Inches(0.3),
                 "2025年年报", Pt(11), ACCENT_BLUE, bold=True)
    annual_items = [
        ("营收", "64.97亿", "+453.21%", ACCENT_GREEN),
        ("归母净利润", "20.59亿", "首次盈利", ACCENT_GOLD),
        ("扣非净利润", "17.70亿", "扭亏", ACCENT_GREEN),
        ("毛利率", "55.22%", "高毛利", ACCENT_BLUE),
    ]
    for i, (label, val, change, color) in enumerate(annual_items):
        y = Inches(2.0) + Inches(i * 0.4)
        add_text_box(slide, Inches(5.3), y, Inches(1.3), Inches(0.3),
                     label, Pt(10), TEXT_DIM)
        add_text_box(slide, Inches(6.6), y, Inches(1.0), Inches(0.3),
                     val, Pt(11), TEXT_WHITE, bold=True, font_name=FONT_EN)
        add_text_box(slide, Inches(7.7), y, Inches(1.0), Inches(0.3),
                     change, Pt(9), color)

    # 最右侧：2026Q1数据
    add_rounded_rect(slide, Inches(9.2), Inches(1.5), Inches(3.7), Inches(2.2),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(9.5), Inches(1.6), Inches(3.2), Inches(0.3),
                 "2026年一季报", Pt(11), ACCENT_BLUE, bold=True)
    q1_items = [
        ("营收", "28.85亿", "+159.56%", ACCENT_GREEN),
        ("归母净利润", "10.13亿", "+185.04%", ACCENT_GREEN),
        ("扣非净利润", "—", "+238.56%", ACCENT_GREEN),
        ("股价", "1699.96元", "历史新高", ACCENT_GOLD),
    ]
    for i, (label, val, change, color) in enumerate(q1_items):
        y = Inches(2.0) + Inches(i * 0.4)
        add_text_box(slide, Inches(9.5), y, Inches(1.2), Inches(0.3),
                     label, Pt(10), TEXT_DIM)
        add_text_box(slide, Inches(10.7), y, Inches(1.1), Inches(0.3),
                     val, Pt(11), TEXT_WHITE, bold=True, font_name=FONT_EN)
        add_text_box(slide, Inches(11.8), y, Inches(1.0), Inches(0.3),
                     change, Pt(9), color)

    # 底部：产品线概览
    add_rounded_rect(slide, Inches(0.6), Inches(4.0), Inches(12.3), Inches(2.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(4.1), Inches(3), Inches(0.3),
                 "三大产品线", Pt(11), TEXT_WHITE, bold=True)

    products = [
        ("云端产品线（思元系列）", "营收占比最大\nMLU370等旗舰产品\n已适配365款主流大模型\nDeepSeek/Qwen3/混元等", ACCENT_BLUE),
        ("智能计算集群系统", "毛利率约70%\n政府算力项目核心供应商\n智算中心一体化方案\n软硬一体交付", ACCENT_GOLD),
        ("边缘产品线（思元220）", "累计销量超百万片\n面向推理和端侧部署\n低功耗高性价比\nIoT/终端AI场景", ACCENT_GREEN),
    ]
    for i, (name, desc, color) in enumerate(products):
        left = Inches(0.9) + Inches(i * 4.1)
        add_section_bar(slide, left, Inches(4.5), Inches(0.04), Inches(0.35), color)
        add_text_box(slide, left + Inches(0.15), Inches(4.5), Inches(3.5), Inches(0.3),
                     name, Pt(11), TEXT_WHITE, bold=True)
        for j, line in enumerate(desc.split("\n")):
            add_text_box(slide, left + Inches(0.15), Inches(4.85) + Inches(j * 0.3), Inches(3.5), Inches(0.3),
                         f"▸ {line}", Pt(9), TEXT_GRAY)

    add_page_number(slide, 5)


# ============================================================
# Slide 6: 寒武纪 - 核心优势与风险
# ============================================================
def make_cambricon_analysis():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "寒武纪：核心竞争力深度分析", "全品类覆盖 + 开源生态 + 规模化部署")

    # 左侧：核心优势
    add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(7.2), Inches(5.3),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(1.6), Inches(3), Inches(0.35),
                 "核心竞争优势", Pt(13), ACCENT_GREEN, bold=True)

    advantages = [
        ("云边端全品类覆盖",
         "全球少数同时具备云端训练/推理、边缘计算、终端AI芯片研发能力的公司。"
         "产品矩阵完整，可满足从数据中心到终端设备的全场景算力需求，"
         "形成软硬件协同的生态壁垒。"),
        ("开源生态深度适配",
         "紧密跟进开源大模型生态，已适配365款主流大模型，包括DeepSeek、Qwen3、混元等。"
         "显著降低客户从英伟达生态迁移的成本，提高切换意愿。"),
        ("规模化行业部署",
         "产品在运营商、金融、互联网等重点行业实现规模化部署。"
         "智能计算集群系统毛利率约70%，在政府算力项目中成为核心供应商。"),
        ("盈利拐点确认",
         "2025年首次实现年度盈利（归母净利润20.59亿），"
         "2026Q1延续高增长态势（营收+159.56%，净利润+185.04%），"
         "标志着从研发投入期进入收获期。"),
    ]
    for i, (title, desc) in enumerate(advantages):
        y = Inches(2.1) + Inches(i * 1.15)
        # 编号
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.03), Inches(0.28), Inches(0.28))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT_BLUE
        num_shape.line.fill.background()
        add_text_box(slide, Inches(0.9), y + Inches(0.03), Inches(0.28), Inches(0.28),
                     str(i+1), Pt(10), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        # 标题
        add_text_box(slide, Inches(1.35), y, Inches(6), Inches(0.3),
                     title, Pt(11), TEXT_WHITE, bold=True)
        # 描述
        add_text_box(slide, Inches(1.35), y + Inches(0.3), Inches(6.2), Inches(0.7),
                     desc, Pt(9), TEXT_GRAY, line_spacing=1.25)

    # 右侧：风险提示
    add_rounded_rect(slide, Inches(8.1), Inches(1.5), Inches(4.8), Inches(5.3),
                     fill_color=BG_CARD, border_color=ACCENT_RED)
    add_text_box(slide, Inches(8.4), Inches(1.6), Inches(3), Inches(0.35),
                 "风险提示", Pt(13), ACCENT_RED, bold=True)

    risks = [
        ("客户集中度较高",
         "前几大客户贡献营收占比偏高，若大客户订单波动，将直接影响业绩表现。"),
        ("先进制程代工受限",
         "国内先进制程产能（5-7nm）与台积电（2-3nm）存在代差，影响芯片性能竞争力。"),
        ("供应链依赖",
         "芯片设计与制造分离，代工产能、封装测试等环节存在不确定性。"),
        ("高估值风险",
         "当前市值超7000亿元，PE估值极高，市场预期已较为充分，需持续超预期才能支撑。"),
    ]
    for i, (title, desc) in enumerate(risks):
        y = Inches(2.1) + Inches(i * 1.15)
        add_text_box(slide, Inches(8.4), y, Inches(4.2), Inches(0.3),
                     f"⚠ {title}", Pt(10), ACCENT_RED, bold=True)
        add_text_box(slide, Inches(8.4), y + Inches(0.3), Inches(4.2), Inches(0.7),
                     desc, Pt(9), TEXT_GRAY, line_spacing=1.25)

    add_page_number(slide, 6)


# ============================================================
# Slide 7: 海光信息 - 公司概况
# ============================================================
def make_hygon_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "海光信息（688041）—— CPU+DCU双轮驱动", "中国版AMD，x86生态稀缺标的，2025年营收首破百亿")

    # 核心数据卡片行
    add_kpi_card(slide, Inches(0.6), Inches(1.5), Inches(2.9), Inches(1.35),
                 "2025年营收", "143.77亿", "YoY +56.92%", ACCENT_BLUE)
    add_kpi_card(slide, Inches(3.8), Inches(1.5), Inches(2.9), Inches(1.35),
                 "2025年归母净利润", "25.45亿", "YoY +31.79%", ACCENT_GREEN)
    add_kpi_card(slide, Inches(7.0), Inches(1.5), Inches(2.9), Inches(1.35),
                 "2026Q1营收", "40.34亿", "YoY +68.06%", ACCENT_BLUE)
    add_kpi_card(slide, Inches(10.2), Inches(1.5), Inches(2.7), Inches(1.35),
                 "总市值", "5782亿", "PE ~207x", ACCENT_GOLD)

    # 左侧：双产品线
    add_rounded_rect(slide, Inches(0.6), Inches(3.1), Inches(6.2), Inches(3.7),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(3.2), Inches(5), Inches(0.35),
                 "双产品线矩阵", Pt(12), TEXT_WHITE, bold=True)

    # CPU卡片
    add_rounded_rect(slide, Inches(0.9), Inches(3.7), Inches(5.6), Inches(1.4),
                     fill_color=BG_CARD2, border_color=ACCENT_BLUE_DIM)
    add_text_box(slide, Inches(1.1), Inches(3.75), Inches(2), Inches(0.3),
                 "CPU处理器（7000/5000/3000系列）", Pt(10), ACCENT_BLUE, bold=True)
    cpu_lines = [
        "▸ 兼容x86指令集，客户零迁移成本",
        "▸ 中国唯一自主可控的x86 CPU",
        "▸ 面向服务器/工作站/桌面场景",
    ]
    for i, line in enumerate(cpu_lines):
        add_text_box(slide, Inches(1.1), Inches(4.1) + Inches(i * 0.28), Inches(5), Inches(0.25),
                     line, Pt(9), TEXT_GRAY)

    # DCU卡片
    add_rounded_rect(slide, Inches(0.9), Inches(5.2), Inches(5.6), Inches(1.4),
                     fill_color=BG_CARD2, border_color=ACCENT_GOLD)
    add_text_box(slide, Inches(1.1), Inches(5.25), Inches(2), Inches(0.3),
                 "DCU加速器（深算系列）", Pt(10), ACCENT_GOLD, bold=True)
    dcu_lines = [
        "▸ 兼容CUDA生态，算子覆盖度超99%",
        "▸ 支持千亿级大模型训练与推理",
        "▸ 与CPU形成C86+GPGPU自研矩阵",
    ]
    for i, line in enumerate(dcu_lines):
        add_text_box(slide, Inches(1.1), Inches(5.6) + Inches(i * 0.28), Inches(5), Inches(0.25),
                     line, Pt(9), TEXT_GRAY)

    # 右侧：关键经营指标
    add_rounded_rect(slide, Inches(7.1), Inches(3.1), Inches(5.8), Inches(3.7),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(7.4), Inches(3.2), Inches(5), Inches(0.35),
                 "2025年经营指标", Pt(12), TEXT_WHITE, bold=True)

    metrics = [
        ("毛利率", "57.78%", "同比降5.92pct，因上游涨价", TEXT_DIM),
        ("经营现金流", "20.97亿", "同比+114.61%，盈利质量提升", ACCENT_GREEN),
        ("扣非净利润", "23.05亿", "核心业务盈利能力扎实", ACCENT_BLUE),
        ("前五大客户集中度", ">90%", "中科曙光为最大客户(56.68%)", ACCENT_RED),
    ]
    for i, (label, value, note, color) in enumerate(metrics):
        y = Inches(3.7) + Inches(i * 0.75)
        add_text_box(slide, Inches(7.4), y, Inches(2), Inches(0.25),
                     label, Pt(10), TEXT_DIM)
        add_text_box(slide, Inches(7.4), y + Inches(0.25), Inches(2), Inches(0.3),
                     value, Pt(14), TEXT_WHITE, bold=True, font_name=FONT_EN)
        add_text_box(slide, Inches(9.5), y + Inches(0.25), Inches(3.2), Inches(0.4),
                     note, Pt(9), color)

    add_page_number(slide, 7)


# ============================================================
# Slide 8: 海光信息 - 核心优势
# ============================================================
def make_hygon_advantages():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "海光信息：生态壁垒与核心优势", "光合组织 + 中科曙光协同 + x86稀缺性")

    # 三大优势卡片
    cards = [
        {
            "title": "中国唯一的 C86+GPGPU 自研矩阵",
            "color": ACCENT_BLUE,
            "points": [
                "CPU兼容x86指令集，DCU兼容CUDA生态",
                "算子覆盖度超99%，客户迁移成本极低",
                "覆盖从通用计算到AI加速全场景",
                "国内唯一同时拥有CPU+GPU自研能力的上市公司",
            ]
        },
        {
            "title": "光合组织生态联盟",
            "color": ACCENT_GOLD,
            "points": [
                "聚合超6,000家合作伙伴",
                "完成15,000+项软硬件适配",
                "构建国产算力生态护城河",
                "从芯片到应用的全栈生态支持",
            ]
        },
        {
            "title": "中科曙光协同效应",
            "color": ACCENT_GREEN,
            "points": [
                "中科曙光为第一大股东（持股27.96%）",
                "中科曙光同时为最大客户（占营收56.68%）",
                "智算中心一体化方案协同",
                "从芯片设计到系统集成的产业链闭环",
            ]
        }
    ]

    for i, card in enumerate(cards):
        left = Inches(0.6) + Inches(i * 4.2)
        add_rounded_rect(slide, left, Inches(1.5), Inches(3.8), Inches(3.5),
                         fill_color=BG_CARD, border_color=BORDER_COLOR)
        add_rect(slide, left, Inches(1.5), Inches(3.8), Pt(3), fill_color=card["color"])
        add_text_box(slide, left + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.6),
                     card["title"], Pt(12), card["color"], bold=True, line_spacing=1.3)
        for j, point in enumerate(card["points"]):
            y = Inches(2.5) + Inches(j * 0.55)
            add_text_box(slide, left + Inches(0.2), y, Inches(3.4), Inches(0.5),
                         f"▸ {point}", Pt(9), TEXT_GRAY, line_spacing=1.2)

    # 底部风险提示
    add_rounded_rect(slide, Inches(0.6), Inches(5.3), Inches(12.3), Inches(1.6),
                     fill_color=BG_CARD, border_color=ACCENT_RED)
    add_text_box(slide, Inches(0.9), Inches(5.4), Inches(3), Inches(0.3),
                 "风险提示", Pt(11), ACCENT_RED, bold=True)

    risk_cols = [
        ("估值偏高", "PE约207倍，反映稀缺性溢价，\n需持续高增长消化估值"),
        ("客户集中度", "前五大客户集中度超90%，\n对中科曙光依赖度极高"),
        ("代工产能", "先进制程代工受限，\n上游产能直接影响交付"),
        ("毛利率承压", "2025年毛利率同比下降5.92pct，\n上游涨价压力持续"),
    ]
    for i, (title, desc) in enumerate(risk_cols):
        left = Inches(0.9) + Inches(i * 3.05)
        add_text_box(slide, left, Inches(5.8), Inches(2.8), Inches(0.25),
                     f"⚠ {title}", Pt(10), ACCENT_RED, bold=True)
        add_text_box(slide, left, Inches(6.05), Inches(2.8), Inches(0.6),
                     desc, Pt(9), TEXT_GRAY, line_spacing=1.2)

    add_page_number(slide, 8)


# ============================================================
# Slide 9: 景嘉微 - 公司概况
# ============================================================
def make_jingjia_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "景嘉微（300474）—— GPU+边端AI SoC双线转型", "从军用GPU向全栈AI算力转型，高研发投入期")

    # 核心数据
    add_kpi_card(slide, Inches(0.6), Inches(1.5), Inches(3.0), Inches(1.3),
                 "2025年营收", "7.20亿", "YoY +54.41%", ACCENT_BLUE)
    add_kpi_card(slide, Inches(3.9), Inches(1.5), Inches(3.0), Inches(1.3),
                 "2025年净利润", "-1.65亿", "连续两年亏损", ACCENT_RED)
    add_kpi_card(slide, Inches(7.2), Inches(1.5), Inches(3.0), Inches(1.3),
                 "研发投入占比", "59.39%", "研发费用4.28亿", ACCENT_GOLD)
    add_kpi_card(slide, Inches(10.5), Inches(1.5), Inches(2.4), Inches(1.3),
                 "2026Q1营收", "8414万", "YoY -17.87%", ACCENT_RED)

    # 左侧：产品线
    add_rounded_rect(slide, Inches(0.6), Inches(3.1), Inches(6.2), Inches(3.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(3.2), Inches(5), Inches(0.3),
                 "三大产品线", Pt(12), TEXT_WHITE, bold=True)

    prods = [
        ("JM11系列GPU", "已进入小批量交付阶段",
         "政务云桌面、工业设计等场景应用\n国产GPU虚拟化方案已在电网等关键基础设施部署\n图形显控模块营收4.51亿，毛利率56.90%", ACCENT_BLUE),
        ("CH37系列边端AI SoC", "面向具身智能时代",
         "64TOPS@INT8峰值算力\n面向机器人/无人机/具身智能场景\n控股诚恒微切入边端AI芯片领域", ACCENT_GOLD),
        ("图形显控模块", "传统业务基石",
         "营收4.51亿元，毛利率56.90%\n军用GPU市场龙头地位\n持续贡献现金流支持研发", ACCENT_GREEN),
    ]
    for i, (name, subtitle, desc, color) in enumerate(prods):
        y = Inches(3.6) + Inches(i * 1.1)
        add_section_bar(slide, Inches(0.9), y, Inches(0.04), Inches(0.35), color)
        add_text_box(slide, Inches(1.05), y, Inches(2.5), Inches(0.25),
                     name, Pt(10), TEXT_WHITE, bold=True)
        add_text_box(slide, Inches(3.6), y, Inches(3), Inches(0.25),
                     subtitle, Pt(9), color)
        add_text_box(slide, Inches(1.05), y + Inches(0.3), Inches(5.5), Inches(0.7),
                     desc, Pt(9), TEXT_DIM, line_spacing=1.2)

    # 右侧：2026Q1数据
    add_rounded_rect(slide, Inches(7.1), Inches(3.1), Inches(5.8), Inches(3.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(7.4), Inches(3.2), Inches(5), Inches(0.3),
                 "2026年一季报数据", Pt(12), TEXT_WHITE, bold=True)

    q1_data = [
        ("营收", "8414万元", "同比下降17.87%", ACCENT_RED),
        ("净利润", "亏损5608万元", "亏损幅度扩大", ACCENT_RED),
        ("研发费用", "1.42亿元", "同比增长102.67%", ACCENT_GOLD),
        ("研发占营收比", "168.5%", "研发投入强度极高", ACCENT_GOLD),
    ]
    for i, (label, val, note, color) in enumerate(q1_data):
        y = Inches(3.7) + Inches(i * 0.7)
        add_text_box(slide, Inches(7.4), y, Inches(1.8), Inches(0.25),
                     label, Pt(10), TEXT_DIM)
        add_text_box(slide, Inches(7.4), y + Inches(0.25), Inches(1.8), Inches(0.3),
                     val, Pt(13), TEXT_WHITE, bold=True, font_name=FONT_EN)
        add_text_box(slide, Inches(9.3), y + Inches(0.25), Inches(3.4), Inches(0.3),
                     note, Pt(9), color)

    # 底部提示
    add_text_box(slide, Inches(7.4), Inches(6.5), Inches(5.2), Inches(0.3),
                 "提示：高研发投入导致短期亏损，关注JM11量产节奏和CH37商业化进展",
                 Pt(9), ACCENT_GOLD)

    add_page_number(slide, 9)


# ============================================================
# Slide 10: 景嘉微 - 战略分析
# ============================================================
def make_jingjia_strategy():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "景嘉微：转型战略与投资逻辑", "云-边-端全栈布局，短期阵痛换长期空间")

    # 战略路径图
    add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(12.3), Inches(2.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5), Inches(0.3),
                 "转型路径：从军用GPU到全栈AI算力", Pt(12), TEXT_WHITE, bold=True)

    path_stages = [
        ("Phase 1", "军用GPU\n图形显控", "传统优势业务\n营收基石\n毛利率56.90%", ACCENT_GREEN, "已实现"),
        ("Phase 2", "国产GPU\nJM11系列", "小批量交付\n政务/工业场景\n虚拟化方案落地", ACCENT_BLUE, "进行中"),
        ("Phase 3", "边端AI SoC\nCH37系列", "64TOPS算力\n机器人/无人机\n具身智能", ACCENT_GOLD, "研发中"),
        ("Phase 4", "全栈AI\n算力平台", "云-边-端一体\n全场景覆盖\n生态闭环", RGBColor(0xA7, 0x8B, 0xFA), "远期目标"),
    ]
    for i, (phase, name, desc, color, status) in enumerate(path_stages):
        left = Inches(0.9) + Inches(i * 3.0)
        # 阶段卡片
        add_rounded_rect(slide, left, Inches(2.1), Inches(2.6), Inches(1.9),
                         fill_color=BG_CARD2, border_color=color)
        add_text_box(slide, left + Inches(0.15), Inches(2.15), Inches(2.3), Inches(0.25),
                     f"{phase}: {status}", Pt(9), color, bold=True)
        add_text_box(slide, left + Inches(0.15), Inches(2.4), Inches(2.3), Inches(0.5),
                     name, Pt(11), TEXT_WHITE, bold=True, line_spacing=1.2)
        add_text_box(slide, left + Inches(0.15), Inches(2.95), Inches(2.3), Inches(0.9),
                     desc, Pt(9), TEXT_DIM, line_spacing=1.3)
        # 箭头
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           left + Inches(2.65), Inches(2.85), Inches(0.3), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE_DIM
            arrow.line.fill.background()

    # 投资逻辑 vs 风险
    add_rounded_rect(slide, Inches(0.6), Inches(4.6), Inches(5.9), Inches(2.4),
                     fill_color=BG_CARD, border_color=ACCENT_GREEN)
    add_text_box(slide, Inches(0.9), Inches(4.7), Inches(5), Inches(0.3),
                 "投资逻辑（看多）", Pt(12), ACCENT_GREEN, bold=True)
    bull_points = [
        "国产GPU稀缺标的，全栈布局具长期价值",
        "CH37边端AI芯片面向具身智能蓝海市场",
        "高研发投入将转化为产品竞争力",
        "军用业务提供现金流安全垫",
        "控股诚恒微，边端AI芯片版图扩展",
    ]
    for i, point in enumerate(bull_points):
        add_text_box(slide, Inches(0.9), Inches(5.1) + Inches(i * 0.35), Inches(5.3), Inches(0.3),
                     f"✦ {point}", Pt(9), TEXT_GRAY)

    add_rounded_rect(slide, Inches(6.8), Inches(4.6), Inches(5.9), Inches(2.4),
                     fill_color=BG_CARD, border_color=ACCENT_RED)
    add_text_box(slide, Inches(7.1), Inches(4.7), Inches(5), Inches(0.3),
                 "风险因素（看空）", Pt(12), ACCENT_RED, bold=True)
    bear_points = [
        "连续两年亏损，盈利拐点尚未确认",
        "2026Q1营收同比下降，转型阵痛持续",
        "研发投入占营收168.5%，盈利遥遥无期",
        "前五大客户集中度64.53%",
        "JM11量产进度和CH37商业化存在不确定性",
    ]
    for i, point in enumerate(bear_points):
        add_text_box(slide, Inches(7.1), Inches(5.1) + Inches(i * 0.35), Inches(5.3), Inches(0.3),
                     f"✦ {point}", Pt(9), TEXT_GRAY)

    add_page_number(slide, 10)


# ============================================================
# Slide 11: 工业富联
# ============================================================
def make_foxconn():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "工业富联（601138）—— AI服务器制造龙头", "产业链中游核心受益者，受益于全球AI算力基建浪潮")

    # 大数字KPI
    add_kpi_card(slide, Inches(0.6), Inches(1.5), Inches(3.8), Inches(1.4),
                 "2025年营收", "9,028.87亿", "YoY +48.22%", ACCENT_BLUE)
    add_kpi_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(1.4),
                 "2025年净利润", "352.86亿", "YoY +51.99%", ACCENT_GREEN)
    add_kpi_card(slide, Inches(8.8), Inches(1.5), Inches(4.1), Inches(1.4),
                 "2026Q1净利润", "105.95亿", "YoY +102.55%", ACCENT_GOLD)

    # 2026Q1亮点
    add_rounded_rect(slide, Inches(0.6), Inches(3.2), Inches(6.0), Inches(3.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(3.3), Inches(5), Inches(0.3),
                 "2026年一季报：超预期增长", Pt(12), TEXT_WHITE, bold=True)

    q1_points = [
        ("营收 2,511亿元", "同比增长56.52%，AI服务器需求持续旺盛", ACCENT_BLUE),
        ("净利润 105.95亿元", "同比翻倍增长（+102.55%），利润率显著提升", ACCENT_GREEN),
        ("AI服务器占比提升", "AI相关业务已成为增长核心驱动力", ACCENT_GOLD),
        ("全球AI算力基建受益", "北美云厂商资本开支合计超7000亿美元", ACCENT_BLUE),
        ("产能持续扩张", "AI服务器产能利用率维持高位", ACCENT_GREEN),
    ]
    for i, (title, desc, color) in enumerate(q1_points):
        y = Inches(3.8) + Inches(i * 0.6)
        add_text_box(slide, Inches(0.9), y, Inches(5.5), Inches(0.25),
                     f"▸ {title}", Pt(10), color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.25), Inches(5.3), Inches(0.25),
                     desc, Pt(9), TEXT_GRAY)

    # 右侧：投资要点
    add_rounded_rect(slide, Inches(6.9), Inches(3.2), Inches(6.0), Inches(3.8),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(7.2), Inches(3.3), Inches(5), Inches(0.3),
                 "投资要点", Pt(12), TEXT_WHITE, bold=True)

    invest_points = [
        ("全球AI服务器代工龙头", "占据全球AI服务器代工最大市场份额，深度绑定英伟达GPU供应链"),
        ("营收规模碾压同行", "9,028亿营收规模远超其他A股AI概念公司，体量优势明显"),
        ("业绩确定性最高", "作为中游制造环节，直接受益于全球AI算力投资，业绩确定性优于上游芯片设计"),
        ("估值相对合理", "相比寒武纪/海光的200倍+PE，工业富联估值更具安全边际"),
    ]
    for i, (title, desc) in enumerate(invest_points):
        y = Inches(3.8) + Inches(i * 0.8)
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y + Inches(0.05), Inches(0.22), Inches(0.22))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT_BLUE
        num_shape.line.fill.background()
        add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(0.22), Inches(0.22),
                     str(i+1), Pt(9), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, Inches(7.55), y, Inches(5.2), Inches(0.25),
                     title, Pt(10), TEXT_WHITE, bold=True)
        add_text_box(slide, Inches(7.55), y + Inches(0.28), Inches(5.2), Inches(0.45),
                     desc, Pt(9), TEXT_GRAY, line_spacing=1.2)

    add_page_number(slide, 11)


# ============================================================
# Slide 12: 四家公司财务对比
# ============================================================
def make_financial_comparison():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "四大标的财务数据对比", "2025年年报核心指标横向比较")

    # 对比表格
    headers = ["指标", "寒武纪", "海光信息", "景嘉微", "工业富联"]
    col_widths = [Inches(2.5), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.6)]
    rows = [
        ["2025营收(亿元)", "64.97", "143.77", "7.20", "9,028.87"],
        ["营收YoY", "+453.21%", "+56.92%", "+54.41%", "+48.22%"],
        ["归母净利润(亿元)", "20.59", "25.45", "-1.65", "352.86"],
        ["净利润YoY", "扭亏为盈", "+31.79%", "亏损扩大", "+51.99%"],
        ["毛利率", "55.22%", "57.78%", "—", "—"],
        ["扣非净利润(亿元)", "17.70", "23.05", "—", "—"],
        ["2026Q1营收(亿元)", "28.85", "40.34", "0.84", "2,511"],
        ["2026Q1营收YoY", "+159.56%", "+68.06%", "-17.87%", "+56.52%"],
        ["2026Q1净利润(亿元)", "10.13", "6.87", "-0.56", "105.95"],
        ["核心标签", "AI芯片龙头", "CPU+DCU", "GPU转型", "服务器代工"],
    ]
    add_table(slide, Inches(0.6), Inches(1.5), Inches(12.3), col_widths, headers, rows)

    # 底部解读
    add_rounded_rect(slide, Inches(0.6), Inches(5.7), Inches(12.3), Inches(1.2),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(5.8), Inches(11.7), Inches(1.0),
                 "关键发现：寒武纪营收增速最快（+453%）且首次盈利，海光营收体量最大（143.77亿元），"
                 "工业富联绝对规模碾压（9028亿元），景嘉微仍处转型阵痛期。"
                 "增速排序：寒武纪 > 海光 > 景嘉微 > 工业富联；规模排序：工业富联 >> 海光 > 寒武纪 > 景嘉微。",
                 Pt(10), TEXT_GRAY, line_spacing=1.3)

    add_page_number(slide, 12)


# ============================================================
# Slide 13: 竞争格局
# ============================================================
def make_competition():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "AI芯片竞争格局全景", "国产替代加速，多赛道并行推进")

    # 四个赛道
    tracks = [
        {
            "title": "云端训练",
            "subtitle": "英伟达绝对霸主",
            "leader": "英伟达 H100/B200",
            "challengers": "寒武纪思元 / 海光DCU\n华为昇腾910B",
            "color": ACCENT_RED,
            "gap": "代差明显，追赶中"
        },
        {
            "title": "云端推理",
            "subtitle": "三足鼎立格局",
            "leader": "寒武纪 / 海光DCU / 昇腾",
            "challengers": "壁仞科技 / 摩尔线程\n天数智芯",
            "color": ACCENT_GOLD,
            "gap": "差距缩小，生态为王"
        },
        {
            "title": "边端AI",
            "subtitle": "百花齐放蓝海",
            "leader": "景嘉微CH37 / 瑞芯微",
            "challengers": "全志科技 / 寒武纪思元220\n华为昇腾310",
            "color": ACCENT_GREEN,
            "gap": "格局未定，快速迭代"
        },
        {
            "title": "GPU通用计算",
            "subtitle": "国产替代核心赛道",
            "leader": "海光（x86兼容）",
            "challengers": "景嘉微JM11 / 摩尔线程\n壁仞科技",
            "color": ACCENT_BLUE,
            "gap": "CUDA兼容是关键"
        }
    ]

    for i, track in enumerate(tracks):
        left = Inches(0.6) + Inches(i * 3.15)
        # 赛道卡片
        add_rounded_rect(slide, left, Inches(1.5), Inches(2.9), Inches(5.3),
                         fill_color=BG_CARD, border_color=BORDER_COLOR)
        # 顶部色条
        add_rect(slide, left, Inches(1.5), Inches(2.9), Pt(3), fill_color=track["color"])

        # 赛道标题
        add_text_box(slide, left + Inches(0.15), Inches(1.7), Inches(2.6), Inches(0.3),
                     track["title"], Pt(13), TEXT_WHITE, bold=True)
        add_text_box(slide, left + Inches(0.15), Inches(2.0), Inches(2.6), Inches(0.25),
                     track["subtitle"], Pt(9), track["color"])

        # 领先者
        add_text_box(slide, left + Inches(0.15), Inches(2.5), Inches(2.6), Inches(0.25),
                     "领先者", Pt(9), TEXT_DIM, bold=True)
        add_rounded_rect(slide, left + Inches(0.15), Inches(2.8), Inches(2.6), Inches(0.4),
                         fill_color=BG_CARD2, border_color=track["color"])
        add_text_box(slide, left + Inches(0.25), Inches(2.85), Inches(2.4), Inches(0.3),
                     track["leader"], Pt(10), TEXT_WHITE, anchor=MSO_ANCHOR.MIDDLE)

        # 追赶者
        add_text_box(slide, left + Inches(0.15), Inches(3.4), Inches(2.6), Inches(0.25),
                     "国产追赶者", Pt(9), TEXT_DIM, bold=True)
        add_rounded_rect(slide, left + Inches(0.15), Inches(3.7), Inches(2.6), Inches(1.0),
                         fill_color=BG_CARD2, border_color=ACCENT_BLUE_DIM)
        add_text_box(slide, left + Inches(0.25), Inches(3.8), Inches(2.4), Inches(0.8),
                     track["challengers"], Pt(9), TEXT_GRAY, line_spacing=1.4)

        # 差距评估
        add_text_box(slide, left + Inches(0.15), Inches(4.9), Inches(2.6), Inches(0.25),
                     "国产化差距", Pt(9), TEXT_DIM, bold=True)
        add_text_box(slide, left + Inches(0.15), Inches(5.2), Inches(2.6), Inches(0.4),
                     track["gap"], Pt(10), track["color"], bold=True)

    add_page_number(slide, 13)


# ============================================================
# Slide 14: 行业趋势
# ============================================================
def make_trends():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "行业四大发展趋势", "国产替代加速，推理放量，端侧崛起，生态卡位")

    trends = [
        {
            "icon": "01",
            "title": "国产替代加速",
            "subtitle": "信创从政务/金融向教育等更多行业渗透",
            "detail": "政策驱动下，信创从党政机关向金融、电信、教育等行业加速渗透。"
                      "国产AI芯片在政府智算中心项目中占比持续提升，"
                      "2025年国产AI芯片在政府算力采购中占比已超60%。",
            "color": ACCENT_BLUE,
        },
        {
            "icon": "02",
            "title": "推理端放量",
            "subtitle": "从训练转向推理应用，对算力成本敏感度提升",
            "detail": "全球AI产业正从训练军备竞赛向推理应用落地范式切换。"
                      "推理场景对算力成本敏感度更高，国产芯片凭借性价比优势，"
                      "在推理端有望实现对英伟达的更大替代。",
            "color": ACCENT_GOLD,
        },
        {
            "icon": "03",
            "title": "端侧AI崛起",
            "subtitle": "2026年端侧AI产品规模化放量关键元年",
            "detail": "弗若斯特沙利文预测2025-2029年全球端侧AI市场规模从3219亿增至1.22万亿元，"
                      "年复合增速39.6%。AI手机、AI PC、机器人等终端设备对边端AI芯片需求爆发。",
            "color": ACCENT_GREEN,
        },
        {
            "icon": "04",
            "title": "生态卡位战",
            "subtitle": "下半场从'有的用'转向'好用'",
            "detail": "AI芯片竞争已从单纯的硬件性能比拼转向生态体系建设。"
                      "软件栈完善度、开发者社区活跃度、客户迁移成本成为竞争关键。"
                      "海光的光合组织、寒武纪的大模型适配都是生态卡位的体现。",
            "color": RGBColor(0xA7, 0x8B, 0xFA),
        }
    ]

    for i, trend in enumerate(trends):
        col = i % 2
        row = i // 2
        left = Inches(0.6) + Inches(col * 6.3)
        top = Inches(1.5) + Inches(row * 2.8)

        add_rounded_rect(slide, left, top, Inches(5.9), Inches(2.5),
                         fill_color=BG_CARD, border_color=BORDER_COLOR)

        # 编号
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2), Inches(0.4), Inches(0.4))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = trend["color"]
        num_shape.line.fill.background()
        add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.4), Inches(0.4),
                     trend["icon"], Pt(12), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # 标题
        add_text_box(slide, left + Inches(0.75), top + Inches(0.2), Inches(4.8), Inches(0.3),
                     trend["title"], Pt(14), TEXT_WHITE, bold=True)
        # 副标题
        add_text_box(slide, left + Inches(0.75), top + Inches(0.55), Inches(4.8), Inches(0.25),
                     trend["subtitle"], Pt(9), trend["color"])
        # 详情
        add_text_box(slide, left + Inches(0.2), top + Inches(0.95), Inches(5.5), Inches(1.3),
                     trend["detail"], Pt(9), TEXT_GRAY, line_spacing=1.35)

    add_page_number(slide, 14)


# ============================================================
# Slide 15: 风险提示
# ============================================================
def make_risks():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "核心风险提示", "高估值、制程受限、地缘政治、客户集中、供应链压力")

    risks = [
        {
            "title": "估值普遍偏高",
            "level": "高",
            "color": ACCENT_RED,
            "detail": "寒武纪PE超200倍，海光PE约207倍，市场预期已较为充分。"
                      "若业绩增速放缓或不及预期，估值回调风险较大。",
            "impact": "寒武纪、海光信息",
        },
        {
            "title": "先进制程代工受限",
            "level": "高",
            "color": ACCENT_RED,
            "detail": "国内先进制程产能（5-7nm）与台积电（2-3nm）存在代差。"
                      "芯片性能直接受制于代工工艺水平，短期内难以突破。",
            "impact": "全行业",
        },
        {
            "title": "地缘政治与出口管制",
            "level": "高",
            "color": ACCENT_RED,
            "detail": "中美科技博弈持续升级，美国对华AI芯片出口管制不断加码。"
                      "可能进一步限制EDA工具、IP授权和先进制程代工，"
                      "对国产AI芯片产业链形成系统性风险。",
            "impact": "全行业",
        },
        {
            "title": "客户集中度高",
            "level": "中",
            "color": ACCENT_GOLD,
            "detail": "海光前五大客户集中度超90%，景嘉微为64.53%。"
                      "大客户订单波动将直接影响业绩表现，议价能力受限。",
            "impact": "海光信息、景嘉微",
        },
        {
            "title": "供应链涨价压力",
            "level": "中",
            "color": ACCENT_GOLD,
            "detail": "存储芯片、封测产能持续紧张，上游涨价压力传导。"
                      "海光2025年毛利率已同比下降5.92个百分点。",
            "impact": "全行业",
        },
    ]

    for i, risk in enumerate(risks):
        left = Inches(0.6) + Inches(i * 2.52)
        add_rounded_rect(slide, left, Inches(1.5), Inches(2.32), Inches(4.5),
                         fill_color=BG_CARD, border_color=risk["color"])

        # 风险等级标签
        add_rounded_rect(slide, left + Inches(0.15), Inches(1.65), Inches(0.8), Inches(0.25),
                         fill_color=risk["color"])
        add_text_box(slide, left + Inches(0.15), Inches(1.65), Inches(0.8), Inches(0.25),
                     f"{risk['level']}风险", Pt(9), BG_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # 标题
        add_text_box(slide, left + Inches(0.15), Inches(2.05), Inches(2.0), Inches(0.4),
                     risk["title"], Pt(12), TEXT_WHITE, bold=True)

        # 详情
        add_text_box(slide, left + Inches(0.15), Inches(2.5), Inches(2.0), Inches(2.2),
                     risk["detail"], Pt(9), TEXT_GRAY, line_spacing=1.4)

        # 影响范围
        add_text_box(slide, left + Inches(0.15), Inches(5.0), Inches(2.0), Inches(0.2),
                     "影响范围", Pt(9), TEXT_DIM, bold=True)
        add_text_box(slide, left + Inches(0.15), Inches(5.25), Inches(2.0), Inches(0.3),
                     risk["impact"], Pt(10), risk["color"], bold=True)

    # 底部风险总结
    add_rounded_rect(slide, Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.7),
                     fill_color=BG_CARD, border_color=ACCENT_RED)
    add_text_box(slide, Inches(0.9), Inches(6.3), Inches(11.7), Inches(0.5),
                 "综合风险评估：当前A股AI芯片板块整体估值偏高，短期回调风险不可忽视。"
                 "建议关注业绩兑现能力（寒武纪、海光）、技术突破进展（景嘉微JM11量产）"
                 "以及估值安全边际（工业富联）。",
                 Pt(10), TEXT_GRAY, line_spacing=1.3)

    add_page_number(slide, 15)


# ============================================================
# Slide 16: 投资建议 - 综合评分
# ============================================================
def make_investment_score():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    slide_title(slide, "综合评估与投资建议", "四维度评分体系，量化对比各标的投资价值")

    # 评分表格
    headers = ["评估维度", "寒武纪", "海光信息", "景嘉微", "工业富联"]
    col_widths = [Inches(2.5), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.6)]
    rows = [
        ["业绩确定性", "★★★★☆", "★★★★☆", "★★☆☆☆", "★★★★★"],
        ["成长性", "★★★★★", "★★★★☆", "★★★☆☆", "★★★★☆"],
        ["估值安全边际", "★★☆☆☆", "★★☆☆☆", "★★★☆☆", "★★★★☆"],
        ["生态壁垒", "★★★★☆", "★★★★★", "★★★☆☆", "★★★☆☆"],
        ["技术领先性", "★★★★★", "★★★★☆", "★★★☆☆", "★★★☆☆"],
        ["综合评分", "19/25", "19/25", "13/25", "19/25"],
        ["投资评级", "增持", "增持", "中性", "增持"],
    ]
    add_table(slide, Inches(0.6), Inches(1.5), Inches(12.3), col_widths, headers, rows)

    # 底部投资建议
    add_rounded_rect(slide, Inches(0.6), Inches(4.4), Inches(12.3), Inches(2.6),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(4.5), Inches(5), Inches(0.3),
                 "分标的投资建议", Pt(12), TEXT_WHITE, bold=True)

    suggestions = [
        ("寒武纪", "国产AI芯片龙头，业绩拐点确认，长期看好。但7000亿市值+极高PE需警惕短期回调。"
         "适合风险承受能力较强的长期投资者。", ACCENT_BLUE),
        ("海光信息", "x86稀缺性+光合组织生态，基本面扎实。PE偏高但确定性较强，"
         "中科曙光协同效应持续释放。适合中长期配置。", ACCENT_GOLD),
        ("景嘉微", "转型期标的，高风险高回报。JM11量产和CH37商业化是关键催化剂。"
         "适合对具身智能赛道有深度认知的投资者。", ACCENT_GREEN),
        ("工业富联", "AI算力基建最确定的受益者，估值最具安全边际。业绩确定性最高，"
         "适合作为AI板块的底仓配置。", RGBColor(0xA7, 0x8B, 0xFA)),
    ]
    for i, (name, desc, color) in enumerate(suggestions):
        y = Inches(4.9) + Inches(i * 0.5)
        add_text_box(slide, Inches(0.9), y, Inches(1.5), Inches(0.4),
                     name, Pt(10), color, bold=True)
        add_text_box(slide, Inches(2.5), y, Inches(10.2), Inches(0.45),
                     desc, Pt(9), TEXT_GRAY, line_spacing=1.25)

    add_page_number(slide, 16)


# ============================================================
# Slide 17: 数据来源与免责声明
# ============================================================
def make_disclaimer():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_top_accent(slide)
    add_footer_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(10), Inches(0.55),
                 "数据来源与免责声明", Pt(24), TEXT_WHITE, bold=True)
    add_line(slide, Inches(0.6), Inches(1.0), Inches(12.5), Inches(1.0), ACCENT_BLUE_DIM, Pt(1))

    # 数据来源
    add_rounded_rect(slide, Inches(0.6), Inches(1.3), Inches(12.3), Inches(2.5),
                     fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text_box(slide, Inches(0.9), Inches(1.4), Inches(5), Inches(0.3),
                 "数据来源", Pt(12), ACCENT_BLUE, bold=True)

    sources = [
        "Gartner - 2027年全球AI芯片市场规模预测",
        "IDC - 2027年中国AI总投资规模预测",
        "弗若斯特沙利文 - 2025-2029年全球端侧AI市场预测",
        "各上市公司2025年年报及2026年一季报（公开披露数据）",
        "北美四大云服务商2026年资本开支指引（公开信息）",
    ]
    for i, src in enumerate(sources):
        add_text_box(slide, Inches(0.9), Inches(1.85) + Inches(i * 0.35), Inches(11.5), Inches(0.3),
                     f"▸ {src}", Pt(10), TEXT_GRAY)

    # 免责声明
    add_rounded_rect(slide, Inches(0.6), Inches(4.1), Inches(12.3), Inches(2.8),
                     fill_color=BG_CARD, border_color=ACCENT_RED)
    add_text_box(slide, Inches(0.9), Inches(4.2), Inches(5), Inches(0.3),
                 "免责声明", Pt(12), ACCENT_RED, bold=True)

    disclaimer_text = (
        "本报告所提供的信息仅供参考，不构成任何投资建议。报告中的数据来源于公开信息，"
        "我们尽力保证数据的准确性和完整性，但不对其准确性和完整性做出任何保证。\n\n"
        "股市有风险，投资需谨慎。过往业绩不代表未来表现。投资者应根据自身风险承受能力、"
        "投资目标和财务状况，独立做出投资决策。\n\n"
        "报告中的任何内容均不应被视为对未来表现的预测或承诺。"
        "本报告的版权归属编制方所有，未经许可不得转载或引用。"
    )
    add_text_box(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0),
                 disclaimer_text, Pt(10), TEXT_GRAY, line_spacing=1.4)

    add_page_number(slide, 17)


# ============================================================
# Slide 18: 封底
# ============================================================
def make_back_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 背景光晕
    add_rect(slide, Inches(3), Inches(0), Inches(7.333), Inches(7.5),
             fill_color=RGBColor(0x07, 0x1A, 0x2E))

    # 主标题
    add_text_box(slide, Inches(2), Inches(2.2), Inches(9.333), Inches(0.8),
                 "THANK YOU", Pt(44), TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name=FONT_EN)
    add_text_box(slide, Inches(2), Inches(3.1), Inches(9.333), Inches(0.5),
                 "感谢阅读", Pt(20), ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # 分隔线
    add_line(slide, Inches(5.5), Inches(3.8), Inches(7.8), Inches(3.8), ACCENT_BLUE_DIM, Pt(1))

    # 信息
    add_text_box(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(0.3),
                 "A股AI芯片行业深度研究报告", Pt(12), TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(4.6), Inches(9.333), Inches(0.3),
                 "Weaver Research  |  2026年5月", Pt(11), TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # 底部装饰
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Pt(3), fill_color=ACCENT_BLUE)


# ============================================================
# 生成PPT
# ============================================================
if __name__ == "__main__":
    print("开始生成PPT...")
    print("  [1/18] 封面")
    make_cover()
    print("  [2/18] 目录")
    make_toc()
    print("  [3/18] 行业概览-市场空间")
    make_market_overview()
    print("  [4/18] 产业链全景图")
    make_industry_chain()
    print("  [5/18] 寒武纪-公司概况")
    make_cambricon_overview()
    print("  [6/18] 寒武纪-核心优势")
    make_cambricon_analysis()
    print("  [7/18] 海光信息-公司概况")
    make_hygon_overview()
    print("  [8/18] 海光信息-核心优势")
    make_hygon_advantages()
    print("  [9/18] 景嘉微-公司概况")
    make_jingjia_overview()
    print("  [10/18] 景嘉微-战略分析")
    make_jingjia_strategy()
    print("  [11/18] 工业富联")
    make_foxconn()
    print("  [12/18] 财务数据对比")
    make_financial_comparison()
    print("  [13/18] 竞争格局")
    make_competition()
    print("  [14/18] 行业趋势")
    make_trends()
    print("  [15/18] 风险提示")
    make_risks()
    print("  [16/18] 投资建议")
    make_investment_score()
    print("  [17/18] 数据来源与免责声明")
    make_disclaimer()
    print("  [18/18] 封底")
    make_back_cover()

    output_path = "/home/admin/.openclaw/workspace/output/AI芯片行业深度研究报告.pptx"
    prs.save(output_path)
    print(f"\nPPT生成完成: {output_path}")
    print(f"总页数: {len(prs.slides)}")
