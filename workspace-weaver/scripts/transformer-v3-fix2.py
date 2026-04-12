#!/usr/bin/env python3
"""
Transformer v3 PPT — 第二轮修复脚本
修复质检报告v2和读者评审v2发现的问题。

修复清单:
  P0-1: 全局字体 <9pt → 提升到 ≥9pt（含footer、Step标签、时间线标签等）
  P0-2: 第19页布局偏空（54.5%覆盖率）→ 添加对比表格
  P1-1: 第20页树形图缺少连接线
  P1-2: 第21页辐射图缺少连接线
  P1-3: 第3页软回车残留（_x000B_, \x0b）
  P1-4: 第11页右侧文字信息密度偏高 → 精简
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

INPUT_FILE = "/home/admin/.openclaw/workspace-weaver/output/Transformer_架构深度解析_v3.pptx"
OUTPUT_FILE = INPUT_FILE  # 覆盖原文件

# 颜色方案
BLUE_DARK = RGBColor(0x1E, 0x3A, 0x5F)
BLUE_MED = RGBColor(0x34, 0x98, 0xDB)
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
GREEN_LIGHT = RGBColor(0xD5, 0xF5, 0xE3)
GREEN_MED = RGBColor(0x2E, 0xCC, 0x71)
ORANGE_MED = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# AML namespace
AML = "http://schemas.openxmlformats.org/drawingml/2006/main"


def fix_small_fonts_all_slides(prs):
    """P0-1: 全局修复所有 <9pt 的字体 → 9pt"""
    count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    # Check p.font.size (explicitly set)
                    if p.font.size and p.font.size < Pt(9):
                        p.font.size = Pt(9)
                        count += 1
                    # Also check run-level font sizes
                    for run in p.runs:
                        if run.font.size and run.font.size < Pt(9):
                            run.font.size = Pt(9)
                            count += 1
                    # Check XML-level sz attributes (for shapes with fs=None in python-pptx)
                    for rPr in p._p.findall(f'.//{{{AML}}}rPr'):
                        sz = rPr.get('sz')
                        if sz and int(sz) < 900:
                            rPr.set('sz', '900')  # 9pt = 900 hundredths of a point
                            count += 1
                    # Also check default text body properties
                    for rPr in p._p.findall(f'.//{{{AML}}}endParaRPr'):
                        sz = rPr.get('sz')
                        if sz and int(sz) < 900:
                            rPr.set('sz', '900')
                            count += 1
    return count


def add_table_slide19(slide):
    """P0-2: 第19页添加对比表格（6维度）"""
    # 根据ppt-plan，第19页应该有6个对比维度的表格
    # 当前只有4个文本段落。在上方添加一个表格。

    # 表格参数
    cols = 3  # 维度 | 训练 | 推理
    rows = 7  # 表头 + 6个维度
    table_left = Inches(0.5)
    table_top = Inches(1.10)
    table_width = Inches(12.3)
    table_height = Inches(2.80)

    # 创建表格shape
    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # 设置列宽
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.9)
    table.columns[2].width = Inches(4.9)

    # 表格数据
    headers = ["Dimension", "Training (Open-book Exam)", "Inference (Closed-book Exam)"]
    data = [
        ["Decoder Input", "Full target sequence (Teacher Forcing)\nModel sees the correct previous words",
         "Model's own generated words\nEach step adds the previous prediction"],
        ["Computation", "All positions in PARALLEL\nGPU utilization is high",
         "Word-by-word SERIAL generation\nMust complete one step before starting next"],
        ["Loss & Gradients", "Cross-entropy loss computed\nBackpropagation updates parameters",
         "No loss, no gradients\nParameters are fixed, forward pass only"],
        ["KV Cache", "Not needed (parallel computation)\nAll K, V computed simultaneously",
         "Required for efficiency\nStores: K, V for all previous tokens"],
        ["Data Flow", "Encoder + Decoder process full sequences\nBidirectional context available",
         "Encoder runs once → cache\nDecoder expands input step by step"],
        ["Goal", "Learn to predict next word accurately\nMinimize loss across training set",
         "Generate useful, coherent responses\nSelect best word at each step"],
    ]

    def set_cell_text(cell, text, font_size=9, bold=False, color=BLACK, fill_color=None, align=PP_ALIGN.LEFT):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.alignment = align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill_color:
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '%02X%02X%02X' % (fill_color[0], fill_color[1], fill_color[2]))

    # 表头行
    for j, header in enumerate(headers):
        set_cell_text(table.cell(0, j), header, font_size=10, bold=True, color=WHITE,
                      fill_color=(0x1E, 0x3A, 0x5F), align=PP_ALIGN.CENTER)

    # 数据行
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            if j == 0:
                fill = (0xD6, 0xEA, 0xF8)  # 浅蓝
                bold = True
            elif j == 1:
                fill = (0xE8, 0xF5, 0xE9)  # 浅绿
                bold = False
            else:
                fill = (0xFF, 0xF8, 0xE1)  # 浅黄
                bold = False
            set_cell_text(table.cell(i + 1, j), cell_text, font_size=9, bold=bold, fill_color=fill)

    # 移动下方现有文字到表格下方
    # 当前文字从 top=4.20 开始，表格到 top=3.90。把文字移到 top=4.05 附近
    text_shapes_to_move = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top > Inches(3.5) and shape.name.startswith('TextBox'):
            text_shapes_to_move.append(shape)

    # 实际上我们把下方4个对比段落删除（已在表格中体现）
    # 但保留footer
    shapes_to_delete = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if shape.top >= Inches(4.0) and shape.top < Inches(7.0):
                if text and not text.startswith('Transformer Architecture Deep Dive'):
                    shapes_to_delete.append(shape)

    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    # 调整表格下方添加一句话总结
    summary = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.05), Inches(12.0), Inches(0.40)
    )
    tf = summary.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Insight: Training is about learning patterns from data; inference is about applying those patterns to generate new content."
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.font.name = "Arial"

    return table_shape


def add_connector_line(slide, x1, y1, x2, y2, color=RGBColor(0x1E, 0x3A, 0x5F), width=Pt(1.5)):
    """添加一条自由线条（带箭头）"""
    # 使用自由线条连接两点
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width

    # 添加箭头尾部
    line = connector._element.find('.//' + qn('a:ln'))
    if line is not None:
        tailEnd = etree.SubElement(line, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')

    return connector


def add_connectors_slide20(slide):
    """P1-1: 第20页三大家族树形图添加连接线"""
    # Transformer 2017 (Oval 11): center = (2.30+1.20, 1.20+0.325) = (3.50, 1.525)
    # Decoder-Only (Rounded Rectangle 12): center-top = (0.30+1.0, 2.30) = (1.30, 2.30)
    # Encoder-Decoder (Rounded Rectangle 15): center-top = (2.60+1.0, 2.30) = (3.60, 2.30)
    # Encoder-Only (Rounded Rectangle 18): center-top = (4.90+1.0, 2.30) = (5.90, 2.30)

    # Transformer oval center
    oval_cx = Inches(2.30) + Inches(1.20)  # = 3.50
    oval_cy = Inches(1.20) + Inches(0.325)  # = 1.525
    oval_bottom = Inches(1.20) + Inches(0.65)  # = 1.85

    # Connect to each branch
    branches = [
        (Inches(1.30), Inches(2.30), BLUE_MED),   # Decoder-Only
        (Inches(3.60), Inches(2.30), GREEN_MED),   # Encoder-Decoder
        (Inches(5.90), Inches(2.30), RGBColor(0x9B, 0x59, 0xB6)),  # Encoder-Only
    ]

    for bx, by, color in branches:
        add_connector_line(slide, oval_cx, oval_bottom, bx, by, color=color, width=Pt(2))


def add_connectors_slide21(slide):
    """P1-2: 第21页辐射图添加连接线"""
    # Transformer center (Oval 11): center = (3.00+0.80, 1.50+0.35) = (3.80, 1.85)
    center_x = Inches(3.00) + Inches(0.80)
    center_y = Inches(1.50) + Inches(0.35)

    # 6 domain nodes - connect from center to each
    domains = [
        (Inches(0.30) + Inches(0.90), Inches(1.00) + Inches(0.325)),  # NLP (1.20, 1.325)
        (Inches(5.50) + Inches(0.90), Inches(1.00) + Inches(0.325)),  # Computer Vision (6.40, 1.325)
        (Inches(7.00) + Inches(0.90), Inches(1.80) + Inches(0.325)),  # Speech (7.90, 2.125)
        (Inches(6.00) + Inches(0.90), Inches(2.70) + Inches(0.325)),  # Multimodal (6.90, 3.025)
        (Inches(0.80) + Inches(0.90), Inches(2.70) + Inches(0.325)),  # Protein (1.70, 3.025)
        (Inches(0.00) + Inches(0.90), Inches(1.80) + Inches(0.325)),  # Time Series (0.90, 2.125)
    ]

    colors = [
        BLUE_MED,
        RGBColor(0xE7, 0x4C, 0x3C),
        RGBColor(0xF3, 0x9C, 0x12),
        RGBColor(0x9B, 0x59, 0xB6),
        RGBColor(0x27, 0xAE, 0x60),
        RGBColor(0x16, 0xA0, 0x85),
    ]

    for (dx, dy), color in zip(domains, colors):
        add_connector_line(slide, center_x, center_y, dx, dy, color=color, width=Pt(1.5))


def fix_soft_returns_slide3(slide):
    """P1-3: 第3页修复软回车残留"""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                original = p.text
                # Replace _x000B_ and \x0b with space
                cleaned = original.replace('_x000B_', ' ').replace('\x0b', ' ')
                if cleaned != original:
                    # Replace in XML runs
                    for run in p.runs:
                        run.text = run.text.replace('_x000B_', ' ').replace('\x0b', ' ')
                    count += 1
    return count


def simplify_slide11_text(slide):
    """P1-4: 第11页精简右侧文字，删除与流程图重复的步骤描述"""
    # 当前的6个Step TextBox (44-53) 重复了左侧流程图的内容
    # 保留标题(Step N)，但精简描述为关键insight

    replacements = {
        'TextBox 44': 'Step 1: Generate Q, K, V',
        'TextBox 45': 'Key insight: Every word gets 3 "search tools" (Q=what I seek, K=what I offer, V=what I know). 8 heads = 8 independent perspectives.',
        'TextBox 46': 'Step 2: Q × Kᵀ = Score Matrix',
        'TextBox 47': 'Key insight: Dot product measures "relevance". Each word compares its Query against all Keys simultaneously → T×T score matrix.',
        'TextBox 48': 'Step 3: Divide by √dₖ',
        'TextBox 49': 'Key insight: Prevents "winner-take-all" saturation. √64=8 calibrates scores so Softmax produces smooth distributions.',
        'TextBox 50': 'Step 4: Softmax = Attention Weights',
        'TextBox 51': 'Key insight: Converts raw scores → probabilities (each row sums to 1). Highest match gets most "attention time".',
        'TextBox 52': 'Step 5: Weights × V = Output',
        'TextBox 53': 'Key insight: Weighted sum of all Values. "it" absorbs "cat"\'s meaning → solves the coreference problem in one pass.',
    }

    # Also adjust the size of the description text boxes
    for shape in slide.shapes:
        if shape.name in replacements:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    new_text = replacements[shape.name]
                    p.text = new_text
                    # Keep bold for Step titles
                    if shape.name.startswith('TextBox 4') and int(shape.name.split()[-1]) % 2 == 0:
                        p.font.bold = True


def self_check(prs):
    """自检：列出所有 <9pt 的 shape"""
    violations = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.font.size and p.font.size < Pt(9):
                        violations.append((i + 1, shape.name, p.text[:40], p.font.size / 12700))
                    for run in p.runs:
                        if run.font.size and run.font.size < Pt(9):
                            violations.append((i + 1, shape.name, run.text[:40], run.font.size / 12700))
                    for rPr in p._p.findall(f'.//{{{AML}}}rPr'):
                        sz = rPr.get('sz')
                        if sz and int(sz) < 900:
                            violations.append((i + 1, shape.name, p.text[:40], int(sz) / 100))
    return violations


def main():
    print("=" * 60)
    print("Transformer v3 PPT — 第二轮修复")
    print("=" * 60)

    prs = Presentation(INPUT_FILE)

    # ========================================
    # P0-1: 全局修复小字体
    # ========================================
    print("\n[P0-1] 修复全局小字体 (<9pt → 9pt)...")
    font_count = fix_small_fonts_all_slides(prs)
    print(f"  → 修复了 {font_count} 处小字体")

    # ========================================
    # P0-2: 第19页丰富内容
    # ========================================
    print("\n[P0-2] 丰富第19页内容（添加6维度对比表格）...")
    slide19 = prs.slides[18]
    add_table_slide19(slide19)
    print("  → 添加了6×3对比表格，删除了重复文字段落")

    # ========================================
    # P1-1: 第20页添加连接线
    # ========================================
    print("\n[P1-1] 第20页树形图添加连接线...")
    slide20 = prs.slides[19]
    add_connectors_slide20(slide20)
    print("  → 添加了3条连接线（Transformer → 3个分支）")

    # ========================================
    # P1-2: 第21页添加连接线
    # ========================================
    print("\n[P1-2] 第21页辐射图添加连接线...")
    slide21 = prs.slides[20]
    add_connectors_slide21(slide21)
    print("  → 添加了6条连接线（Transformer → 6个领域）")

    # ========================================
    # P1-3: 第3页修复软回车
    # ========================================
    print("\n[P1-3] 修复第3页软回车残留...")
    slide3 = prs.slides[2]
    sr_count = fix_soft_returns_slide3(slide3)
    print(f"  → 修复了 {sr_count} 处软回车")

    # ========================================
    # P1-4: 第11页精简文字
    # ========================================
    print("\n[P1-4] 精简第11页右侧文字（减少重复）...")
    slide11 = prs.slides[10]
    simplify_slide11_text(slide11)
    print("  → 6个Step描述改为Key Insight格式，删除重复公式描述")

    # ========================================
    # 自检
    # ========================================
    print("\n" + "=" * 60)
    print("自检：扫描所有 <9pt 的字体...")
    violations = self_check(prs)
    if violations:
        print(f"  ⚠️ 仍有 {len(violations)} 处 <9pt 字体：")
        for slide_num, name, text, size in violations:
            print(f"    Slide {slide_num}, {name}: {text} ({size}pt)")
    else:
        print("  ✅ 所有字体均 ≥ 9pt，无违规！")

    # ========================================
    # 保存
    # ========================================
    print(f"\n保存到: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("✅ 修复完成！")

    # ========================================
    # 修复清单
    # ========================================
    print("\n" + "=" * 60)
    print("修复清单：")
    print("=" * 60)
    print(f"  ✅ P0-1: 全局小字体修复 ({font_count}处)")
    print(f"  ✅ P0-2: 第19页添加6维度对比表格")
    print(f"  ✅ P1-1: 第20页树形图添加3条连接线")
    print(f"  ✅ P1-2: 第21页辐射图添加6条连接线")
    print(f"  ✅ P1-3: 第3页软回车修复 ({sr_count}处)")
    print(f"  ✅ P1-4: 第11页右侧文字精简")
    print(f"  字体自检: {'✅ 全部通过' if not violations else f'⚠️ {len(violations)}处违规'}")


if __name__ == "__main__":
    main()
